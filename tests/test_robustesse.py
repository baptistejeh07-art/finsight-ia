"""
tests/test_robustesse.py — FinSight IA
Audit automatise de robustesse — 10 cas x 6 generateurs = 60 tests.
Aucun appel API. Aucune dependance externe. 100 % offline.

Usage : python tests/test_robustesse.py
"""
from __future__ import annotations

import io
import logging
import os
import statistics
import sys
import tempfile
import time
import types
from pathlib import Path

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

# ─── CONSTANTES ───────────────────────────────────────────────────────────────

GICS_SECTORS = [
    "Technology", "Health Care", "Financials",
    "Consumer Discretionary", "Consumer Staples", "Industrials",
    "Energy", "Materials", "Utilities",
    "Communication Services", "Real Estate",
]

TIMEOUT_S = 30.0          # durée max par test

# Pages / slides attendus (min, max) ou exact
EXPECTED = {
    "sector_pdf":      (6,  25),   # PDF sectoriel
    "sectoral_pptx":   (20, 20),   # PPTX sectoriel — toujours 20 slides
    "indice_pdf":      (5,  25),   # PDF indice
    "screening_excel": (1,  None), # Excel — juste "existe et > taille min"
    "pdf_report":      (4,  30),   # PDF societe individuelle
    "pptx_pitchbook":  (6,  6),    # PPTX pitchbook — toujours 6 slides
}

MIN_SIZE_BYTES = {
    "sector_pdf":      20_000,
    "sectoral_pptx":   50_000,
    "indice_pdf":      15_000,
    "screening_excel": 5_000,
    "pdf_report":      15_000,
    "pptx_pitchbook":  20_000,
}

CASE_LABELS = [
    "01_nominal",
    "02_minimum",
    "03_large",
    "04_none_partout",
    "05_strings_longs",
    "06_valeurs_neg",
    "07_extremes",
    "08_devise_inconnue",
    "09_historique_court",
    "10_sous_ponderer",
]

# ─── FIXTURES COMMUNES ────────────────────────────────────────────────────────

_LONG_COMPANY = "A" * 80          # nom 80 caracteres
_LONG_TEXT    = " ".join(["analyse"] * 500)   # texte ~500 mots
_LONG_SECTOR  = "S" * 40


def _ticker(ticker: str, company: str, sector: str,
            score: float = 60.0, **kw) -> dict:
    """Record ticker complet — tous les champs optionnels presents par defaut."""
    base = {
        "ticker":         ticker,
        "company":        company,
        "sector":         sector,
        "score_global":   score,
        "score_value":    score * 0.25,
        "score_growth":   score * 0.25,
        "score_quality":  score * 0.25,
        "score_momentum": score * 0.25,
        "ev_ebitda":      12.0,
        "ev_revenue":     3.0,
        "pe_ratio":       18.0,
        "ebitda_margin":  25.0,
        "gross_margin":   40.0,
        "net_margin":     12.0,
        "roe":            18.0,
        "revenue_growth": 0.08,
        "momentum_52w":   12.0,
        "altman_z":       3.2,
        "beneish_m":      -2.5,
        "beta":           1.0,
        "price":          100.0,
        "market_cap":     1e10,
        "revenue_ltm":    1e9,
        "currency":       "EUR",
        "sentiment_score": 0.1,
    }
    base.update(kw)
    return base


def _make_tickers(n: int, sector: str, base_score: float = 65.0, **kw) -> list:
    """Genere n tickers dans un meme secteur."""
    result = []
    for i in range(n):
        sc = max(5, min(95, base_score + (i - n // 2) * 2))
        t = _ticker(
            ticker   = f"{sector[:3].upper()}{i+1}",
            company  = f"{sector} Corp {i+1}",
            sector   = sector,
            score    = sc,
            ev_ebitda = max(0.5, 8.0 + i * 1.2),
            ebitda_margin = 20.0 + i * 0.5,
            revenue_growth = 0.05 + i * 0.01,
            momentum_52w   = 5.0 + i * 2.0,
            price          = 80.0 + i * 5,
        )
        t.update(kw)
        result.append(t)
    return result


def _make_tickers_multi(n: int, base_score: float = 65.0, **kw) -> list:
    """Genere n tickers repartis sur les 11 secteurs GICS."""
    result = []
    for i in range(n):
        sec = GICS_SECTORS[i % len(GICS_SECTORS)]
        sc = max(5, min(95, base_score + (i - n // 2) * 2))
        t = _ticker(
            ticker   = f"T{i+1:02d}",
            company  = f"Corp {i+1}",
            sector   = sec,
            score    = sc,
            ev_ebitda = max(0.5, 8.0 + (i % 11) * 1.2),
            ebitda_margin = 20.0 + (i % 11) * 0.5,
            revenue_growth = 0.05 + (i % 11) * 0.01,
            momentum_52w   = 5.0 + (i % 11) * 2.0,
            price          = 80.0 + i * 2,
        )
        t.update(kw)
        result.append(t)
    return result


# ─── BUILDER INDICE DATA ─────────────────────────────────────────────────────

def _make_indice_data(tickers: list, universe: str) -> dict:
    """Construit le dict data pour IndicePDFWriter depuis une liste de tickers."""
    # Grouper par secteur
    sectors: dict = {}
    for t in tickers:
        sec = t.get("sector") or "Autre"
        sectors.setdefault(sec, []).append(t)

    # Construire secteurs_list: tuples (nom, nb, score, signal, ev_str, mg, croi_str, mom_str)
    secteurs_list = []
    for sec_name, items in sectors.items():
        if sec_name == "Autre":
            continue
        nb  = len(items)
        sc  = sum((x.get("score_global") or 0) for x in items) / nb
        evs = [x["ev_ebitda"] for x in items
               if x.get("ev_ebitda") is not None and 1.0 < x["ev_ebitda"] < 100]
        ev_str = f"{statistics.median(evs):.1f}x" if evs else "—"
        mgs = [x.get("ebitda_margin") for x in items if x.get("ebitda_margin") is not None]
        mg  = round(statistics.median(mgs), 1) if mgs else 0.0
        revg = [x.get("revenue_growth") for x in items if x.get("revenue_growth") is not None]
        croi = round(statistics.median(revg) * 100, 1) if revg else 0.0
        croi_str = f"+{croi:.1f}%" if croi >= 0 else f"{croi:.1f}%"
        moms = [x.get("momentum_52w") for x in items if x.get("momentum_52w") is not None]
        mom  = round(statistics.median(moms), 1) if moms else 0.0
        mom_str = f"+{mom:.1f}%" if mom >= 0 else f"{mom:.1f}%"
        sig = ("Surpondérer" if sc >= 60 else ("Sous-pondérer" if sc < 40 else "Neutre"))
        secteurs_list.append((sec_name, nb, int(sc), sig, ev_str, mg, croi_str, mom_str))

    secteurs_list.sort(key=lambda x: x[2], reverse=True)
    if not secteurs_list:
        # fallback: creer un secteur minimal
        secteurs_list = [("Autre", len(tickers), 50, "Neutre", "—", 0.0, "+0.0%", "+0.0%")]

    all_scores = [x.get("score_global") or 0 for x in tickers]
    avg_score  = sum(all_scores) / len(all_scores) if all_scores else 50
    signal_global = ("Surpondérer" if avg_score >= 60
                     else ("Sous-pondérer" if avg_score < 40 else "Neutre"))
    conviction = min(95, max(50, int(avg_score)))

    # Top3 secteurs
    _SIG_S = "Surpondérer"
    surp = [s for s in secteurs_list if s[3] == _SIG_S]
    other = [s for s in secteurs_list if s[3] != _SIG_S]
    top3_secs = (surp[:3] + other[:max(0, 3 - len(surp))])[:3]
    if not top3_secs:
        top3_secs = secteurs_list[:1]

    def _build_top3_entry(s):
        sec_items = sectors.get(s[0], tickers[:3])
        top_tkrs = sorted(sec_items, key=lambda x: x.get("score_global") or 0, reverse=True)[:3]
        socs = []
        for t in top_tkrs:
            er = t.get("ev_ebitda")
            ev_t = f"{er:.1f}x" if (er is not None and 1.0 < er < 100) else "—"
            sig_t = ("Surpondérer" if (t.get("score_global") or 0) >= 60
                     else ("Sous-pondérer" if (t.get("score_global") or 0) < 40 else "Neutre"))
            socs.append((t.get("ticker", "?"), sig_t, ev_t, t.get("score_global") or 0))
        if not socs:
            socs = [("—", "Neutre", "—", 0)]
        return {
            "nom": s[0], "signal": s[3], "score": s[2], "ev_ebitda": s[4],
            "catalyseur": f"Score {s[2]}/100 · momentum {s[7]} · EV/EBITDA {s[4]}",
            "risque":     f"Croiss. revenus {s[6]} — compression multiple si cycle se retourne",
            "societes":   socs,
        }

    top3 = [_build_top3_entry(s) for s in top3_secs]
    rotation = [(s[0], "Expansion", "Moderee", "Moderee", s[3]) for s in secteurs_list]

    top_noms = " / ".join(s[0] for s in secteurs_list[:3])
    nb_sec   = len(secteurs_list)

    return {
        "indice":         universe,
        "code":           universe[:6].upper().replace(" ", ""),
        "nb_secteurs":    nb_sec,
        "nb_societes":    len(tickers),
        "signal_global":  signal_global,
        "conviction_pct": conviction,
        "date_analyse":   "25 mars 2026",
        "cours":          "—",
        "variation_ytd":  "—",
        "pe_forward":     "—",
        "score_global":   int(avg_score),
        "secteurs":       secteurs_list,
        "texte_macro":    f"Signal {signal_global} (conviction {conviction}%) — {len(tickers)} societes sur {nb_sec} secteurs.",
        "texte_signal":   f"Signal global {signal_global}. {nb_sec} secteurs analyses.",
        "texte_rotation": f"Signal {signal_global} — selectivite recommandee. Top : {top_noms}.",
        "surp_noms":      " / ".join(s[0] for s in secteurs_list if s[3] == "Surpondérer") or "aucun",
        "sous_noms":      " / ".join(s[0] for s in secteurs_list if s[3] == "Sous-pondérer") or "aucun",
        "top3_secteurs":  top3,
        "rotation":       rotation,
        "catalyseurs": [
            ("Momentum fondamental", f"Score moyen {int(avg_score)}/100 — leaders solides", "3-6 mois"),
            ("Consolidation M&A",    "Fusions potentielles — prime de valorisation",        "6-12 mois"),
            ("Innovation",           "IA/automatisation — reduction couts structurelle",    "12-18 mois"),
        ],
        "risques": [
            ("Macro",  "Recession potentielle — choc demande",  "50", "Eleve"),
            ("Taux",   "Hausse taux directeurs BCE/Fed",         "35", "Modere"),
            ("Geo",    "Tensions geopolitiques — supply chain",  "25", "Modere"),
        ],
        "scenarios": [
            ("Bull case",   "Score > 65 + momentum positif",  "Surpondérer",  "3-6 mois"),
            ("Bear case",   "Score < 40 via recession",       "Sous-pondérer", "6-12 mois"),
            ("Stagflation", "CPI > 3,5 % + PIB < 1 %",       "Sous-pondérer", "6-9 mois"),
        ],
        "finbert": {
            "nb_articles":  0,
            "score_agrege": 0.0,
            "positif": {"nb": 0, "score": "N/A", "themes": "Donnees non disponibles"},
            "neutre":  {"nb": 0, "score": "N/A", "themes": "Donnees non disponibles"},
            "negatif": {"nb": 0, "score": "N/A", "themes": "Donnees non disponibles"},
            "par_secteur": [(s[0], "N/A", "Neutre") for s in secteurs_list],
        },
        "methodologie": [
            ("Score sectoriel", "Composite 0-100 : 40 % momentum, 30 % BPA, 30 % valorisation"),
            ("Signal",          "Surponderer (>60) / Neutre (40-60) / Sous-ponderer (<40)"),
            ("Valorisation",    "EV/EBITDA median LTM — source FMP / yfinance"),
        ],
        "perf_history": None,
    }


# ─── BUILDER PDF REPORT DATA ──────────────────────────────────────────────────

def _make_pdf_data(ticker: str, company: str, sector: str,
                   recommendation: str = "BUY",
                   currency: str = "EUR",
                   **kw) -> dict:
    """Construit un dict data minimal valide pour pdf_writer.generate_report()."""
    base = {
        "company_name":    company,
        "ticker":          ticker,
        "ticker_exchange": ticker,
        "sector":          sector,
        "currency":        currency,
        "recommendation":  recommendation,
        "price_str":       f"100,00 {currency}",
        "target_price_full": f"125,00 {currency}",
        "upside_str":      "+25 %",
        "market_cap_str":  "10 Md EUR",
        "conviction_str":  "85 %",
        "dividend_yield_str": "2,1 %",
        "pe_ntm_str":      "18,0x",
        "ev_ebitda_str":   "12,0x",
        "date_analyse":    "25 mars 2026",
        # Performance chart
        "perf_months":  ["Jan","Fev","Mar","Avr","Mai","Jun","Jul","Aou","Sep","Oct","Nov","Dec","Jan'25"],
        "perf_ticker":  [100, 102, 105, 108, 110, 107, 112, 115, 118, 120, 122, 125, 128],
        "perf_index":   [100, 101, 103, 106, 108, 105, 110, 113, 115, 118, 120, 123, 126],
        "perf_start_label": "janv 2024",
        "index_name":   "CAC 40",
        # Football field
        "ff_methods":   ["DCF", "Comparables"],
        "ff_lows":      [90.0, 95.0],
        "ff_highs":     [130.0, 135.0],
        "ff_colors":    ["#1B3A6B", "#2A5298"],
        "ff_course":    100.0,
        "ff_course_str": f"100 {currency}",
        # Pie comparables
        "pie_labels":   ["Peer A", "Peer B", "Peer C"],
        "pie_sizes":    [30.0, 25.0, 45.0],
        "pie_ticker":   ticker,
        "pie_pct_str":  "30 %",
        "pie_sector_name": sector,
        "pie_cap_label": "EV",
        # Revenue area
        "area_quarters":   ["T1 24", "T2 24", "T3 24", "T4 24"],
        "area_segments":   {"Segment A": [250, 260, 270, 280],
                            "Segment B": [200, 210, 220, 230]},
        "area_year_labels": ["2024"],
        "area_annual_fallback": [("2022", 900), ("2023", 1000), ("2024", 1100)],
        # IS tableau
        "is_col_headers": ["2021", "2022", "2023", "2024", "LTM", "2025F"],
        "is_data": [
            ["Chiffre d'affaires", "900", "1000", "1100", "1150", "1200", "1300"],
            ["EBITDA",             "225", "250",  "275",  "288",  "300",  "325"],
            ["Marge EBITDA",       "25%", "25%",  "25%",  "25%",  "25%",  "25%"],
            ["EBIT",               "180", "200",  "220",  "230",  "240",  "260"],
            ["Resultat net",       "108", "120",  "132",  "138",  "144",  "156"],
            ["BPA",                "1,08","1,20", "1,32", "1,38", "1,44", "1,56"],
            ["FCF",                "90",  "100",  "110",  "115",  "120",  "130"],
            ["Dette nette",        "400", "380",  "360",  "340",  "320",  "290"],
        ],
        # Scenarios — dicts avec scenario/price/upside/prob/hypothesis
        "scenarios": [
            {"scenario": "Bull", "price": "150 EUR", "upside": "+50 %", "prob": "30 %", "hypothesis": "Croissance acceleree"},
            {"scenario": "Base", "price": "125 EUR", "upside": "+25 %", "prob": "50 %", "hypothesis": "Scenario central"},
            {"scenario": "Bear", "price":  "85 EUR", "upside": "-15 %", "prob": "20 %", "hypothesis": "Recession"},
        ],
        # bear_args — dicts avec name/text
        "bear_args": [
            {"name": "Risque macro",       "text": "Recession possible en 2026"},
            {"name": "Pression marges",    "text": "Hausse des couts de production"},
            {"name": "Concurrence",        "text": "Nouveaux entrants agressifs"},
        ],
        # catalysts — dicts avec num/name/analysis
        "catalysts": [
            {"num": "1", "name": "Nouveaux contrats",         "analysis": "Signature attendue Q2 2026"},
            {"num": "2", "name": "Expansion internationale",  "analysis": "Ouverture marches Asie"},
            {"num": "3", "name": "Hausse marges",             "analysis": "Economies d echelle"},
        ],
        # invalidation_data — dicts avec axe/condition/horizon
        "invalidation_data": [
            {"axe": "BPA",     "condition": "BPA < -10 % YoY sur 2 trimestres", "horizon": "6 mois"},
            {"axe": "Contrat", "condition": "Perte contrat cle > 15 % CA",      "horizon": "12 mois"},
        ],
        # sentiment_data — dicts avec orientation/articles/score/themes
        "sentiment_data": [
            {"orientation": "Positif", "articles": "62", "score": "0,42", "themes": "Resultats, dividende"},
            {"orientation": "Neutre",  "articles": "25", "score": "0,05", "themes": "Macro, taux"},
            {"orientation": "Negatif", "articles": "13", "score": "-0,38", "themes": "Concurrence, dette"},
        ],
        # revision_data — dicts avec revision/trigger/target/style
        "revision_data": [
            {"revision": "Hausse",  "trigger": "BPA NTM > +5 % sur 2 trimestres", "target": "140 EUR", "style": "buy"},
            {"revision": "Baisse",  "trigger": "BPA NTM < -10 % ou perte contrat", "target": "80 EUR",  "style": "sell"},
        ],
        # DCF sensitivity — 5 wacc x 5 tgr (doit correspondre aux valeurs par defaut)
        "wacc_rows":  ["8,4%", "9,4%", "10,4%", "11,4%", "12,4%"],
        "tgr_cols":   ["2,0%", "2,5%", "3,0%", "3,5%", "4,0%"],
        "wacc_base":  "10,4%",
        "tgr_base":   "3,0%",
        "dcf_sensitivity": [
            ["145", "150", "156", "162", "170"],
            ["132", "137", "142", "148", "155"],
            ["120", "125", "130", "136", "143"],
            ["110", "114", "119", "124", "130"],
            ["101", "105", "109", "114", "120"],
        ],
        # ratios_vs_peers — dicts avec label/value/reference/lecture
        "ratios_vs_peers": [
            {"label": "EV/EBITDA", "value": "12,0x", "reference": "14,0x secteur", "lecture": "Decote"},
            {"label": "P/E NTM",   "value": "18,0x", "reference": "20,0x secteur", "lecture": "En ligne"},
        ],
    }
    base.update(kw)
    return base


# ─── BUILDER PPTX OBJECTS ────────────────────────────────────────────────────

def _make_pptx_objs(ticker: str, company: str, sector: str,
                    recommendation: str = "BUY",
                    currency: str = "EUR",
                    score: float = 65.0, **kw):
    """Construit snapshot/ratios/synthesis/qa/devil comme SimpleNamespace."""
    yr = types.SimpleNamespace(
        gross_margin   = kw.get("gross_margin",   40.0),
        ebitda_margin  = kw.get("ebitda_margin",  25.0),
        ebit_margin    = kw.get("ebit_margin",    18.0),
        net_margin     = kw.get("net_margin",     12.0),
        roe            = kw.get("roe",            18.0),
        roic           = kw.get("roic",           15.0),
        net_debt_ebitda = kw.get("net_debt_ebitda", 2.0),
        ev_ebitda      = kw.get("ev_ebitda",      12.0),
        pe_ratio       = kw.get("pe_ratio",       18.0),
        fcf_yield      = kw.get("fcf_yield",       5.0),
        current_ratio  = kw.get("current_ratio",   1.8),
        altman_z       = kw.get("altman_z",        3.2),
        beneish_m      = kw.get("beneish_m",      -2.5),
        revenue_growth = kw.get("revenue_growth",  0.08),
    )
    ci = types.SimpleNamespace(
        company_name=company, ticker=ticker,
        sector=sector, currency=currency,
    )
    mkt      = types.SimpleNamespace(share_price=kw.get("price", 100.0))
    snapshot = types.SimpleNamespace(ticker=ticker, company_info=ci, market=mkt)
    ratios   = types.SimpleNamespace(latest_year="2024", years={"2024": yr})

    strengths = kw.get("strengths", ["Marges elevees", "Leader de marche", "Revenus recurrents"])
    risks     = kw.get("risks",     ["Concurrence", "Risque macro", "Reglementation"])

    synthesis = types.SimpleNamespace(
        recommendation      = recommendation,
        conviction          = kw.get("conviction", 0.80),
        confidence_score    = kw.get("confidence_score", 0.85),
        target_bear         = kw.get("target_bear",  85.0),
        target_base         = kw.get("target_base", 125.0),
        target_bull         = kw.get("target_bull", 150.0),
        summary             = kw.get("summary", "Thesis summary."),
        valuation_comment   = kw.get("valuation_comment", "Fairly valued."),
        invalidation_conditions = kw.get("invalidation_conditions", "If revenue misses."),
        strengths           = strengths,
        risks               = risks,
    )

    qa_flag = types.SimpleNamespace(level="INFO", message="Data complete")
    qa = types.SimpleNamespace(
        qa_score = kw.get("qa_score", 0.85),
        passed   = kw.get("qa_passed", True),
        flags    = kw.get("qa_flags",  [qa_flag]),
    )
    devil = types.SimpleNamespace(
        conviction_delta = kw.get("conviction_delta", 0.1),
        original_reco    = recommendation,
        counter_reco     = "HOLD",
        counter_thesis   = kw.get("counter_thesis", "Marche surévalue."),
        key_assumptions  = kw.get("key_assumptions", ["P/E stable", "Croissance soutenue"]),
    )
    return snapshot, ratios, synthesis, qa, devil


# ─── DEFINITION DES 10 CAS ───────────────────────────────────────────────────

def _build_cases() -> list[dict]:
    """Retourne la liste des 10 cas de test sous forme de dicts."""
    cases = []

    # ── Cas 1 : nominal ────────────────────────────────────────────────────
    td1 = _make_tickers(6, "Technology", base_score=68.0)
    cases.append({
        "label":      CASE_LABELS[0],
        "tickers":    td1,
        "sector":     "Technology",
        "universe":   "CAC 40",
        "indice":     _make_indice_data(td1, "CAC 40"),
        "pdf_data":   _make_pdf_data("SAP.DE", "SAP SE", "Technology"),
        "pptx_objs":  _make_pptx_objs("SAP.DE", "SAP SE", "Technology"),
    })

    # ── Cas 2 : minimum (1 société) ────────────────────────────────────────
    td2 = _make_tickers(1, "Energy", base_score=55.0)
    cases.append({
        "label":      CASE_LABELS[1],
        "tickers":    td2,
        "sector":     "Energy",
        "universe":   "CAC 40",
        "indice":     _make_indice_data(td2, "CAC 40"),
        "pdf_data":   _make_pdf_data("TTE.PA", "TotalEnergies", "Energy"),
        "pptx_objs":  _make_pptx_objs("TTE.PA", "TotalEnergies", "Energy", recommendation="HOLD"),
    })

    # ── Cas 3 : univers large (30 sociétés, 11 secteurs) ───────────────────
    td3 = _make_tickers_multi(30, base_score=60.0)
    cases.append({
        "label":      CASE_LABELS[2],
        "tickers":    td3,
        "sector":     "Technology",
        "universe":   "S&P 500",
        "indice":     _make_indice_data(td3, "S&P 500"),
        "pdf_data":   _make_pdf_data("AAPL", "Apple Inc.", "Technology"),
        "pptx_objs":  _make_pptx_objs("AAPL", "Apple Inc.", "Technology"),
    })

    # ── Cas 4 : None partout (champs optionnels) ───────────────────────────
    td4 = [
        _ticker(f"N{i+1}", f"Corp {i+1}", GICS_SECTORS[i % 11], score=55.0,
                ev_ebitda=None, ev_revenue=None, pe_ratio=None,
                ebitda_margin=None, gross_margin=None, net_margin=None,
                roe=None, revenue_growth=None, momentum_52w=None,
                altman_z=None, beneish_m=None, beta=None,
                market_cap=None, revenue_ltm=None, sentiment_score=None)
        for i in range(6)
    ]
    cases.append({
        "label":      CASE_LABELS[3],
        "tickers":    td4,
        "sector":     "Technology",
        "universe":   "CAC 40",
        "indice":     _make_indice_data(td4, "CAC 40"),
        "pdf_data":   _make_pdf_data("NONE1", "Corp A", "Technology",
                                     ff_lows=[], ff_highs=[], ff_methods=[],
                                     perf_ticker=[], perf_index=[], perf_months=[],
                                     pie_sizes=[], pie_labels=[],
                                     area_segments={}, is_data=[]),
        "pptx_objs":  _make_pptx_objs("NONE1", "Corp A", "Technology",
                                       ev_ebitda=None, altman_z=None,
                                       strengths=None, risks=None,
                                       key_assumptions=None, qa_flags=[]),
    })

    # ── Cas 5 : strings très longs ─────────────────────────────────────────
    td5 = _make_tickers(6, "Technology", base_score=70.0)
    for t in td5:
        t["company"] = _LONG_COMPANY
        t["sector"]  = _LONG_SECTOR
    cases.append({
        "label":      CASE_LABELS[4],
        "tickers":    td5,
        "sector":     "Technology",
        "universe":   "CAC 40",
        "indice":     _make_indice_data(td5, "CAC 40"),
        "pdf_data":   _make_pdf_data("LONG1", _LONG_COMPANY, "Technology",
                                     summary=_LONG_TEXT, valuation_comment=_LONG_TEXT),
        "pptx_objs":  _make_pptx_objs("LONG1", _LONG_COMPANY[:40], "Technology",
                                       summary=_LONG_TEXT,
                                       strengths=[_LONG_TEXT[:120]] * 5,
                                       risks=[_LONG_TEXT[:120]] * 5,
                                       key_assumptions=[_LONG_TEXT[:120]] * 4),
    })

    # ── Cas 6 : valeurs négatives ──────────────────────────────────────────
    td6 = _make_tickers(6, "Energy", base_score=0.0,
                        ebitda_margin=-10.0, gross_margin=-5.0,
                        net_margin=-8.0, revenue_growth=-0.50,
                        momentum_52w=-40.0, roe=-15.0)
    cases.append({
        "label":      CASE_LABELS[5],
        "tickers":    td6,
        "sector":     "Energy",
        "universe":   "CAC 40",
        "indice":     _make_indice_data(td6, "CAC 40"),
        "pdf_data":   _make_pdf_data("NEG1", "Corp Neg", "Energy",
                                     recommendation="SELL",
                                     ff_lows=[-10.0, -5.0], ff_highs=[50.0, 40.0]),
        "pptx_objs":  _make_pptx_objs("NEG1", "Corp Neg", "Energy",
                                       recommendation="SELL",
                                       ebitda_margin=-10.0, net_margin=-8.0,
                                       roe=-15.0, revenue_growth=-0.50,
                                       altman_z=1.2, conviction_delta=-0.5),
    })

    # ── Cas 7 : valeurs extrêmes ───────────────────────────────────────────
    td7 = _make_tickers(6, "Technology", base_score=90.0,
                        ev_ebitda=200.0, revenue_growth=5.0,
                        momentum_52w=-99.0, pe_ratio=500.0)
    cases.append({
        "label":      CASE_LABELS[6],
        "tickers":    td7,
        "sector":     "Technology",
        "universe":   "CAC 40",
        "indice":     _make_indice_data(td7, "CAC 40"),
        "pdf_data":   _make_pdf_data("EXT1", "Corp Extreme", "Technology",
                                     ff_lows=[0.01, 0.01], ff_highs=[9999.0, 8888.0]),
        "pptx_objs":  _make_pptx_objs("EXT1", "Corp Extreme", "Technology",
                                       ev_ebitda=200.0, pe_ratio=500.0,
                                       revenue_growth=5.0, altman_z=0.1),
    })

    # ── Cas 8 : devise et ticker non standard ──────────────────────────────
    td8 = _make_tickers(6, "Financials", base_score=62.0, currency="XBT")
    for i, t in enumerate(td8):
        t["ticker"] = f"UNKN{i+1}.XX"
    cases.append({
        "label":      CASE_LABELS[7],
        "tickers":    td8,
        "sector":     "Financials",
        "universe":   "TSX",
        "indice":     _make_indice_data(td8, "TSX"),
        "pdf_data":   _make_pdf_data("UNKN.XX", "Unknown Corp", "Financials", currency="XBT"),
        "pptx_objs":  _make_pptx_objs("UNKN.XX", "Unknown Corp", "Financials", currency="XBT"),
    })

    # ── Cas 9 : données historiques courtes (4 semaines) ───────────────────
    td9 = _make_tickers(6, "Health Care", base_score=65.0)
    cases.append({
        "label":      CASE_LABELS[8],
        "tickers":    td9,
        "sector":     "Health Care",
        "universe":   "CAC 40",
        "indice":     _make_indice_data(td9, "CAC 40"),
        # Fournir 4 points au lieu de 13 pour les charts
        "pdf_data":   _make_pdf_data("HLT1", "Health Corp", "Health Care",
                                     perf_months=["W1", "W2", "W3", "W4"],
                                     perf_ticker=[100, 102, 99, 104],
                                     perf_index= [100, 101, 100, 103],
                                     area_quarters=["W1", "W2", "W3", "W4"],
                                     area_segments={"Segment A": [250, 255, 260, 265]},
                                     area_annual_fallback=[("2024", 1000)]),
        "pptx_objs":  _make_pptx_objs("HLT1", "Health Corp", "Health Care"),
    })

    # ── Cas 10 : signal Sous-pondérer (tous scores < 40) ───────────────────
    td10 = _make_tickers(6, "Utilities", base_score=20.0)
    cases.append({
        "label":      CASE_LABELS[9],
        "tickers":    td10,
        "sector":     "Utilities",
        "universe":   "CAC 40",
        "indice":     _make_indice_data(td10, "CAC 40"),
        "pdf_data":   _make_pdf_data("UTL1", "Utilities Corp", "Utilities",
                                     recommendation="SELL"),
        "pptx_objs":  _make_pptx_objs("UTL1", "Utilities Corp", "Utilities",
                                       score=20.0, recommendation="SELL",
                                       conviction_delta=-0.6),
    })

    return cases


# ─── HELPERS D'INSPECTION ────────────────────────────────────────────────────

def _count_pdf_pages(path: str) -> int:
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            return len(pdf.pages)
    except Exception:
        with open(path, "rb") as f:
            content = f.read()
        c1 = content.count(b"/Type /Page")
        c2 = content.count(b"/Type/Page")
        return max(c1, c2) or -1


def _count_pptx_slides(path: str) -> int:
    from pptx import Presentation
    return len(Presentation(path).slides)


def _has_none_text_pdf(path: str) -> list[str]:
    try:
        import pdfplumber
        hits = []
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages):
                txt = page.extract_text() or ""
                for line in txt.split("\n"):
                    if "None" in line:
                        hits.append(f"p{i+1}: {line[:80]}")
        return hits
    except Exception:
        return []


def _has_none_text_pptx(path: str) -> list[str]:
    from pptx import Presentation
    hits = []
    prs = Presentation(path)
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if "None" in para.text:
                        hits.append(f"slide {si+1}: {para.text[:80]}")
    return hits


def _capture_log_handler() -> logging.Handler:
    """Handler qui stocke les records en memoire."""
    class _MemH(logging.Handler):
        def __init__(self):
            super().__init__()
            self.records = []
        def emit(self, record):
            self.records.append(record)
    h = _MemH()
    h.setLevel(logging.DEBUG)
    return h


# ─── RUNNER PAR GENERATEUR ───────────────────────────────────────────────────

def _run_sector_pdf(case: dict, tmpdir: str) -> tuple[bool, str]:
    from outputs.sector_pdf_writer import generate_sector_report
    path = os.path.join(tmpdir, "out.pdf")
    generate_sector_report(case["sector"], case["tickers"], path,
                           universe=case["universe"])
    assert os.path.exists(path), "Fichier non genere"
    assert os.path.getsize(path) >= MIN_SIZE_BYTES["sector_pdf"], \
        f"Fichier trop petit ({os.path.getsize(path)} bytes)"
    n = _count_pdf_pages(path)
    lo, hi = EXPECTED["sector_pdf"]
    assert lo <= n <= (hi or 9999), f"{n} pages (attendu {lo}–{hi})"
    hits = _has_none_text_pdf(path)
    assert not hits, f"'None' visible : {hits[:3]}"


def _run_sectoral_pptx(case: dict, tmpdir: str) -> None:
    from outputs.sectoral_pptx_writer import SectoralPPTXWriter
    path = os.path.join(tmpdir, "out.pptx")
    SectoralPPTXWriter.generate(case["tickers"], case["sector"],
                                 case["universe"], path)
    assert os.path.exists(path), "Fichier non genere"
    assert os.path.getsize(path) >= MIN_SIZE_BYTES["sectoral_pptx"], \
        f"Fichier trop petit ({os.path.getsize(path)} bytes)"
    n = _count_pptx_slides(path)
    lo, hi = EXPECTED["sectoral_pptx"]
    assert lo <= n <= (hi or 9999), f"{n} slides (attendu {lo}–{hi})"
    hits = _has_none_text_pptx(path)
    assert not hits, f"'None' visible : {hits[:3]}"


def _run_indice_pdf(case: dict, tmpdir: str) -> None:
    from outputs.indice_pdf_writer import IndicePDFWriter
    path = os.path.join(tmpdir, "out.pdf")
    IndicePDFWriter.generate(case["indice"], path)
    assert os.path.exists(path), "Fichier non genere"
    assert os.path.getsize(path) >= MIN_SIZE_BYTES["indice_pdf"], \
        f"Fichier trop petit ({os.path.getsize(path)} bytes)"
    n = _count_pdf_pages(path)
    lo, hi = EXPECTED["indice_pdf"]
    assert lo <= n <= (hi or 9999), f"{n} pages (attendu {lo}–{hi})"
    hits = _has_none_text_pdf(path)
    assert not hits, f"'None' visible : {hits[:3]}"


def _run_screening_excel(case: dict, tmpdir: str) -> None:
    from outputs.screening_writer import ScreeningWriter
    path = os.path.join(tmpdir, "out.xlsx")
    ScreeningWriter.generate(case["tickers"], case["universe"], path,
                              template_path=None)
    assert os.path.exists(path), "Fichier non genere"
    assert os.path.getsize(path) >= MIN_SIZE_BYTES["screening_excel"], \
        f"Fichier trop petit ({os.path.getsize(path)} bytes)"
    # Pas de check "None" pour Excel — l'openpyxl gere les None en cellules vides


def _run_pdf_report(case: dict, tmpdir: str) -> None:
    from outputs.pdf_writer import generate_report
    path = os.path.join(tmpdir, "out.pdf")
    generate_report(case["pdf_data"], path)
    assert os.path.exists(path), "Fichier non genere"
    assert os.path.getsize(path) >= MIN_SIZE_BYTES["pdf_report"], \
        f"Fichier trop petit ({os.path.getsize(path)} bytes)"
    n = _count_pdf_pages(path)
    lo, hi = EXPECTED["pdf_report"]
    assert lo <= n <= (hi or 9999), f"{n} pages (attendu {lo}–{hi})"
    hits = _has_none_text_pdf(path)
    assert not hits, f"'None' visible : {hits[:3]}"


def _run_pptx_pitchbook(case: dict, tmpdir: str) -> None:
    from outputs.pptx_builder import PPTXBuilder
    path = Path(tmpdir) / "out.pptx"
    snapshot, ratios, synthesis, qa, devil = case["pptx_objs"]
    PPTXBuilder().build(snapshot, ratios, synthesis, qa, devil, output_path=path)
    assert path.exists(), "Fichier non genere"
    assert path.stat().st_size >= MIN_SIZE_BYTES["pptx_pitchbook"], \
        f"Fichier trop petit ({path.stat().st_size} bytes)"
    n = _count_pptx_slides(str(path))
    lo, hi = EXPECTED["pptx_pitchbook"]
    assert lo <= n <= (hi or 9999), f"{n} slides (attendu {lo}–{hi})"
    hits = _has_none_text_pptx(str(path))
    assert not hits, f"'None' visible : {hits[:3]}"


RUNNERS = {
    "sector_pdf":      _run_sector_pdf,
    "sectoral_pptx":   _run_sectoral_pptx,
    "indice_pdf":      _run_indice_pdf,
    "screening_excel": _run_screening_excel,
    "pdf_report":      _run_pdf_report,
    "pptx_pitchbook":  _run_pptx_pitchbook,
}

GEN_ORDER = list(RUNNERS.keys())


# ─── MOTEUR D'EXECUTION ───────────────────────────────────────────────────────

class TestResult:
    def __init__(self):
        self.ok      = False
        self.error   = ""
        self.elapsed = 0.0
        self.errors  = []    # log ERROR records

    def __str__(self):
        if self.ok:
            return f"OK ({self.elapsed:.1f}s)"
        return f"FAIL {self.error[:60]}"


def _run_one(gen_key: str, case: dict) -> TestResult:
    res  = TestResult()
    root = logging.getLogger()
    handler = _capture_log_handler()
    root.addHandler(handler)

    with tempfile.TemporaryDirectory() as tmpdir:
        t0 = time.perf_counter()
        try:
            RUNNERS[gen_key](case, tmpdir)
            res.elapsed = time.perf_counter() - t0
            res.ok      = res.elapsed <= TIMEOUT_S
            if not res.ok:
                res.error = f"Timeout ({res.elapsed:.1f}s > {TIMEOUT_S}s)"
        except Exception as exc:
            res.elapsed = time.perf_counter() - t0
            res.ok      = False
            # Inclure le type et le message court
            import traceback
            tb_lines = traceback.format_exc().splitlines()
            # Garder la derniere ligne (exception) + la ligne du fichier concernee
            relevant = [l for l in tb_lines if "outputs/" in l or "tests/" in l]
            loc = relevant[-1].strip() if relevant else ""
            res.error = f"{type(exc).__name__}: {exc}  |  {loc}"
        finally:
            # Ne garder que les ERROR de nos propres modules (pas yfinance, etc.)
            res.errors = [
                r for r in handler.records
                if r.levelno >= logging.ERROR
                and (r.name or "").startswith("outputs.")
            ]
            root.removeHandler(handler)

    return res


# ─── AFFICHAGE ────────────────────────────────────────────────────────────────

def _print_table(results: dict[str, dict[str, TestResult]]) -> None:
    col_w   = 9   # largeur colonne cas
    gen_w   = 22  # largeur colonne generateur

    # En-tete
    header = f"{'Output':<{gen_w}}"
    for i in range(10):
        header += f" {'C'+str(i+1):^{col_w}}"
    header += f"  {'Score':>8}"
    print("\n" + "=" * len(header))
    print(header)
    print("=" * len(header))

    # Lignes
    failures = []
    total_ok = total_all = 0

    for gen in GEN_ORDER:
        row = f"{gen:<{gen_w}}"
        ok_count = 0
        for ci, case_label in enumerate(CASE_LABELS):
            res = results[gen][case_label]
            total_all += 1
            if res.ok:
                ok_count += 1
                total_ok  += 1
                cell = f"{'OK':^{col_w}}"
            else:
                cell = f"{'FAIL':^{col_w}}"
                failures.append((gen, case_label, res))
            row += f" {cell}"
        row += f"  {ok_count:>3}/10"
        print(row)

    print("=" * len(header))
    print(f"TOTAL : {total_ok}/{total_all} tests passes\n")

    if failures:
        print("--- DETAILS DES ECHECS -----------------------------------------------")
        for gen, case_label, res in failures:
            print(f"  FAIL [{gen}] [{case_label}]")
            print(f"    {res.error}")
        print()


# ─── MAIN ─────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    logging.basicConfig(level=logging.WARNING)

    print("Construction des fixtures...")
    cases = _build_cases()
    print(f"  {len(cases)} cas, {len(RUNNERS)} generateurs -> {len(cases)*len(RUNNERS)} tests\n")

    results: dict[str, dict[str, TestResult]] = {g: {} for g in GEN_ORDER}

    total = len(cases) * len(GEN_ORDER)
    done  = 0

    for case in cases:
        label = case["label"]
        for gen in GEN_ORDER:
            done += 1
            sys.stdout.write(f"\r  [{done:2d}/{total}] {gen:22s} | {label}   ")
            sys.stdout.flush()
            results[gen][label] = _run_one(gen, case)

    sys.stdout.write("\n\n")
    _print_table(results)

    any_fail = any(
        not results[g][c].ok
        for g in GEN_ORDER for c in CASE_LABELS
    )
    sys.exit(1 if any_fail else 0)
