"""
tests/test_outputs.py — FinSight IA
Tests de régression sur les générateurs d'outputs.

3 cas couverts :
  Cas 1 — Nominal   : Technology CAC40, 6 sociétés, signal Surpondérer
  Cas 2 — Limite bas: Energy CAC40,     1 société,  signal Neutre
  Cas 3 — Univers large: Technology S&P500, 30 sociétés, signal Sous-pondérer

Lancer : python tests/test_outputs.py
"""
from __future__ import annotations

import sys, os, tempfile, logging
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

logging.basicConfig(level=logging.WARNING)

# ─── FIXTURES ─────────────────────────────────────────────────────────────────

def _ticker_record(ticker: str, company: str, sector: str,
                   score: float = 60.0, ev: float = 12.0,
                   mg: float = 25.0, rev_g: float = 0.08,
                   mom: float = 12.0, price: float = 100.0) -> dict:
    """Minimal tickers_data record covering all fields used by writers."""
    return {
        "ticker":         ticker,
        "company":        company,
        "sector":         sector,
        "score_global":   score,
        "score_value":    score * 0.25,
        "score_growth":   score * 0.25,
        "score_quality":  score * 0.25,
        "score_momentum": score * 0.25,
        "ev_ebitda":      ev,
        "ev_revenue":     ev / 4.0,
        "pe_ratio":       ev * 1.5,
        "ebitda_margin":  mg,
        "gross_margin":   mg * 1.4,
        "net_margin":     mg * 0.5,
        "roe":            18.0,
        "revenue_growth": rev_g,
        "momentum_52w":   mom,
        "altman_z":       3.2,
        "beneish_m":      -2.5,
        "beta":           1.0,
        "price":          price,
        "market_cap":     price * 1e8,
        "revenue_ltm":    1e10,
        "currency":       "EUR",
        "sentiment_score": 0.1,
    }


def _make_tickers(sector: str, n: int, base_score: float = 60.0) -> list[dict]:
    tickers = []
    for i in range(n):
        score = max(20, min(95, base_score + (i - n // 2) * 3))
        tickers.append(_ticker_record(
            ticker   = f"{sector[:3].upper()}{i+1}",
            company  = f"{sector} Corp {i+1}",
            sector   = sector,
            score    = score,
            ev       = 8.0 + i * 1.2,
            mg       = 20.0 + i * 0.5,
            rev_g    = 0.05 + i * 0.01,
            mom      = 5.0 + i * 2.0,
            price    = 80.0 + i * 5,
        ))
    return tickers


# 3 univers de test
CASE1_NAME     = "Technology"
CASE1_UNIVERSE = "CAC 40"
CASE1_TICKERS  = _make_tickers("Technology", 6, base_score=76.0)  # Surpondérer

CASE2_NAME     = "Energy"
CASE2_UNIVERSE = "CAC 40"
CASE2_TICKERS  = _make_tickers("Energy", 1, base_score=58.0)       # Neutre, 1 seule société

CASE3_NAME     = "Technology"
CASE3_UNIVERSE = "S&P 500"
CASE3_TICKERS  = _make_tickers("Technology", 30, base_score=40.0)  # Sous-pondérer


# ─── HELPERS ──────────────────────────────────────────────────────────────────

def _count_pdf_pages(path: str) -> int:
    """Compte les pages d'un PDF via pdfplumber (déjà en dépendances)."""
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            return len(pdf.pages)
    except Exception:
        # Fallback minimal via bytes
        with open(path, "rb") as f:
            content = f.read()
        return content.count(b"/Type /Page") or content.count(b"/Type/Page") or -1


def _count_pptx_slides(path: str) -> int:
    from pptx import Presentation
    prs = Presentation(path)
    return len(prs.slides)


def _pptx_has_none_text(path: str) -> list[str]:
    """Retourne les textes 'None' visibles dans les slides."""
    from pptx import Presentation
    matches = []
    prs = Presentation(path)
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = para.text
                    if "None" in txt:
                        matches.append(f"slide {slide_idx+1}: {txt[:80]}")
    return matches


def _pdf_has_none_text(path: str) -> list[str]:
    """Retourne les occurrences de 'None' dans un PDF."""
    try:
        import pdfplumber
        matches = []
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages):
                txt = page.extract_text() or ""
                if "None" in txt:
                    for line in txt.split("\n"):
                        if "None" in line:
                            matches.append(f"p{i+1}: {line[:80]}")
        return matches
    except Exception:
        return []


# ─── TESTS PPTX SECTORIEL ─────────────────────────────────────────────────────

def test_sectoral_pptx(tickers: list, sector: str, universe: str, case_name: str):
    from outputs.sectoral_pptx_writer import SectoralPPTXWriter

    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, f"test_{case_name}_sectoral.pptx")
        result = SectoralPPTXWriter.generate(tickers, sector, universe, path)

        assert os.path.exists(path), f"[{case_name}] PPTX sectoriel non généré"
        assert os.path.getsize(path) > 50_000, \
            f"[{case_name}] PPTX sectoriel trop petit ({os.path.getsize(path)} bytes)"

        n_slides = _count_pptx_slides(path)
        assert n_slides == 20, \
            f"[{case_name}] PPTX sectoriel : {n_slides} slides (attendu 20)"

        nones = _pptx_has_none_text(path)
        assert not nones, \
            f"[{case_name}] PPTX sectoriel : 'None' visible dans {len(nones)} slide(s) : {nones[:3]}"

    print(f"  OK sectoral_pptx  [{case_name}] : 20 slides, aucun None")


# ─── TESTS PDF SECTORIEL ──────────────────────────────────────────────────────

def test_sector_pdf(tickers: list, sector: str, universe: str, case_name: str):
    from outputs.sector_pdf_writer import generate_sector_report

    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, f"test_{case_name}_sector.pdf")
        generate_sector_report(sector, tickers, path, universe=universe)

        assert os.path.exists(path), f"[{case_name}] PDF sectoriel non généré"
        assert os.path.getsize(path) > 20_000, \
            f"[{case_name}] PDF sectoriel trop petit ({os.path.getsize(path)} bytes)"

        n_pages = _count_pdf_pages(path)
        assert n_pages >= 6, \
            f"[{case_name}] PDF sectoriel : {n_pages} pages (attendu >= 6)"

        nones = _pdf_has_none_text(path)
        assert not nones, \
            f"[{case_name}] PDF sectoriel : 'None' visible : {nones[:3]}"

    print(f"  OK sector_pdf     [{case_name}] : {n_pages} pages, aucun None")


# ─── TESTS PDF INDICE ─────────────────────────────────────────────────────────

def _make_indice_data(tickers: list, sector: str, universe: str) -> dict:
    """Construit le dict data pour IndicePDFWriter — reproduit _build_indice_data minimal."""
    import statistics

    scores = [t.get("score_global", 50) for t in tickers]
    avg_score = statistics.mean(scores) if scores else 50
    conviction = min(95, max(50, int(avg_score)))

    signal_global = (
        "Surpondérer" if avg_score >= 60 else
        "Neutre"      if avg_score >= 40 else
        "Sous-pondérer"
    )

    nb    = len(tickers)
    sc    = int(avg_score)
    evs   = [t["ev_ebitda"] for t in tickers
             if t.get("ev_ebitda") is not None and 1.0 < t["ev_ebitda"] < 100]
    ev_str = f"{statistics.median(evs):.1f}x" if evs else "—"
    mgs   = [t.get("ebitda_margin", 0) for t in tickers if t.get("ebitda_margin") is not None]
    mg    = round(statistics.median(mgs), 1) if mgs else 0.0
    revgs = [t.get("revenue_growth", 0) for t in tickers if t.get("revenue_growth") is not None]
    croi  = round(statistics.median(revgs) * 100, 1) if revgs else 0.0
    croi_str = f"+{croi:.1f}%" if croi >= 0 else f"{croi:.1f}%"
    moms  = [t.get("momentum_52w", 0) for t in tickers if t.get("momentum_52w") is not None]
    mom   = round(statistics.median(moms), 1) if moms else 0.0
    mom_str = f"+{mom:.1f}%" if mom >= 0 else f"{mom:.1f}%"

    # secteurs: list of tuples (nom, nb, score_int, signal, ev_str, mg_float, croi_str, mom_str)
    secteurs_list = [(sector, nb, sc, signal_global, ev_str, mg, croi_str, mom_str)]

    # top3_secteurs: list of dicts with societes: list of (ticker, signal, ev_str, score_int)
    top_tickers = sorted(tickers, key=lambda t: t.get("score_global", 0), reverse=True)[:3]
    societes = []
    for t in top_tickers:
        ev_raw = t.get("ev_ebitda")
        ev_t = f"{ev_raw:.1f}x" if (ev_raw is not None and 1.0 < ev_raw < 100) else "—"
        sig_t = ("Surpondérer" if (t.get("score_global") or 0) >= 60
                 else ("Sous-pondérer" if (t.get("score_global") or 0) < 40 else "Neutre"))
        societes.append((t.get("ticker", "?"), sig_t, ev_t, t.get("score_global") or 0))
    if not societes:
        societes = [("—", "Neutre", "—", 0)]

    top3 = [{
        "nom":       sector,
        "signal":    signal_global,
        "score":     sc,
        "ev_ebitda": ev_str,
        "catalyseur": f"Score {sc}/100 · momentum {mom_str} · EV/EBITDA {ev_str}",
        "risque":    f"Croiss. revenus {croi_str} sur LTM — compression multiple si cycle se retourne",
        "societes":  societes,
    }]

    rotation = [(sector, "Expansion", "Moderee", "Moderee", signal_global)]

    return {
        "indice":         universe,
        "code":           universe[:6].upper(),
        "nb_secteurs":    1,
        "nb_societes":    nb,
        "signal_global":  signal_global,
        "conviction_pct": conviction,
        "date_analyse":   "25 mars 2026",
        "cours":          "—",
        "variation_ytd":  "—",
        "pe_forward":     "—",
        "score_global":   sc,
        "secteurs":       secteurs_list,
        "texte_macro":    f"Signal {signal_global} (conviction {conviction}%) — {nb} societes analysees.",
        "texte_signal":   f"Signal global {signal_global}. 1 secteur analyse.",
        "texte_rotation": f"Signal {signal_global} — selectivite recommandee sur {sector}.",
        "surp_noms":      sector if signal_global == "Surpondérer" else "aucun",
        "sous_noms":      sector if signal_global == "Sous-pondérer" else "aucun",
        "top3_secteurs":  top3,
        "rotation":       rotation,
        "catalyseurs": [
            ("Momentum fondamental", f"Score moyen {sc}/100 — leaders sectoriels solides", "3-6 mois"),
            ("Consolidation M&A",    "Vague de fusions potentielle — prime de valorisation", "6-12 mois"),
            ("Digitalisation",       "Automatisation IA — reduction couts structurelle",    "12-18 mois"),
        ],
        "risques": [
            ("Macro", "Recession potentielle — choc demande", "50", "Eleve"),
            ("Taux",  "Hausse taux directeurs BCE/Fed",        "35", "Modere"),
            ("Geo",   "Tensions geopolitiques — supply chain", "25", "Modere"),
        ],
        "scenarios": [
            ("Bull case",   "Score > 65 + momentum positif",   "Surpondérer",  "3-6 mois"),
            ("Bear case",   "Score < 40 via recession",        "Sous-pondérer", "6-12 mois"),
            ("Stagflation", "CPI > 3,5% + PIB < 1%",          "Sous-pondérer", "6-9 mois"),
        ],
        "finbert": {
            "nb_articles":   0,
            "score_agrege":  0.0,
            "positif": {"nb": 0, "score": "N/A", "themes": "Donnees non disponibles"},
            "neutre":  {"nb": 0, "score": "N/A", "themes": "Donnees non disponibles"},
            "negatif": {"nb": 0, "score": "N/A", "themes": "Donnees non disponibles"},
            "par_secteur": [(sector, "N/A", "Neutre")],
        },
        "methodologie": [
            ("Score sectoriel",  "Composite 0-100 : 40% momentum, 30% rev. BPA, 30% valorisation"),
            ("Signal",           "Surponderer (>60) / Neutre (40-60) / Sous-ponderer (<40)"),
            ("Valorisation",     "EV/EBITDA median LTM — source FMP / yfinance"),
        ],
        "perf_history": None,
    }


def test_indice_pdf(tickers: list, sector: str, universe: str, case_name: str):
    from outputs.indice_pdf_writer import IndicePDFWriter

    data = _make_indice_data(tickers, sector, universe)

    with tempfile.TemporaryDirectory() as tmpdir:
        path = os.path.join(tmpdir, f"test_{case_name}_indice.pdf")
        IndicePDFWriter.generate(data, path)

        assert os.path.exists(path), f"[{case_name}] PDF indice non généré"
        assert os.path.getsize(path) > 15_000, \
            f"[{case_name}] PDF indice trop petit ({os.path.getsize(path)} bytes)"

        n_pages = _count_pdf_pages(path)
        assert n_pages >= 5, \
            f"[{case_name}] PDF indice : {n_pages} pages (attendu >= 5)"

        nones = _pdf_has_none_text(path)
        assert not nones, \
            f"[{case_name}] PDF indice : 'None' visible : {nones[:3]}"

    print(f"  OK indice_pdf     [{case_name}] : {n_pages} pages, aucun None")


# ─── RUNNER ───────────────────────────────────────────────────────────────────

CASES = [
    (CASE1_NAME, CASE1_UNIVERSE, CASE1_TICKERS, "cas1_nominal"),
    (CASE2_NAME, CASE2_UNIVERSE, CASE2_TICKERS, "cas2_monosecteur"),
    (CASE3_NAME, CASE3_UNIVERSE, CASE3_TICKERS, "cas3_large"),
]

if __name__ == "__main__":
    errors = []
    total = 0

    for sector, universe, tickers, case_name in CASES:
        print(f"\n[{case_name}] — {sector} / {universe} / {len(tickers)} société(s)")

        tests = [
            ("PPTX sectoriel", test_sectoral_pptx),
            ("PDF sectoriel",  test_sector_pdf),
            ("PDF indice",     test_indice_pdf),
        ]

        for test_label, test_fn in tests:
            total += 1
            try:
                test_fn(tickers, sector, universe, case_name)
            except AssertionError as e:
                errors.append(f"  FAIL [{case_name}] {test_label}: {e}")
                print(f"  FAIL {test_label}: {e}")
            except Exception as e:
                errors.append(f"  ERROR [{case_name}] {test_label}: {type(e).__name__}: {e}")
                print(f"  ERROR {test_label}: {type(e).__name__}: {e}")

    print(f"\n{'='*60}")
    if errors:
        print(f"RESULTAT : {len(errors)} echec(s) sur {total} tests\n")
        for e in errors:
            print(e)
        sys.exit(1)
    else:
        print(f"RESULTAT : {total}/{total} tests passes — OK")
        sys.exit(0)
