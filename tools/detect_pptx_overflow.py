#!/usr/bin/env python
"""Détecteur d'overflow PPTX — analyse text boxes vs capacité réelle.

Utilise python-pptx pour lire la géométrie + PIL pour mesurer le texte
avec les vraies métriques de la police Calibri (police par défaut Office).

Sortie : rapport par slide, identifie OVERFLOW / TIGHT / OK.

Usage:
    python tools/detect_pptx_overflow.py path/to/file.pptx
    python tools/detect_pptx_overflow.py path/to/file.pptx --threshold 0.85
    python tools/detect_pptx_overflow.py path/to/file.pptx --only-overflow
"""

from __future__ import annotations
import sys
import argparse
import io
from pathlib import Path
from typing import Optional

# Force UTF-8 stdout pour caractères spéciaux (Windows cp1252 par défaut)
try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

# pip install python-pptx Pillow
try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    print("ERREUR : pip install python-pptx", file=sys.stderr)
    sys.exit(1)

try:
    from PIL import ImageFont
except ImportError:
    print("ERREUR : pip install Pillow", file=sys.stderr)
    sys.exit(1)


# ---------------------------------------------------------------------------
# Constantes typographie Calibri (police Office par défaut)
# ---------------------------------------------------------------------------
# 1 EMU = 1/914400 inch = 1/360000 cm
EMU_PER_CM = 360000
EMU_PER_PT = 12700  # 1 pt = 1/72 inch ≈ 12700 EMU

# Marges internes typiques d'une text box PowerPoint
# Defaults officiels PowerPoint : 0.1 inch L/R, 0.05 inch T/B
# Convertis en cm : L/R = 0.254 cm, T/B = 0.127 cm
DEFAULT_MARGIN_LEFT_CM = 0.254
DEFAULT_MARGIN_RIGHT_CM = 0.254
DEFAULT_MARGIN_TOP_CM = 0.127
DEFAULT_MARGIN_BOTTOM_CM = 0.127

# Calibri à 10pt : line_height ≈ 1.2 * font_size
LINE_HEIGHT_FACTOR = 1.20


# ---------------------------------------------------------------------------
# Police cache
# ---------------------------------------------------------------------------

_FONT_CACHE: dict = {}


def _get_font(size_pt: float) -> ImageFont.FreeTypeFont:
    """Charge Calibri à la taille demandée (cache)."""
    key = round(size_pt, 1)
    if key not in _FONT_CACHE:
        # Cherche calibri ou fallback Arial
        for name in ("calibri.ttf", "Calibri.ttf", "arial.ttf", "Arial.ttf"):
            try:
                _FONT_CACHE[key] = ImageFont.truetype(name, int(size_pt * 96 / 72))
                break
            except (OSError, IOError):
                continue
        else:
            _FONT_CACHE[key] = ImageFont.load_default()
    return _FONT_CACHE[key]


# ---------------------------------------------------------------------------
# Mesure de texte avec wrap simulation
# ---------------------------------------------------------------------------

def _wrap_text(text: str, max_width_px: float, font: ImageFont.FreeTypeFont) -> list[str]:
    """Simule le word wrap PowerPoint : coupe au mot quand la largeur dépasse."""
    if not text:
        return []
    lines: list[str] = []
    for paragraph in text.split("\n"):
        if not paragraph.strip():
            lines.append("")
            continue
        words = paragraph.split(" ")
        current = ""
        for word in words:
            test = (current + " " + word).strip() if current else word
            if font.getlength(test) <= max_width_px:
                current = test
            else:
                # Le mot ne tient pas --> flush ligne courante
                if current:
                    lines.append(current)
                # Si même le mot seul est plus large que la box, on l'ajoute quand même
                current = word
        if current:
            lines.append(current)
    return lines


def _measure_text_lines(
    text: str,
    width_cm: float,
    font_size_pt: float,
    margin_left_cm: float = DEFAULT_MARGIN_LEFT_CM,
    margin_right_cm: float = DEFAULT_MARGIN_RIGHT_CM,
) -> int:
    """Retourne le nombre de lignes nécessaires pour afficher `text`."""
    usable_width_cm = width_cm - margin_left_cm - margin_right_cm
    if usable_width_cm <= 0:
        return 999
    # Convertit cm --> pixels (96 DPI standard)
    usable_width_px = usable_width_cm * (96 / 2.54)
    font = _get_font(font_size_pt)
    return len(_wrap_text(text, usable_width_px, font))


def _max_lines_in_box(
    height_cm: float,
    font_size_pt: float,
    margin_top_cm: float = DEFAULT_MARGIN_TOP_CM,
    margin_bottom_cm: float = DEFAULT_MARGIN_BOTTOM_CM,
) -> int:
    """Calcule combien de lignes tiennent dans une box de hauteur donnée."""
    usable_height_cm = height_cm - margin_top_cm - margin_bottom_cm
    if usable_height_cm <= 0:
        return 0
    # line_height en cm : font_size_pt * 1.2 * (2.54 / 72)
    line_height_cm = font_size_pt * LINE_HEIGHT_FACTOR * (2.54 / 72)
    return max(1, int(usable_height_cm / line_height_cm))


# ---------------------------------------------------------------------------
# Extraction shape attributes
# ---------------------------------------------------------------------------

def _get_font_size(text_frame, default: float = 10.0) -> float:
    """Récupère la taille de police effective (premier run trouvé)."""
    try:
        for para in text_frame.paragraphs:
            for run in para.runs:
                if run.font.size is not None:
                    return run.font.size.pt
    except Exception:
        pass
    return default


def _emu_to_cm(emu: int) -> float:
    return emu / EMU_PER_CM


# ---------------------------------------------------------------------------
# Analyse principale
# ---------------------------------------------------------------------------

def analyze_pptx(
    pptx_path: Path,
    threshold: float = 1.10,
    only_overflow: bool = False,
    min_words: int = 5,
) -> dict:
    """Analyse chaque text box du PPTX, retourne un rapport.

    threshold : ratio fill (lignes_nécessaires / lignes_max) au-delà duquel
                on signale. 1.0 = overflow strict, 0.85 = TIGHT.
    only_overflow : ne reporte que les boxes en overflow (>= threshold).
    min_words : ignore les texts plus courts (titres, KPI numériques).
    """
    prs = Presentation(str(pptx_path))
    sw_cm = _emu_to_cm(prs.slide_width)
    sh_cm = _emu_to_cm(prs.slide_height)

    report = {
        "file": str(pptx_path),
        "slide_size_cm": (round(sw_cm, 2), round(sh_cm, 2)),
        "n_slides": len(prs.slides),
        "issues": [],  # liste de dict (slide, shape, text, ratio, lines_need, lines_max)
        "summary": {"overflow": 0, "tight": 0, "ok": 0, "ignored": 0},
    }

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape_idx, shape in enumerate(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue
            tf = shape.text_frame
            text = (tf.text or "").strip()
            if not text:
                continue
            n_words = len(text.split())
            if n_words < min_words:
                report["summary"]["ignored"] += 1
                continue

            try:
                w_cm = _emu_to_cm(shape.width)
                h_cm = _emu_to_cm(shape.height)
                left_cm = _emu_to_cm(shape.left)
                top_cm = _emu_to_cm(shape.top)
            except Exception:
                continue

            # Lecture des marges réelles si fournies
            try:
                ml = _emu_to_cm(tf.margin_left) if tf.margin_left else DEFAULT_MARGIN_LEFT_CM
                mr = _emu_to_cm(tf.margin_right) if tf.margin_right else DEFAULT_MARGIN_RIGHT_CM
                mt = _emu_to_cm(tf.margin_top) if tf.margin_top else DEFAULT_MARGIN_TOP_CM
                mb = _emu_to_cm(tf.margin_bottom) if tf.margin_bottom else DEFAULT_MARGIN_BOTTOM_CM
            except Exception:
                ml, mr, mt, mb = (DEFAULT_MARGIN_LEFT_CM, DEFAULT_MARGIN_RIGHT_CM,
                                  DEFAULT_MARGIN_TOP_CM, DEFAULT_MARGIN_BOTTOM_CM)

            font_size = _get_font_size(tf)
            lines_need = _measure_text_lines(text, w_cm, font_size, ml, mr)
            lines_max = _max_lines_in_box(h_cm, font_size, mt, mb)

            ratio = lines_need / max(lines_max, 1)
            # Cas "tient pile sur 1 ligne" (ratio=1.00) = pas un overflow réel
            # car le wrap PowerPoint utilise la full width sans scroll.
            if lines_need == 1 and ratio == 1.0:
                status = "OK"
                report["summary"]["ok"] += 1
            elif ratio >= threshold:
                status = "OVERFLOW"
                report["summary"]["overflow"] += 1
            elif ratio >= 0.85:
                status = "TIGHT"
                report["summary"]["tight"] += 1
            else:
                status = "OK"
                report["summary"]["ok"] += 1

            if only_overflow and status != "OVERFLOW":
                continue

            report["issues"].append({
                "slide": slide_idx,
                "shape_idx": shape_idx,
                "shape_name": shape.name,
                "pos_cm": (round(left_cm, 2), round(top_cm, 2)),
                "size_cm": (round(w_cm, 2), round(h_cm, 2)),
                "font_pt": round(font_size, 1),
                "n_words": n_words,
                "n_chars": len(text),
                "lines_need": lines_need,
                "lines_max": lines_max,
                "ratio": round(ratio, 2),
                "status": status,
                "preview": text[:90].replace("\n", " ") + ("…" if len(text) > 90 else ""),
            })

    return report


# ---------------------------------------------------------------------------
# Affichage console
# ---------------------------------------------------------------------------

def _color(s: str, code: str) -> str:
    """ANSI colors for terminal."""
    return f"\033[{code}m{s}\033[0m"


def print_report(report: dict, verbose: bool = False) -> None:
    print()
    print(f"{'='*80}")
    print(f"  Détection overflow PPTX — {report['file']}")
    print(f"{'='*80}")
    sw, sh = report["slide_size_cm"]
    print(f"  Slide : {sw}cm × {sh}cm  .  {report['n_slides']} slides")
    s = report["summary"]
    n_over = s["overflow"]
    n_tight = s["tight"]
    n_ok = s["ok"]
    n_ign = s["ignored"]
    print(f"  Bilan : {_color(f'{n_over} OVERFLOW', '31')}  .  "
          f"{_color(f'{n_tight} TIGHT', '33')}  .  "
          f"{_color(f'{n_ok} OK', '32')}  .  "
          f"{n_ign} ignorés")
    print()

    # Group issues by slide
    by_slide: dict[int, list] = {}
    for issue in report["issues"]:
        by_slide.setdefault(issue["slide"], []).append(issue)

    for slide_num in sorted(by_slide.keys()):
        print(f"--- SLIDE {slide_num} ---")
        for it in by_slide[slide_num]:
            color = "31" if it["status"] == "OVERFLOW" else ("33" if it["status"] == "TIGHT" else "32")
            tag = _color(f"[{it['status']:8s}]", color)
            print(f"  {tag} shape{it['shape_idx']:>2} | "
                  f"pos=({it['pos_cm'][0]:5.1f},{it['pos_cm'][1]:5.1f})cm | "
                  f"size=({it['size_cm'][0]:5.1f}×{it['size_cm'][1]:4.1f})cm | "
                  f"font={it['font_pt']:4.1f}pt | "
                  f"{it['n_words']:>3}m {it['n_chars']:>4}c | "
                  f"lignes {it['lines_need']:>2}/{it['lines_max']:>2} | "
                  f"ratio={it['ratio']:.2f}")
            if verbose or it["status"] == "OVERFLOW":
                print(f"             --> {it['preview']!r}")
        print()


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("pptx", help="Chemin du fichier PPTX")
    parser.add_argument("--threshold", type=float, default=1.10,
                        help="Seuil ratio overflow (1.10 = défaut, tolère 10%% "
                             "d'imprécision de la simulation wrap)")
    parser.add_argument("--only-overflow", action="store_true",
                        help="N'affiche que les boxes en OVERFLOW")
    parser.add_argument("--min-words", type=int, default=5,
                        help="Ignore les textes < N mots (défaut: 5)")
    parser.add_argument("--verbose", "-v", action="store_true",
                        help="Affiche le preview pour TOUTES les boxes")
    args = parser.parse_args()

    p = Path(args.pptx)
    if not p.exists():
        print(f"ERREUR : fichier introuvable : {p}", file=sys.stderr)
        return 1

    report = analyze_pptx(
        p, threshold=args.threshold,
        only_overflow=args.only_overflow, min_words=args.min_words,
    )
    print_report(report, verbose=args.verbose)
    # Exit code = nb d'overflows (0 = clean)
    return min(report["summary"]["overflow"], 1)


if __name__ == "__main__":
    sys.exit(main())
