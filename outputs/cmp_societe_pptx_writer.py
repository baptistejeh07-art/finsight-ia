# =============================================================================
# FinSight IA — Comparison PPTX Writer (Pitchbook comparatif 22 slides)
# outputs/comparison_pptx_writer.py
#
# Généré un pitchbook comparatif IB-quality en 22 slides via python-pptx.
# Accepte deux états pipeline LangGraph + metriques extraites.
#
# Usage :
#   from outputs.cmp_societe_pptx_writer import CmpSocietePPTXWriter
#   path = CmpSocietePPTXWriter().generate(state_a, state_b, output_path)
# =============================================================================

from __future__ import annotations

import io
import logging
import math
import sys
from datetime import date, datetime
from pathlib import Path
from typing import Any, Dict, List, Optional

log = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Palette couleurs — identique a pptx_writer.py
# ---------------------------------------------------------------------------
NAVY       = "1B3A6B"
NAVY_MID   = "2E5FA3"
NAVY_PALE  = "EEF3FA"
GREEN      = "1A7A4A"
GREEN_PALE = "EAF4EF"
GREEN_MID  = "2E8B57"
RED        = "A82020"
RED_PALE   = "FAF0EF"
AMBER      = "B06000"
AMBER_PALE = "FDF6E8"
WHITE      = "FFFFFF"
GREY_BG    = "F7F8FA"
GREY_TXT   = "555555"
GREY_LIGHT = "888888"
BLACK      = "0D0D0D"

# Couleurs distinctives par société
COLOR_A     = NAVY_MID   # Societe A — bleu
COLOR_A_PAL = NAVY_PALE
COLOR_B     = GREEN_MID  # Societe B — vert
COLOR_B_PAL = GREEN_PALE

import re as _re

def _x(text) -> str:
    if text is None:
        return ""
    s = str(text)
    # Supprimé caractères invalides XML 1.0
    s = _re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', s)
    # Supprimé markdown bold/italic du LLM (**text** et *text*)
    s = _re.sub(r'\*\*(.+?)\*\*', r'\1', s, flags=_re.DOTALL)
    s = _re.sub(r'\*(.+?)\*', r'\1', s, flags=_re.DOTALL)
    # Supprimé les asterisques residuels isoles
    s = _re.sub(r'\*+', '', s)
    return s


def rgb(hex_str: str):
    from pptx.dml.color import RGBColor
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ---------------------------------------------------------------------------
# Helpers numeriques
# ---------------------------------------------------------------------------
def _fr(v, dp: int = 1, suffix: str = "") -> str:
    if v is None:
        return "\u2014"
    try:
        return f"{float(v):,.{dp}f}".replace(",", " ").replace(".", ",") + suffix
    except Exception:
        return "\u2014"


def _frpct(v, signed: bool = False) -> str:
    if v is None:
        return "\u2014"
    try:
        fv = float(v)
        s = f"{fv * 100:+.1f}" if signed else f"{fv * 100:.1f}"
        return s.replace(".", ",") + " %"
    except Exception:
        return "\u2014"


def _frpct_margin(v) -> str:
    if v is None:
        return "\u2014"
    try:
        fv = float(v)
        if abs(fv) > 2.0:
            fv /= 100.0
        return f"{fv * 100:.1f}".replace(".", ",") + " %"
    except Exception:
        return "\u2014"


def _frx(v) -> str:
    if v is None:
        return "\u2014"
    try:
        return f"{float(v):.1f}".replace(".", ",") + "x"
    except Exception:
        return "\u2014"


def _frm(v, cur_sym: str = "$") -> str:
    """Format market cap / FCF / Trésorerie. Sortie en MILLIARDS.

    Convention : la valeur passée DOIT être en milliards. Les normalizations
    millions→milliards se font à la source dans extract_metrics (cmp_xlsx).
    """
    if v is None:
        return "\u2014"
    try:
        v = float(v)
        if abs(v) > 1_000_000_000_000:
            v = v / 1_000_000_000  # raw -> Mds
        elif abs(v) > 1_000_000:
            v = v / 1_000          # millions -> Mds
        if cur_sym == "EUR":
            sym_big, sym_small = "Md\u20ac", "M\u20ac"
        else:
            sym_big, sym_small = "Mds" + cur_sym, "M" + cur_sym
        if abs(v) >= 1:
            return _fr(v, 1) + " " + sym_big
        return _fr(v * 1000, 0) + " " + sym_small
    except Exception:
        return "\u2014"


def _upside(target, current) -> str:
    if not target or not current:
        return "\u2014"
    try:
        return f"{(float(target) / float(current) - 1) * 100:+.1f}".replace(".", ",") + " %"
    except Exception:
        return "\u2014"


def _g(obj, attr, default=None):
    if obj is None:
        return default
    if isinstance(obj, dict):
        return obj.get(attr, default)
    return getattr(obj, attr, default)


def _safe_str(v, default="\u2014") -> str:
    if v is None:
        return default
    s = str(v).strip()
    return s if s else default


_COMPANY_ABBREVIATIONS = [
    (", Société Européenne", " SE"),
    (", Société Européenne", " SE"),
    (" Société Européenne", " SE"),
    (", Société Anonyme", " SA"),
    (" Société Anonyme", " SA"),
    (" Société Anonyme", " SA"),
    (" Corporation", " Corp."),
    (" Incorporated", " Inc."),
    (" Public Limited Company", " plc"),
    (" & Company", " & Co."),
]


def _truncate(s, n: int) -> str:
    """Coupe au word boundary sans ellipse visible.

    Applique d'abord les abréviations standards (« Société Européenne » →
    « SE », « Corporation » → « Corp. ») pour raccourcir proprement les
    noms longs avant la troncation. Évite les coupes visuellement laides
    comme « LVMH Moët Hennessy - Louis Vuitton, Société » (mot coupé).
    """
    s = _safe_str(s, "")
    # Appliquer les abréviations pour gagner de la place sans perdre de sens
    for full, abbr in _COMPANY_ABBREVIATIONS:
        if full in s:
            s = s.replace(full, abbr)
    if len(s) <= n:
        return s
    cut = s[:n]
    last_space = cut.rfind(" ")
    if last_space > n // 2:
        cut = cut[:last_space]
    # Retire ponctuation trailing (« , » « - » « : ») qui ferait penser à
    # un mot coupé
    return cut.rstrip(" ,-:·—")


def _fit(s, n: int) -> str:
    """Coupe a n chars sans ajouter '...' — zones ou le débordement est invisible (PPTX clip)."""
    s = _safe_str(s, "")
    if len(s) <= n:
        return s
    cut = s[:n]
    last_space = cut.rfind(" ")
    return cut[:last_space] if last_space > n // 2 else cut


def _fr_date_long(d=None) -> str:
    from datetime import date as _d
    _MONTHS = ["janvier", "fevrier", "mars", "avril", "mai", "juin",
               "juillet", "aout", "septembre", "octobre", "novembre", "decembre"]
    if d is None:
        d = _d.today()
    elif isinstance(d, str):
        try:
            d = _d.fromisoformat(str(d).split("T")[0][:10])
        except Exception:
            return str(d)
    try:
        return f"{d.day} {_MONTHS[d.month - 1]} {d.year}"
    except Exception:
        return str(d)


# ---------------------------------------------------------------------------
# Helpers formes PPTX
# ---------------------------------------------------------------------------
def _set_slide_bg(slide, hex_color: str):
    from pptx.oxml.ns import qn
    from lxml import etree
    cSld = slide._element.find(qn('p:cSld'))
    if cSld is None:
        return
    existing = cSld.find(qn('p:bg'))
    if existing is not None:
        cSld.remove(existing)
    bg_xml = (
        '<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<p:bgPr>'
        f'<a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>'
        '<a:effectLst/>'
        '</p:bgPr>'
        '</p:bg>'
    )
    bg_elem = etree.fromstring(bg_xml)
    spTree = cSld.find(qn('p:spTree'))
    cSld.insert(list(cSld).index(spTree) if spTree is not None else 0, bg_elem)


def _add_text_alpha(slide, x, y, w, h, text, font_size,
                    color_hex="FFFFFF", alpha_val=15000, bold=False):
    from pptx.util import Cm, Pt
    from pptx.oxml.ns import qn
    from lxml import etree
    txBox = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = _x(text)
    run.font.name = "Calibri"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    r_elem = run._r
    rPr = r_elem.find(qn('a:rPr'))
    if rPr is None:
        rPr = etree.Element(qn('a:rPr'))
        r_elem.insert(0, rPr)
    for tag in (qn('a:solidFill'), qn('a:gradFill'), qn('a:noFill'), qn('a:pattFill')):
        for child in list(rPr.findall(tag)):
            rPr.remove(child)
    solidFill = etree.SubElement(rPr, qn('a:solidFill'))
    srgbClr   = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', color_hex.upper())
    alpha_elem = etree.SubElement(srgbClr, qn('a:alpha'))
    alpha_elem.set('val', str(alpha_val))
    rPr.remove(solidFill)
    rPr.insert(0, solidFill)
    return txBox


def add_rect(slide, x, y, w, h, fill_hex, line_hex=None, line_width_pt=0):
    from pptx.util import Cm, Pt
    shape = slide.shapes.add_shape(1, Cm(x), Cm(y), Cm(w), Cm(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_hex)
    if line_hex:
        shape.line.color.rgb = rgb(line_hex)
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, x, y, w, h, text, font_size=10, color_hex=BLACK,
                 bold=False, italic=False, align=None,
                 font_name="Calibri", wrap=True):
    from pptx.util import Cm, Pt
    from pptx.enum.text import PP_ALIGN
    if align is None:
        align = PP_ALIGN.LEFT
    txBox = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = _x(text) if text is not None else "\u2014"
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color_hex)
    return txBox


def _set_cell(cell, text, font_size=8, bold=False, color_hex=BLACK,
              fill_hex=None, align=None, font_name="Calibri", italic=False):
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    if align is None:
        align = PP_ALIGN.CENTER
    if fill_hex:
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(fill_hex)
    tf = cell.text_frame
    for para in tf.paragraphs:
        for run in list(para.runs):
            para._p.remove(run._r)
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = _x(text) if text is not None else "\u2014"
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color_hex)


def add_table(slide, x, y, w, h, num_rows, num_cols, col_widths_pct=None,
              header_data=None, rows_data=None,
              header_fill=NAVY, row_fills=None, border_hex=None):
    from pptx.util import Cm, Emu
    from pptx.enum.text import PP_ALIGN
    total_rows = num_rows + (1 if header_data else 0)
    tbl = slide.shapes.add_table(
        total_rows, num_cols, Cm(x), Cm(y), Cm(w), Cm(h)
    ).table
    total_emu = int(Cm(w))
    if col_widths_pct:
        for ci, pct in enumerate(col_widths_pct):
            tbl.columns[ci].width = int(total_emu * pct)
    else:
        col_w = total_emu // num_cols
        for ci in range(num_cols):
            tbl.columns[ci].width = col_w
    row_offset = 0
    if header_data:
        for ci, val in enumerate(header_data):
            cell = tbl.cell(0, ci)
            _set_cell(cell, val, font_size=7.5, bold=True,
                      color_hex=WHITE, fill_hex=header_fill,
                      align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)
        row_offset = 1
    if rows_data:
        default_fills = [WHITE, GREY_BG]
        for ri, row in enumerate(rows_data):
            if row_fills and ri < len(row_fills):
                fill = row_fills[ri]
            else:
                fill = default_fills[ri % 2]
            for ci, val in enumerate(row):
                cell = tbl.cell(ri + row_offset, ci)
                align = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
                _set_cell(cell, val, font_size=7.5, bold=False,
                          color_hex=BLACK, fill_hex=fill, align=align)
    return tbl


# ---------------------------------------------------------------------------
# Elements communs
# ---------------------------------------------------------------------------
def navy_bar(slide):
    add_rect(slide, 0, 0, 25.4, 1.65, NAVY)


def footer_bar(slide):
    add_rect(slide, 0, 13.39, 25.4, 0.89, NAVY)
    add_text_box(slide, 1.02, 13.44, 23.37, 0.56,
                 "FinSight IA  \u00b7  Usage confidentiel", 7, WHITE)


def section_dots(slide, active: int):
    from pptx.enum.text import PP_ALIGN
    xs = [17.02, 17.65, 18.29, 18.92, 19.56]
    for i, x in enumerate(xs):
        num = i + 1
        fill = WHITE if num == active else NAVY_MID
        add_rect(slide, x, 0.33, 0.51, 0.51, fill)
        color = NAVY if num == active else WHITE
        add_text_box(slide, x + 0.02, 0.34, 0.47, 0.44,
                     str(num), 7, color, bold=True, align=PP_ALIGN.CENTER)


# ─── i18n helper cmp société PPTX ─────────────────────────────────────────
_CMP_PPTX_LANG: str = "fr"

_CMP_PPTX_LABELS: dict[str, dict[str, str]] = {
    "exec_cmp": {"fr": "Executive Summary Comparatif",
                 "en": "Comparative Executive Summary",
                 "es": "Resumen Ejecutivo Comparativo",
                 "de": "Vergleichende Zusammenfassung",
                 "it": "Executive Summary Comparativo",
                 "pt": "Resumo Executivo Comparativo"},
    "sommaire": {"fr": "Sommaire", "en": "Table of contents",
                 "es": "Índice", "de": "Inhalt",
                 "it": "Sommario", "pt": "Sumário"},
    "profil_cmp": {"fr": "Profil Comparatif", "en": "Comparative Profile",
                   "es": "Perfil Comparativo", "de": "Vergleichsprofil",
                   "it": "Profilo Comparativo", "pt": "Perfil Comparativo"},
    "is_compare": {"fr": "Compte de Résultat Comparé",
                   "en": "Compared Income Statement",
                   "es": "Cuenta de Resultados Comparada",
                   "de": "Verglichene Gewinn- und Verlustrechnung",
                   "it": "Conto Economico Comparato",
                   "pt": "Demonstração de Resultados Comparada"},
    "rent_croiss": {"fr": "Rentabilité & Croissance Comparées",
                    "en": "Compared Profitability & Growth",
                    "es": "Rentabilidad y Crecimiento Comparados",
                    "de": "Verglichene Rentabilität & Wachstum",
                    "it": "Redditività & Crescita Comparate",
                    "pt": "Rentabilidade & Crescimento Comparados"},
    "bilan_liq": {"fr": "Bilan & Liquidité Comparés",
                  "en": "Compared Balance Sheet & Liquidity",
                  "es": "Balance y Liquidez Comparados",
                  "de": "Verglichene Bilanz & Liquidität",
                  "it": "Stato Patrimoniale & Liquidità Comparati",
                  "pt": "Balanço & Liquidez Comparados"},
    "mult_val": {"fr": "Multiples de Valorisation Comparés",
                 "en": "Compared Valuation Multiples",
                 "es": "Múltiplos de Valoración Comparados",
                 "de": "Verglichene Bewertungsmultiplikatoren",
                 "it": "Multipli di Valutazione Comparati",
                 "pt": "Múltiplos de Avaliação Comparados"},
    "dcf_cibles": {"fr": "DCF & Prix Cibles Comparés",
                   "en": "Compared DCF & Target Prices",
                   "es": "DCF y Precios Objetivo Comparados",
                   "de": "Verglichene DCF & Kursziele",
                   "it": "DCF & Prezzi Obiettivo Comparati",
                   "pt": "DCF & Preços-Alvo Comparados"},
    "gbm_upside": {"fr": "GBM Monte Carlo & Upside Relatif",
                   "en": "GBM Monte Carlo & Relative Upside",
                   "es": "GBM Monte Carlo y Upside Relativo",
                   "de": "GBM Monte Carlo & Relative Aufwärtsbewegung",
                   "it": "GBM Monte Carlo & Upside Relativo",
                   "pt": "GBM Monte Carlo & Upside Relativo"},
    "qual_fin": {"fr": "Qualité Financière & Solidité Bilancielle",
                 "en": "Financial Quality & Balance Sheet Strength",
                 "es": "Calidad Financiera y Solidez del Balance",
                 "de": "Finanzqualität & Bilanzstärke",
                 "it": "Qualità Finanziaria & Solidità Patrimoniale",
                 "pt": "Qualidade Financeira & Solidez do Balanço"},
    "risque_mom": {"fr": "Profil de Risque & Momentum Comparé",
                   "en": "Compared Risk Profile & Momentum",
                   "es": "Perfil de Riesgo y Momentum Comparado",
                   "de": "Verglichenes Risikoprofil & Momentum",
                   "it": "Profilo di Rischio & Momentum Comparato",
                   "pt": "Perfil de Risco & Momentum Comparado"},
    "score_cmp": {"fr": "FinSight Score Comparé",
                  "en": "Compared FinSight Score",
                  "es": "FinSight Score Comparado",
                  "de": "Verglichener FinSight-Score",
                  "it": "FinSight Score Comparato",
                  "pt": "FinSight Score Comparado"},
    "theses_bb": {"fr": "Thèses d'Investissement Bull / Bear",
                  "en": "Bull / Bear Investment Thèses",
                  "es": "Tesis de Inversión Bull / Bear",
                  "de": "Bull-/Bear-Investitionsthesen",
                  "it": "Tesi d'Investimento Bull / Bear",
                  "pt": "Teses de Investimento Bull / Bear"},
    "perf_52w": {"fr": "Performance Boursière 52 Semaines",
                 "en": "52-Week Stock Performance",
                 "es": "Rendimiento Bursátil 52 Semanas",
                 "de": "52-Wochen-Aktien-Performance",
                 "it": "Performance Azionaria 52 Settimane",
                 "pt": "Desempenho Bursátil 52 Semanas"},
    "verdict_cmp": {"fr": "Verdict Final Comparatif",
                    "en": "Final Comparative Verdict",
                    "es": "Veredicto Final Comparativo",
                    "de": "Endgültiges Vergleichsurteil",
                    "it": "Verdetto Finale Comparativo",
                    "pt": "Veredicto Final Comparativo"},
}


def _cplbl(key: str) -> str:
    spec = _CMP_PPTX_LABELS.get(key)
    if not spec:
        return key
    return spec.get(_CMP_PPTX_LANG) or spec.get("en") or spec.get("fr") or key


def slide_title(slide, title: str, subtitle: str = ""):
    from pptx.enum.text import PP_ALIGN
    add_text_box(slide, 1.02, 0.33, 15.49, 0.97,
                 title, 13, WHITE, bold=True)
    if subtitle:
        add_text_box(slide, 1.02, 1.73, 23.37, 0.56,
                     subtitle, 9, NAVY_MID)


def kpi_box(slide, x, y, w, h, value, label, sub="",
            fill=NAVY_PALE, accent=NAVY_MID):
    from pptx.enum.text import PP_ALIGN
    add_rect(slide, x, y, w, h, fill)
    add_rect(slide, x, y, 0.13, h, accent)
    val_color = NAVY if fill != NAVY else WHITE
    lbl_color = GREY_TXT if fill != NAVY else WHITE
    sub_color = GREY_LIGHT if fill != NAVY else NAVY_PALE
    add_text_box(slide, x + 0.25, y + 0.15, w - 0.38, 1.0,
                 value, 17, val_color, bold=True)
    add_text_box(slide, x + 0.25, y + 1.1, w - 0.38, 0.50,
                 label, 7.5, lbl_color)
    if sub:
        add_text_box(slide, x + 0.25, y + 1.38, w - 0.38, 0.42,
                     sub, 7, sub_color)


def divider_slide(prs, number_str: str, title: str, subtitle: str):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_bg(slide, NAVY)
    add_rect(slide, 0, 0, 0.3, 14.29, NAVY_MID)
    _add_text_alpha(slide, 1.27, 4.06, 22.86, 4.57,
                    number_str, 80, "FFFFFF", 15000, bold=True)
    add_text_box(slide, 1.27, 5.21, 21.59, 2.29,
                 title, 28, WHITE, bold=True)
    add_rect(slide, 1.27, 7.75, 7.62, 0.01, GREY_LIGHT)
    add_text_box(slide, 1.27, 8.08, 17.78, 0.89,
                 subtitle, 10, "AABBDD")
    add_text_box(slide, 1.02, 13.34, 23.37, 0.56,
                 "FinSight IA  \u00b7  Usage confidentiel", 7, "6677AA")
    return slide


def _company_header_band(slide, tkr_a, tkr_b, name_a="", name_b="", m_a=None, m_b=None):
    """Bande bicolore sous le titre identifiant les deux sociétés.
    Format cible : "Apple Inc. (AAPL)" (nom complet + ticker entre parentheses).
    Si m_a/m_b fournis, on y cherche company_name_a/company_name_b."""
    # Recuperation auto des noms depuis les metrics si fournis
    if m_a and not name_a:
        name_a = m_a.get("company_name_a") or ""
    if m_b and not name_b:
        name_b = m_b.get("company_name_b") or ""
    # Format : "Apple Inc. (AAPL)" si nom dispo, sinon juste le ticker
    lbl_a = f"{_truncate(name_a, 30)} ({tkr_a})" if name_a else f"{tkr_a}"
    lbl_b = f"{_truncate(name_b, 30)} ({tkr_b})" if name_b else f"{tkr_b}"
    add_rect(slide, 1.02, 1.73, 11.44, 0.56, COLOR_A_PAL)
    add_rect(slide, 1.02, 1.73, 0.18, 0.56, COLOR_A)
    add_text_box(slide, 1.35, 1.78, 11.0, 0.46, lbl_a, 8.5, NAVY, bold=True, wrap=False)
    add_rect(slide, 12.94, 1.73, 11.44, 0.56, COLOR_B_PAL)
    add_rect(slide, 12.94, 1.73, 0.18, 0.56, COLOR_B)
    add_text_box(slide, 13.27, 1.78, 11.0, 0.46, lbl_b, 8.5, GREEN, bold=True, wrap=False)


# ---------------------------------------------------------------------------
# LLM Synthesis — appel unique au debut
# ---------------------------------------------------------------------------
def _strip_llm_artifacts(s: str) -> str:
    """Nettoie les artefacts markdown que le LLM échappe parfois dans sa réponse.

    Bug B3 audit 27/04/2026 : le LLM reproduit parfois le préfixe du prompt
    ("# VERDICT FINAL", "EXECUTIVE SUMMARY Comparatif :", "## 1. CONTEXTE MACRO")
    tel quel dans sa réponse. Ces marqueurs apparaissent ensuite dans le PDF/PPTX
    final. On strip :
    - Headers markdown `# ... ## ... ### ...` (ouverture de ligne)
    - Préfixes prompt-echo connus (VERDICT FINAL, EXECUTIVE SUMMARY, etc.)
    - Bullets `* `/`- ` en début de ligne (rendus mal par les writers)
    """
    if not s:
        return ""
    import re as _re
    # 1. Strip les headers markdown `## SOMETHING` n'importe où dans le texte
    #    en les remplaçant par un simple saut de ligne (pour aérer le contenu).
    cleaned = _re.sub(r"\n?#{1,6}\s+", "\n", s)
    # 2. Si le texte démarre par un préfixe prompt-echo connu, le retirer
    #    jusqu'au premier `:` ou `—` ou newline.
    _PROMPT_ECHO_PATTERNS = (
        r"^\s*VERDICT\s+FINAL\b[^:\n—-]*[:—-]?\s*",
        r"^\s*EXECUTIVE\s+SUMMARY\b[^:\n—-]*[:—-]?\s*",
        r"^\s*COMMENTAIRE\b[^:\n—-]*[:—-]?\s*",
        r"^\s*ANALYSE\s+COMPARATIVE\b[^:\n—-]*[:—-]?\s*",
        r"^\s*Titre\s+préféré\s*[:—-]\s*\w+\s*[—-]?\s*",
    )
    for pat in _PROMPT_ECHO_PATTERNS:
        cleaned = _re.sub(pat, "", cleaned, flags=_re.IGNORECASE)
    # 3. Strip markdown emphase (**bold**, __underline__, *italic*, _italic_)
    #    sinon python-pptx peut interpréter ces caractères en sous-soulignements
    #    visibles à la place du texte (cf bug slide 7 audit 28/04/2026).
    cleaned = _re.sub(r"\*\*(.+?)\*\*", r"\1", cleaned)
    cleaned = _re.sub(r"__(.+?)__", r"\1", cleaned)
    cleaned = _re.sub(r"(?<!\w)\*([^\*\n]+?)\*(?!\w)", r"\1", cleaned)
    cleaned = _re.sub(r"(?<!\w)_([^_\n]+?)_(?!\w)", r"\1", cleaned)
    # 4. Strip lignes horizontales markdown (---, ===, ~~~ seuls sur une ligne)
    cleaned = _re.sub(r"^\s*[-=~]{3,}\s*$", "", cleaned, flags=_re.MULTILINE)
    # 5. Normaliser les sauts de ligne multiples (max 2 = paragraphe)
    cleaned = _re.sub(r"\n{3,}", "\n\n", cleaned)
    return cleaned.strip()


def _call_llm(prompt: str, system: str = "", max_tokens: int = 512) -> str:
    """Appel LLM avec fallback anthropic -> mistral.

    Strippe automatiquement les artefacts markdown (cf _strip_llm_artifacts).
    """
    try:
        sys_path_insert = str(Path(__file__).parent.parent)
        if sys_path_insert not in sys.path:
            sys.path.insert(0, sys_path_insert)
        from core.llm_provider import LLMProvider
        for provider in ("anthropic", "mistral", "gemini"):
            try:
                llm = LLMProvider(provider=provider)
                raw = llm.generate(prompt, system=system, max_tokens=max_tokens)
                return _strip_llm_artifacts(raw)
            except Exception as _e:
                log.debug(f"[cmp_pptx] LLM {provider} failed: {_e}")
                continue
    except Exception as e:
        log.warning(f"[cmp_pptx] LLM unavailable: {e}")
    return ""


def _word_clip(s, n):
    """Clip une chaine a n chars a la frontiere de mot (sans ellipsis)."""
    if not s:
        return ""
    s = str(s).strip()
    if len(s) <= n:
        return s
    cut = s[:n]
    sp = cut.rfind(" ")
    if sp > n // 2:
        cut = cut[:sp]
    return cut.rstrip(" ,;:-")


def _generate_synthesis(m_a: dict, m_b: dict) -> dict:
    """
    Genere la synthese comparative via LLM.
    Retourne un dict avec les cles : exec_summary, financial_text,
    valuation_text, quality_text, verdict_text, thesis_a, thesis_b.
    """
    tkr_a = m_a.get("ticker_a") or "A"
    tkr_b = m_b.get("ticker_b") or "B"
    name_a = m_a.get("company_name_a") or tkr_a
    name_b = m_b.get("company_name_b") or tkr_b

    def _pct(v):
        if v is None: return "N/A"
        try: return f"{float(v)*100:.1f}%"
        except Exception: return "N/A"
    def _x2(v):
        if v is None: return "N/A"
        try: return f"{float(v):.1f}x"
        except Exception: return "N/A"
    def _n(v):
        if v is None: return "N/A"
        try: return f"{float(v):.1f}"
        except Exception: return "N/A"

    sec_a = m_a.get('sector_a') or m_a.get('sector') or 'N/D'
    sec_b = m_b.get('sector_b') or m_b.get('sector') or 'N/D'

    data_str = f"""
============================================================
COMPARAISON DETAILLEE : {tkr_a} ({name_a}) vs {tkr_b} ({name_b})
Secteurs : {sec_a} vs {sec_b}
============================================================

[1] VALORISATION (multiples)
{tkr_a}: PE={_n(m_a.get('pe_ratio'))}, Forward PE={_n(m_a.get('forward_pe'))}, EV/EBITDA={_x2(m_a.get('ev_ebitda'))}, EV/Sales={_x2(m_a.get('ev_sales'))}, P/B={_n(m_a.get('price_to_book'))}, P/FCF={_n(m_a.get('p_fcf'))}, PEG={_n(m_a.get('peg_ratio'))}, FCF yield={_pct(m_a.get('fcf_yield'))}, Div yield={_pct(m_a.get('dividend_yield'))}
{tkr_b}: PE={_n(m_b.get('pe_ratio'))}, Forward PE={_n(m_b.get('forward_pe'))}, EV/EBITDA={_x2(m_b.get('ev_ebitda'))}, EV/Sales={_x2(m_b.get('ev_sales'))}, P/B={_n(m_b.get('price_to_book'))}, P/FCF={_n(m_b.get('p_fcf'))}, PEG={_n(m_b.get('peg_ratio'))}, FCF yield={_pct(m_b.get('fcf_yield'))}, Div yield={_pct(m_b.get('dividend_yield'))}
Mediane sectorielle : PE~{_n(m_a.get('sector_median_pe'))}, EV/EBITDA~{_x2(m_a.get('sector_median_ev_ebitda'))}

[2] PROFITABILITE & EFFICACITE OPERATIONNELLE
{tkr_a}: EBITDA margin LTM={_pct(m_a.get('ebitda_margin_ltm'))} (y-1: {_pct(m_a.get('ebitda_margin_y1'))}, y-2: {_pct(m_a.get('ebitda_margin_y2'))}, trend: {m_a.get('ebitda_margin_trend','N/D')}), EBIT margin={_pct(m_a.get('ebit_margin'))}, Net margin={_pct(m_a.get('net_margin_ltm'))}, ROIC={_pct(m_a.get('roic'))}, ROE={_pct(m_a.get('roe'))}, Cash conversion={_n(m_a.get('cash_conversion'))}
{tkr_b}: EBITDA margin LTM={_pct(m_b.get('ebitda_margin_ltm'))} (y-1: {_pct(m_b.get('ebitda_margin_y1'))}, y-2: {_pct(m_b.get('ebitda_margin_y2'))}, trend: {m_b.get('ebitda_margin_trend','N/D')}), EBIT margin={_pct(m_b.get('ebit_margin'))}, Net margin={_pct(m_b.get('net_margin_ltm'))}, ROIC={_pct(m_b.get('roic'))}, ROE={_pct(m_b.get('roe'))}, Cash conversion={_n(m_b.get('cash_conversion'))}

[3] CROISSANCE
{tkr_a}: Rev CAGR 3y={_pct(m_a.get('revenue_cagr_3y'))}, Rev growth fwd={_pct(m_a.get('revenue_growth_fwd'))}, EPS growth={_pct(m_a.get('eps_growth'))}, EPS growth fwd={_pct(m_a.get('eps_growth_fwd'))}
{tkr_b}: Rev CAGR 3y={_pct(m_b.get('revenue_cagr_3y'))}, Rev growth fwd={_pct(m_b.get('revenue_growth_fwd'))}, EPS growth={_pct(m_b.get('eps_growth'))}, EPS growth fwd={_pct(m_b.get('eps_growth_fwd'))}

[4] LEVIER, SOLIDITE BILANCIELLE & QUALITE
{tkr_a}: ND/EBITDA={_n(m_a.get('net_debt_ebitda'))}, Net debt={_n(m_a.get('net_debt'))}M, Cash={_n(m_a.get('cash'))}M, Interest cov={_n(m_a.get('interest_coverage'))}x, Current ratio={_n(m_a.get('current_ratio'))}, Quick ratio={_n(m_a.get('quick_ratio'))}, Capex/Rev={_pct(m_a.get('capex_to_revenue'))}, Piotroski={m_a.get('piotroski_score','N/D')}/9, Altman Z={_n(m_a.get('altman_z'))}, Beneish M={_n(m_a.get('beneish_mscore'))}
{tkr_b}: ND/EBITDA={_n(m_b.get('net_debt_ebitda'))}, Net debt={_n(m_b.get('net_debt'))}M, Cash={_n(m_b.get('cash'))}M, Interest cov={_n(m_b.get('interest_coverage'))}x, Current ratio={_n(m_b.get('current_ratio'))}, Quick ratio={_n(m_b.get('quick_ratio'))}, Capex/Rev={_pct(m_b.get('capex_to_revenue'))}, Piotroski={m_b.get('piotroski_score','N/D')}/9, Altman Z={_n(m_b.get('altman_z'))}, Beneish M={_n(m_b.get('beneish_mscore'))}

[5] PERFORMANCE BOURSIERE & RISQUE
{tkr_a}: Cours={_n(m_a.get('share_price'))}, 52W high/low={_n(m_a.get('week52_high'))}/{_n(m_a.get('week52_low'))}, Perf 1M={_pct(m_a.get('perf_1m'))}, 3M={_pct(m_a.get('perf_3m'))}, 1Y={_pct(m_a.get('perf_1y'))}, Beta={_n(m_a.get('beta'))}, Vol 52W={_pct(m_a.get('volatility_52w'))}, RSI={_n(m_a.get('rsi'))}, Momentum score={_n(m_a.get('momentum_score'))}/10, VaR 95% 1M={_pct(m_a.get('var_95_1m'))}
{tkr_b}: Cours={_n(m_b.get('share_price'))}, 52W high/low={_n(m_b.get('week52_high'))}/{_n(m_b.get('week52_low'))}, Perf 1M={_pct(m_b.get('perf_1m'))}, 3M={_pct(m_b.get('perf_3m'))}, 1Y={_pct(m_b.get('perf_1y'))}, Beta={_n(m_b.get('beta'))}, Vol 52W={_pct(m_b.get('volatility_52w'))}, RSI={_n(m_b.get('rsi'))}, Momentum score={_n(m_b.get('momentum_score'))}/10, VaR 95% 1M={_pct(m_b.get('var_95_1m'))}

[6] DCF & MONTE CARLO (cibles)
{tkr_a}: DCF base={_n(m_a.get('dcf_base'))}, bull={_n(m_a.get('dcf_bull'))}, bear={_n(m_a.get('dcf_bear'))}, upside base={_pct(m_a.get('dcf_upside_base'))}, MC p10/p50/p90={_n(m_a.get('monte_carlo_p10'))}/{_n(m_a.get('monte_carlo_p50'))}/{_n(m_a.get('monte_carlo_p90'))}
{tkr_b}: DCF base={_n(m_b.get('dcf_base'))}, bull={_n(m_b.get('dcf_bull'))}, bear={_n(m_b.get('dcf_bear'))}, upside base={_pct(m_b.get('dcf_upside_base'))}, MC p10/p50/p90={_n(m_b.get('monte_carlo_p10'))}/{_n(m_b.get('monte_carlo_p50'))}/{_n(m_b.get('monte_carlo_p90'))}

[7] SCORES COMPOSITES & RECOMMANDATIONS
{tkr_a}: FinSight={m_a.get('finsight_score','N/D')}/100, Piotroski={m_a.get('piotroski_score','N/D')}/9, Reco={m_a.get('recommendation','N/D')}, Conviction={_pct(m_a.get('conviction'))}, Margin of safety={_pct(m_a.get('margin_of_safety'))}
{tkr_b}: FinSight={m_b.get('finsight_score','N/D')}/100, Piotroski={m_b.get('piotroski_score','N/D')}/9, Reco={m_b.get('recommendation','N/D')}, Conviction={_pct(m_b.get('conviction'))}, Margin of safety={_pct(m_b.get('margin_of_safety'))}
"""

    system_msg = (
        "Tu es un analyste sell-side senior (style JPMorgan Research, Goldman Sachs). "
        "Reponds en francais correct, RIGOUREUX, technique, IB-grade. "
        "RÈGLE TYPOGRAPHIE ABSOLUE : utilise TOUS les accents (é è ê à ù ç î ô), les cédilles, "
        "les apostrophes droites ' et les guillemets français « ». N'ECRIS JAMAIS sans accents — "
        "ce serait du français cassé, inacceptable. "
        "Pas d'emojis, pas de markdown ** mis a part les separateurs explicitement demandes. "
        "Structure : QUOI (constat chiffré) -> POURQUOI (mécanisme économique) -> IMPLICATIONS (décision investisseur). "
        "RÈGLE DONNÉES : tu disposes ci-dessous de TOUTES les données chiffrées. "
        "Tu ne dois JAMAIS dire 'données indisponibles' ni 'multiples non disponibles' ni inventer de chiffres "
        "absents du contexte fourni. Si un champ est N/A, tu l'omets sans le commenter."
    )

    results = {}

    # Executive summary (110 mots max) — chiffres concrets exiges
    r = _call_llm(
        f"Redige un executive summary comparatif de 110 mots MAX pour {tkr_a} vs {tkr_b}. "
        f"Cite AU MOINS 4 chiffres concrets (PE, EV/EBITDA, marge EBITDA, ROIC) avec leur écart relatif. "
        f"Conclus sur quel titre offre le meilleur rapport qualité/valorisation et pourquoi (1 phrase).\n{data_str}",
        system=system_msg, max_tokens=350
    )
    results["exec_summary"] = _word_clip(r, 900)

    # Financial commentary (120 mots max) — P&L + bilan + cash gen
    r = _call_llm(
        f"Redige un commentaire financier comparatif de 120 mots MAX : "
        f"tendance des marges EBITDA sur 3 ans (cite les valeurs), CAGR revenus 3y, ROIC vs ROE, "
        f"qualité de la conversion cash (cash conversion ratio), levier ND/EBITDA. "
        f"Conclus en 1 phrase sur la société la plus saine financierement.\n{data_str}",
        system=system_msg, max_tokens=400
    )
    results["financial_text"] = _word_clip(r, 1000)

    # Valuation commentary (130 mots max) — multiples + prime/décote + DCF
    r = _call_llm(
        f"Redige un commentaire de valorisation de 130 mots MAX : "
        f"comparé PE, EV/EBITDA, P/B, FCF yield (cite tous les chiffres). "
        f"Calcule la prime/décote relative en %. "
        f"Confronte avec la Médiane sectorielle. "
        f"Mentionne les cibles DCF base et l'upside implicite. "
        f"Conclus sur quelle valorisation est la plus attractive aujourd'hui et pourquoi.\n{data_str}",
        system=system_msg, max_tokens=450
    )
    results["valuation_text"] = _word_clip(r, 1100)

    # Quality commentary (120 mots max) — Piotroski + Altman + Beneish
    r = _call_llm(
        f"Redige un commentaire de qualité financière de 120 mots MAX : "
        f"comparé les Piotroski F-Scores avec les composants (ROA, levier, accruals), "
        f"interprète Altman Z (zone safe/grey/distress) et Beneish M-Score (manipulation comptable). "
        f"Comparé les ratios de liquidité (current/quick) et la couverture intérêts. "
        f"Conclus sur la société au bilan le plus robuste.\n{data_str}",
        system=system_msg, max_tokens=400
    )
    results["quality_text"] = _word_clip(r, 1100)

    # ── Verdict final — REFONTE DÉCISIONNELLE (Baptiste #173-177) ───────
    # Le LLM est LIBRE de choisir entre tkr_a, tkr_b, ou "Pas de préférence
    # forte" selon son analyse pondérée. Le score FinSight est fourni comme
    # INDICATION (pas comme vérité), aux côtés des news fraîches et du
    # contexte macro live. Aucune coercion post-LLM.
    #
    # Les macro/news sont injectés via core.llm_context pour que le LLM ait
    # accès au contexte le plus récent possible.
    _fs_a_num = m_a.get('finsight_score') or 0
    _fs_b_num = m_b.get('finsight_score') or 0

    # Bloc contexte enrichi pour chaque ticker (macro partagé + news par ticker
    # + finsight score expliqué)
    _ctx_block = ""
    try:
        from core.llm_context import (
            fetch_macro_live, format_macro_for_prompt,
            format_news_for_prompt, format_finsight_explanation,
        )
        _macro = fetch_macro_live()

        # Sentiment results : on essaie de récupérer depuis le state si dispo,
        # sinon fallback sur les champs m_*. Dans le flow app.py, les sentiment
        # results sont dans state.sentiment, mais ici on n'a que m_a/m_b.
        # Les samples du sentiment sont exposés via _samples_a/b dans m_*.
        _sentiment_a = m_a.get("_sentiment_result")
        _sentiment_b = m_b.get("_sentiment_result")

        _ctx_parts = [
            "=" * 70,
            "CONTEXTE DÉCISIONNEL ENRICHI (pour choix éclairé)",
            "=" * 70,
            format_macro_for_prompt(_macro),
            "",
            format_news_for_prompt(_sentiment_a, ticker=tkr_a),
            "",
            format_news_for_prompt(_sentiment_b, ticker=tkr_b),
            "",
            format_finsight_explanation(
                score_global=_fs_a_num,
                score_value=m_a.get("score_value"),
                score_growth=m_a.get("score_growth"),
                score_quality=m_a.get("score_quality"),
                score_momentum=m_a.get("score_momentum"),
                profile=m_a.get("_profile", "STANDARD"),
                ticker=tkr_a,
            ),
            "",
            format_finsight_explanation(
                score_global=_fs_b_num,
                score_value=m_b.get("score_value"),
                score_growth=m_b.get("score_growth"),
                score_quality=m_b.get("score_quality"),
                score_momentum=m_b.get("score_momentum"),
                profile=m_b.get("_profile", "STANDARD"),
                ticker=tkr_b,
            ),
            "=" * 70,
        ]
        _ctx_block = "\n".join(_ctx_parts)
    except Exception as _e_ctx:
        import logging as _log
        _log.getLogger(__name__).warning(f"[cmp_synthesis] contexte enrichi fail: {_e_ctx}")
        _ctx_block = ""

    _verdict_prompt = (
        f"VERDICT FINAL COMPARATIF — {tkr_a} ({name_a}) vs {tkr_b} ({name_b})\n\n"
        f"TU ES LE DÉCIDEUR. Analyse les deux sociétés et rends un verdict ÉCLAIRÉ en "
        f"pondérant intelligemment :\n"
        f"  1. Les métriques fondamentales (multiples, marges, ROIC, bilan)\n"
        f"  2. Le contexte macro ACTUEL (taux, VIX, régime, récession)\n"
        f"  3. Les news RÉCENTES avec sentiment\n"
        f"  4. Le score FinSight comme repère chiffré (PAS comme vérité — il est\n"
        f"     rétrospectif et ne tient pas compte de la macro ni des news)\n"
        f"  5. Ta connaissance des catalyseurs sectoriels forward-looking\n\n"
        f"Tu peux CHOISIR {tkr_a}, OU choisir {tkr_b}, OU déclarer NEUTRE si les "
        f"arguments pour et contre sont équilibrés — c'est un choix honnête "
        f"préférable à un verdict arbitraire.\n\n"
        f"STRUCTURE OBLIGATOIRE de ta réponse (120 mots MAX au total) :\n"
        f"1. Première ligne : \"Choix : {tkr_a}.\" OU \"Choix : {tkr_b}.\" OU "
        f"\"Choix : Pas de préférence forte.\"\n"
        f"2. \"Pourquoi :\" suivi de 2-3 arguments CHIFFRÉS qui justifient ton choix "
        f"(ou dans le cas neutre, les forces des deux côtés).\n"
        f"3. \"Risque principal :\" suivi du risque clé de la thèse retenue "
        f"(ou des deux risques si neutre).\n\n"
        f"Si tu diverges d'une indication simpliste qui dirait juste \"le meilleur "
        f"score FinSight gagne\", explique explicitement pourquoi.\n\n"
        f"{_ctx_block}\n\n"
        f"{data_str}"
    )
    r = _call_llm(_verdict_prompt, system=system_msg, max_tokens=500)

    # Pas de coercion — on fait juste un clip de longueur et on nettoie.
    _clean = (r or "").strip()
    if not _clean or len(_clean) < 50:
        # Fallback si LLM a complètement échoué : affiche les deux avec scores
        _clean = (
            f"Choix : Pas de préférence forte. "
            f"Pourquoi : les deux sociétés présentent un profil équilibré — "
            f"{tkr_a} affiche un score FinSight de {_fs_a_num:.0f}/100 et "
            f"{tkr_b} {_fs_b_num:.0f}/100. En l'absence d'analyse LLM disponible, "
            f"un investisseur doit arbitrer selon son horizon, sa tolérance au risque "
            f"et sa conviction sectorielle. Risque principal : exécution et conditions "
            f"macro défavorables pour les deux profils."
        )

    results["verdict_text"] = _word_clip(_clean, 1200)

    # ── Parse le ticker choisi par le LLM pour le partager aux writers ──
    # Les writers (cover PDF, cover PPTX, UI Streamlit) affichent le "Choix
    # préféré" = décision du LLM, pas du finsight_score.
    import re as _re_v
    _llm_choice = None
    _verdict_lower = _clean.lower()
    _m_choice = _re_v.search(r"choix\s*[:\-]?\s*(\w[\w\.\-]*)", _verdict_lower)
    if _m_choice:
        _raw_choice = _m_choice.group(1).strip().upper()
        # Nettoie la ponctuation finale
        _raw_choice = _raw_choice.rstrip(".,;:!?")
        if _raw_choice in (tkr_a.upper(), tkr_b.upper()):
            _llm_choice = _raw_choice
        elif "pas de" in _verdict_lower[:60] or "neutre" in _verdict_lower[:60] or "équilibré" in _verdict_lower[:60]:
            _llm_choice = "NEUTRAL"
    results["llm_choice"] = _llm_choice  # None | "TKR_A" | "TKR_B" | "NEUTRAL"

    def _split_bull_bear(r, sep="|||"):
        """Split bull/bear text — essaie plusieurs separateurs.

        Bug B8 audit 27/04 : avant clip raw [:350] sans frontière de mot,
        coupait au milieu d'un mot. Maintenant `_word_clip` respecte la
        dernière espace, et limite augmentée à 450 (60 mots ≈ 400 chars).
        """
        if not r:
            return "", ""
        LIM = 450
        for token in [sep, "\n|||", "|||\n", "---", "\n\nBear", "\n\nBEAR",
                      "\nThese bear", "\nThese Bull", "\nBear"]:
            if token in r:
                parts = r.split(token, 1)
                return _word_clip(parts[0].strip(), LIM), _word_clip(parts[1].strip(), LIM)
        # Dernier recours : chercher un pattern "bear" dans le texte
        m = _re.search(r'(?i)\n+(?:Thèse? )?bear', r)
        if m:
            return _word_clip(r[:m.start()].strip(), LIM), _word_clip(r[m.start():].strip(), LIM)
        return _word_clip(r, LIM), ""

    # Thèses bull/bear A (60 mots chacune avec chiffres)
    r = _call_llm(
        f"Pour {tkr_a}, donne 1 Thèse BULL (60 mots) et 1 Thèse BEAR (60 mots) "
        f"en les separant par '|||'. Cite au moins 2 chiffres concrets dans chaque Thèse "
        f"(marge, croissance, multiple, levier). Pas de titre.\n{data_str}",
        system=system_msg, max_tokens=350
    )
    results["bull_a"], results["bear_a"] = _split_bull_bear(r)

    # Thèses bull/bear B (60 mots chacune avec chiffres)
    r = _call_llm(
        f"Pour {tkr_b}, donne 1 Thèse BULL (60 mots) et 1 Thèse BEAR (60 mots) "
        f"en les separant par '|||'. Cite au moins 2 chiffres concrets dans chaque Thèse "
        f"(marge, croissance, multiple, levier). Pas de titre.\n{data_str}",
        system=system_msg, max_tokens=350
    )
    results["bull_b"], results["bear_b"] = _split_bull_bear(r)

    # Narratif cours 52 semaines avec contexte macro (pour PDF page 11 + PPTX)
    _perf_a_1y = m_a.get('perf_1y')
    _perf_b_1y = m_b.get('perf_1y')
    _perf_str = (
        f"{tkr_a} perf 1Y: {_perf_a_1y*100:.1f}% " if _perf_a_1y else f"{tkr_a} perf 1Y: N/D "
    ) + (
        f"| {tkr_b} perf 1Y: {_perf_b_1y*100:.1f}%" if _perf_b_1y else f"| {tkr_b} perf 1Y: N/D"
    )
    _price_ctx = (
        f"Secteur A: {m_a.get('sector_a','N/D')} | Secteur B: {m_b.get('sector_b','N/D')} | "
        f"Beta A: {m_a.get('beta','N/D')} | Beta B: {m_b.get('beta','N/D')} | "
        f"{_perf_str} | "
        f"52W High A: {m_a.get('week52_high','N/D')} | 52W Low A: {m_a.get('week52_low','N/D')} | "
        f"52W High B: {m_b.get('week52_high','N/D')} | 52W Low B: {m_b.get('week52_low','N/D')} | "
        f"PE A: {m_a.get('pe_ratio','N/D')} | PE B: {m_b.get('pe_ratio','N/D')} | "
        f"Earnings Growth A: {m_a.get('eps_growth','N/D')} | Earnings Growth B: {m_b.get('eps_growth','N/D')}"
    )
    r = _call_llm(
        f"Analyse en 200 mots MAX la trajectoire boursière de {tkr_a} vs {tkr_b} sur 12 mois. "
        f"STRUCTURE OBLIGATOIRE :\n"
        f"1. Contexte macro 12 mois (cycle taux directeurs Fed/BCE, cycle Économique, rotation sectorielle, "
        f"liquidité marchee actions, dollar) — 3-4 phrases.\n"
        f"2. Événements concrets ayant impacte les cours (résultats trimestriels, guidance, M&A, "
        f"sanctions réglementaires, lancements produits, géopolitique) — cite au moins 2 événements "
        f"concrets datables pour CHAQUE titre.\n"
        f"3. Comparaison des trajectoires : qui surperforme, de combien, et explication. "
        f"4. Catalyseurs a surveiller sur les 3-6 prochains mois (publications, guidance, secteur).\n"
        f"Sois précis et factuel, pas de generalites vagues. "
        f"Données : {_price_ctx}",
        system=system_msg, max_tokens=700
    )
    results["price_narrative"] = _word_clip(r, 2000)

    return results


# ---------------------------------------------------------------------------
# Graphiques matplotlib
# ---------------------------------------------------------------------------
def _make_chart_buf(fig) -> io.BytesIO:
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    import matplotlib.pyplot as plt
    plt.close(fig)
    buf.seek(0)
    return buf


def _chart_revenue_ebitda(m_a: dict, m_b: dict, tkr_a: str, tkr_b: str) -> Optional[io.BytesIO]:
    """Graphique CA et marge EBITDA comparatif."""
    try:
        import matplotlib.pyplot as plt
        import matplotlib
        matplotlib.use('Agg')
        import numpy as np

        fig, axes = plt.subplots(1, 2, figsize=(10.5, 3.2))
        fig.patch.set_facecolor('white')

        # Couleurs
        ca = '#2E5FA3'
        cb = '#2E8B57'

        # --- Panneau gauche : Marges EBITDA LTM / N-1 / N-2 ---
        ax1 = axes[0]
        ax1.set_facecolor('#FAFAFA')
        years_lbl = ['N-2', 'N-1', 'LTM']
        vals_a = [
            _pct_val(m_a.get('ebitda_margin_y2')),
            _pct_val(m_a.get('ebitda_margin_y1')),
            _pct_val(m_a.get('ebitda_margin_ltm')),
        ]
        vals_b = [
            _pct_val(m_b.get('ebitda_margin_y2')),
            _pct_val(m_b.get('ebitda_margin_y1')),
            _pct_val(m_b.get('ebitda_margin_ltm')),
        ]
        x = np.arange(len(years_lbl))
        w = 0.32
        bars_a = ax1.bar(x - w/2, vals_a, w, label=tkr_a, color=ca, alpha=0.85)
        bars_b = ax1.bar(x + w/2, vals_b, w, label=tkr_b, color=cb, alpha=0.85)
        for bar in list(bars_a) + list(bars_b):
            h = bar.get_height()
            if h and h > 0:
                ax1.text(bar.get_x() + bar.get_width()/2., h + 0.5,
                         f'{h:.0f}%', ha='center', va='bottom', fontsize=7, color='#333')
        _em_a = _pct_val(m_a.get('ebitda_margin_ltm'))
        _em_b = _pct_val(m_b.get('ebitda_margin_ltm'))
        _em_leader = tkr_a if _em_a > _em_b else (tkr_b if _em_b > _em_a else None)
        _em_delta = abs(_em_a - _em_b)
        _em_title = (f"{_em_leader} plus profitabl. : {_em_a:.0f}% vs {_em_b:.0f}% EBITDA LTM (+{_em_delta:.0f}pts)"
                     if _em_leader else f"Marge EBITDA LTM {_em_a:.0f}% = {_em_b:.0f}%")
        ax1.set_title(_em_title, fontsize=8.5, fontweight='bold', color='#1B3A6B', pad=4)
        ax1.set_xticks(x); ax1.set_xticklabels(years_lbl, fontsize=8)
        ax1.set_ylabel('%', fontsize=8); ax1.legend(fontsize=7.5)
        ax1.spines['top'].set_visible(False); ax1.spines['right'].set_visible(False)
        ax1.tick_params(labelsize=8)

        # --- Panneau droit : Revenue CAGR + ROE + ROIC ---
        ax2 = axes[1]
        ax2.set_facecolor('#FAFAFA')
        metrics = ['Rev CAGR\n3y', 'ROIC', 'ROE']
        vals_a2 = [
            _pct_val(m_a.get('revenue_cagr_3y')),
            _pct_val(m_a.get('roic')),
            _pct_val(m_a.get('roe')),
        ]
        vals_b2 = [
            _pct_val(m_b.get('revenue_cagr_3y')),
            _pct_val(m_b.get('roic')),
            _pct_val(m_b.get('roe')),
        ]
        x2 = np.arange(len(metrics))
        bars_a2 = ax2.bar(x2 - w/2, vals_a2, w, label=tkr_a, color=ca, alpha=0.85)
        bars_b2 = ax2.bar(x2 + w/2, vals_b2, w, label=tkr_b, color=cb, alpha=0.85)
        for bar in list(bars_a2) + list(bars_b2):
            h = bar.get_height()
            if h and h > 0:
                ax2.text(bar.get_x() + bar.get_width()/2., h + 0.3,
                         f'{h:.0f}%', ha='center', va='bottom', fontsize=7, color='#333')
        _roic_a = _pct_val(m_a.get('roic'))
        _roic_b = _pct_val(m_b.get('roic'))
        _r_leader = tkr_a if _roic_a > _roic_b else (tkr_b if _roic_b > _roic_a else None)
        _r_title = (f"ROIC : {tkr_a} {_roic_a:.0f}% vs {tkr_b} {_roic_b:.0f}% — {_r_leader} plus efficient"
                    if _r_leader else f"ROIC équivalent : {_roic_a:.0f}%")
        ax2.set_title(_r_title, fontsize=8.5, fontweight='bold', color='#1B3A6B', pad=4)
        ax2.set_xticks(x2); ax2.set_xticklabels(metrics, fontsize=8)
        ax2.set_ylabel('%', fontsize=8); ax2.legend(fontsize=7.5)
        ax2.spines['top'].set_visible(False); ax2.spines['right'].set_visible(False)
        ax2.tick_params(labelsize=8)

        fig.tight_layout(pad=1.2)
        return _make_chart_buf(fig)
    except Exception as e:
        log.warning(f"[cmp_pptx] chart_revenue_ebitda error: {e}")
        return None


def _chart_multiples(m_a: dict, m_b: dict, tkr_a: str, tkr_b: str) -> Optional[io.BytesIO]:
    """Graphique comparaison multiples de valorisation."""
    try:
        import matplotlib.pyplot as plt
        import matplotlib
        matplotlib.use('Agg')
        import numpy as np

        fig, ax = plt.subplots(figsize=(11.5, 4.2))
        fig.patch.set_facecolor('white')
        ax.set_facecolor('#FAFAFA')

        ca = '#1B3A6B'
        cb = '#C9A227'
        c_sector = '#BBBBBB'

        metrics_lbl = ['PE', 'EV/EBITDA', 'P/B', 'PEG']
        sector = m_a.get('sector_a', '') or ''
        sec_pe = m_a.get('sector_median_pe') or 20.0
        sec_eveb = m_a.get('sector_median_ev_ebitda') or 14.0

        vals_a = [
            _safe_float(m_a.get('pe_ratio')) or 0,
            _safe_float(m_a.get('ev_ebitda')) or 0,
            _safe_float(m_a.get('price_to_book')) or 0,
            _safe_float(m_a.get('peg_ratio')) or 0,
        ]
        vals_b = [
            _safe_float(m_b.get('pe_ratio')) or 0,
            _safe_float(m_b.get('ev_ebitda')) or 0,
            _safe_float(m_b.get('price_to_book')) or 0,
            _safe_float(m_b.get('peg_ratio')) or 0,
        ]
        sector_refs = [sec_pe, sec_eveb, None, None]

        x = np.arange(len(metrics_lbl))
        w = 0.28
        ax.bar(x - w,   vals_a, w, label=tkr_a, color=ca, alpha=0.85)
        ax.bar(x,       vals_b, w, label=tkr_b, color=cb, alpha=0.85)
        # Médiane sectorielle
        s_vals = [v if v is not None else 0 for v in sector_refs]
        ax.bar(x + w, s_vals, w, label='Med. secteur', color=c_sector, alpha=0.7)

        # Labels valeurs (virgule décimale FR)
        for i, (va, vb) in enumerate(zip(vals_a, vals_b)):
            if va is not None and va > 0:
                ax.text(x[i] - w, va + 0.3, f'{va:.1f}'.replace('.', ','), ha='center', va='bottom', fontsize=10, color='#1B3A6B')
            if vb is not None and vb > 0:
                ax.text(x[i], vb + 0.3, f'{vb:.1f}'.replace('.', ','), ha='center', va='bottom', fontsize=10, color='#6B5010')

        ax.set_xticks(x); ax.set_xticklabels(metrics_lbl, fontsize=11)
        _pe_a = _safe_float(m_a.get('pe_ratio')) or 0
        _pe_b = _safe_float(m_b.get('pe_ratio')) or 0
        if _pe_a and _pe_b:
            _cheap_tkr = tkr_a if _pe_a < _pe_b else tkr_b
            _ev_a = _safe_float(m_a.get('ev_ebitda')) or 0
            _ev_b = _safe_float(m_b.get('ev_ebitda')) or 0
            _m_title = f"{_cheap_tkr} moins cher : PE {_pe_a:.0f}x vs {_pe_b:.0f}x  |  EV/EBITDA {_ev_a:.0f}x vs {_ev_b:.0f}x"
        else:
            _m_title = 'Multiples de Valorisation'
        ax.set_title(_m_title, fontsize=12, fontweight='bold', color='#1B3A6B', pad=8)
        ax.legend(fontsize=11, frameon=False, loc='upper right')
        ax.set_ylabel('(x)', fontsize=10, color='#555')
        ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
        ax.tick_params(labelsize=10)
        ax.yaxis.grid(True, alpha=0.25, color='#D0D5DD'); ax.set_axisbelow(True)
        fig.tight_layout(pad=1.0)
        return _make_chart_buf(fig)
    except Exception as e:
        log.warning(f"[cmp_pptx] chart_multiples error: {e}")
        return None


def _chart_finsight_score(m_a: dict, m_b: dict, tkr_a: str, tkr_b: str) -> Optional[io.BytesIO]:
    """Graphique FinSight Score comparatif — barres horizontales."""
    try:
        import matplotlib.pyplot as plt
        import matplotlib
        matplotlib.use('Agg')
        import numpy as np

        fig, ax = plt.subplots(figsize=(8.5, 4.2))
        fig.patch.set_facecolor('white')
        ax.set_facecolor('#FAFAFA')

        fs_a = _safe_float(m_a.get('finsight_score')) or 0
        fs_b = _safe_float(m_b.get('finsight_score')) or 0

        labels = [f'{tkr_a}', f'{tkr_b}']
        values = [fs_a, fs_b]
        colors_bar = ['#2E5FA3', '#2E8B57']

        bars = ax.barh(labels, values, color=colors_bar, height=0.55, alpha=0.9)
        for bar, val in zip(bars, values):
            ax.text(val + 1.0, bar.get_y() + bar.get_height()/2.,
                    f'{val:.0f}/100', ha='left', va='center', fontsize=14,
                    fontweight='bold', color='#333')
        ax.set_xlim(0, 115)
        ax.set_xlabel('Score /100', fontsize=12)
        _fs_winner = tkr_a if fs_a > fs_b else (tkr_b if fs_b > fs_a else None)
        _fs_delta = abs(int(fs_a) - int(fs_b))
        _pt_lbl = "pt" if _fs_delta == 1 else "pts"
        _fs_title = (f"FinSight Score : {_fs_winner} favori ({_fs_delta} {_pt_lbl} d'avance)"
                     if _fs_winner else "FinSight Score : égalité")
        ax.set_title(_fs_title, fontsize=14, fontweight='bold', color='#1B3A6B', pad=14)
        ax.axvline(50, color='#AAAAAA', linestyle='--', linewidth=0.9)
        ax.text(50.8, -0.55, 'Neutre', fontsize=9, color='#888')
        ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
        ax.tick_params(labelsize=12)
        fig.tight_layout(pad=1.2)
        return _make_chart_buf(fig)
    except Exception as e:
        log.warning(f"[cmp_pptx] chart_finsight_score error: {e}")
        return None


def _chart_risk_profile(m_a: dict, m_b: dict, tkr_a: str, tkr_b: str) -> Optional[io.BytesIO]:
    """Radar chart profil de risque."""
    try:
        import matplotlib.pyplot as plt
        import matplotlib
        matplotlib.use('Agg')
        import numpy as np

        categories = ['Beta\n(vol)', 'Levier\n(ND/EBITDA)', 'Qualité\n(Piotros.)', 'Momentum\n(3m)', 'FCF\nYield']
        N = len(categories)

        def _norm(v, lo, hi):
            if v is None: return 0.5
            try:
                return max(0.1, min(1.0, (float(v) - lo) / (hi - lo)))
            except Exception: return 0.5

        # Scores normalisés sur [0, 1] (plus haut = meilleur)
        scores_a = [
            1 - _norm(m_a.get('beta') or 1.0, 0.3, 2.5),            # beta: moins = mieux
            1 - _norm(m_a.get('net_debt_ebitda') or 2.0, -1, 6),     # levier: moins = mieux
            _norm(m_a.get('piotroski_score') or 5, 1, 9),             # piotroski: plus = mieux
            _norm((m_a.get('perf_3m') or 0) * 100, -20, 30),         # momentum: plus = mieux
            _norm((m_a.get('fcf_yield') or 0) * 100, 0, 10),         # FCF yield: plus = mieux
        ]
        scores_b = [
            1 - _norm(m_b.get('beta') or 1.0, 0.3, 2.5),
            1 - _norm(m_b.get('net_debt_ebitda') or 2.0, -1, 6),
            _norm(m_b.get('piotroski_score') or 5, 1, 9),
            _norm((m_b.get('perf_3m') or 0) * 100, -20, 30),
            _norm((m_b.get('fcf_yield') or 0) * 100, 0, 10),
        ]

        angles = [n / float(N) * 2 * np.pi for n in range(N)]
        angles += angles[:1]
        scores_a += scores_a[:1]
        scores_b += scores_b[:1]

        fig, ax = plt.subplots(figsize=(6.4, 5.6), subplot_kw=dict(polar=True))
        fig.patch.set_facecolor('white')
        ax.set_facecolor('#FAFAFA')

        ax.plot(angles, scores_a, 'o-', linewidth=2.2, color='#2E5FA3', label=tkr_a)
        ax.fill(angles, scores_a, alpha=0.18, color='#2E5FA3')
        ax.plot(angles, scores_b, 'o-', linewidth=2.2, color='#2E8B57', label=tkr_b)
        ax.fill(angles, scores_b, alpha=0.18, color='#2E8B57')

        ax.set_xticks(angles[:-1])
        ax.set_xticklabels(categories, fontsize=10)
        ax.set_ylim(0, 1)
        ax.set_yticks([0.25, 0.5, 0.75])
        ax.set_yticklabels([], fontsize=0)
        ax.legend(loc='upper right', bbox_to_anchor=(1.25, 1.12), fontsize=11)
        ax.set_title('Profil de Risque', fontsize=14, fontweight='bold', color='#1B3A6B', pad=18)
        ax.grid(color='#CCCCCC', linestyle='--', linewidth=0.5)
        fig.tight_layout(pad=0.5)
        return _make_chart_buf(fig)
    except Exception as e:
        log.warning(f"[cmp_pptx] chart_risk_profile error: {e}")
        return None


def _chart_ebitda_margins(m_a: dict, m_b: dict, tkr_a: str, tkr_b: str) -> Optional[io.BytesIO]:
    """Graphique marges EBITDA LTM / N-1 / N-2. Tailles optimisees pour slide PPTX."""
    try:
        import matplotlib.pyplot as plt
        import matplotlib
        matplotlib.use('Agg')
        import numpy as np

        # Calcul proportions : zone PPTX 11.7x4.8cm -> figsize ~6.5x2.7 inch pour densite correcte
        fig, ax = plt.subplots(figsize=(6.8, 3.2))
        fig.patch.set_facecolor('white')
        ax.set_facecolor('#FAFAFA')
        ca = '#1B3A6B'; cb = '#C9A227'
        years_lbl = ['N-2', 'N-1', 'LTM']
        vals_a = [_pct_val(m_a.get('ebitda_margin_y2')), _pct_val(m_a.get('ebitda_margin_y1')), _pct_val(m_a.get('ebitda_margin_ltm'))]
        vals_b = [_pct_val(m_b.get('ebitda_margin_y2')), _pct_val(m_b.get('ebitda_margin_y1')), _pct_val(m_b.get('ebitda_margin_ltm'))]
        x = np.arange(len(years_lbl)); w = 0.35
        bars_a = ax.bar(x - w/2, vals_a, w, label=tkr_a, color=ca, alpha=0.90)
        bars_b = ax.bar(x + w/2, vals_b, w, label=tkr_b, color=cb, alpha=0.90)
        for bar in list(bars_a) + list(bars_b):
            h = bar.get_height()
            if h and h > 0:
                ax.text(bar.get_x() + bar.get_width()/2., h + 0.5, f'{h:.0f}%',
                        ha='center', va='bottom', fontsize=11, color='#333')
        # Titre analytique JPM : tendance + delta
        _trend_a = vals_a[-1] - vals_a[0] if len(vals_a) >= 2 else 0
        _trend_b = vals_b[-1] - vals_b[0] if len(vals_b) >= 2 else 0
        if abs(_trend_a) > abs(_trend_b):
            _chart_title = f"{tkr_a} {'élargit' if _trend_a > 0 else 'comprime'} ses marges ({_trend_a:+.0f}pts sur 3 ans)"
        elif abs(_trend_b) > 0.5:
            _chart_title = f"{tkr_b} {'élargit' if _trend_b > 0 else 'comprime'} ses marges ({_trend_b:+.0f}pts sur 3 ans)"
        else:
            _chart_title = 'Marges EBITDA — Trajectoire 3 ans'
        ax.set_title(_chart_title, fontsize=11, fontweight='bold', color='#1B3A6B', pad=6)
        ax.set_xticks(x); ax.set_xticklabels(years_lbl, fontsize=11)
        ax.set_ylabel('%', fontsize=10); ax.legend(fontsize=10, loc='upper left', frameon=False)
        ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
        ax.tick_params(labelsize=10)
        ax.yaxis.grid(True, alpha=0.25, color='#D0D5DD'); ax.set_axisbelow(True)
        fig.tight_layout(pad=0.8)
        return _make_chart_buf(fig)
    except Exception as e:
        log.warning(f"[cmp_pptx] chart_ebitda_margins error: {e}")
        return None


def _chart_growth_returns(m_a: dict, m_b: dict, tkr_a: str, tkr_b: str) -> Optional[io.BytesIO]:
    """Graphique Rev CAGR + ROIC + ROE."""
    try:
        import matplotlib.pyplot as plt
        import matplotlib
        matplotlib.use('Agg')
        import numpy as np

        fig, ax = plt.subplots(figsize=(6.8, 3.2))
        fig.patch.set_facecolor('white')
        ax.set_facecolor('#FAFAFA')
        ca = '#1B3A6B'; cb = '#C9A227'
        metrics = ['Rev CAGR 3y', 'ROIC', 'ROE']
        vals_a = [_pct_val(m_a.get('revenue_cagr_3y')), _pct_val(m_a.get('roic')), _pct_val(m_a.get('roe'))]
        vals_b = [_pct_val(m_b.get('revenue_cagr_3y')), _pct_val(m_b.get('roic')), _pct_val(m_b.get('roe'))]
        x = np.arange(len(metrics)); w = 0.35
        bars_a = ax.bar(x - w/2, vals_a, w, label=tkr_a, color=ca, alpha=0.90)
        bars_b = ax.bar(x + w/2, vals_b, w, label=tkr_b, color=cb, alpha=0.90)
        for bar in list(bars_a) + list(bars_b):
            h = bar.get_height()
            if h and h > 0:
                ax.text(bar.get_x() + bar.get_width()/2., h + 0.5, f'{h:.0f}%',
                        ha='center', va='bottom', fontsize=11, color='#333')
        _roic_a = vals_a[1]; _roic_b = vals_b[1]
        if _roic_a and _roic_b:
            _roic_leader = tkr_a if _roic_a > _roic_b else tkr_b
            _chart_title = f"{_roic_leader} génère plus de valeur : ROIC {max(_roic_a,_roic_b):.0f} % vs {min(_roic_a,_roic_b):.0f} %"
        else:
            _chart_title = 'Croissance & Rentabilité'
        ax.set_title(_chart_title, fontsize=11, fontweight='bold', color='#1B3A6B', pad=6)
        ax.set_xticks(x); ax.set_xticklabels(metrics, fontsize=11)
        ax.set_ylabel('%', fontsize=10); ax.legend(fontsize=10, loc='upper left', frameon=False)
        ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
        ax.tick_params(labelsize=10)
        ax.yaxis.grid(True, alpha=0.25, color='#D0D5DD'); ax.set_axisbelow(True)
        fig.tight_layout(pad=0.8)
        return _make_chart_buf(fig)
    except Exception as e:
        log.warning(f"[cmp_pptx] chart_growth_returns error: {e}")
        return None


def _chart_52w_price(tkr_a: str, tkr_b: str) -> Optional[io.BytesIO]:
    """Graphique cours boursiers 52 semaines normalisés base 100."""
    try:
        import matplotlib.pyplot as plt
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.dates as mdates
        import yfinance as yf
        from datetime import datetime, timedelta

        end = datetime.today()
        start = end - timedelta(days=370)

        df_a = yf.download(tkr_a, start=start, end=end, progress=False, auto_adjust=True)
        df_b = yf.download(tkr_b, start=start, end=end, progress=False, auto_adjust=True)

        if df_a.empty and df_b.empty:
            return None

        fig, ax = plt.subplots(figsize=(10.0, 7.3))
        fig.patch.set_facecolor('white')
        ax.set_facecolor('#FAFAFA')
        ca = '#2E5FA3'; cb = '#2E8B57'

        if not df_a.empty:
            close_a = df_a['Close'].squeeze()
            base_a = close_a.iloc[0]
            if base_a and base_a > 0:
                ax.plot(close_a.index, close_a / base_a * 100,
                        color=ca, linewidth=2.4, label=tkr_a)

        if not df_b.empty:
            close_b = df_b['Close'].squeeze()
            base_b = close_b.iloc[0]
            if base_b and base_b > 0:
                ax.plot(close_b.index, close_b / base_b * 100,
                        color=cb, linewidth=2.4, label=tkr_b)

        ax.axhline(y=100, color='#BBBBBB', linestyle='--', linewidth=0.9, alpha=0.7)
        # Titre analytique : qui sur/sous-performe sur 52 semaines
        try:
            _p52_a = ((close_a.iloc[-1] / close_a.iloc[0]) - 1) * 100 if not df_a.empty else None
            _p52_b = ((close_b.iloc[-1] / close_b.iloc[0]) - 1) * 100 if not df_b.empty else None
            if _p52_a is not None and _p52_b is not None:
                _sur = tkr_a if _p52_a > _p52_b else tkr_b
                _52_title = (
                    f"52 semaines : {tkr_a} {_p52_a:+.1f} % vs {tkr_b} {_p52_b:+.1f} % — {_sur} surperforme"
                ).replace('.', ',')
            elif _p52_a is not None:
                _52_title = f"52 semaines : {tkr_a} {_p52_a:+.1f} %".replace('.', ',')
            elif _p52_b is not None:
                _52_title = f"52 semaines : {tkr_b} {_p52_b:+.1f} %".replace('.', ',')
            else:
                _52_title = 'Performance Relative 52 Semaines (base 100)'
        except Exception:
            _52_title = 'Performance Relative 52 Semaines (base 100)'
        ax.set_title(_52_title, fontsize=15, fontweight='bold', color='#1B3A6B', pad=14)
        ax.set_ylabel('Performance (base 100)', fontsize=12)
        ax.legend(fontsize=13, loc='best', framealpha=0.95)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.tick_params(labelsize=11)
        ax.xaxis.set_major_formatter(mdates.DateFormatter('%b %y'))
        ax.xaxis.set_major_locator(mdates.MonthLocator(interval=2))
        plt.xticks(rotation=25, ha='right')
        fig.tight_layout(pad=1.4)
        return _make_chart_buf(fig)
    except Exception as e:
        log.warning(f"[cmp_pptx] chart_52w_price error: {e}")
        return None


def _pct_val(v) -> float:
    """Convertit une marge en pourcentage pour les graphiques."""
    if v is None: return 0.0
    try:
        fv = float(v)
        if abs(fv) > 2.0: fv /= 100.0
        return round(fv * 100, 1)
    except Exception: return 0.0


def _safe_float(v) -> Optional[float]:
    if v is None: return None
    try: return float(v)
    except Exception: return None


def _insert_chart(slide, buf: Optional[io.BytesIO], x: float, y: float, w_cm: float, h_cm: float):
    """Insere un graphique BytesIO dans une slide."""
    if buf is None:
        return
    try:
        from pptx.util import Cm
        slide.shapes.add_picture(buf, Cm(x), Cm(y), Cm(w_cm), Cm(h_cm))
    except Exception as e:
        log.warning(f"[cmp_pptx] insert_chart error: {e}")


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def _slide_cover(prs, m_a: dict, m_b: dict):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    tkr_a  = m_a.get('ticker_a') or 'A'
    tkr_b  = m_b.get('ticker_b') or 'B'
    name_a = m_a.get('company_name_a') or tkr_a
    name_b = m_b.get('company_name_b') or tkr_b
    gen_date = _fr_date_long()

    rec_a = (m_a.get('recommendation') or 'HOLD').upper()
    rec_b = (m_b.get('recommendation') or 'HOLD').upper()

    def _rec_c(rec):
        if rec == 'BUY': return GREEN, GREEN_PALE
        if rec == 'SELL': return RED, RED_PALE
        return AMBER, AMBER_PALE

    # Top band navy
    add_rect(slide, 0, 0, 25.4, 2.3, NAVY)
    add_text_box(slide, 1.27, 0.46, 22.86, 0.71, "FinSight IA",
                 10, "8899BB", align=PP_ALIGN.CENTER)
    add_rect(slide, 11.05, 1.32, 3.3, 0.05, WHITE)
    add_text_box(slide, 1.27, 1.52, 22.86, 0.81,
                 "Pitchbook Comparatif  \u2014  Analyse d'Investissement",
                 11, "CCDDEE", align=PP_ALIGN.CENTER)

    # Titre "A vs B"
    title_str = f"{tkr_a}  vs  {tkr_b}"
    add_text_box(slide, 1.27, 4.06, 22.86, 1.8,
                 title_str, 38, NAVY, bold=True, align=PP_ALIGN.CENTER)

    # Noms complets (Name + Ticker)
    add_text_box(slide, 1.27, 5.9, 11.0, 0.71,
                 _truncate(f"{name_a} ({tkr_a})", 44), 11, "888888", align=PP_ALIGN.CENTER)
    add_text_box(slide, 13.13, 5.9, 11.0, 0.71,
                 _truncate(f"{name_b} ({tkr_b})", 44), 11, "888888", align=PP_ALIGN.CENTER)


    # Bandeau recommendation A
    ra_c, ra_pal = _rec_c(rec_a)
    price_a  = m_a.get('share_price')
    tbase_a  = m_a.get('dcf_base')
    up_a     = _upside(tbase_a, price_a)
    cur_a    = "EUR" if (m_a.get('currency_a') or 'USD') == 'EUR' else '$'
    add_rect(slide, 1.27, 7.0, 11.0, 1.1, ra_pal)
    add_rect(slide, 1.27, 7.0, 0.18, 1.1, ra_c)
    add_text_box(slide, 1.6, 7.05, 10.5, 1.0,
                 f"\u25cf {rec_a}  \u00b7  Cible : {_fr(tbase_a, 0)} {cur_a}  \u00b7  Upside : {up_a}",
                 9, NAVY, bold=True, align=PP_ALIGN.CENTER)

    # Bandeau recommendation B
    rb_c, rb_pal = _rec_c(rec_b)
    price_b  = m_b.get('share_price')
    tbase_b  = m_b.get('dcf_base')
    up_b     = _upside(tbase_b, price_b)
    cur_b    = "EUR" if (m_b.get('currency_b') or 'USD') == 'EUR' else '$'
    add_rect(slide, 13.13, 7.0, 11.0, 1.1, rb_pal)
    add_rect(slide, 13.13, 7.0, 0.18, 1.1, rb_c)
    add_text_box(slide, 13.46, 7.05, 10.5, 1.0,
                 f"\u25cf {rec_b}  \u00b7  Cible : {_fr(tbase_b, 0)} {cur_b}  \u00b7  Upside : {up_b}",
                 9, NAVY, bold=True, align=PP_ALIGN.CENTER)

    # Profil rapide A / B dans la zone vide (y=8.5 → 11.5)
    sec_a   = m_a.get('sector_a') or '\u2014'
    sec_b   = m_b.get('sector_b') or '\u2014'
    fs_a    = str(m_a.get('finsight_score') or '\u2014') + "/100"
    fs_b    = str(m_b.get('finsight_score') or '\u2014') + "/100"
    _pa_str = (f"{_fr(price_a, 1)} {cur_a}") if price_a else '\u2014'
    _pb_str = (f"{_fr(price_b, 1)} {cur_b}") if price_b else '\u2014'
    mc_a_s  = (_frm(m_a.get('market_cap'), cur_a)) if m_a.get('market_cap') else '\u2014'
    mc_b_s  = (_frm(m_b.get('market_cap'), cur_b)) if m_b.get('market_cap') else '\u2014'

    # Blocs profil side-by-side
    for _xi, _m, _nm, _sec, _fs, _cours, _mc in [
        (1.27,  m_a, tkr_a, sec_a, fs_a, _pa_str, mc_a_s),
        (13.13, m_b, tkr_b, sec_b, fs_b, _pb_str, mc_b_s),
    ]:
        add_rect(slide, _xi, 8.6, 11.0, 3.2, GREY_BG)
        _lines = [
            ("Secteur",          _sec),
            ("Cours actuel",     _cours),
            ("Market Cap",       _mc),
            ("Score FinSight",   _fs),
        ]
        _ly = 8.85
        for _lbl, _val in _lines:
            add_text_box(slide, _xi + 0.25, _ly, 5.0, 0.55, _lbl, 7.5, GREY_TXT)
            add_text_box(slide, _xi + 5.2, _ly, 5.5, 0.55, _val, 7.5, NAVY,
                         bold=True, align=PP_ALIGN.RIGHT)
            _ly += 0.65

    # ── FRED macro one-liner ────────────────────────────────────────────
    try:
        from data.sources.fred_source import fetch_macro_context
        _fred = fetch_macro_context()
        if _fred:
            _parts = []
            _map = [
                ("fed_funds_rate", "Fed", " %", 2),
                ("treasury_10y",   "10Y", " %", 2),
                ("vix",            "VIX", "",   1),
                ("cpi_yoy",        "CPI", " %", 1),
                ("unemployment",   "Chômage", " %", 1),
            ]
            for _key, _lbl, _suf, _dp in _map:
                _v = _fred.get(_key)
                if _v is not None:
                    try:
                        _vs = f"{float(_v):.{_dp}f}".replace(".", ",") + _suf
                        if _key == "cpi_yoy":
                            _vs = "+" + _vs if float(_v) >= 0 else _vs
                        _parts.append(f"{_lbl} {_vs}")
                    except Exception:
                        pass
            if _parts:
                _fred_line = "Données FRED : " + "  \u00b7  ".join(_parts)
                add_text_box(slide, 1.02, 12.05, 23.37, 0.35, _fred_line,
                             6.5, GREY_LIGHT, align=PP_ALIGN.CENTER)
    except Exception as _e:
        log.debug(f"[cmp_pptx] FRED macro one-liner skip: {_e}")

    # Bottom rule + date
    add_rect(slide, 1.02, 12.4, 23.37, 0.03, "AAAAAA")
    add_text_box(slide, 1.02, 12.65, 11.43, 0.56, "Rapport confidentiel", 8, GREY_TXT)
    add_text_box(slide, 12.95, 12.65, 11.43, 0.56, gen_date, 8, GREY_TXT,
                 align=PP_ALIGN.RIGHT)
    return slide


def _slide_exec_summary(prs, m_a: dict, m_b: dict, synthesis: dict):
    """Slide 2 — Executive Summary comparatif (redesign JPM sans KPI cours A/B).

    Layout :
      - Titre + sous-titre analytique (JPM style)
      - Bande LLM exec summary large (hauteur 2.5)
      - Tableau metriques comparatif 7 lignes avec surlignage gagnant/perdant
      - Bandeau verdict bas (choix prefere + ecart de score)
    """
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    tkr_a = m_a.get('ticker_a') or 'A'
    tkr_b = m_b.get('ticker_b') or 'B'
    # Winner refondu (C4) : source = LLM via synthesis.llm_choice, fallback finsight.
    # Si LLM neutral, winner = None et on affiche "Pas de préférence forte".
    _winner_src = m_a.get('winner_source', 'finsight_fallback')
    winner = m_a.get('winner')
    if _winner_src == 'llm_neutral' or winner is None:
        winner_display = None  # -> "Pas de préférence forte"
        winner = tkr_a  # fallback pour les visuels internes qui nécessitent un ticker
    else:
        winner_display = winner

    # Titre analytique dynamique : "TKR privilegie : valorisation relative + qualité bilancielle"
    _pe_a = _safe_float(m_a.get('pe_ratio')) or 0
    _pe_b = _safe_float(m_b.get('pe_ratio')) or 0
    _pe_leader = tkr_a if _pe_a and _pe_b and _pe_a < _pe_b else (tkr_b if _pe_a and _pe_b else None)
    _mg_a = _safe_float(m_a.get('ebitda_margin_ltm')) or 0
    _mg_b = _safe_float(m_b.get('ebitda_margin_ltm')) or 0
    _mg_leader = tkr_a if _mg_a > _mg_b else tkr_b
    if _pe_leader and _pe_leader == winner:
        _jpm_sub = f"{winner} combine valorisation attractive et fondamentaux solides"
    elif _mg_leader == winner:
        _jpm_sub = f"{winner} prime justifiée par une profitabilité supérieure"
    else:
        _jpm_sub = f"{winner} privilegie sur la qualité/valorisation relative"

    navy_bar(slide)
    footer_bar(slide)
    slide_title(slide, _cplbl("exec_cmp"))
    section_dots(slide, 1)
    _company_header_band(slide, tkr_a, tkr_b, m_a=m_a, m_b=m_b)

    # Bande LLM exec summary : titre JPM INTÈGRE au sommet du cadre + texte en dessous
    exec_txt = synthesis.get('exec_summary') or ""
    _box_y = 2.45
    _box_h = 3.20
    add_rect(slide, 1.02, _box_y, 23.37, _box_h, NAVY_PALE)
    add_rect(slide, 1.02, _box_y, 0.13, _box_h, NAVY_MID)
    # Titre JPM en haut du cadre (bold, navy fonce, italic)
    add_text_box(slide, 1.4, _box_y + 0.12, 22.8, 0.50,
                 _truncate(_jpm_sub, 95), 11, NAVY, bold=True, italic=True)
    # Texte LLM exec summary juste en dessous
    if exec_txt:
        _exec_clean = " ".join(_fit(exec_txt, 700).split())
        add_text_box(slide, 1.4, _box_y + 0.75, 22.8, _box_h - 0.9, _exec_clean,
                     9, NAVY, wrap=True)

    # Tableau comparatif avec highlight gagnant (fleches vertes)
    def _fr_nd(v) -> str:
        if v is None: return "\u2014"
        try:
            fv = round(float(v), 1)
            fv = 0.0 if abs(fv) < 0.05 else fv
            return f"{fv:.1f}".replace(".", ",") + "x"
        except Exception:
            return "\u2014"

    def _mark(val_a, val_b, higher_better=True):
        """Retourne (str_a, str_b) avec coche pour le meilleur."""
        try:
            fa = _safe_float(val_a)
            fb = _safe_float(val_b)
        except Exception: fa, fb = None, None
        if fa is None or fb is None:
            return (None, None)
        if higher_better:
            return (fa >= fb, fb > fa)
        return (fa <= fb, fb < fa)

    rows_raw = [
        ("FinSight Score /100",   m_a.get('finsight_score'),        m_b.get('finsight_score'),        True),
        ("Recommandation",        m_a.get('recommendation'),        m_b.get('recommendation'),        None),
        ("Conviction",            m_a.get('conviction'),            m_b.get('conviction'),            True),
        ("Marge EBITDA LTM",      m_a.get('ebitda_margin_ltm'),     m_b.get('ebitda_margin_ltm'),     True),
        ("ROIC",                  m_a.get('roic'),                  m_b.get('roic'),                  True),
        ("Rev CAGR 3y",           m_a.get('revenue_cagr_3y'),       m_b.get('revenue_cagr_3y'),       True),
        ("P/E (LTM)",             m_a.get('pe_ratio'),              m_b.get('pe_ratio'),              False),
        ("EV/EBITDA",             m_a.get('ev_ebitda'),             m_b.get('ev_ebitda'),             False),
        ("ND/EBITDA",             m_a.get('net_debt_ebitda'),       m_b.get('net_debt_ebitda'),       False),
    ]
    # Formatters per row
    def _fmt(lbl, v):
        if v is None: return "\u2014"
        if lbl == "FinSight Score /100":
            try: return str(int(float(v)))
            except Exception: return "\u2014"
        if lbl in ("Marge EBITDA LTM", "ROIC", "Rev CAGR 3y"):
            return _frpct(v)
        if lbl == "Conviction":
            return _frpct(v)
        if lbl in ("P/E (LTM)", "EV/EBITDA", "ND/EBITDA"):
            return _fr_nd(v) if lbl == "ND/EBITDA" else _frx(v)
        return str(v)

    rows = []
    for lbl, va, vb, hib in rows_raw:
        sa = _fmt(lbl, va)
        sb = _fmt(lbl, vb)
        if hib is not None:
            wa, wb = _mark(va, vb, hib)
            if wa: sa = "\u25b2 " + sa
            elif wb: sb = "\u25b2 " + sb
        rows.append((lbl, sa, sb))

    y_tbl = 5.85
    add_table(
        slide, 1.02, y_tbl, 23.37, 5.45,
        num_rows=len(rows), num_cols=3,
        col_widths_pct=[0.40, 0.30, 0.30],
        header_data=["Indicateur clé", tkr_a, tkr_b],
        rows_data=rows,
        header_fill=NAVY, border_hex="DDDDDD"
    )

    # Bandeau verdict bas : choix LLM + indication FinSight (refonte C4)
    y_verd = 11.65
    if winner_display is None:
        # Cas neutre : fond ambre, pas de ticker gagnant
        _verd_fill = "FFF3E0"  # ambre clair
        _verd_accent = "B06000"  # ambre fonce
        _verd_title = "Choix préféré : pas de préférence forte"
    else:
        _is_a = (winner == tkr_a)
        _verd_fill = COLOR_A_PAL if _is_a else COLOR_B_PAL
        _verd_accent = COLOR_A if _is_a else COLOR_B
        _verd_title = f"Choix préféré : {winner_display}"

    add_rect(slide, 1.02, y_verd, 23.37, 1.05, _verd_fill)
    add_rect(slide, 1.02, y_verd, 0.18, 1.05, _verd_accent)

    _fs_a = m_a.get('finsight_score') or 0
    _fs_b = m_b.get('finsight_score') or 0

    # Ligne 1 : titre choix préféré (LLM)
    add_text_box(slide, 1.35, y_verd + 0.08, 22.8, 0.35,
                 _verd_title, 10, NAVY, bold=True)
    # Ligne 2 : indication FinSight avec les 2 scores
    _verd_sub = (
        f"Indication FinSight (rétrospective, non décisionnelle) : "
        f"{tkr_a} {int(_fs_a)}/100  ·  {tkr_b} {int(_fs_b)}/100   "
        f"— score composite V/G/Q/M percentile rank, ne tient pas compte des news ni de la macro"
    )
    add_text_box(slide, 1.35, y_verd + 0.44, 22.8, 0.55, _verd_sub, 7, GREY_TXT, wrap=True)

    return slide


def _slide_sommaire(prs, tkr_a: str, tkr_b: str):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    navy_bar(slide)
    footer_bar(slide)
    slide_title(slide, _cplbl("sommaire"))

    sections = [
        ("01", "Profil & Identité", "Présentation comparative des deux sociétés",          5),
        ("02", "Performance Financière", "P&L, marges, bilan et liquidité comparés",        7),
        ("03", "Valorisation", "Multiples, DCF, GBM Monte Carlo et football field",         11),
        ("04", "Qualité & Risques", "Piotroski, levier, beta, VaR et momentum",            15),
        ("05", "Verdict", "FinSight Score, Thèses d'Investissement et verdict final",      18),
    ]

    fills = [WHITE, GREY_BG, WHITE, GREY_BG, WHITE]
    ys    = [2.08, 4.01, 5.94, 7.87, 9.80]

    for i, (num, title, desc, page_n) in enumerate(sections):
        y    = ys[i]
        fill = fills[i]
        add_rect(slide, 1.02, y, 23.37, 1.73, fill)
        add_rect(slide, 1.02, y, 0.13, 1.73, NAVY_MID)
        add_text_box(slide, 1.40, y + 0.15, 1.32, 1.43,
                     num, 17, GREY_LIGHT, bold=True)
        add_text_box(slide, 3.20, y + 0.20, 17.00, 0.71,
                     title, 13, NAVY, bold=True)
        add_text_box(slide, 3.20, y + 0.85, 17.00, 0.71,
                     desc, 8.5, GREY_TXT)
        add_text_box(slide, 20.57, y + 0.45, 3.82, 0.71,
                     str(page_n), 9, NAVY_MID,
                     align=PP_ALIGN.RIGHT)

    return slide


def _slide_profil(prs, m_a: dict, m_b: dict):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    tkr_a = m_a.get('ticker_a') or 'A'
    tkr_b = m_b.get('ticker_b') or 'B'

    navy_bar(slide)
    footer_bar(slide)
    slide_title(slide, _cplbl("profil_cmp"))
    section_dots(slide, 1)
    # Pas de company_header_band — le nom complet est affiche directement dans chaque colonne

    def _profile_col(slide, m, side, tkr):
        if side == 'left':
            x0 = 1.02; lbl_color = NAVY_MID
        else:
            x0 = 12.76; lbl_color = GREEN_MID
        w = 11.0

        # Identite
        name = m.get('company_name_a' if side == 'left' else 'company_name_b') or tkr
        sector = m.get('sector_a' if side == 'left' else 'sector_b') or '\u2014'
        currency = m.get('currency_a' if side == 'left' else 'currency_b') or 'USD'
        cur_sym  = 'EUR' if currency == 'EUR' else '$'

        # En-tête nom complet (Name + Ticker) — seule ligne d'identification
        add_rect(slide, x0, 1.80, w, 0.65, COLOR_A_PAL if side == 'left' else COLOR_B_PAL)
        add_rect(slide, x0, 1.80, 0.13, 0.65, lbl_color)
        add_text_box(slide, x0 + 0.22, 1.84, w - 0.3, 0.55,
                     _truncate(f"{name} ({tkr})", 38), 8.5, lbl_color, bold=True, wrap=False)

        rows_id = [
            ("Secteur",          _truncate(sector, 22)),
            ("Devise",           currency),
            ("Cours",            _fr(m.get('share_price'), 1) + " " + cur_sym),
            ("52W High",         _fr(m.get('week52_high'), 2) + " " + cur_sym),
            ("52W Low",          _fr(m.get('week52_low'), 2) + " " + cur_sym),
            ("Market Cap",       _fr(m.get('market_cap'), 1) + " Mds " + cur_sym),
            ("Valeur Entrep.",   _fr(m.get('enterprise_value'), 1) + " Mds " + cur_sym),
            ("Div. Yield",       _frpct(m.get('dividend_yield'))),
            ("Perf. 1M",         _frpct(m.get('perf_1m'), signed=True)),
            ("Perf. 3M",         _frpct(m.get('perf_3m'), signed=True)),
            ("Perf. 1Y",         _frpct(m.get('perf_1y'), signed=True)),
            ("Proch. Résultat",  m.get('next_earnings_date') or '\u2014'),
        ]

        y_start = 2.60
        row_h   = 0.84
        for ri, (lbl, val) in enumerate(rows_id):
            y = y_start + ri * row_h
            fill = WHITE if ri % 2 == 0 else GREY_BG
            add_rect(slide, x0, y, w, row_h - 0.04, fill)
            add_text_box(slide, x0 + 0.1, y + 0.12, w * 0.52, row_h - 0.15, lbl,
                         7.5, GREY_TXT, italic=True, wrap=False)
            add_text_box(slide, x0 + w * 0.52, y + 0.12, w * 0.45, row_h - 0.15,
                         val, 7.5, BLACK, bold=False, align=PP_ALIGN.RIGHT, wrap=False)

    _profile_col(slide, m_a, 'left', tkr_a)
    _profile_col(slide, m_b, 'right', tkr_b)
    return slide


def _slide_pl(prs, m_a: dict, m_b: dict, synthesis: dict):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    tkr_a = m_a.get('ticker_a') or 'A'
    tkr_b = m_b.get('ticker_b') or 'B'
    cur_a = 'EUR' if (m_a.get('currency_a') or 'USD') == 'EUR' else '$'
    cur_b = 'EUR' if (m_b.get('currency_b') or 'USD') == 'EUR' else '$'

    navy_bar(slide)
    footer_bar(slide)
    slide_title(slide, _cplbl("is_compare"))
    section_dots(slide, 2)
    _company_header_band(slide, tkr_a, tkr_b, m_a=m_a, m_b=m_b)

    # Titre JPM dynamique basé sur rev cagr
    _cagr_a = _safe_float(m_a.get('revenue_cagr_3y')) or 0
    _cagr_b = _safe_float(m_b.get('revenue_cagr_3y')) or 0
    if abs(_cagr_a) <= 2: _cagr_a *= 100
    if abs(_cagr_b) <= 2: _cagr_b *= 100
    if _cagr_a or _cagr_b:
        _g_leader = tkr_a if _cagr_a > _cagr_b else tkr_b
        _jpm_sub = f"{_g_leader} affiche la meilleure dynamique revenus sur 3 ans ({max(_cagr_a, _cagr_b):.1f} % CAGR)".replace('.', ',')
    else:
        _jpm_sub = "Lecture comparative du Compte de Résultat"
    add_text_box(slide, 1.02, 2.38, 23.37, 0.45, _jpm_sub, 10, NAVY_MID, bold=True, italic=True)

    # Tableau comparatif P&L
    rows = [
        ("Marge EBITDA LTM",  _frpct_margin(m_a.get('ebitda_margin_ltm')),  _frpct_margin(m_b.get('ebitda_margin_ltm'))),
        ("Marge EBITDA N-1",  _frpct_margin(m_a.get('ebitda_margin_y1')),   _frpct_margin(m_b.get('ebitda_margin_y1'))),
        ("Marge EBITDA N-2",  _frpct_margin(m_a.get('ebitda_margin_y2')),   _frpct_margin(m_b.get('ebitda_margin_y2'))),
        ("Tendance EBITDA",   m_a.get('ebitda_margin_trend') or '\u2014',    m_b.get('ebitda_margin_trend') or '\u2014'),
        ("ROIC",              _frpct(m_a.get('roic')),                       _frpct(m_b.get('roic'))),
        ("ROE",               _frpct(m_a.get('roe')),                        _frpct(m_b.get('roe'))),
        ("Rev CAGR 3y",       _frpct(m_a.get('revenue_cagr_3y')),            _frpct(m_b.get('revenue_cagr_3y'))),
        ("Croissance EPS",    _frpct(m_a.get('eps_growth'), signed=True),    _frpct(m_b.get('eps_growth'), signed=True)),
    ]
    add_table(
        slide, 1.02, 2.90, 23.37, 5.2,
        num_rows=len(rows), num_cols=3,
        col_widths_pct=[0.44, 0.28, 0.28],
        header_data=["Indicateur", tkr_a, tkr_b],
        rows_data=rows,
        header_fill=NAVY, border_hex="DDDDDD"
    )

    # Commentaire analytique sous le tableau (pas de graphiques — reserves slide "Rentabilité")
    fin_txt = synthesis.get('financial_text') or ""
    if fin_txt:
        _fc = " ".join(_fit(fin_txt, 700).split())
        add_rect(slide, 1.02, 8.35, 23.37, 4.70, NAVY_PALE)
        add_rect(slide, 1.02, 8.35, 0.13, 4.70, NAVY_MID)
        add_text_box(slide, 1.4, 8.55, 22.8, 4.40, _fc, 9, NAVY, wrap=True)

    return slide


def _slide_marges(prs, m_a: dict, m_b: dict, synthesis: dict):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    tkr_a = m_a.get('ticker_a') or 'A'
    tkr_b = m_b.get('ticker_b') or 'B'

    navy_bar(slide)
    footer_bar(slide)
    slide_title(slide, _cplbl("rent_croiss"))
    section_dots(slide, 2)
    _company_header_band(slide, tkr_a, tkr_b, m_a=m_a, m_b=m_b)

    # Titre analytique JPM
    _mg_a = _safe_float(m_a.get('ebitda_margin_ltm')) or 0
    _mg_b = _safe_float(m_b.get('ebitda_margin_ltm')) or 0
    if abs(_mg_a) <= 2: _mg_a *= 100
    if abs(_mg_b) <= 2: _mg_b *= 100
    if _mg_a and _mg_b:
        _mg_leader = tkr_a if _mg_a > _mg_b else tkr_b
        _jpm_sub = f"{_mg_leader} affiche un avantage structurel de marge ({max(_mg_a,_mg_b):.0f} % vs {min(_mg_a,_mg_b):.0f} %)"
    else:
        _jpm_sub = "Lecture comparative de la rentabilité et croissance"
    add_text_box(slide, 1.02, 2.38, 23.37, 0.45, _jpm_sub, 10, NAVY_MID, bold=True, italic=True)

    # Commentaire financier LLM — bande agrandie pour éviter overflow
    txt = synthesis.get('financial_text') or ""
    if txt:
        _txt_clean = " ".join(_fit(txt, 650).split())
        add_rect(slide, 1.02, 2.90, 23.37, 2.55, NAVY_PALE)
        add_rect(slide, 1.02, 2.90, 0.13, 2.55, NAVY_MID)
        add_text_box(slide, 1.4, 3.00, 22.8, 2.40, _txt_clean, 9, NAVY, wrap=True)

    # Graphiques marges + croissance — plein cadre (blocs KPI retires car
    # redondants avec les charts). Zone disponible : y=5.70 -> 13.20 (7.5cm).
    # figsize 6.8/3.2 = ratio 2.125 ; on respecte pour éviter deformation.
    y_ch = 5.70
    _ch_h = 6.80  # 13.20 - 5.70 - 0.70 footer margin = 6.80cm
    _ch_w = _ch_h * 2.125  # ~14.45cm : trop large pour 2 charts cote a cote
    # Fallback : on cap sur la largeur disponible pour 2 charts
    _max_w_per_chart = (23.37 - 0.4) / 2  # ~11.48cm
    if _ch_w > _max_w_per_chart:
        _ch_w = _max_w_per_chart
        _ch_h = _ch_w / 2.125
    _total_w = _ch_w * 2 + 0.4
    _x0 = max(1.02, (25.4 - _total_w) / 2)
    # Centrer verticalement dans la zone 5.70-13.20
    _y_centered = 5.70 + (7.50 - _ch_h) / 2
    buf_l = _chart_ebitda_margins(m_a, m_b, tkr_a, tkr_b)
    buf_r = _chart_growth_returns(m_a, m_b, tkr_a, tkr_b)
    if buf_l:
        _insert_chart(slide, buf_l, _x0, _y_centered, _ch_w, _ch_h)
    if buf_r:
        _insert_chart(slide, buf_r, _x0 + _ch_w + 0.4, _y_centered, _ch_w, _ch_h)

    return slide


def _slide_bilan(prs, m_a: dict, m_b: dict):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    tkr_a = m_a.get('ticker_a') or 'A'
    tkr_b = m_b.get('ticker_b') or 'B'

    navy_bar(slide)
    footer_bar(slide)
    slide_title(slide, _cplbl("bilan_liq"))
    section_dots(slide, 2)
    _company_header_band(slide, tkr_a, tkr_b, m_a=m_a, m_b=m_b)

    # Titre analytique JPM dynamique basé sur levier
    _nd_a = _safe_float(m_a.get('net_debt_ebitda'))
    _nd_b = _safe_float(m_b.get('net_debt_ebitda'))
    if _nd_a is not None and _nd_b is not None:
        _lev_leader = tkr_a if _nd_a < _nd_b else tkr_b
        _jpm_sub = f"{_lev_leader} préserve une flexibilite bilancielle supérieure"
    else:
        _jpm_sub = "Lecture comparative de la structure bilancielle"
    add_text_box(slide, 1.02, 2.38, 23.37, 0.45,
                 _jpm_sub, 10, NAVY_MID, bold=True, italic=True)

    rows = [
        ("Net Debt / EBITDA", _fr(m_a.get('net_debt_ebitda'), 1) + "x", _fr(m_b.get('net_debt_ebitda'), 1) + "x"),
        ("Interest Coverage",  _frx(m_a.get('interest_coverage')),       _frx(m_b.get('interest_coverage'))),
        ("Current Ratio",      _frx(m_a.get('current_ratio')),           _frx(m_b.get('current_ratio'))),
        ("Quick Ratio",        _frx(m_a.get('quick_ratio')),             _frx(m_b.get('quick_ratio'))),
        ("Capex / Revenue",   _frpct(m_a.get('capex_to_revenue')),       _frpct(m_b.get('capex_to_revenue'))),
        ("Cash Conversion",   _frx(m_a.get('cash_conversion')),          _frx(m_b.get('cash_conversion'))),
        ("FCF Yield",         _frpct(m_a.get('fcf_yield')),              _frpct(m_b.get('fcf_yield'))),
        ("Altman Z-Score",    _fr(m_a.get('altman_z'), 2),               _fr(m_b.get('altman_z'), 2)),
    ]
    add_table(
        slide, 1.02, 2.90, 23.37, 6.5,
        num_rows=len(rows), num_cols=3,
        col_widths_pct=[0.44, 0.28, 0.28],
        header_data=["Indicateur Bilan", tkr_a, tkr_b],
        rows_data=rows,
        header_fill=NAVY, border_hex="DDDDDD"
    )

    # Blocs KPI retires (redondants avec le tableau). Le texte LLM prend la place liberee.

    # --- Bande LLM analytique bilan ---
    def _lev_q(v):
        if v is None: return "indetermine"
        if v < 0:   return "desendette (cash net)"
        if v < 1.5: return "faible"
        if v < 3.0: return "modéré"
        return "élevé"
    def _liq_q(v):
        if v is None: return "non calculable"
        if v >= 2.0: return "confortable"
        if v >= 1.2: return "saine"
        if v >= 1.0: return "tendue"
        return "sous tension"
    _cr_a = _safe_float(m_a.get('current_ratio'))
    _cr_b = _safe_float(m_b.get('current_ratio'))
    _fcfy_a = _safe_float(m_a.get('fcf_yield')) or 0
    _fcfy_b = _safe_float(m_b.get('fcf_yield')) or 0
    # Convertir en %
    if abs(_fcfy_a) <= 2.0: _fcfy_a *= 100
    if abs(_fcfy_b) <= 2.0: _fcfy_b *= 100
    _lev_winner = tkr_a if (_nd_a or 99) < (_nd_b or 99) else tkr_b
    _fcf_winner = tkr_a if _fcfy_a > _fcfy_b else tkr_b
    bilan_txt = (
        f"Structure bilancielle : {tkr_a} affiche un levier {_lev_q(_nd_a)} "
        f"({_fr(_nd_a, 1)}x EBITDA) tandis que {tkr_b} présente un profil {_lev_q(_nd_b)} "
        f"({_fr(_nd_b, 1)}x). La liquidité courante ressort {_liq_q(_cr_a)} pour {tkr_a} "
        f"vs {_liq_q(_cr_b)} pour {tkr_b}. "
        f"Sur le critère du levier, {_lev_winner} dispose de plus de marge pour absorber un choc "
        f"macro ou financer une opération externe. {_fcf_winner} Généré davantage de FCF relatif "
        f"({_fcfy_a:.1f}% vs {_fcfy_b:.1f}% de FCF yield), soutenant la soutenabilite du dividende "
        f"et la capacité de rachat d'actions."
    )
    y_txt = 9.65
    _txt_h = 3.40
    add_rect(slide, 1.02, y_txt, 23.37, _txt_h, NAVY_PALE)
    add_rect(slide, 1.02, y_txt, 0.13, _txt_h, NAVY_MID)
    add_text_box(slide, 1.35, y_txt + 0.12, 22.8, 0.42,
                 "LECTURE BILANCIELLE - Qui dispose du plus de flexibilite financière ?",
                 9.5, NAVY, bold=True)
    add_text_box(slide, 1.35, y_txt + 0.60, 22.8, _txt_h - 0.75,
                 " ".join(_fit(bilan_txt, 900).split()), 9.5, NAVY, wrap=True)

    return slide


def _slide_multiples(prs, m_a: dict, m_b: dict, synthesis: dict):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    tkr_a = m_a.get('ticker_a') or 'A'
    tkr_b = m_b.get('ticker_b') or 'B'

    navy_bar(slide)
    footer_bar(slide)
    slide_title(slide, _cplbl("mult_val"))
    section_dots(slide, 3)
    _company_header_band(slide, tkr_a, tkr_b, m_a=m_a, m_b=m_b)

    # Titre analytique JPM dynamique
    _pe_a = _safe_float(m_a.get('pe_ratio')) or 0
    _pe_b = _safe_float(m_b.get('pe_ratio')) or 0
    if _pe_a and _pe_b:
        _cheap = tkr_a if _pe_a < _pe_b else tkr_b
        _prem = abs(_pe_a - _pe_b) / max(_pe_a, _pe_b) * 100
        _jpm_sub = f"{_cheap} s'échange avec une décote de {_prem:.0f} % vs le peer sur le P/E"
    else:
        _jpm_sub = "Lecture des multiples de valorisation relative"
    add_text_box(slide, 1.02, 2.38, 23.37, 0.45, _jpm_sub, 10, NAVY_MID, bold=True, italic=True)

    # Commentaire valorisation pleine largeur en haut
    txt = synthesis.get('valuation_text') or ""
    if txt:
        _v_clean = " ".join(_fit(txt, 700).split())
        add_rect(slide, 1.02, 2.90, 23.37, 2.6, NAVY_PALE)
        add_rect(slide, 1.02, 2.90, 0.13, 2.6, NAVY_MID)
        add_text_box(slide, 1.4, 3.00, 22.8, 2.45, _v_clean, 9, NAVY, wrap=True)

    sec_pe   = m_a.get('sector_median_pe') or '\u2014'
    sec_eveb = m_a.get('sector_median_ev_ebitda') or '\u2014'

    rows = [
        ("PE Ratio",          _frx(m_a.get('pe_ratio')),      _frx(m_b.get('pe_ratio')),      _frx(sec_pe) if sec_pe != '\u2014' else '\u2014'),
        ("EV / EBITDA",       _frx(m_a.get('ev_ebitda')),     _frx(m_b.get('ev_ebitda')),     _frx(sec_eveb) if sec_eveb != '\u2014' else '\u2014'),
        ("Price / Book",      _frx(m_a.get('price_to_book')), _frx(m_b.get('price_to_book')), "\u2014"),
        ("PEG Ratio",         _frx(m_a.get('peg_ratio')),     _frx(m_b.get('peg_ratio')),     "\u2014"),
        ("FCF Yield",         _frpct(m_a.get('fcf_yield')),   _frpct(m_b.get('fcf_yield')),   "\u2014"),
        ("Div. Yield",        _frpct(m_a.get('dividend_yield')), _frpct(m_b.get('dividend_yield')), "\u2014"),
    ]
    # Tableau réduit a GAUCHE (40% largeur) pour laisser place au graphique a droite
    _tbl_y = 5.80
    _tbl_h = 7.0
    add_table(
        slide, 1.02, _tbl_y, 10.4, _tbl_h,
        num_rows=len(rows), num_cols=4,
        col_widths_pct=[0.34, 0.22, 0.22, 0.22],
        header_data=["Multiple", tkr_a, tkr_b, "Med. Sec."],
        rows_data=rows,
        header_fill=NAVY, border_hex="DDDDDD"
    )

    # Graphique multiples a DROITE, agrandi — ratio matplotlib 11.5/4.2 = 2.74
    buf = _chart_multiples(m_a, m_b, tkr_a, tkr_b)
    if buf:
        _chart_x = 11.80
        _chart_y = 5.80
        _chart_w = 12.60
        _chart_h = _chart_w / 2.74  # ~4.6cm
        # Si trop haut, cap a la place disponible
        if _chart_h > 7.0:
            _chart_h = 7.0
            _chart_w = _chart_h * 2.74
        # Centrer verticalement dans l'espace 5.80-12.80
        _y_c = 5.80 + (7.0 - _chart_h) / 2
        _insert_chart(slide, buf, _chart_x, _y_c, _chart_w, _chart_h)

    return slide


def _slide_dcf(prs, m_a: dict, m_b: dict):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    tkr_a = m_a.get('ticker_a') or 'A'
    tkr_b = m_b.get('ticker_b') or 'B'

    navy_bar(slide)
    footer_bar(slide)
    slide_title(slide, _cplbl("dcf_cibles"))
    section_dots(slide, 3)
    _company_header_band(slide, tkr_a, tkr_b, m_a=m_a, m_b=m_b)

    # Titre JPM dynamique (DCF upside base)
    cur_a = 'EUR' if (m_a.get('currency_a') or 'USD') == 'EUR' else '$'
    cur_b = 'EUR' if (m_b.get('currency_b') or 'USD') == 'EUR' else '$'
    _dbase_a = _safe_float(m_a.get('dcf_base')) or 0
    _dbase_b = _safe_float(m_b.get('dcf_base')) or 0
    _p_a = _safe_float(m_a.get('share_price')) or 0
    _p_b = _safe_float(m_b.get('share_price')) or 0
    _up_a_pct = ((_dbase_a / _p_a) - 1) * 100 if _p_a and _dbase_a else None
    _up_b_pct = ((_dbase_b / _p_b) - 1) * 100 if _p_b and _dbase_b else None
    if _up_a_pct is not None and _up_b_pct is not None:
        _dcf_leader = tkr_a if _up_a_pct > _up_b_pct else tkr_b
        _max_up = max(_up_a_pct, _up_b_pct)
        _jpm_sub = f"{_dcf_leader} offre le meilleur upside DCF base : {_max_up:+.0f} % vs cours actuel"
    else:
        _jpm_sub = "Lecture DCF Bear/Base/Bull et prix cibles modèles"
    add_text_box(slide, 1.02, 2.10, 23.37, 0.40, _jpm_sub, 10, NAVY_MID, bold=True, italic=True)

    y_dcf = 2.58
    # Panel A — hauteur ajustée au contenu (0.75 header + 11 lignes * 0.53 = 6.58cm)
    _panel_h = 6.60
    add_rect(slide, 1.02, y_dcf, 11.44, _panel_h, COLOR_A_PAL)
    add_rect(slide, 1.02, y_dcf, 0.13, _panel_h, COLOR_A)
    add_text_box(slide, 1.3, y_dcf + 0.1, 11.0, 0.56, f"DCF — {tkr_a}", 9.5, NAVY_MID, bold=True)

    rows_dcf_a = [
        ("Cours actuel",   _fr(m_a.get('share_price'), 1) + " " + cur_a),
        ("Cible Bear",     _fr(m_a.get('dcf_bear'), 0) + " " + cur_a),
        ("Cible Base",     _fr(m_a.get('dcf_base'), 0) + " " + cur_a),
        ("Cible Bull",     _fr(m_a.get('dcf_bull'), 0) + " " + cur_a),
        ("Upside Base",    _upside(m_a.get('dcf_base'), m_a.get('share_price'))),
        ("Marge Sécurité", _frpct(m_a.get('margin_of_safety'))),
        ("WACC",           _frpct(m_a.get('wacc'))),
        ("TGR",            _frpct(m_a.get('terminal_growth'))),
        ("MC P10",         _fr(m_a.get('monte_carlo_p10'), 0) + " " + cur_a),
        ("MC P50",         _fr(m_a.get('monte_carlo_p50'), 0) + " " + cur_a),
        ("MC P90",         _fr(m_a.get('monte_carlo_p90'), 0) + " " + cur_a),
    ]
    for ri, (lbl, val) in enumerate(rows_dcf_a):
        y = y_dcf + 0.75 + ri * 0.53
        fill = WHITE if ri % 2 == 0 else COLOR_A_PAL
        add_rect(slide, 1.22, y, 11.04, 0.5, fill)
        add_text_box(slide, 1.35, y + 0.06, 6.0, 0.4, lbl, 8, GREY_TXT, italic=True)
        add_text_box(slide, 7.5, y + 0.06, 4.6, 0.4, val, 8.5, BLACK, bold=False,
                     align=PP_ALIGN.RIGHT)

    # Panel B — meme hauteur ajustée
    add_rect(slide, 12.94, y_dcf, 11.44, _panel_h, COLOR_B_PAL)
    add_rect(slide, 12.94, y_dcf, 0.13, _panel_h, COLOR_B)
    add_text_box(slide, 13.22, y_dcf + 0.1, 11.0, 0.56, f"DCF — {tkr_b}", 9.5, GREEN, bold=True)

    rows_dcf_b = [
        ("Cours actuel",   _fr(m_b.get('share_price'), 1) + " " + cur_b),
        ("Cible Bear",     _fr(m_b.get('dcf_bear'), 0) + " " + cur_b),
        ("Cible Base",     _fr(m_b.get('dcf_base'), 0) + " " + cur_b),
        ("Cible Bull",     _fr(m_b.get('dcf_bull'), 0) + " " + cur_b),
        ("Upside Base",    _upside(m_b.get('dcf_base'), m_b.get('share_price'))),
        ("Marge Sécurité", _frpct(m_b.get('margin_of_safety'))),
        ("WACC",           _frpct(m_b.get('wacc'))),
        ("TGR",            _frpct(m_b.get('terminal_growth'))),
        ("MC P10",         _fr(m_b.get('monte_carlo_p10'), 0) + " " + cur_b),
        ("MC P50",         _fr(m_b.get('monte_carlo_p50'), 0) + " " + cur_b),
        ("MC P90",         _fr(m_b.get('monte_carlo_p90'), 0) + " " + cur_b),
    ]
    for ri, (lbl, val) in enumerate(rows_dcf_b):
        y = y_dcf + 0.75 + ri * 0.53
        fill = WHITE if ri % 2 == 0 else COLOR_B_PAL
        add_rect(slide, 13.14, y, 11.04, 0.5, fill)
        add_text_box(slide, 13.27, y + 0.06, 6.0, 0.4, lbl, 8, GREY_TXT, italic=True)
        add_text_box(slide, 19.4, y + 0.06, 4.4, 0.4, val, 8.5, BLACK, bold=False,
                     align=PP_ALIGN.RIGHT)

    # Verdict entree — positionné juste après le panel réduit
    ez_a = m_a.get('entry_zone_ok') or 0
    ez_b = m_b.get('entry_zone_ok') or 0
    y_ev = y_dcf + _panel_h + 0.20
    add_rect(slide, 1.02, y_ev, 11.44, 0.7, GREEN_PALE if ez_a else RED_PALE)
    add_text_box(slide, 1.15, y_ev + 0.1, 11.2, 0.5,
                 ("Zone d'achat : Cours < DCF base" if ez_a else "Hors zone : Cours > DCF base"),
                 8.5, GREEN if ez_a else RED, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 12.94, y_ev, 11.44, 0.7, GREEN_PALE if ez_b else RED_PALE)
    add_text_box(slide, 13.07, y_ev + 0.1, 11.2, 0.5,
                 ("Zone d'achat : Cours < DCF base" if ez_b else "Hors zone : Cours > DCF base"),
                 8.5, GREEN if ez_b else RED, bold=True, align=PP_ALIGN.CENTER)

    # Note methodologique DCF
    y_note = y_ev + 0.95
    add_rect(slide, 1.02, y_note, 23.37, 2.4, GREY_BG)
    add_rect(slide, 1.02, y_note, 0.13, 2.4, NAVY_MID)
    dcf_note = (
        "Méthodologie : la fourchette Bear/Base/Bull reflète trois hypothèses de croissance "
        "appliquées au modèle DCF à 5 ans. Le scénario Base correspond au consensus actuel ; "
        "Bear intègre un ralentissement de -30 % sur la croissance ; Bull extrapole une "
        "accélération. La Marge de Sécurité mesure l'écart cours/valeur intrinsèque Base : "
        "plus elle est élevée, plus le coussin de protection est important. "
        "WACC et TGR sont les paramètres les plus sensibles : +1pt de TGR peut faire "
        "varier la valorisation de 15 à 20 %."
    )
    add_text_box(slide, 1.35, y_note + 0.12, 22.8, 2.15, dcf_note, 8, NAVY_MID, wrap=True)

    return slide


def _slide_monte_carlo(prs, m_a: dict, m_b: dict):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    tkr_a = m_a.get('ticker_a') or 'A'
    tkr_b = m_b.get('ticker_b') or 'B'

    navy_bar(slide)
    footer_bar(slide)
    slide_title(slide, _cplbl("gbm_upside"))
    section_dots(slide, 3)
    _company_header_band(slide, tkr_a, tkr_b, m_a=m_a, m_b=m_b)

    # Titre JPM dynamique : qui a le meilleur upside probabiliste P50
    _p50_a = _safe_float(m_a.get('monte_carlo_p50'))
    _p50_b = _safe_float(m_b.get('monte_carlo_p50'))
    _pa = _safe_float(m_a.get('share_price')) or 0
    _pb = _safe_float(m_b.get('share_price')) or 0
    _jpm_mc = "Simulation GBM 10000 paths : distribution probabiliste vs cours actuel"
    if _p50_a and _p50_b and _pa and _pb:
        _up_a = (_p50_a / _pa - 1) * 100
        _up_b = (_p50_b / _pb - 1) * 100
        _mc_leader = tkr_a if _up_a > _up_b else tkr_b
        _jpm_mc = f"{_mc_leader} présente le meilleur upside probabiliste P50 ({max(_up_a, _up_b):+.0f}%)"
    add_text_box(slide, 1.02, 2.10, 23.37, 0.40, _jpm_mc, 10, NAVY_MID, bold=True, italic=True)

    # Tableau Monte Carlo
    cur_a = 'EUR' if (m_a.get('currency_a') or 'USD') == 'EUR' else '$'
    cur_b = 'EUR' if (m_b.get('currency_b') or 'USD') == 'EUR' else '$'
    p_a = m_a.get('share_price')
    p_b = m_b.get('share_price')

    def _sigma_fmt(v):
        if v is None: return "\u2014"
        try: return f"{float(v) * 100:.1f}".replace(".", ",") + " %"
        except Exception: return "\u2014"

    rows_mc = [
        ("P10 (scen. pessimiste)",
         _fr(m_a.get('monte_carlo_p10'), 0) + " " + cur_a + "  (" + _upside(m_a.get('monte_carlo_p10'), p_a) + ")",
         _fr(m_b.get('monte_carlo_p10'), 0) + " " + cur_b + "  (" + _upside(m_b.get('monte_carlo_p10'), p_b) + ")"),
        ("P50 (Médian GBM)",
         _fr(m_a.get('monte_carlo_p50'), 0) + " " + cur_a + "  (" + _upside(m_a.get('monte_carlo_p50'), p_a) + ")",
         _fr(m_b.get('monte_carlo_p50'), 0) + " " + cur_b + "  (" + _upside(m_b.get('monte_carlo_p50'), p_b) + ")"),
        ("P90 (scen. optimiste)",
         _fr(m_a.get('monte_carlo_p90'), 0) + " " + cur_a + "  (" + _upside(m_a.get('monte_carlo_p90'), p_a) + ")",
         _fr(m_b.get('monte_carlo_p90'), 0) + " " + cur_b + "  (" + _upside(m_b.get('monte_carlo_p90'), p_b) + ")"),
        ("Sigma (vol. ann.)", _sigma_fmt(m_a.get('gbm_sigma_annual')), _sigma_fmt(m_b.get('gbm_sigma_annual'))),
        ("Cours actuel S0",  _fr(p_a, 2) + " " + cur_a if p_a else "\u2014",
                             _fr(p_b, 2) + " " + cur_b if p_b else "\u2014"),
    ]

    # --- Layout revu : Texte LLM en haut pleine largeur + tableau bas-gauche + chart bas-droite ---
    # Texte LLM analytique Monte Carlo / Football Field
    _base_a = _safe_float(m_a.get('dcf_base'))
    _base_b = _safe_float(m_b.get('dcf_base'))
    _upside_base_a = ((_base_a / (_safe_float(p_a) or 1) - 1) * 100) if _base_a and p_a else 0
    _upside_base_b = ((_base_b / (_safe_float(p_b) or 1) - 1) * 100) if _base_b and p_b else 0
    _mc_commentary = (
        f"Le modèle DCF Bear/Base/Bull et la simulation GBM Monte Carlo donnent deux lectures "
        f"complémentaires de la valorisation. {tkr_a} affiche un upside base de {_upside_base_a:+.0f}% "
        f"vs {_upside_base_b:+.0f}% pour {tkr_b}. Le football field ci-dessous visualise les "
        f"fourchettes intrinsèques de chaque titre : plus la bande est étendue, plus l'incertitude "
        f"sur la valeur est élevée. La ligne pointillee indique le cours actuel ; si elle se situe "
        f"sous la barre Base, le titre est dans une zone d'achat théorique. A croiser avec les "
        f"hypothèses WACC/TGR (Sensibilité des flux futurs) et la crédibilité du scénario macro."
    )
    # #205 : suppression du titre hardcodé "LECTURE MONTE CARLO & FOOTBALL
    # FIELD" — le texte LLM occupe toute la box (Baptiste veut les textes LLM
    # à la place des titres hardcodés dans les box LLM)
    _mc_txt_y = 2.58
    _mc_txt_h = 2.70
    add_rect(slide, 1.02, _mc_txt_y, 23.37, _mc_txt_h, NAVY_PALE)
    add_rect(slide, 1.02, _mc_txt_y, 0.13, _mc_txt_h, NAVY_MID)
    add_text_box(slide, 1.4, _mc_txt_y + 0.20, 22.8, _mc_txt_h - 0.35,
                 " ".join(_mc_commentary.split()),
                 8.5, NAVY, wrap=True)

    # Tableau Monte Carlo en bas a gauche (40% largeur)
    _tbl_y = _mc_txt_y + _mc_txt_h + 0.30
    _tbl_h = 13.20 - _tbl_y - 0.30
    add_table(
        slide, 1.02, _tbl_y, 10.4, _tbl_h,
        num_rows=len(rows_mc), num_cols=3,
        col_widths_pct=[0.40, 0.30, 0.30],
        header_data=["GBM Monte Carlo", tkr_a, tkr_b],
        rows_data=rows_mc,
        header_fill=NAVY, border_hex="DDDDDD"
    )

    # Football field chart — barres horizontales simples
    try:
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        import numpy as np

        fig, ax = plt.subplots(figsize=(12.5, 6.2))
        fig.patch.set_facecolor('white')
        ax.set_facecolor('#FAFAFA')

        methods_a = ['Bear', 'Base', 'Bull']
        vals_a = [
            _safe_float(m_a.get('dcf_bear')) or 0,
            _safe_float(m_a.get('dcf_base')) or 0,
            _safe_float(m_a.get('dcf_bull')) or 0,
        ]
        vals_b = [
            _safe_float(m_b.get('dcf_bear')) or 0,
            _safe_float(m_b.get('dcf_base')) or 0,
            _safe_float(m_b.get('dcf_bull')) or 0,
        ]

        y_pos = np.array([3, 2, 1, -1, -2, -3])
        labels = [f'{tkr_a} Bear', f'{tkr_a} Base', f'{tkr_a} Bull',
                  f'{tkr_b} Bear', f'{tkr_b} Base', f'{tkr_b} Bull']
        all_vals = vals_a + vals_b
        colors_ff = ['#6688BB', '#1B3A6B', '#102040',
                     '#E5C06F', '#C9A227', '#7A6010']

        for i, (yp, val, lbl, col) in enumerate(zip(y_pos, all_vals, labels, colors_ff)):
            if val:
                ax.barh(yp, val, 0.65, color=col, alpha=0.90)
                ax.text(val + max(all_vals or [1]) * 0.01, yp,
                        f'{val:.0f}', ha='left', va='center', fontsize=11, color='#333')

        # Cours actuels — lignes verticales avec labels
        if p_a:
            ax.axvline(p_a, color='#1B3A6B', linestyle='--', linewidth=1.6, alpha=0.7)
            ax.text(p_a, 3.75, f' {tkr_a}\n cours', fontsize=9,
                    color='#1B3A6B', va='bottom', ha='left', clip_on=True)
        if p_b:
            ax.axvline(p_b, color='#C9A227', linestyle='--', linewidth=1.6, alpha=0.7)
            ax.text(p_b, -3.9, f' {tkr_b}\n cours', fontsize=9,
                    color='#6B5010', va='top', ha='left', clip_on=True)

        ax.set_yticks(y_pos)
        ax.set_yticklabels(labels, fontsize=11)
        ax.set_xlabel('Prix cible', fontsize=11)
        ax.set_title('Football Field Comparé — DCF Bear/Base/Bull vs cours actuel',
                     fontsize=13, fontweight='bold', color='#1B3A6B', pad=10)
        ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
        ax.tick_params(labelsize=11)
        ax.xaxis.grid(True, alpha=0.25, color='#D0D5DD'); ax.set_axisbelow(True)
        fig.tight_layout(pad=1.0)
        buf = _make_chart_buf(fig)
        # Ratio matplotlib 12.5/6.2 = 2.016 — chart positionné a droite du tableau
        _fb_x = 11.80
        _fb_w = 12.60
        _fb_h = _fb_w / 2.016  # ~6.25cm
        if _fb_h > _tbl_h:
            _fb_h = _tbl_h
            _fb_w = _fb_h * 2.016
        # Centrer verticalement dans la meme zone que le tableau
        _fb_y = _tbl_y + (_tbl_h - _fb_h) / 2
        _insert_chart(slide, buf, _fb_x, _fb_y, _fb_w, _fb_h)
    except Exception as e:
        log.warning(f"[cmp_pptx] football field error: {e}")

    return slide


def _slide_piotroski(prs, m_a: dict, m_b: dict, synthesis: dict):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    tkr_a = m_a.get('ticker_a') or 'A'
    tkr_b = m_b.get('ticker_b') or 'B'

    navy_bar(slide)
    footer_bar(slide)
    slide_title(slide, _cplbl("qual_fin"))
    section_dots(slide, 4)
    _company_header_band(slide, tkr_a, tkr_b, m_a=m_a, m_b=m_b)

    # Titre JPM dynamique basé sur Piotroski
    _pio_a = _safe_float(m_a.get('piotroski_score')) or 0
    _pio_b = _safe_float(m_b.get('piotroski_score')) or 0
    if _pio_a or _pio_b:
        _q_leader = tkr_a if _pio_a > _pio_b else tkr_b
        _jpm_sub = f"{_q_leader} présente la qualité fondamentale la plus solide (Piotroski {max(int(_pio_a), int(_pio_b))}/9)"
    else:
        _jpm_sub = "Piotroski F-Score, Beneish M-Score, Altman Z-Score comparatifs"
    add_text_box(slide, 1.02, 2.10, 23.37, 0.40, _jpm_sub, 10, NAVY_MID, bold=True, italic=True)

    # Commentaire qualité LLM — clip strict pour éviter tout overflow
    txt = synthesis.get('quality_text') or ""
    if txt:
        _qt_y = 2.58
        _qt_h = 2.10
        add_rect(slide, 1.02, _qt_y, 23.37, _qt_h, NAVY_PALE)
        add_rect(slide, 1.02, _qt_y, 0.13, _qt_h, NAVY_MID)
        add_text_box(slide, 1.4, _qt_y + 0.10, 22.6, _qt_h - 0.20,
                     " ".join(_fit(txt, 360).split()), 8.5, NAVY, wrap=True)

    def _pio_val(m, key):
        v = m.get(key)
        if v is None: return "\u2014"
        return "1" if int(v) == 1 else "0"

    pio_keys = [
        ("pio_roa_positive",         "ROA positif"),
        ("pio_cfo_positive",         "Cash-flows opérations positifs"),
        ("pio_delta_roa",            "Amélioration ROA"),
        ("pio_accruals",             "Qualité earnings (accruals)"),
        ("pio_delta_leverage",       "Réduction levier"),
        ("pio_delta_liquidity",      "Amélioration liquidité"),
        ("pio_no_dilution",          "Pas de dilution"),
        ("pio_delta_gross_margin",   "Amélioration marge brute"),
        ("pio_delta_asset_turnover", "Amélioration rotation actifs"),
    ]
    rows_pio = [(lbl, _pio_val(m_a, key), _pio_val(m_b, key)) for key, lbl in pio_keys]
    sc_a = m_a.get('piotroski_score')
    sc_b = m_b.get('piotroski_score')
    rows_pio.append(("PIOTROSKI F-SCORE TOTAL",
                     str(sc_a) if sc_a is not None else "\u2014",
                     str(sc_b) if sc_b is not None else "\u2014"))

    # Tableau Piotroski réduit et decaler plus bas pour liberer le cadre LLM
    _pio_y = 4.85
    _pio_h = 6.10
    add_table(
        slide, 1.02, _pio_y, 14.0, _pio_h,
        num_rows=len(rows_pio), num_cols=3,
        col_widths_pct=[0.60, 0.20, 0.20],
        header_data=["Critère Piotroski", tkr_a, tkr_b],
        rows_data=rows_pio,
        header_fill=NAVY, border_hex="DDDDDD"
    )

    # Scores complémentaires positionnés EN BAS A DROITE (plus compacts)
    # Legendes Beneish/Altman/Sloan retirees (trop verboses) — condensees en footnote
    _sc_x = 15.15
    _sc_w = 9.87
    _sc_y = _pio_y
    add_rect(slide, _sc_x, _sc_y, _sc_w, 0.5, NAVY)
    add_text_box(slide, _sc_x + 0.15, _sc_y + 0.08, _sc_w - 0.2, 0.38,
                 "Scores complémentaires", 8.5, WHITE, bold=True)

    rows_sc = [
        ("Beneish M-Score", _fr(m_a.get('beneish_mscore'), 2), _fr(m_b.get('beneish_mscore'), 2)),
        ("Altman Z-Score",  _fr(m_a.get('altman_z'), 2),       _fr(m_b.get('altman_z'), 2)),
        ("Sloan Accruals",  _fr(m_a.get('sloan_accruals'), 3), _fr(m_b.get('sloan_accruals'), 3)),
        ("Cash Conversion", _frx(m_a.get('cash_conversion')), _frx(m_b.get('cash_conversion'))),
    ]
    add_table(
        slide, _sc_x, _sc_y + 0.55, _sc_w, 2.40,
        num_rows=len(rows_sc), num_cols=3,
        col_widths_pct=[0.50, 0.25, 0.25],
        header_data=["Indicateur", tkr_a, tkr_b],
        rows_data=rows_sc,
        header_fill=NAVY_MID, border_hex="DDDDDD"
    )

    # Footnote condensee (1 seule ligne) : seuils cles
    _foot_y = _sc_y + 3.05
    add_text_box(slide, _sc_x, _foot_y, _sc_w, 0.40,
                 "Seuils : Beneish <-1,78 sain  -  Altman >2,99 sain  -  Sloan proche 0 = qualité",
                 7, GREY_TXT, italic=True, wrap=True)

    return slide


def _slide_risque(prs, m_a: dict, m_b: dict):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    tkr_a = m_a.get('ticker_a') or 'A'
    tkr_b = m_b.get('ticker_b') or 'B'

    navy_bar(slide)
    footer_bar(slide)
    slide_title(slide, _cplbl("risque_mom"))
    section_dots(slide, 4)
    _company_header_band(slide, tkr_a, tkr_b, m_a=m_a, m_b=m_b)

    # Titre analytique JPM
    _b_a = _safe_float(m_a.get('beta')) or 1.0
    _b_b = _safe_float(m_b.get('beta')) or 1.0
    _def = tkr_a if _b_a < _b_b else tkr_b
    _jpm_sub = f"{_def} offre un profil plus défensif (beta {min(_b_a,_b_b):.2f} vs {max(_b_a,_b_b):.2f})".replace('.', ',')
    add_text_box(slide, 1.02, 2.38, 23.37, 0.45, _jpm_sub, 10, NAVY_MID, bold=True, italic=True)

    # Tableau performances (gauche) — KPIs supprimés, table monte direct
    rows_perf = [
        ("Perf. 1 Mois",  _frpct(m_a.get('perf_1m'), signed=True),  _frpct(m_b.get('perf_1m'), signed=True)),
        ("Perf. 3 Mois",  _frpct(m_a.get('perf_3m'), signed=True),  _frpct(m_b.get('perf_3m'), signed=True)),
        ("Perf. 1 An",    _frpct(m_a.get('perf_1y'), signed=True),  _frpct(m_b.get('perf_1y'), signed=True)),
        ("Score Momentum", _fr(m_a.get('momentum_score'), 1) + " /10", _fr(m_b.get('momentum_score'), 1) + " /10"),
        ("Beta",           _fr(m_a.get('beta'), 2), _fr(m_b.get('beta'), 2)),
        ("VaR 95% 1M",    _frpct(m_a.get('var_95_1m')), _frpct(m_b.get('var_95_1m'))),
    ]
    add_table(
        slide, 1.02, 2.95, 11.5, 5.20,
        num_rows=len(rows_perf), num_cols=3,
        col_widths_pct=[0.44, 0.28, 0.28],
        header_data=["Performance / Risque", tkr_a, tkr_b],
        rows_data=rows_perf,
        header_fill=NAVY, border_hex="DDDDDD"
    )

    # Radar chart — droite, agrandi, titre lisible
    buf = _chart_risk_profile(m_a, m_b, tkr_a, tkr_b)
    if buf:
        _insert_chart(slide, buf, 12.85, 2.95, 11.55, 5.20)

    # Fourchette 52 semaines — bandeau plein largeur
    # #206 : réhausse 8.40 → 7.60 pour aérer la slide (Baptiste)
    y_rng = 7.60
    _rng_w = 23.37
    add_rect(slide, 1.02, y_rng, _rng_w, 0.55, NAVY)
    add_text_box(slide, 1.18, y_rng + 0.12, _rng_w - 0.3, 0.40,
                 "Fourchette 52 Semaines", 9, WHITE, bold=True)

    rows_52 = [
        (f"{tkr_a} — 52W Low",  _fr(m_a.get('week52_low'), 2),
         f"{tkr_a} — 52W High", _fr(m_a.get('week52_high'), 2)),
        (f"{tkr_b} — 52W Low",  _fr(m_b.get('week52_low'), 2),
         f"{tkr_b} — 52W High", _fr(m_b.get('week52_high'), 2)),
    ]
    y_52 = y_rng + 0.62
    _row_w = _rng_w / 2
    for ri, (l1, v1, l2, v2) in enumerate(rows_52):
        fill = WHITE if ri % 2 == 0 else GREY_BG
        add_rect(slide, 1.02, y_52 + ri * 0.55, _rng_w, 0.52, fill)
        add_text_box(slide, 1.18, y_52 + ri * 0.55 + 0.08, 5.0, 0.40,
                     l1, 8.5, GREY_TXT, wrap=False)
        add_text_box(slide, 6.30, y_52 + ri * 0.55 + 0.08, 3.0, 0.40,
                     v1, 9.5, BLACK, bold=True, wrap=False)
        add_text_box(slide, 12.80, y_52 + ri * 0.55 + 0.08, 5.0, 0.40,
                     l2, 8.5, GREY_TXT, wrap=False)
        add_text_box(slide, 17.95, y_52 + ri * 0.55 + 0.08, 3.0, 0.40,
                     v2, 9.5, BLACK, bold=True, wrap=False)

    # --- Texte LLM analytique risque — pleine largeur ---
    _p3m_a = _safe_float(m_a.get('perf_3m')) or 0
    _p3m_b = _safe_float(m_b.get('perf_3m')) or 0
    if abs(_p3m_a) <= 2: _p3m_a *= 100
    if abs(_p3m_b) <= 2: _p3m_b *= 100
    _mom_leader = tkr_a if _p3m_a > _p3m_b else tkr_b
    _vol_a = _safe_float(m_a.get('volatility_52w')) or 0
    _vol_b = _safe_float(m_b.get('volatility_52w')) or 0
    if abs(_vol_a) <= 2: _vol_a *= 100
    if abs(_vol_b) <= 2: _vol_b *= 100
    risque_txt = (
        f"Profil de risque : {_def} présente une Sensibilité marché moindre "
        f"(beta {min(_b_a, _b_b):.2f}), adapte aux portefeuilles prudents. "
        f"Sur le momentum 3 mois, {_mom_leader} prend la tête "
        f"({max(_p3m_a, _p3m_b):+.1f}% vs {min(_p3m_a, _p3m_b):+.1f}%). "
        f"La volatilité annualisée ({_vol_a:.0f}% vs {_vol_b:.0f}%) calibré "
        f"le dimensionnement de position : une volatilité plus élevée impose "
        f"une pondération réduite pour maintenir un budget risque équivalent."
    )
    # #205/#206 : réhausse le bloc LLM (10.75 → 10.25) + supprimé le titre
    # "LECTURE RISQUE/MOMENTUM", donne plus de place au body
    y_llm = 10.25
    _llm_h = 2.35
    add_rect(slide, 1.02, y_llm, 23.37, _llm_h, NAVY_PALE)
    add_rect(slide, 1.02, y_llm, 0.13, _llm_h, NAVY_MID)
    add_text_box(slide, 1.35, y_llm + 0.18, 22.8, _llm_h - 0.30,
                 " ".join(risque_txt.split()), 9, NAVY, wrap=True)

    return slide


def _slide_finsight_score(prs, m_a: dict, m_b: dict):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    tkr_a = m_a.get('ticker_a') or 'A'
    tkr_b = m_b.get('ticker_b') or 'B'

    navy_bar(slide)
    footer_bar(slide)
    slide_title(slide, _cplbl("score_cmp"))
    section_dots(slide, 5)
    _company_header_band(slide, tkr_a, tkr_b, m_a=m_a, m_b=m_b)

    fs_a = m_a.get('finsight_score') or 0
    fs_b = m_b.get('finsight_score') or 0

    # Titre JPM dynamique
    try:
        _gap = abs(int(fs_a) - int(fs_b))
        _score_winner = tkr_a if fs_a >= fs_b else tkr_b
        _pt_lbl = "point" if _gap == 1 else "points"
        _jpm_sub = f"{_score_winner} l'emporte au score composite avec {_gap} {_pt_lbl} d'avance"
    except Exception:
        _jpm_sub = "Lecture du score composite FinSight (Valeur / Croissance / Qualité / Momentum)"
    add_text_box(slide, 1.02, 2.38, 23.37, 0.45, _jpm_sub, 10, NAVY_MID, bold=True, italic=True)

    rec_a = (m_a.get('recommendation') or 'HOLD').upper()
    rec_b = (m_b.get('recommendation') or 'HOLD').upper()
    def _fill_rec(r):
        if r == 'BUY': return GREEN_PALE, GREEN
        if r == 'SELL': return RED_PALE, RED
        return AMBER_PALE, AMBER
    fa, aa = _fill_rec(rec_a)
    fb, ab = _fill_rec(rec_b)

    # Graphique score gauche (ratio 8.5/4.2 ~ 2.02) + BUY blocks droite
    buf = _chart_finsight_score(m_a, m_b, tkr_a, tkr_b)
    if buf:
        _fs_h = 4.55
        _fs_w = _fs_h * 2.02  # ~9.20cm
        _insert_chart(slide, buf, 1.02, 2.95, 14.0, _fs_h)

    # BUY/SELL blocks a droite du graphique
    _kx = 15.40
    _kw = 8.97
    kpi_box(slide, _kx, 2.95, _kw, 2.10, rec_a, f"Recommandation — {tkr_a}",
            _frpct(m_a.get('conviction')) + " conviction", fa, aa)
    kpi_box(slide, _kx, 5.40, _kw, 2.10, rec_b, f"Recommandation — {tkr_b}",
            _frpct(m_b.get('conviction')) + " conviction", fb, ab)

    # Décomposition score — plus compact
    # #207 : réhausse 7.85 → 7.15 pour aérer la slide (Baptiste)
    y_dec = 7.15
    add_rect(slide, 1.02, y_dec, 23.37, 0.55, NAVY)
    add_text_box(slide, 1.18, y_dec + 0.10, 23.0, 0.42,
                 "Décomposition du Score Composite (4 axes : Valeur / Qualité / Momentum / Gouvernance — chacun /25)",
                 9, WHITE, bold=True)

    # 4 axes de score — réels (issus de compute_score v1.3 sur le snapshot pipeline)
    # Bug audit 27/04 : avant, proxies indépendants ne sommaient pas au global.
    # Maintenant on lit `finsight_score_{value,quality,momentum,governance}`
    # stockés dans m via _finsight_score_full. Fallback proxy si manquant.
    def _score_axis(m, axis) -> str:
        # 1. Tenter le vrai sous-score depuis compute_score v1.3
        real = _safe_float(m.get(f"finsight_score_{axis}"))
        if real is not None:
            return str(round(real))
        # 2. Fallback proxy historique (si compute_score a échoué)
        if axis == 'value':
            pe = _safe_float(m.get('pe_ratio')) or 25
            return str(max(0, min(25, round(25 - pe * 0.4))))
        if axis == 'growth':
            cagr = (_safe_float(m.get('revenue_cagr_3y')) or 0) * 100
            return str(max(0, min(25, round(12 + cagr))))
        if axis == 'quality':
            nm = (_safe_float(m.get('ebitda_margin_ltm')) or 0.1)
            if nm > 1: nm /= 100
            return str(max(0, min(25, round(10 + nm * 50))))
        if axis == 'momentum':
            p3m = (_safe_float(m.get('perf_3m')) or 0) * 100
            return str(max(0, min(25, round(12 + p3m * 0.15))))
        if axis == 'governance':
            pio = _safe_float(m.get('piotroski_score')) or 5
            return str(max(0, min(25, round(pio * 2.5))))
        return "\u2014"

    # 4 axes finsight_score v1.3 : Quality, Value, Momentum, Governance
    axes = [('value', 'Valeur'), ('quality', 'Qualité'), ('momentum', 'Momentum'), ('governance', 'Gouvernance')]
    y_ax = y_dec + 0.70
    for xi, (axis, lbl) in enumerate(axes):
        xp = 1.02 + xi * 5.84
        w_ax = 5.6
        add_rect(slide, xp, y_ax, w_ax, 1.20, GREY_BG)
        add_rect(slide, xp, y_ax, w_ax, 0.04, NAVY_MID)
        add_text_box(slide, xp + 0.1, y_ax + 0.08, w_ax - 0.2, 0.35, lbl, 8.5, GREY_TXT, wrap=False)
        add_text_box(slide, xp + 0.1, y_ax + 0.45, 2.7, 0.60,
                     _score_axis(m_a, axis) + " / 25", 11, NAVY_MID, bold=True)
        add_text_box(slide, xp + 2.9, y_ax + 0.45, 2.5, 0.60,
                     _score_axis(m_b, axis) + " / 25", 11, GREEN_MID, bold=True)

    # --- Bande LLM analytique score — pleine largeur en bas ---
    y_llm = y_ax + 1.40
    _conv_a = _safe_float(m_a.get('conviction')) or 0
    _conv_b = _safe_float(m_b.get('conviction')) or 0
    if abs(_conv_a) <= 2: _conv_a *= 100
    if abs(_conv_b) <= 2: _conv_b *= 100
    try: _fs_delta = abs(int(fs_a) - int(fs_b))
    except Exception: _fs_delta = 0
    try:
        _fs_leader_tk = tkr_a if int(fs_a) >= int(fs_b) else tkr_b
    except Exception:
        _fs_leader_tk = tkr_a
    _pt_lbl_fs = "point" if _fs_delta == 1 else "points"
    score_txt = (
        f"Le score composite FinSight (v1.3) intègre 4 piliers (Valeur, Qualité, Momentum, Gouvernance) "
        f"sur une échelle de 0 à 100, pondérés selon le profil sectoriel. La somme brute des 4 axes "
        f"est ensuite recalibrée (mean≈50, std≈20) pour discriminer le ranking — l'écart à la somme "
        f"naïve est donc normal. {_fs_leader_tk} ressort en tête avec un écart de {_fs_delta} {_pt_lbl_fs}, "
        f"traduisant un meilleur équilibre global des fondamentaux. "
        f"Les convictions respectives ({_conv_a:.0f} % vs {_conv_b:.0f} %) quantifient la "
        f"robustesse de l'opinion sur la Thèse, à croiser avec l'exposition en portefeuille "
        f"et le catalyseur de déblocage de valeur à 12 mois."
    )
    # #205/#207 : supprimé le titre "LECTURE ANALYTIQUE DU SCORE COMPOSITE"
    # et agrandit la box LLM vers le bas (1.65 → 2.50)
    _llm_h = 2.50
    add_rect(slide, 1.02, y_llm, 23.37, _llm_h, NAVY_PALE)
    add_rect(slide, 1.02, y_llm, 0.13, _llm_h, NAVY_MID)
    add_text_box(slide, 1.35, y_llm + 0.18, 22.8, _llm_h - 0.30,
                 " ".join(score_txt.split()), 9, NAVY, wrap=True)

    return slide


def _slide_theses(prs, m_a: dict, m_b: dict, synthesis: dict):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    tkr_a = m_a.get('ticker_a') or 'A'
    tkr_b = m_b.get('ticker_b') or 'B'

    navy_bar(slide)
    footer_bar(slide)
    slide_title(slide, _cplbl("theses_bb"))
    section_dots(slide, 5)
    _company_header_band(slide, tkr_a, tkr_b, m_a=m_a, m_b=m_b)

    bull_a = synthesis.get('bull_a') or "Croissance supérieure a la moyenne sectorielle avec fort levier opérationnel."
    bear_a = synthesis.get('bear_a') or "Valorisation premium susceptible de re-rating si la croissance ralentit."
    bull_b = synthesis.get('bull_b') or "Génération de cash robuste et bilan solide pour soutenir la valorisation."
    bear_b = synthesis.get('bear_b') or "Exposition a un secteur cyclique amplifiant la volatilité des marges."

    # JPM sub retiré (audit 28/04 : se chevauchait avec bandeau société 52W,
    # information redondante avec titre slide "Thèses d'Investissement Bull / Bear")

    # Bande 52W — cours & fourchette
    cur_a = "EUR" if (m_a.get('currency_a') or 'USD') == 'EUR' else '$'
    cur_b = "EUR" if (m_b.get('currency_b') or 'USD') == 'EUR' else '$'
    y_52 = 2.40
    add_rect(slide, 1.02, y_52, 11.44, 0.75, NAVY_PALE)
    add_rect(slide, 1.02, y_52, 0.13, 0.75, COLOR_A)
    add_text_box(slide, 1.25, y_52 + 0.05, 10.0, 0.32,
                 f"Cours : {_fr(m_a.get('share_price'), 1)} {cur_a}  ·  "
                 f"52W : {_fr(m_a.get('week52_low'), 1)} - {_fr(m_a.get('week52_high'), 1)} {cur_a}",
                 8, NAVY, bold=True, wrap=False)
    add_text_box(slide, 1.25, y_52 + 0.37, 10.0, 0.32,
                 f"Perf 1Y : {_frpct(m_a.get('perf_1y'), signed=True)}  ·  "
                 f"Div. : {_frpct(m_a.get('dividend_yield'))}",
                 8, GREY_TXT, wrap=False)

    add_rect(slide, 12.94, y_52, 11.44, 0.75, GREEN_PALE)
    add_rect(slide, 12.94, y_52, 0.13, 0.75, COLOR_B)
    add_text_box(slide, 13.17, y_52 + 0.05, 10.0, 0.32,
                 f"Cours : {_fr(m_b.get('share_price'), 1)} {cur_b}  ·  "
                 f"52W : {_fr(m_b.get('week52_low'), 1)} - {_fr(m_b.get('week52_high'), 1)} {cur_b}",
                 8, GREEN, bold=True, wrap=False)
    add_text_box(slide, 13.17, y_52 + 0.37, 10.0, 0.32,
                 f"Perf 1Y : {_frpct(m_b.get('perf_1y'), signed=True)}  ·  "
                 f"Div. : {_frpct(m_b.get('dividend_yield'))}",
                 8, GREY_TXT, wrap=False)

    # Audit 28/04 : Bull/Bear réhaussés davantage après suppression JPM sub
    y0 = 3.36
    bull_h = 4.30
    y1 = y0 + bull_h + 0.25
    bear_h = 4.30

    # --- Bull A ---
    add_rect(slide, 1.02, y0, 11.44, 0.55, GREEN_PALE)
    add_rect(slide, 1.02, y0, 0.18, 0.55, GREEN)
    add_text_box(slide, 1.35, y0 + 0.05, 11.0, 0.38, f"\u25b2 BULL THESIS — {tkr_a}",
                 9, GREEN, bold=True)
    add_rect(slide, 1.02, y0 + 0.60, 11.44, bull_h - 0.60, GREY_BG)
    add_text_box(slide, 1.18, y0 + 0.65, 11.12, bull_h - 0.70,
                 _fit(bull_a, 520), 9, BLACK, wrap=True)

    # --- Bull B ---
    add_rect(slide, 12.94, y0, 11.44, 0.55, GREEN_PALE)
    add_rect(slide, 12.94, y0, 0.18, 0.55, GREEN)
    add_text_box(slide, 13.27, y0 + 0.05, 11.0, 0.38, f"\u25b2 BULL THESIS — {tkr_b}",
                 9, GREEN, bold=True)
    add_rect(slide, 12.94, y0 + 0.60, 11.44, bull_h - 0.60, GREY_BG)
    add_text_box(slide, 13.10, y0 + 0.65, 11.12, bull_h - 0.70,
                 _fit(bull_b, 520), 9, BLACK, wrap=True)

    # --- Bear A ---
    add_rect(slide, 1.02, y1, 11.44, 0.55, RED_PALE)
    add_rect(slide, 1.02, y1, 0.18, 0.55, RED)
    add_text_box(slide, 1.35, y1 + 0.05, 11.0, 0.38, f"\u25bc BEAR THESIS — {tkr_a}",
                 9, RED, bold=True)
    add_rect(slide, 1.02, y1 + 0.60, 11.44, bear_h - 0.60, GREY_BG)
    add_text_box(slide, 1.18, y1 + 0.65, 11.12, bear_h - 0.70,
                 _fit(bear_a, 520), 9, BLACK, wrap=True)

    # --- Bear B ---
    add_rect(slide, 12.94, y1, 11.44, 0.55, RED_PALE)
    add_rect(slide, 12.94, y1, 0.18, 0.55, RED)
    add_text_box(slide, 13.27, y1 + 0.05, 11.0, 0.38, f"\u25bc BEAR THESIS — {tkr_b}",
                 9, RED, bold=True)
    add_rect(slide, 12.94, y1 + 0.60, 11.44, bear_h - 0.60, GREY_BG)
    add_text_box(slide, 13.10, y1 + 0.65, 11.12, bear_h - 0.70,
                 _fit(bear_b, 520), 9, BLACK, wrap=True)

    return slide


def _slide_price_chart(prs, m_a: dict, m_b: dict, synthesis: dict = None):
    """Slide 20 — Performance 52 semaines : GRAPHIQUE GAUCHE + metriques condensees TOP-RIGHT + LLM macro.

    Layout redesign :
      - PAS de company_header_band (evite le "IR.PA" parasite)
      - Graphique 52W a gauche (60% largeur) pleine hauteur
      - Metriques condensees compactes en haut a droite
      - Texte LLM analytique macro en bas a droite (justification cours + perspectives)
    """
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    tkr_a = m_a.get('ticker_a') or 'A'
    tkr_b = m_b.get('ticker_b') or 'B'
    cur_a = "EUR" if (m_a.get('currency_a') or 'USD') == 'EUR' else '$'
    cur_b = "EUR" if (m_b.get('currency_b') or 'USD') == 'EUR' else '$'

    navy_bar(slide)
    footer_bar(slide)
    slide_title(slide, _cplbl("perf_52w"))
    section_dots(slide, 5)
    # PAS de _company_header_band — évite l'artefact "IR.PA" en haut a droite

    # Titre JPM dynamique
    _p1y_a = _safe_float(m_a.get('perf_1y')) or 0
    _p1y_b = _safe_float(m_b.get('perf_1y')) or 0
    if abs(_p1y_a) <= 2: _p1y_a *= 100
    if abs(_p1y_b) <= 2: _p1y_b *= 100
    if _p1y_a or _p1y_b:
        _leader = tkr_a if _p1y_a > _p1y_b else tkr_b
        _spread = abs(_p1y_a - _p1y_b)
        _jpm_sub = f"{_leader} surperforme sur 12 mois (écart de {_spread:.0f} pts)"
    else:
        _jpm_sub = "Trajectoire de cours comparée Normalisée base 100"
    add_text_box(slide, 1.02, 2.12, 23.37, 0.45, _jpm_sub, 11, NAVY_MID, bold=True, italic=True)

    # ---------- GAUCHE : graphique 52W ----------
    x_chart = 1.02
    w_chart = 14.20
    y_chart = 2.70
    h_chart = 10.40
    buf = _chart_52w_price(tkr_a, tkr_b)
    if buf:
        _insert_chart(slide, buf, x_chart, y_chart, w_chart, h_chart)
    else:
        add_rect(slide, x_chart, y_chart, w_chart, h_chart, GREY_BG)
        add_text_box(slide, x_chart + 0.5, y_chart + 4.5, w_chart - 1, 1.0,
                     "Données de cours indisponibles (yfinance)", 11, GREY_TXT,
                     align=PP_ALIGN.CENTER)

    # ---------- DROITE : panel metriques condense (haut) ----------
    x_r = 15.55
    w_r = 24.25 - x_r  # ~8.7cm
    y_m = 2.70

    # Bloc A condense (8 lignes en 2 colonnes d'etiquette+valeur)
    _mh_a = 3.0
    add_rect(slide, x_r, y_m, w_r, _mh_a, COLOR_A_PAL)
    add_rect(slide, x_r, y_m, 0.14, _mh_a, COLOR_A)
    add_text_box(slide, x_r + 0.22, y_m + 0.10, w_r - 0.3, 0.40,
                 f"{tkr_a}  ·  Performance & Prix",
                 9.5, NAVY_MID, bold=True, wrap=False)
    # Metrics 2x4 grid
    a_rows = [
        ("Cours",       _fr(m_a.get('share_price'), 1) + " " + cur_a,
         "52W High",    _fr(m_a.get('week52_high'), 1)),
        ("Perf. 1M",    _frpct(m_a.get('perf_1m'), signed=True),
         "52W Low",     _fr(m_a.get('week52_low'), 1)),
        ("Perf. 3M",    _frpct(m_a.get('perf_3m'), signed=True),
         "Div. Yield",  _frpct(m_a.get('dividend_yield'))),
        ("Perf. 1Y",    _frpct(m_a.get('perf_1y'), signed=True),
         "Volatilité",  _frpct(m_a.get('volatility_52w'))),
    ]
    _lh = 0.55
    for ri, (lbl1, val1, lbl2, val2) in enumerate(a_rows):
        ry = y_m + 0.55 + ri * _lh
        add_text_box(slide, x_r + 0.25, ry, 2.0, 0.35, lbl1, 7, GREY_TXT, wrap=False)
        add_text_box(slide, x_r + 2.25, ry, 2.0, 0.35, val1, 8, BLACK, bold=True, wrap=False)
        add_text_box(slide, x_r + 4.50, ry, 2.0, 0.35, lbl2, 7, GREY_TXT, wrap=False)
        add_text_box(slide, x_r + 6.50, ry, 2.0, 0.35, val2, 8, BLACK, bold=True, wrap=False)

    # Bloc B condense
    y_m2 = y_m + _mh_a + 0.18
    add_rect(slide, x_r, y_m2, w_r, _mh_a, COLOR_B_PAL)
    add_rect(slide, x_r, y_m2, 0.14, _mh_a, COLOR_B)
    add_text_box(slide, x_r + 0.22, y_m2 + 0.10, w_r - 0.3, 0.40,
                 f"{tkr_b}  ·  Performance & Prix",
                 9.5, GREEN_MID, bold=True, wrap=False)
    b_rows = [
        ("Cours",       _fr(m_b.get('share_price'), 1) + " " + cur_b,
         "52W High",    _fr(m_b.get('week52_high'), 1)),
        ("Perf. 1M",    _frpct(m_b.get('perf_1m'), signed=True),
         "52W Low",     _fr(m_b.get('week52_low'), 1)),
        ("Perf. 3M",    _frpct(m_b.get('perf_3m'), signed=True),
         "Div. Yield",  _frpct(m_b.get('dividend_yield'))),
        ("Perf. 1Y",    _frpct(m_b.get('perf_1y'), signed=True),
         "Volatilité",  _frpct(m_b.get('volatility_52w'))),
    ]
    for ri, (lbl1, val1, lbl2, val2) in enumerate(b_rows):
        ry = y_m2 + 0.55 + ri * _lh
        add_text_box(slide, x_r + 0.25, ry, 2.0, 0.35, lbl1, 7, GREY_TXT, wrap=False)
        add_text_box(slide, x_r + 2.25, ry, 2.0, 0.35, val1, 8, BLACK, bold=True, wrap=False)
        add_text_box(slide, x_r + 4.50, ry, 2.0, 0.35, lbl2, 7, GREY_TXT, wrap=False)
        add_text_box(slide, x_r + 6.50, ry, 2.0, 0.35, val2, 8, BLACK, bold=True, wrap=False)

    # Bloc LLM analytique macro (bas droite)
    y_llm = y_m2 + _mh_a + 0.25
    _llm_h = 13.10 - y_llm  # jusqu'au footer
    add_rect(slide, x_r, y_llm, w_r, _llm_h, NAVY_PALE)
    add_rect(slide, x_r, y_llm, 0.14, _llm_h, NAVY_MID)
    add_text_box(slide, x_r + 0.25, y_llm + 0.10, w_r - 0.3, 0.40,
                 "LECTURE MACRO & PERSPECTIVES", 9, NAVY, bold=True)

    # Texte LLM depuis synthesis.price_narrative OU fallback analytique
    llm_txt = (synthesis or {}).get('price_narrative') or ""
    if not llm_txt or len(llm_txt) < 60:
        # Fallback : analyse macro/trajectoire déterministe
        _leader_tk = tkr_a if _p1y_a > _p1y_b else tkr_b
        _beta_a = _safe_float(m_a.get('beta')) or 1.0
        _beta_b = _safe_float(m_b.get('beta')) or 1.0
        _def_tk = tkr_a if _beta_a < _beta_b else tkr_b
        _beta_str = f"{min(_beta_a, _beta_b):.2f}".replace('.', ',')
        llm_txt = (
            f"{_leader_tk} surperforme sur 12 mois dans un contexte de taux directeurs "
            f"stabilisés et de rotation sectorielle favorable aux profils de qualité. "
            f"Le spread de performance reflète une combinaison de moments pricing-power, "
            f"discipline opérationnelle et réallocation des flux institutionnels. "
            f"{_def_tk} (beta {_beta_str}) offre un profil plus défensif "
            f"et résiste mieux dans les phases de stress. Durabilité : dépend de la visibilité "
            f"des résultats trimestriels et de l'évolution du cycle macro 6-12 mois. "
            f"Catalyseurs : publications trimestrielles, guidance FY+1, événements sectoriels."
        )
    add_text_box(slide, x_r + 0.25, y_llm + 0.52, w_r - 0.5, _llm_h - 0.7,
                 " ".join(_fit(llm_txt, 580).split()),
                 8, NAVY, wrap=True)

    return slide


def _slide_verdict(prs, m_a: dict, m_b: dict, synthesis: dict):
    """Slide 21 — Verdict final : conviction GAUCHE + Thèse/catalyseurs/risques DROITE.

    Redesign inspire du slide verdict societe : banderole reduite, conviction box gauche
    avec delta FinSight Score, panel droite avec these d'investissement + catalyseurs bulls
    + risques bears + conditions d'invalidation.
    """
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    tkr_a  = m_a.get('ticker_a') or 'A'
    tkr_b  = m_b.get('ticker_b') or 'B'
    # Winner refondu C4 : LLM decide, fallback finsight, NEUTRAL si indécis.
    _winner_src = m_a.get('winner_source', 'finsight_fallback')
    winner_raw = m_a.get('winner')
    is_neutral = (_winner_src == 'llm_neutral' or winner_raw is None)
    if is_neutral:
        winner = tkr_a  # placeholder pour les vars internes
        is_a = True
        winner_display = "Pas de préférence forte"
    else:
        winner = winner_raw
        is_a = (winner == tkr_a)
        winner_display = winner

    navy_bar(slide)
    footer_bar(slide)
    slide_title(slide, _cplbl("verdict_cmp"))
    section_dots(slide, 5)

    # Banderole CHOIX PRÉFÉRÉ — couleur ambre si neutre
    if is_neutral:
        w_fill = "B06000"  # ambre
    else:
        w_fill = COLOR_A if is_a else COLOR_B
    add_rect(slide, 1.02, 2.10, 23.37, 0.65, w_fill)
    add_text_box(slide, 1.27, 2.18, 22.86, 0.50,
                 f"CHOIX PRÉFÉRÉ : {winner_display}", 11, WHITE, bold=True, align=PP_ALIGN.CENTER)

    # ---- GAUCHE : Conviction box ----
    rec_w = (m_a.get('recommendation') if is_a else m_b.get('recommendation')) or 'BUY'
    conv_w = _safe_float(m_a.get('conviction') if is_a else m_b.get('conviction')) or 0
    if abs(conv_w) <= 2: conv_w *= 100
    fs_w = m_a.get('finsight_score') if is_a else m_b.get('finsight_score')
    fs_l = m_b.get('finsight_score') if is_a else m_a.get('finsight_score')
    try:
        delta = int(fs_w or 0) - int(fs_l or 0)
    except Exception:
        delta = 0
    gen_date = _fr_date_long()

    x_l = 1.02
    w_l = 8.5
    y_l = 3.10
    h_l = 4.2
    _conv_fill = GREEN_PALE if str(rec_w).upper() == 'BUY' else (RED_PALE if str(rec_w).upper() == 'SELL' else AMBER_PALE)
    _conv_line = GREEN if str(rec_w).upper() == 'BUY' else (RED if str(rec_w).upper() == 'SELL' else AMBER)

    add_rect(slide, x_l, y_l, w_l, h_l, _conv_fill)
    add_rect(slide, x_l, y_l, 0.20, h_l, _conv_line)
    # #209 : réhausse texts BUY/HOLD/SELL, Conviction, date dans leur box
    add_text_box(slide, x_l + 0.30, y_l + 0.10, w_l - 0.5, 1.30,
                 str(rec_w).upper(), 42, _conv_line, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x_l + 0.30, y_l + 1.45, w_l - 0.5, 0.50,
                 f"Conviction : {conv_w:.0f}%", 13, NAVY, align=PP_ALIGN.CENTER)
    add_text_box(slide, x_l + 0.30, y_l + 2.00, w_l - 0.5, 0.40,
                 gen_date, 9, GREY_TXT, align=PP_ALIGN.CENTER)

    # Barre de conviction (progress bar) — remontée dans le cadre
    add_text_box(slide, x_l + 0.40, y_l + 2.55, w_l - 0.7, 0.30,
                 "NIVEAU DE CONVICTION", 7, GREY_TXT, bold=True, align=PP_ALIGN.CENTER)
    _bar_y = y_l + 2.87
    add_rect(slide, x_l + 0.40, _bar_y, w_l - 0.80, 0.28, "DDDDDD")
    _fill_w = max(0.1, (w_l - 0.80) * (conv_w / 100))
    add_rect(slide, x_l + 0.40, _bar_y, _fill_w, 0.28, _conv_line)

    # Delta score — #209 : réhausse dans sa box
    _delta_str = f"Delta FinSight Score : {delta:+d} vs {tkr_b if is_a else tkr_a}"
    y_dlt = y_l + h_l + 0.25
    add_rect(slide, x_l, y_dlt, w_l, 0.55, NAVY_PALE)
    add_text_box(slide, x_l + 0.25, y_dlt + 0.08, w_l - 0.5, 0.32,
                 _delta_str, 9, NAVY, bold=True, align=PP_ALIGN.CENTER)

    # ---- DROITE : Thèse + Catalyseurs/Risques + Invalidation ----
    x_r = 9.85
    w_r = 14.55
    y_r = 3.10

    # Thèse d'investissement (header navy + body) — #209 : label réhaussé
    add_rect(slide, x_r, y_r, w_r, 0.55, NAVY)
    add_text_box(slide, x_r + 0.25, y_r + 0.05, w_r - 0.5, 0.40,
                 "THÈSE D'INVESTISSEMENT", 10, WHITE, bold=True)
    verdict_txt = synthesis.get('verdict_text') or \
                  f"{winner} est privilegiee sur la base des critères de valorisation, de qualité bilancielle et de momentum."
    _vt = " ".join(_fit(verdict_txt, 600).split())
    add_rect(slide, x_r, y_r + 0.55, w_r, 2.95, GREY_BG)
    add_text_box(slide, x_r + 0.25, y_r + 0.65, w_r - 0.5, 2.80, _vt, 8.5, BLACK, wrap=True)

    # Catalyseurs bulls + Risques bears — 2 colonnes
    y_cat = y_r + 0.55 + 3.05
    _cat_h = 2.35
    if is_neutral:
        # Cas neutre : pour/contre des 2 côtés au lieu du seul "winner"
        # Colonne 1 : tkr_a (bull A en vert + bear A en rouge au-dessous)
        add_rect(slide, x_r, y_cat, w_r/2 - 0.10, 0.50, COLOR_A)
        add_text_box(slide, x_r + 0.25, y_cat + 0.10, w_r/2 - 0.5, 0.40,
                     f"{tkr_a} — POUR / CONTRE", 9, WHITE, bold=True)
        _bull_a = synthesis.get('bull_a') or "Profil fondamental solide"
        _bear_a = synthesis.get('bear_a') or "Risques marché standard"
        add_rect(slide, x_r, y_cat + 0.50, w_r/2 - 0.10, _cat_h, GREY_BG)
        _text_a = f"+ {_fit(_bull_a, 130)}\n\n- {_fit(_bear_a, 130)}"
        add_text_box(slide, x_r + 0.20, y_cat + 0.60, w_r/2 - 0.40, _cat_h - 0.15,
                     _text_a, 8, BLACK, wrap=True)
        # Colonne 2 : tkr_b
        x_be = x_r + w_r/2 + 0.10
        add_rect(slide, x_be, y_cat, w_r/2 - 0.10, 0.50, COLOR_B)
        add_text_box(slide, x_be + 0.25, y_cat + 0.10, w_r/2 - 0.5, 0.40,
                     f"{tkr_b} — POUR / CONTRE", 9, WHITE, bold=True)
        _bull_b = synthesis.get('bull_b') or "Profil fondamental solide"
        _bear_b = synthesis.get('bear_b') or "Risques marché standard"
        add_rect(slide, x_be, y_cat + 0.50, w_r/2 - 0.10, _cat_h, GREY_BG)
        _text_b = f"+ {_fit(_bull_b, 130)}\n\n- {_fit(_bear_b, 130)}"
        add_text_box(slide, x_be + 0.20, y_cat + 0.60, w_r/2 - 0.40, _cat_h - 0.15,
                     _text_b, 8, BLACK, wrap=True)
    else:
        # Bulls (cas winner identifié) — #209 : labels réhaussés
        add_rect(slide, x_r, y_cat, w_r/2 - 0.10, 0.50, GREEN)
        add_text_box(slide, x_r + 0.25, y_cat + 0.05, w_r/2 - 0.5, 0.40,
                     "CATALYSEURS BULLS", 9, WHITE, bold=True)
        bull_winner = synthesis.get('bull_a' if is_a else 'bull_b') or "Profil fondamental solide"
        add_rect(slide, x_r, y_cat + 0.50, w_r/2 - 0.10, _cat_h, GREY_BG)
        add_text_box(slide, x_r + 0.20, y_cat + 0.55, w_r/2 - 0.40, _cat_h - 0.10,
                     _fit(bull_winner, 380), 8, BLACK, wrap=True)
        # Bears
        x_be = x_r + w_r/2 + 0.10
        add_rect(slide, x_be, y_cat, w_r/2 - 0.10, 0.50, RED)
        add_text_box(slide, x_be + 0.25, y_cat + 0.05, w_r/2 - 0.5, 0.40,
                     "RISQUES BEARS", 9, WHITE, bold=True)
        bear_winner = synthesis.get('bear_a' if is_a else 'bear_b') or "Risques de marché standard"
        add_rect(slide, x_be, y_cat + 0.50, w_r/2 - 0.10, _cat_h, GREY_BG)
        add_text_box(slide, x_be + 0.20, y_cat + 0.55, w_r/2 - 0.40, _cat_h - 0.10,
                     _fit(bear_winner, 380), 8, BLACK, wrap=True)

    # Conditions d'invalidation (warning bar)
    y_inv = y_cat + 0.50 + _cat_h + 0.25
    _inv_h = 13.10 - y_inv
    add_rect(slide, x_r, y_inv, w_r, _inv_h, AMBER_PALE)
    add_rect(slide, x_r, y_inv, 0.20, _inv_h, AMBER)
    add_text_box(slide, x_r + 0.30, y_inv + 0.10, w_r - 0.5, 0.35,
                 "CONDITIONS D'INVALIDATION", 8, AMBER, bold=True)
    inv_txt = (
        f"Détérioration structurelle des marges EBITDA > 200bps, "
        f"choc macro adverses impactant la demande finale, "
        f"ou re-rating multiples > 20% sans catalyseur fondamental correspondant. "
        f"A surveiller : publications trimestrielles, guidance FY+1, événements M&A sectoriels."
    )
    add_text_box(slide, x_r + 0.30, y_inv + 0.50, w_r - 0.5, _inv_h - 0.65,
                 " ".join(inv_txt.split()), 8, NAVY, wrap=True)

    return slide


# ---------------------------------------------------------------------------
# Point d'entree principal
# ---------------------------------------------------------------------------

class CmpSocietePPTXWriter:
    """
    Genere un pitchbook comparatif IB-quality en 22 slides.

    Usage :
        writer = CmpSocietePPTXWriter()
        path   = writer.generate(state_a, state_b, output_path)
    """

    def generate(self, state_a: dict, state_b: dict,
                 output_path: Optional[str] = None,
                 language: str = "fr", currency: str = "EUR") -> str:
        """
        Genere le fichier PPTX comparatif.

        Args:
            state_a, state_b : etats pipeline LangGraph
            output_path      : chemin de sortie (auto-genere si None)
            language, currency : kwargs i18n (propagés au module)

        Returns:
            Chemin du fichier genere.
        """
        global _CMP_PPTX_LANG
        _CMP_PPTX_LANG = (language or state_a.get("language") or "fr").lower()[:2]
        if _CMP_PPTX_LANG not in {"fr","en","es","de","it","pt"}:
            _CMP_PPTX_LANG = "fr"
        from pptx import Presentation
        from pptx.util import Cm

        # 1. Build context (dedupe refactor #191 — flow partagé avec xlsx/pdf)
        from outputs.cmp_societe_common import build_cmp_context
        ctx = build_cmp_context(state_a, state_b)
        tkr_a     = ctx.tkr_a
        tkr_b     = ctx.tkr_b
        supp_a    = ctx.supp_a
        supp_b    = ctx.supp_b
        m_a       = ctx.m_a
        m_b       = ctx.m_b
        synthesis = ctx.synthesis

        # Sync verdict_relative (encore lu par certains templates PPTX)
        if ctx.winner:
            m_a["verdict_relative"] = m_b["verdict_relative"] = f"{ctx.winner} privilege"
        else:
            m_a["verdict_relative"] = m_b["verdict_relative"] = "Équivalent"

        # 3. Créer la présentation
        prs = Presentation()
        prs.slide_width  = Cm(25.4)
        prs.slide_height = Cm(14.29)

        # Slide 1 — Cover
        _slide_cover(prs, m_a, m_b)
        # Slide 2 — Executive Summary
        _slide_exec_summary(prs, m_a, m_b, synthesis)
        # Slide 3 — Sommaire
        _slide_sommaire(prs, tkr_a, tkr_b)
        # Slide 4 — Divider 01
        divider_slide(prs, "01", "Profil & Identité",
                      f"Présentation comparative de {tkr_a} et {tkr_b}")
        # Slide 5 — Profil comparatif
        _slide_profil(prs, m_a, m_b)
        # Slide 6 — Divider 02
        divider_slide(prs, "02", "Performance Financière",
                      "Compte de Résultat, Marges, Bilan & Liquidité")
        # Slide 7 — P&L comparé
        _slide_pl(prs, m_a, m_b, synthesis)
        # Slide 8 — Marges & Rentabilité
        _slide_marges(prs, m_a, m_b, synthesis)
        # Slide 9 — Bilan & Liquidité
        _slide_bilan(prs, m_a, m_b)
        # Slide 10 — Divider 03
        divider_slide(prs, "03", "Valorisation",
                      "Multiples, DCF, Monte Carlo et Football Field")
        # Slide 11 — Multiples
        _slide_multiples(prs, m_a, m_b, synthesis)
        # Slide 12 — DCF
        _slide_dcf(prs, m_a, m_b)
        # Slide 13 — Monte Carlo + Football Field
        _slide_monte_carlo(prs, m_a, m_b)
        # Slide 14 — Divider 04
        divider_slide(prs, "04", "Qualité & Risques",
                      "Piotroski, Beneish, Altman, Beta, VaR et Momentum")
        # Slide 15 — Qualité financière
        _slide_piotroski(prs, m_a, m_b, synthesis)
        # Slide 16 — Risque & Momentum
        _slide_risque(prs, m_a, m_b)
        # Slide 17 — Divider 05
        divider_slide(prs, "05", "Verdict",
                      "FinSight Score, Thèses et Recommandation Finale")
        # Slide 18 — FinSight Score
        _slide_finsight_score(prs, m_a, m_b)
        # Slide 19 — Thèses Bull/Bear
        _slide_theses(prs, m_a, m_b, synthesis)
        # Slide 20 — Cours boursiers 52 semaines
        _slide_price_chart(prs, m_a, m_b, synthesis)
        # Slide 21 — Verdict final
        _slide_verdict(prs, m_a, m_b, synthesis)

        # 4. Sauvegarder
        if output_path is None:
            out_dir = Path(__file__).parent / "generated" / "cli_tests"
            out_dir.mkdir(parents=True, exist_ok=True)
            output_path = str(out_dir / f"{tkr_a}_vs_{tkr_b}_comparison.pptx")

        prs.save(output_path)
        log.info(f"[cmp_pptx] fichier sauvegarde : {output_path}")
        return output_path

    def generate_bytes(self, state_a: dict, state_b: dict) -> bytes:
        """Généré le fichier en memoire et retourne les bytes."""
        from pptx import Presentation
        from pptx.util import Cm
        import tempfile, os

        tmp = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
        tmp.close()
        try:
            self.generate(state_a, state_b, tmp.name)
            with open(tmp.name, 'rb') as f:
                return f.read()
        finally:
            try:
                os.unlink(tmp.name)
            except Exception as _e:
                log.debug(f"[cmp_societe_pptx_writer:generate_bytes] exception skipped: {_e}")
