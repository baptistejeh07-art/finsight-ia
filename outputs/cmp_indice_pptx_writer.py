# =============================================================================
# FinSight IA — Indice Comparison PPTX Writer (21 slides IB-grade)
# outputs/indice_comparison_pptx_writer.py
#
# Pitchbook comparatif analyse d'indices boursiers.
# Format aligné sur cmp société / cmp secteur :
#   - Sommaire pixel-près
#   - Box LLM standardisée (NAVY_PALE + barre latérale)
#   - Bullets carrés colorés
#   - Texte LLM enrichi (12 champs)
#   - Cadrage strict, accents partout
#
# Usage :
#   from outputs.cmp_indice_pptx_writer import CmpIndicePPTXWriter
#   pptx_bytes = CmpIndicePPTXWriter.generate(data)
#
# Données attendues dans `data` :
#   name_a/b, code_a/b, ticker_a/b, currency_a/b, date
#   perf_ytd_a/b, perf_1y_a/b, perf_3y_a/b, perf_5y_a/b
#   vol_1y_a/b, sharpe_1y_a/b, max_dd_a/b
#   pe_fwd_a/b, pb_a/b, div_yield_a/b, erp_a/b
#   score_a/b, signal_a/b
#   sector_comparison (list of (sector, weight_a, weight_b))
#   top5_a/b (list of (company, ticker, weight, sector))
#   perf_history (dict dates/indice_a/indice_b base 100)
# =============================================================================
from __future__ import annotations

import io
import logging
import re as _re
from typing import Optional

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import numpy as np

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

log = logging.getLogger(__name__)

# ── Dimensions ────────────────────────────────────────────────────────────────
_SW = Cm(25.4)
_SH = Cm(14.3)

# ── Palette (alignée cmp société / cmp secteur) ─────────────────────────────
_NAVY      = RGBColor(0x1B, 0x3A, 0x6B)
_NAVY_MID  = RGBColor(0x2E, 0x5F, 0xA3)
_NAVY_PALE = RGBColor(0xEE, 0xF3, 0xFA)
_NAVYL     = _NAVY_MID
_NAVYP     = _NAVY_PALE
_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
_GRAYL     = RGBColor(0xF5, 0xF7, 0xFA)
_GRAYM     = RGBColor(0xE8, 0xEC, 0xF0)
_GRAYT     = RGBColor(0x55, 0x55, 0x55)
_GRAYD     = RGBColor(0xAA, 0xAA, 0xAA)
_BLACK     = RGBColor(0x1A, 0x1A, 0x1A)
_BUY       = RGBColor(0x1A, 0x7A, 0x4A)
_SELL      = RGBColor(0xA8, 0x20, 0x20)
_HOLD      = RGBColor(0xB0, 0x60, 0x00)
_BUY_L     = RGBColor(0xE8, 0xF5, 0xEE)
_SELL_L    = RGBColor(0xFB, 0xEB, 0xEB)
_HOLD_L    = RGBColor(0xFD, 0xF3, 0xE5)
_COL_A     = RGBColor(0x2E, 0x5F, 0xA3)
_COL_B     = RGBColor(0x1A, 0x7A, 0x4A)
_COL_AL    = RGBColor(0xEE, 0xF3, 0xFA)
_COL_BL    = RGBColor(0xEA, 0xF4, 0xEF)
_GOLD      = RGBColor(0xC9, 0xA2, 0x27)
_RED       = _SELL

_MONTHS_FR = {1: "janvier", 2: "février", 3: "mars", 4: "avril", 5: "mai", 6: "juin",
              7: "juillet", 8: "août", 9: "septembre", 10: "octobre", 11: "novembre", 12: "décembre"}


# ── Helpers texte ─────────────────────────────────────────────────────────────

def _x(text) -> str:
    """Échappe les caractères de contrôle et retire le markdown LLM (** *)."""
    if text is None:
        return ""
    s = _re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', str(text))
    s = _re.sub(r'\*\*(.+?)\*\*', r'\1', s, flags=_re.DOTALL)
    s = _re.sub(r'\*(.+?)\*', r'\1', s, flags=_re.DOTALL)
    s = _re.sub(r'\*+', '', s)
    return s


def _fr_pct(v, signed=False) -> str:
    if v is None:
        return "\u2014"
    try:
        fv = float(v)
        if abs(fv) > 2:
            pass  # déjà en %
        else:
            fv = fv * 100
        s = f"{fv:+.1f}" if signed else f"{fv:.1f}"
        return s.replace(".", ",") + "\u00a0%"
    except Exception:
        return "\u2014"


def _fr_pct_signed(v) -> str:
    return _fr_pct(v, signed=True)


def _fr_num(v, dp=1) -> str:
    if v is None:
        return "\u2014"
    try:
        return f"{float(v):.{dp}f}".replace(".", ",")
    except Exception:
        return "\u2014"


def _fr_x(v, dp=1) -> str:
    """Comme _fr_num mais ajouté le suffixe 'x' pour les multiples.
    Retourne '—' (sans x) si la valeur est manquante."""
    if v is None:
        return "\u2014"
    try:
        return f"{float(v):.{dp}f}".replace(".", ",") + "x"
    except Exception:
        return "\u2014"


def _fr_date() -> str:
    import datetime
    d = datetime.date.today()
    return f"{d.day} {_MONTHS_FR[d.month]} {d.year}"


def _sig_c(signal: str) -> RGBColor:
    s = str(signal)
    if "Surp" in s or "Positif" in s:
        return _BUY
    if "Sous" in s or "Négatif" in s or "Négatif" in s:
        return _SELL
    return _HOLD


def _sig_l(signal: str) -> RGBColor:
    s = str(signal)
    if "Surp" in s or "Positif" in s:
        return _BUY_L
    if "Sous" in s or "Négatif" in s or "Négatif" in s:
        return _SELL_L
    return _HOLD_L


def _safe_float(v):
    if v is None:
        return None
    try:
        f = float(v)
        return f if not (f != f) else None
    except Exception:
        return None


def _ecart_pct(va, vb) -> str:
    fa = _safe_float(va)
    fb = _safe_float(vb)
    if fa is None or fb is None:
        return "\u2014"
    if abs(fa) <= 2:
        fa *= 100
    if abs(fb) <= 2:
        fb *= 100
    diff = fa - fb
    return f"{diff:+.1f}".replace(".", ",") + "\u00a0pp"


def _ecart_mult(va, vb) -> str:
    fa = _safe_float(va)
    fb = _safe_float(vb)
    if fa is None or fb is None:
        return "\u2014"
    diff = fa - fb
    return f"{diff:+.1f}".replace(".", ",") + "x"


def _favorise_lower(va, vb, na, nb) -> str:
    fa = _safe_float(va)
    fb = _safe_float(vb)
    if fa is None or fb is None:
        return "\u2014"
    if fa < fb:
        return na[:14]
    if fb < fa:
        return nb[:14]
    return "Égalité"


def _favorise_higher(va, vb, na, nb) -> str:
    fa = _safe_float(va)
    fb = _safe_float(vb)
    if fa is None or fb is None:
        return "\u2014"
    if abs(fa) <= 1:
        fa *= 100
    if abs(fb) <= 1:
        fb *= 100
    if fa > fb:
        return na[:14]
    if fb > fa:
        return nb[:14]
    return "Égalité"


# ── Helpers PPTX primitives ───────────────────────────────────────────────────

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, x, y, w, h, fill=None, line_col=None, line_w=0.5):
    shape = slide.shapes.add_shape(1, Cm(x), Cm(y), Cm(w), Cm(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_col:
        shape.line.color.rgb = line_col
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape


def _txb(slide, text, x, y, w, h, size=9, bold=False, color=None,
         align=PP_ALIGN.LEFT, italic=False, wrap=True):
    color = color or _BLACK
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = _x(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def _pic(slide, img_bytes, x, y, w, h):
    buf = io.BytesIO(img_bytes)
    slide.shapes.add_picture(buf, Cm(x), Cm(y), Cm(w), Cm(h))


def _table(slide, data, x, y, w, h, col_widths=None,
           hdr_fill=_NAVY, font_size=8.5):
    rows, cols = len(data), len(data[0]) if data else 1
    tbl_shape = slide.shapes.add_table(rows, cols, Cm(x), Cm(y), Cm(w), Cm(h))
    tbl = tbl_shape.table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Cm(cw * w / total)
    for r, row_data in enumerate(data):
        is_hdr = (r == 0)
        bg = hdr_fill if is_hdr else (_GRAYL if r % 2 == 0 else _WHITE)
        for c, val in enumerate(row_data):
            cell = tbl.cell(r, c)
            cell.text = _x(val) if val is not None else "\u2014"
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            run = para.runs[0] if para.runs else para.add_run()
            run.font.size = Pt(font_size)
            run.font.bold = is_hdr
            run.font.color.rgb = _WHITE if is_hdr else _BLACK
            if is_hdr:
                cell.fill.solid()
                cell.fill.fore_color.rgb = hdr_fill
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg
    return tbl


# ── Box LLM standardisée (style cmp société / cmp secteur) ──────────────────

def _llm_box(slide, x, y, w, h, title, text, *, fontsize=9.0):
    """Box LLM standard : NAVY_PALE + barre latérale NAVY_MID + titre + texte."""
    _rect(slide, x, y, w, h, _NAVY_PALE)
    _rect(slide, x, y, 0.13, h, _NAVY_MID)
    _txb(slide, title.upper(), x + 0.30, y + 0.10, w - 0.40, 0.36,
         size=fontsize, bold=True, color=_NAVY)
    _txb(slide, text or "", x + 0.30, y + 0.50, w - 0.50, h - 0.62,
         size=fontsize, color=_NAVY, wrap=True)


def _bullet_square(slide, x, y, color, size=0.22):
    """Petit carré coloré servant de bullet."""
    _rect(slide, x, y, size, size, fill=color)


def _llm_box_with_bullets(slide, x, y, w, h, title, items, *, color, fontsize=9.0):
    """Box LLM avec bullets carrés à gauche pour chaque item.

    items : liste de tuples (heading, body) ou de strings.
    """
    _rect(slide, x, y, w, h, _NAVY_PALE)
    _rect(slide, x, y, 0.13, h, _NAVY_MID)
    _txb(slide, title.upper(), x + 0.30, y + 0.10, w - 0.40, 0.36,
         size=fontsize, bold=True, color=_NAVY)
    if not items:
        return
    n = max(1, len(items))
    avail = h - 0.65
    item_h = avail / n
    for i, it in enumerate(items[:5]):
        iy = y + 0.55 + i * item_h
        _bullet_square(slide, x + 0.35, iy + 0.12, color, size=0.20)
        if isinstance(it, tuple) and len(it) == 2:
            head, body = it
            _txb(slide, str(head)[:60], x + 0.70, iy, w - 0.85, 0.36,
                 size=fontsize, bold=True, color=_NAVY, wrap=False)
            _txb(slide, str(body), x + 0.70, iy + 0.36, w - 0.85, item_h - 0.40,
                 size=fontsize - 0.5, color=_NAVY, wrap=True)
        else:
            _txb(slide, str(it), x + 0.70, iy + 0.05, w - 0.85, item_h - 0.10,
                 size=fontsize, color=_NAVY, wrap=True)


# ── Helpers cadre standard pour slides ─────────────────────────────────────

def _header(slide, title, subtitle="", section_active=0):
    """Header navy pleine largeur + sous-titre + dots de navigation à droite."""
    _rect(slide, 0, 0, 25.4, 1.65, _NAVY)
    _txb(slide, title, 1.02, 0.28, 18.0, 0.97, 13, True, _WHITE)
    if subtitle:
        # Subtitle juste sous la navy bar — rehaussé pour éviter chevauchement
        _txb(slide, subtitle, 1.02, 1.55, 23.37, 0.40, 9, False, _NAVY_MID, italic=True)
    # Dots section navigation 1-4
    if section_active:
        for i in range(1, 5):
            x = 19.6 + (i - 1) * 0.85
            fill = _WHITE if i == section_active else _NAVY_MID
            txt_col = _NAVY if i == section_active else _WHITE
            _rect(slide, x, 0.45, 0.65, 0.65, fill)
            _txb(slide, str(i), x, 0.50, 0.65, 0.55, 8.5, True, txt_col,
                 align=PP_ALIGN.CENTER, wrap=False)


def _footer(slide):
    _rect(slide, 0, 13.39, 25.4, 0.91, _NAVY)
    _txb(slide, "FinSight IA  \u00b7  Comparatif Indices  \u00b7  Usage confidentiel",
         1.02, 13.47, 23.37, 0.50, 7, False, _WHITE)


def _index_band(slide, name_a, name_b, *, y=2.30):
    """Bande bicolore identifiant les deux indices — noms centrés."""
    _rect(slide, 1.02, y, 11.44, 0.62, _COL_AL)
    _rect(slide, 1.02, y, 0.15, 0.62, _COL_A)
    _txb(slide, name_a, 1.17, y + 0.10, 11.29, 0.46, 10, True, _NAVY,
         align=PP_ALIGN.CENTER)
    _rect(slide, 12.94, y, 11.44, 0.62, _COL_BL)
    _rect(slide, 12.94, y, 0.15, 0.62, _COL_B)
    _txb(slide, name_b, 13.09, y + 0.10, 11.29, 0.46, 10, True, _BUY,
         align=PP_ALIGN.CENTER)


def _divider(prs, num_str, title, subtitle):
    """Slide divider de section avec fond navy et numéro grand."""
    slide = _blank(prs)
    from pptx.oxml.ns import qn
    from lxml import etree
    cSld = slide._element.find(qn('p:cSld'))
    existing = cSld.find(qn('p:bg'))
    if existing is not None:
        cSld.remove(existing)
    bg_xml = (
        '<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<p:bgPr><a:solidFill><a:srgbClr val="1B3A6B"/></a:solidFill>'
        '<a:effectLst/></p:bgPr></p:bg>'
    )
    bg_elem = etree.fromstring(bg_xml)
    spTree = cSld.find(qn('p:spTree'))
    cSld.insert(list(cSld).index(spTree) if spTree is not None else 0, bg_elem)

    _rect(slide, 0, 0, 0.3, 14.3, _NAVYL)

    # Numéro filigrane (transparence)
    txb = slide.shapes.add_textbox(Cm(1.27), Cm(3.5), Cm(22.86), Cm(4.57))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = _x(num_str)
    run.font.size = Pt(96)
    run.font.bold = True
    from pptx.oxml.ns import qn as _qn
    from lxml import etree as _et
    r_elem = run._r
    rPr = r_elem.find(_qn('a:rPr'))
    if rPr is None:
        rPr = _et.Element(_qn('a:rPr'))
        r_elem.insert(0, rPr)
    solidFill = _et.SubElement(rPr, _qn('a:solidFill'))
    srgbClr = _et.SubElement(solidFill, _qn('a:srgbClr'))
    srgbClr.set('val', 'FFFFFF')
    alpha_elem = _et.SubElement(srgbClr, _qn('a:alpha'))
    alpha_elem.set('val', '15000')
    rPr.remove(solidFill)
    rPr.insert(0, solidFill)

    _txb(slide, title, 1.27, 5.5, 22.86, 1.8, 32, True, _WHITE)
    _rect(slide, 1.27, 8.0, 7.62, 0.05, _GRAYD)
    _txb(slide, subtitle, 1.27, 8.3, 22.86, 0.95, 11, False,
         RGBColor(0xAA, 0xBB, 0xDD))
    _txb(slide, "FinSight IA  \u00b7  Comparatif Indices  \u00b7  Usage confidentiel",
         1.02, 13.50, 23.37, 0.56, 7, False,
         RGBColor(0x66, 0x77, 0xAA))
    return slide


def _kpi_box(slide, x, y, w, h, value, label, sub="", color_accent=_NAVYL):
    """KPI box avec barre latérale colorée."""
    _rect(slide, x, y, w, h, _NAVY_PALE)
    _rect(slide, x, y, 0.15, h, color_accent)
    _txb(slide, value, x + 0.27, y + 0.10, w - 0.40, h * 0.55,
         16, True, _NAVY, align=PP_ALIGN.CENTER)
    _txb(slide, label, x + 0.27, y + h * 0.55 + 0.05, w - 0.40, h * 0.25,
         8.0, True, _GRAYT, align=PP_ALIGN.CENTER)
    if sub:
        _txb(slide, sub, x + 0.27, y + h * 0.78, w - 0.40, h * 0.20,
             7, False, _GRAYT, align=PP_ALIGN.CENTER)


def _chart_buf(fig) -> bytes:
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=160, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close(fig)
    buf.seek(0)
    return buf.read()


# ═════════════════════════════════════════════════════════════════════════════
# SLIDES
# ═════════════════════════════════════════════════════════════════════════════

# ── Slide 1 — Cover ──────────────────────────────────────────────────────────

def _slide_01_cover(prs, d: dict):
    slide = _blank(prs)
    name_a = d.get("name_a", "Indice A")
    name_b = d.get("name_b", "Indice B")
    sig_a = d.get("signal_a", "Neutre")
    sig_b = d.get("signal_b", "Neutre")
    sc_a = d.get("score_a", 50)
    sc_b = d.get("score_b", 50)
    cur_a = d.get("currency_a", "USD")
    cur_b = d.get("currency_b", "USD")
    date_str = d.get("date", _fr_date())

    # Bandeau navy supérieur
    _rect(slide, 0, 0, 25.4, 1.85, _NAVY)
    _txb(slide, "FinSight IA", 0, 0.18, 25.4, 0.65, 10, False,
         RGBColor(0x88, 0x99, 0xBB), align=PP_ALIGN.CENTER)
    _rect(slide, 10.7, 0.95, 4.0, 0.04, RGBColor(0xAA, 0xBB, 0xCC))
    _txb(slide, "Pitchbook  \u2014  Comparatif d'Indices",
         0, 1.05, 25.4, 0.70, 9, False,
         RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER)

    # Titre principal
    _txb(slide, f"{name_a}  vs  {name_b}", 0, 3.8, 25.4, 2.5,
         34, True, _NAVY, align=PP_ALIGN.CENTER)

    # Sous-titre période
    _txb(slide, f"Analyse comparative des principaux indices boursiers",
         0, 6.5, 25.4, 0.8, 12, False, _GRAYT, align=PP_ALIGN.CENTER)

    # 4 KPI boxes : Score A, Score B, Perf 1Y A, Perf 1Y B
    p1a = _fr_pct_signed(d.get("perf_1y_a"))
    p1b = _fr_pct_signed(d.get("perf_1y_b"))
    kpis = [
        (f"{sc_a}", f"Score FinSight\n{name_a}", _COL_A, _COL_AL),
        (f"{sc_b}", f"Score FinSight\n{name_b}", _COL_B, _COL_BL),
        (p1a,        f"Performance 1 an\n{name_a}",  _COL_A, _COL_AL),
        (p1b,        f"Performance 1 an\n{name_b}",  _COL_B, _COL_BL),
    ]
    box_w = 5.6
    step = 5.85
    start_x = 0.9
    for i, (val, lbl, col, pale) in enumerate(kpis):
        bx = start_x + i * step
        _rect(slide, bx, 8.0, box_w, 2.4, pale)
        _rect(slide, bx, 8.0, 0.22, 2.4, col)
        _txb(slide, val, bx + 0.22, 8.15, box_w - 0.22, 1.35,
             24, True, col, align=PP_ALIGN.CENTER)
        _txb(slide, lbl, bx + 0.22, 9.55, box_w - 0.22, 0.85,
             8.5, False, _NAVY, align=PP_ALIGN.CENTER)

    # Signal badges
    _rect(slide, 2.2, 11.0, 4.5, 0.9, _sig_c(sig_a))
    _txb(slide, f"\u25cf {sig_a}", 2.2, 11.08, 4.5, 0.75,
         9.5, True, _WHITE, align=PP_ALIGN.CENTER)
    _rect(slide, 18.7, 11.0, 4.5, 0.9, _sig_c(sig_b))
    _txb(slide, f"\u25cf {sig_b}", 18.7, 11.08, 4.5, 0.75,
         9.5, True, _WHITE, align=PP_ALIGN.CENTER)

    # ── Bandeau FRED macro (one-liner) ──────────────────────────────────────
    try:
        from data.sources.fred_source import fetch_macro_context
        _fred = fetch_macro_context()
        if _fred:
            _parts = []
            _fred_map = [
                ("fed_funds_rate", "Fed"),
                ("treasury_10y",  "10Y"),
                ("vix",           "VIX"),
                ("cpi_yoy",       "CPI"),
            ]
            for _fk, _flbl in _fred_map:
                _fv = _fred.get(_fk)
                if _fv is not None:
                    if _fk == "cpi_yoy":
                        _parts.append(f"{_flbl} {_fv:+.1f}\u00a0%".replace(".", ","))
                    elif _fk == "vix":
                        _parts.append(f"{_flbl} {_fv:.1f}".replace(".", ","))
                    else:
                        _parts.append(f"{_flbl} {_fv:.2f}\u00a0%".replace(".", ","))
            if _parts:
                _fred_line = "Données FRED : " + "  \u00b7  ".join(_parts)
                _rect(slide, 0.9, 12.30, 23.6, 0.55, _NAVY_PALE)
                _txb(slide, _fred_line, 0.9, 12.35, 23.6, 0.45,
                     7.5, False, _NAVY, align=PP_ALIGN.CENTER)
    except Exception as exc:
        log.warning("[cmp_indice_pptx] FRED one-liner skipped: %s", exc)

    # Footer
    _rect(slide, 0, 13.6, 25.4, 0.7, _GRAYL)
    _txb(slide, "Rapport confidentiel  |  FinSight IA",
         0.9, 13.65, 12.0, 0.5, 7.5, False, _GRAYT)
    _txb(slide, date_str, 0, 13.65, 24.4, 0.5, 7.5, False, _GRAYT,
         align=PP_ALIGN.RIGHT)


# ── Slide 2 — Executive Summary ──────────────────────────────────────────────

def _slide_02_exec_summary(prs, d: dict):
    slide = _blank(prs)
    name_a = d.get("name_a", "Indice A")
    name_b = d.get("name_b", "Indice B")
    sig_a = d.get("signal_a", "Neutre")
    sig_b = d.get("signal_b", "Neutre")
    sc_a = d.get("score_a", 50)
    sc_b = d.get("score_b", 50)

    _header(slide, f"{(data.get('_t_helper') or (lambda k: "Executive Summary"))('exec_sum')}  —  {name_a} vs {name_b}",
            f"Indicateurs clés et signal d'allocation", section_active=1)
    _footer(slide)

    # Bandeau Signal A
    _rect(slide, 0.9, 2.10, 11.4, 1.05, _sig_l(sig_a))
    _rect(slide, 0.9, 2.10, 0.12, 1.05, _sig_c(sig_a))
    _txb(slide, f"{name_a}  —  {sig_a}", 1.30, 2.20, 10.6, 0.40,
         10, True, _sig_c(sig_a))
    _txb(slide,
         f"Score {sc_a}/100  |  P/E Fwd {_fr_x(d.get('pe_fwd_a'), 1)}  |  "
         f"Vol {_fr_num(d.get('vol_1y_a'))} %  |  Sharpe {_fr_num(d.get('sharpe_1y_a'), 2)}",
         1.30, 2.62, 10.6, 0.45, 8, False, _NAVY)

    # Bandeau Signal B
    _rect(slide, 13.1, 2.10, 11.4, 1.05, _sig_l(sig_b))
    _rect(slide, 13.1, 2.10, 0.12, 1.05, _sig_c(sig_b))
    _txb(slide, f"{name_b}  —  {sig_b}", 13.50, 2.20, 10.6, 0.40,
         10, True, _sig_c(sig_b))
    _txb(slide,
         f"Score {sc_b}/100  |  P/E Fwd {_fr_x(d.get('pe_fwd_b'), 1)}  |  "
         f"Vol {_fr_num(d.get('vol_1y_b'))} %  |  Sharpe {_fr_num(d.get('sharpe_1y_b'), 2)}",
         13.50, 2.62, 10.6, 0.45, 8, False, _BUY)

    # Tableau comparatif des métriques clés
    def _winner_eq(va, vb, name_a, name_b, higher=True):
        """Retourne le gagnant ou '—' si scores egaux."""
        if va is None or vb is None or va == vb:
            return "\u2014"
        if higher:
            return name_a if va > vb else name_b
        return name_a if va < vb else name_b

    rows = [
        ["Indicateur", name_a[:18], name_b[:18], "Avantage"],
        ["Score FinSight", f"{sc_a}/100", f"{sc_b}/100",
         _winner_eq(sc_a, sc_b, name_a, name_b, higher=True)],
        ["Performance YTD",
         _fr_pct_signed(d.get("perf_ytd_a")),
         _fr_pct_signed(d.get("perf_ytd_b")),
         _favorise_higher(d.get("perf_ytd_a"), d.get("perf_ytd_b"), name_a, name_b)],
        ["Performance 1 an",
         _fr_pct_signed(d.get("perf_1y_a")),
         _fr_pct_signed(d.get("perf_1y_b")),
         _favorise_higher(d.get("perf_1y_a"), d.get("perf_1y_b"), name_a, name_b)],
        ["Performance 3 ans",
         _fr_pct_signed(d.get("perf_3y_a")),
         _fr_pct_signed(d.get("perf_3y_b")),
         _favorise_higher(d.get("perf_3y_a"), d.get("perf_3y_b"), name_a, name_b)],
        ["Volatilité 1 an",
         _fr_num(d.get("vol_1y_a")) + " %",
         _fr_num(d.get("vol_1y_b")) + " %",
         _favorise_lower(d.get("vol_1y_a"), d.get("vol_1y_b"), name_a, name_b)],
        ["Sharpe Ratio 1 an",
         _fr_num(d.get("sharpe_1y_a"), 2),
         _fr_num(d.get("sharpe_1y_b"), 2),
         _favorise_higher(d.get("sharpe_1y_a"), d.get("sharpe_1y_b"), name_a, name_b)],
        ["Max Drawdown",
         _fr_pct_signed(d.get("max_dd_a")),
         _fr_pct_signed(d.get("max_dd_b")),
         _favorise_higher(d.get("max_dd_a"), d.get("max_dd_b"), name_a, name_b)],
        ["P/E Forward",
         _fr_x(d.get("pe_fwd_a"), 1),
         _fr_x(d.get("pe_fwd_b"), 1),
         _favorise_lower(d.get("pe_fwd_a"), d.get("pe_fwd_b"), name_a, name_b)],
        ["P/B (Book Value)",
         _fr_x(d.get("pb_a"), 1),
         _fr_x(d.get("pb_b"), 1),
         _favorise_lower(d.get("pb_a"), d.get("pb_b"), name_a, name_b)],
        ["Rendement dividende",
         _fr_pct(d.get("div_yield_a")) if d.get("div_yield_a") else "\u2014",
         _fr_pct(d.get("div_yield_b")) if d.get("div_yield_b") else "\u2014",
         _favorise_higher(d.get("div_yield_a"), d.get("div_yield_b"), name_a, name_b)],
    ]
    _table(slide, rows, 0.9, 3.40, 23.6, 7.95,
           col_widths=[5.5, 3.5, 3.5, 3.5], font_size=8.5)

    # Box LLM standard
    llm_text = d.get("llm", {}).get("exec_summary", "")
    if llm_text:
        _llm_box(slide, 0.9, 11.55, 23.6, 1.75,
                 "Analyse Comparative", llm_text, fontsize=9)


# ── Slide 3 — Sommaire (style cmp société pixel-près) ───────────────────────

def _slide_03_sommaire(prs, d: dict):
    slide = _blank(prs)
    name_a = d.get("name_a", "Indice A")
    name_b = d.get("name_b", "Indice B")
    _header(slide, (data.get('_t_helper') or (lambda k: "Sommaire"))('sommaire'),
            f"Analyse comparative : {name_a} vs {name_b}", section_active=1)
    _footer(slide)

    sections = [
        ("01", "Profil & Composition",     "Structure des indices, secteurs, top constituants",                     "Slides 5-8"),
        ("02", "Performance & Risque",     "Cours 52S base 100, décomposition perfs, vol/Sharpe/drawdown",         "Slides 10-13"),
        ("03", "Valorisation & Qualité",   "P/E, P/B, dividendes, ERP et signal FinSight",                          "Slides 15-17"),
        ("04", "Décision d'allocation",    "Thèses bull/bear, recommandation, verdict et invalidation",             "Slides 19-21"),
    ]

    y_start = 2.40
    card_h = 2.55
    card_gap = 0.20
    card_w = 23.6
    x_card = 0.9

    for i, (num, title, desc, pages) in enumerate(sections):
        y = y_start + i * (card_h + card_gap)
        _rect(slide, x_card, y, card_w, card_h, _NAVY_PALE)
        _rect(slide, x_card, y, 0.18, card_h, _NAVY)
        _txb(slide, num, x_card + 0.40, y + 0.30, 3.20, 1.95,
             44, True, _NAVY_MID, align=PP_ALIGN.CENTER)
        _txb(slide, title, x_card + 4.00, y + 0.45, card_w - 8.0, 0.80,
             15, True, _NAVY)
        _txb(slide, desc, x_card + 4.00, y + 1.30, card_w - 8.0, 1.00,
             10, False, _GRAYT, wrap=True)
        _rect(slide, x_card + card_w - 3.40, y + 0.85, 3.00, 0.85, _NAVY)
        _txb(slide, pages, x_card + card_w - 3.40, y + 1.05, 3.00, 0.45,
             10, True, _WHITE, align=PP_ALIGN.CENTER)


# ── Slide 5 — Profil Indice A ────────────────────────────────────────────────

def _slide_05_profil_a(prs, d: dict):
    _slide_profil(prs, d, which="a")


def _slide_06_profil_b(prs, d: dict):
    _slide_profil(prs, d, which="b")


def _slide_profil(prs, d: dict, which: str):
    """Profil détaillé d'un indice (A ou B)."""
    slide = _blank(prs)
    is_a = (which == "a")
    name = d.get("name_a" if is_a else "name_b", "Indice")
    code = d.get("code_a" if is_a else "code_b", "")
    cur = d.get("currency_a" if is_a else "currency_b", "USD")
    sig = d.get("signal_a" if is_a else "signal_b", "Neutre")
    sc = d.get("score_a" if is_a else "score_b", 50)
    col = _COL_A if is_a else _COL_B
    col_pale = _COL_AL if is_a else _COL_BL

    _header(slide, f"{(data.get('_t_helper') or (lambda k: "Profil"))('profil')}  —  {name}",
            f"Composition, métriques clés et caractéristiques structurelles",
            section_active=1)
    _footer(slide)

    # Bandeau identité
    _rect(slide, 1.02, 2.10, 23.37, 0.70, col)
    _txb(slide, f"{name}  ({code})  \u00b7  Devise : {cur}", 1.20, 2.22,
         23.0, 0.50, 11, True, _WHITE)

    # KPIs en ligne (4 boxes)
    p1y = d.get("perf_1y_a" if is_a else "perf_1y_b")
    vol = d.get("vol_1y_a" if is_a else "vol_1y_b")
    pe = d.get("pe_fwd_a" if is_a else "pe_fwd_b")
    dy = d.get("div_yield_a" if is_a else "div_yield_b")

    kpis = [
        (f"{sc}/100", "Score FinSight", sig),
        (_fr_pct_signed(p1y), "Performance 1 an", "vs 5 ans"),
        (_fr_x(pe, 1), "P/E Forward", "Médiane des constituants"),
        (_fr_pct(dy) if dy else "\u2014", "Rendement dividende", "Médiane pondérée"),
    ]
    xs = [1.02, 6.92, 12.81, 18.70]
    for i, (val, lbl, sub) in enumerate(kpis):
        _kpi_box(slide, xs[i], 3.10, 5.6, 1.95, val, lbl, sub, color_accent=col)

    # Top 5 holdings
    top5 = d.get("top5_a" if is_a else "top5_b", [])
    rows = [["Société (Ticker)", "Poids (%)", "Secteur"]]
    for item in top5[:5]:
        company = str(item[0] if len(item) > 0 else "\u2014")[:34]
        ticker = str(item[1] if len(item) > 1 else "")
        weight = item[2] if len(item) > 2 else None
        sector = str(item[3] if len(item) > 3 else "")[:24]
        wt_str = _fr_num(weight, 1) + "\u00a0%" if weight is not None else "\u2014"
        label = f"{company} ({ticker})" if ticker else company
        rows.append([label, wt_str, sector])
    while len(rows) < 6:
        rows.append(["\u2014", "\u2014", "\u2014"])
    _table(slide, rows, 1.02, 5.50, 13.0, 4.30,
           col_widths=[7, 2.5, 3.5], hdr_fill=col, font_size=9)

    # Box LLM "Profil sectoriel"
    llm_text = d.get("llm", {}).get("profil_a" if is_a else "profil_b", "")
    if not llm_text:
        llm_text = (
            f"{name} regroupe les sociétés cotées sur l'univers {code} et constitue "
            f"un baromètre représentatif de l'économie sous-jacente. Ses caractéristiques "
            f"structurelles (composition sectorielle, biais croissance/value, liquidité) "
            f"déterminent son comportement face aux cycles macroéconomiques. "
            f"L'horizon recommandé pour une allocation cœur de portefeuille est de 3 à 5 ans."
        )
    # JPM-style dynamic title : KPI signature de l'indice
    _car_title = f"{name} ({code}) — perf 1Y {_fr_pct_signed(p1y)}, P/E {_fr_num(pe, 1)}x"
    _llm_box(slide, 14.30, 5.50, 10.10, 4.30,
             _car_title, llm_text, fontsize=9)

    # Bande footer "Lecture du profil"
    profil_read = d.get("llm", {}).get("profil_a_read" if is_a else "profil_b_read", "")
    if not profil_read:
        profil_read = (
            f"L'évaluation du profil de {name} doit être croisée avec les conditions de marché "
            f"actuelles : niveau des taux directeurs, cycle économique, et flux de capitaux. "
            f"Le score FinSight {sc}/100 traduit un signal {sig.lower()} sur un horizon 6-12 mois."
        )
    # JPM-style dynamic title
    _prof_title = f"{name} — Score FinSight {sc}/100 : signal {sig.lower()} sur 6-12 mois"
    _llm_box(slide, 1.02, 10.10, 23.37, 3.20,
             _prof_title, profil_read, fontsize=9)


# ── Slide 7 — Composition sectorielle comparée ──────────────────────────────

def _slide_07_secteurs(prs, d: dict):
    slide = _blank(prs)
    name_a = d.get("name_a", "Indice A")
    name_b = d.get("name_b", "Indice B")
    _header(slide, (data.get('_t_helper') or (lambda k: "Composition Sectorielle"))('compo_sec'),
            f"Poids sectoriels (%) — {name_a} vs {name_b}", section_active=1)
    _footer(slide)
    _index_band(slide, name_a, name_b)

    sector_cmp = d.get("sector_comparison", [])

    if sector_cmp:
        try:
            sects = [s[0][:22] for s in sector_cmp[:11]]
            wa = [float(s[1] or 0) for s in sector_cmp[:11]]
            wb = [float(s[2] or 0) for s in sector_cmp[:11]]

            fig, ax = plt.subplots(figsize=(13, 6.5))
            fig.patch.set_facecolor('white')
            ax.set_facecolor('#FAFBFD')

            y = np.arange(len(sects))
            bh = 0.36

            ax.barh(y + bh / 2, wa, bh, label=name_a[:22],
                    color='#2E5FA3', alpha=0.88)
            ax.barh(y - bh / 2, wb, bh, label=name_b[:22],
                    color='#1A7A4A', alpha=0.88)

            ax.set_yticks(y)
            ax.set_yticklabels(sects, fontsize=11)
            ax.invert_yaxis()
            ax.set_xlabel("Poids (%)", fontsize=11)
            ax.legend(loc='lower right', fontsize=11, framealpha=0.95)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.grid(axis='x', alpha=0.3)

            for i, (va, vb) in enumerate(zip(wa, wb)):
                if va > 0.5:
                    ax.text(va + 0.2, i + bh / 2, f"{va:.1f}%",
                            fontsize=9, va='center', color='#2E5FA3')
                if vb > 0.5:
                    ax.text(vb + 0.2, i - bh / 2, f"{vb:.1f}%",
                            fontsize=9, va='center', color='#1A7A4A')

            fig.tight_layout(pad=0.5)
            img = _chart_buf(fig)
            _pic(slide, img, 1.02, 3.10, 14.5, 7.50)
        except Exception as e:
            log.warning(f"[indice_cmp_pptx] sector chart error: {e}")
    else:
        _txb(slide, "Données de composition sectorielle non disponibles.",
             1.02, 5.0, 14.5, 1.0, 9, False, _GRAYT)

    # Box LLM "Lecture sectorielle" (droite)
    sectoral_read = d.get("llm", {}).get("sectoral_read", "")
    if not sectoral_read:
        sectoral_read = (
            f"La composition sectorielle révèle les biais structurels de chaque indice. "
            f"Les écarts de pondération entre {name_a} et {name_b} traduisent des "
            f"différences fondamentales : exposition aux cycles, sensibilité aux taux, "
            f"orientation croissance vs value. Un investisseur peut combiner les deux "
            f"indices pour obtenir une exposition équilibrée ou en privilégier un selon "
            f"sa conviction sectorielle."
        )
    # JPM-style dynamic title : plus gros écart sectoriel
    _sect_title = f"{name_a} vs {name_b} — composition sectorielle"
    try:
        if sector_cmp:
            _spreads = [(s[0], float(s[1] or 0) - float(s[2] or 0)) for s in sector_cmp[:11]]
            _max = max(_spreads, key=lambda x: abs(x[1]))
            _winner = name_a if _max[1] > 0 else name_b
            _sect_title = f"{_winner} surpondère {_max[0]} (+{abs(_max[1]):.1f} pts) — biais structurel principal"
    except Exception as _e:
        log.debug(f"[cmp_indice_pptx_writer:_slide_07_secteurs] exception skipped: {_e}")
    _llm_box(slide, 16.00, 3.10, 8.40, 7.50,
             _sect_title, sectoral_read, fontsize=9)

    # Bande footer "Implications allocation"
    impl = d.get("llm", {}).get("sectoral_impl", "")
    if not impl:
        impl = (
            f"Les surpondérations sectorielles déterminent la sensibilité de l'indice "
            f"aux cycles macroéconomiques : un indice tech-heavy bénéficie des phases "
            f"de baisse des taux et d'innovation, tandis qu'un indice financials-heavy "
            f"profite des hausses de taux et d'une courbe pentue. À pondérer selon le régime macro."
        )
    _llm_box(slide, 1.02, 10.85, 23.37, 2.45,
             "Implications pour l'allocation", impl, fontsize=9)


# ── Slide 8 — Top 5 / Top 10 holdings comparatif ────────────────────────────

def _slide_08_top_holdings(prs, d: dict):
    slide = _blank(prs)
    name_a = d.get("name_a", "Indice A")
    name_b = d.get("name_b", "Indice B")
    _header(slide, (data.get('_t_helper') or (lambda k: "Top Constituants Comparés"))('top_const'),
            f"Principales valeurs de {name_a} et {name_b}", section_active=1)
    _footer(slide)

    top5_a = d.get("top5_a", [])
    top5_b = d.get("top5_b", [])

    # Bandeau A
    _rect(slide, 1.02, 2.10, 11.44, 0.70, _COL_A)
    _txb(slide, f"Top 5  —  {name_a}", 1.30, 2.22, 11.0, 0.50,
         11, True, _WHITE)
    rows_a = [["Société (Ticker)", "Poids", "Secteur"]]
    for item in top5_a[:5]:
        company = str(item[0] if len(item) > 0 else "\u2014")[:30]
        ticker = str(item[1] if len(item) > 1 else "")
        weight = item[2] if len(item) > 2 else None
        sector = str(item[3] if len(item) > 3 else "")[:18]
        wt_str = _fr_num(weight, 1) + " %" if weight is not None else "\u2014"
        label = f"{company} ({ticker})" if ticker else company
        rows_a.append([label, wt_str, sector])
    while len(rows_a) < 6:
        rows_a.append(["\u2014", "\u2014", "\u2014"])
    _table(slide, rows_a, 1.02, 2.95, 11.44, 4.50,
           col_widths=[6.5, 1.8, 3.14], hdr_fill=_COL_A, font_size=8.5)

    # Bandeau B
    _rect(slide, 12.94, 2.10, 11.44, 0.70, _COL_B)
    _txb(slide, f"Top 5  —  {name_b}", 13.20, 2.22, 11.0, 0.50,
         11, True, _WHITE)
    rows_b = [["Société (Ticker)", "Poids", "Secteur"]]
    for item in top5_b[:5]:
        company = str(item[0] if len(item) > 0 else "\u2014")[:30]
        ticker = str(item[1] if len(item) > 1 else "")
        weight = item[2] if len(item) > 2 else None
        sector = str(item[3] if len(item) > 3 else "")[:18]
        wt_str = _fr_num(weight, 1) + " %" if weight is not None else "\u2014"
        label = f"{company} ({ticker})" if ticker else company
        rows_b.append([label, wt_str, sector])
    while len(rows_b) < 6:
        rows_b.append(["\u2014", "\u2014", "\u2014"])
    _table(slide, rows_b, 12.94, 2.95, 11.44, 4.50,
           col_widths=[6.5, 1.8, 3.14], hdr_fill=_COL_B, font_size=8.5)

    # Box LLM concentration
    concentration_read = d.get("llm", {}).get("concentration_read", "")
    if not concentration_read:
        concentration_read = (
            f"La concentration des principaux constituants détermine le risque idiosyncrasique "
            f"de l'indice. Les top 5 valeurs de {name_a} pèsent typiquement 20-30 % du total, "
            f"contre 10-15 % pour {name_b}. Une concentration élevée amplifie l'impact des "
            f"publications individuelles sur la performance de l'indice (effet 'magnificent 7' "
            f"sur le S&P 500). À l'inverse, une dispersion plus large offre une diversification "
            f"intrinsèque mais réduit la sensibilité aux leaders technologiques."
        )
    _llm_box(slide, 1.02, 7.65, 23.37, 5.65,
             "Concentration des principaux constituants",
             concentration_read, fontsize=9)


# ── Slide 10 — Performance historique 52W (chart base 100 + macro) ──────────

def _slide_10_perf_chart(prs, d: dict):
    slide = _blank(prs)
    name_a = d.get("name_a", "Indice A")
    name_b = d.get("name_b", "Indice B")
    _header(slide, (data.get('_t_helper') or (lambda k: "Performance Historique"))('perf_52w'),
            f"Évolution normalisée base 100 — {name_a} vs {name_b}",
            section_active=2)
    _footer(slide)
    _index_band(slide, name_a, name_b)

    ph = d.get("perf_history") or {}
    has_data = bool(ph.get("dates") and ph.get("indice_a") and ph.get("indice_b"))

    if has_data:
        try:
            fig, ax = plt.subplots(figsize=(13, 7.5))
            fig.patch.set_facecolor('white')
            ax.set_facecolor('#FAFBFD')

            dates_raw = ph["dates"]
            ys_a = ph["indice_a"]
            ys_b = ph["indice_b"]
            n = len(dates_raw)
            step = max(1, n // 100)
            idx = list(range(0, n, step))
            if idx and idx[-1] != n - 1:
                idx.append(n - 1)
            xs = list(range(len(idx)))
            lbls = [dates_raw[i][:7] for i in idx]

            ax.plot(xs, [ys_a[i] for i in idx], color='#2E5FA3', linewidth=2.5,
                    label=name_a[:24], zorder=3)
            ax.plot(xs, [ys_b[i] for i in idx], color='#1A7A4A', linewidth=2.5,
                    label=name_b[:24], zorder=3)
            ax.axhline(100, color='#CCCCCC', linewidth=1.0, linestyle='--', zorder=1)
            ax.fill_between(xs, [ys_a[i] for i in idx], 100, alpha=0.10, color='#2E5FA3')
            ax.fill_between(xs, [ys_b[i] for i in idx], 100, alpha=0.10, color='#1A7A4A')

            tick_step = max(1, len(xs) // 8)
            ax.set_xticks(xs[::tick_step])
            ax.set_xticklabels(lbls[::tick_step], rotation=25, ha='right', fontsize=10)
            ax.set_ylabel("Base 100", fontsize=11)
            ax.tick_params(labelsize=10)
            ax.legend(loc='upper left', fontsize=12, framealpha=0.95)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.grid(axis='y', alpha=0.3)

            # Titre analytique en haut du chart
            ret_a = ys_a[-1] - 100
            ret_b = ys_b[-1] - 100
            leader = name_a if ret_a > ret_b else name_b
            spread = abs(ret_a - ret_b)
            ax.set_title(
                f"52 semaines : {name_a} {ret_a:+.1f}% vs {name_b} {ret_b:+.1f}% "
                f"— {leader} surperforme de {spread:.1f} pts",
                fontsize=13, fontweight='bold', color='#1B3A6B', pad=14
            )

            fig.tight_layout(pad=0.5)
            img = _chart_buf(fig)
            _pic(slide, img, 1.02, 3.10, 16.50, 7.50)
        except Exception as e:
            log.warning(f"[indice_cmp_pptx] perf chart error: {e}")
    else:
        _txb(slide, "Données de performance historique indisponibles.",
             1.02, 6.0, 16.50, 1.0, 9, False, _GRAYT, italic=True)

    # Tableau perfs YTD/1Y/3Y/5Y droite
    rows = [
        ["Horizon", name_a[:14], name_b[:14], "Écart"],
        ["YTD", _fr_pct_signed(d.get("perf_ytd_a")),
         _fr_pct_signed(d.get("perf_ytd_b")),
         _ecart_pct(d.get("perf_ytd_a"), d.get("perf_ytd_b"))],
        ["1 an", _fr_pct_signed(d.get("perf_1y_a")),
         _fr_pct_signed(d.get("perf_1y_b")),
         _ecart_pct(d.get("perf_1y_a"), d.get("perf_1y_b"))],
        ["3 ans", _fr_pct_signed(d.get("perf_3y_a")),
         _fr_pct_signed(d.get("perf_3y_b")),
         _ecart_pct(d.get("perf_3y_a"), d.get("perf_3y_b"))],
        ["5 ans", _fr_pct_signed(d.get("perf_5y_a")),
         _fr_pct_signed(d.get("perf_5y_b")),
         _ecart_pct(d.get("perf_5y_a"), d.get("perf_5y_b"))],
    ]
    _table(slide, rows, 17.85, 3.10, 6.55, 4.50,
           col_widths=[2.5, 2, 2, 2], font_size=8.5)

    # Box LLM macro
    perf_macro = d.get("llm", {}).get("perf_macro_read", "")
    if not perf_macro:
        perf_macro = (
            f"Sur les 52 dernières semaines, le contexte macro a été marqué par la stabilisation "
            f"des taux directeurs (Fed 4,25-4,50 %, BCE 3,00 % après le pivot mi-2024), un "
            f"atterrissage en douceur de l'économie américaine et une rotation sectorielle au "
            f"profit des valeurs de qualité. Le dollar est resté fort vs euro (EUR/USD ~1,06), "
            f"pénalisant les revenus européens convertis. La trajectoire des deux indices reflète "
            f"ces dynamiques croisées et leur exposition différenciée."
        )
    _llm_box(slide, 17.85, 7.85, 6.55, 2.75,
             "Contexte macro 12 mois", perf_macro, fontsize=8.5)

    # Bande footer LLM lecture analytique
    perf_lecture = d.get("llm", {}).get("perf_chart_read", "")
    if not perf_lecture:
        perf_lecture = (
            f"La courbe normalisée base 100 permet d'isoler la performance pure de chaque indice, "
            f"indépendamment de leur niveau absolu. La zone ombragée matérialise l'écart relatif "
            f"sur la période. Une divergence persistante signale une rotation structurelle entre "
            f"les deux univers, à confronter aux flux institutionnels et aux révisions de consensus."
        )
    _p1y_a = d.get("perf_1y_a") or 0
    _p1y_b = d.get("perf_1y_b") or 0
    _perf_winner = name_a if _p1y_a >= _p1y_b else name_b
    _traj_title = f"{_perf_winner} surperforme sur 52S ({max(_p1y_a, _p1y_b):+.1f}%) — divergence structurelle"
    _llm_box(slide, 1.02, 10.85, 23.37, 2.45,
             _traj_title, perf_lecture, fontsize=9)


# ── Slide 11 — Décomposition perfs ──────────────────────────────────────────

def _slide_11_perf_decomposition(prs, d: dict):
    slide = _blank(prs)
    name_a = d.get("name_a", "Indice A")
    name_b = d.get("name_b", "Indice B")
    _header(slide, (data.get('_t_helper') or (lambda k: "Décomposition Performance"))('decompo_perf'),
            f"YTD, 1 an, 3 ans, 5 ans  —  écarts relatifs et lecture comparative",
            section_active=2)
    _footer(slide)
    _index_band(slide, name_a, name_b)

    # Graphique barres comparatif horizons multiples
    horizons = ["YTD", "1 an", "3 ans", "5 ans"]
    keys = ["perf_ytd", "perf_1y", "perf_3y", "perf_5y"]
    vals_a = []
    vals_b = []
    for k in keys:
        va = _safe_float(d.get(k + "_a"))
        vb = _safe_float(d.get(k + "_b"))
        if va is not None and abs(va) <= 2:
            va *= 100
        if vb is not None and abs(vb) <= 2:
            vb *= 100
        vals_a.append(va or 0)
        vals_b.append(vb or 0)

    try:
        x = np.arange(len(horizons))
        bw = 0.36
        fig, ax = plt.subplots(figsize=(11.5, 5.5))
        fig.patch.set_facecolor('white')
        ax.set_facecolor('#FAFBFD')
        b1 = ax.bar(x - bw / 2, vals_a, bw, label=name_a[:24],
                    color='#2E5FA3', alpha=0.88)
        b2 = ax.bar(x + bw / 2, vals_b, bw, label=name_b[:24],
                    color='#1A7A4A', alpha=0.88)
        ax.set_xticks(x)
        ax.set_xticklabels(horizons, fontsize=12)
        ax.set_ylabel("Performance (%)", fontsize=11)
        ax.legend(fontsize=11, framealpha=0.95)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.grid(axis='y', alpha=0.3)
        ax.axhline(0, color='#666666', linewidth=0.8)
        for bar in list(b1) + list(b2):
            h = bar.get_height()
            ax.annotate(f"{h:+.1f}%",
                        xy=(bar.get_x() + bar.get_width() / 2, h),
                        xytext=(0, 4 if h >= 0 else -12),
                        textcoords="offset points",
                        ha='center', va='bottom' if h >= 0 else 'top',
                        fontsize=9, fontweight='bold')
        fig.tight_layout(pad=0.5)
        img = _chart_buf(fig)
        _pic(slide, img, 1.02, 3.10, 14.5, 6.50)
    except Exception as e:
        log.warning(f"[indice_cmp_pptx] perf Décomposition error: {e}")

    # Tableau écarts (droite)
    rows = [
        ["Horizon", name_a[:14], name_b[:14], "Écart (pp)"],
    ]
    for h, k in zip(horizons, keys):
        va = d.get(k + "_a")
        vb = d.get(k + "_b")
        rows.append([h, _fr_pct_signed(va), _fr_pct_signed(vb), _ecart_pct(va, vb)])
    _table(slide, rows, 16.0, 3.10, 8.40, 4.50,
           col_widths=[2.5, 2, 2, 2], font_size=9)

    # Box LLM lecture
    perf_decomp = d.get("llm", {}).get("perf_decomp_read", "")
    if not perf_decomp:
        perf_decomp = (
            f"La comparaison sur 4 horizons révèle si la surperformance d'un indice est récente "
            f"(momentum court terme) ou structurelle (avantage long terme). Un indice qui surperforme "
            f"YTD mais sous-performe sur 5 ans peut traduire une rotation tactique passagère, "
            f"tandis qu'une surperformance constante signale un avantage compétitif durable de "
            f"l'écosystème sous-jacent (innovation, démographie, productivité). À croiser avec "
            f"les valorisations relatives pour identifier les opportunités de mean reversion."
        )
    _llm_box(slide, 16.0, 7.85, 8.40, 5.45,
             "Lecture décomposition perfs", perf_decomp, fontsize=8.5)

    # Bande analytique full width
    perf_horizon = d.get("llm", {}).get("perf_horizon_read", "")
    if not perf_horizon:
        perf_horizon = (
            f"Pour un investisseur long terme, c'est la performance 5 ans qui importe : elle "
            f"capture les cycles complets et lisse les effets de marché temporaires. Pour un "
            f"investisseur tactique, l'horizon YTD/1Y est plus pertinent. Le bon mix dépend du "
            f"mandat de gestion et du benchmark de référence."
        )
    _llm_box(slide, 1.02, 9.85, 14.5, 3.45,
             "Implications pour l'horizon d'investissement",
             perf_horizon, fontsize=9)


# ── Slide 12 — Risque comparatif ────────────────────────────────────────────

def _slide_12_risque(prs, d: dict):
    slide = _blank(prs)
    name_a = d.get("name_a", "Indice A")
    name_b = d.get("name_b", "Indice B")
    _header(slide, (data.get('_t_helper') or (lambda k: "Risque Comparatif"))('risque_cmp'),
            f"Volatilité, Sharpe ratio et Max Drawdown sur 12 mois",
            section_active=2)
    _footer(slide)
    _index_band(slide, name_a, name_b)

    vol_a = d.get("vol_1y_a")
    vol_b = d.get("vol_1y_b")
    sha_a = d.get("sharpe_1y_a")
    sha_b = d.get("sharpe_1y_b")
    dd_a = d.get("max_dd_a")
    dd_b = d.get("max_dd_b")

    # 3 KPIs comparatifs
    kpis = [
        (_fr_num(vol_a) + " / " + _fr_num(vol_b) + "\u00a0%",
         "Volatilité annualisée (A / B)",
         "Écart-type des rendements quotidiens annualisé"),
        (_fr_num(sha_a, 2) + " / " + _fr_num(sha_b, 2),
         "Sharpe ratio 1 an (A / B)",
         "Rendement ajusté du risque — taux sans risque 4 %"),
        (_fr_pct_signed(dd_a) + " / " + _fr_pct_signed(dd_b),
         "Max Drawdown (A / B)",
         "Perte maximale du plus haut au plus bas sur la période"),
    ]
    xs = [1.02, 9.22, 17.42]
    for i, (val, lbl, sub) in enumerate(kpis):
        _rect(slide, xs[i], 3.10, 7.8, 2.5, _NAVY_PALE)
        _rect(slide, xs[i], 3.10, 0.15, 2.5, _NAVY_MID)
        _txb(slide, val, xs[i] + 0.30, 3.20, 7.4, 1.10, 16, True, _NAVY)
        _txb(slide, lbl, xs[i] + 0.30, 4.25, 7.4, 0.45, 8, True, _GRAYT)
        _txb(slide, sub, xs[i] + 0.30, 4.70, 7.4, 0.80, 7.5, False, _GRAYT, wrap=True)

    # Graphique barres comparatif
    try:
        cats = ["Volatilité", "Sharpe", "Max Drawdown"]
        va_arr = [
            float(vol_a or 0),
            float(sha_a or 0),
            abs(float(dd_a or 0)) * 100 if dd_a and abs(float(dd_a)) < 1 else abs(float(dd_a or 0))
        ]
        vb_arr = [
            float(vol_b or 0),
            float(sha_b or 0),
            abs(float(dd_b or 0)) * 100 if dd_b and abs(float(dd_b)) < 1 else abs(float(dd_b or 0))
        ]
        x = np.arange(len(cats))
        bw = 0.36
        fig, ax = plt.subplots(figsize=(13, 4.0))
        fig.patch.set_facecolor('white')
        ax.set_facecolor('#FAFBFD')
        b1 = ax.bar(x - bw / 2, va_arr, bw, label=name_a[:24], color='#2E5FA3', alpha=0.88)
        b2 = ax.bar(x + bw / 2, vb_arr, bw, label=name_b[:24], color='#1A7A4A', alpha=0.88)
        ax.set_xticks(x)
        ax.set_xticklabels(cats, fontsize=12)
        ax.legend(fontsize=11, framealpha=0.95)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.grid(axis='y', alpha=0.3)
        for bar in list(b1) + list(b2):
            h = bar.get_height()
            ax.annotate(f"{h:.1f}",
                        xy=(bar.get_x() + bar.get_width() / 2, h),
                        xytext=(0, 4),
                        textcoords="offset points",
                        ha='center', va='bottom', fontsize=10)
        fig.tight_layout(pad=0.5)
        img = _chart_buf(fig)
        _pic(slide, img, 1.02, 5.85, 23.37, 4.50)
    except Exception as e:
        log.warning(f"[indice_cmp_pptx] risque chart error: {e}")

    # Box LLM risque
    risque_read = d.get("llm", {}).get("risque_read", "")
    if not risque_read:
        risque_read = (
            f"Le profil de risque détermine la taille de position en portefeuille : un indice "
            f"plus volatil impose une pondération réduite pour maintenir un budget risque "
            f"équivalent. Le Sharpe ratio mesure l'efficience du couple rendement/risque. "
            f"Un Sharpe > 1 est considéré comme excellent, 0,5-1 comme correct, < 0,5 comme "
            f"décevant. Le Max Drawdown capture le pire scénario réalisé sur la période et "
            f"informe sur la résilience en phase de stress."
        )
    # JPM title : choisi sur Sharpe (proxy correct du couple rendement/risque)
    _sh_a = d.get("sharpe_1y_a")
    _sh_b = d.get("sharpe_1y_b")
    if _sh_a is not None and _sh_b is not None and _sh_a != _sh_b:
        _better = name_a if _sh_a > _sh_b else name_b
        _better_sh = max(_sh_a, _sh_b)
        _risk_title = f"{_better} — Sharpe {_better_sh:.2f} : meilleur couple rendement/risque"
    else:
        _risk_title = f"{name_a} vs {name_b} — profil risque comparé"
    _llm_box(slide, 1.02, 10.50, 23.37, 2.80,
             _risk_title, risque_read, fontsize=9)


# ── Slide 13 — Risk/Return scatter ──────────────────────────────────────────

def _slide_13_risk_return(prs, d: dict):
    slide = _blank(prs)
    name_a = d.get("name_a", "Indice A")
    name_b = d.get("name_b", "Indice B")
    _header(slide, (data.get('_t_helper') or (lambda k: "Profil Rendement / Risque"))('rend_risque'),
            f"Positionnement {name_a} vs {name_b} sur le plan vol-rendement",
            section_active=2)
    _footer(slide)
    _index_band(slide, name_a, name_b)

    vol_a = _safe_float(d.get("vol_1y_a")) or 0
    vol_b = _safe_float(d.get("vol_1y_b")) or 0
    p1y_a = _safe_float(d.get("perf_1y_a")) or 0
    p1y_b = _safe_float(d.get("perf_1y_b")) or 0
    if abs(p1y_a) <= 2:
        p1y_a *= 100
    if abs(p1y_b) <= 2:
        p1y_b *= 100

    try:
        fig, ax = plt.subplots(figsize=(11, 6.5))
        fig.patch.set_facecolor('white')
        ax.set_facecolor('#FAFBFD')

        # Scatter
        ax.scatter([vol_a], [p1y_a], s=550, color='#2E5FA3', alpha=0.85,
                   edgecolors='white', linewidths=2, label=name_a[:24], zorder=5)
        ax.scatter([vol_b], [p1y_b], s=550, color='#1A7A4A', alpha=0.85,
                   edgecolors='white', linewidths=2, label=name_b[:24], zorder=5)

        # Annotations
        ax.annotate(name_a[:18], (vol_a, p1y_a),
                    xytext=(8, 8), textcoords="offset points",
                    fontsize=11, fontweight='bold', color='#2E5FA3')
        ax.annotate(name_b[:18], (vol_b, p1y_b),
                    xytext=(8, 8), textcoords="offset points",
                    fontsize=11, fontweight='bold', color='#1A7A4A')

        # Quadrants
        x_max = max(vol_a, vol_b) * 1.5 + 5
        y_max = max(p1y_a, p1y_b, 5) * 1.5
        y_min = min(p1y_a, p1y_b, -5) * 1.5
        ax.axhline(0, color='#999', linewidth=0.7, linestyle='--', alpha=0.7)
        ax.axvline((vol_a + vol_b) / 2, color='#999', linewidth=0.7,
                   linestyle='--', alpha=0.7)
        ax.set_xlim(0, x_max)
        ax.set_ylim(y_min - 5, y_max + 5)

        ax.set_xlabel("Volatilité annualisée (%)", fontsize=12)
        ax.set_ylabel("Performance 1 an (%)", fontsize=12)
        ax.set_title("Profil rendement / risque sur 12 mois",
                     fontsize=14, fontweight='bold', color='#1B3A6B', pad=12)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.grid(alpha=0.3)
        ax.legend(loc='lower right', fontsize=11, framealpha=0.95)

        fig.tight_layout(pad=0.5)
        img = _chart_buf(fig)
        _pic(slide, img, 1.02, 3.10, 14.5, 7.50)
    except Exception as e:
        log.warning(f"[indice_cmp_pptx] risk-return error: {e}")

    # Box LLM risk-return
    rr_read = d.get("llm", {}).get("risk_return_read", "")
    if not rr_read:
        rr_read = (
            f"Le plan rendement-risque permet de visualiser instantanément l'efficience relative "
            f"des deux indices. L'indice idéal se situe en haut à gauche (rendement élevé, "
            f"volatilité basse), zone des allocations défensives premium. Inversement, le "
            f"quadrant bas-droite (rendement faible, volatilité élevée) signale un mauvais "
            f"couple risk/reward à éviter. Pour un mandat balanced, viser le quadrant haut-droite "
            f"(rendement élevé même au prix d'une volatilité supérieure) reste rationnel si "
            f"l'horizon est suffisamment long."
        )
    _llm_box(slide, 16.0, 3.10, 8.40, 7.50,
             "Lecture rendement / risque", rr_read, fontsize=9)

    # Bande implications
    rr_impl = d.get("llm", {}).get("risk_return_impl", "")
    if not rr_impl:
        rr_impl = (
            f"La position relative de chaque indice oriente l'allocation : si {name_a} domine "
            f"sur l'efficience, le surpondérer en cœur de portefeuille. Si {name_b} offre une "
            f"meilleure protection en phase de stress (Max DD plus contenu), l'utiliser comme "
            f"hedge. Une combinaison équilibrée 60/40 ou 70/30 capture le meilleur des deux mondes."
        )
    _llm_box(slide, 1.02, 10.85, 14.5, 2.45,
             "Implications pour la construction de portefeuille",
             rr_impl, fontsize=9)


# ── Slide 15 — Valorisation agrégée ─────────────────────────────────────────

def _slide_15_valorisation(prs, d: dict):
    slide = _blank(prs)
    name_a = d.get("name_a", "Indice A")
    name_b = d.get("name_b", "Indice B")
    _header(slide, (data.get('_t_helper') or (lambda k: "Valorisation Comparée"))('valo_cmp'),
            f"P/E Forward, P/B, Rendement du dividende",
            section_active=3)
    _footer(slide)
    _index_band(slide, name_a, name_b)

    pe_a = d.get("pe_fwd_a")
    pe_b = d.get("pe_fwd_b")
    pb_a = d.get("pb_a")
    pb_b = d.get("pb_b")
    dy_a = d.get("div_yield_a")
    dy_b = d.get("div_yield_b")

    rows = [
        ["Indicateur", name_a[:18], name_b[:18], "Écart", "Avantage"],
        ["P/E Forward",
         _fr_x(pe_a, 1), _fr_x(pe_b, 1),
         _ecart_mult(pe_a, pe_b),
         _favorise_lower(pe_a, pe_b, name_a, name_b)],
        ["P/B (Book Value)",
         _fr_x(pb_a, 1), _fr_x(pb_b, 1),
         _ecart_mult(pb_a, pb_b),
         _favorise_lower(pb_a, pb_b, name_a, name_b)],
        ["Rendement dividende",
         _fr_pct(dy_a) if dy_a else "\u2014",
         _fr_pct(dy_b) if dy_b else "\u2014",
         _ecart_pct(dy_a, dy_b),
         _favorise_higher(dy_a, dy_b, name_a, name_b)],
    ]
    _table(slide, rows, 1.02, 3.10, 14.5, 4.50,
           col_widths=[4.5, 2.5, 2.5, 2.5, 2.5], font_size=9)

    # Chart barres P/E + P/B
    try:
        cats = []
        va = []
        vb = []
        if pe_a is not None or pe_b is not None:
            cats.append("P/E Forward")
            va.append(float(pe_a or 0))
            vb.append(float(pe_b or 0))
        if pb_a is not None or pb_b is not None:
            cats.append("P/B")
            va.append(float(pb_a or 0))
            vb.append(float(pb_b or 0))

        if cats:
            x = np.arange(len(cats))
            bw = 0.36
            fig, ax = plt.subplots(figsize=(8, 4.5))
            fig.patch.set_facecolor('white')
            ax.set_facecolor('#FAFBFD')
            b1 = ax.bar(x - bw / 2, va, bw, label=name_a[:24], color='#2E5FA3', alpha=0.88)
            b2 = ax.bar(x + bw / 2, vb, bw, label=name_b[:24], color='#1A7A4A', alpha=0.88)
            ax.set_xticks(x)
            ax.set_xticklabels(cats, fontsize=12)
            ax.set_ylabel("Multiple (x)", fontsize=11)
            ax.legend(fontsize=11, framealpha=0.95)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.grid(axis='y', alpha=0.3)
            for bar in list(b1) + list(b2):
                h = bar.get_height()
                if h > 0:
                    ax.annotate(f"{h:.1f}x",
                                xy=(bar.get_x() + bar.get_width() / 2, h),
                                xytext=(0, 4),
                                textcoords="offset points",
                                ha='center', va='bottom', fontsize=10)
            fig.tight_layout(pad=0.5)
            img = _chart_buf(fig)
            _pic(slide, img, 1.02, 7.85, 14.5, 5.45)
    except Exception as e:
        log.warning(f"[indice_cmp_pptx] valuation chart error: {e}")

    # Box LLM
    val_read = d.get("llm", {}).get("valorisation_read", "")
    if not val_read:
        val_read = (
            f"Les multiples agrégés des indices reflètent la prime de valorisation que le marché "
            f"attribue à chaque univers. Un P/E élevé peut traduire (1) des anticipations de "
            f"croissance supérieures, (2) une qualité bilancielle premium, ou (3) une exubérance "
            f"de marché. Le P/B complète l'analyse en mesurant la prime sur les actifs nets. "
            f"Le rendement du dividende informe sur la maturité bilancielle et la politique de "
            f"distribution agrégée. À comparer aux moyennes historiques 5-10 ans pour identifier "
            f"les anomalies de mean reversion."
        )
    _pe_a = d.get("pe_fwd_a")
    _pe_b = d.get("pe_fwd_b")
    if _pe_a is not None and _pe_b is not None and _pe_a != _pe_b:
        _cheap = name_a if _pe_a < _pe_b else name_b
        _val_title = f"{_cheap} décoté ({min(_pe_a, _pe_b):.1f}x P/E fwd)"
    else:
        _val_title = f"{name_a} vs {name_b} — valorisation comparée"
    _llm_box(slide, 16.0, 3.10, 8.40, 10.20,
             _val_title, val_read, fontsize=9)


# ── Slide 16 — ERP focus ────────────────────────────────────────────────────

def _slide_16_erp_focus(prs, d: dict):
    slide = _blank(prs)
    name_a = d.get("name_a", "Indice A")
    name_b = d.get("name_b", "Indice B")
    _header(slide, (data.get('_t_helper') or (lambda k: "Equity Risk Premium"))('erp'),
            f"Comparaison du couple rendement-prime",
            section_active=3)
    _footer(slide)
    _index_band(slide, name_a, name_b)

    erp_a = d.get("erp_a", "\u2014")
    erp_b = d.get("erp_b", "\u2014")

    def _parse_erp(s):
        try:
            return float(str(s).replace(",", ".").replace("%", "").replace("+", "").strip())
        except Exception:
            return None

    erp_a_val = _parse_erp(erp_a)
    erp_b_val = _parse_erp(erp_b)

    def _erp_signal(v):
        if v is None:
            return ("N/D", _GRAYT)
        if v >= 6.0:
            return ("Très attractif", _BUY)
        if v >= 4.0:
            return ("Attractif", _BUY)
        if v >= 2.0:
            return ("Neutre", _HOLD)
        if v >= 0.0:
            return ("Tendu", _HOLD)
        return ("Négatif", _SELL)

    sig_a, col_a = _erp_signal(erp_a_val)
    sig_b, col_b = _erp_signal(erp_b_val)

    # 2 grandes boxes ERP
    # A
    _rect(slide, 1.02, 3.10, 11.44, 4.20, _COL_AL)
    _rect(slide, 1.02, 3.10, 0.15, 4.20, _COL_A)
    _txb(slide, name_a, 1.30, 3.25, 11.0, 0.50, 12, True, _NAVY)
    _txb(slide, str(erp_a), 1.30, 3.95, 11.0, 1.30, 32, True, _COL_A)
    _txb(slide, "ERP estimé sur 12 mois", 1.30, 5.40, 11.0, 0.45, 9, False, _GRAYT)
    _rect(slide, 1.30, 6.00, 6.0, 0.90, col_a)
    _txb(slide, sig_a, 1.30, 6.10, 6.0, 0.70, 11, True, _WHITE,
         align=PP_ALIGN.CENTER)

    # B
    _rect(slide, 12.94, 3.10, 11.44, 4.20, _COL_BL)
    _rect(slide, 12.94, 3.10, 0.15, 4.20, _COL_B)
    _txb(slide, name_b, 13.22, 3.25, 11.0, 0.50, 12, True, _BUY)
    _txb(slide, str(erp_b), 13.22, 3.95, 11.0, 1.30, 32, True, _COL_B)
    _txb(slide, "ERP estimé sur 12 mois", 13.22, 5.40, 11.0, 0.45, 9, False, _GRAYT)
    _rect(slide, 13.22, 6.00, 6.0, 0.90, col_b)
    _txb(slide, sig_b, 13.22, 6.10, 6.0, 0.70, 11, True, _WHITE,
         align=PP_ALIGN.CENTER)

    # Box LLM ERP
    erp_read = d.get("llm", {}).get("erp_read", "")
    if not erp_read:
        erp_read = (
            f"L'Equity Risk Premium (ERP) mesure la prime exigée par les investisseurs pour "
            f"détenir des actions plutôt que des obligations souveraines (proxy taux 10 ans local). "
            f"Un ERP élevé (>4 %) traduit une valorisation actions historiquement attractive : le "
            f"surplus de rendement compense largement le risque additionnel. Un ERP faible (<2 %) "
            f"signale au contraire une prime insuffisanté, exposant à un re-rating baissier en "
            f"cas de hausse des taux ou de déception sur les bénéfices. La comparaison ERP entre "
            f"deux indices permet d'identifier le mieux rémunéré relativement au risque."
        )
    _erp_title = f"ERP {name_a} {erp_a} vs {name_b} {erp_b} — prime de risque actions comparée"
    _llm_box(slide, 1.02, 7.65, 23.37, 5.65,
             _erp_title, erp_read, fontsize=9)


# ── Slide 17 — Score FinSight comparatif ────────────────────────────────────

def _slide_17_score(prs, d: dict):
    slide = _blank(prs)
    name_a = d.get("name_a", "Indice A")
    name_b = d.get("name_b", "Indice B")
    _header(slide, (data.get('_t_helper') or (lambda k: "Signal FinSight Comparatif"))('signal_fs'),
            "Score 0-100 et recommandation d'allocation", section_active=3)
    _footer(slide)
    _index_band(slide, name_a, name_b)

    sc_a = int(d.get("score_a", 50))
    sc_b = int(d.get("score_b", 50))
    sig_a = d.get("signal_a", "Neutre")
    sig_b = d.get("signal_b", "Neutre")

    # 2 gauges textuelles
    for is_a in (True, False):
        x0 = 1.02 if is_a else 13.37
        sc = sc_a if is_a else sc_b
        sig = sig_a if is_a else sig_b
        name = name_a if is_a else name_b
        col = _sig_c(sig)
        col_l = _sig_l(sig)

        _rect(slide, x0, 3.10, 11.0, 5.30, col_l)
        _rect(slide, x0, 3.10, 0.18, 5.30, col)
        _txb(slide, str(sc), x0 + 0.4, 3.40, 10.2, 2.40,
             58, True, col, align=PP_ALIGN.CENTER)
        _txb(slide, "/ 100", x0 + 0.4, 5.40, 10.2, 0.60,
             16, False, _GRAYT, align=PP_ALIGN.CENTER)
        _txb(slide, sig, x0 + 0.4, 6.10, 10.2, 0.85,
             14, True, col, align=PP_ALIGN.CENTER)
        _txb(slide, f"Signal FinSight  —  {name}",
             x0 + 0.4, 7.10, 10.2, 0.50, 8, False, _GRAYT,
             align=PP_ALIGN.CENTER)

    # Box LLM score
    score_read = d.get("llm", {}).get("score_read", "")
    if not score_read:
        leader = name_a if sc_a > sc_b else (name_b if sc_b > sc_a else "Égalité")
        score_read = (
            f"Le score FinSight est un indicateur propriétaire 0-100 agrégeant 4 dimensions : "
            f"valorisation, croissance, qualité bilancielle et momentum. Les seuils d'allocation "
            f"sont : ≥ 65 = Surpondérer (signal positif fort), 45-64 = Neutre (équilibré), "
            f"< 45 = Sous-pondérer (signal négatif). Le différentiel de score entre les deux "
            f"indices ({sc_a} vs {sc_b}) traduit l'avantage relatif structurel. {leader} ressort "
            f"avec le meilleur signal global, méritant une surpondération tactique sur un horizon "
            f"6-12 mois sous réserve de stabilité macroéconomique."
        )
    _llm_box(slide, 1.02, 8.70, 23.37, 4.60,
             "Lecture du signal FinSight", score_read, fontsize=9)


# ── Slide 19 — Thèses bull/bear comparées ───────────────────────────────────

def _slide_19_theses(prs, d: dict):
    slide = _blank(prs)
    name_a = d.get("name_a", "Indice A")
    name_b = d.get("name_b", "Indice B")
    _header(slide, (data.get('_t_helper') or (lambda k: "Thèses Bull/Bear"))('theses_bb'),
            f"Arguments structurels et risques majeurs : {name_a} vs {name_b}",
            section_active=4)
    _footer(slide)

    # Bandeaux titres
    _rect(slide, 0.9, 2.10, 11.44, 0.65, _COL_A)
    _txb(slide, f"THÈSE LONG  —  {name_a.upper()}", 1.10, 2.20, 11.0, 0.50,
         10, True, _WHITE)

    _rect(slide, 13.04, 2.10, 11.44, 0.65, _GOLD)
    _txb(slide, f"THÈSE LONG  —  {name_b.upper()}", 13.24, 2.20, 11.0, 0.50,
         10, True, _WHITE)

    # Items A (bull avec carrés bleus)
    bull_a_items = d.get("llm", {}).get("bull_a_items") or [
        ("Avantage structurel", "Composition sectorielle alignée avec les moteurs de croissance long terme."),
        ("Profil rendement-risque", "Couple performance/volatilité historiquement attractif sur 5 ans."),
        ("Liquidité institutionnelle", "Univers profond et liquide, accessible via ETF à coût réduit."),
    ]
    y_a = 2.95
    for j, (title, body) in enumerate(bull_a_items[:3]):
        _bullet_square(slide, 0.95, y_a + j * 2.40 + 0.10, _COL_A, size=0.22)
        _txb(slide, title[:48], 1.30, y_a + j * 2.40, 11.0, 0.45,
             10, True, _NAVY)
        _txb(slide, body[:200], 1.30, y_a + j * 2.40 + 0.50, 11.0, 1.75,
             8.5, False, _GRAYT, italic=True, wrap=True)

    # Items B (bull avec carrés or)
    bull_b_items = d.get("llm", {}).get("bull_b_items") or [
        ("Diversification", "Indice complémentaire offrant une exposition différenciée."),
        ("Valorisation relative", "Multiples plus attractifs vs moyennes historiques."),
        ("Catalyseurs sectoriels", "Bénéficiaire des rotations à venir sur cycle macro."),
    ]
    y_b = 2.95
    for j, (title, body) in enumerate(bull_b_items[:3]):
        _bullet_square(slide, 13.09, y_b + j * 2.40 + 0.10, _GOLD, size=0.22)
        _txb(slide, title[:48], 13.44, y_b + j * 2.40, 11.0, 0.45,
             10, True, _BUY)
        _txb(slide, body[:200], 13.44, y_b + j * 2.40 + 0.50, 11.0, 1.75,
             8.5, False, _GRAYT, italic=True, wrap=True)

    # Risques principaux (rouge, full width bottom)
    _rect(slide, 0.9, 10.40, 23.6, 0.55, _RED)
    _txb(slide, "RISQUES PRINCIPAUX", 1.10, 10.46, 23.2, 0.45,
         10, True, _WHITE, align=PP_ALIGN.CENTER)

    bear_text = d.get("llm", {}).get("bear_combined", "")
    if not bear_text:
        bear_text = (
            f"{name_a} : exposition aux cycles macro et risques de re-rating en cas de hausse "
            f"des taux. Concentration sur les leaders amplifie l'impact des publications individuelles. "
            f"{name_b} : sensibilité aux flux de capitaux et aux décisions des banques centrales. "
            f"Les deux indices sont corrélés sur les chocs systémiques (récession, crise de liquidité)."
        )
    _llm_box(slide, 0.9, 11.00, 23.6, 2.30,
             f"{name_a} vs {name_b} — risques structurels et points d'invalidation", bear_text, fontsize=9)


# ── Slide 20 — Verdict final / Recommandation ───────────────────────────────

def _slide_20_verdict(prs, d: dict):
    slide = _blank(prs)
    name_a = d.get("name_a", "Indice A")
    name_b = d.get("name_b", "Indice B")
    sc_a = int(d.get("score_a", 50))
    sc_b = int(d.get("score_b", 50))
    sig_a = d.get("signal_a", "Neutre")
    sig_b = d.get("signal_b", "Neutre")
    winner = name_a if sc_a >= sc_b else name_b
    winner_sig = sig_a if sc_a >= sc_b else sig_b
    winner_col = _COL_A if sc_a >= sc_b else _GOLD

    _header(slide, (data.get('_t_helper') or (lambda k: "Verdict Comparatif"))('verdict'),
            f"Indice à privilégier, arguments décisionnels et conditions d'invalidation",
            section_active=4)
    _footer(slide)

    # Bandeau verdict principal
    _rect(slide, 0.9, 2.10, 23.6, 1.40, winner_col)
    _txb(slide, f"INDICE PRIVILÉGIÉ : {winner.upper()}",
         1.0, 2.20, 23.4, 0.75, 17, True, _WHITE, align=PP_ALIGN.CENTER)
    _txb(slide,
         f"Signal : {winner_sig}  |  Score {max(sc_a, sc_b)}/100 vs "
         f"{min(sc_a, sc_b)}/100  |  Écart {abs(sc_a - sc_b)} pts",
         1.0, 2.97, 23.4, 0.4, 9, False, _WHITE, align=PP_ALIGN.CENTER)

    # Thèse d'allocation (LLM)
    alloc_thesis = d.get("llm", {}).get("verdict_read", "")
    if not alloc_thesis:
        alloc_thesis = (
            f"{winner} présente le profil risque/rendement le plus attractif sur l'horizon 6-12 mois. "
            f"Le score FinSight de {max(sc_a, sc_b)}/100 traduit une supériorité structurelle "
            f"sur les dimensions performance, risque et valorisation. L'écart de {abs(sc_a - sc_b)} pts "
            f"vs l'autre indice ({min(sc_a, sc_b)}/100) légitime une surpondération tactique, "
            f"sous réserve de stabilité macroéconomique et d'absence de choc réglementaire majeur."
        )
    _gap = abs(sc_a - sc_b)
    _verdict_title = f"{winner} privilégié ({max(sc_a, sc_b)}/100) — écart {_gap} pts : conviction {winner_sig.lower()}"
    _llm_box(slide, 0.9, 3.70, 23.6, 2.30,
             _verdict_title, alloc_thesis[:520], fontsize=9)

    # 3 colonnes : Catalyseurs / Risques / Invalidation
    col_w3 = 7.4
    xs = [0.9, 8.6, 16.3]

    # Catalyseurs
    _rect(slide, xs[0], 6.20, col_w3, 0.50, _BUY)
    _txb(slide, "CATALYSEURS CLÉS", xs[0], 6.25, col_w3, 0.40,
         9, True, _WHITE, align=PP_ALIGN.CENTER)
    cats = d.get("llm", {}).get("verdict_catalyseurs") or [
        ("Stabilisation des taux",
         "Pivot dovish des banques centrales soutient les actions risquées."),
        ("Révisions positives",
         "Consensus EPS attendu en hausse sur les 2 prochains trimestres."),
    ]
    y_cat = 6.85
    for ct, cb in cats[:2]:
        _bullet_square(slide, xs[0], y_cat, _BUY, size=0.20)
        _txb(slide, str(ct)[:38], xs[0] + 0.32, y_cat - 0.02, col_w3 - 0.4, 0.45,
             8.5, True, _NAVY)
        _txb(slide, str(cb)[:120], xs[0] + 0.32, y_cat + 0.45, col_w3 - 0.4, 1.40,
             7.5, False, _GRAYT, wrap=True)
        y_cat += 1.95

    # Risques
    _rect(slide, xs[1], 6.20, col_w3, 0.50, _RED)
    _txb(slide, "RISQUES PRINCIPAUX", xs[1], 6.25, col_w3, 0.40,
         9, True, _WHITE, align=PP_ALIGN.CENTER)
    risks = d.get("llm", {}).get("verdict_risques") or [
        ("Choc géopolitique",
         "Tension commerciale ou conflit majeur déclenchant un risk-off généralisé."),
        ("Inflation persistante",
         "Repli du pivot dovish, hausse des taux longs et compression des multiples."),
    ]
    y_risk = 6.85
    for rt, rb in risks[:2]:
        _bullet_square(slide, xs[1], y_risk, _RED, size=0.20)
        _txb(slide, str(rt)[:38], xs[1] + 0.32, y_risk - 0.02, col_w3 - 0.4, 0.45,
             8.5, True, _NAVY)
        _txb(slide, str(rb)[:120], xs[1] + 0.32, y_risk + 0.45, col_w3 - 0.4, 1.40,
             7.5, False, _GRAYT, wrap=True)
        y_risk += 1.95

    # Conditions d'invalidation
    _rect(slide, xs[2], 6.20, col_w3, 0.50, _HOLD)
    _txb(slide, "CONDITIONS D'INVALIDATION", xs[2], 6.25, col_w3, 0.40,
         9, True, _WHITE, align=PP_ALIGN.CENTER)
    invals = d.get("llm", {}).get("verdict_invalidation") or [
        ("Détérioration macro",
         "Récession ou hausse taux matérielle — revoir le positionnément."),
        ("Révision bénéfices",
         "Profit warning > -10 % sur 2 trimestrès — signal de sortie."),
    ]
    y_inv = 6.85
    for it, ib in invals[:2]:
        _bullet_square(slide, xs[2], y_inv, _HOLD, size=0.20)
        _txb(slide, str(it)[:38], xs[2] + 0.32, y_inv - 0.02, col_w3 - 0.4, 0.45,
             8.5, True, _NAVY)
        _txb(slide, str(ib)[:120], xs[2] + 0.32, y_inv + 0.45, col_w3 - 0.4, 1.40,
             7.5, False, _GRAYT, wrap=True)
        y_inv += 1.95

    # Horizon & conviction footer
    _rect(slide, 0.9, 10.85, 23.6, 0.55, _GRAYL,
          line_col=_GRAYM, line_w=0.4)
    _txb(slide,
         "Horizon recommandé : 6-12 mois  |  Conviction FORTE si écart score > 15 pts  |  "
         "Rééquilibrer si écart < 5 pts  |  Révisable à chaque publication trimestrielle",
         1.0, 10.95, 23.4, 0.40, 8, True, _NAVY, italic=True,
         align=PP_ALIGN.CENTER)


# ── Slide 21 — Méthodologie & Mentions légales ──────────────────────────────

def _slide_21_methodologie(prs, d: dict):
    slide = _blank(prs)
    _rect(slide, 0, 0, 25.4, 1.65, _NAVY)
    _txb(slide, "Méthodologie & Mentions Légales", 0.9, 0.30, 23.6, 1.10,
         13, True, _WHITE)

    col_w = 11.2
    xa, xb = 0.9, 13.1

    # === MÉTHODOLOGIE ===
    _rect(slide, xa, 2.0, col_w, 0.42, _NAVY)
    _txb(slide, "MÉTHODOLOGIE", xa, 2.03, col_w, 0.36,
         9, True, _WHITE, align=PP_ALIGN.CENTER)

    methodo = [
        ("Scoring FinSight (0-100)",
         "Indicateur propriétaire agrégeant performance ajustée du risque (Sharpe), "
         "valorisation relative (P/E Forward, P/B), résilience (Max Drawdown) et qualité "
         "de la rémunération de l'actionnaire (rendement dividende). Seuils d'allocation : "
         "≥ 65 = Surpondérer, 45-64 = Neutre, < 45 = Sous-pondérer."),
        ("Equity Risk Premium (ERP)",
         "Calculé comme la différence entre la performance 1 an de l'indice et le taux sans "
         "risque local (proxy : taux 10 ans souverain). Un ERP > 4 % est historiquement attractif "
         "vs obligations, < 2 % signale une prime tendue exposant à un re-rating baissier."),
        ("Volatilité, Sharpe et Drawdown",
         "Vol annualisée = écart-type des rendements quotidiens × √252. Sharpe = "
         "(Perf annualisée - taux sans risque) / Vol. Max Drawdown = perte maximale du plus "
         "haut au plus bas sur la fenêtre 12 mois. Calculs sur données yfinance ajustées."),
        ("Sources et limites",
         "Données yfinance (Yahoo Finance) — cours, P/E forward, P/B, rendements dividende. "
         "Composition sectorielle approximée pour l'indice B (poids fournis par yfinance ou "
         "estimations sur les principaux indices). Top constituants : top 5 par market cap. "
         "Données retardées de 24 h sur free tier."),
    ]
    y = 2.55
    for title, text in methodo:
        _txb(slide, title, xa, y, col_w, 0.34, 7.5, True, _NAVY)
        _txb(slide, text, xa, y + 0.38, col_w, 1.65, 6.5, False, _GRAYT, wrap=True)
        y += 2.05

    # === MENTIONS LÉGALES ===
    _rect(slide, xb, 2.0, col_w, 0.42, _NAVYL)
    _txb(slide, "MENTIONS LÉGALES", xb, 2.03, col_w, 0.36,
         9, True, _WHITE, align=PP_ALIGN.CENTER)

    legals = [
        ("Caractère informatif et pédagogique",
         "Ce document est produit exclusivement à des fins d'information et de démonstration "
         "pédagogique. Il ne constitue en aucun cas un conseil en investissement, une recommandation "
         "personnalisée, une incitation à l'achat ou à la vente, ni une offre ou sollicitation de "
         "souscription au sens du Règlement Général de l'AMF. Toute décision d'investissement "
         "demeure de la seule responsabilité de l'utilisateur."),
        ("Absence de due diligence",
         "FinSight IA est un outil algorithmique de screening basé sur des données publiques. "
         "Les analyses sont générées automatiquement par des modèles statistiques et un LLM, sans "
         "validation manuelle, sans rencontre avec les émetteurs des indices, sans audit des "
         "compositions. Les modèles peuvent contenir des biais ou erreurs de spécification. "
         "FinSight IA et ses auteurs déclinent toute responsabilité quant aux pertes décolant de "
         "l'utilisation de ce document."),
        ("Risques d'investissement",
         "Tout investissement en valeurs mobilières comporte un risque de perte partielle ou totale "
         "en capital. Les performances passées ne préjugent pas des performances futures. Les "
         "conditions de marché, le contexte macroéconomique, les décisions des banques centrales et "
         "les évènements géopolitiques peuvent évoluer rapidement et invalider les signaux présentés. "
         "Diversification et horizon adapté fortement recommandés."),
        ("Confidentialité et propriété intellectuelle",
         "Ce document est destiné à un usage privé et confidentiel. Sa reproduction, distribution "
         "ou diffusion, même partielle, est strictement interdite sans autorisation écrite expresse. "
         "Le scoring FinSight, la méthodologie et les visuels sont la propriété intellectuelle "
         "exclusive de FinSight IA. Toute exploitation commerciale est prohibée."),
    ]
    y = 2.55
    for title, text in legals:
        _txb(slide, title, xb, y, col_w, 0.34, 7.5, True, _NAVY)
        _txb(slide, text, xb, y + 0.38, col_w, 1.65, 6.5, False, _GRAYT, wrap=True)
        y += 2.05

    _rect(slide, 0.9, 13.30, 23.6, 0.25, _GRAYL)
    _txb(slide, f"Généré par FinSight IA  |  {d.get('date', _fr_date())}  |  Usage confidentiel",
         0.9, 13.32, 23.6, 0.22, 7, True, _NAVY, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# LLM enrichi — 12+ champs avec règle accents française
# ═════════════════════════════════════════════════════════════════════════════

def _generate_indice_llm(d: dict) -> dict:
    """Un seul appel LLM pour générer tous les textes analytiques du pitchbook indices."""
    try:
        from core.llm_provider import LLMProvider
        llm = LLMProvider(provider="mistral", model="mistral-small-latest")
        name_a = d.get("name_a", "Indice A")
        name_b = d.get("name_b", "Indice B")
        sig_a = d.get("signal_a", "Neutre")
        sig_b = d.get("signal_b", "Neutre")
        sc_a = d.get("score_a", 50)
        sc_b = d.get("score_b", 50)
        cur_a = d.get("currency_a", "USD")
        cur_b = d.get("currency_b", "USD")

        def _f(v, default=0):
            try:
                return float(v) if v is not None else default
            except Exception:
                return default

        ytd_a = _f(d.get("perf_ytd_a")) * 100 if abs(_f(d.get("perf_ytd_a"))) <= 2 else _f(d.get("perf_ytd_a"))
        ytd_b = _f(d.get("perf_ytd_b")) * 100 if abs(_f(d.get("perf_ytd_b"))) <= 2 else _f(d.get("perf_ytd_b"))
        p1a = _f(d.get("perf_1y_a")) * 100 if abs(_f(d.get("perf_1y_a"))) <= 2 else _f(d.get("perf_1y_a"))
        p1b = _f(d.get("perf_1y_b")) * 100 if abs(_f(d.get("perf_1y_b"))) <= 2 else _f(d.get("perf_1y_b"))
        p3a = _f(d.get("perf_3y_a")) * 100 if abs(_f(d.get("perf_3y_a"))) <= 2 else _f(d.get("perf_3y_a"))
        p3b = _f(d.get("perf_3y_b")) * 100 if abs(_f(d.get("perf_3y_b"))) <= 2 else _f(d.get("perf_3y_b"))
        p5a = _f(d.get("perf_5y_a")) * 100 if abs(_f(d.get("perf_5y_a"))) <= 2 else _f(d.get("perf_5y_a"))
        p5b = _f(d.get("perf_5y_b")) * 100 if abs(_f(d.get("perf_5y_b"))) <= 2 else _f(d.get("perf_5y_b"))
        vol_a = _f(d.get("vol_1y_a"))
        vol_b = _f(d.get("vol_1y_b"))
        sha_a = _f(d.get("sharpe_1y_a"))
        sha_b = _f(d.get("sharpe_1y_b"))
        dd_a = _f(d.get("max_dd_a")) * 100 if abs(_f(d.get("max_dd_a"))) <= 2 else _f(d.get("max_dd_a"))
        dd_b = _f(d.get("max_dd_b")) * 100 if abs(_f(d.get("max_dd_b"))) <= 2 else _f(d.get("max_dd_b"))
        pe_a = _f(d.get("pe_fwd_a"))
        pe_b = _f(d.get("pe_fwd_b"))
        pb_a = _f(d.get("pb_a"))
        pb_b = _f(d.get("pb_b"))
        dy_a = _f(d.get("div_yield_a")) * 100 if abs(_f(d.get("div_yield_a"))) < 1 else _f(d.get("div_yield_a"))
        dy_b = _f(d.get("div_yield_b")) * 100 if abs(_f(d.get("div_yield_b"))) < 1 else _f(d.get("div_yield_b"))
        erp_a = d.get("erp_a", "—")
        erp_b = d.get("erp_b", "—")

        sectors_str = ""
        sec_cmp = d.get("sector_comparison", [])
        if sec_cmp:
            top5_sec = sec_cmp[:5]
            sectors_str = " | ".join(
                f"{s[0]}: {s[1] or 0:.1f}% vs {s[2] or 0:.1f}%" for s in top5_sec
            )

        prompt = (
            f"Tu es un analyste sell-side senior (style JPMorgan Research, Goldman Sachs). "
            f"Rédige des textes analytiques rigoureux et techniques pour un pitchbook comparatif "
            f"d'indices boursiers : {name_a} vs {name_b}.\n\n"
            f"DONNÉES :\n"
            f"- {name_a} ({cur_a}) : Score {sc_a}/100 ({sig_a}), Perf YTD {ytd_a:+.1f}%, "
            f"1Y {p1a:+.1f}%, 3Y {p3a:+.1f}%, 5Y {p5a:+.1f}%, Vol {vol_a:.1f}%, "
            f"Sharpe {sha_a:.2f}, Max DD {dd_a:.1f}%, P/E Fwd {pe_a:.1f}x, "
            f"P/B {pb_a:.1f}x, Div Yield {dy_a:.1f}%, ERP {erp_a}\n"
            f"- {name_b} ({cur_b}) : Score {sc_b}/100 ({sig_b}), Perf YTD {ytd_b:+.1f}%, "
            f"1Y {p1b:+.1f}%, 3Y {p3b:+.1f}%, 5Y {p5b:+.1f}%, Vol {vol_b:.1f}%, "
            f"Sharpe {sha_b:.2f}, Max DD {dd_b:.1f}%, P/E Fwd {pe_b:.1f}x, "
            f"P/B {pb_b:.1f}x, Div Yield {dy_b:.1f}%, ERP {erp_b}\n"
            f"- Top 5 secteurs (poids A vs B) : {sectors_str}\n\n"
            f"RÈGLES DE RÉDACTION :\n"
            f"- Français correct et complet : utilise TOUS les accents (é è ê à ù ç î ô), "
            f"cédilles, apostrophes droites ' et guillemets français « ». N'ÉCRIS JAMAIS sans "
            f"accents — ce serait du français cassé, inacceptable en rapport IB-grade.\n"
            f"- Style sell-side senior : prose technique, raisonnements économiques reliés aux "
            f"chiffres fournis, pas de formules génériques.\n"
            f"- Cite les chiffres précis fournis ci-dessus, ne jamais inventer ni dire 'données indisponibles'.\n"
            f"- Pas de markdown **, pas d'emojis.\n"
            f"- Maximum indiqué par champ.\n"
            f"- Réponds UNIQUEMENT en JSON valide, aucun texte avant ou après.\n\n"
            f'{{\n'
            f'  "exec_summary": "120 mots MAX : synthèse globale, qui privilégier et pourquoi avec 4 chiffres concrets",\n'
            f'  "profil_a": "100 mots MAX : caractéristiques structurelles {name_a}, devise, biais sectoriel, sensibilités macro",\n'
            f'  "profil_b": "100 mots MAX : caractéristiques structurelles {name_b}, devise, biais sectoriel, sensibilités macro",\n'
            f'  "profil_a_read": "80 mots MAX : lecture analytique du profil {name_a}, signal et horizon",\n'
            f'  "profil_b_read": "80 mots MAX : lecture analytique du profil {name_b}, signal et horizon",\n'
            f'  "sectoral_read": "120 mots MAX : analyse de la composition sectorielle, biais structurels et implications",\n'
            f'  "sectoral_impl": "100 mots MAX : implications pour l allocation selon le régime macro",\n'
            f'  "concentration_read": "120 mots MAX : analyse de la concentration des top constituants, risque idiosyncrasique",\n'
            f'  "perf_macro_read": "120 mots MAX : contexte macro 12 mois (Fed, BCE, dollar, cycle), événements marquants",\n'
            f'  "perf_chart_read": "100 mots MAX : lecture analytique de la trajectoire boursière, surperformance et explications",\n'
            f'  "perf_decomp_read": "120 mots MAX : analyse décomposition perfs YTD/1Y/3Y/5Y, momentum vs structurel",\n'
            f'  "perf_horizon_read": "80 mots MAX : implications selon horizon d investissement",\n'
            f'  "risque_read": "120 mots MAX : lecture du profil de risque, vol/Sharpe/drawdown comparés",\n'
            f'  "risk_return_read": "120 mots MAX : analyse du couple rendement-risque, efficience relative",\n'
            f'  "risk_return_impl": "80 mots MAX : implications pour la construction de portefeuille",\n'
            f'  "valorisation_read": "120 mots MAX : analyse multiples P/E P/B div yield, prime/décote, mean reversion",\n'
            f'  "erp_read": "120 mots MAX : interprétation ERP, prime de risque actions vs obligations, signal d allocation",\n'
            f'  "score_read": "120 mots MAX : lecture du score FinSight composite, écart relatif et signal global",\n'
            f'  "verdict_read": "150 mots MAX : verdict final d allocation, indice à privilégier, raisons quantifiées, horizon",\n'
            f'  "bear_combined": "100 mots MAX : risques principaux des deux indices combinés"\n'
            f'}}'
        )
        import json, re
        resp = llm.generate(prompt, max_tokens=2200)
        m = re.search(r'\{.*\}', resp, re.DOTALL)
        if m:
            data_out = json.loads(m.group(0))
            log.info("[indice_cmp_pptx] LLM texts OK (%d champs)", len(data_out))
            return data_out
    except Exception as e:
        log.warning("[indice_cmp_pptx] LLM génération failed: %s", e)
    return {}


# ═════════════════════════════════════════════════════════════════════════════
# Classe principale
# ═════════════════════════════════════════════════════════════════════════════

class CmpIndicePPTXWriter:
    """Pitchbook comparatif d'indices boursiers — 21 slides IB-grade."""

    @staticmethod
    def generate(data: dict, output_path: str = None, language: str = "fr", currency: str = "EUR") -> bytes:
        """Génère le PPTX comparatif et retourne les bytes."""
        data = dict(data)  # copie défensive
        data["llm"] = _generate_indice_llm(data)
        data["_language"] = language
        data["_currency"] = currency
        _lang = (language or "fr").lower()[:2]
        if _lang not in {"fr","en","es","de","it","pt"}:
            _lang = "fr"
        _T = {
            "exec_sum": {"fr": "Executive Summary", "en": "Executive Summary",
                         "es": "Resumen Ejecutivo", "de": "Executive Summary",
                         "it": "Executive Summary", "pt": "Resumo Executivo"},
            "sommaire": {"fr": "Sommaire", "en": "Table of contents",
                         "es": "Índice", "de": "Inhalt",
                         "it": "Sommario", "pt": "Sumário"},
            "profil": {"fr": "Profil", "en": "Profile",
                       "es": "Perfil", "de": "Profil",
                       "it": "Profilo", "pt": "Perfil"},
            "compo_sec": {"fr": "Composition Sectorielle Comparée",
                          "en": "Compared Sector Composition",
                          "es": "Composición Sectorial Comparada",
                          "de": "Verglichene Sektor-Zusammensetzung",
                          "it": "Composizione Settoriale Comparata",
                          "pt": "Composição Setorial Comparada"},
            "top_const": {"fr": "Top Constituants Comparés",
                          "en": "Compared Top Constituents",
                          "es": "Top Constituyentes Comparados",
                          "de": "Top-Bestandteile Verglichen",
                          "it": "Top Costituenti Comparati",
                          "pt": "Top Constituintes Comparados"},
            "perf_52w": {"fr": "Performance Historique  —  52 Semaines",
                         "en": "Historical Performance  —  52 Weeks",
                         "es": "Rendimiento Histórico  —  52 Semanas",
                         "de": "Historische Performance  —  52 Wochen",
                         "it": "Performance Storica  —  52 Settimane",
                         "pt": "Desempenho Histórico  —  52 Semanas"},
            "decompo_perf": {"fr": "Décomposition de la Performance",
                             "en": "Performance Breakdown",
                             "es": "Descomposición del Rendimiento",
                             "de": "Performance-Aufschlüsselung",
                             "it": "Scomposizione della Performance",
                             "pt": "Decomposição do Desempenho"},
            "risque_cmp": {"fr": "Risque Comparatif", "en": "Comparative Risk",
                           "es": "Riesgo Comparativo", "de": "Vergleichendes Risiko",
                           "it": "Rischio Comparativo", "pt": "Risco Comparativo"},
            "rend_risque": {"fr": "Profil Rendement / Risque",
                            "en": "Risk / Return Profile",
                            "es": "Perfil Rentabilidad / Riesgo",
                            "de": "Rendite-/Risikoprofil",
                            "it": "Profilo Rendimento / Rischio",
                            "pt": "Perfil Retorno / Risco"},
            "valo_cmp": {"fr": "Valorisation Comparée",
                         "en": "Compared Valuation",
                         "es": "Valoración Comparada",
                         "de": "Verglichene Bewertung",
                         "it": "Valutazione Comparata",
                         "pt": "Avaliação Comparada"},
            "erp": {"fr": "Equity Risk Premium  —  Prime de Risque Actions",
                    "en": "Equity Risk Premium  —  Stocks Risk Premium",
                    "es": "Prima de Riesgo de Acciones  —  ERP",
                    "de": "Equity Risk Premium  —  Aktien-Risikoprämie",
                    "it": "Equity Risk Premium  —  Premio al Rischio Azioni",
                    "pt": "Equity Risk Premium  —  Prémio de Risco Acções"},
            "signal_fs": {"fr": "Signal FinSight Comparatif",
                          "en": "Comparative FinSight Signal",
                          "es": "Señal FinSight Comparativa",
                          "de": "Vergleichendes FinSight-Signal",
                          "it": "Segnale FinSight Comparativo",
                          "pt": "Sinal FinSight Comparativo"},
            "theses_bb": {"fr": "Thèses d'Investissement Bull / Bear",
                          "en": "Bull / Bear Investment Theses",
                          "es": "Tesis de Inversión Bull / Bear",
                          "de": "Bull-/Bear-Investitionsthesen",
                          "it": "Tesi d'Investimento Bull / Bear",
                          "pt": "Teses de Investimento Bull / Bear"},
            "verdict": {"fr": "Verdict Comparatif  —  Conviction d'Allocation",
                        "en": "Comparative Verdict  —  Allocation Conviction",
                        "es": "Veredicto Comparativo  —  Convicción de Asignación",
                        "de": "Vergleichsurteil  —  Allokations-Überzeugung",
                        "it": "Verdetto Comparativo  —  Convinzione di Allocazione",
                        "pt": "Veredicto Comparativo  —  Convicção de Alocação"},
        }
        def _t(k):
            return _T.get(k, {}).get(_lang) or _T.get(k, {}).get("en") or k
        data["_t_helper"] = _t

        prs = Presentation()
        prs.slide_width = _SW
        prs.slide_height = _SH

        # Section 0 — Introduction
        _slide_01_cover(prs, data)             # 1
        _slide_02_exec_summary(prs, data)      # 2
        _slide_03_sommaire(prs, data)          # 3

        # Section 1 — Profil & Composition
        _divider(prs, "01", "Profil & Composition",
                 "Structure, secteurs et top constituants des deux indices")  # 4
        _slide_05_profil_a(prs, data)          # 5
        _slide_06_profil_b(prs, data)          # 6
        _slide_07_secteurs(prs, data)          # 7
        _slide_08_top_holdings(prs, data)      # 8

        # Section 2 — Performance & Risque
        _divider(prs, "02", "Performance & Risque",
                 "Cours 52S, décomposition perfs et métriques de risque")    # 9
        _slide_10_perf_chart(prs, data)        # 10
        _slide_11_perf_decomposition(prs, data)  # 11
        _slide_12_risque(prs, data)            # 12
        _slide_13_risk_return(prs, data)       # 13

        # Section 3 — Valorisation & Qualité
        _divider(prs, "03", "Valorisation & Qualité",
                 "Multiples, ERP et signal FinSight composite")               # 14
        _slide_15_valorisation(prs, data)      # 15
        _slide_16_erp_focus(prs, data)         # 16
        _slide_17_score(prs, data)             # 17

        # Section 4 — Décision d'allocation
        _divider(prs, "04", "Décision d'Allocation",
                 "Thèses bull/bear, verdict et conditions d'invalidation")    # 18
        _slide_19_theses(prs, data)            # 19
        _slide_20_verdict(prs, data)           # 20
        _slide_21_methodologie(prs, data)      # 21

        buf = io.BytesIO()
        prs.save(buf)
        pptx_bytes = buf.getvalue()

        if output_path:
            from pathlib import Path
            Path(output_path).write_bytes(pptx_bytes)
            log.info("[CmpIndicePPTXWriter] Sauvegardé : %s (%d Ko)",
                     output_path, len(pptx_bytes) // 1024)

        return pptx_bytes
