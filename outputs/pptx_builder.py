# =============================================================================
# FinSight IA — PPTX Builder (Pitchbook)
# outputs/pptx_builder.py
#
# Génère un pitchbook IB en 6 slides via python-pptx.
# Style : fond blanc, accents bleu marine, police Calibri.
#
# Slides :
#   1. Cover         — société, ticker, date, recommandation badge
#   2. Financials    — tableau ratios clés (année la plus récente)
#   3. Recommandation — BUY/HOLD/SELL, conviction, target prices, summary
#   4. Forces & Risques — 3+3 bullets
#   5. QA + Devil    — score QA, thèse inverse
#   6. Disclaimer
#
# Usage :
#   builder = PPTXBuilder()
#   path    = builder.build(snapshot, ratios, synthesis, qa_python, devil)
# =============================================================================

from __future__ import annotations

import logging
import os
from datetime import date
from pathlib import Path
from typing import Optional

log = logging.getLogger(__name__)

_OUTPUT_DIR = Path(__file__).parent / "generated"

# ---------------------------------------------------------------------------
# CONFIG — seuils et limites d'affichage
# ---------------------------------------------------------------------------
MAX_STRENGTHS_DISPLAYED  = 3
MAX_RISKS_DISPLAYED      = 3
MAX_QA_FLAGS_DISPLAYED   = 5
MAX_ASSUMPTIONS_DISPLAYED = 2

ALTMAN_Z_SAINE           = 2.99   # z >= seuil → santé financière saine
ALTMAN_Z_GRISE           = 1.81   # z >= seuil → zone grise
CONVICTION_DELTA_ROBUSTE = 0.2    # delta > seuil → thèse robuste
CONVICTION_DELTA_FRAGILE = -0.2   # delta < seuil → thèse fragile

# ---------------------------------------------------------------------------
# Palette couleurs IB
# ---------------------------------------------------------------------------

_NAVY   = "1B2A4A"   # bleu marine — titres
_WHITE  = "FFFFFF"
_LIGHT  = "F0F4F8"   # fond slides pair
_ACCENT = "2E6DA4"   # bleu accent — badges
_GREEN  = "1A7A4A"   # BUY
_RED    = "B22222"   # SELL
_ORANGE = "CC6600"   # HOLD
_GREY   = "6C757D"   # texte secondaire
_DARK   = "212529"   # texte principal


def _rgb(hex_str: str):
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    r = int(hex_str[0:2], 16)
    g = int(hex_str[2:4], 16)
    b = int(hex_str[4:6], 16)
    return RGBColor(r, g, b)


# ---------------------------------------------------------------------------
# PPTXBuilder
# ---------------------------------------------------------------------------

class PPTXBuilder:
    """
    Pitchbook IB — 6 slides, python-pptx, zéro LLM.
    """

    def build(
        self,
        snapshot,
        ratios,
        synthesis,
        qa_python=None,
        devil=None,
        output_path: Optional[Path] = None,
    ) -> Path:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
        except ImportError:
            raise RuntimeError("python-pptx requis : pip install python-pptx")

        ticker = snapshot.ticker.replace(".", "_")
        today  = date.today().isoformat()

        if output_path is None:
            _OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
            output_path = _OUTPUT_DIR / f"{ticker}_{today}_pitchbook.pptx"

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        blank = prs.slide_layouts[6]  # layout vide

        ci   = snapshot.company_info
        mkt  = snapshot.market
        latest = ratios.latest_year
        yr     = ratios.years.get(latest)

        # ------------------------------------------------------------------
        # Slide 1 — Cover
        # ------------------------------------------------------------------
        sl = prs.slides.add_slide(blank)
        _fill_bg(sl, _NAVY)

        # Bande accent
        _rect(sl, 0, Inches(4.8), Inches(13.33), Inches(0.06), _ACCENT)

        _text_box(sl, Inches(0.8), Inches(1.0), Inches(11.0), Inches(1.2),
                  ci.company_name, 36, _WHITE, bold=True)

        subtitle = f"{ci.ticker}  |  {ci.sector or 'Analyse FinSight IA'}  |  {ci.currency}"
        _text_box(sl, Inches(0.8), Inches(2.3), Inches(11.0), Inches(0.6),
                  subtitle, 16, _ACCENT)

        if synthesis:
            reco_color = {"BUY": _GREEN, "SELL": _RED}.get(synthesis.recommendation, _ORANGE)
            _rect(sl, Inches(0.8), Inches(3.2), Inches(2.4), Inches(0.8), reco_color)
            _text_box(sl, Inches(0.8), Inches(3.2), Inches(2.4), Inches(0.8),
                      synthesis.recommendation, 24, _WHITE, bold=True, center=True)
            _conv = synthesis.conviction or 0
            _conf = synthesis.confidence_score or 0
            _text_box(sl, Inches(3.4), Inches(3.3), Inches(5.0), Inches(0.6),
                      f"Conviction : {_conv:.0%}  |  Confiance : {_conf:.0%}",
                      14, _LIGHT)

        _text_box(sl, Inches(0.8), Inches(6.5), Inches(11.0), Inches(0.4),
                  f"FinSight IA  |  {today}  |  Confidentiel — usage interne",
                  11, _GREY)

        # ------------------------------------------------------------------
        # Slide 2 — Ratios financiers clés
        # ------------------------------------------------------------------
        sl = prs.slides.add_slide(blank)
        _fill_bg(sl, _WHITE)
        _slide_header(sl, f"Ratios Financiers Clés  —  {latest}", ci.company_name)

        if yr:
            def _p(v): return f"{v*100:.1f}%" if v is not None else "N/A"
            def _x(v): return f"{v:.2f}x"     if v is not None else "N/A"
            def _n(v): return f"{v:.2f}"      if v is not None else "N/A"

            rows = [
                ("Gross Margin",      _p(yr.gross_margin)),
                ("EBITDA Margin",     _p(yr.ebitda_margin)),
                ("EBIT Margin",       _p(yr.ebit_margin)),
                ("Net Margin",        _p(yr.net_margin)),
                ("ROE",               _p(yr.roe)),
                ("ROIC",              _p(yr.roic)),
                ("Net Debt / EBITDA", _x(yr.net_debt_ebitda)),
                ("EV / EBITDA",       _x(yr.ev_ebitda)),
                ("P/E Ratio",         _x(yr.pe_ratio)),
                ("FCF Yield",         _p(yr.fcf_yield)),
                ("Current Ratio",     _n(yr.current_ratio)),
                ("Altman Z-Score",    _n(yr.altman_z)),
            ]
            if yr.beneish_m is not None:
                rows.append(("Beneish M-Score", f"{yr.beneish_m:.3f}"))
            if yr.revenue_growth is not None:
                rows.append(("Revenue Growth YoY", _p(yr.revenue_growth)))

            # 2 colonnes
            mid = (len(rows) + 1) // 2
            col1, col2 = rows[:mid], rows[mid:]

            _ratio_column(sl, Inches(0.5), Inches(1.4), col1)
            _ratio_column(sl, Inches(7.0), Inches(1.4), col2)

            # Altman Z flag
            if yr.altman_z is not None:
                z = yr.altman_z
                z_label = "Sante financière : SAINE" if z >= ALTMAN_Z_SAINE else (
                    "Sante financière : ZONE GRISE" if z >= ALTMAN_Z_GRISE else "DETRESSE FINANCIERE")
                z_color = _GREEN if z >= ALTMAN_Z_SAINE else (_ORANGE if z >= ALTMAN_Z_GRISE else _RED)
                _rect(sl, Inches(0.5), Inches(6.6), Inches(4.5), Inches(0.45), z_color)
                _text_box(sl, Inches(0.5), Inches(6.6), Inches(4.5), Inches(0.45),
                          z_label, 11, _WHITE, bold=True)

            # Cours actuel
            if mkt.share_price:
                price_str = f"Cours : {mkt.share_price:.2f} {ci.currency}"
                _text_box(sl, Inches(7.0), Inches(6.6), Inches(5.5), Inches(0.45),
                          price_str, 12, _DARK)

        # ------------------------------------------------------------------
        # Slide 3 — Recommandation
        # ------------------------------------------------------------------
        sl = prs.slides.add_slide(blank)
        _fill_bg(sl, _WHITE)
        _slide_header(sl, "Recommandation & Valorisation", ci.company_name)

        if synthesis:
            reco_color = {"BUY": _GREEN, "SELL": _RED}.get(synthesis.recommendation, _ORANGE)

            # Badge reco
            _rect(sl, Inches(0.5), Inches(1.3), Inches(3.2), Inches(1.2), reco_color)
            _text_box(sl, Inches(0.5), Inches(1.3), Inches(3.2), Inches(1.2),
                      synthesis.recommendation, 36, _WHITE, bold=True, center=True)

            # Conviction / Confiance
            _conv = synthesis.conviction or 0
            _conf = synthesis.confidence_score or 0
            _text_box(sl, Inches(4.0), Inches(1.4), Inches(8.5), Inches(0.4),
                      f"Conviction : {_conv:.0%}   |   Confiance IA : {_conf:.0%}",
                      14, _DARK)

            # Target prices
            targets = []
            if synthesis.target_bear is not None:
                targets.append(f"Bear : {synthesis.target_bear:.0f} {ci.currency}")
            if synthesis.target_base is not None:
                targets.append(f"Base : {synthesis.target_base:.0f} {ci.currency}")
            if synthesis.target_bull is not None:
                targets.append(f"Bull : {synthesis.target_bull:.0f} {ci.currency}")
            if targets:
                _text_box(sl, Inches(4.0), Inches(1.9), Inches(8.5), Inches(0.4),
                          "   |   ".join(targets), 13, _ACCENT, bold=True)

            # Summary
            if synthesis.summary:
                _text_box(sl, Inches(0.5), Inches(2.8), Inches(12.3), Inches(1.4),
                          synthesis.summary, 13, _DARK)

            # Valuation comment
            if synthesis.valuation_comment:
                _text_box(sl, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.8),
                          f"Valorisation : {synthesis.valuation_comment}", 12, _GREY)

            # Invalidation
            if synthesis.invalidation_conditions:
                _text_box(sl, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.6),
                          f"Invalidation : {synthesis.invalidation_conditions[:200]}", 10, _RED)

        # ------------------------------------------------------------------
        # Slide 4 — Forces & Risques
        # ------------------------------------------------------------------
        sl = prs.slides.add_slide(blank)
        _fill_bg(sl, _WHITE)
        _slide_header(sl, "Points Forts & Risques", ci.company_name)

        if synthesis:
            # Forces (gauche)
            _rect(sl, Inches(0.5), Inches(1.3), Inches(5.8), Inches(0.45), _GREEN)
            _text_box(sl, Inches(0.5), Inches(1.3), Inches(5.8), Inches(0.45),
                      "POINTS FORTS", 13, _WHITE, bold=True)
            _strengths = synthesis.strengths or []
            if len(_strengths) > MAX_STRENGTHS_DISPLAYED:
                log.warning("pptx_builder slide4: %d strengths disponibles — "
                            "affichage tronque a MAX_STRENGTHS_DISPLAYED=%d",
                            len(_strengths), MAX_STRENGTHS_DISPLAYED)
            for i, s in enumerate(_strengths[:MAX_STRENGTHS_DISPLAYED]):
                _text_box(sl, Inches(0.6), Inches(1.9 + i * 0.9), Inches(5.6), Inches(0.75),
                          f"+ {s}", 12, _DARK)

            # Risques (droite)
            _rect(sl, Inches(7.0), Inches(1.3), Inches(5.8), Inches(0.45), _RED)
            _text_box(sl, Inches(7.0), Inches(1.3), Inches(5.8), Inches(0.45),
                      "RISQUES IDENTIFIES", 13, _WHITE, bold=True)
            _risks = synthesis.risks or []
            if len(_risks) > MAX_RISKS_DISPLAYED:
                log.warning("pptx_builder slide4: %d risks disponibles — "
                            "affichage tronque a MAX_RISKS_DISPLAYED=%d",
                            len(_risks), MAX_RISKS_DISPLAYED)
            for i, r in enumerate(_risks[:MAX_RISKS_DISPLAYED]):
                _text_box(sl, Inches(7.1), Inches(1.9 + i * 0.9), Inches(5.6), Inches(0.75),
                          f"- {r}", 12, _DARK)

            # Ligne séparatrice verticale
            _rect(sl, Inches(6.6), Inches(1.2), Inches(0.04), Inches(4.5), _LIGHT)

        # ------------------------------------------------------------------
        # Slide 5 — QA & Devil's Advocate
        # ------------------------------------------------------------------
        sl = prs.slides.add_slide(blank)
        _fill_bg(sl, _WHITE)
        _slide_header(sl, "Validation QA & Avocat du Diable", ci.company_name)

        y = Inches(1.3)

        if qa_python:
            score_color = _GREEN if qa_python.qa_score >= 0.85 else (
                _ORANGE if qa_python.qa_score >= 0.65 else _RED)
            _rect(sl, Inches(0.5), y, Inches(3.5), Inches(0.5), score_color)
            _text_box(sl, Inches(0.5), y, Inches(3.5), Inches(0.5),
                      f"QA Score : {qa_python.qa_score:.0%}  {'VALIDE' if qa_python.passed else 'ECHEC'}",
                      13, _WHITE, bold=True)
            y += Inches(0.6)

            _flags = qa_python.flags or []
            if len(_flags) > MAX_QA_FLAGS_DISPLAYED:
                log.warning("pptx_builder slide5: %d QA flags disponibles — "
                            "affichage tronque a MAX_QA_FLAGS_DISPLAYED=%d",
                            len(_flags), MAX_QA_FLAGS_DISPLAYED)
            for fl in _flags[:MAX_QA_FLAGS_DISPLAYED]:
                sym = "[E]" if fl.level == "ERROR" else "[W]" if fl.level == "WARNING" else "[i]"
                col = _RED if fl.level == "ERROR" else (_ORANGE if fl.level == "WARNING" else _GREY)
                _text_box(sl, Inches(0.5), y, Inches(12.0), Inches(0.35),
                          f"{sym}  {fl.message[:120]}", 10, col)
                y += Inches(0.38)

        if devil:
            y += Inches(0.2)
            delta_str = f"{devil.conviction_delta:+.2f}"
            solidity  = ("These fragile" if devil.conviction_delta < CONVICTION_DELTA_FRAGILE
                         else "These robuste" if devil.conviction_delta > CONVICTION_DELTA_ROBUSTE
                         else "These moderement solide")
            _rect(sl, Inches(0.5), y, Inches(12.3), Inches(0.45), _NAVY)
            _text_box(sl, Inches(0.5), y, Inches(12.3), Inches(0.45),
                      f"Avocat du Diable : {devil.original_reco} → {devil.counter_reco}  |  "
                      f"Delta conviction : {delta_str}  ({solidity})",
                      12, _WHITE, bold=True)
            y += Inches(0.55)

            if devil.counter_thesis:
                _text_box(sl, Inches(0.5), y, Inches(12.3), Inches(0.7),
                          devil.counter_thesis[:250], 11, _DARK)
                y += Inches(0.8)

            _assumptions = devil.key_assumptions or []
            if len(_assumptions) > MAX_ASSUMPTIONS_DISPLAYED:
                log.warning("pptx_builder slide5: %d key_assumptions disponibles — "
                            "affichage tronque a MAX_ASSUMPTIONS_DISPLAYED=%d",
                            len(_assumptions), MAX_ASSUMPTIONS_DISPLAYED)
            for a in _assumptions[:MAX_ASSUMPTIONS_DISPLAYED]:
                _text_box(sl, Inches(0.5), y, Inches(12.3), Inches(0.35),
                          f"? Hypothese fragile : {a[:100]}", 10, _ORANGE)
                y += Inches(0.38)

        # ------------------------------------------------------------------
        # Slide 6 — Disclaimer
        # ------------------------------------------------------------------
        sl = prs.slides.add_slide(blank)
        _fill_bg(sl, _LIGHT)

        _text_box(sl, Inches(1.5), Inches(1.5), Inches(10.3), Inches(0.6),
                  "AVERTISSEMENT IMPORTANT", 18, _NAVY, bold=True)

        disclaimer = (
            "Ce document a été généré automatiquement par FinSight IA à titre d'information "
            "et d'aide à la décision uniquement. Il ne constitue pas un conseil en investissement, "
            "une recommandation d'achat ou de vente de valeurs mobilières, ni une offre ou "
            "sollicitation de transaction. Les informations contenues dans ce document proviennent "
            "de sources publiques considérées comme fiables, mais leur exactitude n'est pas garantie. "
            "Les performances passées ne préjugent pas des performances futures. "
            "FinSight IA décline toute responsabilité quant aux décisions prises sur la base "
            "de ce document. Ce document est strictement confidentiel et destiné à un usage interne."
        )
        _text_box(sl, Inches(1.5), Inches(2.5), Inches(10.3), Inches(3.0),
                  disclaimer, 12, _DARK)

        _text_box(sl, Inches(1.5), Inches(6.2), Inches(10.3), Inches(0.4),
                  f"FinSight IA  |  {today}  |  Généré automatiquement",
                  10, _GREY)

        # ------------------------------------------------------------------
        # Sauvegarde
        # ------------------------------------------------------------------
        prs.save(str(output_path))
        log.info(f"[PPTXBuilder] Pitchbook généré : {output_path.name}")
        return output_path


# ---------------------------------------------------------------------------
# Helpers shapes
# ---------------------------------------------------------------------------

def _fill_bg(slide, hex_color: str) -> None:
    from pptx.util import Inches
    from pptx.enum.dml import MSO_THEME_COLOR
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(hex_color)


def _rect(slide, left, top, width, height, hex_color: str):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(hex_color)
    shape.line.fill.background()
    return shape


def _text_box(slide, left, top, width, height, text: str,
              font_size: int, hex_color: str,
              bold: bool = False, center: bool = False) -> None:
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = True

    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.color.rgb = _rgb(hex_color)
    run.font.bold  = bold
    run.font.name  = "Calibri"


def _slide_header(slide, title: str, company: str) -> None:
    from pptx.util import Inches
    _rect(slide, 0, 0, Inches(13.33), Inches(0.9), _NAVY)
    _text_box(slide, Inches(0.4), Inches(0.12), Inches(9.0), Inches(0.6),
              title, 18, _WHITE, bold=True)
    _text_box(slide, Inches(9.5), Inches(0.2), Inches(3.5), Inches(0.5),
              company, 11, _LIGHT)


def _ratio_column(slide, left, top, rows: list) -> None:
    from pptx.util import Inches
    for i, (label, value) in enumerate(rows):
        y = top + Inches(i * 0.38)
        bg = _LIGHT if i % 2 == 0 else _WHITE
        _rect(slide, left, y, Inches(6.0), Inches(0.36), bg)
        _text_box(slide, left + Inches(0.1), y, Inches(3.8), Inches(0.36),
                  label, 11, _DARK)
        _text_box(slide, left + Inches(4.0), y, Inches(1.9), Inches(0.36),
                  value, 11, _ACCENT, bold=True)
