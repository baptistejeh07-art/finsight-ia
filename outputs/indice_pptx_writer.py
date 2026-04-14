"""
indice_pptx_writer.py — FinSight IA
Pitchbook PPTX 20 slides pour analyse d'indice boursier.
Interface : IndicePPTXWriter.generate(data, output_path) -> bytes
Modele reference : FinSight_IA_Indice_SP500 (4).pptx
"""
from __future__ import annotations
import io, logging, re as _re
from typing import Optional

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.ticker as mticker
import numpy as np

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

log = logging.getLogger(__name__)

def _x(text) -> str:
    """Supprimé les caractères invalides XML 1.0 qui corrompent le PPTX."""
    if text is None:
        return ""
    return _re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', str(text))

# ── Dimensions (25.4 x 14.3 cm) ───────────────────────────────────────────────
_SW = Cm(25.4)
_SH = Cm(14.3)

# ── Palette ────────────────────────────────────────────────────────────────────
_NAVY  = RGBColor(0x1B, 0x3A, 0x6B)
_NAVYL = RGBColor(0x2A, 0x52, 0x98)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_GRAYL = RGBColor(0xF5, 0xF7, 0xFA)
_GRAYM = RGBColor(0xE8, 0xEC, 0xF0)
_GRAYT = RGBColor(0x55, 0x55, 0x55)
_GRAYD = RGBColor(0xAA, 0xAA, 0xAA)
_BLACK = RGBColor(0x1A, 0x1A, 0x1A)
_BUY   = RGBColor(0x1A, 0x7A, 0x4A)
_SELL  = RGBColor(0xA8, 0x20, 0x20)
_HOLD  = RGBColor(0xB0, 0x60, 0x00)
_BUY_L = RGBColor(0xE8, 0xF5, 0xEE)
_SELL_L= RGBColor(0xFB, 0xEB, 0xEB)
_HOLD_L= RGBColor(0xFD, 0xF3, 0xE5)
_NAVY_LIGHT = RGBColor(0xD0, 0xDC, 0xF5)

_fr_months = {1:"janvier",2:"fevrier",3:"mars",4:"avril",5:"mai",6:"juin",
              7:"juillet",8:"aout",9:"septembre",10:"octobre",11:"novembre",12:"decembre"}

# Couleurs lignes ETF (identiques a sectoral_pptx_writer)
_LINE_COLORS_HEX = ['#1B3A6B','#1A7A4A','#A82020','#B06000',
                    '#2A5298','#6B3A1B','#3A6B1B','#6B1B3A']
_LINE_COLORS_RGB = [
    RGBColor(0x1B,0x3A,0x6B), RGBColor(0x1A,0x7A,0x4A),
    RGBColor(0xA8,0x20,0x20), RGBColor(0xB0,0x60,0x00),
    RGBColor(0x2A,0x52,0x98), RGBColor(0x6B,0x3A,0x1B),
    RGBColor(0x3A,0x6B,0x1B), RGBColor(0x6B,0x1B,0x3A),
]


# ── Helpers generiques ─────────────────────────────────────────────────────────

def _abbrev_sector(name: str, maxlen: int = 16) -> str:
    """Retourne l'abreviation FR courte du secteur (i18n via core/sector_labels)."""
    if not name:
        return ""
    try:
        from core.sector_labels import slug_from_any
    except ImportError:
        return str(name)[:maxlen]
    slug = slug_from_any(name)
    if slug is None:
        s = str(name)
        return s[:maxlen] if len(s) > maxlen else s
    SHORT_FR = {
        "TECHNOLOGY":       "Technologie",
        "HEALTHCARE":       "Santé",
        "FINANCIALS":       "Finance",
        "CONSUMERCYCLICAL": "Conso. Cycl.",
        "CONSUMERDEFENSIVE":"Conso. Déf.",
        "ENERGY":           "Énergie",
        "INDUSTRIALS":      "Industrie",
        "MATERIALS":        "Matériaux",
        "REALESTATE":       "Immobilier",
        "UTILITIES":        "Serv. Publ.",
        "COMMUNICATION":    "Télécoms",
    }
    s = SHORT_FR.get(slug, str(name))
    return s[:maxlen] if len(s) > maxlen else s

def _trunc(text: str, n: int) -> str:
    """Truncate at word boundary."""
    if len(text) <= n:
        return text
    return text[:n].rsplit(' ', 1)[0] + '...'

def _fr_date():
    import datetime
    d = datetime.date.today()
    return f"{d.day} {_fr_months[d.month]} {d.year}"


def _sig_color(signal: str) -> RGBColor:
    s = str(signal)
    if "Surp" in s: return _BUY
    if "Sous" in s: return _SELL
    return _HOLD


def _sig_light(signal: str) -> RGBColor:
    s = str(signal)
    if "Surp" in s: return _BUY_L
    if "Sous" in s: return _SELL_L
    return _HOLD_L


def _sig_hex(signal: str) -> str:
    s = str(signal)
    if "Surp" in s: return '#1A7A4A'
    if "Sous" in s: return '#A82020'
    return '#B06000'


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, x, y, w, h, fill=None, line=False, line_col=None, line_w=0.5):
    shape = slide.shapes.add_shape(1, Cm(x), Cm(y), Cm(w), Cm(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line and line_col:
        shape.line.color.rgb = line_col
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape


def _txb(slide, text, x, y, w, h, size=9, bold=False, color=None,
         align=PP_ALIGN.LEFT, italic=False, wrap=True, vcenter=False):
    color = color or _BLACK
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    if vcenter:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = _x(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def _pic(slide, img_bytes, x, y, w, h):
    buf = io.BytesIO(img_bytes)
    slide.shapes.add_picture(buf, Cm(x), Cm(y), Cm(w), Cm(h))


def _add_table(slide, data, x, y, w, h, col_widths=None,
               header_fill=_NAVY, header_color=_WHITE,
               alt_fill=None, font_size=7.5, header_size=7.5):
    rows, cols = len(data), len(data[0]) if data else 1
    tbl_shape = slide.shapes.add_table(rows, cols, Cm(x), Cm(y), Cm(w), Cm(h))
    tbl = tbl_shape.table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Cm(cw * w / total)
    for r, row_data in enumerate(data):
        is_hdr = (r == 0)
        for c, val in enumerate(row_data):
            cell = tbl.cell(r, c)
            cell.text = _x(val) if val is not None else "—"
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            run = para.runs[0] if para.runs else para.add_run()
            run.font.size = Pt(header_size if is_hdr else font_size)
            run.font.bold = is_hdr
            if is_hdr:
                run.font.color.rgb = header_color
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
            else:
                run.font.color.rgb = _BLACK
                if alt_fill and r % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = alt_fill
                else:
                    cell.fill.background()
    return tbl


def _color_cell(tbl, row, col, fill, text_color=None):
    cell = tbl.cell(row, col)
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    if text_color:
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = text_color


# ── Navigation / header / footer ───────────────────────────────────────────────

_NAV = ["1","2","3","4","5"]


def _header(slide, title, subtitle, active=1):
    _rect(slide, 0, 0, 25.4, 1.4, fill=_NAVY)
    _txb(slide, title, 0.9, 0.05, 19.1, 1.3, size=13, bold=True, color=_WHITE)
    for i, lbl in enumerate(_NAV):
        dx = 21.8 + i * 0.7
        fill = _WHITE if (i+1) == active else _NAVYL
        c = _BLACK if (i+1) == active else _GRAYD
        _rect(slide, dx, 0.35, 0.55, 0.55, fill=fill)
        _txb(slide, lbl, dx, 0.35, 0.55, 0.55, size=7, bold=True, color=c,
             align=PP_ALIGN.CENTER)
    _txb(slide, subtitle, 0.9, 1.6, 23.6, 0.6, size=8, color=_GRAYT)


def _footer(slide):
    _rect(slide, 0, 13.75, 25.4, 0.5, fill=_GRAYL)
    _txb(slide, "FinSight IA  ·  Usage confidentiel  ·  Ne constitue pas un conseil en investissement MiFID II",
         0.9, 13.8, 23.6, 0.4, size=7, color=_GRAYD)


def _chapter_divider(prs, num_str, title, subtitle):
    slide = _blank(prs)
    _rect(slide, 0, 0, 25.4, 14.3, fill=_NAVY)
    _txb(slide, num_str, 1.0, 3.5, 8.0, 4.5, size=72, bold=True, color=_NAVYL)
    _txb(slide, title, 1.0, 7.0, 23.0, 2.0, size=28, bold=True, color=_WHITE)
    _rect(slide, 1.0, 9.1, 15.0, 0.05, fill=_GRAYD)
    _txb(slide, subtitle, 1.0, 9.4, 22.9, 0.8, size=11, color=_GRAYD)
    _txb(slide, "FinSight IA  ·  Usage confidentiel", 0.9, 13.75, 23.6, 0.4,
         size=7, color=_GRAYD)
    return slide


def _lecture_box(slide, title, text, y_top=9.5, height=3.8, max_chars=800):
    """Boite lecture analytique en bas de slide."""
    _rect(slide, 0.9, y_top, 23.6, height, fill=_GRAYL)
    _rect(slide, 0.9, y_top, 0.12, height, fill=_NAVY)
    _txb(slide, title, 1.2, y_top + 0.15, 22.8, 0.6,
         size=8, bold=True, color=_NAVY)
    _txb(slide, _trunc(text, max_chars), 1.2, y_top + 0.8, 22.8, height - 0.9,
         size=7.5, color=_GRAYT, wrap=True)


# ── Parseurs ────────────────────────────────────────────────────────────────────

def _parse_x(s) -> float:
    """Parse '28,4x' → 28.4"""
    try:
        return float(str(s).replace(',', '.').replace('x', '').replace('*', '').strip())
    except:
        return 0.0


def _parse_pct(s) -> float:
    """Parse '+12,3 %' → 12.3"""
    try:
        return float(str(s).replace(',', '.').replace('%', '').replace(' ', '').strip())
    except:
        return 0.0


# ── Charts ─────────────────────────────────────────────────────────────────────

def _chart_scatter(secteurs: list) -> bytes:
    """Scatter EV/EBITDA vs croissance BPA par secteur. Fallback Mg.EBITDA si EV absent."""
    fig, ax = plt.subplots(figsize=(8.5, 5.5))
    fig.patch.set_facecolor('#FFFFFF')
    ax.set_facecolor('#F8F9FA')

    ev_vals  = [_parse_x(s[4])   for s in secteurs]
    bpa_vals = [_parse_pct(s[6]) for s in secteurs]
    signals  = [s[3] for s in secteurs]
    noms     = [s[0] for s in secteurs]

    # Fallback : Mg.EBITDA quand EV/EBITDA indisponible
    use_mg = not ev_vals or not any(v > 0 for v in ev_vals)
    if use_mg:
        mg_vals = [float(s[5]) if isinstance(s[5], (int, float)) and s[5] else 0.0 for s in secteurs]
        if not any(v > 0 for v in mg_vals):
            # Aucune Donnée disponible
            ax.text(0.5, 0.5, "Données EV/EBITDA et Mg.EBITDA\nnon disponibles pour cet indice",
                    ha='center', va='center', transform=ax.transAxes, color='#999999', fontsize=10)
            ax.set_xticks([]); ax.set_yticks([])
            for sp in ax.spines.values(): sp.set_visible(False)
            buf = io.BytesIO()
            fig.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='white')
            plt.close(fig); buf.seek(0); return buf.read()
        y_vals  = mg_vals
        y_label = "Mg. EBITDA (%)"
        q_labels = ("Mg. élevée + Fort.", "Mg. faible + Fort.", "Mg. élevée + Faible", "Mg. faible + Faible")
    else:
        y_vals  = ev_vals
        y_label = "EV/EBITDA (x)"
        q_labels = ("Premium Justifie", "Value Trap ?", "Opportunite", "Risque")

    med_y   = float(np.median([v for v in y_vals if v > 0]))
    med_bpa = float(np.median(bpa_vals))

    ax.axhline(med_y,   color='#CCCCCC', linewidth=0.8, linestyle='--', alpha=0.7)
    ax.axvline(med_bpa, color='#CCCCCC', linewidth=0.8, linestyle='--', alpha=0.7)

    for i, (yv, bpa, sig, nom) in enumerate(zip(y_vals, bpa_vals, signals, noms)):
        col = _sig_hex(sig)
        ax.scatter(bpa, yv, s=100, color=col, alpha=0.85, zorder=5, edgecolors='white', linewidths=0.5)
        short = _abbrev_sector(nom, 14)
        ax.annotate(short, (bpa, yv), textcoords='offset points', xytext=(5, 4),
                    fontsize=8, color='#333333')

    kw = dict(transform=ax.transAxes, fontsize=8.5, alpha=0.7)
    ax.text(0.97, 0.93, q_labels[0], ha='right', color='#555555', **kw)
    ax.text(0.03, 0.93, q_labels[1], ha='left',  color='#555555', **kw)
    ax.text(0.97, 0.04, q_labels[2], ha='right', color='#1A7A4A', **kw)
    ax.text(0.03, 0.04, q_labels[3], ha='left',  color='#A82020', **kw)

    ax.set_xlabel("Croissance BPA (%)", fontsize=10, color='#555555')
    ax.set_ylabel(y_label, fontsize=10, color='#555555')
    ax.tick_params(labelsize=9, colors='#777777')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#DDDDDD')
    ax.spines['bottom'].set_color('#DDDDDD')
    ax.grid(True, alpha=0.3, linestyle=':')
    plt.tight_layout()

    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def _chart_score_bars(secteurs: list) -> bytes:
    """Barres horizontales de scores par secteur, tries descroissant."""
    sorted_s = sorted(secteurs, key=lambda s: s[2], reverse=True)
    noms   = [_abbrev_sector(s[0], 18) for s in sorted_s]
    scores = [s[2] for s in sorted_s]
    cols   = [_sig_hex(s[3]) for s in sorted_s]

    fig, ax = plt.subplots(figsize=(9, max(3.5, len(noms) * 0.42)))
    fig.patch.set_facecolor('#FFFFFF')
    ax.set_facecolor('#F8F9FA')

    y = np.arange(len(noms))
    bars = ax.barh(y, scores, color=cols, alpha=0.85, height=0.62, edgecolor='white', linewidth=0.5)
    for bar, val in zip(bars, scores):
        ax.text(bar.get_width() + 0.5, bar.get_y() + bar.get_height()/2,
                f"{val}", va='center', ha='left', fontsize=7, color='#333333', fontweight='bold')

    ax.axvline(50, color='#FF8C00', linewidth=1.0, linestyle='--', alpha=0.8, label='Seuil 50')
    ax.axvline(65, color='#1A7A4A', linewidth=0.8, linestyle=':', alpha=0.6)
    ax.set_yticks(y)
    ax.set_yticklabels(noms, fontsize=7.5)
    ax.set_xlim(0, 105)
    ax.set_xlabel("Score FinSight (0-100)", fontsize=8, color='#555555')
    ax.tick_params(labelsize=7, colors='#777777')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_visible(False)
    ax.grid(True, alpha=0.3, linestyle=':', axis='x')
    plt.tight_layout()

    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def _chart_ev_distribution(secteurs: list) -> bytes:
    """Barres horizontales EV/EBITDA par secteur, colore par signal.
    Fallback Mg.EBITDA si EV/EBITDA absent (ex: CAC40, DAX40)."""
    sorted_all = sorted(secteurs, key=lambda s: _parse_x(s[4]), reverse=True)
    sorted_s = [s for s in sorted_all if _parse_x(s[4]) > 0]
    if not sorted_s:
        sorted_s = sorted_all  # fallback si tous vides
    noms  = [_abbrev_sector(s[0], 18) for s in sorted_s]
    evs   = [_parse_x(s[4]) for s in sorted_s]
    cols  = [_sig_hex(s[3]) for s in sorted_s]
    filtered_evs = [v for v in evs if v > 0]

    fig, ax = plt.subplots(figsize=(8.5, max(3.5, len(noms) * 0.42)))
    fig.patch.set_facecolor('#FFFFFF')
    ax.set_facecolor('#F8F9FA')

    if not filtered_evs:
        # Fallback Mg.EBITDA : barres horizontales de marge EBITDA par secteur
        mg_data = [(s[0], s[5], s[3]) for s in secteurs
                   if isinstance(s[5], (int, float)) and s[5] > 0]
        if not mg_data:
            # Dernier recours : placeholder texte
            ax.text(0.5, 0.5, "Données non disponibles\npour cet indice",
                    ha='center', va='center', transform=ax.transAxes,
                    color='#999999', fontsize=11)
            ax.set_xticks([]); ax.set_yticks([])
            for sp in ax.spines.values(): sp.set_visible(False)
            buf = io.BytesIO()
            fig.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='white')
            plt.close(fig)
            buf.seek(0)
            return buf.read()
        plt.close(fig)
        mg_sorted = sorted(mg_data, key=lambda x: x[1], reverse=True)
        fig2, ax2 = plt.subplots(figsize=(8.5, max(3.5, len(mg_sorted) * 0.42)))
        fig2.patch.set_facecolor('#FFFFFF')
        ax2.set_facecolor('#F8F9FA')
        _noms2 = [_abbrev_sector(r[0], 18) for r in mg_sorted]
        _vals2 = [r[1] for r in mg_sorted]
        _cols2 = [_sig_hex(r[2]) for r in mg_sorted]
        _med2  = float(np.median(_vals2))
        _y2    = np.arange(len(_noms2))
        _bars2 = ax2.barh(_y2, _vals2, color=_cols2, alpha=0.85, height=0.62,
                          edgecolor='white', linewidth=0.5)
        for _bar, _val in zip(_bars2, _vals2):
            ax2.text(_bar.get_width() + 0.3, _bar.get_y() + _bar.get_height() / 2,
                     f"{_val:.1f}%", va='center', ha='left', fontsize=7, color='#333333')
        ax2.axvline(_med2, color='#1B3A6B', linewidth=1.0, linestyle='--', alpha=0.7,
                    label=f"Med. {_med2:.1f}%")
        ax2.legend(fontsize=7, framealpha=0.7, loc='lower right')
        ax2.set_yticks(_y2)
        ax2.set_yticklabels(_noms2, fontsize=7.5)
        ax2.set_xlabel("Marge EBITDA (%)", fontsize=8, color='#555555')
        ax2.tick_params(labelsize=7, colors='#777777')
        ax2.spines['top'].set_visible(False)
        ax2.spines['right'].set_visible(False)
        ax2.spines['left'].set_visible(False)
        ax2.grid(True, alpha=0.3, linestyle=':', axis='x')
        plt.tight_layout()
        _buf2 = io.BytesIO()
        fig2.savefig(_buf2, format='png', dpi=150, bbox_inches='tight',
                     facecolor='white', edgecolor='none')
        plt.close(fig2)
        _buf2.seek(0)
        return _buf2.read()

    med = float(np.median(filtered_evs))
    y = np.arange(len(noms))
    bars = ax.barh(y, evs, color=cols, alpha=0.85, height=0.62, edgecolor='white', linewidth=0.5)
    for bar, val in zip(bars, evs):
        if val > 0:
            ax.text(bar.get_width() + 0.2, bar.get_y() + bar.get_height()/2,
                    f"{val:.1f}x", va='center', ha='left', fontsize=7, color='#333333')

    ax.axvline(med, color='#1B3A6B', linewidth=1.0, linestyle='--', alpha=0.7,
               label=f"Med. {med:.1f}x")
    ax.legend(fontsize=7, framealpha=0.7, loc='lower right')
    ax.set_yticks(y)
    ax.set_yticklabels(noms, fontsize=7.5)
    ax.set_xlabel("EV/EBITDA (x)", fontsize=8, color='#555555')
    ax.tick_params(labelsize=7, colors='#777777')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_visible(False)
    ax.grid(True, alpha=0.3, linestyle=':', axis='x')
    plt.tight_layout()

    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def _chart_zone_entree(data: dict) -> bytes:
    """PE actuel vs PE médiane 10 ans par secteur — chart dotplot."""
    # Recupere les Données PE depuis top3 + secteurs si dispo
    secteurs = data.get("secteurs", [])
    top3     = data.get("top3_secteurs", [])

    # PE global de l'indice comme ancre (calcule en premier)
    try:
        _pe_global = float(str(data.get("pe_forward","17")).replace("x","").strip())
        if _pe_global < 5 or _pe_global > 50: _pe_global = 17.0
    except:
        _pe_global = 17.0
    _pe_med_global = data.get("pe_Médiane_10y", 17.0) or 17.0
    if not isinstance(_pe_med_global, (int, float)): _pe_med_global = 17.0

    # Construire mapping nom -> (pe_fwd, pe_med)
    pe_data = {}
    for s in top3:
        pe_fwd = s.get("pe_forward_raw", 0) or 0
        if pe_fwd < 5:
            # Estimation ancree sur pe_global (identique aux secteurs)
            pe_fwd = round(_pe_global + (s.get("score", 50) - 50) * 0.05, 1)
        # Cap pe_med au réaliste (évite bandes trop larges)
        pe_med_val = min(s.get("pe_Médiane_10y", _pe_med_global), _pe_global * 1.3)
        pe_data[s["nom"]] = (pe_fwd, float(pe_med_val))

    for s in secteurs:
        if s[0] not in pe_data:
            score = s[2]
            pe_est = round(_pe_global + (score - 50) * 0.05, 1)
            pe_data[s[0]] = (pe_est, float(_pe_med_global))

    if not pe_data:
        fig, ax = plt.subplots(figsize=(9, 4))
        ax.text(0.5, 0.5, "Données PE non disponibles", ha='center', va='center',
                transform=ax.transAxes)
        buf = io.BytesIO()
        fig.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='white')
        plt.close(fig)
        buf.seek(0)
        return buf.read()

    # Map signal per sector name
    sig_map = {s[0]: s[3] for s in secteurs}

    noms    = list(pe_data.keys())
    pe_fwds = [pe_data[n][0] for n in noms]
    pe_meds = [pe_data[n][1] for n in noms]
    sigs    = [sig_map.get(n, "Neutre") for n in noms]

    # figsize matches PPTX slot ratio (15.0 x 10.5 cm = 1.43:1) to avoid distortion
    _fig_h = max(9.5 * 10.5 / 15.0, len(noms) * 0.55)
    fig, ax = plt.subplots(figsize=(9.5, _fig_h))
    fig.patch.set_facecolor('#FFFFFF')
    ax.set_facecolor('#F8F9FA')

    y = np.arange(len(noms))
    # Bande "zone normale" autour de Médiane
    for i, (med, fwd, sig, nom) in enumerate(zip(pe_meds, pe_fwds, sigs, noms)):
        ax.barh(i, med * 0.3, left=med * 0.85, height=0.3,
                color='#E8ECF0', alpha=0.7, zorder=3)
        # Dot PE actuel
        col = _sig_hex(sig)
        ax.scatter(fwd, i, s=80, color=col, zorder=6, edgecolors='white', linewidths=0.5)
        # Ligne med
        ax.scatter(med, i, s=40, color='#AAAAAA', marker='|', zorder=5)

    ax.set_yticks(y)
    ax.set_yticklabels([_abbrev_sector(n, 18) for n in noms], fontsize=7.5)
    ax.set_xlabel("P/E Forward", fontsize=8, color='#555555')
    ax.tick_params(labelsize=7, colors='#777777')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_visible(False)
    ax.grid(True, alpha=0.3, linestyle=':', axis='x')
    # Limiter l'axe X au range réel des dots (pas d'espace vide a droite)
    _xmin = max(0, min(pe_fwds) - 1.5)
    _xmax = max(pe_fwds) + 1.5
    ax.set_xlim(_xmin, _xmax)

    from matplotlib.lines import Line2D
    legend_el = [
        Line2D([0],[0], marker='o', color='w', markerfacecolor='#1A7A4A', markersize=7, label='Surpondérer'),
        Line2D([0],[0], marker='o', color='w', markerfacecolor='#B06000',  markersize=7, label='Neutre'),
        Line2D([0],[0], marker='o', color='w', markerfacecolor='#A82020',  markersize=7, label='Sous-pondérer'),
        Line2D([0],[0], marker='|', color='#AAAAAA', markersize=8, label='Médiane 10Y'),
    ]
    ax.legend(handles=legend_el, fontsize=6.5, framealpha=0.7, loc='lower right')
    plt.tight_layout()

    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def _chart_sentiment_bars(sentiment_agg: dict) -> bytes:
    """Barres sentiment par secteur."""
    par_sec = sentiment_agg.get("par_secteur", [])
    if not par_sec:
        fig, ax = plt.subplots(figsize=(8, 3))
        ax.text(0.5, 0.5, "Données sentiment non disponibles", ha='center', va='center',
                transform=ax.transAxes)
        buf = io.BytesIO()
        fig.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='white')
        plt.close(fig)
        buf.seek(0)
        return buf.read()

    sorted_ps = sorted(par_sec, key=lambda x: x[1], reverse=True)
    noms   = [_abbrev_sector(p[0], 20) for p in sorted_ps]
    scores = [float(p[1]) for p in sorted_ps]
    cols   = ['#1A7A4A' if v >= 0.05 else ('#A82020' if v <= -0.05 else '#B06000') for v in scores]

    # Figsize cale sur le ratio de la zone d'affichage PPTX (13.5cm x 11.0cm = 1.23:1)
    # Layout slide 19 : graphique a gauche w=13.5cm h=11.0cm
    _disp_w_cm, _disp_h_cm = 13.5, 11.0
    _fig_w_in = 8.0
    _fig_h_in = _fig_w_in * (_disp_h_cm / _disp_w_cm)  # = 8.0 * 0.815 = 6.52 in
    fig, ax = plt.subplots(figsize=(_fig_w_in, _fig_h_in))
    fig.patch.set_facecolor('#FFFFFF')
    ax.set_facecolor('#F8F9FA')

    y = np.arange(len(noms))
    ax.barh(y, scores, color=cols, alpha=0.85, height=0.60, edgecolor='white', linewidth=0.5)
    ax.axvline(0, color='#333333', linewidth=0.8)
    ax.set_yticks(y)
    ax.set_yticklabels(noms, fontsize=8)
    ax.set_xlabel("Score composite (-0.25 a +0.25)", fontsize=8, color='#555555')
    ax.tick_params(labelsize=8, colors='#777777')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.grid(True, alpha=0.3, linestyle=':', axis='x')
    # Legende formelle — positionnee en haut a droite pour éviter chevauchement
    patches_leg = [
        mpatches.Patch(color='#1A7A4A', label='Surpondérer (> 0.05)'),
        mpatches.Patch(color='#B06000', label='Neutre'),
        mpatches.Patch(color='#A82020', label='Sous-pondérer (< -0.05)'),
    ]
    ax.legend(handles=patches_leg, fontsize=7.5, loc='upper right',
              framealpha=0.9, frameon=True, ncol=3,
              bbox_to_anchor=(1.0, 1.15))
    ax.set_title("Score composite par secteur — derive du scoring FinSight",
                 fontsize=9, color='#1B3A6B', fontweight='bold', pad=4)
    plt.tight_layout(pad=0.4)

    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def _chart_etf_perf(data: dict) -> bytes:
    """Performance 52 semaines des ETF sectoriels, indexe a 100."""
    import datetime as dt
    etf_perf = data.get("etf_perf", {})
    colors_line = _LINE_COLORS_HEX

    fig, ax = plt.subplots(figsize=(9.5, 5.5))
    fig.patch.set_facecolor('#FFFFFF')
    ax.set_facecolor('#F8F9FA')

    plotted = 0
    try:
        import yfinance as yf
        etf_list = list(etf_perf.keys())[:8]
        for i, etf in enumerate(etf_list):
            try:
                hist_df = yf.Ticker(etf).history(period='1y', interval='1wk')
                if hist_df.empty or 'Close' not in hist_df.columns:
                    continue
                hist = hist_df['Close']
                if len(hist) < 4:
                    continue
                norm = (hist / hist.iloc[0]) * 100
                label = f"{etf} — {etf_perf[etf].get('nom','')[:12]}"
                ax.plot(norm.index, norm.values,
                        color=colors_line[i % len(colors_line)],
                        linewidth=1.6, label=label, alpha=0.9)
                plotted += 1
            except Exception:
                pass
    except Exception:
        pass

    if plotted == 0:
        # Fallback simule basé sur return_1y
        np.random.seed(42)
        x = np.linspace(0, 52, 53)
        for i, (etf, info) in enumerate(list(etf_perf.items())[:8]):
            ret_ann = info.get("return_1y", 5.0) / 100.0
            ret_wk = ret_ann / 52
            np.random.seed(i * 13)
            noise = np.random.randn(53) * 1.2
            y = 100 * np.cumprod(1 + ret_wk + noise * 0.01)
            y = y / y[0] * 100
            nom = info.get("nom", etf)[:12]
            ax.plot(x, y, color=colors_line[i % len(colors_line)],
                    linewidth=1.6, label=f"{etf} — {nom}", alpha=0.9)
        ax.set_xlabel("Semaines (illustratif)", fontsize=8, color='#555555')
    else:
        ax.set_xlabel("Date", fontsize=8, color='#555555')

    ax.axhline(100, color='#CCCCCC', linewidth=1.0, linestyle='-')
    ax.set_ylabel("Performance indexee (base 100)", fontsize=8, color='#555555')
    ax.tick_params(labelsize=7, colors='#777777')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#DDDDDD')
    ax.spines['bottom'].set_color('#DDDDDD')
    ax.grid(True, alpha=0.3, linestyle=':')
    ax.legend(fontsize=6.5, framealpha=0.7, ncol=2, loc='upper left')
    plt.tight_layout()

    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def _chart_index_perf(data: dict) -> bytes:
    """Graphique performance 52 semaines de l'indice vs obligations et or."""
    import matplotlib.dates as mdates
    from datetime import datetime as _dt2
    ph = data.get("perf_history")
    if not ph or not ph.get("dates") or not ph.get("indice"):
        return _make_empty_chart()
    try:
        dates   = [_dt2.strptime(d[:10], "%Y-%m-%d") for d in ph["dates"]]
        i_perf  = ph.get("indice", [])
        b_perf  = ph.get("bonds", [])
        g_perf  = ph.get("gold", [])
        sp_perf = ph.get("sp500", [])
        if not dates or not i_perf:
            return _make_empty_chart()

        fig, ax = plt.subplots(figsize=(11, 4.5))
        fig.patch.set_facecolor('#FFFFFF')
        ax.set_facecolor('#F8F9FA')

        _indice_name_perf = ph.get("indice_name", "Indice")
        ax.plot(dates, i_perf, color='#1B3A6B', lw=2.5,
                label=_indice_name_perf, zorder=3)
        if len(sp_perf) == len(dates):
            # Éviter doublon legende si l'indice est déjà le S&P 500
            _sp_label = "S&P 500 (ref.)" if "S&P 500" in _indice_name_perf else "S&P 500"
            ax.plot(dates, sp_perf, color='#555555', lw=1.5, linestyle='-.',
                    label=_sp_label, alpha=0.75, zorder=2)
        if len(b_perf) == len(dates):
            ax.plot(dates, b_perf, color='#A82020', lw=1.5, linestyle='--',
                    label="Obligations (^TNX)", alpha=0.8, zorder=2)
        if len(g_perf) == len(dates):
            ax.plot(dates, g_perf, color='#B06000', lw=1.5, linestyle=':',
                    label="Or (GC=F)", alpha=0.8, zorder=2)

        ax.axhline(100, color='#CCCCCC', lw=1.0, linestyle='-', zorder=1)
        ax.xaxis.set_major_formatter(mdates.DateFormatter('%b %y'))
        ax.xaxis.set_major_locator(mdates.MonthLocator(interval=2))
        plt.xticks(rotation=25, fontsize=8)
        ax.tick_params(axis='y', labelsize=8)
        ax.set_ylabel("Base 100", fontsize=8, color='#555555')
        ax.legend(fontsize=8, loc='upper left', framealpha=0.8)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.spines['left'].set_color('#DDDDDD')
        ax.spines['bottom'].set_color('#DDDDDD')
        ax.grid(axis='y', alpha=0.3, linestyle=':')

        if i_perf:
            final_pct = i_perf[-1] - 100
            _col_ann  = '#1A7A4A' if final_pct >= 0 else '#A82020'
            ax.annotate(
                f"{final_pct:+.1f}%",
                xy=(dates[-1], i_perf[-1]),
                xytext=(-5, 6), textcoords='offset points',
                fontsize=9, color=_col_ann, fontweight='bold',
            )

        plt.tight_layout()
        buf = io.BytesIO()
        fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                    facecolor='white', edgecolor='none')
        plt.close(fig)
        buf.seek(0)
        return buf.read()
    except Exception as _e:
        log.warning("_chart_index_perf: %s", _e)
        return _make_empty_chart()


# ── Slides 1-20 ─────────────────────────────────────────────────────────────────

def _s01_cover(prs, D):
    slide = _blank(prs)

    # ── Bandeau header navy (identique pptx_writer société) ───────────────────
    _rect(slide, 0, 0, 25.4, 3.81, fill=_NAVY)
    _txb(slide, "FinSight IA", 1.27, 0.46, 22.86, 0.71,
         size=10, bold=False, color=RGBColor(0x88, 0x99, 0xBB),
         align=PP_ALIGN.CENTER)
    _rect(slide, 11.05, 1.32, 3.3, 0.05, fill=_WHITE)
    _txb(slide, "Pitchbook  \u00b7  Analyse d'Indice", 1.27, 1.52, 22.86, 0.81,
         size=11, color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER)

    # ── Corps blanc ───────────────────────────────────────────────────────────
    # Nom indice (large, centre)
    _txb(slide, D.get("indice", ""), 1.27, 4.2, 22.86, 1.6,
         size=32, bold=True, color=_NAVY, align=PP_ALIGN.CENTER)

    # Tagline code + secteurs + sociétés
    code = D.get("code", "")
    nb_s = D.get("nb_secteurs", 0)
    nb_c = D.get("nb_societes", 0)
    parts = [p for p in [
        code,
        (f"{nb_s} secteurs analys\u00e9s" if nb_s else ""),
        (f"{nb_c} soci\u00e9t\u00e9s couvertes" if nb_c else ""),
    ] if p]
    _txb(slide, "  \u00b7  ".join(parts), 1.27, 5.9, 22.86, 0.71,
         size=11, color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)

    # Badge signal (pattern recommandation société)
    _SIG_NORM_COV = {
        "Surpondérer":  "SURPOND\u00c9RER", "surpond\u00e9rer": "SURPOND\u00c9RER",
        "Sous-pondérer":"SOUS-POND\u00c9RER","sous-pond\u00e9rer":"SOUS-POND\u00c9RER",
        "neutre":       "NEUTRE",
    }
    sig = D.get("signal_global", "Neutre")
    sig_col   = _sig_color(sig)
    sig_light = _sig_light(sig)
    sig_txt   = _SIG_NORM_COV.get(sig.lower(), sig.upper())
    conv      = D.get("conviction_pct", 50)
    erp_cov   = D.get("erp") or None
    erp_part  = f"  \u00b7  ERP {erp_cov}" if erp_cov else ""
    _rect(slide, 1.27, 6.8, 22.86, 1.32, fill=sig_light)
    _rect(slide, 1.27, 6.8, 0.13,  1.32, fill=sig_col)
    _txb(slide, f"\u25cf {sig_txt}  \u00b7  Conviction {conv}\u00a0%{erp_part}",
         1.60, 6.85, 22.40, 1.22,
         size=9, bold=True, color=sig_col, align=PP_ALIGN.CENTER)

    # KPI boxes (fond gris clair, texte navy — lisibles sur blanc)
    secteurs_cov = D.get("secteurs", [])
    score_med = D.get("score_Médian")
    if not score_med and secteurs_cov:
        scores_l = [s[2] for s in secteurs_cov if len(s) > 2
                    and isinstance(s[2], (int, float))]
        score_med = round(sum(scores_l) / len(scores_l)) if scores_l else None
    sm_str   = f"{score_med}/100" if isinstance(score_med, (int, float)) else "\u2014"
    nb_surp  = sum(1 for s in secteurs_cov if len(s) > 3 and "surp"   in str(s[3]).lower())
    nb_neut  = sum(1 for s in secteurs_cov if len(s) > 3 and "neutre" in str(s[3]).lower())
    nb_sousp = sum(1 for s in secteurs_cov if len(s) > 3 and "sous"   in str(s[3]).lower())
    _cours_cov = D.get("cours", "\u2014")
    _ytd_cov   = D.get("variation_ytd", "\u2014")
    metrics = [
        ("Cours",              _cours_cov),
        ("Variation YTD",      _ytd_cov),
        ("Surpond\u00e9rer",   f"{nb_surp} sect."),
        ("Neutre",             f"{nb_neut} sect."),
        ("Sous-pond\u00e9rer", f"{nb_sousp} sect."),
        ("Score m\u00e9dian",  sm_str),
    ]
    box_w = 3.7
    step  = 22.86 / 6
    for i, (lbl, val) in enumerate(metrics):
        xpos = 1.27 + i * step
        _rect(slide, xpos, 8.5, box_w, 1.6, fill=_GRAYL)
        _txb(slide, lbl, xpos + 0.15, 8.6,  box_w - 0.3, 0.5, size=7, color=_GRAYT)
        _txb(slide, val, xpos + 0.15, 9.05, box_w - 0.3, 0.9, size=11, bold=True, color=_NAVY)

    # Ligne de separation + crédits (identique société)
    _rect(slide, 1.02, 11.8, 23.37, 0.03, fill=_GRAYD)
    date_str = D.get("date_analyse") or _fr_date()
    _txb(slide, "Rapport confidentiel", 1.02, 12.05, 11.43, 0.56,
         size=8, color=_GRAYT)
    _txb(slide, date_str, 12.95, 12.05, 11.43, 0.56,
         size=8, color=_GRAYT, align=PP_ALIGN.RIGHT)

    return slide


def _s02_exec_summary(prs, D):
    slide = _blank(prs)
    indice = D.get("indice","")
    code   = D.get("code","")
    nb_s   = D.get("nb_secteurs",0)
    nb_c   = D.get("nb_societes",0)
    _header(slide, "Executive Summary",
            f"{indice} ({code})  ·  {nb_s} secteurs  ·  {nb_c} sociétés  ·  Horizon 12 mois",
            active=1)

    # Signal badge
    _SIG_NORM = {
        "Surpondérer": "SURPONDÉRER", "surpondérer": "SURPONDÉRER",
        "Sous-pondérer": "SOUS-PONDÉRER", "sous-pondérer": "SOUS-PONDÉRER",
        "neutre": "NEUTRE",
    }
    sig     = D.get("signal_global","Neutre")
    sig_col = _sig_color(sig)
    sig_txt = _SIG_NORM.get(sig.lower(), sig.upper())
    _rect(slide, 0.9, 2.3, 4.8, 0.75, fill=sig_col)
    _txb(slide, sig_txt, 1.1, 2.35, 4.4, 0.55, size=10, bold=True, color=_WHITE)

    # Metriques ligne
    _erp_s02 = D.get("erp", "—")
    _esig_s02 = D.get("erp_signal", "")
    _erp_display_s02 = f"ERP : {_erp_s02}" + (f" ({_esig_s02})" if _esig_s02 else "")
    mets = [
        f"Cours : {D.get('cours','—')}",
        f"YTD : {D.get('variation_ytd','—')}",
        f"P/E Forward : {D.get('pe_forward','—')}",
        _erp_display_s02,
        f"Conviction {D.get('conviction_pct',50)} %",
    ]
    _txb(slide, "  ·  ".join(mets), 0.9, 3.15, 23.6, 0.55, size=8, color=_GRAYT)

    # Régime de Marché — badges individuels
    _mac = D.get("macro") or {}
    _reg = _mac.get("regime", "")
    if _reg and _reg != "Inconnu":
        _RSIG_COL = {"Bull": _BUY, "Bear": _SELL, "Volatile": _HOLD, "Transition": _HOLD}
        _r_col = _RSIG_COL.get(_reg, _NAVYL)
        _rec_6 = _mac.get("recession_prob_6m")
        _vx    = _mac.get("vix")
        _sp_s  = _mac.get("yield_spread_10y_3m")
        _badges = [(f"Régime {_reg}", _r_col, _WHITE)]
        if _vx:    _badges.append((f"VIX {_vx:.0f}", _GRAYL, _NAVY))
        if _sp_s is not None: _badges.append((f"Spread {_sp_s:+.1f}%", _GRAYL, _NAVY))
        if _rec_6 is not None: _badges.append((f"Rec. 6M {_rec_6}%", _GRAYL, _NAVY))
        _bx = 0.9
        for _blbl, _bfill, _btxt in _badges:
            _bw = max(2.2, len(_blbl) * 0.13 + 0.5)
            _bw = min(_bw, 5.0)
            _rect(slide, _bx, 3.8, _bw, 0.45, fill=_bfill)
            _txb(slide, _blbl, _bx + 0.1, 3.8, _bw - 0.15, 0.45,
                 size=7.5, bold=(_bfill == _r_col), color=_btxt, align=PP_ALIGN.CENTER, vcenter=True)
            _bx += _bw + 0.15

    # Catalyseurs macro (y fixe suffisant après regime_v strip eventuel)
    _rect(slide, 0.9, 4.4, 23.6, 0.5, fill=_NAVY)
    _txb(slide, "CATALYSEURS MACRO", 0.9, 4.4, 23.6, 0.5,
         size=7.5, bold=True, color=_WHITE, align=PP_ALIGN.CENTER, vcenter=True)

    cats = D.get("catalyseurs", [])
    for j, cat in enumerate(cats[:3]):
        yy = 5.05 + j * 0.95
        _rect(slide, 0.9, yy, 0.08, 0.9, fill=_BUY)
        nom  = cat[0][:50] if isinstance(cat, (list,tuple)) else str(cat)[:50]
        desc = cat[1][:120] if isinstance(cat,(list,tuple)) and len(cat) > 1 else ""
        _txb(slide, nom,  1.15, yy,       22.3, 0.45, size=8, bold=True, color=_NAVY)
        _txb(slide, desc, 1.15, yy + 0.45, 22.3, 0.55, size=7.5, color=_GRAYT, wrap=True)

    # Synthèse signal en bas — enrichie avec metriques calculés
    _texte_raw = D.get("texte_signal", "")
    _cours_s02 = D.get("cours", "—")
    _ytd_s02   = D.get("variation_ytd", "—")
    _pe_f_s02  = D.get("pe_forward", "—")
    _pe_m_s02  = D.get("pe_Médiane_10y", "—")
    _erp_s02b  = D.get("erp", "—")
    _erp_sig_s02 = D.get("erp_signal", "")
    _scr_s02   = D.get("score_Médian")
    if not isinstance(_scr_s02, (int, float)):
        _scores_l02 = [s[2] for s in D.get("secteurs", []) if len(s) > 2 and isinstance(s[2], (int, float))]
        _scr_s02 = round(sum(_scores_l02) / len(_scores_l02)) if _scores_l02 else "—"
    if not isinstance(_scr_s02, (int, float)):
        _scr_s02 = "—"
    _conv_s02  = D.get("conviction_pct", 50)
    _surp_s02  = D.get("surp_noms", "") or ""
    _sous_s02  = D.get("sous_noms", "") or ""
    _pm_s02    = f"{_pe_m_s02}x" if isinstance(_pe_m_s02, (int, float)) else str(_pe_m_s02)
    _prime_s02 = D.get("prime_decote", "")
    _prime_lbl = "prime de valorisation significative" if isinstance(_prime_s02, str) and "+" in str(_prime_s02) else "valorisation en ligne avec l'historique"
    _erp_lbl   = " (prime insuffisanté — prudence sur les entrées)" if _erp_sig_s02 in ("Tendu", "Comprime") else (" (adéquat)" if _erp_sig_s02 else "")
    _suppl_s02 = (
        f"Le {D.get('indice', '')} affiche un cours de {_cours_s02} (YTD : {_ytd_s02}). "
        f"P/E Forward {_pe_f_s02} vs médiane historique 10 ans {_pm_s02} — {_prime_lbl}. "
        f"ERP Damodaran {_erp_s02b}{_erp_lbl}. "
        f"Score composite moyen : {_scr_s02}/100 — conviction {_conv_s02} %. "
        f"Secteurs a Surpondérer : {_surp_s02 or 'aucun'}. "
        f"Secteurs a éviter : {_sous_s02 or 'aucun'}."
    )
    # Ajouter horizon seulement si pas déjà dans le texte IA
    if "Horizon" not in _texte_raw:
        _suppl_s02 += " Horizon d'allocation recommande : 12 mois."
    texte = _trunc((_texte_raw + "  " + _suppl_s02).strip() if _texte_raw else _suppl_s02, 700)
    _y_synth = 8.3
    _h_synth = 13.3 - _y_synth
    _rect(slide, 0.9, _y_synth, 23.6, _h_synth, fill=_GRAYL)
    _rect(slide, 0.9, _y_synth, 0.12, _h_synth, fill=_NAVY)
    _txb(slide, "SYNTHÈSE DU SIGNAL", 1.2, _y_synth + 0.1, 22.8, 0.5, size=8, bold=True, color=_NAVY)
    _txb(slide, texte, 1.2, _y_synth + 0.7, 22.8, _h_synth - 0.8, size=8, color=_GRAYT, wrap=True)

    _footer(slide)
    return slide


def _s03_sommaire(prs, D):
    slide = _blank(prs)
    indice = D.get("indice","")
    code   = D.get("code","")
    _header(slide, "Sommaire",
            f"{indice}  ·  {code}  ·  Structure de l'analyse macro institutionnelle",
            active=1)

    chapitres = [
        ("01", "Synthèse Macro & Signal Global",
         "Description indice · P/E Forward · ERP Damodaran · Catalyseurs & risques · Cycle économique",
         "p. 5-7"),
        ("02", "Cartographie des Secteurs",
         f"{D.get('nb_secteurs',11)} secteurs GICS · Scores FinSight · Scatter EV/EBITDA vs croissance · Décomposition scores",
         "p. 9-11"),
        ("03", "Top 3 Secteurs Recommandés",
         "Synthèse signal · Sociétés représentatives · Distribution valorisations · Zone d'entrée",
         "p. 13-15"),
        ("04", "Risques, Rotation & Sentiment",
         "Scénarios alternatifs · Rotation sectorielle · FinBERT · Performance ETF",
         "p. 17-20"),
    ]

    for i, (num, titre, desc, pages) in enumerate(chapitres):
        yy = 2.5 + i * 2.75
        _rect(slide, 0.9, yy, 23.6, 2.4, fill=_GRAYL)
        _rect(slide, 0.9, yy, 1.2, 2.4, fill=_NAVY)
        _txb(slide, num, 0.9, yy + 0.7, 1.2, 1.0, size=14, bold=True,
             color=_WHITE, align=PP_ALIGN.CENTER)
        _txb(slide, titre, 2.3, yy + 0.15, 18.5, 0.75, size=10, bold=True, color=_NAVY)
        _txb(slide, desc,  2.3, yy + 0.9,  18.5, 1.1,  size=8,  color=_GRAYT, wrap=True)
        _txb(slide, pages, 21.5, yy + 0.75, 2.7, 0.6,  size=9, bold=True,
             color=_NAVY, align=PP_ALIGN.RIGHT)

    _footer(slide)
    return slide


def _s05_description(prs, D):
    slide = _blank(prs)
    indice = D.get("indice","")
    code   = D.get("code","")
    _header(slide, f"{indice}  —  Caractéristiques & Valorisation",
            f"{code}  ·  Analyse institutionnelle FinSight  ·  Données yfinance / FMP",
            active=1)

    # Texte description gauche — wrapp\u00e9 dans une box analytique avec titre
    # (audit Baptiste 2026-04-14 : slide 5 texte sans box LLM)
    desc = D.get("texte_description","")[:1400]
    _rect(slide, 0.9, 2.3, 13.4, 7.1, fill=_GRAYL)
    _rect(slide, 0.9, 2.3, 0.12, 7.1, fill=_NAVY)
    _txb(slide, "DESCRIPTION DE L'INDICE & CONTEXTE", 1.2, 2.4, 13.2, 0.45,
         size=8.5, bold=True, color=_NAVY)
    _txb(slide, desc, 1.2, 3.0, 13.1, 6.3, size=8.5, color=_GRAYT, wrap=True)

    # Tableau metriques droite
    cours     = D.get("cours","—")
    ytd       = D.get("variation_ytd","—")
    pe_fwd    = D.get("pe_forward","—")
    pe_med    = D.get("pe_Médiane_10y","—")
    erp       = D.get("erp","—")
    bpa       = D.get("bpa_growth","—")
    nb_s      = D.get("nb_secteurs",11)
    nb_c      = D.get("nb_societes","—")

    rows = [
        ["MÉTRIQUE",              "VALEUR"],
        ["Cours actuel",          cours],
        ["Variation YTD",         ytd],
        ["P/E Forward",           pe_fwd],
        ["P/E Médiane 10 ans",    f"{pe_med}x" if isinstance(pe_med, (int,float)) else str(pe_med)],
        ["ERP (Damodaran)",       erp],
        ["Croissance BPA",        bpa],
        ["Secteurs analysés",     str(nb_s)],
        ["Sociétés couvertes",    str(nb_c)],
    ]
    _add_table(slide, rows, 15.1, 2.3, 9.4, 7.0,
               col_widths=[6, 3], font_size=8, header_size=8, alt_fill=_GRAYL)

    # Lecture analytique — contexte sectoriel et valorisation
    secteurs_d = D.get("secteurs",[])
    surp = [s[0] for s in secteurs_d if "Surp" in str(s[3])]
    sous = [s[0] for s in secteurs_d if "Sous" in str(s[3])]
    nb_n = len([s for s in secteurs_d if "Surp" not in str(s[3]) and "Sous" not in str(s[3])])
    _surp_str = "  ·  ".join([_abbrev_sector(x,18) for x in surp[:3]]) or "aucun"
    _sous_str = "  ·  ".join([_abbrev_sector(x,18) for x in sous[:3]]) or "aucun"
    _pe_f5  = D.get("pe_forward","—")
    _pe_m5  = D.get("pe_Médiane_10y","—")
    _erp5   = D.get("erp","—")
    _erp_s5 = D.get("erp_signal","")
    _pm5    = f"{_pe_m5}x" if isinstance(_pe_m5,(int,float)) else str(_pe_m5)
    _prime5 = D.get("prime_decote","")
    _sig5   = D.get("signal_global","Neutre")
    _scr5   = D.get("score_Médian")
    if not isinstance(_scr5, (int, float)):
        _scores_l5 = [s[2] for s in D.get("secteurs", []) if len(s) > 2 and isinstance(s[2], (int, float))]
        _scr5 = round(sum(_scores_l5) / len(_scores_l5)) if _scores_l5 else "—"
    if not isinstance(_scr5, (int, float)):
        _scr5 = "—"
    lec_txt = (
        f"Signal global : {_sig5} (score composite {_scr5}/100). "
        f"P/E Forward {_pe_f5} vs médiane historique {_pm5} — "
        f"{'prime de valorisation, entrees a calibrer' if isinstance(_prime5,str) and '+' in str(_prime5) else 'valorisation proche des normes historiques'}. "
        f"ERP Damodaran : {_erp5}{' (' + _erp_s5 + ')' if _erp_s5 else ''}. "
        f"Secteurs a Surpondérer : {_surp_str}. "
        f"Secteurs a Sous-pondérer : {_sous_str}. "
        f"{nb_n} secteur(s) en position Neutre. "
        f"Horizon d'allocation recommande : 12 mois."
    )
    _lecture_box(slide, "Lecture — Signal global & positionnément sectoriel",
                 _trunc(lec_txt, 700), y_top=9.8, height=3.5)

    _footer(slide)
    return slide


def _s06_valorisation(prs, D):
    slide = _blank(prs)
    indice = D.get("indice","")
    _header(slide, "Valorisation Macro Top-Down",
            f"{indice}  ·  BPA forward x P/E cible historique  ·  ERP Damodaran  ·  Signaux FRED",
            active=1)

    # Table valorisation
    pe_fwd    = D.get("pe_forward","—")
    pe_med    = D.get("pe_Médiane_10y","—")
    prime     = D.get("prime_decote","—")
    erp       = D.get("erp","—")
    bpa       = D.get("bpa_growth","—")
    rows = [
        ["INDICATEUR",         "VALEUR",     "vs HISTORIQUE",    "INTERPRÉTATION"],
        ["P/E Forward",        pe_fwd,       prime,              "Prime vs médiane 10Y"],
        ["P/E Médiane 10 ans", f"{pe_med}x" if isinstance(pe_med,(int,float)) else str(pe_med),
                               "Reference",  "Niveau historique normalisé"],
        ["ERP (Damodaran)",    erp,          "Correct",          "Rendement excess equites"],
        ["Croissance BPA",     bpa,          "+",                "Consensu analystes 12M"],
        ["Prime/Décote",       prime,        "Surévalué si > 20 %","Signal d'alerte valorisation"],
    ]
    tbl = _add_table(slide, rows, 0.9, 2.3, 23.6, 5.0,
               col_widths=[5, 3, 3, 12.6], font_size=8, header_size=8, alt_fill=_GRAYL)

    # Colorer prime
    for r in range(2, len(rows)):
        prime_val = rows[r][2]
        if isinstance(prime_val, str) and '+' in prime_val:
            _color_cell(tbl, r, 2, _HOLD_L, _BLACK)

    # Lecture analytique — toujours regeneree pour la qualité
    _pe_m6 = D.get("pe_Médiane_10y","—")
    _pe_f6 = D.get("pe_forward","—")
    _erp6  = D.get("erp","—")
    _sig6  = D.get("signal_global","Neutre")
    _scr6  = D.get("score_Médian")
    if not isinstance(_scr6, (int, float)):
        _scores_l6 = [s[2] for s in D.get("secteurs", []) if len(s) > 2 and isinstance(s[2], (int, float))]
        _scr6 = round(sum(_scores_l6) / len(_scores_l6)) if _scores_l6 else "—"
    if not isinstance(_scr6, (int, float)):
        _scr6 = "—"
    if True:  # toujours construire la lecture analytique enrichie
        _pm_str = f"{_pe_m6}x" if isinstance(_pe_m6,(int,float)) else str(_pe_m6)
        _prime6  = D.get("prime_decote","")
        _erp_s6  = D.get("erp_signal","")
        _ytd6    = D.get("variation_ytd","—")
        _cours6  = D.get("cours","—")
        _conv6   = D.get("conviction_pct",50)
        _bpa6    = D.get("bpa_growth","—")
        _nb_surp6 = sum(1 for s in D.get("secteurs",[]) if len(s)>3 and "Surp" in str(s[3]))
        _nb_sous6 = sum(1 for s in D.get("secteurs",[]) if len(s)>3 and "Sous" in str(s[3]))
        # Diagnostic valorisation
        try:
            _pe_num = float(str(_pe_f6).replace("x","").replace(",",".").strip())
            _pm_num = float(str(_pe_m6).replace("x","").replace(",",".").strip()) if isinstance(_pe_m6,(int,float)) else 17.0
            _prime_pct = (_pe_num - _pm_num) / _pm_num * 100
            if _prime_pct > 25:
                _val_diag = (f"Le P/E Forward de {_pe_f6} se situe {_prime6} vs la Médiane historique "
                             f"({_pm_str}) — valorisation tendue qui exige une croissance BPA soutenue "
                             f"pour justifier le multiple. Le moindre miss sur les BPA NTM serait "
                             f"susceptible de comprimer le multiple de facon significative.")
            elif _prime_pct > 5:
                _val_diag = (f"Le P/E Forward de {_pe_f6} traduit une prime modeste de {_prime6} "
                             f"vs la Médiane historique ({_pm_str}). La valorisation reste defensible "
                             f"si les BPA NTM sont livres comme attendu par le consensus.")
            elif _prime_pct < -10:
                _val_diag = (f"Le P/E Forward de {_pe_f6} offre une décote de {_prime6} vs l'historique "
                             f"({_pm_str}) — opportunité si les fondamentaux restent solides. "
                             f"Un re-rating est possible en cas de révision haussière des BPA ou de pivot Fed.")
            else:
                _val_diag = (f"Le P/E Forward de {_pe_f6} est en ligne avec la Médiane historique "
                             f"({_pm_str}) — valorisation neutre, le Marché pricant un scénario central "
                             f"sans excès dans un sens ni dans l'autre.")
        except Exception:
            _val_diag = f"P/E Forward {_pe_f6} vs Médiane historique {_pm_str}."
        # Lecture ERP
        if _erp_s6 in ("Tendu","Comprime"):
            _erp_impl = (f"L'ERP Damodaran de {_erp6} signale une prime de risque comprimee — "
                         "les actions sont peu renunerees vs les taux. Dans ce contexte, la sélectivité "
                         "sectorielle prime : se concentrer sur les secteurs a forte visibilité BPA "
                         "et pricing power, éviter les dossiers a multiples etires.")
        elif _erp_s6 in ("Favorable","Attractif"):
            _erp_impl = (f"L'ERP Damodaran de {_erp6} signale une prime de risque attractive — "
                         "les equites offrent un surplus de rendement justifiant un positionnément "
                         "actif. Le contexte macro soutient les surponderations sectorielles ciblees.")
        else:
            _erp_impl = (f"L'ERP Damodaran de {_erp6} reste dans la zone de reference. "
                         "Ni comprime ni attractif — privilegier la rotation sectorielle et "
                         "le stock-picking plutot qu'une exposition beta pure.")
        # Croissance implicite requise pour justifier le P/E actuel vs historique
        try:
            _pe_num_impl = float(str(_pe_f6).replace("x","").replace(",",".").strip())
            _pm_num_impl = float(str(_pe_m6).replace("x","").replace(",",".").strip()) if isinstance(_pe_m6,(int,float)) else 17.0
            _pct_impl = (_pe_num_impl / _pm_num_impl - 1) * 100
            if abs(_pct_impl) > 5:
                _miss_consequence = ("Entraînerait une contraction de multiple significative"
                                     if _pct_impl > 0
                                     else "serait partiellement absorbe par la décote existante")
                _impl_txt = (
                    f"A {_pe_f6}, le Marché price implicitement une croissance BPA supérieure "
                    f"de {abs(_pct_impl):.0f}% au scénario central — tout ralentissement "
                    f"{_miss_consequence}. "
                )
            else:
                _impl_txt = ""
        except Exception:
            _impl_txt = ""
        # Implication allocation
        try:
            _pe_n2 = float(str(_pe_f6).replace("x","").replace(",",".").strip())
            if _pe_n2 > 22:
                _alloc_impl = (
                    f"Implication allocation : favoriser les secteurs a forte visibilité BPA "
                    f"et pricing power (cf. slide 11 — top scoreurs). Réduire l'exposition "
                    f"aux cycliques purs dont la valorisation intègre déjà un scénario optimiste."
                )
            elif _pe_n2 < 14:
                _alloc_impl = (
                    f"Implication allocation : valorisation deprimee offre un point d'entree "
                    f"favorable. Renforcer les secteurs a score FinSight élevé (>65) — "
                    f"la décote compense le risque de normalisation des multiples."
                )
            else:
                _alloc_impl = (
                    f"Implication allocation : valorisation équilibrée — privilegier la "
                    f"sélectivité sectorielle (slides 11-12) et le momentum BPA vs momentum prix. "
                    f"Éviter les secteurs en declin de visibilité BPA malgre des multiples apparemment raisonnables."
                )
        except Exception:
            _alloc_impl = ""
        texte_val = (
            f"Cours {_cours6} (YTD {_ytd6}) — Croissance BPA forward {_bpa6}. "
            f"{_val_diag} "
            f"{_impl_txt}"
            f"{_erp_impl} "
            f"{_alloc_impl} "
            f"Score composite FinSight : {_scr6}/100 — signal {_sig6} (conviction {_conv6}%). "
            f"{_nb_surp6} secteur(s) en Surpondérer, {_nb_sous6} en Sous-pondérer."
        )
    _val_title = f"Score FinSight {_scr6}/100 — signal {_sig6}, {_nb_surp6} secteurs Surpondérer / {_nb_sous6} Sous-pondérer"
    _lecture_box(slide, _val_title, _trunc(texte_val, 1100), y_top=7.8, height=5.55)

    _footer(slide)
    return slide


def _s07_cycle(prs, D):
    slide = _blank(prs)
    indice = D.get("indice","")
    phase  = D.get("phase_cycle","Expansion avancée")
    _header(slide, "Positionnement dans le Cycle Économique",
            f"{indice}  ·  Phase actuelle : {phase}  ·  Signaux FRED  ·  Modèle 4 phases",
            active=1)

    # Carte phase cycle gauche
    _rect(slide, 0.9, 2.3, 8.2, 6.5, fill=_NAVY)
    _rect(slide, 0.9, 2.3, 8.2, 0.7, fill=_NAVYL)
    _txb(slide, "PHASE ACTUELLE", 1.1, 2.35, 7.8, 0.6,
         size=8, bold=True, color=_WHITE)
    phase_parts = phase.split()
    if len(phase_parts) >= 2:
        _txb(slide, " ".join(phase_parts[:-1]), 1.1, 3.2, 7.8, 1.0,
             size=18, bold=True, color=_WHITE)
        _txb(slide, phase_parts[-1], 1.1, 4.2, 7.8, 0.85,
             size=13, bold=False, color=_GRAYD)
    else:
        _txb(slide, phase, 1.1, 3.2, 7.8, 1.5, size=16, bold=True, color=_WHITE)

    detail = D.get("cycle_detail","PIB positif · Taux restrictifs\nMarges sous pression")
    _txb(slide, detail, 1.1, 5.3, 7.8, 2.0, size=8, color=_GRAYD, wrap=True)

    # Première phrase du texte_cycle seulement pour éviter le débordement
    texte_cycle_full = D.get("texte_cycle","")
    texte_cycle = _trunc(texte_cycle_full.split(".")[0] + "." if texte_cycle_full else "", 120)
    _txb(slide, texte_cycle, 1.1, 7.4, 7.8, 1.3, size=7.5, color=_GRAYD, wrap=True)

    # Tableau FRED signals droite
    fred = D.get("fred_signals", [])
    fred_rows = [["INDICATEUR", "VALEUR", "INTERPRÉTATION"]]
    for sig in fred[:7]:
        if isinstance(sig, (list,tuple)) and len(sig) >= 3:
            fred_rows.append([sig[0], str(sig[1]), str(sig[2])[:40]])
    if len(fred_rows) == 1:
        fred_rows.append(["PMI Composite", "51,2", "Expansion modérée"])
        fred_rows.append(["Courbe 10Y-2Y", "-0,12 %", "Légèrement inversée"])
        fred_rows.append(["ISM Manuf.", "49,8", "Sous expansion"])
        fred_rows.append(["Chômage", "3,9 %", "Marché solide"])
        fred_rows.append(["CPI YoY", "3,1 %", "Au-dessus cible"])

    tbl = _add_table(slide, fred_rows, 9.5, 2.3, 15.0, len(fred_rows) * 0.85 + 0.1,
               col_widths=[5.0, 2.5, 7.5], font_size=8, header_size=8, alt_fill=_GRAYL)

    # Allocation recommandée
    alloc = _trunc(D.get("texte_cycle",""), 420)
    _rect(slide, 9.5, 9.3, 15.0, 2.9, fill=_GRAYL)
    _rect(slide, 9.5, 9.3, 0.1, 2.9, fill=_BUY)
    _txb(slide, "Allocation recommandée selon le positionnément de cycle",
         9.8, 9.4, 14.5, 0.6, size=8, bold=True, color=_NAVY)
    _txb(slide, alloc, 9.8, 10.0, 14.5, 2.0, size=7.5, color=_GRAYT, wrap=True)

    _footer(slide)
    return slide


def _s09_cartographie(prs, D):
    slide = _blank(prs)
    indice  = D.get("indice","")
    nb_s    = D.get("nb_secteurs",11)
    nb_c    = D.get("nb_societes","")
    _header(slide, "Cartographie des Secteurs",
            f"{nb_s} secteurs GICS  ·  {nb_c} sociétés  ·  Tri par score FinSight décroissant",
            active=2)

    secteurs = D.get("secteurs", [])
    sorted_s = sorted(secteurs, key=lambda s: s[2], reverse=True)

    _SIG_LABEL = {
        "Surpondérer": "Surpondérer", "surpondérer": "Surpondérer",
        "Sous-pondérer": "Sous-pondérer", "sous-pondérer": "Sous-pondérer",
        "neutre": "Neutre",
    }
    # Table unifiée : Score + Signal + Métriques opérationnelles + P/Book + Div.Yield
    # EV/EBITDA est déjà sur slide 10 scatter — on le remplace par P/B + DivYield ici
    pb_by_sector = D.get("pb_by_sector", {}) or {}
    dy_by_sector = D.get("dy_by_sector", {}) or {}
    # Alias : noms secteurs yfinance EU vs noms ETF SPDR (ex: "Healthcare" -> "Health Care")
    _SECT_ALIAS_S09 = {
        "Financial Services": "Financials",
        "Healthcare":         "Health Care",
        "Basic Materials":    "Materials",
        "Consumer Défensive": "Consumer Staples",
        "Consumer Défensive": "Consumer Staples",
        "Consumer Cyclical":  "Consumer Discretionary",
    }
    rows = [["Rg", "Secteur", "Score", "Signal", "Mg.EBITDA", "Croiss.", "P/Book", "Div.Yield"]]
    for rang, s in enumerate(sorted_s, 1):
        raw_sig = str(s[3])[:15]
        norm_sig = _SIG_LABEL.get(raw_sig.strip().lower(), raw_sig)
        _sect_key = s[0]
        _sect_alias = _SECT_ALIAS_S09.get(_sect_key, _sect_key)
        _pb = pb_by_sector.get(_sect_key) if pb_by_sector.get(_sect_key) is not None else pb_by_sector.get(_sect_alias)
        _dy = dy_by_sector.get(_sect_key) if dy_by_sector.get(_sect_key) is not None else dy_by_sector.get(_sect_alias)
        rows.append([
            str(rang),
            _abbrev_sector(s[0], 22),
            str(s[2]),
            norm_sig,
            f"{s[5]:.1f}%" if isinstance(s[5], (int, float)) and s[5] else "—",
            str(s[6]) if len(s) > 6 else "—",
            f"{float(_pb):.1f}x" if _pb is not None else "—",
            f"{float(_dy):.1f}%" if _dy is not None else "—",
        ])
    tbl = _add_table(slide, rows, 0.9, 2.3, 23.6,
                     min(7.5, 0.65 + len(rows) * 0.65),
                     col_widths=[1.0, 4.5, 1.5, 3.2, 2.5, 2.5, 2.5, 2.5],
                     font_size=8, header_size=8, alt_fill=_GRAYL)
    sig_col = 3

    # Colorer signal cells
    for r in range(1, len(rows)):
        sig = sorted_s[r-1][3] if r-1 < len(sorted_s) else "Neutre"
        _color_cell(tbl, r, sig_col, _sig_light(sig), _sig_color(sig))

    # Lecture analytique enrichie — noms secteurs traduits en francais
    nb_surp = sum(1 for s in secteurs if len(s) > 3 and "Surp" in str(s[3]))
    nb_sous = sum(1 for s in secteurs if len(s) > 3 and "Sous" in str(s[3]))
    nb_neut = nb_s - nb_surp - nb_sous
    surp_noms = " · ".join(_abbrev_sector(s[0], 22) for s in sorted_s if len(s) > 3 and "Surp" in str(s[3]))
    sous_noms = " · ".join(_abbrev_sector(s[0], 22) for s in sorted_s if len(s) > 3 and "Sous" in str(s[3]))
    top_s  = sorted_s[0]  if sorted_s else None
    bot_s  = sorted_s[-1] if sorted_s else None
    top_nom = _abbrev_sector(top_s[0], 22) if top_s else "—"
    bot_nom = _abbrev_sector(bot_s[0], 22) if bot_s else "—"
    score_top = top_s[2] if top_s else "—"
    score_bot = bot_s[2] if bot_s else "—"
    score_spread = (top_s[2] - bot_s[2]) if (top_s and bot_s) else 0
    # Secteurs proches du seuil Surpondérer (score 55-65) — a surveiller
    _near_surp = [s for s in sorted_s if 55 <= s[2] < 65]
    _near_str  = " · ".join(_abbrev_sector(s[0], 22) for s in _near_surp[:3]) if _near_surp else ""
    # Meilleure marge EBITDA parmi les surponderes
    _top_mg    = max((s[5] for s in sorted_s if "Surp" in str(s[3]) and isinstance(s[5], (int,float)) and s[5] > 0), default=None)
    _top_mg_s  = f"{_top_mg:.1f}%" if _top_mg else "—"
    # Strategie selon dispersion
    if score_spread > 35:
        _strat = ("La dispersion élevée justifié une approche concentrée : surexposer les secteurs verts, "
                  "éviter les secteurs rouges. Un portefeuille également pondéré sous-performerait le benchmark.")
    elif score_spread > 20:
        _strat = ("La bifurcation modérée autorise une allocation selective : privilegier les leaders de score "
                  "tout en conservant une exposition diversifiée aux secteurs Neutrès proches du seuil de bascule.")
    else:
        _strat = ("La faible dispersion des scores signale une convergence sectorielle — "
                  "le stock-picking prime sur l'allocation sectorielle dans ce contexte.")
    lecture = (
        f"{nb_surp} secteur(s) sur {nb_s} franchissent le seuil Surpondérer (score > 65) : "
        f"{surp_noms or 'aucun'}. "
        f"{nb_sous} secteur(s) en Sous-pondérer : {sous_noms or 'aucun'}. "
        f"{nb_neut} secteur(s) en Neutre. "
        f"L'Écart leader/retardataire ({top_nom} {score_top}pts / "
        f"{bot_nom} {score_bot}pts) = {score_spread} pts de spread — "
        f"bifurcation {'marquée' if score_spread > 30 else 'modérée'}. "
        + (f"Secteurs proches du seuil de bascule (55-65) : {_near_str}. " if _near_str else "")
        + (f"Marge EBITDA Médiane des surpondérés : {_top_mg_s}. " if _top_mg else "")
        + _strat
    )
    y_top = min(10.5, 2.3 + len(rows) * 0.65 + 0.8)
    _carto_title = (
        f"{nb_surp} secteurs Surpondérer / {nb_sous} Sous-pondérer — spread {score_spread} pts "
        f"({'bifurcation marquée' if score_spread > 30 else 'dispersion modérée'})"
    )
    _lecture_box(slide, _carto_title,
                 lecture, y_top=y_top, height=13.35 - y_top)

    _footer(slide)
    return slide


def _s10_scatter(prs, D, chart_bytes: bytes):
    slide = _blank(prs)
    indice = D.get("indice","")
    secteurs = D.get("secteurs",[])
    ev_vals = [_parse_x(s[4]) for s in secteurs]
    use_mg  = not ev_vals or not any(v > 0 for v in ev_vals)
    if use_mg:
        axis_lbl = "Mg. EBITDA vs Croissance BPA"
        sub_lbl  = "Marge EBITDA (%) vs Croissance BPA  ·  4 quadrants  ·  Un point par secteur GICS"
    else:
        axis_lbl = "EV/EBITDA vs Croissance BPA"
        sub_lbl  = "EV/EBITDA Médian vs Croissance BPA  ·  4 quadrants  ·  Un point par secteur GICS"
    _header(slide, f"Valorisation vs Croissance BPA ({axis_lbl})", sub_lbl, active=2)

    _pic(slide, chart_bytes, 0.9, 2.3, 14.2, 8.6)

    # Lecture droite
    sorted_s = sorted(secteurs, key=lambda s: s[2], reverse=True)
    surp_s = [s for s in sorted_s if "Surp" in str(s[3])]
    sous_s = [s for s in sorted_s if "Sous" in str(s[3])]
    if use_mg:
        best = [_abbrev_sector(s[0],18) for s in surp_s if isinstance(s[5],(int,float)) and s[5] > 20][:2]
        opport = [_abbrev_sector(s[0],18) for s in surp_s if isinstance(s[5],(int,float)) and s[5] <= 20][:2]
        txt = (
            f"{'  ·  '.join(best) or 'N/A'} : forte marge EBITDA + score Surpondérer — "
            f"profil qualité avec visibilité élevée sur les benefices.\n\n"
            f"{'  ·  '.join(opport) or 'N/A'} : zone Opportunite — "
            f"croissance correcte, marge adéquate, valorisation attractive.\n\n"
            f"EV/EBITDA non disponible pour cet indice (EU) via yfinance. "
            f"La Mg. EBITDA est utilisee comme proxy de qualité opérationnelle. "
            f"Favoriser les secteurs a forte marge et croissance BPA positive."
        )
    else:
        premium = [_abbrev_sector(s[0],18) for s in surp_s if _parse_x(s[4]) > 20][:2]
        opport  = [_abbrev_sector(s[0],18) for s in surp_s if 0 < _parse_x(s[4]) <= 20][:2]
        # Fallback: si premium vide, top Surpondérer par score
        if not premium and surp_s:
            premium = [_abbrev_sector(s[0],18) for s in surp_s[:2]]
            _prem_label = "en zone Opportunite (EV/EBITDA attractif) — forte conviction score Surpondérer"
        else:
            _prem_label = "en quadrant Premium Justifie — valorisation élevée supportee par forte croissance BPA"
        if not opport and surp_s:
            opport = [_abbrev_sector(s[0],18) for s in surp_s[len(premium):len(premium)+2]]
        _has_proxy = any("*" in str(s[4]) for s in secteurs)
        _proxy_note = ("\n\n* EV/EBITDA estime via ETF proxy (indice). "
                       "Valeurs individuelles non disponibles." if _has_proxy else "")
        txt = (
            f"{'  ·  '.join(premium) or 'Aucun secteur Surpondérer'} {_prem_label if premium else ''}.\n\n"
            f"{'  ·  '.join(opport) or 'Aucun'} : valorisation attractive — "
            f"croissance correcte, meilleur ratio risque/rendement.\n\n"
            f"Les pointilles representent les Médianes sectorielles (EV/EBITDA et BPA growth). "
            f"Favoriser les secteurs en bas droite (forte croissance, valorisation modérée)."
            f"{_proxy_note}"
        )
    # Enrichissement LLM : analyse approfondie (audit Baptiste slide 10 2026-04-14)
    _llm_extra_s10 = ""
    try:
        from core.llm_provider import llm_call
        _surp_names = ", ".join(_abbrev_sector(s[0], 22) for s in surp_s[:4]) or "aucun"
        _sous_names = ", ".join(_abbrev_sector(s[0], 22) for s in sous_s[:4]) or "aucun"
        _prompt_s10 = (
            f"Tu es analyste sell-side senior. Redige une analyse (200-240 mots) "
            f"de la carte valorisation vs croissance pour {indice}.\n\n"
            f"Secteurs Surponderer : {_surp_names}\n"
            f"Secteurs Sous-ponderer : {_sous_names}\n\n"
            f"Structure en 2 paragraphes :\n"
            f"1. Lecture du scatter : quadrants dominants, couple "
            f"valorisation/croissance, signaux d'inflexion a surveiller\n"
            f"2. Implications allocation : sur/sous-pond\u00e9ration recommand\u00e9e, "
            f"condition de basculement vers un profil plus d\u00e9fensif ou "
            f"offensif, horizon 12 mois\n\n"
            f"Francais correct avec accents. Pas de markdown. Pas d'emojis. "
            f"Cite les secteurs en francais."
        )
        _llm_extra_s10 = llm_call(_prompt_s10, phase="long", max_tokens=700) or ""
    except Exception:
        pass
    _full_txt_s10 = (txt + "\n\n" + _llm_extra_s10).strip() if _llm_extra_s10 else txt

    _rect(slide, 15.6, 2.3, 8.9, 11.0, fill=_GRAYL)
    _rect(slide, 15.6, 2.3, 0.12, 11.0, fill=_NAVY)
    _txb(slide, "Lecture du Scatter & Implications", 15.9, 2.4, 8.4, 0.6,
         size=8.5, bold=True, color=_NAVY)
    _txb(slide, _full_txt_s10[:2000], 15.9, 3.1, 8.4, 10.0,
         size=7.5, color=_GRAYT, wrap=True)

    _footer(slide)
    return slide


def _s11_decomposition(prs, D):
    slide = _blank(prs)
    nb_s = D.get("nb_secteurs",11)
    _header(slide, "Décomposition des Scores FinSight",
            f"{nb_s} secteurs GICS  ·  Score 0-100  ·  Momentum 40 % · Révisions BPA 30 % · Valorisation 30 %",
            active=2)

    secteurs  = D.get("secteurs",[])
    top3      = D.get("top3_secteurs",[])
    sorted_s  = sorted(secteurs, key=lambda s: s[2], reverse=True)

    # Construire un mapping de sous-scores depuis top3
    sub_map = {}
    for t in top3:
        sub_map[t["nom"]] = (
            t.get("score_momentum", round(t.get("score",50) * 0.4)),
            t.get("score_revisions", round(t.get("score",50) * 0.3)),
            t.get("score_valorisation", round(t.get("score",50) * 0.3)),
        )

    _SIG_NORM_TBL = {
        "Surpondérer": "Surpondérer", "surpondérer": "Surpondérer",
        "Sous-pondérer": "Sous-pondérer", "sous-pondérer": "Sous-pondérer",
        "neutre": "Neutre",
    }
    rows = [["Secteur", "Score Total", "Momentum (40%)", "Révisions (30%)", "Valorisation (30%)", "Signal"]]
    for s in sorted_s:
        nom   = s[0] if len(s) > 0 else "—"
        score = s[2] if len(s) > 2 else 0
        sig   = s[3] if len(s) > 3 else "Neutre"
        if nom in sub_map:
            sm, sr, sv = sub_map[nom]
        else:
            sm = round(score * 0.40)
            sr = round(score * 0.30)
            sv = round(score * 0.30)
        sig_norm = _SIG_NORM_TBL.get(str(sig).strip().lower(), str(sig)[:15])
        rows.append([_abbrev_sector(nom, 20), str(score), str(sm), str(sr), str(sv), sig_norm])

    tbl = _add_table(slide, rows, 0.9, 2.3, 23.6,
                     min(7.0, 0.6 + len(rows) * 0.62),
                     col_widths=[4.5, 2.5, 3.5, 3.5, 3.5, 3.5],
                     font_size=8, header_size=8, alt_fill=_GRAYL)

    # Colorer signal
    for r in range(1, len(rows)):
        sig = sorted_s[r-1][3] if r-1 < len(sorted_s) else "Neutre"
        _color_cell(tbl, r, 5, _sig_light(sig), _sig_color(sig))
    # Colorer score
    for r in range(1, len(rows)):
        score = int(rows[r][1]) if rows[r][1].isdigit() else 50
        fill = _BUY_L if score >= 65 else (_SELL_L if score < 45 else _HOLD_L)
        _color_cell(tbl, r, 1, fill, _BLACK)

    note = "Lecture : Score >= 65 = vert · 45-64 = neutre · < 45 = rouge · Méthode : composite 40 % momentum + 30 % révisions BPA + 30 % valorisation relative"
    y_note = 2.3 + len(rows) * 0.62 + 0.2
    _txb(slide, note, 0.9, y_note, 23.6, 0.6, size=7, color=_GRAYD, italic=True)

    # Lecture enrichie — noms secteurs traduits en francais
    top_nom_raw = sorted_s[0][0] if sorted_s else "—"
    bot_nom_raw = sorted_s[-1][0] if sorted_s else "—"
    top_nom = _abbrev_sector(top_nom_raw, 22) if sorted_s else "—"
    bot_nom = _abbrev_sector(bot_nom_raw, 22) if sorted_s else "—"
    top_scores = sub_map.get(top_nom_raw, (0,0,0))
    bot_scores = sub_map.get(bot_nom_raw, (0,0,0))
    # Dimension dominante overall
    all_mom = [sub_map.get(s[0],(0,0,0))[0] for s in sorted_s if s[0] in sub_map]
    all_rev = [sub_map.get(s[0],(0,0,0))[1] for s in sorted_s if s[0] in sub_map]
    all_val = [sub_map.get(s[0],(0,0,0))[2] for s in sorted_s if s[0] in sub_map]
    _dom = "momentum" if (sum(all_mom) >= sum(all_rev) and sum(all_mom) >= sum(all_val)) else ("révisions BPA" if sum(all_rev) >= sum(all_val) else "valorisation relative")
    # Secteurs unanimes sur 3 dimensions (scores > 50 partout) — noms FR
    _unanimes = [_abbrev_sector(s[0], 22) for s in sorted_s if s[0] in sub_map and all(v > 50 for v in sub_map[s[0]])]
    _unani_str = " · ".join(_unanimes[:3]) if _unanimes else "aucun"
    _top_sc = sorted_s[0][2] if sorted_s and len(sorted_s[0]) > 2 else "—"
    _bot_sc = sorted_s[-1][2] if sorted_s and len(sorted_s[-1]) > 2 else "—"
    lecture = (
        f"{top_nom} affiche le profil le plus robuste (score {_top_sc}/100) "
        f"avec momentum={top_scores[0]}pts, révisions={top_scores[1]}pts, valorisation={top_scores[2]}pts. "
        f"\u00c0 l'oppos\u00e9, {bot_nom} (score {_bot_sc}/100) cumule les handicaps : "
        f"momentum={bot_scores[0]}, révisions={bot_scores[1]}, valorisation={bot_scores[2]}. "
        f"Dimension dominante sur l'ensemble de l'univers : {_dom}. "
        f"Secteurs avec profil composite fort (>50 sur les 3 dimensions) : {_unani_str}. "
        f"Ces secteurs repr\u00e9sentent les meilleures opportunit\u00e9s d'entr\u00e9e dans le contexte actuel."
    )
    y_lec = min(10.5, y_note + 0.8)
    _lecture_box(slide, "Analyse des scores — Points saillants",
                 lecture, y_top=y_lec, height=13.35 - y_lec)

    _footer(slide)
    return slide


def _s13_top3(prs, D):
    slide = _blank(prs)
    _header(slide, "Top 3 Secteurs — Synthèse",
            "Secteurs Surpondérer  ·  Signal · Score · EV/EBITDA · Sociétés représentatives · Catalyseur · Risque",
            active=3)

    top3 = D.get("top3_secteurs", [])
    if not top3:
        _txb(slide, "Données top3 non disponibles", 1.0, 5.0, 23.0, 1.0,
             size=12, color=_GRAYT, align=PP_ALIGN.CENTER)
        _footer(slide)
        return slide

    # Padding to always render 3 panels (avoids blank slots when <3 sectors)
    while len(top3) < 3:
        top3 = list(top3) + [{"nom": "—", "signal": "Neutre", "score": 0,
                               "ev_ebitda": "—", "catalyseur": "Données insuffisantés",
                               "risque": "—", "societes": []}]

    panel_w = 7.5
    for i, sect in enumerate(top3[:3]):
        xoff = 0.9 + i * (panel_w + 0.35)
        nom   = sect.get("nom","—")
        sig   = sect.get("signal","Neutre")
        score = sect.get("score",50)
        ev    = sect.get("ev_ebitda","—")
        cat   = sect.get("catalyseur","—")[:180]
        rsk   = sect.get("risque","—")[:180]
        socs  = sect.get("societes",[])

        # Panel bg
        _rect(slide, xoff, 2.3, panel_w, 11.1, fill=_GRAYL)
        _rect(slide, xoff, 2.3, panel_w, 0.8, fill=_NAVY)

        # Nom + signal badge
        _txb(slide, nom[:18], xoff + 0.2, 2.35, panel_w - 2.5, 0.7,
             size=10, bold=True, color=_WHITE)
        _rect(slide, xoff + panel_w - 2.0, 2.35, 1.8, 0.65, fill=_sig_color(sig))
        _txb(slide, "Surp." if "Surp" in sig else ("Sous." if "Sous" in sig else "Neutre"),
             xoff + panel_w - 1.9, 2.4, 1.7, 0.55, size=7.5, bold=True,
             color=_WHITE, align=PP_ALIGN.CENTER)

        # Score + metrique dispo (EV/EBITDA ou Mg.EBITDA en fallback)
        mg_eb = sect.get("mg_ebitda", sect.get("marge_ebitda", "—"))
        if str(ev) in ("—", "", "None") and mg_eb and str(mg_eb) not in ("—","","None"):
            _mg_disp = f"{float(mg_eb):.1f}%" if isinstance(mg_eb,(int,float)) else str(mg_eb)
            _met_disp = f"Score : {score}/100  ·  Mg.EBITDA : {_mg_disp}"
        else:
            _met_disp = f"Score : {score}/100  ·  EV/EBITDA : {ev}"
        _txb(slide, _met_disp, xoff + 0.2, 3.25, panel_w - 0.3, 0.5, size=7.5, color=_NAVY)

        # Sociétés
        _txb(slide, "Sociétés représentatives", xoff + 0.2, 3.85, panel_w - 0.3, 0.45,
             size=7.5, bold=True, color=_NAVY)
        for j, soc in enumerate(socs[:3]):
            yy = 4.4 + j * 0.82
            _rect(slide, xoff + 0.2, yy, panel_w - 0.3, 0.72, fill=_WHITE)
            tk  = soc[0] if isinstance(soc,(list,tuple)) else str(soc)
            sg  = soc[1] if isinstance(soc,(list,tuple)) and len(soc)>1 else "—"
            evs = soc[2] if isinstance(soc,(list,tuple)) and len(soc)>2 else "—"
            sc2 = soc[3] if isinstance(soc,(list,tuple)) and len(soc)>3 else "—"
            # Si EV/EBITDA absent au niveau société, utiliser le secteur comme proxy
            if str(evs) in ("—", "", "None") and str(ev) not in ("—", "", "None"):
                evs = f"~{ev}"
            # Abrev signal
            sg_abbr = "Surp." if "Surp" in str(sg) else ("Sous." if "Sous" in str(sg) else "Neutre")
            _ev_absent = str(evs) in ("—", "", "None", "\u2014")
            if _ev_absent:
                # Layout 3 col : ticker | signal | score (pas d'EV)
                _txb(slide, tk,      xoff + 0.4,       yy + 0.1, 2.4, 0.55,
                     size=8, bold=True, color=_NAVY)
                _txb(slide, sg_abbr, xoff + 0.4 + 2.4, yy + 0.1, 2.2, 0.55,
                     size=8, color=_sig_color(sg))
                _sc_lbl = f"Score : {sc2}" if str(sc2) not in ("—","","None") else ""
                _txb(slide, _sc_lbl, xoff + 0.4 + 2.4 + 2.2, yy + 0.1, 2.2, 0.55,
                     size=8, color=_GRAYT)
            else:
                _txb(slide, tk,       xoff + 0.4,             yy + 0.1, 1.9, 0.55,
                     size=8, bold=True, color=_NAVY)
                _txb(slide, sg_abbr,  xoff + 0.4 + 1.9,       yy + 0.1, 1.5, 0.55,
                     size=8, color=_sig_color(sg))
                _txb(slide, str(evs), xoff + 0.4 + 1.9 + 1.5, yy + 0.1, 1.5, 0.55,
                     size=8, color=_GRAYT)
                _txb(slide, str(sc2), xoff + 0.4 + 1.9 + 1.5 + 1.5, yy + 0.1, 1.5, 0.55,
                     size=8, color=_GRAYT)

        # Catalyseur
        _txb(slide, "Catalyseur", xoff + 0.2, 6.9, panel_w - 0.3, 0.45,
             size=7.5, bold=True, color=_NAVY)
        _rect(slide, xoff + 0.2, 7.4, panel_w - 0.3, 2.0, fill=_BUY_L)
        _txb(slide, cat, xoff + 0.4, 7.5, panel_w - 0.6, 1.85, size=7.5, color=_GRAYT, wrap=True)

        # Risque
        _txb(slide, "Risque principal", xoff + 0.2, 9.55, panel_w - 0.3, 0.45,
             size=7.5, bold=True, color=_NAVY)
        _rect(slide, xoff + 0.2, 10.05, panel_w - 0.3, 2.0, fill=_SELL_L)
        _txb(slide, rsk, xoff + 0.4, 10.15, panel_w - 0.6, 1.85, size=7.5, color=_GRAYT, wrap=True)

    _footer(slide)
    return slide


def _s14_allocation(prs, D):
    """Slide 14 — Allocation Optimale Markowitz (S&P 500) ou message limitation."""
    slide = _blank(prs)
    _header(slide, "Allocation Optimale — Markowitz",
            "Min-Variance · Tangency (Max Sharpe) · Equal Risk Contribution  ·  ETF SPDR sectoriels 52S",
            active=3)

    opt = D.get("optimal_portfolios")
    if not opt or not opt.get("sectors"):
        # Indice non-S&P : message explicatif
        _rect(slide, 0.9, 2.5, 23.6, 4.5, fill=_GRAYL)
        _rect(slide, 0.9, 2.5, 0.12, 4.5, fill=_NAVY)
        _txb(slide, "Optimisation non disponible pour cet indice",
             1.2, 2.6, 23.0, 0.7, size=11, bold=True, color=_NAVY)
        _txb(slide,
             "L'optimisation de portefeuille Markowitz repose sur les ETF SPDR sectoriels (XLK, XLV, XLF...) "
             "qui ne couvrent que le Marché US (S&P 500). Pour les indices européens (CAC 40, DAX, Euro Stoxx), "
             "les proxies sectoriels équivalents ne sont pas disponibles via yfinance. "
             "L'optimisation est donc calculée uniquement lors d'une analyse S&P 500.",
             1.2, 3.4, 23.0, 3.3, size=9, color=_GRAYT, wrap=True)

        # Afficher un rappel des secteurs recommandes comme substitut
        secteurs = D.get("secteurs", [])
        surp = [(s[0], s[2]) for s in secteurs if "surp" in str(s[3]).lower()]
        sous = [(s[0], s[2]) for s in secteurs if "sous" in str(s[3]).lower()]
        if surp or sous:
            _rect(slide, 0.9, 7.5, 23.6, 0.6, fill=_NAVY)
            _txb(slide, "POSITIONNEMENT SECTORIEL RECOMMANDE (sans optimisation quantitative)",
                 1.1, 7.55, 23.0, 0.5, size=8, bold=True, color=_WHITE)
            _rect(slide, 0.9, 8.15, 11.4, 4.0, fill=_BUY_L)
            _rect(slide, 0.9, 8.15, 0.1, 4.0, fill=_BUY)
            _txb(slide, "SURPONDÉRER", 1.1, 8.2, 10.9, 0.45, size=7.5, bold=True, color=_BUY)
            for i, (nm, sc) in enumerate(surp[:4]):
                _txb(slide, f"• {_abbrev_sector(nm, 30)} — score {sc}/100",
                     1.2, 8.7 + i * 0.75, 10.8, 0.65, size=8, color=_GRAYT)
            _rect(slide, 13.1, 8.15, 11.4, 4.0, fill=_SELL_L)
            _rect(slide, 13.1, 8.15, 0.1, 4.0, fill=_SELL)
            _txb(slide, "SOUS-PONDÉRER", 13.3, 8.2, 10.9, 0.45, size=7.5, bold=True, color=_SELL)
            for i, (nm, sc) in enumerate(sous[:4]):
                _txb(slide, f"• {_abbrev_sector(nm, 30)} — score {sc}/100",
                     13.4, 8.7 + i * 0.75, 10.8, 0.65, size=8, color=_GRAYT)
        _footer(slide)
        return slide

    # S&P 500 — afficher les 3 portefeuilles
    sectors  = opt.get("sectors", [])
    _mv_raw  = opt.get("min_var",  [])
    _tang_raw= opt.get("tangency", [])
    _erc_raw = opt.get("erc",      [])

    # Structure dict {"weights":[...], "sharpe":...} ou liste directe
    if isinstance(_mv_raw, dict):
        mv   = _mv_raw.get("weights", [])
        tang = _tang_raw.get("weights", []) if isinstance(_tang_raw, dict) else _tang_raw
        erc  = _erc_raw.get("weights",  []) if isinstance(_erc_raw,  dict) else _erc_raw
        sharpe = {
            "mv":       _mv_raw.get("sharpe", 0),
            "tangency": _tang_raw.get("sharpe", 0) if isinstance(_tang_raw, dict) else 0,
            "erc":      _erc_raw.get("sharpe",  0) if isinstance(_erc_raw,  dict) else 0,
        }
    else:
        mv, tang, erc = _mv_raw, _tang_raw, _erc_raw
        sharpe = opt.get("sharpe", {})

    # Les poids sont déjà en % (ex: 15.3 signifie 15.3%), pas en fraction decimale
    def _fmt_w(lst, i):
        if i >= len(lst):
            return "—"
        v = lst[i]
        try:
            fv = float(v)
            # Si fraction decimale (0-1), convertir en %
            return f"{fv*100:.1f} %" if fv <= 1.0 else f"{fv:.1f} %"
        except (TypeError, ValueError):
            return "—"

    # Table centrale : secteur | Min-Var | Tangency | ERC
    rows = [["Secteur", "Min-Variance", "Tangency (Max Sharpe)", "ERC"]]
    for i, sec in enumerate(sectors):
        rows.append([_abbrev_sector(sec, 25),
                     _fmt_w(mv, i), _fmt_w(tang, i), _fmt_w(erc, i)])

    # Row Sharpe
    rows.append([
        "Sharpe ratio",
        f"{sharpe.get('mv', 0):.2f}",
        f"{sharpe.get('tangency', 0):.2f}",
        f"{sharpe.get('erc', 0):.2f}",
    ])

    tbl_h = min(9.5, 0.7 + len(rows) * 0.65)
    tbl = _add_table(slide, rows, 0.9, 2.4, 15.5, tbl_h,
                     col_widths=[5.5, 3.0, 3.5, 3.0],
                     font_size=8, header_size=8, alt_fill=_GRAYL)

    # Colorer les lignes selon le poids Tangency (portefeuille prioritaire)
    _tang_floats = []
    for v in tang:
        try:
            fv = float(v)
            _tang_floats.append(fv if fv > 1.0 else fv * 100)
        except (TypeError, ValueError):
            _tang_floats.append(0.0)
    for r in range(1, len(rows) - 1):  # skip header and sharpe row
        if r - 1 < len(_tang_floats):
            w = _tang_floats[r - 1]
            if w >= 20.0:
                try:
                    from pptx.util import Pt
                    from pptx.dml.color import RGBColor as _RGB
                    for ci in range(4):
                        cell = tbl.cell(r, ci)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = _RGB(0xEA, 0xF4, 0xEF)
                except Exception:
                    pass

    # Panel droit — lecture analytique
    try:
        mv_floats  = [float(v) for v in mv]
        t_floats   = [float(v) for v in tang]
        erc_floats = [float(v) for v in erc]
        mv_max  = sectors[mv_floats.index(max(mv_floats))]   if mv_floats  and sectors else "—"
        t_max   = sectors[t_floats.index(max(t_floats))]     if t_floats   and sectors else "—"
        erc_max = sectors[erc_floats.index(max(erc_floats))] if erc_floats and sectors else "—"
    except (ValueError, IndexError):
        mv_max = t_max = erc_max = "—"
    _rect(slide, 17.0, 2.4, 7.5, tbl_h, fill=_GRAYL)
    _rect(slide, 17.0, 2.4, 0.12, tbl_h, fill=_NAVY)
    _txb(slide, "LECTURE QUANTITATIVE", 17.3, 2.5, 7.0, 0.55, size=8, bold=True, color=_NAVY)
    lecture = (
        f"Portefeuille Tangency (Max Sharpe = {sharpe.get('tangency', 0):.2f}) : "
        f"sur-allocation sur {t_max} — secteur offrant le meilleur ratio rendement/risque sur 52 semaines. "
        f"Ce portefeuille est le plus agressif des trois.\n\n"
        f"Min-Variance (Sharpe = {sharpe.get('mv', 0):.2f}) : poids dominant sur {mv_max} — "
        f"secteur le moins volatil sur la période. Approche défensive, adaptee a un contexte d'incertitude élevée.\n\n"
        f"ERC (Sharpe = {sharpe.get('erc', 0):.2f}) : allocation équilibrée, {erc_max} ressort dominant. "
        f"Chaque secteur contribue également au risque total du portefeuille.\n\n"
        f"Note : Basé sur rendements journaliers 52S ETF SPDR. Contraintes : poids 0-40% par secteur."
    )
    _txb(slide, lecture, 17.3, 3.1, 7.0, tbl_h - 0.8, size=7.5, color=_GRAYT, wrap=True)

    _footer(slide)
    return slide


def _s15_zone_entree(prs, D, chart_bytes: bytes):
    slide = _blank(prs)
    _header(slide, "Zone d'Entrée Optimale par Secteur",
            "P/E actuel vs médiane historique 10 ans  ·  Signal d'entrée Accumuler / Neutre / Alléger",
            active=3)

    # Graphique GAUCHE (60 % de la largeur)
    _pic(slide, chart_bytes, 0.9, 2.3, 15.0, 10.5)

    # Colonne analytique DROITE
    top3    = D.get("top3_secteurs",[])
    secteurs = D.get("secteurs",[])
    _pe_g   = D.get("pe_forward","—")
    _pe_med = D.get("pe_Médiane_10y","—")

    noms_b  = []
    noms_h  = []
    for t in top3:
        pef = t.get("pe_forward_raw",0) or 0
        pem = t.get("pe_Médiane_10y",18.0) or 18.0
        if pef > 0 and pef < pem * 1.05:
            noms_b.append(t["nom"])
        elif pef > 0 and pef > pem * 1.15:
            noms_h.append(t["nom"])

    # Secteurs tous confondus sous/sur Médiane
    all_sous = [s[0] for s in secteurs if True]  # placeholder
    _pm_str = f"{_pe_med}x" if isinstance(_pe_med,(int,float)) else str(_pe_med)

    txt_col = (
        f"PE médiane 10 ans : {_pm_str}\n"
        f"PE Forward indice : {_pe_g}\n\n"
        f"Zone d'entrée favorable (PE < Médiane) :\n"
        f"{'  ·  '.join([_abbrev_sector(n,18) for n in noms_b]) or 'Aucun secteur top3'}\n\n"
        f"Zone de prudence (PE > Médiane +15 %) :\n"
        f"{'  ·  '.join([_abbrev_sector(n,18) for n in noms_h]) or 'Aucun'}\n\n"
        f"Méthodologie : le PE Forward est ancre sur le PE global de l'indice "
        f"ajust\u00e9 du score sectoriel. Les secteurs \u00e0 gauche de la m\u00e9diane "
        f"10 ans offrent un meilleur point d'entr\u00e9e."
    )
    # Enrichissement LLM (audit Baptiste slide 15 2026-04-14)
    _llm_zone_s15 = ""
    try:
        from core.llm_provider import llm_call
        _indice_name = D.get("indice", "l'indice")
        _buy_sec  = ", ".join(_abbrev_sector(n, 22) for n in noms_b) or "aucun"
        _high_sec = ", ".join(_abbrev_sector(n, 22) for n in noms_h) or "aucun"
        _prompt_s15 = (
            f"Tu es analyste sell-side senior. Redige une analyse (200-240 mots) "
            f"sur les zones d'entree par secteur pour {_indice_name}.\n\n"
            f"PE Forward indice : {_pe_g}. PE median 10 ans : {_pm_str}.\n"
            f"Secteurs en zone d'entree favorable : {_buy_sec}\n"
            f"Secteurs en zone de prudence : {_high_sec}\n\n"
            f"Structure en 2 paragraphes :\n"
            f"1. Interpretation des zones d'entree : conditions pour que le PE "
            f"se re-rate, triggers macro et microstructure a surveiller\n"
            f"2. Tactique d'accumulation : phasing recommande pour les secteurs "
            f"en zone favorable, sizing initial vs rampe de conviction, stop-loss "
            f"analytique base sur la deterioration des fondamentaux\n\n"
            f"Francais correct avec accents. Pas de markdown. Pas d'emojis."
        )
        _llm_zone_s15 = llm_call(_prompt_s15, phase="long", max_tokens=700) or ""
    except Exception:
        pass
    _full_s15 = (txt_col + "\n\n" + _llm_zone_s15).strip() if _llm_zone_s15 else txt_col

    _rect(slide, 16.3, 2.3, 8.1, 11.0, fill=_GRAYL)
    _rect(slide, 16.3, 2.3, 0.12, 11.0, fill=_NAVY)
    _txb(slide, "Lecture — Zone d'entr\u00e9e & tactique", 16.6, 2.4, 7.6, 0.6,
         size=8.5, bold=True, color=_NAVY)
    _txb(slide, _full_s15[:2000], 16.6, 3.1, 7.6, 10.0,
         size=7.5, color=_GRAYT, wrap=True)

    _footer(slide)
    return slide


def _s17_risques(prs, D):
    slide = _blank(prs)
    indice = D.get("indice","")
    _SIG_NORM_R = {
        "Surpondérer": "Surpondérer", "surpondérer": "Surpondérer",
        "Sous-pondérer": "Sous-pondérer", "sous-pondérer": "Sous-pondérer",
        "neutre": "Neutre",
    }
    sig    = D.get("signal_global","Neutre")
    sig_disp = _SIG_NORM_R.get(sig.strip().lower(), sig)
    _header(slide, "Risques Macro & Scénarios",
            f"Analyse adversariale  ·  3 scénarios alternatifs  ·  Conditions d'invalidation du signal {sig_disp}",
            active=4)

    scenarios = D.get("scenarios",[])
    if not scenarios:
        scenarios = [
            {"titre":"Récession technique","prob":"18 %",
             "desc":"Deux trimestrès PIB < 0 % — révision BPA -15/-20 %. Signal passerait Sous-pondérer."},
            {"titre":"Resserrement Fed prolongé","prob":"35 %",
             "desc":"Fed Funds > 4,5 % jusqu'en 2027 — compression multiples growth et rotation sectorielle."},
            {"titre":"Choc géopolitique","prob":"20 %",
             "desc":"Escalade géopolitique — spike VIX > 35, rotation défensive vers Consumer Staples et Health Care."},
        ]

    # 3 blocs scénarios
    box_w = 7.5
    for i, sc in enumerate(scenarios[:3]):
        xoff = 0.9 + i * (box_w + 0.35)
        # Handle both tuple format (titre, condition, signal, horizon) and dict format
        if isinstance(sc, (list, tuple)):
            _sc = list(sc) + ["", "", "", ""]
            sc_titre  = str(_sc[0] or "—")
            sc_desc   = str(_sc[1] or "")   # condition = desc
            sc_signal = str(_sc[2] or "")   # signal résultat
            # Couleur et probabilite selon signal du scénario
            _sl = sc_signal.lower()
            if "surp" in _sl:
                hdr_col  = _BUY
                prob_str = "25 %"
            elif "sous" in _sl:
                hdr_col  = _SELL
                prob_str = "20 %"
            else:
                hdr_col  = _HOLD
                prob_str = "35 %"
        else:
            prob_str = str(sc.get("prob","—"))
            try:
                prob_int = int(prob_str.replace('%','').replace(' ',''))
            except:
                prob_int = 20
            sc_titre = sc.get("titre","—")
            sc_desc  = sc.get("desc","")
            hdr_col = _SELL if prob_int >= 35 else (_HOLD if prob_int >= 25 else _BUY)

        _rect(slide, xoff, 2.3, box_w, 4.3, fill=_GRAYL)
        _rect(slide, xoff, 2.3, box_w, 0.85, fill=hdr_col)
        _txb(slide, sc_titre[:40], xoff + 0.2, 2.35, box_w - 2.5, 0.75,
             size=9, bold=True, color=_WHITE)
        _txb(slide, f"Prob. {prob_str}", xoff + box_w - 2.2, 2.4, 2.0, 0.65,
             size=9, bold=True, color=_WHITE, align=PP_ALIGN.RIGHT)
        _txb(slide, sc_desc[:220], xoff + 0.2, 3.25, box_w - 0.3, 3.2,
             size=8, color=_GRAYT, wrap=True)

    # Conditions invalidation
    conds = D.get("conditions_invalidation",[])
    if not conds:
        conds = [
            f"{indice} casse le support clé — signal passe Sous-pondérer",
            "Fed pivot dovish confirmé + CPI < 2,5 % — signal passe Surpondérer",
            "Révisions BPA agrégées < -5 % sur 2 trimestrès consécutifs",
        ]

    _rect(slide, 0.9, 7.0, 23.6, 0.55, fill=_NAVY)
    _txb(slide, "Conditions d'invalidation du signal " + sig_disp.upper(),
         1.1, 7.05, 22.8, 0.45, size=8, bold=True, color=_WHITE)

    for j, cond in enumerate(conds[:4]):
        yy = 7.7 + j * 1.1
        _rect(slide, 0.9, yy, 0.1, 0.85, fill=_HOLD)
        _txb(slide, str(cond)[:120], 1.15, yy, 22.3, 0.9, size=8.5, color=_GRAYT, wrap=True)

    # Lecture analytique scénarios
    _sig17 = D.get("signal_global","Neutre")
    _conv17 = D.get("conviction_pct", 50)
    _scen_prob_bull = next((int(str(sc.get("prob","25")).replace('%','').strip()) for sc in scenarios[:1] if isinstance(sc,dict) and "Surp" not in str(sc.get("titre",""))), 25)
    _scen_prob_bear = next((int(str(sc.get("prob","20")).replace('%','').strip()) for sc in scenarios if isinstance(sc,dict) and ("Bear" in str(sc.get("titre","")) or "ecession" in str(sc.get("titre","")))), 20)
    _residuel = 100 - _scen_prob_bull - _scen_prob_bear
    _y17_lec = 7.7 + min(len(conds), 4) * 1.1 + 0.2
    if _y17_lec < 11.5:
        _h17_lec = 13.35 - _y17_lec
        # LLM analyse des risques (audit Baptiste 2026-04-14 : chasse hardcoding #90)
        _lec17_llm = ""
        try:
            from core.llm_provider import llm_call
            _scen_desc = " | ".join(
                f"{sc.get('titre','?') if isinstance(sc,dict) else (sc[0] if isinstance(sc,(list,tuple)) and len(sc)>0 else '?')} "
                f"({sc.get('prob','?') if isinstance(sc,dict) else '?'})"
                for sc in scenarios[:3]
            )
            _prompt_s17 = (
                f"Tu es strategist buy-side senior. Redige une analyse (200-260 mots) "
                f"des risques macro et scenarios alternatifs pour l'indice {indice}.\n\n"
                f"Signal central : {_sig17} (conviction {_conv17}%).\n"
                f"Scenarios identifies : {_scen_desc}\n\n"
                f"Structure en 3 paragraphes :\n"
                f"1. Lecture des scenarios alternatifs : ce qu'implique la probabilite "
                f"de chaque scenario, les triggers respectifs a surveiller\n"
                f"2. Conditions d'invalidation du signal central : ce qui ferait "
                f"basculer la these et sur quel horizon\n"
                f"3. Couverture recommandee : secteurs defensifs a renforcer, "
                f"instruments de hedge possibles (options, obligations, or)\n\n"
                f"Francais correct avec accents. Pas de markdown. Pas d'emojis."
            )
            _lec17_llm = llm_call(_prompt_s17, phase="long", max_tokens=700) or ""
        except Exception:
            pass
        if not _lec17_llm.strip():
            _lec17 = (
                f"Signal {_sig17} (conviction {_conv17}%) — sc\u00e9nario central. "
                f"Sc\u00e9nario haussier ({_scen_prob_bull}% de probabilit\u00e9) : catalyseurs macro "
                f"suffisants pour franchir le seuil Surpond\u00e9rer. "
                f"Sc\u00e9nario baissier ({_scen_prob_bear}%) : d\u00e9t\u00e9rioration des fondamentaux, "
                f"passage Sous-pond\u00e9rer. Sc\u00e9nario r\u00e9siduel ({_residuel}%) : stagnation ou "
                f"choc exog\u00e8ne non anticip\u00e9. Surveiller les conditions d'invalidation."
            )
        else:
            _lec17 = _lec17_llm[:1400]
        _lecture_box(slide, "Analyse des risques — Probabilités et implications", _lec17,
                     y_top=_y17_lec, height=_h17_lec)

    _footer(slide)
    return slide


def _s18_rotation(prs, D):
    slide = _blank(prs)
    indice = D.get("indice","")
    phase  = D.get("phase_cycle","Expansion avancée")
    _header(slide, "Rotation Sectorielle & Cycle Économique",
            f"{indice}  ·  Phase actuelle : {phase}  ·  Sensibilités taux/PIB  ·  Signal de rotation",
            active=4)

    rotation = D.get("rotation",[])
    if not rotation:
        rotation = [
            ("Technologie",             "Expansion",  "Faible",     "Forte",    "Surpondérer"),
            ("Santé",                   "Tous cycles","Modérée",    "Modérée",  "Surpondérer"),
            ("Finance",                 "Expansion",  "Haute",      "Haute",    "Surpondérer"),
            ("Conso. Cyclique",         "Expansion",  "Modérée",    "Haute",    "Neutre"),
            ("Télécoms",                "Expansion",  "Faible",     "Modérée",  "Neutre"),
            ("Industrie",               "Expansion",  "Modérée",    "Forte",    "Neutre"),
            ("Conso. Défensive",        "Contraction","Modérée",    "Faible",   "Neutre"),
            ("Énergie",                 "Tous cycles","Faible",     "Modérée",  "Neutre"),
            ("Immobilier",              "Contraction","Très haute", "Faible",   "Sous-pondérer"),
            ("Services Publics",        "Contraction","Très haute", "Faible",   "Sous-pondérer"),
            ("Matériaux",               "Expansion",  "Faible",     "Forte",    "Neutre"),
        ]

    rows = [["Secteur", "Phase favorisée", "Sens. Taux", "Sens. PIB", "Signal Rotation"]]
    for r in rotation:
        row = [_abbrev_sector(str(r[0]), 20)] + [str(r[i])[:20] if i < len(r) else "—" for i in range(1, 5)]
        rows.append(row)

    tbl = _add_table(slide, rows, 0.9, 2.3, 23.6,
                     min(7.8, 0.6 + len(rows) * 0.65),
                     col_widths=[4.5, 3.5, 3.0, 3.0, 3.5],
                     font_size=8, header_size=8, alt_fill=_GRAYL)

    # Colorer signal rotation
    for r in range(1, len(rows)):
        sig = rows[r][4] if len(rows[r]) > 4 else "Neutre"
        _color_cell(tbl, r, 4, _sig_light(sig), _sig_color(sig))

    # Lecture enrichie
    texte_rot = _trunc(D.get("texte_rotation",""), 350)
    if not texte_rot:
        surp_rots = [r[0] for r in rotation if len(r) > 4 and "Surp" in str(r[4])]
        sous_rots = [r[0] for r in rotation if len(r) > 4 and "Sous" in str(r[4])]
        acc_rots  = [r[0] for r in rotation if len(r) > 4 and "Accum" in str(r[4])]
        _taux_fort = [r[0] for r in rotation if len(r) > 2 and "lev" in str(r[2]).lower() or "haute" in str(r[2]).lower()]
        _pib_fort  = [r[0] for r in rotation if len(r) > 3 and "lev" in str(r[3]).lower() or "forte" in str(r[3]).lower()]
        texte_rot = (
            f"En phase {phase}, la grille de rotation favorise : "
            f"{' · '.join((surp_rots + acc_rots)[:3]) or 'N/A'} "
            f"(visibilité BPA forte, faible Sensibilité aux taux). "
            + (f"Alléger : {' · '.join(sous_rots[:2])}. " if sous_rots else "")
            + (f"Secteurs a forte Sensibilité taux a surveiller si la BCE pivote : {' · '.join(_taux_fort[:2])}. " if _taux_fort else "")
            + f"La rotation suit le cycle avec un decalage de 2-3 trimestrès — "
            f"un signal Surpondérer qui emerge maintenant anticipe une surperformance sur 6-12 mois. "
            f"Croiser avec les scores FinSight (slide 11) pour valider la cohérence momentum/fondamentaux."
        )

    y_lec = min(10.8, 2.3 + len(rows) * 0.65 + 0.3)
    _rect(slide, 0.9, y_lec, 23.6, 13.35 - y_lec, fill=_GRAYL)
    _rect(slide, 0.9, y_lec, 0.12, 13.35 - y_lec, fill=_NAVY)
    _txb(slide, "Positionnement de cycle — Synthèse et catalyseur de rotation",
         1.2, y_lec + 0.15, 22.8, 0.6, size=8, bold=True, color=_NAVY)
    _txb(slide, texte_rot, 1.2, y_lec + 0.8, 22.8, 13.35 - y_lec - 0.9,
         size=7.5, color=_GRAYT, wrap=True)

    _footer(slide)
    return slide


def _s19_sentiment(prs, D, chart_bytes: bytes):
    slide = _blank(prs)
    indice  = D.get("indice","")
    _header(slide, "Sentiment Sectoriel Composite",
            f"Score proxy dérivé des scores sectoriels  ·  {D.get('nb_secteurs',11)} secteurs  ·  Signal composite",
            active=4)

    sa = D.get("sentiment_agg",{})
    score = sa.get("score", 0.0)
    label = sa.get("label","Neutre")
    p_nb  = sa.get("positif_nb", 0)
    p_pct = sa.get("positif_pct", 0)
    n_nb  = sa.get("neutre_nb",  0)
    n_pct = sa.get("neutre_pct", 0)
    m_nb  = sa.get("negatif_nb",  0)
    m_pct = sa.get("negatif_pct", 0)
    t_pos = sa.get("themes_pos", [])
    t_neg = sa.get("themes_neg", [])

    # Layout: graphique gauche (x=0.9, w=13.5) + panneau texte droite (x=14.8, w=9.7)
    _score_col = _BUY if score > 0.05 else (_SELL if score < -0.05 else _HOLD)
    _score_str = (f"+{score:.2f}" if score > 0.01 else
                  (f"{score:.2f}" if score < -0.01 else "~0.00"))

    # === Graphique gauche — chart propre sans distorsion ===
    # Le chart est Généré avec le bon ratio pour la zone 13.5 x 11.0
    _pic(slide, chart_bytes, 0.9, 2.3, 13.5, 11.0)

    # === Panneau texte droite ===
    _rx = 14.9
    _rw = 9.5

    # Score composite + signal
    _rect(slide, _rx, 2.3, _rw, 3.2, fill=_GRAYL)
    _rect(slide, _rx, 2.3, 0.12, 3.2, fill=_score_col)
    _txb(slide, _score_str, _rx + 0.3, 2.4, _rw - 0.5, 1.5,
         size=28, bold=True, color=_score_col)
    _txb(slide, "Score composite secteurs", _rx + 0.3, 3.95, _rw - 0.5, 0.45, size=7.5, color=_GRAYT)
    _txb(slide, label, _rx + 0.3, 4.4, _rw - 0.5, 0.55, size=9, bold=True, color=_score_col)

    # Distribution Surpondérer / Neutre / Sous-pondérer
    dist_items = [
        ("Surpondérer",  p_nb, p_pct, _BUY),
        ("Neutre",       n_nb, n_pct, _HOLD),
        ("Sous-pondérer",m_nb, m_pct, _SELL),
    ]
    for j, (lbl, nb, pct, col) in enumerate(dist_items):
        yy = 5.7 + j * 1.05
        _rect(slide, _rx, yy, _rw, 0.9, fill=_GRAYL)
        _bar_w = pct * (_rw / 100.0)
        if _bar_w > 0:
            _rect(slide, _rx, yy, _bar_w, 0.9, fill=col)
        txt_col = _WHITE if pct >= 35 else _BLACK
        _txb(slide, lbl, _rx + 0.15, yy + 0.05, _rw - 0.3, 0.4, size=8, bold=True, color=txt_col)
        _txb(slide, f"{nb} sect. ({pct} %)", _rx + 0.15, yy + 0.47, _rw - 0.3, 0.38, size=7.5, color=txt_col)

    # Secteurs leaders / retardataires — texte traduit en FR + LLM enrichi
    _t_pos_fr = [_abbrev_sector(t, 22) if len(t) < 30 else t for t in t_pos]
    _t_neg_fr = [_abbrev_sector(t, 22) if len(t) < 30 else t for t in t_neg]
    pos_txt = "  ·  ".join(_t_pos_fr[:3]) if _t_pos_fr else "—"
    neg_txt = "  ·  ".join(_t_neg_fr[:3]) if _t_neg_fr else "—"
    _rect(slide, _rx, 8.80, _rw, 0.4, fill=_NAVY)
    _txb(slide, "SECTEURS LEADERS / RETARDATAIRES", _rx + 0.1, 8.82, _rw - 0.2, 0.36,
         size=7, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)
    _rect(slide, _rx, 9.25, _rw, 1.05, fill=_BUY_L)
    _txb(slide, "Surpond\u00e9rer : " + pos_txt[:250], _rx + 0.15, 9.30, _rw - 0.3, 0.95,
         size=7.5, color=_GRAYT, wrap=True)
    _rect(slide, _rx, 10.40, _rw, 1.05, fill=_SELL_L)
    _txb(slide, "Sous-pond\u00e9rer : " + neg_txt[:250], _rx + 0.15, 10.45, _rw - 0.3, 0.95,
         size=7.5, color=_GRAYT, wrap=True)

    # Enrichissement LLM (audit Baptiste slide 19 2026-04-14) : analyse sentiment
    _llm_sent_s19 = ""
    try:
        from core.llm_provider import llm_call
        _indice_name = D.get("indice", "l'indice")
        _prompt_s19 = (
            f"Tu es analyste sentiment senior. Redige une analyse (150-180 mots) "
            f"du sentiment sectoriel composite de {_indice_name}.\n\n"
            f"Score composite : {_score_str} ({label}). "
            f"Surponderer : {p_nb} sect. ({p_pct}%). "
            f"Neutre : {n_nb} sect. ({n_pct}%). "
            f"Sous-ponderer : {m_nb} sect. ({m_pct}%).\n"
            f"Secteurs leaders : {pos_txt}\n"
            f"Secteurs retardataires : {neg_txt}\n\n"
            f"Structure : (1) lecture du signal composite et ce qu'il implique "
            f"pour les 3 prochains mois, (2) divergences a surveiller entre "
            f"sentiment et fondamentaux (convergence/divergence), (3) impact "
            f"sur le positionnement actuel.\n"
            f"Francais correct avec accents. Pas de markdown."
        )
        _llm_sent_s19 = llm_call(_prompt_s19, phase="long", max_tokens=500) or ""
    except Exception:
        pass
    if _llm_sent_s19.strip():
        _rect(slide, _rx, 11.55, _rw, 1.75, fill=_GRAYL)
        _rect(slide, _rx, 11.55, 0.08, 1.75, fill=_NAVY)
        _txb(slide, "LECTURE SENTIMENT & IMPLICATIONS", _rx + 0.15, 11.60, _rw - 0.3, 0.35,
             size=7, bold=True, color=_NAVY)
        _txb(slide, _llm_sent_s19[:900], _rx + 0.15, 11.98, _rw - 0.3, 1.30,
             size=6.5, color=_GRAYT, wrap=True)

    _footer(slide)
    return slide


def _s20_etf_perf(prs, D, chart_bytes: bytes):
    slide = _blank(prs)
    indice = D.get("indice","")
    etf_perf_check = D.get("etf_perf",{})

    if not etf_perf_check:
        # Fallback : graphique performance indice 52 semaines + scorecard sectoriel
        _ph = D.get("perf_history")
        _has_chart = bool(_ph and _ph.get("dates") and _ph.get("indice"))

        if _has_chart:
            _header(slide, f"Performance de l'Indice — 52 Semaines",
                    f"{indice}  ·  vs S&P 500, Obligations & Or  ·  Indexe a 100  ·  Source yfinance",
                    active=5)
            # Générer le graphique index perf
            _idx_bytes = _chart_index_perf(D)
            _pic(slide, _idx_bytes, 0.9, 2.2, 23.6, 6.5)

            # Lecture analytique sous le graphique — narratif explicatif
            _i_perf  = _ph.get("indice", [])
            _b_perf  = _ph.get("bonds", [])
            _g_perf  = _ph.get("gold", [])
            _sp_perf = _ph.get("sp500", [])
            _dates   = _ph.get("dates", [])
            _ytd_val   = D.get("variation_ytd","—")
            _cours_val = D.get("cours","—")
            _final_idx = (_i_perf[-1]  - 100) if _i_perf  else 0
            _final_b   = (_b_perf[-1]  - 100) if _b_perf  else None
            _final_g   = (_g_perf[-1]  - 100) if _g_perf  else None
            _final_sp  = (_sp_perf[-1] - 100) if _sp_perf else None
            _pstart = _dates[0][:7] if _dates else "—"
            _pend   = _dates[-1][:7] if _dates else "—"

            # --- Identifier les mouvements cles pour le narratif ---
            _narrative_parts = []
            if len(_i_perf) >= 10 and len(_dates) >= 10:
                # Pic et creux sur la periode
                _peak_val = max(_i_perf)
                _peak_idx = _i_perf.index(_peak_val)
                _trough_val = min(_i_perf)
                _trough_idx = _i_perf.index(_trough_val)
                _peak_date = _dates[_peak_idx][:7]
                _trough_date = _dates[_trough_idx][:7]
                _drawdown = (_trough_val - _peak_val) / _peak_val * 100 if _peak_val and _peak_idx < _trough_idx else None

                # Tendance globale
                _trend = "haussière" if _final_idx > 5 else ("baissière" if _final_idx < -5 else "laterale")
                _narrative_parts.append(
                    f"Tendance {_trend} sur 52 semaines ({_final_idx:+.1f}%). "
                    f"Cours actuel : {_cours_val} (YTD : {_ytd_val})."
                )

                # Pic
                _peak_perf = _peak_val - 100
                if abs(_peak_perf) > 2:
                    _narrative_parts.append(
                        f"Plus haut atteint en {_peak_date} ({_peak_perf:+.1f}% vs base) "
                        f"— porté par l'appetit pour le risque et les flux entrants sur les equites européennes."
                    )

                # Drawdown depuis le pic
                if _drawdown is not None and abs(_drawdown) > 3:
                    _narrative_parts.append(
                        f"Correction de {_drawdown:.1f}% entre {_peak_date} et {_trough_date} "
                        f"— pression vendeuse liée aux incertitudes macro (taux directeurs, tensions commerciales) "
                        f"et rotation sectorielle vers les actifs refuges."
                    )
                elif _trough_idx > 0:
                    _trough_perf = _trough_val - 100
                    _narrative_parts.append(
                        f"Point bas en {_trough_date} ({_trough_perf:+.1f}% vs base)."
                    )

                # Rebond eventuel après le creux
                if _trough_idx < len(_i_perf) - 5:
                    _rebound = _i_perf[-1] - _trough_val
                    if _rebound > 3:
                        _narrative_parts.append(
                            f"Rebond de {_rebound:+.1f}pts depuis le creux de {_trough_date} "
                            f"— soutenu par les publications de résultats et les anticipations de politique monetaire accommodante."
                        )

            # --- Comparaisons cross-asset ---
            if _final_sp is not None:
                _rel = _final_idx - _final_sp
                _rel_lbl = f"surperformance de {_rel:+.1f}pt" if _rel >= 0 else f"sous-performance de {abs(_rel):.1f}pt"
                _narrative_parts.append(
                    f"Vs S&P 500 ({_final_sp:+.1f}%) : {_rel_lbl} — "
                    f"{'les equites européennes captent les flux de rebalancement' if _rel >= 0 else 'le Marché US reste porté par la tech et le momentum des megacaps'}."
                )
            if _final_g is not None and abs(_final_g) > 2:
                _narrative_parts.append(
                    f"Or ({_final_g:+.1f}%) {'surperforme nettement' if _final_g > _final_idx + 10 else 'sous-performe'} "
                    f"{'— signal d aversion au risque et demande de protection contre l inflation' if _final_g > _final_idx else '— appetit pour le risque dominant'}."
                )

            if not _narrative_parts:
                _narrative_parts = [
                    f"Performance {indice} : {_final_idx:+.1f}% sur 52 semaines. "
                    f"Cours actuel : {_cours_val} (YTD : {_ytd_val})."
                ]

            _hist_title = f"{indice} {_final_idx:+.1f}% sur 52S — cours {_cours_val}, YTD {_ytd_val}"
            _lecture_box(slide, _hist_title,
                         "  ".join(_narrative_parts), y_top=9.0, height=4.35, max_chars=900)
        else:
            # Pas de perf_history : tableau scores + note
            _header(slide, "Scores Sectoriels — Synthèse Allocataire",
                    f"{indice}  ·  Score composite FinSight  ·  Données yfinance",
                    active=5)
            _txb(slide, "Scores sectoriels FinSight — synthèse allocataire",
                 0.9, 2.1, 23.6, 0.55, size=8.5, bold=True, color=_NAVY)
            secteurs_etf = D.get("secteurs", [])
            sorted_etf   = sorted(secteurs_etf, key=lambda s: s[2], reverse=True)
            _SIG_ETF = {"Surpondérer":"Surpondérer","neutre":"Neutre","Sous-pondérer":"Sous-pondérer"}
            rows_etf = [["Rang","Secteur","Score","Signal","EV/EBITDA","Mg.EBITDA","Croiss."]]
            for rg, s in enumerate(sorted_etf, 1):
                sig_raw = str(s[3]).strip().lower()
                rows_etf.append([
                    str(rg), _abbrev_sector(s[0], 22), str(s[2]),
                    _SIG_ETF.get(sig_raw, str(s[3])),
                    str(s[4]),
                    f"{s[5]:.1f}%" if isinstance(s[5],(int,float)) and s[5] else "—",
                    str(s[6]) if len(s) > 6 else "—",
                ])
            tbl_etf = _add_table(slide, rows_etf, 0.9, 2.75, 23.6,
                                 min(8.5, 0.65 + len(rows_etf)*0.65),
                                 col_widths=[1.2, 5, 1.5, 3.5, 2.5, 2.5, 2.5],
                                 font_size=8, header_size=8, alt_fill=_GRAYL)
            for r in range(1, len(rows_etf)):
                sig_r = sorted_etf[r-1][3] if r-1 < len(sorted_etf) else "Neutre"
                _color_cell(tbl_etf, r, 3, _sig_light(sig_r), _sig_color(sig_r))
            _rect(slide, 0.9, 11.8, 23.6, 1.5, fill=_GRAYL)
            _rect(slide, 0.9, 11.8, 0.12, 1.5, fill=_GRAYD)
            _txb(slide, "Note — ETF sectoriels SPDR/iShares",
                 1.2, 11.9, 22.8, 0.5, size=8, bold=True, color=_GRAYT)
            _txb(slide,
                 "Les ETF sectoriels sont mapp\u00e9s via core/sector_etfs.py : "
                 "SPDR Select Sector pour le S&P 500/Nasdaq, iShares STOXX Europe 600 "
                 "pour CAC 40/DAX 40/FTSE 100 (UCITS, suffixe .DE sur Xetra). "
                 "Le fetch du pipeline indice n'est pas encore branche sur cette matrice "
                 "pour les indices europeens : l'analyse sectorielle directe reste "
                 "disponible via 'Analyse sectorielle' pour acceder aux holdings reels.",
                 1.2, 12.45, 22.8, 0.8, size=7.5, color=_GRAYT, wrap=True)
        _footer(slide)
        return slide

    _header(slide, "Performance des ETF Sectoriels — 52 Semaines",
            f"{indice}  ·  ETF sectoriels SPDR/iShares  ·  Indexe a 100 au debut de la période  ·  Données yfinance",
            active=5)

    _pic(slide, chart_bytes, 0.9, 2.3, 15.0, 11.1)

    # Panel droit legende
    etf_perf = D.get("etf_perf",{})
    _rect(slide, 16.4, 2.3, 8.1, 11.1, fill=_GRAYL)
    _rect(slide, 16.4, 2.3, 8.1, 0.65, fill=_NAVY)
    _txb(slide, "Legende ETF SPDR", 16.6, 2.35, 7.8, 0.55, size=8, bold=True, color=_WHITE)

    # Trier par return desc
    etf_list = sorted(etf_perf.items(), key=lambda x: x[1].get("return_1y",0), reverse=True)
    for i, (etf, info) in enumerate(etf_list[:8]):
        yy = 3.15 + i * 1.2
        col = _LINE_COLORS_RGB[i % len(_LINE_COLORS_RGB)]
        _rect(slide, 16.5, yy + 0.15, 0.5, 0.35, fill=col)
        nom = info.get("nom","")[:18]
        ret = info.get("return_1y", 0)
        ret_str = f"{ret:+.1f} %"
        _txb(slide, f"{etf}  —  {nom}", 17.2, yy,      7.1, 0.55,
             size=7.5, bold=True, color=_NAVY)
        _txb(slide, ret_str, 17.2, yy + 0.55, 7.1, 0.55,
             size=8, bold=True,
             color=_BUY if ret > 0 else _SELL)

    # Lecture analytique bas droit — enrichie "pourquoi ca bouge"
    _cours_s20   = D.get("cours", "—")
    _ytd_s20     = D.get("variation_ytd", "—")
    _erp_s20     = D.get("erp", "—")
    _bpa_s20     = D.get("bpa_growth", "—")
    _pe_s20      = D.get("pe_forward", "—")
    _phase_s20   = D.get("phase_cycle", "")
    _sig_s20     = D.get("signal_global", "Neutre")
    _cours_ctx   = f"Cours {indice} : {_cours_s20} (YTD : {_ytd_s20}). " if _cours_s20 not in ("—", "") else ""
    if etf_list:
        best_etf,  best_info  = etf_list[0]
        worst_etf, worst_info = etf_list[-1]
        _best_nom   = best_info.get("nom",  "")[:16]
        _worst_nom  = worst_info.get("nom", "")[:16]
        _best_ret   = best_info.get("return_1y", 0)
        _worst_ret  = worst_info.get("return_1y", 0)
        _dispersion = abs(_best_ret - _worst_ret)
        # Pourquoi ca bouge — macro context
        _why_best = (
            "porté par la révision haussière des BPA et la réouverture des marges Opérationnelles"
            if _best_ret > 20 else
            "beneficiant du repositionnément des flux en fin de cycle"
        )
        _why_worst = (
            "pénalisé par la hausse des taux longs, comprimant les multiples de valorisation"
            if _worst_ret < -10 else
            "en sous-performance relative face aux secteurs a forte visibilité BPA"
        )
        _rate_ctx = (
            f"Dans un environnement ERP {_erp_s20}, les secteurs a duration longue subissent "
            f"la concurrence des obligations — la sélectivité sectorielle est cle."
            if isinstance(_erp_s20, str) and _erp_s20 not in ("—", "") else
            f"La croissance BPA consensus de {_bpa_s20} justifié un P/E Forward de {_pe_s20}."
        )
        _disp_comment = (
            f"Dispersion {_dispersion:.0f}pt entre extremes — fort levier d'alpha sectoriel."
            if _dispersion > 15 else
            f"Dispersion modeste ({_dispersion:.0f}pt) — Marché peu discriminant, rotation prudente."
        )
        commentary = (
            f"{_cours_ctx}"
            f"{best_etf} ({_best_nom}) en Tête (+{_best_ret:.1f}%) {_why_best}. "
            f"{worst_etf} ({_worst_nom}) en retrait ({_worst_ret:+.1f}%) {_why_worst}. "
            f"{_rate_ctx} "
            f"{_disp_comment} "
            f"Signal global FinSight : {_sig_s20}."
        )
    else:
        commentary = (
            f"{_cours_ctx}"
            f"Performance relative des ETF SPDR sur 52 semaines — source yfinance. "
            f"BPA consensus : {_bpa_s20} | P/E Forward : {_pe_s20} | ERP : {_erp_s20}. "
            f"Signal global FinSight : {_sig_s20}."
        )

    # Enrichissement LLM avec événements datés impactant le cours
    # Refonte 2026-04-14 : llm_call(phase=long) -> Mistral primary + prompt etendu
    _llm_macro_comment = ""
    try:
        from core.llm_provider import llm_call
        _llm_prompt_s20 = (
            f"Tu es analyste macro sell-side senior. Redige un commentaire approfondi "
            f"(300-360 mots) sur les facteurs macro et catalyseurs sectoriels ayant "
            f"impacte le cours de l'indice {indice} sur les 12 derniers mois.\n\n"
            f"Donnees : YTD {_ytd_s20}, ERP {_erp_s20}, P/E Forward {_pe_s20}, "
            f"BPA growth {_bpa_s20}, phase cycle {_phase_s20}.\n\n"
            f"Structure en 3 paragraphes (120 mots chacun) :\n"
            f"1. Environnement macro dates : politique monetaire (Fed/BCE avec dates "
            f"cles), inflation, taux longs, cycle economique. Cite les evenements "
            f"datables (pivot Fed, reunion BCE, publication ISM, CPI) et leur impact "
            f"chiffre sur l'indice.\n"
            f"2. Catalyseurs sectoriels specifiques : quels secteurs ont porte/freine "
            f"la performance, revisions de consensus, publications trimestrielles "
            f"marquantes, M&A de reference.\n"
            f"3. Implications forward : catalyseurs a surveiller sur les 6 prochains "
            f"mois, niveaux techniques critiques, conditions de bascule du regime.\n\n"
            f"Francais correct avec accents. Pas de markdown. Pas d'emojis. "
            f"Cite des chiffres precis et des dates."
        )
        _llm_macro_comment = llm_call(_llm_prompt_s20, phase="long", max_tokens=900) or ""
    except Exception as _llm_e:
        log.warning("[indice_pptx] s20 LLM macro: %s", _llm_e)

    if _llm_macro_comment.strip():
        commentary = _llm_macro_comment.strip()

    # Box elargie pour accueillir le texte LLM etendu (1.75 -> 2.60 cm)
    _rect(slide, 16.4, 11.0, 8.1, 0.05, fill=_GRAYD)
    _txb(slide, "Pourquoi ca bouge — environnement macro et catalyseurs",
         16.6, 11.1, 7.8, 0.5, size=7.5, bold=True, color=_NAVY)
    _txb(slide, commentary[:2000], 16.6, 11.65, 7.8, 2.15,
         size=6.8, color=_GRAYT, wrap=True)

    _footer(slide)
    return slide


def _s21_disclaimer(prs, D):
    """Slide 21 — Avertissement légal et méthodologie."""
    slide = _blank(prs)
    _header(slide, "Avertissement & Méthodologie",
            "Informations réglementaires et sources de données", 4)

    indice = D.get("indice", "Indice")
    date_str = D.get("date_analyse", "")
    nb_soc = D.get("nb_societes", 0)
    nb_sec = D.get("nb_secteurs", 0)

    # Avertissement légal
    _rect(slide, 0.9, 2.10, 23.6, 2.80, _GRAYL)
    _rect(slide, 0.9, 2.10, 0.12, 2.80, _SELL)
    _txb(slide, "AVERTISSEMENT LÉGAL", 1.15, 2.15, 23.2, 0.35,
         size=9, bold=True, color=_SELL)
    disc_text = (
        f"Ce document a été généré automatiquement par FinSight IA v1.0 le {date_str}. "
        f"Il est produit intégralement par un système d'intelligence artificielle et "
        f"ne constitue pas un conseil en investissement au sens de la directive européenne "
        f"MiFID II (2014/65/UE). FinSight IA n'est pas un prestataire de services "
        f"d'investissement agréé. Tout investisseur est invité à procéder à sa propre "
        f"diligence et à consulter un professionnel qualifié avant toute décision. "
        f"Les performances passées ne préjugent pas des performances futures."
    )
    _txb(slide, disc_text, 1.15, 2.55, 23.2, 2.30,
         size=7.5, color=_GRAYT, wrap=True)

    # Méthodologie
    _rect(slide, 0.9, 5.20, 23.6, 0.40, _NAVY)
    _txb(slide, "MÉTHODOLOGIE & SOURCES", 1.05, 5.22, 23.3, 0.35,
         size=8.5, bold=True, color=_WHITE)

    metho_items = [
        (f"Univers", f"{indice} — {nb_soc} sociétés, {nb_sec} secteurs GICS"),
        ("Score composite", "Moyenne pondérée : 40% momentum prix 3M, 30% révision BPA, 30% valorisation relative"),
        ("Signal", "Surpondérer (score >= 60) · Neutre (40-60) · Sous-pondérer (< 40)"),
        ("Allocation Markowitz", "Optimisation mean-variance sur rendements ETF SPDR 52S. Contrainte max 40% par secteur."),
        ("ERP sectoriel", "Earnings Yield (1/PE forward) - Taux 10Y US. Seuils : Tendu < 2%, Neutre 2-4%, Favorable > 4%."),
        ("Rotation sectorielle", "Modèle 4 phases (Expansion, Ralentissement, Récession, Reprise). ISM, courbe taux, Leading indicators."),
        ("Sentiment", "FinBERT (ProsusAI/finbert) sur articles Finnhub 7 jours. Score agrégé par secteur."),
        ("Biais de cadrage", "Score FinSight composite, pas de backtest performance. Les ratios utilisés (momentum/révisions BPA/valorisation) sont de natures différentes et leur pondération peut être débattue."),
        ("Limites sectorielles", "Les sous-secteurs (banques vs assurance, E&P vs raffinage) ne sont pas traités indépendamment — le profil dominant de chaque secteur GICS guide l'analyse."),
        ("Horizon d'allocation", "12 mois glissants. Les signaux sont révisés mensuellement selon la dynamique des ratios et l'évolution du contexte macro (réunions BCE/Fed, publications trimestrielles)."),
        ("Sources données", "yfinance (cours, fondamentaux), FMP (multiples), Finnhub (news), ETF SPDR (sectoriels US)."),
    ]
    y = 5.80
    # Spacing elargi (0.42 -> 0.62) pour remplir l'espace vertical disponible
    for label, desc in metho_items:
        _txb(slide, f"{label} :", 1.05, y, 5.0, 0.40, size=7.5, bold=True, color=_NAVY)
        _txb(slide, desc, 6.05, y, 18.4, 0.55, size=7.5, color=_GRAYT, wrap=True)
        y += 0.62

    # Confidentialité
    _txb(slide, "Document confidentiel — Diffusion restreinte. © 2026 FinSight IA.",
         0.9, 12.60, 23.6, 0.40, size=7, bold=True, color=_NAVY,
         align=PP_ALIGN.CENTER)

    _footer(slide)


# ── Classe principale ──────────────────────────────────────────────────────────

class IndicePPTXWriter:

    @staticmethod
    def generate(data: dict, output_path: Optional[str] = None) -> bytes:
        """
        Genere un pitchbook PPTX 20 slides pour l'analyse d'indice.
        Retourne les bytes du fichier PPTX.
        Si output_path fourni, sauvegarde aussi sur disque.
        """
        log.info("IndicePPTXWriter: Génération pour %s", data.get("indice","—"))

        # Macro regime_v (si pas déjà calculé)
        if not data.get("macro"):
            try:
                import sys as _sys, os as _os
                _sys.path.insert(0, _os.path.dirname(_os.path.dirname(__file__)))
                from agents.agent_macro import AgentMacro
                data["macro"] = AgentMacro().analyze()
            except Exception as _me:
                log.warning("IndicePPTXWriter: AgentMacro: %s", _me)
                data.setdefault("macro", {})

        # Pre-calcul des charts (peut echouer silencieusement)
        try:
            secteurs = data.get("secteurs",[])
            scatter_bytes   = _chart_scatter(secteurs)
            score_bytes     = _chart_score_bars(secteurs)
            ev_bytes        = _chart_ev_distribution(secteurs)
            zone_bytes      = _chart_zone_entree(data)
            sent_bytes      = _chart_sentiment_bars(data.get("sentiment_agg",{}))
            etf_bytes       = _chart_etf_perf(data)
        except Exception as e:
            log.warning("IndicePPTXWriter: erreur pre-calcul charts: %s", e)
            # Fallback: image vide
            _empty = _make_empty_chart()
            scatter_bytes = score_bytes = ev_bytes = zone_bytes = sent_bytes = etf_bytes = _empty

        prs = Presentation()
        prs.slide_width  = _SW
        prs.slide_height = _SH

        def _safe(fn, *args, label="slide"):
            try:
                fn(*args)
            except Exception as _e:
                log.warning("IndicePPTXWriter: %s: %s", label, _e)
                try:
                    sl = prs.slides.add_slide(prs.slide_layouts[6])
                    _txb(sl, label, Cm(0.5), Cm(0.5), Cm(8), Cm(1), 11, bold=False)
                except Exception:
                    pass

        # Slide 1 — Cover
        _safe(_s01_cover, prs, data, label="s01_cover")
        # Slide 2 — Executive Summary
        _safe(_s02_exec_summary, prs, data, label="s02_exec_summary")
        # Slide 3 — Sommaire
        _safe(_s03_sommaire, prs, data, label="s03_sommaire")
        # Slide 4 — Chapter 01 divider
        _safe(_chapter_divider, prs, "01", "Synth\u00e8se Macro & Signal Global",
              "Valorisation top-down \u00b7 ERP \u00b7 Cycle \u00e9conomique \u00b7 Catalyseurs & risques",
              label="s04_divider")
        # Slide 5 — Description de l'Indice
        _safe(_s05_description, prs, data, label="s05_description")
        # Slide 6 — Valorisation Macro Top-Down
        _safe(_s06_valorisation, prs, data, label="s06_valorisation")
        # Slide 7 — Positionnement dans le Cycle
        _safe(_s07_cycle, prs, data, label="s07_cycle")
        # Slide 8 — Chapter 02 divider
        _safe(_chapter_divider, prs, "02", "Cartographie des Secteurs",
              f"{data.get('nb_secteurs',11)} secteurs GICS \u00b7 Scores \u00b7 Scatter valorisation \u00b7 D\u00e9composition",
              label="s08_divider")
        # Slide 9 — Cartographie des Secteurs
        _safe(_s09_cartographie, prs, data, label="s09_cartographie")
        # Slide 10 — Valorisation vs Croissance BPA
        _safe(_s10_scatter, prs, data, scatter_bytes, label="s10_scatter")
        # Slide 11 — Décomposition des Scores
        _safe(_s11_decomposition, prs, data, label="s11_Décomposition")
        # Slide 12 — Chapter 03 divider
        _safe(_chapter_divider, prs, "03", "Top 3 Secteurs Recommand\u00e9s",
              "Synth\u00e8se signal \u00b7 Soci\u00e9t\u00e9s repr\u00e9sentatives \u00b7 Zone d'entr\u00e9e \u00b7 Allocation Markowitz",
              label="s12_divider")
        # Slide 13 — Top 3 Secteurs
        _safe(_s13_top3, prs, data, label="s13_top3")
        # Slide 14 — Allocation Optimale Markowitz
        _safe(_s14_allocation, prs, data, label="s14_allocation")
        # Slide 15 — Zone d'Entree
        _safe(_s15_zone_entree, prs, data, zone_bytes, label="s15_zone_entree")
        # Slide 16 — Chapter 04 divider
        _safe(_chapter_divider, prs, "04", "Risques, Rotation & Sentiment",
              "Sc\u00e9narios alternatifs \u00b7 Rotation sectorielle \u00b7 FinBERT \u00b7 Performance ETF",
              label="s16_divider")
        # Slide 17 — Risques Macro & Scénarios
        _safe(_s17_risques, prs, data, label="s17_risques")
        # Slide 18 — Rotation Sectorielle
        _safe(_s18_rotation, prs, data, label="s18_rotation")
        # Slide 19 — Sentiment FinBERT
        _safe(_s19_sentiment, prs, data, sent_bytes, label="s19_sentiment")
        # Slide 20 — Performance ETF
        _safe(_s20_etf_perf, prs, data, etf_bytes, label="s20_etf_perf")
        # Slide 21 — Avertissement & Méthodologie
        _safe(_s21_disclaimer, prs, data, label="s21_disclaimer")

        buf = io.BytesIO()
        prs.save(buf)
        raw = buf.getvalue()

        if output_path:
            import os
            os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
            with open(output_path, 'wb') as f:
                f.write(raw)
            log.info("IndicePPTXWriter: sauvegarde -> %s (%d Ko)", output_path, len(raw)//1024)

        return raw


def _make_empty_chart() -> bytes:
    """Graphique vide de secours."""
    fig, ax = plt.subplots(figsize=(8, 4))
    ax.text(0.5, 0.5, "Graphique non disponible", ha='center', va='center',
            transform=ax.transAxes, color='#999999', fontsize=12)
    ax.axis('off')
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=100, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    buf.seek(0)
    return buf.read()
