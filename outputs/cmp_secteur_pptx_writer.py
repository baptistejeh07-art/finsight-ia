"""
cmp_secteur_pptx_writer.py — FinSight IA
Pitchbook comparatif sectoriel 20 slides via python-pptx.
Compare deux secteurs cote a cote (memes ratios, scoring, top acteurs, risques).

Usage :
    from outputs.cmp_secteur_pptx_writer import CmpSecteurPPTXWriter
    bytes_ = CmpSecteurPPTXWriter.generate(
        tickers_a, sector_a, universe_a,
        tickers_b, sector_b, universe_b,
        output_path="cmp_tech_vs_sante.pptx"
    )

tickers_a / tickers_b : list[dict] tel que retourne par _fetch_real_sector_data
    Cles utiles : ticker, company, score_global, score_value, score_growth,
                  score_quality, score_momentum, pe_ratio, ev_ebitda, ev_revenue,
                  ebitda_margin, gross_margin, net_margin, roe, roic,
                  revenue_growth, momentum_52w, altman_z, piotroski_f,
                  peg_ratio, fcf_yield, beta, price, market_cap, currency
"""
from __future__ import annotations

import io
import logging
import math
from datetime import datetime
from pathlib import Path
from typing import Optional

import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.dates as mdates

log = logging.getLogger(__name__)

from pptx import Presentation
from pptx.util import Cm, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ──────────────────────────────────────────────────────────────────
_NAVY   = RGBColor(0x1B, 0x3A, 0x6B)
_NAVYL  = RGBColor(0x2A, 0x52, 0x98)
_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
_BLACK  = RGBColor(0x1A, 0x1A, 0x1A)
_GRAYL  = RGBColor(0xF5, 0xF7, 0xFA)
_GRAYM  = RGBColor(0xE8, 0xEC, 0xF0)
_GRAYT  = RGBColor(0x55, 0x55, 0x55)
_GRAYD  = RGBColor(0xAA, 0xAA, 0xAA)
_GREEN  = RGBColor(0x1A, 0x7A, 0x4A)
_RED    = RGBColor(0xA8, 0x20, 0x20)
_AMBER  = RGBColor(0xB8, 0x92, 0x2A)

# Couleurs differenciantes A vs B
_COL_A      = RGBColor(0x1B, 0x3A, 0x6B)   # navy   → Secteur A
_COL_A_PALE = RGBColor(0xEE, 0xF3, 0xFA)
_COL_B      = RGBColor(0x1A, 0x7A, 0x4A)   # vert   → Secteur B
_COL_B_PALE = RGBColor(0xEA, 0xF4, 0xEF)
_GOLD       = RGBColor(0xC9, 0xA2, 0x27)   # or     → Secteur B (THESE slides)
_GOLD_PALE  = RGBColor(0xFD, 0xF6, 0xE3)

# ── Dimensions ───────────────────────────────────────────────────────────────
_SW = Inches(10.0)
_SH = Inches(5.625)

# Matplotlib hex (pour les charts)
_HEX_A = "#1B3A6B"
_HEX_B = "#1A7A4A"
_HEX_AG = "#5b7fb5"
_HEX_BG = "#5aab7a"


# ── Utilitaires safe ─────────────────────────────────────────────────────────
def _safe_float(v, default=0.0) -> float:
    try:
        f = float(v)
        return f if math.isfinite(f) else default
    except Exception:
        return default


def _med(vals: list, default=0.0) -> float:
    clean = [_safe_float(v) for v in vals if v is not None]
    if not clean:
        return default
    clean.sort()
    n = len(clean)
    return clean[n // 2] if n % 2 else (clean[n // 2 - 1] + clean[n // 2]) / 2


def _mean(vals: list, default=0.0) -> float:
    clean = [_safe_float(v) for v in vals if v is not None]
    return sum(clean) / len(clean) if clean else default


# ── Helpers de cadrage standardisé (style cmp société) ──────────────────────
_NAVY_PALE = RGBColor(0xEE, 0xF3, 0xFA)
_NAVY_MID  = RGBColor(0x2E, 0x5F, 0xA3)


def _llm_box(slide, x: float, y: float, w: float, h: float,
             title: str, text: str, *, fontsize: float = 9.0):
    """Box LLM standardisée : fond bleu pale + barre laterale navy + titre + texte.

    Reproduit le pattern utilise dans comparison_pptx_writer.py (cmp societe).
    """
    _rect(slide, x, y, w, h, fill=_NAVY_PALE)
    _rect(slide, x, y, 0.13, h, fill=_NAVY_MID)
    _txb(slide, title.upper(), x + 0.30, y + 0.10, w - 0.40, 0.36,
         size=fontsize, bold=True, color=_NAVY)
    _txb(slide, text or "", x + 0.30, y + 0.50, w - 0.50, h - 0.62,
         size=fontsize, color=_NAVY, wrap=True)


def _bullet_square(slide, x: float, y: float, color, size: float = 0.22):
    """Petit carre colore servant de bullet pour paragraphes (slides 15/16/18)."""
    _rect(slide, x, y, size, size, fill=color)


def _llm_box_with_bullets(slide, x: float, y: float, w: float, h: float,
                          title: str, items: list, *, color, fontsize: float = 9.0):
    """Box LLM avec une liste a puces carrees colorees a gauche.

    items : liste de tuples (heading, body) ou de strings.
    """
    _rect(slide, x, y, w, h, fill=_NAVY_PALE)
    _rect(slide, x, y, 0.13, h, fill=_NAVY_MID)
    _txb(slide, title.upper(), x + 0.30, y + 0.10, w - 0.40, 0.36,
         size=fontsize, bold=True, color=_NAVY)
    if not items:
        return
    n = max(1, len(items))
    avail = h - 0.65
    item_h = avail / n
    for i, it in enumerate(items[:5]):
        iy = y + 0.55 + i * item_h
        _bullet_square(slide, x + 0.35, iy + 0.12, color, size=0.20)
        if isinstance(it, tuple) and len(it) == 2:
            head, body = it
            _txb(slide, str(head)[:60], x + 0.70, iy, w - 0.85, 0.36,
                 size=fontsize, bold=True, color=_NAVY, wrap=False)
            _txb(slide, str(body), x + 0.70, iy + 0.36, w - 0.85, item_h - 0.40,
                 size=fontsize - 0.5, color=_NAVY, wrap=True)
        else:
            _txb(slide, str(it), x + 0.70, iy + 0.05, w - 0.85, item_h - 0.10,
                 size=fontsize, color=_NAVY, wrap=True)


# ── PPTX primitives ──────────────────────────────────────────────────────────
def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, x, y, w, h, fill=None, line=False, line_col=None, line_w=0.5):
    shape = slide.shapes.add_shape(1, Cm(x), Cm(y), Cm(w), Cm(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line and line_col:
        shape.line.color.rgb = line_col
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape


def _txb(slide, text, x, y, w, h, size=9, bold=False, color=None, align=PP_ALIGN.LEFT,
         italic=False, wrap=True):
    color = color or _BLACK
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def _pic(slide, img_bytes, x, y, w, h):
    buf = io.BytesIO(img_bytes)
    slide.shapes.add_picture(buf, Cm(x), Cm(y), Cm(w), Cm(h))


def _add_table(slide, data, x, y, w, h, col_widths=None,
               header_fill=_NAVY, header_color=_WHITE,
               alt_fill=None, font_size=7.5, header_size=7.5):
    rows = len(data)
    cols = len(data[0]) if data else 1
    tbl_shape = slide.shapes.add_table(rows, cols, Cm(x), Cm(y), Cm(w), Cm(h))
    tbl = tbl_shape.table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Cm(cw * w / total)
    for r_idx, row_data in enumerate(data):
        is_header = (r_idx == 0)
        for c_idx, cell_text in enumerate(row_data):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = str(cell_text) if cell_text is not None else "—"
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            run = para.runs[0] if para.runs else para.add_run()
            run.font.size = Pt(header_size if is_header else font_size)
            run.font.bold = is_header
            if is_header:
                run.font.color.rgb = header_color
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
            else:
                run.font.color.rgb = _BLACK
                if alt_fill and r_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = alt_fill
                else:
                    cell.fill.background()
    return tbl


# ── Header / footer / chapter divider ────────────────────────────────────────
_NAV_LABELS = ["1", "2", "3", "4", "5"]


def _header(slide, title, subtitle, active_section=1):
    _rect(slide, 0, 0, 25.4, 1.4, fill=_NAVY)
    _txb(slide, title, 0.9, 0.05, 19.0, 1.3, size=12, bold=True, color=_WHITE)
    for i, lbl in enumerate(_NAV_LABELS):
        dot_x = 21.8 + i * 0.7
        fill = _WHITE if (i + 1) == active_section else _NAVYL
        _rect(slide, dot_x, 0.35, 0.55, 0.55, fill=fill)
        c = _BLACK if (i + 1) == active_section else _GRAYD
        _txb(slide, lbl, dot_x, 0.35, 0.55, 0.55, size=7, bold=True, color=c, align=PP_ALIGN.CENTER)
    _txb(slide, subtitle, 0.9, 1.55, 23.6, 0.55, size=7.5, color=_GRAYT)


def _footer(slide, D):
    _rect(slide, 0, 13.75, 25.4, 0.5, fill=_GRAYL)
    _txb(slide, "FinSight IA  |  Comparatif Sectoriel  |  Usage confidentiel", 0.9, 13.8, 14.0, 0.4, size=7, color=_GRAYD)
    _txb(slide, D["date_str"], 0, 13.8, 24.4, 0.4, size=7, color=_GRAYD, align=PP_ALIGN.RIGHT)


def _chapter_divider(prs, num_str, chapter_title, subtitle):
    slide = _blank(prs)
    _rect(slide, 0, 0, 25.4, 14.3, fill=_NAVY)
    _txb(slide, num_str, 1.0, 3.5, 8.0, 4.5, size=72, bold=True, color=_NAVYL)
    _txb(slide, chapter_title, 1.0, 7.0, 23.0, 2.0, size=28, bold=True, color=_WHITE)
    _rect(slide, 1.0, 9.1, 15.0, 0.05, fill=_GRAYD)
    _txb(slide, subtitle, 1.0, 9.4, 22.9, 0.8, size=11, color=_GRAYD)
    _txb(slide, "FinSight IA  |  Usage confidentiel", 0.9, 13.75, 23.6, 0.4, size=7, color=_GRAYD)


# ── Legend badge ─────────────────────────────────────────────────────────────
def _legend_badge(slide, x, y):
    _rect(slide, x, y, 0.4, 0.22, fill=_COL_A)
    _txb(slide, "A", x + 0.45, y, 1.5, 0.22, size=7, color=_COL_A)
    _rect(slide, x + 2.0, y, 0.4, 0.22, fill=_COL_B)
    _txb(slide, "B", x + 2.45, y, 1.5, 0.22, size=7, color=_COL_B)


# ── Preparation des données ───────────────────────────────────────────────────
def _prepare_data(tickers_a, sector_a, universe_a, tickers_b, sector_b, universe_b):
    """Calcule les Médianes et Agrégats pour les deux secteurs."""
    def _stats(td):
        if not td:
            return {}
        pe     = _med([t.get("pe_ratio") for t in td if t.get("pe_ratio")])
        ev_eb  = _med([t.get("ev_ebitda") for t in td if t.get("ev_ebitda")])
        ev_rev = _med([t.get("ev_revenue") for t in td if t.get("ev_revenue")])
        gm     = _med([t.get("gross_margin") for t in td])
        em     = _med([t.get("ebitda_margin") for t in td])
        nm     = _med([t.get("net_margin") for t in td])
        roe    = _med([t.get("roe") for t in td])
        roic   = _med([t.get("roic") for t in td if t.get("roic")], default=None)
        _rg_vals = [v for v in [t.get("revenue_growth") for t in td] if v is not None and abs(float(v)) <= 2.0]
        revg   = _med(_rg_vals if _rg_vals else [0.0]) * 100
        mom    = _med([t.get("momentum_52w", 0) for t in td])
        beta   = _med([t.get("beta", 1.0) for t in td if t.get("beta")])
        score  = int(_mean([t.get("score_global", 50) for t in td]))
        # score_value peut être en 0-25 (fetch_real) ou 0-100 (compute_screening)
        # Normaliser vers 0-25 : si la valeur > 30, c'est du 0-100 -> diviser par 4
        def _sub(t, key, default=12):
            v = t.get(key)
            if v is None:
                return default
            return round(float(v) / 4.0, 1) if float(v) > 30 else float(v)
        s_val  = _mean([_sub(t, "score_value")    for t in td])
        s_gro  = _mean([_sub(t, "score_growth")   for t in td])
        s_qua  = _mean([_sub(t, "score_quality")  for t in td])
        s_mom  = _mean([_sub(t, "score_momentum") for t in td])
        pf     = _med([t.get("piotroski_f") for t in td if t.get("piotroski_f") is not None])
        az     = _med([t.get("altman_z") for t in td if t.get("altman_z") is not None], default=None)
        fcfy   = _med([t.get("fcf_yield") for t in td if t.get("fcf_yield") is not None])
        peg    = _med([t.get("peg_ratio") for t in td if t.get("peg_ratio")])
        div_yield = _med([t.get("div_yield") for t in td if t.get("div_yield") and float(t.get("div_yield") or 0) > 0], default=None)
        payout = _med([t.get("payout_ratio") for t in td if t.get("payout_ratio") and float(t.get("payout_ratio") or 0) > 0], default=None)
        return dict(
            pe=pe, ev_eb=ev_eb, ev_rev=ev_rev, gm=gm, em=em, nm=nm,
            roe=roe, roic=roic, revg=revg, mom=mom, beta=beta,
            score=score, s_val=s_val, s_gro=s_gro, s_qua=s_qua, s_mom=s_mom,
            pf=pf, az=az, fcfy=fcfy, peg=peg, div_yield=div_yield, payout=payout,
        )

    sa = _stats(tickers_a)
    sb = _stats(tickers_b)

    # Signal couleur selon score
    def _sig(score):
        if score >= 65:
            return _GREEN, "Surpondérer"
        if score >= 45:
            return _AMBER, "Neutre"
        return _RED, "Sous-pondérer"

    sig_a_col, sig_a_lbl = _sig(sa.get("score", 50))
    sig_b_col, sig_b_lbl = _sig(sb.get("score", 50))

    # Universe commun ou combine
    universes = set(filter(None, [universe_a, universe_b]))
    universe_label = " / ".join(sorted(universes)) if universe_a != universe_b else (universe_a or "Global")

    _MOIS_FR = ["janvier","fevrier","mars","avril","mai","juin",
                "juillet","aout","septembre","octobre","novembre","decembre"]
    _now = datetime.now()
    _date_str = f"{_now.day} {_MOIS_FR[_now.month - 1]} {_now.year}"

    return dict(
        sector_a=sector_a, sector_b=sector_b,
        universe_a=universe_a, universe_b=universe_b,
        universe_label=universe_label,
        na=len(tickers_a), nb=len(tickers_b),
        td_a=tickers_a, td_b=tickers_b,
        sa=sa, sb=sb,
        sig_a_col=sig_a_col, sig_a_lbl=sig_a_lbl,
        sig_b_col=sig_b_col, sig_b_lbl=sig_b_lbl,
        date_str=_date_str,
        year=_now.year,
        llm={},
    )


# ── Helpers de formatage ──────────────────────────────────────────────────────
def _fmt(v, pct=False, x=False, dp=1):
    """Formate une valeur numerique pour affichage."""
    if v is None or (isinstance(v, float) and not math.isfinite(v)):
        return "—"
    try:
        f = float(v)
        if pct:
            return f"{f:+.{dp}f} %"
        if x:
            return f"{f:.{dp}f}x"
        return f"{f:.{dp}f}"
    except Exception:
        return str(v)


def _fmt_simple(v, pct=False, x=False, dp=1):
    """Formate sans signe +."""
    if v is None or (isinstance(v, float) and not math.isfinite(v)):
        return "—"
    try:
        f = float(v)
        if pct:
            return f"{f:.{dp}f} %"
        if x:
            return f"{f:.{dp}f}x"
        return f"{f:.{dp}f}"
    except Exception:
        return str(v)


# ── Contenu sectoriel (repris de sectoral_pptx_writer) ───────────────────────
from outputs.sectoral_pptx_writer import _SECTOR_CONTENT


def _get_content(sector_name: str) -> dict:
    """Cherche le contenu dans la librairie (insensible a la casse / partiel)."""
    key = sector_name.strip()
    if key in _SECTOR_CONTENT:
        return _SECTOR_CONTENT[key]
    key_lc = key.lower()
    for k, v in _SECTOR_CONTENT.items():
        if k.lower() == key_lc or key_lc in k.lower() or k.lower() in key_lc:
            return v
    # Fallback generique
    return {
        "description": f"Le secteur {sector_name} regroupe les sociétés actives dans ce domaine. "
                       f"Analyse basée sur les données financières réelles des principaux acteurs.",
        "catalyseurs": [
            ("Croissance Structurelle", "Demande sectorielle portée par des tendances long terme favorables."),
            ("Transformation Digitale", "Adoption technologique acceleree — gains de productivite et nouveaux modèles."),
            ("Consolidation", "Fusions-acquisitions — Économies d'échelle et elimination des acteurs faibles."),
        ],
        "risques": [
            ("Ralentissement Macro", "Récession potentielle — contraction de la demande et compression des marges."),
            ("Pression Concurrentielle", "Nouveaux entrants et disruption — erosion des parts de marché."),
            ("Risque Réglementaire", "Évolutions réglementaires — Surcoûts de Conformité."),
        ],
        "drivers": [
            ("up", "Demande Structurelle", "Croissance portée par megatendances"),
            ("up", "Innovation", "R&D et transformation technologique"),
            ("down", "Sensibilité Macro", "Exposition aux cycles Économiques"),
            ("down", "Pression Concurrence", "Nouveaux entrants et disruption"),
        ],
        "cycle_comment": "Positionnement cycle a déterminer selon contexte macro",
        "metriques": [],
        "conditions": [],
    }


# ── Chart builders ────────────────────────────────────────────────────────────

def _chart_valuation(sa: dict, sb: dict, label_a: str, label_b: str) -> bytes:
    """Barres groupees : P/E, EV/EBITDA, EV/Revenue."""
    metrics = ["P/E", "EV/EBITDA", "EV/Revenue"]
    vals_a  = [sa.get("pe", 0) or 0, sa.get("ev_eb", 0) or 0, sa.get("ev_rev", 0) or 0]
    vals_b  = [sb.get("pe", 0) or 0, sb.get("ev_eb", 0) or 0, sb.get("ev_rev", 0) or 0]

    fig, ax = plt.subplots(figsize=(7.5, 4.2))
    x = np.arange(len(metrics))
    w = 0.35
    bars_a = ax.bar(x - w/2, vals_a, w, label=label_a, color=_HEX_A, alpha=0.9)
    bars_b = ax.bar(x + w/2, vals_b, w, label=label_b, color=_HEX_B, alpha=0.9)
    ax.set_xticks(x)
    ax.set_xticklabels(metrics, fontsize=10)
    ax.set_ylabel("Multiple (x)", fontsize=9)
    ax.tick_params(axis="y", labelsize=8)
    ax.yaxis.grid(True, linestyle="--", alpha=0.4)
    ax.set_axisbelow(True)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    for bar in list(bars_a) + list(bars_b):
        h = bar.get_height()
        if h > 0:
            ax.annotate(f"{h:.1f}x", xy=(bar.get_x() + bar.get_width()/2, h),
                        xytext=(0, 3), textcoords="offset points",
                        ha="center", va="bottom", fontsize=8)
    ax.legend(fontsize=9, framealpha=0)
    plt.tight_layout(pad=0.5)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=140, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def _chart_margins(sa: dict, sb: dict, label_a: str, label_b: str) -> bytes:
    """Barres groupees : Marge Brute, EBITDA, Nette."""
    metrics = ["Marge Brute", "Marge EBITDA", "Marge Nette"]
    vals_a  = [sa.get("gm", 0) or 0, sa.get("em", 0) or 0, sa.get("nm", 0) or 0]
    vals_b  = [sb.get("gm", 0) or 0, sb.get("em", 0) or 0, sb.get("nm", 0) or 0]

    fig, ax = plt.subplots(figsize=(7.5, 4.2))
    x = np.arange(len(metrics))
    w = 0.35
    bars_a = ax.bar(x - w/2, vals_a, w, label=label_a, color=_HEX_A, alpha=0.9)
    bars_b = ax.bar(x + w/2, vals_b, w, label=label_b, color=_HEX_B, alpha=0.9)
    ax.set_xticks(x)
    ax.set_xticklabels(metrics, fontsize=10)
    ax.set_ylabel("Marge (%)", fontsize=9)
    ax.tick_params(axis="y", labelsize=8)
    ax.yaxis.grid(True, linestyle="--", alpha=0.4)
    ax.set_axisbelow(True)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    for bar in list(bars_a) + list(bars_b):
        h = bar.get_height()
        if abs(h) > 0.5:
            ax.annotate(f"{h:.1f}%", xy=(bar.get_x() + bar.get_width()/2, max(h, 0)),
                        xytext=(0, 3), textcoords="offset points",
                        ha="center", va="bottom", fontsize=7.5)
    ax.legend(fontsize=9, framealpha=0)
    plt.tight_layout(pad=0.5)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=140, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def _chart_radar(sa: dict, sb: dict, label_a: str, label_b: str) -> bytes:
    """Radar chart scoring 4 axes : Value, Growth, Quality, Momentum."""
    categories = ["Value", "Growth", "Quality", "Momentum"]
    vals_a = [sa.get("s_val", 12), sa.get("s_gro", 12), sa.get("s_qua", 12), sa.get("s_mom", 12)]
    vals_b = [sb.get("s_val", 12), sb.get("s_gro", 12), sb.get("s_qua", 12), sb.get("s_mom", 12)]

    N = len(categories)
    angles = [n / float(N) * 2 * math.pi for n in range(N)]
    angles += angles[:1]
    va = vals_a + vals_a[:1]
    vb = vals_b + vals_b[:1]

    fig, ax = plt.subplots(figsize=(4.5, 4.5), subplot_kw=dict(polar=True))
    ax.plot(angles, va, "o-", linewidth=2, color=_HEX_A, label=label_a)
    ax.fill(angles, va, alpha=0.15, color=_HEX_A)
    ax.plot(angles, vb, "s-", linewidth=2, color=_HEX_B, label=label_b)
    ax.fill(angles, vb, alpha=0.15, color=_HEX_B)
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(categories, fontsize=10)
    ax.set_ylim(0, 25)
    ax.set_yticks([5, 10, 15, 20, 25])
    ax.set_yticklabels(["5", "10", "15", "20", "25"], fontsize=7, color="grey")
    ax.grid(color="grey", linestyle="--", linewidth=0.5, alpha=0.4)
    ax.legend(loc="upper right", bbox_to_anchor=(1.35, 1.1), fontsize=9)
    plt.tight_layout(pad=0.3)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=140, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def _chart_momentum(sa: dict, sb: dict, label_a: str, label_b: str) -> bytes:
    """Barres : Croissance Revenue (%) et Momentum 52S (%)."""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(7.5, 3.0))

    # Revenue growth
    revg_a = sa.get("revg", 0) or 0
    revg_b = sb.get("revg", 0) or 0
    bars1 = ax1.bar([label_a, label_b], [revg_a, revg_b], color=[_HEX_A, _HEX_B], alpha=0.9, width=0.5)
    ax1.axhline(0, color="black", linewidth=0.6)
    ax1.set_title("Croissance Revenue medianne", fontsize=9, fontweight="bold")
    ax1.set_ylabel("%", fontsize=8)
    ax1.spines["top"].set_visible(False)
    ax1.spines["right"].set_visible(False)
    ax1.yaxis.grid(True, linestyle="--", alpha=0.4)
    ax1.set_axisbelow(True)
    for bar in bars1:
        h = bar.get_height()
        ax1.annotate(f"{h:+.1f}%", xy=(bar.get_x() + bar.get_width()/2, h),
                     xytext=(0, 4 if h >= 0 else -12), textcoords="offset points",
                     ha="center", fontsize=9, fontweight="bold")

    # Momentum 52S
    mom_a = sa.get("mom", 0) or 0
    mom_b = sb.get("mom", 0) or 0
    bars2 = ax2.bar([label_a, label_b], [mom_a, mom_b], color=[_HEX_A, _HEX_B], alpha=0.9, width=0.5)
    ax2.axhline(0, color="black", linewidth=0.6)
    ax2.set_title("Performance 52 semaines medianne", fontsize=9, fontweight="bold")
    ax2.set_ylabel("%", fontsize=8)
    ax2.spines["top"].set_visible(False)
    ax2.spines["right"].set_visible(False)
    ax2.yaxis.grid(True, linestyle="--", alpha=0.4)
    ax2.set_axisbelow(True)
    for bar in bars2:
        h = bar.get_height()
        ax2.annotate(f"{h:+.1f}%", xy=(bar.get_x() + bar.get_width()/2, h),
                     xytext=(0, 4 if h >= 0 else -12), textcoords="offset points",
                     ha="center", fontsize=9, fontweight="bold")

    plt.tight_layout(pad=0.5)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=140, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def _chart_rentabilite(sa: dict, sb: dict, label_a: str, label_b: str) -> bytes:
    """Barres groupees : ROE, ROIC, FCF Yield (ROIC omis si indisponible)."""
    _roic_a = sa.get("roic")
    _roic_b = sb.get("roic")
    if _roic_a is None and _roic_b is None:
        metrics = ["ROE", "FCF Yield"]
        vals_a  = [sa.get("roe", 0) or 0, sa.get("fcfy", 0) or 0]
        vals_b  = [sb.get("roe", 0) or 0, sb.get("fcfy", 0) or 0]
    else:
        metrics = ["ROE", "ROIC", "FCF Yield"]
        vals_a  = [sa.get("roe", 0) or 0, _roic_a or 0, sa.get("fcfy", 0) or 0]
        vals_b  = [sb.get("roe", 0) or 0, _roic_b or 0, sb.get("fcfy", 0) or 0]

    fig, ax = plt.subplots(figsize=(7.5, 4.2))
    x = np.arange(len(metrics))
    w = 0.35
    bars_a = ax.bar(x - w/2, vals_a, w, label=label_a, color=_HEX_A, alpha=0.9)
    bars_b = ax.bar(x + w/2, vals_b, w, label=label_b, color=_HEX_B, alpha=0.9)
    ax.set_xticks(x)
    ax.set_xticklabels(metrics, fontsize=10)
    ax.set_ylabel("(%)", fontsize=9)
    ax.tick_params(axis="y", labelsize=8)
    ax.yaxis.grid(True, linestyle="--", alpha=0.4)
    ax.set_axisbelow(True)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    for bar in list(bars_a) + list(bars_b):
        h = bar.get_height()
        if abs(h) > 0.5:
            ax.annotate(f"{h:.1f}%", xy=(bar.get_x() + bar.get_width()/2, max(h, 0)),
                        xytext=(0, 3), textcoords="offset points",
                        ha="center", va="bottom", fontsize=7.5)
    ax.legend(fontsize=9, framealpha=0)
    plt.tight_layout(pad=0.5)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=140, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def _fetch_price_52w_pptx(tickers_a, tickers_b):
    """Composite cours Normalisé base 100 sur 52 semaines pour PPTX."""
    try:
        import yfinance as yf
        top_a = sorted(tickers_a, key=lambda t: t.get("score_global", 0), reverse=True)[:15]
        top_b = sorted(tickers_b, key=lambda t: t.get("score_global", 0), reverse=True)[:15]
        syms_a = [t["ticker"] for t in top_a if t.get("ticker")]
        syms_b = [t["ticker"] for t in top_b if t.get("ticker")]
        all_syms = list(dict.fromkeys(syms_a + syms_b))
        if not all_syms:
            return None, None
        hist = yf.download(all_syms, period="1y", interval="1wk",
                           auto_adjust=True, progress=False, timeout=30)
        if hist.empty:
            return None, None
        if len(all_syms) == 1:
            close = hist[["Close"]]
            close.columns = all_syms
        else:
            close = hist["Close"] if hasattr(hist.columns, "levels") else hist[["Close"]]
        close = close.ffill()

        def _composite(syms):
            cols = [s for s in syms if s in close.columns]
            if not cols:
                return None
            sub = close[cols].dropna(how="all")
            if sub.empty or len(sub) < 5:
                return None
            first = sub.iloc[0].replace(0, np.nan)
            norm = sub.div(first) * 100
            return norm.mean(axis=1)

        return _composite(syms_a), _composite(syms_b)
    except Exception as exc:
        log.warning("[cmp_secteur_pptx] price 52w fetch: %s", exc)
        return None, None


def _chart_price_52w_pptx(perf_a, perf_b, sector_a, sector_b) -> bytes:
    """Courbe Normalisée base 100 sur 52 semaines pour slide PPTX."""
    fig, ax = plt.subplots(figsize=(11.0, 4.5))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("#FAFBFD")
    if perf_a is not None:
        ax.plot(perf_a.index, perf_a.values, color=_HEX_A, linewidth=2.5, label=sector_a)
    if perf_b is not None:
        ax.plot(perf_b.index, perf_b.values, color="#C9A227", linewidth=2.5, label=sector_b)
    if perf_a is not None and perf_b is not None:
        try:
            idx = perf_a.index.intersection(perf_b.index)
            a_vals = perf_a.reindex(idx).values
            b_vals = perf_b.reindex(idx).values
            ax.fill_between(idx, a_vals, b_vals,
                            where=(a_vals >= b_vals), alpha=0.12, color=_HEX_A)
            ax.fill_between(idx, a_vals, b_vals,
                            where=(a_vals < b_vals), alpha=0.12, color="#C9A227")
        except Exception:
            pass
    for perf, col, nm in [(perf_a, _HEX_A, sector_a[:14]), (perf_b, "#C9A227", sector_b[:14])]:
        if perf is not None and len(perf) > 0:
            last_v = float(perf.iloc[-1])
            perf_pct = last_v - 100.0
            ax.annotate(f"{nm}: {perf_pct:+.1f}%",
                        xy=(perf.index[-1], last_v),
                        xytext=(8, 0), textcoords="offset points",
                        fontsize=9, fontweight="bold", color=col, va="center")
    ax.axhline(100, color="#AAAAAA", linestyle="--", linewidth=0.8, alpha=0.6)
    ax.xaxis.set_major_formatter(mdates.DateFormatter("%b %y"))
    ax.xaxis.set_major_locator(mdates.MonthLocator(interval=2))
    plt.setp(ax.xaxis.get_majorticklabels(), rotation=30, ha="right", fontsize=8)
    ax.set_ylabel("Indice base 100", fontsize=9)
    ax.yaxis.grid(True, linestyle="--", alpha=0.4)
    ax.set_axisbelow(True)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.legend(fontsize=9, framealpha=0.9, loc="upper left")
    ax.set_title("Performance comparative 52 semaines — composite Normalisé base 100 (top 15 par secteur)",
                 fontsize=9, color="#555555", pad=6)
    plt.tight_layout(pad=0.5)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=140, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf.read()


# ── Slides ────────────────────────────────────────────────────────────────────

def _s01_cover(prs, D):
    slide = _blank(prs)
    # Barre navy (style société : fond navy + texte grise discret)
    _rect(slide, 0, 0, 25.4, 1.85, fill=_NAVY)
    _txb(slide, "FinSight IA", 0, 0.15, 25.4, 0.7, size=10, bold=False,
         color=RGBColor(0x88, 0x99, 0xBB), align=PP_ALIGN.CENTER)
    _rect(slide, 10.7, 0.9, 4.0, 0.04, fill=RGBColor(0xAA, 0xBB, 0xCC))
    _txb(slide, "Pitchbook  —  Comparatif Sectoriel", 0, 0.98, 25.4, 0.75, size=9,
         color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER)

    # Titre principal
    title = f"{D['sector_a']}  vs  {D['sector_b']}"
    _txb(slide, title, 0, 4.0, 25.4, 2.5, size=34, bold=True, color=_NAVY, align=PP_ALIGN.CENTER)

    # Sous-titre univers
    _txb(slide, D["universe_label"], 0, 6.8, 25.4, 1.0, size=13, color=_GRAYT, align=PP_ALIGN.CENTER)

    # 4 KPI boxes : N sociétés A, Score A, N sociétés B, Score B
    kpis = [
        (f"{D['na']}", f"sociétés\n{D['sector_a']}", _COL_A),
        (f"{D['sa'].get('score', 0)}", f"Score FinSight\n{D['sector_a']}", _COL_A),
        (f"{D['nb']}", f"sociétés\n{D['sector_b']}", _COL_B),
        (f"{D['sb'].get('score', 0)}", f"Score FinSight\n{D['sector_b']}", _COL_B),
    ]
    box_w = 5.6
    step  = 5.85
    start_x = 0.9
    # KPI boxes rehaussees (8.5 -> 8.0) pour plus de marge avec footer
    for i, (val, lbl, col) in enumerate(kpis):
        bx = start_x + i * step
        pale = _COL_A_PALE if col == _COL_A else _COL_B_PALE
        _rect(slide, bx, 8.0, box_w, 2.4, fill=pale)
        _rect(slide, bx, 8.0, 0.22, 2.4, fill=col)
        _txb(slide, val, bx + 0.22, 8.15, box_w - 0.22, 1.35, size=28, bold=True, color=col, align=PP_ALIGN.CENTER)
        _txb(slide, lbl, bx + 0.22, 9.55, box_w - 0.22, 0.85, size=8.5, color=_NAVY, align=PP_ALIGN.CENTER)

    # Signal badges (rehausses 11.5 -> 11.0)
    _rect(slide, 2.2, 11.0, 4.5, 0.9, fill=D["sig_a_col"])
    _txb(slide, f"● {D['sig_a_lbl']}", 2.2, 11.08, 4.5, 0.75, size=9.5, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)
    _rect(slide, 18.7, 11.0, 4.5, 0.9, fill=D["sig_b_col"])
    _txb(slide, f"● {D['sig_b_lbl']}", 18.7, 11.08, 4.5, 0.75, size=9.5, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)

    # Footer
    _rect(slide, 0, 13.6, 25.4, 0.7, fill=_GRAYL)
    _txb(slide, "Rapport confidentiel  |  FinSight IA", 0.9, 13.65, 12.0, 0.5, size=7.5, color=_GRAYT)
    _txb(slide, D["date_str"], 0, 13.65, 24.4, 0.5, size=7.5, color=_GRAYT, align=PP_ALIGN.RIGHT)


def _s02_exec_summary(prs, D):
    slide = _blank(prs)
    sa, sb = D["sa"], D["sb"]
    _header(slide, f"Executive Summary  —  {D['sector_a']} vs {D['sector_b']}",
            f"Univers : {D['universe_label']}  |  {D['na']} + {D['nb']} sociétés analysées", 1)
    _footer(slide, D)

    # Signal A — bandeau h=1.05 pour contenir titre + sous-titre proprement
    _rect(slide, 0.9, 2.1, 11.4, 1.05, fill=_COL_A_PALE)
    _rect(slide, 0.9, 2.1, 0.12, 1.05, fill=_COL_A)
    _txb(slide, f"{D['sector_a']}  —  {D['sig_a_lbl']}", 1.3, 2.20, 10.6, 0.40, size=10, bold=True, color=_COL_A)
    _txb(slide, f"Score {sa.get('score', 0)}/100  |  P/E {_fmt_simple(sa.get('pe'), x=True)}  |  Croiss. {_fmt(sa.get('revg'), pct=True)}  |  Mg. EBITDA {_fmt_simple(sa.get('em'), pct=True)}",
         1.3, 2.62, 10.6, 0.45, size=7.5, color=_NAVY)

    # Signal B — bandeau h=1.05
    _rect(slide, 13.1, 2.1, 11.4, 1.05, fill=_COL_B_PALE)
    _rect(slide, 13.1, 2.1, 0.12, 1.05, fill=_COL_B)
    _txb(slide, f"{D['sector_b']}  —  {D['sig_b_lbl']}", 13.5, 2.20, 10.6, 0.40, size=10, bold=True, color=_COL_B)
    _txb(slide, f"Score {sb.get('score', 0)}/100  |  P/E {_fmt_simple(sb.get('pe'), x=True)}  |  Croiss. {_fmt(sb.get('revg'), pct=True)}  |  Mg. EBITDA {_fmt_simple(sb.get('em'), pct=True)}",
         13.5, 2.62, 10.6, 0.45, size=7.5, color=_GREEN)

    # Tableau metriques côte à côte
    rows = [
        ["Metrique", D["sector_a"], D["sector_b"], "Avantage"],
        ["Score FinSight", f"{sa.get('score', 0)}/100", f"{sb.get('score', 0)}/100",
         D["sector_a"] if sa.get("score", 0) > sb.get("score", 0) else D["sector_b"]],
        ["P/E Médian", _fmt_simple(sa.get("pe"), x=True), _fmt_simple(sb.get("pe"), x=True),
         D["sector_a"] if (sa.get("pe") or 999) < (sb.get("pe") or 999) else D["sector_b"]],
        ["EV/EBITDA Médian", _fmt_simple(sa.get("ev_eb"), x=True), _fmt_simple(sb.get("ev_eb"), x=True),
         D["sector_a"] if (sa.get("ev_eb") or 999) < (sb.get("ev_eb") or 999) else D["sector_b"]],
        ["Croissance rev. med.", _fmt(sa.get("revg"), pct=True), _fmt(sb.get("revg"), pct=True),
         D["sector_a"] if (sa.get("revg") or -999) > (sb.get("revg") or -999) else D["sector_b"]],
        ["Marge EBITDA med.", _fmt_simple(sa.get("em"), pct=True), _fmt_simple(sb.get("em"), pct=True),
         D["sector_a"] if (sa.get("em") or -999) > (sb.get("em") or -999) else D["sector_b"]],
        ["Marge Nette med.", _fmt_simple(sa.get("nm"), pct=True), _fmt_simple(sb.get("nm"), pct=True),
         D["sector_a"] if (sa.get("nm") or -999) > (sb.get("nm") or -999) else D["sector_b"]],
        ["ROE Médian", _fmt_simple(sa.get("roe"), pct=True), _fmt_simple(sb.get("roe"), pct=True),
         D["sector_a"] if (sa.get("roe") or -999) > (sb.get("roe") or -999) else D["sector_b"]],
        ["Beta Médian", _fmt_simple(sa.get("beta"), dp=2), _fmt_simple(sb.get("beta"), dp=2),
         D["sector_a"] if (sa.get("beta") or 999) < (sb.get("beta") or 999) else D["sector_b"]],
        ["Perf. 52S med.", _fmt(sa.get("mom"), pct=True), _fmt(sb.get("mom"), pct=True),
         D["sector_a"] if (sa.get("mom") or -999) > (sb.get("mom") or -999) else D["sector_b"]],
    ]
    _add_table(slide, rows, 0.9, 3.40, 23.6, 7.95,
               col_widths=[5.5, 3.0, 3.0, 3.5],
               header_fill=_NAVY, header_color=_WHITE,
               alt_fill=_GRAYL, font_size=8.5, header_size=8.5)

    # Synthèse IA — box LLM standard
    llm_text = D.get("llm", {}).get("exec_summary", "")
    if llm_text:
        _llm_box(slide, 0.9, 11.55, 23.6, 2.05,
                 "Analyse Comparative", llm_text, fontsize=9)


def _s03_sommaire(prs, D):
    """Slide 3 — Sommaire (style aligné sur cmp société : cards bleu pale numérotées)."""
    slide = _blank(prs)
    _header(slide, "Sommaire", f"Analyse comparative : {D['sector_a']} vs {D['sector_b']}", 1)
    _footer(slide, D)

    sections = [
        ("01", "Profil & Valorisation",  "Structure des secteurs, multiples de marché, qualité des bilans",        "Slides 5-8"),
        ("02", "Performance & Scoring",  "Cours 52 semaines, revenue growth et scoring multidimensionnel",         "Slides 10-11"),
        ("03", "Top Acteurs & Risques",  "Meilleures sociétés et conditions d'invalidation par secteur",           "Slides 13-16"),
        ("04", "Synthèse & Décision",    "Thèses d'Investissement, verdict comparatif et recommandation",          "Slides 18-21"),
    ]

    y_start  = 2.40
    card_h   = 2.55
    card_gap = 0.20
    card_w   = 23.6
    x_card   = 0.9

    for i, (num, title, desc, pages) in enumerate(sections):
        y = y_start + i * (card_h + card_gap)
        # Card bleu pale + barre laterale navy
        _rect(slide, x_card, y, card_w, card_h, fill=_NAVY_PALE)
        _rect(slide, x_card, y, 0.18, card_h, fill=_NAVY)
        # Numero gros bleu navy a gauche
        _txb(slide, num, x_card + 0.40, y + 0.30, 3.20, 1.95,
             size=44, bold=True, color=_NAVY_MID, align=PP_ALIGN.CENTER)
        # Titre bold navy
        _txb(slide, title, x_card + 4.00, y + 0.45, card_w - 8.0, 0.80,
             size=15, bold=True, color=_NAVY)
        # Description en gris fonce
        _txb(slide, desc, x_card + 4.00, y + 1.30, card_w - 8.0, 1.00,
             size=10, color=_GRAYT, wrap=True)
        # Badge "Slides X-Y" a droite
        _rect(slide, x_card + card_w - 3.40, y + 0.85, 3.00, 0.85, fill=_NAVY)
        _txb(slide, pages, x_card + card_w - 3.40, y + 1.05, 3.00, 0.45,
             size=10, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)


def _s05_profil(prs, D):
    """Slide 5 — Profil des 2 secteurs côte à côte."""
    slide = _blank(prs)
    _header(slide, "Profil & Structure des Secteurs",
            f"{D['sector_a']}  vs  {D['sector_b']}  |  Positionnement structurel", 1)
    _footer(slide, D)

    content_a = _get_content(D["sector_a"])
    content_b = _get_content(D["sector_b"])
    col_w = 11.5
    xa = 0.9
    xb = 13.0

    # Secteur A — bandeau titre rehausse h=0.70
    _rect(slide, xa, 2.10, col_w, 0.70, fill=_COL_A)
    _txb(slide, D["sector_a"][:30], xa + 0.20, 2.20, col_w - 0.4, 0.55,
         size=11, bold=True, color=_WHITE)

    # Description A en box LLM standard
    desc_a = content_a.get("description", "")[:480]
    _llm_box(slide, xa, 2.95, col_w, 3.30, "Description sectorielle", desc_a, fontsize=8.5)

    # Cycle A
    cycle_a = content_a.get("cycle_comment", "")
    _rect(slide, xa, 6.40, col_w, 0.50, fill=_COL_A_PALE)
    _txb(slide, f"Cycle : {cycle_a}", xa + 0.15, 6.45, col_w - 0.3, 0.42,
         size=8.5, italic=True, color=_COL_A)

    # Drivers A — bullets carres
    _txb(slide, "Drivers principaux", xa, 7.10, col_w, 0.42,
         size=9.5, bold=True, color=_NAVY)
    drivers_a = content_a.get("drivers", [])[:4]
    for j, drv in enumerate(drivers_a):
        direction, name, desc = (drv[0], drv[1], drv[2]) if len(drv) >= 3 else ("up", str(drv), "")
        bullet_col = _COL_A if direction == "up" else _RED
        _bullet_square(slide, xa, 7.75 + j * 1.45, bullet_col, size=0.22)
        _txb(slide, name[:38], xa + 0.40, 7.70 + j * 1.45, col_w - 0.45, 0.45,
             size=9, bold=True, color=_BLACK)
        _txb(slide, desc[:90], xa + 0.40, 8.18 + j * 1.45, col_w - 0.45, 0.85,
             size=7.5, italic=True, color=_GRAYT, wrap=True)

    # Secteur B — bandeau titre rehausse
    _rect(slide, xb, 2.10, col_w, 0.70, fill=_COL_B)
    _txb(slide, D["sector_b"][:30], xb + 0.20, 2.20, col_w - 0.4, 0.55,
         size=11, bold=True, color=_WHITE)

    desc_b = content_b.get("description", "")[:480]
    _llm_box(slide, xb, 2.95, col_w, 3.30, "Description sectorielle", desc_b, fontsize=8.5)

    cycle_b = content_b.get("cycle_comment", "")
    _rect(slide, xb, 6.40, col_w, 0.50, fill=_COL_B_PALE)
    _txb(slide, f"Cycle : {cycle_b}", xb + 0.15, 6.45, col_w - 0.3, 0.42,
         size=8.5, italic=True, color=_COL_B)

    _txb(slide, "Drivers principaux", xb, 7.10, col_w, 0.42,
         size=9.5, bold=True, color=_COL_B)
    drivers_b = content_b.get("drivers", [])[:4]
    for j, drv in enumerate(drivers_b):
        direction, name, desc = (drv[0], drv[1], drv[2]) if len(drv) >= 3 else ("up", str(drv), "")
        bullet_col = _COL_B if direction == "up" else _RED
        _bullet_square(slide, xb, 7.75 + j * 1.45, bullet_col, size=0.22)
        _txb(slide, name[:38], xb + 0.40, 7.70 + j * 1.45, col_w - 0.45, 0.45,
             size=9, bold=True, color=_BLACK)
        _txb(slide, desc[:90], xb + 0.40, 8.18 + j * 1.45, col_w - 0.45, 0.85,
             size=7.5, italic=True, color=_GRAYT, wrap=True)


def _s06_valorisation(prs, D):
    """Slide 6 — Multiples de valorisation comparatifs."""
    slide = _blank(prs)
    sa, sb = D["sa"], D["sb"]
    _header(slide, "Valorisation Comparée  —  Multiples de Marché",
            f"P/E et EV-multiples Médianes — {D['sector_a']} vs {D['sector_b']}", 1)
    _footer(slide, D)

    # Chart gauche (aspect ratio correct : figsize 7.5x4.2 → placement 12.5x5.5)
    try:
        img = _chart_valuation(sa, sb, D["sector_a"], D["sector_b"])
        _pic(slide, img, 0.9, 2.0, 12.5, 5.5)
    except Exception as e:
        log.warning("[cmp_secteur] _s06 chart: %s", e)
        _txb(slide, f"Graphique indisponible: {e}", 0.9, 5.0, 12.5, 1.0, size=8, color=_GRAYT)

    # Tableau metriques complémentaires sous le chart (gauche)
    rows_extra = [
        ["Indicateurs", D["sector_a"], D["sector_b"]],
        ["FCF Yield med.", _fmt_simple(sa.get("fcfy"), pct=True, dp=1), _fmt_simple(sb.get("fcfy"), pct=True, dp=1)],
        ["Beta Médian", _fmt_simple(sa.get("beta"), dp=2), _fmt_simple(sb.get("beta"), dp=2)],
        ["Piotroski med.", _fmt_simple(sa.get("pf"), dp=1), _fmt_simple(sb.get("pf"), dp=1)],
    ]
    # Altman Z : afficher seulement si au moins un secteur a la donnée
    if sa.get("az") is not None or sb.get("az") is not None:
        rows_extra.insert(3, ["Altman Z med.", _fmt_simple(sa.get("az"), dp=2), _fmt_simple(sb.get("az"), dp=2)])
    _add_table(slide, rows_extra, 0.9, 7.8, 12.5, 2.85,
               col_widths=[4.5, 4.0, 4.0], alt_fill=_GRAYL, font_size=9, header_size=9)

    # Tableau metriques droite (multiples)
    rows = [
        ["Metrique", D["sector_a"], D["sector_b"]],
        ["P/E Médian", _fmt_simple(sa.get("pe"), x=True), _fmt_simple(sb.get("pe"), x=True)],
        ["EV/EBITDA med.", _fmt_simple(sa.get("ev_eb"), x=True), _fmt_simple(sb.get("ev_eb"), x=True)],
        ["EV/Revenue med.", _fmt_simple(sa.get("ev_rev"), x=True), _fmt_simple(sb.get("ev_rev"), x=True)],
        ["PEG ratio med.", _fmt_simple(sa.get("peg"), dp=2), _fmt_simple(sb.get("peg"), dp=2)],
    ]
    _add_table(slide, rows, 14.0, 2.0, 10.5, 6.0,
               col_widths=[4.5, 3.0, 3.0], alt_fill=_GRAYL, font_size=9, header_size=9)

    # Lecture analytique
    pe_a = sa.get("pe") or 0
    pe_b = sb.get("pe") or 0
    cheaper = D["sector_a"] if pe_a < pe_b else D["sector_b"]
    premium = D["sector_a"] if pe_a > pe_b else D["sector_b"]
    pe_spread = abs(pe_a - pe_b)
    fallback_lecture = (
        f"{cheaper} côté a une prime inferieure de {pe_spread:.1f}x P/E "
        f"vs {premium}. Écart reflète "
        f"{'une anticipation de croissance superieure' if (sa.get('revg') or 0) > (sb.get('revg') or 0) else 'un profil de risque différentié'}."
    )
    lecture = D.get("llm", {}).get("valuation_read") or fallback_lecture
    _llm_box(slide, 14.0, 8.50, 10.5, 4.85, "Lecture analytique", lecture, fontsize=8.5)


def _s07_marges(prs, D):
    """Slide 7 — Qualité, Marges & Rentabilité."""
    slide = _blank(prs)
    sa, sb = D["sa"], D["sb"]
    _header(slide, "Qualité & Rentabilité  —  Marges et Retour sur Capital",
            f"Qui généré plus de valeur par euro de chiffre d'affaires ?", 1)
    _footer(slide, D)

    try:
        img = _chart_margins(sa, sb, D["sector_a"], D["sector_b"])
        _pic(slide, img, 0.9, 2.0, 12.5, 5.5)
    except Exception as e:
        log.warning("[cmp_secteur] _s07 chart: %s", e)

    try:
        img2 = _chart_rentabilite(sa, sb, D["sector_a"], D["sector_b"])
        _pic(slide, img2, 13.5, 2.0, 11.0, 5.5)
    except Exception as e:
        log.warning("[cmp_secteur] _s07 chart2: %s", e)

    # Lecture analytique marges
    winner_margin = D["sector_a"] if (sa.get("em") or -999) > (sb.get("em") or -999) else D["sector_b"]
    winner_roe    = D["sector_a"] if (sa.get("roe") or -999) > (sb.get("roe") or -999) else D["sector_b"]
    fallback_diag = (
        f"{winner_margin} affiche une marge EBITDA superieure "
        f"({_fmt_simple(sa.get('em') if winner_margin == D['sector_a'] else sb.get('em'), pct=True)}). "
        f"{winner_roe} domine sur le ROE "
        f"({_fmt_simple(min(sa.get('roe') or 0, 999.9) if winner_roe == D['sector_a'] else min(sb.get('roe') or 0, 999.9), pct=True)})."
    )
    diag = D.get("llm", {}).get("margins_read") or fallback_diag
    _llm_box(slide, 0.9, 7.75, 23.6, 5.55, "Lecture analytique", diag, fontsize=9)


def _s07b_capital_alloc(prs, D):
    """Slide 7b — Capital Allocation & Rémunération de l'Actionnaire."""
    slide = _blank(prs)
    sa, sb = D["sa"], D["sb"]
    _header(slide, "Capital Allocation  —  Dividendes, FCF & Rémunération",
            f"Comparaison politique de distribution — {D['sector_a']} vs {D['sector_b']}", 1)
    _footer(slide, D)

    # Tableau comparatif gauche
    def _fmt_pct(v):
        """div_yield et payout sont en fraction (×100). fcfy est déjà en %, utiliser _fmt_pct_direct."""
        if v is None: return "—"
        try: return f"{float(v)*100:.1f} %"
        except: return "—"
    def _fmt_pct_direct(v):
        """fcfy est déjà en % (1.38 = 1.38%) — pas de ×100."""
        if v is None: return "—"
        try: return f"{float(v):.1f} %"
        except: return "—"

    rows = [
        ["Indicateur", D["sector_a"], D["sector_b"]],
        ["Rendement dividende med.", _fmt_pct(sa.get("div_yield")), _fmt_pct(sb.get("div_yield"))],
        ["FCF Yield Médian",         _fmt_pct_direct(sa.get("fcfy")),  _fmt_pct_direct(sb.get("fcfy"))],
        ["Payout Ratio Médian",      _fmt_pct(sa.get("payout")),    _fmt_pct(sb.get("payout"))],
        # fcfy en % (1.4), div_yield en fraction (0.004) → convertir div_yield en %
        ["FCF Yield / Div Yield",    f"{(sa.get('fcfy') or 0) / max((sa.get('div_yield') or 0)*100, 0.01):.1f}x" if (sa.get("fcfy") and sa.get("div_yield") and (sa.get("div_yield") or 0)*100 > 0.05) else "—",
                                     f"{(sb.get('fcfy') or 0) / max((sb.get('div_yield') or 0)*100, 0.01):.1f}x" if (sb.get("fcfy") and sb.get("div_yield") and (sb.get("div_yield") or 0)*100 > 0.05) else "—"],
    ]
    _add_table(slide, rows, 0.9, 2.0, 13.0, 5.5,
               col_widths=[5.5, 3.7, 3.7], alt_fill=_GRAYL, font_size=9, header_size=9)

    # Graphique barres capital allocation
    try:
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        import numpy as np
        import io as _io

        metrics = []
        vals_a  = []
        vals_b  = []
        # div_yield est en fraction (×100) ; fcfy est déjà en % (pas de ×100)
        for label, key, scale in [("Div Yield", "div_yield", 100), ("FCF Yield", "fcfy", 1)]:
            va = sa.get(key)
            vb = sb.get(key)
            if va is not None or vb is not None:
                metrics.append(label)
                vals_a.append(float(va or 0) * scale)
                vals_b.append(float(vb or 0) * scale)

        if metrics:
            x  = np.arange(len(metrics))
            bw = 0.32
            fig, ax = plt.subplots(figsize=(7.5, 3.2))
            fig.patch.set_facecolor('white')
            ax.set_facecolor('#FAFBFD')
            b1 = ax.bar(x - bw/2, vals_a, bw, label=D["sector_a"][:18],
                        color='#2E5FA3', alpha=0.85)
            b2 = ax.bar(x + bw/2, vals_b, bw, label=D["sector_b"][:18],
                        color='#1A7A4A', alpha=0.85)
            ax.set_xticks(x)
            ax.set_xticklabels(metrics, fontsize=10)
            ax.set_ylabel("%", fontsize=9)
            ax.legend(fontsize=8, framealpha=0.9)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.grid(axis='y', alpha=0.3)
            for bar in list(b1) + list(b2):
                h = bar.get_height()
                if h > 0:
                    ax.annotate(f'{h:.1f}%',
                               xy=(bar.get_x()+bar.get_width()/2, h),
                               xytext=(0, 2), textcoords="offset points",
                               ha='center', va='bottom', fontsize=8)
            fig.tight_layout(pad=0.5)
            buf_chart = _io.BytesIO()
            fig.savefig(buf_chart, format='png', dpi=130, bbox_inches='tight')
            plt.close(fig)
            buf_chart.seek(0)
            _pic(slide, buf_chart.read(), 0.9, 7.8, 13.0, 5.5)
    except Exception as e:
        log.warning("[cmp_secteur_pptx] capital_alloc chart: %s", e)

    # Lecture analytique (toute la moitie droite — Tech/Health blocks supprimés,
    # redondants avec le tableau de gauche). Plus grande box LLM standard.
    dy_a = sa.get("div_yield") or 0
    dy_b = sb.get("div_yield") or 0
    fy_a = sa.get("fcfy") or 0
    fy_b = sb.get("fcfy") or 0
    high_dy = D["sector_a"] if dy_a > dy_b else D["sector_b"]
    high_fy = D["sector_a"] if fy_a > fy_b else D["sector_b"]
    fallback_interp = (
        f"{high_dy} offre le meilleur rendement dividende (med. {_fmt_pct(max(dy_a, dy_b))}), "
        f"ce qui reflète un modèle de maturité bilancielle et une politique de distribution établie. "
        f"{high_fy} généré davantage de FCF ({_fmt_pct_direct(max(fy_a, fy_b))} de FCF Yield), "
        f"signal d'une capacite de rémunération durable et d'une allocation du capital disciplinée. "
        f"Un FCF Yield superieur au dividende versé garantit la soutenabilite et la résilience "
        f"de la distribution meme en phase de contraction des marges. "
        f"L'arbitrage entre dividende versé et réinvestissement traduit la maturité du cycle "
        f"d'investissement : un secteur en phase de croissance préféré allouer son cash a "
        f"l'expansion, un secteur mature monétisé via la distribution."
    )
    interp_txt = D.get("llm", {}).get("capital_alloc_read") or fallback_interp
    _llm_box(slide, 14.5, 2.0, 9.95, 11.30, "Lecture analytique", interp_txt, fontsize=9)


def _s08_croissance(prs, D):
    """Slide 10 — Performance Boursière : cours comparatif 52 semaines."""
    slide = _blank(prs)
    sa, sb = D["sa"], D["sb"]
    _header(slide, "Performance Boursière  —  Cours Comparatif 52 Semaines",
            f"Composite Normalisé base 100 — top 15 sociétés par secteur  |  {D['sector_a']} vs {D['sector_b']}", 2)
    _footer(slide, D)

    perf_a = D.get("perf_a_52w")
    perf_b = D.get("perf_b_52w")

    if perf_a is not None or perf_b is not None:
        try:
            img = _chart_price_52w_pptx(perf_a, perf_b, D["sector_a"], D["sector_b"])
            _pic(slide, img, 0.9, 2.0, 16.5, 7.5)
        except Exception as e:
            log.warning("[cmp_secteur] _s08 52w chart: %s", e)
            _txb(slide, "Graphique cours indisponible", 0.9, 5.0, 16.5, 1.0, size=8, color=_GRAYT)
    else:
        _txb(slide, "Données cours indisponibles — vérifier connexion yfinance",
             0.9, 6.0, 16.5, 1.0, size=8.5, color=_GRAYT, align=PP_ALIGN.CENTER)

    # Stats table on right (colonnes assez larges pour les noms de secteur)
    sa_lbl = D["sector_a"][:9] if len(D["sector_a"]) > 9 else D["sector_a"]
    sb_lbl = D["sector_b"][:9] if len(D["sector_b"]) > 9 else D["sector_b"]
    rows = [
        ["", sa_lbl, sb_lbl],
        ["Perf. 52S med.", _fmt(sa.get("mom"), pct=True), _fmt(sb.get("mom"), pct=True)],
        ["Croiss. rev. med.", _fmt(sa.get("revg"), pct=True), _fmt(sb.get("revg"), pct=True)],
        ["Beta Médian", _fmt_simple(sa.get("beta"), dp=2), _fmt_simple(sb.get("beta"), dp=2)],
        ["FCF Yield med.", _fmt_simple(sa.get("fcfy"), pct=True, dp=1), _fmt_simple(sb.get("fcfy"), pct=True, dp=1)],
        ["Score Global", f"{sa.get('score', 0)}/100", f"{sb.get('score', 0)}/100"],
    ]
    _add_table(slide, rows, 17.5, 2.0, 7.0, 7.5,
               col_widths=[2.8, 2.1, 2.1], alt_fill=_GRAYL, font_size=8.5, header_size=8)

    # Lecture analytique
    faster = D["sector_a"] if (sa.get("mom") or -999) > (sb.get("mom") or -999) else D["sector_b"]
    faster_s = sa if faster == D["sector_a"] else sb
    slower_s = sb if faster == D["sector_a"] else sa
    spread_m = abs((sa.get("mom") or 0) - (sb.get("mom") or 0))
    spread_g = abs((sa.get("revg") or 0) - (sb.get("revg") or 0))
    fallback_lecture = (
        f"{faster} surperforme sur 52 semaines avec un écart de {spread_m:.1f} pts "
        f"({_fmt(faster_s.get('mom'), pct=True)} vs {_fmt(slower_s.get('mom'), pct=True)}). "
        f"Cet écart de performance reflète des dynamiques fondamentales distinctes : "
        f"{'croissance des revenus superieure de ' + str(round(spread_g, 1)) + ' pts, ' if spread_g > 1 else ''}"
        f"qualité bilancielle et momentum sectoriel divergents. "
        f"La courbe Normalisée permet d'isoler la performance pure, independamment des niveaux absolus."
    )
    lecture = D.get("llm", {}).get("growth_read") or fallback_lecture
    _llm_box(slide, 0.9, 9.80, 23.6, 3.55, "Lecture analytique cours 52S & macro",
             lecture, fontsize=9)


def _s09_scoring(prs, D):
    """Slide 9 — Scoring multi-critères (radar)."""
    slide = _blank(prs)
    sa, sb = D["sa"], D["sb"]
    _header(slide, "Scoring Multi-critères  —  Value / Growth / Quality / Momentum",
            f"Score FinSight décompose par dimension — {D['sector_a']} vs {D['sector_b']}", 2)
    _footer(slide, D)

    # Tableau scores — colonne gauche (x=0.9, w=10.5)
    rows_a = [
        ["Dimension", D["sector_a"], D["sector_b"]],
        ["Value", f"{sa.get('s_val', 0):.0f}/25", f"{sb.get('s_val', 0):.0f}/25"],
        ["Growth", f"{sa.get('s_gro', 0):.0f}/25", f"{sb.get('s_gro', 0):.0f}/25"],
        ["Quality", f"{sa.get('s_qua', 0):.0f}/25", f"{sb.get('s_qua', 0):.0f}/25"],
        ["Momentum", f"{sa.get('s_mom', 0):.0f}/25", f"{sb.get('s_mom', 0):.0f}/25"],
        ["TOTAL", f"{sa.get('score', 0)}/100", f"{sb.get('score', 0)}/100"],
    ]
    _add_table(slide, rows_a, 0.9, 2.2, 10.5, 7.0,
               col_widths=[4.0, 3.25, 3.25], alt_fill=_GRAYL, font_size=10, header_size=10)

    # Radar chart — coin droit réduit pour laisser plus d'espace au texte
    try:
        img = _chart_radar(sa, sb, D["sector_a"], D["sector_b"])
        _pic(slide, img, 13.0, 1.8, 11.5, 9.5)
    except Exception as e:
        log.warning("[cmp_secteur] _s09 radar: %s", e)

    # Interprétation globale (pleine largeur sous le tableau)
    winner = D["sector_a"] if sa.get("score", 0) > sb.get("score", 0) else D["sector_b"]
    score_w = sa.get("score", 0) if winner == D["sector_a"] else sb.get("score", 0)
    score_l = sb.get("score", 0) if winner == D["sector_a"] else sa.get("score", 0)
    fallback_interp = (
        f"{winner} ressort en avance sur le scoring FinSight global ({score_w}/100 vs {score_l}/100), "
        f"avec un écart de {abs(score_w - score_l)} pts revelant une superiorite structurelle. "
        f"L'analyse radar permet d'identifier les dimensions de force et de faiblesse relatives : "
        f"un secteur peut dominer sur la qualité tout en Étant pénalisé sur la valorisation. "
        f"L'écart global conditionne le signal d'allocation : >= 65 pts = Surpondérer, "
        f"45-64 pts = neutre, < 45 pts = Sous-pondérer."
    )
    interp = D.get("llm", {}).get("scoring_read") or fallback_interp
    _llm_box(slide, 0.9, 9.45, 11.5, 3.90,
             "Interprétation scoring multidimensionnel", interp, fontsize=9)


def _s10_top_a(prs, D):
    """Slide 10 — Top acteurs Secteur A."""
    slide = _blank(prs)
    _header(slide, f"Top Acteurs  —  {D['sector_a']}",
            f"{D['na']} sociétés analysées — classees par score FinSight", 3)
    _footer(slide, D)
    _barre_secteur(slide, D["sector_a"], D["td_a"], _COL_A, D["universe_a"],
                   llm_text=D.get("llm", {}).get("top_a_read"))


def _s11_top_b(prs, D):
    """Slide 11 — Top acteurs Secteur B."""
    slide = _blank(prs)
    _header(slide, f"Top Acteurs  —  {D['sector_b']}",
            f"{D['nb']} sociétés analysées — classees par score FinSight", 3)
    _footer(slide, D)
    _barre_secteur(slide, D["sector_b"], D["td_b"], _COL_B, D["universe_b"],
                   llm_text=D.get("llm", {}).get("top_b_read"))


def _barre_secteur(slide, sector_name, tickers_data, col, universe, llm_text=None):
    """Tableau top 8 acteurs d'un secteur."""
    sorted_td = sorted(tickers_data, key=lambda x: x.get("score_global", 0), reverse=True)[:8]

    table_h = 7.8 if llm_text else 11.0
    rows = [["Société", "Score", "P/E", "EV/EBITDA", "Croiss.Rev.", "Mg.EBITDA", "Mg.Nette", "ROE", "Beta"]]
    for t in sorted_td:
        rows.append([
            t.get("company", t.get("ticker", ""))[:24],
            f"{t.get('score_global', 0)}/100",
            _fmt_simple(t.get("pe_ratio"), x=True),
            _fmt_simple(t.get("ev_ebitda"), x=True),
            _fmt(t.get("revenue_growth", 0) * 100 if (t.get("revenue_growth") or 0) < 5 else t.get("revenue_growth", 0), pct=True),
            _fmt_simple(t.get("ebitda_margin"), pct=True),
            _fmt_simple(t.get("net_margin"), pct=True),
            ("N/M" if t.get("roe") is not None and abs(t.get("roe")) > 200 else _fmt_simple(t.get("roe"), pct=True)),
            _fmt_simple(t.get("beta"), dp=2),
        ])

    _add_table(slide, rows, 0.9, 2.0, 23.6, table_h,
               col_widths=[4.5, 2.0, 2.0, 2.5, 2.5, 2.5, 2.5, 2.0, 1.8],
               header_fill=col, alt_fill=_GRAYL, font_size=8.5, header_size=8.5)

    if llm_text:
        _llm_box(slide, 0.9, 10.05, 23.6, 3.30,
                 "Lecture analytique", llm_text, fontsize=9)


def _s12_risques_a(prs, D):
    """Slide 12 — Risques & Catalyseurs Secteur A."""
    slide = _blank(prs)
    content = _get_content(D["sector_a"])
    _header(slide, f"Risques & Catalyseurs  —  {D['sector_a']}",
            "Conditions d'investissement et d'invalidation de la Thèse", 3)
    _footer(slide, D)
    _risques_slide(slide, content, _COL_A, _COL_A_PALE,
                   llm_text=D.get("llm", {}).get("Thèse_a"))


def _s13_risques_b(prs, D):
    """Slide 13 — Risques & Catalyseurs Secteur B."""
    slide = _blank(prs)
    content = _get_content(D["sector_b"])
    _header(slide, f"Risques & Catalyseurs  —  {D['sector_b']}",
            "Conditions d'investissement et d'invalidation de la Thèse", 3)
    _footer(slide, D)
    _risques_slide(slide, content, _COL_B, _COL_B_PALE,
                   llm_text=D.get("llm", {}).get("Thèse_b"))


def _risques_slide(slide, content, col, col_pale, llm_text=None):
    cats  = content.get("catalyseurs", [])[:3]
    risks = content.get("risques", [])[:3]
    conds = content.get("conditions", [])[:4]

    # Catalyseurs (bandeau h=0.60 + bullets carres a gauche)
    _rect(slide, 0.9, 2.10, 11.4, 0.60, fill=col)
    _txb(slide, "CATALYSEURS", 0.9, 2.16, 11.4, 0.48, size=9, bold=True,
         color=_WHITE, align=PP_ALIGN.CENTER)
    for i, (title, body) in enumerate(cats):
        y = 2.90 + i * 2.20
        _bullet_square(slide, 0.95, y + 0.10, col, size=0.22)
        _txb(slide, title[:48], 1.30, y, 11.0, 0.45, size=9.5, bold=True, color=col)
        _txb(slide, body[:180], 1.30, y + 0.50, 10.95, 1.55, size=8, color=_BLACK, wrap=True)

    # Risques (bandeau h=0.60 + bullets carres rouges)
    _rect(slide, 13.1, 2.10, 11.4, 0.60, fill=_RED)
    _txb(slide, "RISQUES", 13.1, 2.16, 11.4, 0.48, size=9, bold=True,
         color=_WHITE, align=PP_ALIGN.CENTER)
    for i, (title, body) in enumerate(risks):
        y = 2.90 + i * 2.20
        _bullet_square(slide, 13.15, y + 0.10, _RED, size=0.22)
        _txb(slide, title[:48], 13.50, y, 11.0, 0.45, size=9.5, bold=True, color=_RED)
        _txb(slide, body[:180], 13.50, y + 0.50, 10.95, 1.55, size=8, color=_BLACK, wrap=True)

    # Conditions d'invalidation
    if conds:
        _rect(slide, 0.9, 9.55, 23.6, 0.48, fill=_NAVY)
        _txb(slide, "CONDITIONS D'INVALIDATION", 0.9, 9.60, 23.6, 0.42, size=9, bold=True,
             color=_WHITE, align=PP_ALIGN.CENTER)
        col_w = [3.5, 10.0, 2.5]
        rows_cond = [["Type", "Condition", "Horizon"]] + [[c[0], c[1][:100], c[2]] for c in conds]
        _add_table(slide, rows_cond, 0.9, 10.15, 23.6, 3.20,
                   col_widths=col_w, header_fill=_NAVY, alt_fill=_GRAYL,
                   font_size=8, header_size=8)
    elif llm_text:
        _llm_box(slide, 0.9, 9.55, 23.6, 3.80,
                 "Synthèse FinSight IA", llm_text, fontsize=9)


def _s14_synthese(prs, D):
    """Slide 18 — Synthèse : THÈSE LONG A (navy) + THÈSE LONG B (gold) + RISQUES PRINCIPAUX (red)."""
    slide = _blank(prs)
    sa, sb = D["sa"], D["sb"]
    _header(slide, "Synthèse Comparative  —  Thèses d'Investissement",
            f"Arguments structurels et risques majeurs : {D['sector_a']} vs {D['sector_b']}", 4)
    _footer(slide, D)

    col_w = 11.2
    xa = 0.9
    xb = 13.1

    # === THÈSE LONG SECTEUR A (navy) ===
    _rect(slide, xa, 2.1, col_w, 0.6, fill=_COL_A)
    _txb(slide, f"THÈSE LONG  {D['sector_a'].upper()}", xa + 0.12, 2.16, col_w - 0.25, 0.48,
         size=9.5, bold=True, color=_WHITE)

    forces_a = _build_forces(sa, D["sector_a"])
    these_a_llm = D.get("llm", {}).get("Thèse_long_a", "")
    y_a = 2.9
    items_a = _split_these_items(these_a_llm, forces_a)
    for j, (title, desc) in enumerate(items_a[:3]):
        _rect(slide, xa, y_a + j * 2.5, 0.22, 0.22, fill=_COL_A)
        _txb(slide, title[:42], xa + 0.38, y_a + j * 2.5 - 0.04, col_w - 0.4, 0.48, size=8.5, bold=True, color=_BLACK)
        _txb(slide, desc[:130], xa + 0.38, y_a + j * 2.5 + 0.48, col_w - 0.4, 1.8, size=8, italic=True, color=_GRAYT, wrap=True)

    # === THÈSE LONG SECTEUR B (gold) ===
    _rect(slide, xb, 2.1, col_w, 0.6, fill=_GOLD)
    _txb(slide, f"THÈSE LONG  {D['sector_b'].upper()}", xb + 0.12, 2.16, col_w - 0.25, 0.48,
         size=9.5, bold=True, color=_WHITE)

    forces_b = _build_forces(sb, D["sector_b"])
    these_b_llm = D.get("llm", {}).get("Thèse_long_b", "")
    y_b = 2.9
    items_b = _split_these_items(these_b_llm, forces_b)
    for j, (title, desc) in enumerate(items_b[:3]):
        _rect(slide, xb, y_b + j * 2.5, 0.22, 0.22, fill=_GOLD)
        _txb(slide, title[:42], xb + 0.38, y_b + j * 2.5 - 0.04, col_w - 0.4, 0.48, size=8.5, bold=True, color=_BLACK)
        _txb(slide, desc[:130], xb + 0.38, y_b + j * 2.5 + 0.48, col_w - 0.4, 1.8, size=8, italic=True, color=_GRAYT, wrap=True)

    # === RISQUES PRINCIPAUX (red, pleine largeur) ===
    risks_txt = D.get("llm", {}).get("risques_principaux") or ""
    if not risks_txt:
        weaks_a = _build_weaknesses(sa, D["sector_a"])
        weaks_b = _build_weaknesses(sb, D["sector_b"])
        w_a = weaks_a[0] if weaks_a else ""
        w_b = weaks_b[0] if weaks_b else ""
        risks_txt = (
            f"{D['sector_a']} : {w_a}.  "
            f"{D['sector_b']} : {w_b}."
        )
    # Header navy "RISQUES PRINCIPAUX" puis box LLM standard
    _rect(slide, xa, 10.50, 23.4, 0.55, fill=_RED)
    _txb(slide, "RISQUES PRINCIPAUX", xa + 0.20, 10.56, 23.0, 0.45, size=10, bold=True,
         color=_WHITE, align=PP_ALIGN.CENTER)
    _llm_box(slide, xa, 11.10, 23.4, 2.25, "Lecture des risques", risks_txt[:480], fontsize=9)


def _split_these_items(llm_text: str, fallback_forces: list) -> list:
    """Transforme un texte LLM ou les forces en liste (title, desc) pour THÈSE LONG."""
    if llm_text and len(llm_text) > 40:
        # Tenter de splititer par ". " en phrases
        parts = [p.strip() for p in llm_text.split(". ") if len(p.strip()) > 15]
        result = []
        for p in parts[:3]:
            # Chercher un separateur ":" ou " — "
            if " : " in p:
                t, d = p.split(" : ", 1)
                result.append((t.strip(), d.strip()))
            elif " - " in p:
                t, d = p.split(" - ", 1)
                result.append((t.strip(), d.strip()))
            else:
                # Prendre les 4 premiers mots comme titre
                words = p.split()
                title = " ".join(words[:4]) if len(words) >= 4 else p[:30]
                desc = " ".join(words[4:]) if len(words) > 4 else ""
                result.append((title, desc))
        if result:
            return result
    # Fallback : utiliser les forces en les splittant sur " — "
    result = []
    for force in fallback_forces[:3]:
        if " — " in force:
            t, d = force.split(" — ", 1)
        elif "(" in force:
            t = force[:force.index("(")].strip()
            d = force[force.index("("):].strip()
        else:
            t = force[:35]
            d = force[35:] if len(force) > 35 else ""
        result.append((t.strip(), d.strip()))
    return result


def _build_forces(s: dict, sector_name: str) -> list[str]:
    forces = []
    if (s.get("em") or 0) > 20:
        forces.append(f"Marge EBITDA élevée ({s.get('em', 0):.0f}%) — forte capacite a générer du cash opérationnel")
    if (s.get("revg") or 0) > 5:
        forces.append(f"Croissance revenue soutenue ({s.get('revg', 0):+.1f}%) — secteur en expansion structurelle")
    if (s.get("roe") or 0) > 15:
        forces.append(f"ROE attractif ({s.get('roe', 0):.0f}%) — modèle a forte création de valeur actionnariale")
    if (s.get("roic") or 0) > 10:
        forces.append(f"ROIC > Coût du capital ({s.get('roic', 0):.1f}%) — allocation capital disciplinée")
    if (s.get("fcfy") or 0) > 3:
        forces.append(f"FCF Yield ({s.get('fcfy', 0):.1f}%) — génération de cash visible, base pour dividendes et rachats")
    if (s.get("mom") or 0) > 5:
        forces.append(f"Momentum 52S positif ({s.get('mom', 0):+.1f}%) — confiance des investisseurs confirmee")
    # Padding : toujours 3 items pour équilibrer la mise en page
    _pad = [
        "Bilan sectoriel solide — solidité financière confirmee",
        "Diversification des revenus géographique structurelle",
        "Visibilite sur les flux de trésorerie favorable",
    ]
    i = 0
    while len(forces) < 3:
        forces.append(_pad[i % len(_pad)])
        i += 1
    return forces[:3]


def _build_weaknesses(s: dict, sector_name: str) -> list[str]:
    weaks = []
    if (s.get("pe") or 0) > 30:
        weaks.append(f"Valorisation tendue (P/E {s.get('pe', 0):.0f}x) — peu de marge de sécurité en cas de déception")
    if (s.get("revg") or 0) < 2:
        weaks.append(f"Croissance revenue faible ({s.get('revg', 0):+.1f}%) — risque de re-rating baissier")
    if (s.get("nm") or 0) < 5:
        weaks.append(f"Marge nette comprimee ({s.get('nm', 0):.1f}%) — Sensibilité élevée aux hausses de Coûts")
    if (s.get("beta") or 1) > 1.3:
        weaks.append(f"Beta élevé ({s.get('beta', 0):.2f}) — volatilite superieure au marché, prudence en phase de correction")
    if (s.get("mom") or 0) < -5:
        weaks.append(f"Momentum negatif 52S ({s.get('mom', 0):+.1f}%) — pression vendeuse persistante")
    if (s.get("ev_eb") or 0) > 20:
        weaks.append(f"EV/EBITDA premium ({s.get('ev_eb', 0):.0f}x) — execution sans faute requise")
    # Padding : toujours 2 items minimum pour équilibrer
    _pad = [
        "Surveillance du consensus — révision de valorisation possible",
        "Sensibilité macro a surveiller sur les 12 prochains mois",
        "Execution du management déterminante pour maintenir les multiples",
    ]
    i = 0
    while len(weaks) < 2:
        weaks.append(_pad[i % len(_pad)])
        i += 1
    return weaks[:3]


def _s15_allocation(prs, D):
    """Slide 19 — Recommandation d'allocation."""
    slide = _blank(prs)
    sa, sb = D["sa"], D["sb"]
    _header(slide, "Recommandation  —  Positionnement Portefeuille",
            "Signal FinSight et implications pour la construction de portefeuille", 4)
    _footer(slide, D)

    # Contexte macro
    _rect(slide, 0.9, 2.1, 23.6, 0.55, fill=_GRAYL)
    _txb(slide, "Les signaux FinSight sont calculés sur données fondamentales réelles (yfinance). Ils ne constituent pas un conseil en investissement.",
         1.0, 2.15, 23.4, 0.45, size=7.5, italic=True, color=_GRAYT)

    # Panel A — bandeau h=1.05 pour bien contenir titre + univers
    _rect(slide, 0.9, 2.88, 11.2, 1.05, fill=_COL_A)
    _txb(slide, D["sector_a"][:22], 1.0, 2.96, 11.0, 0.48, size=11, bold=True, color=_WHITE)
    _txb(slide, D["universe_a"], 1.0, 3.50, 11.0, 0.40, size=8.5, color=_WHITE)

    _rect(slide, 0.9, 4.0, 11.2, 1.3, fill=D["sig_a_col"])
    _txb(slide, f"● {D['sig_a_lbl']}", 0.9, 4.2, 11.2, 0.75, size=20, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)

    # KPIs A
    kpi_rows_a = [
        ("Score FinSight", f"{sa.get('score', 0)}/100"),
        ("P/E Médian", _fmt_simple(sa.get("pe"), x=True)),
        ("Croissance Rev.", _fmt(sa.get("revg"), pct=True)),
        ("Marge EBITDA", _fmt_simple(sa.get("em"), pct=True)),
        ("ROE", _fmt_simple(sa.get("roe"), pct=True)),
        ("Perf. 52S", _fmt(sa.get("mom"), pct=True)),
    ]
    for j, (label, val) in enumerate(kpi_rows_a):
        y = 5.5 + j * 0.75
        _rect(slide, 0.9, y, 11.2, 0.68, fill=_GRAYL if j % 2 == 0 else _WHITE)
        _txb(slide, label, 1.0, y + 0.1, 7.0, 0.52, size=8.5, color=_GRAYT)
        _txb(slide, val, 1.0, y + 0.1, 11.0, 0.52, size=8.5, bold=True, color=_NAVY, align=PP_ALIGN.RIGHT)

    # Separateur
    _rect(slide, 12.7, 2.7, 0.04, 7.5, fill=_GRAYM)

    # Panel B
    _rect(slide, 13.1, 2.88, 11.2, 1.05, fill=_COL_B)
    _txb(slide, D["sector_b"][:22], 13.2, 2.96, 11.0, 0.48, size=11, bold=True, color=_WHITE)
    _txb(slide, D["universe_b"], 13.2, 3.50, 11.0, 0.40, size=8.5, color=_WHITE)

    _rect(slide, 13.1, 4.0, 11.2, 1.3, fill=D["sig_b_col"])
    _txb(slide, f"● {D['sig_b_lbl']}", 13.1, 4.2, 11.2, 0.75, size=20, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)

    kpi_rows_b = [
        ("Score FinSight", f"{sb.get('score', 0)}/100"),
        ("P/E Médian", _fmt_simple(sb.get("pe"), x=True)),
        ("Croissance Rev.", _fmt(sb.get("revg"), pct=True)),
        ("Marge EBITDA", _fmt_simple(sb.get("em"), pct=True)),
        ("ROE", _fmt_simple(sb.get("roe"), pct=True)),
        ("Perf. 52S", _fmt(sb.get("mom"), pct=True)),
    ]
    for j, (label, val) in enumerate(kpi_rows_b):
        y = 5.5 + j * 0.75
        _rect(slide, 13.1, y, 11.2, 0.68, fill=_GRAYL if j % 2 == 0 else _WHITE)
        _txb(slide, label, 13.2, y + 0.1, 7.0, 0.52, size=8.5, color=_GRAYT)
        _txb(slide, val, 13.2, y + 0.1, 11.0, 0.52, size=8.5, bold=True, color=_GREEN, align=PP_ALIGN.RIGHT)

    # Recommandation IA — box LLM standard
    alloc_text = D.get("llm", {}).get("allocation_read") or ""
    if alloc_text:
        _llm_box(slide, 0.9, 10.20, 23.6, 3.10, "Recommandation",
                 alloc_text, fontsize=9)

    # Note de bas de page
    _txb(slide, "Score FinSight : indicateur proprietaire 0-100 (value + growth + qualité + momentum, 25 pts chacun). Signal : Surpondérer >=65, Neutre 45-64, Sous-pondérer <45.",
         0.9, 13.38, 23.6, 0.3, size=6.5, italic=True, color=_GRAYD, wrap=True)


def _s15b_verdict(prs, D):
    """Slide 20 — Verdict Comparatif & Conviction d'Allocation."""
    slide = _blank(prs)
    sa, sb = D["sa"], D["sb"]
    _header(slide, "Verdict Comparatif  —  Conviction d'Allocation",
            f"Secteur a privilegier, arguments décisionnels et conditions d'invalidation", 4)
    _footer(slide, D)

    score_a = sa.get("score", 0)
    score_b = sb.get("score", 0)
    winner = D["sector_a"] if score_a >= score_b else D["sector_b"]
    loser  = D["sector_b"] if score_a >= score_b else D["sector_a"]
    winner_sig = D["sig_a_lbl"] if score_a >= score_b else D["sig_b_lbl"]
    winner_col = _COL_A if score_a >= score_b else _GOLD
    winner_score = max(score_a, score_b)
    loser_score  = min(score_a, score_b)
    gap = abs(score_a - score_b)

    # Box verdict principal
    _rect(slide, 0.9, 2.1, 23.6, 1.4, fill=winner_col)
    _txb(slide, f"SECTEUR PRIVILEGIE : {winner.upper()}", 1.0, 2.2, 23.4, 0.75,
         size=17, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)
    _txb(slide, (f"Signal : {winner_sig}  |  Score {winner_score}/100 vs {loser_score}/100 ({loser})"
                 f"  |  Écart {gap} pts"),
         1.0, 2.97, 23.4, 0.4, size=8.5, color=_WHITE, align=PP_ALIGN.CENTER)

    # Thèse d'allocation (LLM)
    alloc_thesis = D.get("llm", {}).get("verdict_read") or (
        f"{winner} présente le profil risque/rendement le plus attractif dans l'univers {D['universe_label']}. "
        f"Avec un score FinSight de {winner_score}/100, la superiorite structurelle s'exprime sur "
        f"les dimensions value, growth et qualité. "
        f"L'écart de {gap} pts face a {loser} ({loser_score}/100) legitime une surponderation "
        f"tactique sur un horizon 6-12 mois, sous reserve de stabilité macro."
    )
    # Thèse d'allocation — box LLM standard
    _llm_box(slide, 0.9, 3.70, 23.6, 2.20, "Thèse d'allocation",
             alloc_thesis[:480], fontsize=9)

    # 3 colonnes : Catalyseurs | Risques | Invalidation
    col_w3 = 7.4
    xs = [0.9, 8.6, 16.3]

    # Catalyseurs (secteur gagnant)
    _rect(slide, xs[0], 6.15, col_w3, 0.48, fill=_GREEN)
    _txb(slide, "CATALYSEURS CLES", xs[0], 6.18, col_w3, 0.4, size=8, bold=True,
         color=_WHITE, align=PP_ALIGN.CENTER)
    winner_content = _get_content(winner)
    cats = winner_content.get("catalyseurs", [])[:2]
    y_cat = 6.8
    for title, body in cats:
        _rect(slide, xs[0], y_cat, 0.2, 0.2, fill=_GREEN)
        _txb(slide, title[:36], xs[0] + 0.32, y_cat - 0.02, col_w3 - 0.32, 0.42, size=8, bold=True, color=_BLACK)
        _txb(slide, body[:110], xs[0] + 0.32, y_cat + 0.44, col_w3 - 0.32, 1.3, size=7.5, color=_GRAYT, wrap=True)
        y_cat += 1.95

    # Risques (communs ou secteur gagnant)
    _rect(slide, xs[1], 6.15, col_w3, 0.48, fill=_RED)
    _txb(slide, "RISQUES PRINCIPAUX", xs[1], 6.18, col_w3, 0.4, size=8, bold=True,
         color=_WHITE, align=PP_ALIGN.CENTER)
    weaks = _build_weaknesses(sa if score_a >= score_b else sb, winner)
    y_risk = 6.8
    for risk in weaks[:2]:
        parts = risk.split(" — ", 1) if " — " in risk else risk.split("—", 1)
        rt = parts[0].strip()[:36]
        rd = parts[1].strip()[:110] if len(parts) > 1 else risk[36:]
        _rect(slide, xs[1], y_risk, 0.2, 0.2, fill=_RED)
        _txb(slide, rt, xs[1] + 0.32, y_risk - 0.02, col_w3 - 0.32, 0.42, size=8, bold=True, color=_BLACK)
        _txb(slide, rd, xs[1] + 0.32, y_risk + 0.44, col_w3 - 0.32, 1.3, size=7.5, color=_GRAYT, wrap=True)
        y_risk += 1.95

    # Conditions d'invalidation
    _rect(slide, xs[2], 6.15, col_w3, 0.48, fill=_AMBER)
    _txb(slide, "CONDITIONS D'INVALIDATION", xs[2], 6.18, col_w3, 0.4, size=8, bold=True,
         color=_WHITE, align=PP_ALIGN.CENTER)
    conditions = [
        ("Détérioration macro", "Récession ou hausse taux materielle — revoir positionnement"),
        ("Révision benefices", "Profit warning > -10% / 2 trimestres — signal de sortie"),
    ]
    y_cond = 6.8
    for ct, cb in conditions:
        _rect(slide, xs[2], y_cond, 0.2, 0.2, fill=_AMBER)
        _txb(slide, ct[:36], xs[2] + 0.32, y_cond - 0.02, col_w3 - 0.32, 0.42, size=8, bold=True, color=_BLACK)
        _txb(slide, cb[:110], xs[2] + 0.32, y_cond + 0.44, col_w3 - 0.32, 1.3, size=7.5, color=_GRAYT, wrap=True)
        y_cond += 1.95

    # Horizon & conviction — bandeau h=0.55 pour cadrage propre
    _rect(slide, 0.9, 10.85, 23.6, 0.55, fill=_GRAYL,
          line=True, line_col=_GRAYM, line_w=0.4)
    _txb(slide, ("Horizon recommande : 6-12 mois  |  Conviction FORTE si écart score > 15 pts"
                 "  |  Rééquilibrer si écart < 5 pts  |  Révisable a chaque publication trimestrielle"),
         1.0, 10.95, 23.4, 0.40, size=8, italic=True, color=_NAVY, align=PP_ALIGN.CENTER)


def _s16_disclaimer(prs, D):
    """Slide 21 — Mentions légales & Méthodologie detaillee."""
    slide = _blank(prs)
    _rect(slide, 0, 0, 25.4, 1.8, fill=_NAVY)
    _txb(slide, "Méthodologie & Mentions Légales", 0.9, 0.3, 23.6, 1.2, size=13, bold=True, color=_WHITE)

    # Deux colonnes : Méthodologie (gauche) + Mentions légales (droite)
    col_w = 11.2
    xa, xb = 0.9, 13.1

    # === MÉTHODOLOGIE (gauche) ===
    _rect(slide, xa, 2.0, col_w, 0.42, fill=_NAVY)
    _txb(slide, "MÉTHODOLOGIE", xa, 2.03, col_w, 0.36, size=8, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)

    methodo = [
        ("Scoring FinSight (0-100)",
         "Agrégation de 4 dimensions également pondérées a 25 pts chacune : "
         "Value (P/E, EV/EBITDA, PEG, FCF Yield), Growth (CAGR revenus, EPS growth, révisions), "
         "Quality (ROE, ROIC, Piotroski F-Score, Altman Z, marges), "
         "Momentum (perf. 52S, RSI relative, révision consensus). "
         "Seuils d'allocation : >= 65 = Surpondérer, 45-64 = Neutre, < 45 = Sous-pondérer."),
        ("Indicateurs de qualité financière",
         "Piotroski F-Score (0-9) : 9 critères binaires (ROA positif, CFO > NI, levier en baisse, etc.). "
         "Altman Z-Score : modèle de detresse financière (Z>2.99 = sain, 1.81-2.99 = grise, <1.81 = distress). "
         "Beneish M-Score : modèle de detection de manipulation comptable (M < -2.22 = sain, M > -1.78 = signal). "
         "Ces scores sont retires des affichages principaux pour préserver la lisibilite mais restent "
         "calculés en arriere-plan et intégrés au scoring composite Quality."),
        ("Construction de l'univers",
         "Univers S&P 500, CAC 40, STOXX 600 ou global selon le parametre sélectionné par l'utilisateur. "
         "Toutes les sociétés du secteur avec données yfinance disponibles (min. 3 ratios renseignés). "
         "Valeurs aberrantes filtrees : P/E > 999x exclus, ROE < -500% exclus, ratios LTM uniquement. "
         "Médianes utilisees plutot que moyennes pour la robustesse aux outliers."),
        ("Sources de données",
         "yfinance (Yahoo Finance) : cours, bilan, compte de résultats, flux de trésorerie — fréquence "
         "trimestrielle ou annuelle selon disponibilite. Finnhub : news et sentiment. FMP : données "
         "supplementaires si disponibles. Perf. 52S : composite top-15 Normalisé base 100."),
        ("Limites & Mises en garde",
         "Données retardees de 24h sur yfinance free tier. Certains ratios (ROIC, Piotroski) peuvent être "
         "indisponibles pour des sociétés hors US. Médianes sectorielles masquent la dispersion intra. "
         "Les signaux sont Mécaniques, statiques (snapshot point-in-time) et non ajustés du cycle. "
         "Aucune analyse qualitative manuelle (management, gouvernance, ESG) n'est réalisée."),
    ]
    y = 2.55
    for title, text in methodo:
        _txb(slide, title, xa, y, col_w, 0.34, size=7.5, bold=True, color=_NAVY)
        _txb(slide, text, xa, y + 0.38, col_w, 1.65, size=6.5, color=_GRAYT, wrap=True)
        y += 2.05

    # === MENTIONS LÉGALES (droite) ===
    _rect(slide, xb, 2.0, col_w, 0.42, fill=_NAVYL)
    _txb(slide, "MENTIONS LÉGALES", xb, 2.03, col_w, 0.36, size=8, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)

    legals = [
        ("Caractère informatif et pedagogique",
         "Ce document est produit exclusivement a des fins d'information et de démonstration pedagogique. "
         "Il ne constitue en aucun cas un conseil en investissement, une recommandation personnalisee, "
         "une incitation a l'achat ou a la vente de valeurs mobilières, ni une offre ou sollicitation "
         "de souscription a un produit financier au sens du Règlement Général de l'AMF. "
         "L'utilisateur est seul responsable de ses décisions d'investissement et doit consulter un "
         "conseiller en investissement financier (CIF) agréé avant toute opération."),
        ("Absence de due diligence et limites de l'analyse",
         "FinSight IA est un outil algorithmique de screening basé sur des données publiques (yfinance, "
         "Finnhub, FMP). Les analyses présentées sont générées automatiquement par des modèles statistiques "
         "et un LLM, sans validation manuelle, sans rencontre avec le management, sans audit des comptes. "
         "Aucune due diligence spécifique, expertise sectorielle approfondie ou vérification croisee "
         "n'est réalisée. Les modèles peuvent contenir des biais, erreurs ou simplifications. "
         "FinSight IA et ses auteurs déclinent toute responsabilité quant aux pertes ou prejudices "
         "découlant de l'utilisation de ce document."),
        ("Risques d'investissement",
         "Tout investissement en valeurs mobilières comporte un risque de perte partielle ou totale "
         "en capital. Les performances passees ne préjugent pas des performances futures. Les "
         "conditions de marché, le contexte macroeconomique, les décisions réglementaires et les "
         "événements géopolitiques peuvent evoluer rapidement et invalider les signaux présentés. "
         "La diversification et un horizon d'investissement adapte sont recommandes."),
        ("Confidentialité et propriete intellectuelle",
         "Ce document est destine a un usage prive et confidentiel. Sa reproduction, distribution, "
         "publication ou diffusion, meme partielle, est strictement interdite sans autorisation "
         "ecrite expresse. Le scoring FinSight, la Méthodologie et les visuels sont la propriete "
         "intellectuelle exclusive de FinSight IA. Toute exploitation commerciale est prohibee. "
         "Ne pas utiliser comme base exclusive pour une décision d'investissement."),
    ]
    y = 2.55
    for title, text in legals:
        _txb(slide, title, xb, y, col_w, 0.34, size=7.5, bold=True, color=_NAVY)
        _txb(slide, text, xb, y + 0.38, col_w, 1.65, size=6.5, color=_GRAYT, wrap=True)
        y += 2.05

    _rect(slide, 0.9, 13.3, 23.6, 0.25, fill=_GRAYL)
    _txb(slide, f"Généré par FinSight IA  |  {D['date_str']}  |  Usage confidentiel",
         0.9, 13.32, 23.6, 0.22, size=7, bold=True, color=_NAVY, align=PP_ALIGN.CENTER)


# ── Generateur de textes analytiques LLM ─────────────────────────────────────

def _generate_llm_texts(D: dict) -> dict:
    """Un seul appel Mistral pour générer tous les textes analytiques du pitchbook."""
    try:
        from core.llm_provider import LLMProvider
        llm = LLMProvider(provider="mistral", model="mistral-small-latest")
        sa, sb = D["sa"], D["sb"]
        sector_a, sector_b = D["sector_a"], D["sector_b"]
        roe_a = min(float(sa.get("roe") or 0), 999.9)
        roe_b = min(float(sb.get("roe") or 0), 999.9)
        dy_a = (sa.get("div_yield") or 0) * 100
        dy_b = (sb.get("div_yield") or 0) * 100
        fy_a = sa.get("fcfy") or 0
        fy_b = sb.get("fcfy") or 0
        pt_a = (sa.get("payout") or 0) * 100
        pt_b = (sb.get("payout") or 0) * 100
        score_a = sa.get("score", 0)
        score_b = sb.get("score", 0)
        winner = sector_a if score_a >= score_b else sector_b
        prompt = (
            f"Tu es analyste financier senior chez une grande banque d'investissement. "
            f"Redige des textes analytiques précis et professionnels (style rapport JPMorgan) "
            f"pour un pitchbook comparatif sectoriel entre {sector_a} et {sector_b} "
            f"(univers : {D['universe_label']}, {D['na']} + {D['nb']} sociétés).\n\n"
            f"DONNÉES MÉDIANES :\n"
            f"- {sector_a} : Score {score_a}/100, Signal {D['sig_a_lbl']}, "
            f"P/E {sa.get('pe',0):.1f}x, EV/EBITDA {sa.get('ev_eb',0):.1f}x, EV/Rev {sa.get('ev_rev',0):.1f}x, "
            f"Croiss.Rev {sa.get('revg',0):+.1f}%, Mg.Brute {sa.get('gm',0):.1f}%, "
            f"Mg.EBITDA {sa.get('em',0):.1f}%, Mg.Nette {sa.get('nm',0):.1f}%, "
            f"ROE {roe_a:.1f}%, Perf.52S {sa.get('mom',0):+.1f}%, Beta {sa.get('beta',1.0):.2f}, "
            f"Div.Yield {dy_a:.1f}%, FCF.Yield {fy_a:.1f}%, Payout {pt_a:.0f}%\n"
            f"- {sector_b} : Score {score_b}/100, Signal {D['sig_b_lbl']}, "
            f"P/E {sb.get('pe',0):.1f}x, EV/EBITDA {sb.get('ev_eb',0):.1f}x, EV/Rev {sb.get('ev_rev',0):.1f}x, "
            f"Croiss.Rev {sb.get('revg',0):+.1f}%, Mg.Brute {sb.get('gm',0):.1f}%, "
            f"Mg.EBITDA {sb.get('em',0):.1f}%, Mg.Nette {sb.get('nm',0):.1f}%, "
            f"ROE {roe_b:.1f}%, Perf.52S {sb.get('mom',0):+.1f}%, Beta {sb.get('beta',1.0):.2f}, "
            f"Div.Yield {dy_b:.1f}%, FCF.Yield {fy_b:.1f}%, Payout {pt_b:.0f}%\n\n"
            f"RÈGLES DE REDACTION :\n"
            f"- Francais correct et complet : utilise TOUS les accents (e e e a u c i o), "
            f"les cedilles, les apostrophes droites ' et les guillemets francais << >>. "
            f"N'ecris JAMAIS sans accents — ce serait du francais casse, inacceptable en rapport IB.\n"
            f"- Style sell-side senior (JPMorgan, Goldman, Morgan Stanley Research). Prose technique, "
            f"chiffres précis, raisonnements Économiques relies aux données fournies.\n"
            f"- Cite les chiffres précis fournis ci-dessus, ne jamais en inventer ni dire 'données indisponibles'.\n"
            f"- Pas de markdown, pas d'emojis, pas de formules generiques.\n"
            f"- Maximum indique par champ.\n"
            f"- Reponds UNIQUEMENT en JSON valide, aucun texte avant ou après.\n\n"
            f'{{\n'
            f'  "exec_summary": "3 phrases max (400 car.) : synthèse globale, qui privilegier et pourquoi, écarts cles",\n'
            f'  "valuation_read": "2 phrases (300 car.) : analyse multiples P/E EV/EBITDA, prime/décote et implications",\n'
            f'  "margins_read": "2-3 phrases (350 car.) : qualité opérationnelle, qui généré plus de valeur et pourquoi, implications pour l investisseur",\n'
            f'  "capital_alloc_read": "2-3 phrases (380 car.) : politique de distribution, FCF yield vs dividende, soutenabilite et implications allocation",\n'
            f'  "growth_read": "2-3 phrases (350 car.) : dynamique de croissance revenue, momentum cours 52S, divergence ou convergence et implications",\n'
            f'  "scoring_read": "2-3 phrases (350 car.) : analyse radar multidimensionnel, forces et faiblesses relatives par dimension, signal résultant",\n'
            f'  "top_a_read": "2-3 phrases (300 car.) : profil des meilleurs acteurs {sector_a}, caractéristiques communes, champion sectoriel",\n'
            f'  "top_b_read": "2-3 phrases (300 car.) : profil des meilleurs acteurs {sector_b}, caractéristiques communes, champion sectoriel",\n'
            f'  "Thèse_a": "2-3 phrases (320 car.) : Thèse investissement {sector_a}, 2-3 arguments structurels avec chiffres",\n'
            f'  "Thèse_b": "2-3 phrases (320 car.) : Thèse investissement {sector_b}, 2-3 arguments structurels avec chiffres",\n'
            f'  "Thèse_long_a": "3 arguments {sector_a} format Titre : Description (380 car. total), separes par . ",\n'
            f'  "Thèse_long_b": "3 arguments {sector_b} format Titre : Description (380 car. total), separes par . ",\n'
            f'  "risques_principaux": "2-3 risques communs ou spécifiques (280 car.) : risque macro, valorisation, execution",\n'
            f'  "verdict_read": "3 phrases (420 car.) : verdict allocation, {winner} recommande, raisons quantifiees, horizon et conditions",\n'
            f'  "allocation_read": "3 phrases (400 car.) : recommandation portefeuille detaillee, surponderation cles, ratio risque/rendement"\n'
            f'}}'
        )
        import json, re
        resp = llm.generate(prompt, max_tokens=2200)
        m = re.search(r'\{.*\}', resp, re.DOTALL)
        if m:
            data = json.loads(m.group(0))
            log.info("[cmp_secteur_pptx] LLM texts OK (%d champs)", len(data))
            return data
    except Exception as e:
        log.warning("[cmp_secteur_pptx] LLM texts génération failed: %s", e)
    return {}


# ── Classe principale ─────────────────────────────────────────────────────────

class CmpSecteurPPTXWriter:
    """Generateur de pitchbook comparatif sectoriel 21 slides."""

    @staticmethod
    def generate(
        tickers_a: list[dict],
        sector_a: str,
        universe_a: str,
        tickers_b: list[dict],
        sector_b: str,
        universe_b: str,
        output_path: Optional[str] = None,
    ) -> bytes:
        """
        Genere le PPTX comparatif et retourne les bytes.
        Si output_path est fourni, sauvegarde aussi sur disque.
        """
        D = _prepare_data(tickers_a, sector_a, universe_a, tickers_b, sector_b, universe_b)
        D["llm"] = _generate_llm_texts(D)

        # Fetch des cours 52S (une seule fois avant les slides)
        log.info("[CmpSecteurPPTXWriter] Fetch cours 52S...")
        D["perf_a_52w"], D["perf_b_52w"] = _fetch_price_52w_pptx(tickers_a, tickers_b)

        prs = Presentation()
        prs.slide_width  = _SW
        prs.slide_height = _SH

        # Section 1 — Vue d'ensemble
        _s01_cover(prs, D)                # Slide 1  — Cover
        _s02_exec_summary(prs, D)         # Slide 2  — Executive Summary
        _s03_sommaire(prs, D)             # Slide 3  — Sommaire

        # Section 2 — Profil & Valorisation
        _chapter_divider(prs, "01", "Profil & Valorisation",
                         "Structure des secteurs, multiples et qualité des bilans")
        _s05_profil(prs, D)               # Slide 5  — Profil cote a cote
        _s06_valorisation(prs, D)         # Slide 6  — P/E, EV/EBITDA
        _s07_marges(prs, D)               # Slide 7  — Marges & Rentabilite
        _s07b_capital_alloc(prs, D)       # Slide 8  — Capital Allocation

        # Section 3 — Performance & Scoring
        _chapter_divider(prs, "02", "Performance & Scoring",
                         "Cours 52 semaines, revenue growth et scoring multidimensionnel")
        _s08_croissance(prs, D)           # Slide 10 — Cours 52S comparatif
        _s09_scoring(prs, D)              # Slide 11 — Radar scoring

        # Section 4 — Top Acteurs & Risques
        _chapter_divider(prs, "03", "Top Acteurs & Risques",
                         "Meilleures sociétés et conditions d'invalidation par secteur")
        _s10_top_a(prs, D)                # Slide 13 — Top acteurs A
        _s11_top_b(prs, D)                # Slide 14 — Top acteurs B
        _s12_risques_a(prs, D)            # Slide 15 — Risques A
        _s13_risques_b(prs, D)            # Slide 16 — Risques B

        # Section 5 — Synthèse & Décision
        _chapter_divider(prs, "04", "Synthèse & Décision",
                         "Thèses d'Investissement, verdict comparatif et recommandation")
        _s14_synthese(prs, D)             # Slide 18 — THESE LONG A+B + RISQUES
        _s15_allocation(prs, D)           # Slide 19 — Recommandation allocation
        _s15b_verdict(prs, D)             # Slide 20 — Verdict Comparatif (NEW)
        _s16_disclaimer(prs, D)           # Slide 21 — Methodologie & Mentions legales

        buf = io.BytesIO()
        prs.save(buf)
        pptx_bytes = buf.getvalue()

        if output_path:
            Path(output_path).write_bytes(pptx_bytes)
            log.info("[CmpSecteurPPTXWriter] Sauvegarde : %s (%d Ko)", output_path, len(pptx_bytes) // 1024)

        return pptx_bytes
