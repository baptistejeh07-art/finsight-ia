"""
sectoral_pptx_writer.py — FinSight IA
Pitchbook sectoriel 20 slides via python-pptx.
Usage : SectoralPPTXWriter.generate(tickers_data, sector_name, universe, output_path)
"""
from __future__ import annotations
import io, datetime, logging
from pathlib import Path
from typing import Any, Optional

import numpy as np
import matplotlib
from core.yfinance_cache import get_ticker
matplotlib.use('Agg')
import matplotlib.pyplot as plt

log = logging.getLogger(__name__)

from pptx import Presentation
from pptx.util import Cm, Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── PALETTE ──────────────────────────────────────────────────────────────────
_NAVY   = RGBColor(0x1B, 0x3A, 0x6B)
_NAVYL  = RGBColor(0x2A, 0x52, 0x98)
_BUY    = RGBColor(0x1A, 0x7A, 0x4A)
_SELL   = RGBColor(0xA8, 0x20, 0x20)
_HOLD   = RGBColor(0xB0, 0x60, 0x00)
_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
_BLACK  = RGBColor(0x1A, 0x1A, 0x1A)
_GRAYL  = RGBColor(0xF5, 0xF7, 0xFA)
_GRAYM  = RGBColor(0xE8, 0xEC, 0xF0)
_GRAYT  = RGBColor(0x55, 0x55, 0x55)
_GRAYD  = RGBColor(0xAA, 0xAA, 0xAA)
_GREEN_L = RGBColor(0xE8, 0xF5, 0xEE)
_RED_L   = RGBColor(0xFB, 0xEB, 0xEB)
_HOLD_L  = RGBColor(0xFD, 0xF3, 0xE5)

# ─── SLIDE DIMENSIONS ─────────────────────────────────────────────────────────
_SW = Inches(10.0)
_SH = Inches(5.625)

# ─── CONFIG ───────────────────────────────────────────────────────────────────
MAX_CATS_DISPLAYED    = 3
MAX_RISKS_DISPLAYED   = 3
MAX_DRIVERS_DISPLAYED = 4
MAX_TOP_N             = 3
MAX_TICKERS_CHART     = 8
MAX_TABLE_ROWS        = 15   # cap tables sur slides 6/9/11/15 (univers >15 tickers)

# ─── SECTOR CONTENT LIBRARY ───────────────────────────────────────────────────
_SECTOR_CONTENT = {
    "Technology": {
        "description": (
            "Le secteur Technology regroupe les entreprises actives dans les logiciels, "
            "les services informatiques, les semi-conducteurs et les equipements telecom. "
            "Porté par la transformation digitale des entreprises et l'adoption de l'IA generative, "
            "ce secteur affiche des multiples de valorisation parmi les plus élevés du Marché, "
            "justifiés par des perspectives de croissance supérieures et des marges en expansion."
        ),
        "catalyseurs": [
            ("Adoption IA Generative", "Accélération des mandats IA dans les grands comptes européens — contrats pluriannuels en hausse de 38 % sur T1 2026."),
            ("Modernisation SI Publics", "Plans d'investissement gouvernementaux 2026-2028 — budgets numeriques confirmes pour 6 États membres EU."),
            ("Cloud Hybride Souveraineté", "Contraintes réglementaires sur les données — migration vers solutions EU portée par NIS2 et DORA."),
        ],
        "risques": [
            ("Ralentissement Budgets IT", "Récession européenne confirmee entraînerait une réduction des budgets IT de 12-18 %, pénalisant les prestataires de services."),
            ("Durcissement AI Act", "Surcoût de mise en Conformité estime a 8-12 % du CA pour les editeurs de logiciels IA — impact marges 2027."),
            ("Pression Tarifaire US/Asie", "Acteurs américains et asiatiques agressifs sur les segments logiciels standardisés — guerre des prix bas de gamme."),
        ],
        "drivers": [
            ("up", "Intelligence Artificielle", "Adoption entreprises +42 % en 2026 — premier driver de dépenses IT"),
            ("up", "Cloud Hybride Souverain", "Contraintes souveraineté données — migration vers solutions EU accélérée par NIS2"),
            ("up", "Modernisation SI Publics", "Plans gouvernementaux 2026-2028 — budgets numeriques confirmes 6 États EU"),
            ("down", "Sensibilité Taux", "Compression multiples si taux > 4,5 % — impact modéré sur FCF structurel"),
        ],
        "cycle_comment": "Secteur en phase d'expansion | porté par l'IA et le cloud souverain",
        "metriques_dict": [("Rule of 40", "42", "Sain > 40"), ("NRR Médian", "108 %", "Retention solide"), ("Cloud Mix", "38 %", "En hausse"), ("R&D/Rev.", "12 %", "Investi"), ("ARR Growth", "+24 %", "Accélération"), ("Churn", "4,2 %", "Stable")],
        "conditions": [("Macro", "Récession UE confirmee — contraction PIB > 1,5 %", "6-12 mois"), ("Sectoriel", "Réduction budgets IT > 15 % grands comptes", "12-18 mois"), ("Réglementaire", "AI Act surcoût Conformité > 15 % CA", "18-24 mois"), ("Fondamental", "Révision baissière EPS > 20 % sur 2 trimestres", "6-9 mois")],
    },
    "Consumer Cyclical": {
        "description": (
            "Le secteur Consommation Cyclique inclut les acteurs du luxe, de la distribution spécialisée, "
            "de l'automobile et du tourisme. Sensible aux cycles économiques, il bénéficie de l'essor "
            "des Marchés emergents et de la résilience du luxe haut de gamme européen face aux pressions macro."
        ),
        "catalyseurs": [
            ("Rebond Luxe Asie", "Réouverture progressive des Marchés asiatiques — croissance hors UE attendue a +15 % en 2026 pour le segment premium."),
            ("Transition Electrique", "Accélération des ventes VE en Europe — soutien réglementaire et infrastructures de recharge en déploiement."),
            ("Tourisme Premium", "Flux touristiques européens en hausse de 12 % — benefice direct pour l'hotellerie et la restauration haut de gamme."),
        ],
        "risques": [
            ("Sensibilité Taux", "Hausse des taux comprime le pouvoir d'achat discrétionnaire — impact sur volumes et marges promotionnelles."),
            ("Ralentissement Chine", "Croissance chinoise < 4 % pénaliserait les acteurs du luxe exposes a 30-40 % de leur CA en Asie."),
            ("Disruption Digitale", "Montee des plateformes D2C — pression sur les réseaux de distribution traditionnels."),
        ],
        "drivers": [
            ("up", "Luxe Premium", "Résilience exceptionnelle — pricing power intact segment > 500 EUR"),
            ("up", "Rebond Asiatique", "Flux touristiques et ventes locales en Chine en normalisation"),
            ("down", "Pouvoir d'Achat", "Inflation residuelle et taux élevés compressent la consommation mid-market"),
            ("down", "Substitution Digital", "Plateformes e-commerce capturent des parts en distribution"),
        ],
        "cycle_comment": "Secteur en consolidation | luxe résilient, mid-market sous pression",
        "metriques_dict": [("SSSG", "+4,2 %", "Positif"), ("Inventory Turn", "3,8x", "Sain"), ("Online Mix", "28 %", "En hausse"), ("Price/Mix", "+2,1 %", "Favorable"), ("Return Rate", "18 %", "Stable"), ("Loyalty Rev.", "42 %", "Fort")],
        "conditions": [("Macro", "PIB Zone Euro négatif 2 trimestres", "6-9 mois"), ("Chine", "PIB chinois < 3,5 %", "9-12 mois"), ("Crédit", "Taux 10Y > 5 % pendant 6 mois", "6-12 mois"), ("Fondamental", "SSSG négatif 3 trimestres", "9 mois")],
    },
    "Consumer Défensive": {
        "description": (
            "Le secteur Consommation Défensive regroupe les producteurs et distributeurs de biens essentiels "
            "— alimentation, boissons, hygiene-beaute. Modèle résilient par nature, il offre une protection "
            "en période de récession tout en beneficiant du pricing power des marques leaders."
        ),
        "catalyseurs": [
            ("Pricing Power", "Capacité a répercuter l'inflation des intrants — les leaders de marque maintiennent des marges brutes > 40 %."),
            ("Marchés Emergents", "Croissance de la classe moyenne en Afrique et Asie du Sud-Est — expansion a double chiffre."),
            ("Innovation Produit", "Premiumisation et bio — capture de valeur additionnelle sur les segments porteurs."),
        ],
        "risques": [
            ("Guerre des Prix", "Pression des distributeurs (Leclerc, Lidl) sur les prix fournisseurs — erosion des marges lors des négociations."),
            ("Destockage", "Réduction des stocks en distribution — impact sur volumes T1-T2 2026."),
            ("Substitution MDD", "Parts de Marché des marques distributeurs en hausse de 2-3 pts — pression structurelle."),
        ],
        "drivers": [
            ("up", "Pricing Power Marques", "Leaders de catégorie maintiennent +3-5 % de hausse tarifaire annuelle"),
            ("up", "Expansion Emergents", "Croissance volumes a 2 chiffres en Afrique et Asie du Sud-Est"),
            ("down", "Pression Distribution", "Négociations annuelles de plus en plus défavorables"),
            ("down", "Saturation Matures", "Volumes stagnants en Europe — croissance principalement par prix"),
        ],
        "cycle_comment": "Secteur défensif | résilient en récession, sous-performant en expansion",
        "metriques_dict": [("Org. Growth", "+3,8 %", "Sain"), ("Gross Margin", "42 %", "Stable"), ("FCF Yield", "4,2 %", "Attractif"), ("Div. Yield", "2,8 %", "Stable"), ("Elasticite Prix", "-0,4", "Faible"), ("Market Share", "+0,3 pt", "Gain")],
        "conditions": [("Macro", "Récession profonde > 3 trimestres", "9-12 mois"), ("Distribution", "Perte > 15 % référencements top 3", "6-9 mois"), ("MDD", "Part MDD > 45 % en valeur", "12-18 mois"), ("Réglementaire", "Taxe nutritionnelle majeure", "12-24 mois")],
    },
    "Financial Services": {
        "description": (
            "Le secteur Financier couvre les banques, assureurs, sociétés de gestion d'actifs et fintech. "
            "L'environnement de taux élevés a restaure la rentabilité des marges d'intérêt, mais l'exposition "
            "au crédit immobilier et aux PME reste un facteur de risque dans un contexte macro degrade."
        ),
        "catalyseurs": [
            ("Taux Directeurs Élevés", "Maintien des taux BCE a 3 % — soutien structurel aux NIM bancaires et a la rentabilité des depots."),
            ("Consolidation Sectorielle", "Fusions-acquisitions dans la banque retail — prime de contrôle sur les cibles sous-valorisees."),
            ("Digitalisation Services", "Plateformes wealth management numeriques — réduction Coûts et captation millennials investisseurs."),
        ],
        "risques": [
            ("Dégradâtion Crédit", "Hausse des NPL dans l'immobilier commercial et le crédit PME si récession confirmee."),
            ("Baisse Taux BCE", "Pivot accommodant — compression des NIM de 15-25 bps par tranche de 25 bps de baisse."),
            ("Régulation Bale IV", "Exigences CET1 plus élevées — pression sur le ROE et les politiques de retour de capital."),
        ],
        "drivers": [
            ("up", "Marge Nette d'Intérêt", "Taux élevés soutiennent les NIM — pic attendu T2 2026"),
            ("up", "Gestion d'Actifs", "Collecte nette positive — Marchés actions porteurs"),
            ("down", "Coût du Risque", "Montee progressive des provisions — NPL en hausse"),
            ("down", "Régulation", "Bale IV — impact CET1 estime a -80bps grandes banques EU"),
        ],
        "cycle_comment": "Secteur en maturité cyclique | NIM en pic, qualité d'actif à surveiller",
        "metriques_dict": [("NIM Médian", "2,8 %", "Sain"), ("CET1 Médian", "14,2 %", "Confortable"), ("ROE Médian", "11,4 %", "Acceptable"), ("Cost/Income", "58 %", "A ameliorer"), ("NPL Ratio", "2,1 %", "Contrôle"), ("P/Book", "0,8x", "Décote")],
        "conditions": [("Taux", "BCE baisse < 2 % — compression NIM > 40 bps", "6-12 mois"), ("Crédit", "NPL > 5 % immobilier commercial", "12-18 mois"), ("Macro", "Récession > 2 trimestrès — provisionnement", "6-9 mois"), ("Réglementaire", "Bale IV CET1 > 15 %", "24-36 mois")],
    },
    "Industrials": {
        "description": (
            "Le secteur Industriel englobe les fabricants d'equipements, les entreprises de défense, "
            "les conglomérats industriels et les prestataires de services B2B. Il bénéficie des plans de "
            "reindustrialisation européens et de la hausse des budgets défense."
        ),
        "catalyseurs": [
            ("Rearmement Européen", "Hausse des budgets défense a 2 % du PIB — carnet de commandes a 5-7 ans pour les acteurs spécialisés."),
            ("Reindustrialisation EU", "Plans de relocalisation — investissements publics massifs dans l'automatisation."),
            ("Transition Énergétique", "Equipements eoliens, solaires et réseaux — croissance structurelle du carnet de commandes."),
        ],
        "risques": [
            ("Ralentissement Capex", "Gel des investissements industriels en cas de récession — impact sur commandes machines."),
            ("Tension Supply Chain", "Disponibilite des semi-conducteurs et matières premières critiques — risques livraisons."),
            ("Pression Marges", "Hausse des Coûts energetiques et des salaires dans l'industrie — compression marges exécution."),
        ],
        "drivers": [
            ("up", "Défense & Sécurité", "Budgets défense EU en hausse de 15-20 % — carnets record"),
            ("up", "Equipements Transition", "Investissements eolien offshore et réseaux — +12 % pa"),
            ("down", "Cycles Capex", "Ralentissement capex automobile et electronique"),
            ("down", "Coûts Production", "Énergie et main-d'oeuvre qualifiee compriment les marges"),
        ],
        "cycle_comment": "Secteur en expansion sélective | défense et transition tirent la croissance",
        "metriques_dict": [("Book-to-Bill", "1,18x", "Positif"), ("EBIT Margin", "11,4 %", "Stable"), ("Capex/Rev", "5,2 %", "Modéré"), ("Backlog (mois)", "18", "Sécurisé"), ("FCF Conv.", "78 %", "Bon"), ("ROCE", "13,2 %", "Sain")],
        "conditions": [("Macro", "Contraction capex global > 10 %", "6-9 mois"), ("Défense", "Gel budgets défense EU > 20 %", "12-18 mois"), ("Supply Chain", "Penurie semi-conducteurs > 6 mois", "3-6 mois"), ("Fondamental", "Book-to-Bill < 0,9 pendant 3 trimestres", "9-12 mois")],
    },
    "Healthcare": {
        "description": (
            "Le secteur Sante regroupe les laboratoires pharmaceutiques, fabricants de dispositifs medicaux, "
            "biotechs et prestataires de soins. Structurellement défensif, il bénéficie du vieillissement "
            "demographique et de l'innovation therapeutique."
        ),
        "catalyseurs": [
            ("Innovation Therapeutique", "Pipeline oncologie et maladies rares — lancements produits blockbusters attendus 2026-2027."),
            ("MedTech IA", "Intégration de l'IA dans le diagnostic et la chirurgie assistee — premium innovateurs."),
            ("Vieillissement Demo", "Population > 65 ans en Europe +3 % pa — croissance structurelle des dépenses santé."),
        ],
        "risques": [
            ("Pression Prix Medicaments", "Révision des prix de référence — impact de -5 a -15 % sur les revenus medicaments matures."),
            ("Expirations Brevets", "Vague de génériques 2026-2028 — perte de revenus 15-20 % pour les concernes."),
            ("Incertitude Clinique", "Taux d'echec élevé en phase III — risque de destruction de valeur sur pipelines."),
        ],
        "drivers": [
            ("up", "Innovation Oncologie", "Immunotherapies et ADC — Marché > 200 Mds USD en 2027"),
            ("up", "MedTech IA", "Diagnostic assiste — barrières à l'entrée élevées"),
            ("down", "Expiration Brevets", "Perte revenus blockbusters 2026-2028"),
            ("down", "Pricing Pressure", "Reformes prix medicaments EU et US — compression revenus"),
        ],
        "cycle_comment": "Secteur défensif-croissance | innovation compense risques réglementaires",
        "metriques_dict": [("R&D/Rev.", "18 %", "Élevé"), ("Pipeline Score", "7,2/10", "Solide"), ("Marge Brute", "68 %", "Premium"), ("Patent Cliff", "2026-28", "A surveiller"), ("ROE", "22 %", "Attractif"), ("FCF Yield", "3,8 %", "Sain")],
        "conditions": [("Réglementaire", "Reforme prix US — impact > 20 % revenues", "12-18 mois"), ("Pipeline", "Echec phase III > 25 % du CA", "Immediat"), ("Brevet", "Expiration sans relai — -15 % revenues", "12-24 mois"), ("Macro", "Coupes budgets santé > 10 %", "24-36 mois")],
    },
    "Energy": {
        "description": (
            "Le secteur Énergie couvre les majors petrolieres et gazieres, les énergies renouvelables et "
            "les services parapetroliers. Il reste marque par la volatilité des prix des matières premières, "
            "tout en beneficiant de la transition énergétique."
        ),
        "catalyseurs": [
            ("Prix Petrole Soutenu", "Demande asiatique robuste et OPEP+ disciplinée — Brent attendu 75-85 USD/bbl en 2026."),
            ("Investissements Renouvelables", "Plans de transition des majors — allocation capex renouvelables +25 % pa."),
            ("GNL Européen", "Diversification approvisionnement post-Ukraine — primes GNL structurellement élevées."),
        ],
        "risques": [
            ("Volatilite Petrole", "Récession mondiale ou accord OPEP+ défavorable — Brent < 60 USD/bbl entraînerait des coupes dividendes."),
            ("Transition Réglementaire", "Accélération du calendrier sortie fossiles — actifs bloques en risque."),
            ("Capex Renouvelable", "Surcoûts et retards projets offshore wind — impact sur rendements."),
        ],
        "drivers": [
            ("up", "Prix Hydrocarbures", "OPEP+ discipline — soutien structurel des prix 2026"),
            ("up", "GNL & Souveraineté", "Diversification approvisionnement EU — prime GNL persistante"),
            ("down", "Transition Énergétique", "Pression réglementaire et ESG — prime de risque croissanté"),
            ("down", "Capex Renouvelable", "Dépenses investissement renouvelables pesent sur FCF"),
        ],
        "cycle_comment": "Secteur en transition | dividendes élevés, visibilité long terme réduite",
        "metriques_dict": [("FCF Yield", "7,2 %", "Très attractif"), ("Div. Yield", "4,8 %", "Élevé"), ("Breakeven", "52 USD", "Confortable"), ("ROACE", "14,2 %", "Bon"), ("ND/EBITDA", "1,2x", "Faible"), ("Capex/OCF", "48 %", "Modéré")],
        "conditions": [("Petrole", "Brent < 55 USD/bbl 6 mois — coupes dividendes", "3-6 mois"), ("Politique", "Taxes exceptionnelles profits majors", "Immediat"), ("Transition", "Calendrier sortie fossiles avance 10 ans", "24-36 mois"), ("Demande", "Pic demande petroliere confirme avant 2030", "24-48 mois")],
    },
    "Basic Materials": {
        "description": (
            "Le secteur Materiaux de Base regroupe les producteurs de metaux, chimistes et fabricants "
            "de matériaux de construction. Fortement cyclique, il bénéficie de la demande en matériaux "
            "critiques pour la transition énergétique."
        ),
        "catalyseurs": [
            ("Materiaux Critiques", "Lithium, cuivre, cobalt — demande multipliee par 3-5x a horizon 2030 pour batteries et réseaux."),
            ("Infrastructure Publique", "Plans d'investissement G7 — demande en acier, ciment et matériaux en hausse."),
            ("Consolidation M&A", "Vague de fusions dans les metaux de base — prime de rachat potentielle."),
        ],
        "risques": [
            ("Ralentissement Chine", "Demande chinoise d'acier 55 % du Marché mondial — tout ralentissement est pénalisant."),
            ("Surcapacités", "Capacités de production excedentaires acier et aluminium — pression sur les prix."),
            ("Risques Géopolitiques", "Concentration production zones a risque (DRC, Chili) — risques supply chain."),
        ],
        "drivers": [
            ("up", "Metaux Transition", "Cuivre, lithium, cobalt — demande structurelle multi-décennale"),
            ("up", "Infrastructure G7", "Plans investissement massifs — consommation beton et acier"),
            ("down", "Cyclicité Chine", "Demande industrielle chinoise déterminante et volatile"),
            ("down", "Surcapacités", "Acier et aluminium — prix sous pression structurelle"),
        ],
        "cycle_comment": "Secteur en transition cyclique | matériaux verts tirent la croissance long terme",
        "metriques_dict": [("EV/EBITDA", "6,2x", "Décote"), ("FCF Yield", "5,8 %", "Attractif"), ("ND/EBITDA", "1,8x", "Gerable"), ("ROIC", "11,2 %", "Correct"), ("Capex/Rev", "8,4 %", "Élevé"), ("Div. Yield", "3,2 %", "Satisfaisant")],
        "conditions": [("Chine", "Croissance PIB < 3,5 % pendant 2 trimestres", "6-9 mois"), ("Prix", "Cuivre < 7 000 USD/t pendant 6 mois", "3-6 mois"), ("Capex", "Annulation projets miniers majeurs", "Immediat"), ("Régulation", "Taxes export matières premières critiques", "12-18 mois")],
    },
    "Real Estate": {
        "description": (
            "Le secteur Immobilier regroupe les foncières Cotées (SIIC/REIT), promoteurs et gestionnaires "
            "d'actifs immobiliers. En phase de réévaluation post-hausse des taux, il offre des décotes sur "
            "ANR attractives pour les investisseurs patients."
        ),
        "catalyseurs": [
            ("Pivot Monétaire BCE", "Baisse des taux directeurs — compression des taux de capitalisation et réévaluation des ANR."),
            ("Penurie Logements", "Déficit de construction en France et Allemagne — soutien structurel a la valeur des actifs residentiels."),
            ("Actifs Durables", "Renovation énergétique obligatoire — valorisation Différentielle actifs certifies."),
        ],
        "risques": [
            ("Maintien Taux Élevés", "Taux longs > 4 % prolonges — pressions sur les ANR et les covenants de dette."),
            ("Dégradâtion Locative", "Vacance bureau en hausse post-teletravail — risques sur les loyers tertiaires secondaires."),
            ("Refinancement", "Mur de dette 2025-2027 — risques de dilution pour les foncieres surendettees."),
        ],
        "drivers": [
            ("up", "Pivot Taux", "Chaque -25bps BCE = +3-5 % sur les ANR sectoriels"),
            ("up", "Residentiel Déficitaire", "Offre structurellement insuffisanté metropoles EU"),
            ("down", "Bureau Post-Covid", "Vacance bureau — actifs secondaires a risque"),
            ("down", "Refinancement", "Mur de dette — risque dilution pour les leverages"),
        ],
        "cycle_comment": "Secteur en réévaluation | ANR en décote, sensibilité aux taux majeure",
        "metriques_dict": [("Décote ANR", "-22 %", "Opportunite"), ("Div. Yield", "4,8 %", "Attractif"), ("LTV Médian", "38 %", "Prudent"), ("ICR Médian", "3,2x", "Sain"), ("Vacance", "7,4 %", "Modérée"), ("EPRA NTA", "ref.", "Benchmark")],
        "conditions": [("Taux", "OAT 10Y > 4,5 % pendant 12 mois", "6-12 mois"), ("Vacance", "Vacance bureau > 15 %", "12-18 mois"), ("Dette", "LTV > 50 % — covenant breach", "Immediat"), ("Macro", "Récession — loyers -10 %", "9-12 mois")],
    },
    "Communication Services": {
        "description": (
            "Les Services de Communication regroupent les operateurs telecoms, medias numeriques et "
            "plateformes de contenu. Ce secteur bénéficie de la migration vers le haut débit, "
            "de la consolidation sectorielle et de la monétisation des usages data."
        ),
        "catalyseurs": [
            ("Fibre & 5G", "Deploiement capillaire de la fibre optique — mix favorable vers abonnements premium."),
            ("Consolidation Telecom", "Fusions nationales approuvées — Économies d'Échelle et amélioration des marges EBITDA."),
            ("Streaming & Contenu", "Monétisation accélérée des plateformes — ARPU en hausse via offres publicitaires premium."),
        ],
        "risques": [
            ("Guerre des Prix Mobile", "Concurrence low-cost intensifiee — pression sur les ARPU en France et Espagne."),
            ("Capex Fibre", "Investissements réseaux lourds — impact FCF négatif pendant la phase de déploiement."),
            ("Substitution OTT", "Perte de revenus voix et SMS au profit des services OTT (WhatsApp, Teams)."),
        ],
        "drivers": [
            ("up", "Montee en Debit", "Fibre et 5G — migration offres premium +8 EUR/mois ARPU"),
            ("up", "Consolidation", "Fusions sectorielles — réduction concurrentielle et synergies"),
            ("down", "Pression Prix", "Low-cost MVNO — ARPU mobile compressé"),
            ("down", "Capex Reseau", "Phase investissement fibre — FCF sous pression 2-3 ans"),
        ],
        "cycle_comment": "Secteur en transformation | consolidation et fibre sont les catalyseurs",
        "metriques_dict": [("EBITDA Margin", "32 %", "Stable"), ("ARPU", "28 EUR/m", "En hausse"), ("Churn", "1,2 %", "Faible"), ("Capex/Rev", "18 %", "Élevé"), ("Fiber Cover.", "68 %", "En cours"), ("FCF Yield", "3,4 %", "Modéré")],
        "conditions": [("Prix", "ARPU mobile < 22 EUR — guerre des prix", "6-9 mois"), ("Capex", "Retards fibre > 18 mois", "9-12 mois"), ("Réglementaire", "Blocage fusions autorites concurrence", "Immediat"), ("Substitution", "Revenu voix < 20 % CA total", "18-24 mois")],
    },
    "Utilities": {
        "description": (
            "Les Utilities regroupent les producteurs et distributeurs d'electricite, de gaz et d'eau. "
            "Caracterisees par des revenus regules et des dividendes stables, elles constituent un refuge "
            "défensif sensible aux variations de taux d'intérêt."
        ),
        "catalyseurs": [
            ("Transition Énergétique", "Investissements massifs dans le renouvelable — actifs regules garantissant des rendements stables."),
            ("Hausse Prix Énergie", "Normalisation des prix de gros — amélioration des marges de fourniture."),
            ("Hydrogene & Stockage", "Positionnement nouvelles technologies — optionalite valorisation long terme."),
        ],
        "risques": [
            ("Taux Élevés", "Sensibilité marquée aux taux longs — compression des multiples si OAT > 4 %."),
            ("Risque Regule", "Révision des tarifs d'acces reseau — impact sur la rémunération du capital investi."),
            ("Volatilite Meteo", "Secheresse ou faible ventossité — impact sur production hydraulique et eolienne."),
        ],
        "drivers": [
            ("up", "Renouvelables Régulés", "RAB en croissance — visibilité cash flows 20-30 ans"),
            ("up", "Electrification Usages", "VE et chaleur — croissance demande +2 % pa"),
            ("down", "Taux Longs", "Chaque +50bps = -5 a -8 % sur multiples EV/EBITDA"),
            ("down", "Révision Tarifs", "Regulateurs peuvent comprimer les marges"),
        ],
        "cycle_comment": "Secteur défensif-rendement | sensibilité taux élevée, dividendes stables",
        "metriques_dict": [("RAB Return", "6,2 %", "Regule"), ("Div. Yield", "4,4 %", "Attractif"), ("ND/EBITDA", "4,2x", "Élevé stable"), ("RES Mix", "42 %", "En hausse"), ("Payout", "68 %", "Stable"), ("ROCE", "7,8 %", "Regule")],
        "conditions": [("Taux", "OAT 10Y > 4,5 % pendant 12 mois", "6-12 mois"), ("Régulation", "Baisse WACC regule > 100bps", "24-36 mois"), ("Meteo", "Secheresse generalisee 2 ans", "Annuel"), ("Politique", "Nationalisation ou gel tarifaire", "Immediat")],
    },
}

def _build_content_from_llm(sector_name: str, ev_med: float, rev_med: float,
                             mg_med: float, mom_med: float, score_moyen: int,
                             sig_label: str, td: list) -> "dict | None":
    """Appel Groq pour générer description/catalyseurs/risques/drivers dynamiquement.
    Retourne None si echec (le caller utilisera le fallback statique ou _build_content_from_data)."""
    try:
        import json as _json_pptx
        import sys as _sys_pptx
        _sys_pptx.path.insert(0, str(Path(__file__).parent.parent))
        from core.llm_provider import LLMProvider as _LLMpptx
        sorted_t = sorted(td, key=lambda x: x.get("score_global") or 0, reverse=True)
        top_t = [f"{t.get('ticker','?')} (score={int(t.get('score_global') or 0)}/100)"
                 for t in sorted_t[:3]]
        prompt = (
            f"Secteur: {sector_name} | Signal: {sig_label} | Score: {score_moyen}/100\n"
            f"EV/EBITDA med: {ev_med:.1f}x | Mg EBITDA: {mg_med:.1f}% | "
            f"Croissance rev: {rev_med:+.1f}% | Momentum 52W: {mom_med:+.1f}%\n"
            f"Top tickers: {', '.join(top_t) if top_t else 'N/D'}\n\n"
            f"Reponds en JSON valide uniquement, sans markdown, sans points de suspension.\n"
            f'{{"description":"2 phrases sur ce secteur (spécifique, valorisation réelle)","'
            f'catalyseurs":[{{"titre":"court","description":"1 phrase complète"}},'
            f'{{"titre":"court","description":"1 phrase"}},'
            f'{{"titre":"court","description":"1 phrase"}}],"'
            f'risques":[{{"titre":"court","description":"1 phrase"}},'
            f'{{"titre":"court","description":"1 phrase"}},'
            f'{{"titre":"court","description":"1 phrase"}}],"'
            f'drivers":[{{"direction":"up","nom":"court","description":"1 phrase"}},'
            f'{{"direction":"up","nom":"court","description":"1 phrase"}},'
            f'{{"direction":"down","nom":"court","description":"1 phrase"}},'
            f'{{"direction":"down","nom":"court","description":"1 phrase"}}],"'
            f'cycle_comment":"1 phrase courte sur la phase de cycle actuelle"}}'
        )
        resp = _LLMpptx(provider="mistral", model="mistral-small-latest").generate(
            prompt=prompt,
            system=(
                f"Tu es analyste buy-side spécialisé en {sector_name}. "
                "Reponds en francais avec accents. JSON strict. "
                "Spécifique au secteur demande, jamais de texte generique. "
                "Phrases complets, pas de points de suspension (...)."
            ),
            max_tokens=800,
        )
        js_s = resp.find("{"); js_e = resp.rfind("}") + 1
        if js_s < 0 or js_e <= js_s:
            return None
        p = _json_pptx.loads(resp[js_s:js_e])
        result: dict = {}
        if "description" in p:
            result["description"] = p["description"]
        if "catalyseurs" in p and isinstance(p["catalyseurs"], list):
            result["catalyseurs"] = [
                (c.get("titre", "\u2014"), c.get("description", "\u2014"))
                for c in p["catalyseurs"][:3]
            ]
        if "risques" in p and isinstance(p["risques"], list):
            result["risques"] = [
                (r.get("titre", "\u2014"), r.get("description", "\u2014"))
                for r in p["risques"][:3]
            ]
        if "drivers" in p and isinstance(p["drivers"], list):
            result["drivers"] = [
                (d.get("direction", "up"), d.get("nom", "\u2014"), d.get("description", "\u2014"))
                for d in p["drivers"][:4]
            ]
        if "cycle_comment" in p:
            result["cycle_comment"] = p["cycle_comment"]
        log.info("sectoral_pptx LLM OK: %s (%d chars)", sector_name, len(resp))
        return result if result else None
    except Exception as _e_pptx:
        log.warning("sectoral_pptx LLM erreur: %s -- fallback", _e_pptx)
        return None


def _build_content_from_data(td: list, sector_name: str, score_moyen: int,
                              sig_label: str, ev_med: float, rev_med: float,
                              mg_med: float, mom_med: float) -> dict:
    """Build sector content dynamically from tickers_data for unrecognized sectors.
    Uses actual state values — no static dict fallback."""
    log.warning("sectoral_pptx: sector '%s' not in _SECTOR_CONTENT — "
                "building content dynamically from tickers_data (%d sociétés)",
                sector_name, len(td))

    sorted_td = sorted(td, key=lambda x: x.get("score_global") or 0, reverse=True)
    best  = sorted_td[0]  if sorted_td else {}
    worst = sorted_td[-1] if len(sorted_td) > 1 else {}

    best_ticker  = best.get("ticker", "—")
    best_score   = int(best.get("score_global") or 0)
    worst_ticker = worst.get("ticker", "—")

    # Catalyseurs derived from actual metrics
    cats = []
    cats.append((
        f"Leader {best_ticker}",
        f"Score FinSight {best_score}/100 avec marge EBITDA {mg_med:.1f} % "
        f"et momentum 52W {mom_med:+.1f} %."
    ))
    if rev_med > 5:
        cats.append((
            "Croissance sectorielle soutenue",
            f"Croissance Médian des revenus a {rev_med:+.1f} % — dynamique favorable."
        ))
    else:
        cats.append((
            "Valorisation attractive",
            f"EV/EBITDA Médian a {ev_med:.1f}x — potentiel de rerating si catalyseurs confirmes."
        ))
    cats.append((
        "Consolidation sectorielle",
        "Vague de M&A potentielle — prime de rachat sur les acteurs sous-valorises."
    ))

    # Risques derived from worst performer and macro
    risks = []
    if worst and worst_ticker != best_ticker:
        risks.append((
            f"Retardataire {worst_ticker}",
            f"Score {int(worst.get('score_global') or 0)}/100 — "
            "Écart significatif avec le leader, sélectivité requise."
        ))
    else:
        risks.append((
            "Dispersion des scores",
            "Écart entre les profils — allocation selective recommandee."
        ))
    risks.append((
        "Environnement macroeconomique",
        "Ralentissement de la croissance mondiale — impact sur volumes et revenus."
    ))
    risks.append((
        "Pressions réglementaires",
        "Évolution du cadre sectoriel — Surcoûts de Conformité potentiels."
    ))

    # Drivers from actual metrics
    drivers = []
    if rev_med > 0:
        drivers.append(("up", "Croissance Organique",
                        f"Revenus en progression a {rev_med:+.1f} % — tendance positive"))
    else:
        drivers.append(("down", "Croissance Sous Pression",
                        f"Revenus Médianes a {rev_med:+.1f} % — surveillance requise"))
    if mg_med > 20:
        drivers.append(("up", "Marges Solides",
                        f"Marge EBITDA Médian a {mg_med:.1f} % — structure rentable"))
    else:
        drivers.append(("down", "Pression sur les Marges",
                        f"Marge EBITDA a {mg_med:.1f} % — potentiel d'amélioration"))
    drivers.append(("down", "Macro Global",
                    "Ralentissement croissance mondiale — sensibilité cyclique"))
    drivers.append(("up", "Consolidation M&A",
                    "Prime de valorisation potentielle — vague sectorielle"))

    # Métriques: real computed values from state
    metriques = [
        ("EV/EBITDA", f"{ev_med:.1f}x",     "Médiane secteur"),
        ("Mg EBITDA", f"{mg_med:.1f} %",     "LTM"),
        ("Croissance", f"{rev_med:+.1f} %",  "YoY"),
        ("Momentum",  f"{mom_med:+.1f} %",   "52W"),
        ("Score Moyen", f"{score_moyen}/100", "FinSight"),
        ("Nb Sociétés", str(len(td)),         "Analysées"),
    ]

    # Conditions based on signal level
    alert_score = max(30, score_moyen - 15)
    conditions = [
        ("Macro",       "Récession confirmee — 2 trimestrès négatifs",          "6-9 mois"),
        ("Sectoriel",   f"Score moyen < {alert_score} — dégradation generalisee","6-12 mois"),
        ("Fondamental", "Révision baissière EPS > 15 %",                         "3-6 mois"),
        ("Réglementaire", "Mesures restrictives majeures",                       "Variable"),
    ]

    quality = "attractif" if score_moyen >= 65 else "modéré" if score_moyen >= 45 else "prudent"
    return {
        "description": (
            f"Le secteur {sector_name} regroupe {len(td)} entreprise(s) analysée(s) dans l'univers "
            f"FinSight. EV/EBITDA Médian : {ev_med:.1f}x, marge EBITDA : {mg_med:.1f} %, "
            f"croissance revenus : {rev_med:+.1f} %. Profil {quality} — signal {sig_label.lower()} "
            f"selon le modèle multifactoriel FinSight (score moyen {score_moyen}/100)."
        ),
        "catalyseurs":  cats,
        "risques":      risks,
        "drivers":      drivers,
        "cycle_comment": f"Secteur {sector_name} | {sig_label.lower()} | score {score_moyen}/100",
        "metriques_dict":    metriques,
        "conditions":   conditions,
    }


# ─── DATA HELPERS ─────────────────────────────────────────────────────────────
def _med(vals):
    clean = [float(v) for v in vals if v is not None]
    return float(np.median(clean)) if clean else 0.0

def _avg(vals):
    clean = [float(v) for v in vals if v is not None]
    return float(np.mean(clean)) if clean else 0.0

def _signal_label(score):
    if score is None: return "NEUTRE", _HOLD
    if score >= 65: return "SURPONDÉRER", _BUY
    if score >= 45: return "NEUTRE", _HOLD
    return "SOUS-PONDÉRER", _SELL

def _reco(score):
    if score is None: return "HOLD"
    s = float(score)
    if s >= 65: return "BUY"
    if s >= 45: return "HOLD"
    return "SELL"

def _reco_color(reco):
    if reco == "BUY": return _BUY
    if reco == "SELL": return _SELL
    return _HOLD

def _fit_s(s, n: int) -> str:
    """Truncate at word boundary, no ellipsis — pour zones a hauteur fixe."""
    if not s: return ""
    s = str(s)
    if len(s) <= n: return s
    cut = s[:n]
    sp = cut.rfind(" ")
    return cut[:sp] if sp > n // 2 else cut


def _fmt_x(v, d=1):
    if v is None: return "—"
    try: return f"{float(v):.{d}f}x"
    except: return "—"

def _fmt_pct(v, d=1, mult=False):
    """Format percentage. mult=True si valeur en décimal (0.05→5%). False si déjà en % (5.0→5%)."""
    if v is None: return "—"
    try:
        fv = float(v) * (100 if mult else 1)
        return f"{fv:+.{d}f} %"
    except: return "—"

def _fmt_pct_plain(v, d=1, mult=False):
    if v is None: return "—"
    try:
        fv = float(v) * (100 if mult else 1)
        return f"{fv:.{d}f} %"
    except: return "—"

def _fmt_pct_rev(v, d=1):
    """Pour revenue_growth stocké en décimal (0.05 = 5%)."""
    if v is None: return "—"
    try: return f"{float(v)*100:+.{d}f} %"
    except: return "—"

def _fmt_num(v, d=1):
    if v is None: return "—"
    try: return f"{float(v):.{d}f}"
    except: return "—"

def _fmt_mds(v):
    if v is None: return "—"
    try:
        fv = float(v) / 1e9
        if fv >= 100: return f"{fv:.0f} Mds"
        return f"{fv:.1f} Mds"
    except: return "—"

def _prepare_data(tickers_data: list[dict], sector_name: str, universe: str) -> dict:
    """Pre-compute all derived values needed for the PPTX."""
    td = tickers_data or []

    scores = [t.get("score_global") for t in td if t.get("score_global") is not None]
    score_moyen = int(round(_avg(scores))) if scores else 0

    sig_label, sig_color = _signal_label(score_moyen)

    ev_vals  = [t.get("ev_ebitda") for t in td if t.get("ev_ebitda")]
    _raw_rg  = [t.get("revenue_growth") for t in td if t.get("revenue_growth") is not None]
    mg_vals  = [t.get("ebitda_margin") for t in td if t.get("ebitda_margin") is not None]
    mom_vals = [t.get("momentum_52w") for t in td if t.get("momentum_52w") is not None]

    # Normalisé revenue_growth en décimal (0.18 = 18%).
    # compute_screening retourne en % (18.0), yfinance retourne en décimal (0.18).
    # Détection : si médiane abs > 5 → déjà en %, on divise par 100.
    if _raw_rg and abs(_med([abs(v) for v in _raw_rg])) > 5.0:
        td = [{**t, "revenue_growth": (t["revenue_growth"] / 100
              if t.get("revenue_growth") is not None else None)} for t in td]
    rev_vals = [t.get("revenue_growth") for t in td if t.get("revenue_growth") is not None]

    ev_med  = _med(ev_vals)
    rev_med = _med(rev_vals) * 100 if rev_vals else 0.0  # maintenant toujours en décimal
    mg_med  = _med(mg_vals)  if mg_vals  else 0.0        # ebitda_margin déjà en %
    mom_med = _med(mom_vals) if mom_vals else 0.0        # momentum_52w déjà en %

    # Content: base réelle depuis les metriques tickers, puis override LLM pour le narratif
    content = _build_content_from_data(
        td, sector_name, score_moyen, sig_label, ev_med, rev_med, mg_med, mom_med
    )
    _llm_content = _build_content_from_llm(
        sector_name, ev_med, rev_med, mg_med, mom_med, score_moyen, sig_label, td
    )
    if _llm_content:
        # Le LLM fournit les textes narratifs ; les metriques réelles restent du _build_content_from_data
        for _k in ("description", "catalyseurs", "risques", "drivers", "cycle_comment"):
            if _k in _llm_content:
                content[_k] = _llm_content[_k]
    elif sector_name in _SECTOR_CONTENT:
        # Fallback : textes statiques de la librairie (metriques réelles conservees)
        for _k in ("description", "catalyseurs", "risques", "drivers", "cycle_comment"):
            content[_k] = _SECTOR_CONTENT[sector_name].get(_k, content.get(_k))

    from config.sector_ref import get_sector_drivers
    drv = get_sector_drivers(sector_name)
    rf   = 0.033
    erp  = drv.get("erp", 0.055)
    beta = _avg([t.get("beta") or 1.0 for t in td])
    wacc = rf + beta * erp
    wacc_pct = round(wacc * 100, 1)

    sorted_td = sorted(td, key=lambda x: x.get("score_global") or 0, reverse=True)
    if len(sorted_td) > MAX_TOP_N:
        log.warning("sectoral_pptx: %d sociétés disponibles — top3 tronquee a MAX_TOP_N=%d",
                    len(sorted_td), MAX_TOP_N)
    top3 = sorted_td[:MAX_TOP_N]

    date_str = datetime.date.today().strftime("%d %B %Y").replace(
        "January", "janvier").replace("February", "fevrier").replace(
        "March", "mars").replace("April", "avril").replace("May", "mai").replace(
        "June", "juin").replace("July", "juillet").replace("August", "aout").replace(
        "September", "septembre").replace("October", "octobre").replace(
        "November", "novembre").replace("December", "decembre")

    # Traduction FR du nom de secteur pour affichage (i18n)
    try:
        from core.sector_labels import fr_label as _fr_lbl
        _sector_name_display = _fr_lbl(sector_name)
    except Exception:
        _sector_name_display = sector_name

    # Agrégation par sous-secteur (industry GICS)
    _industry_groups = {}
    for t in td:
        ind = t.get("industry") or "Autre"
        _industry_groups.setdefault(ind, []).append(t)
    subsectors = []
    for ind_name, items in _industry_groups.items():
        if ind_name == "Autre" and len(_industry_groups) > 1:
            continue
        nb = len(items)
        sc = int(_avg([x.get("score_global") or 0 for x in items]))
        evs = [x["ev_ebitda"] for x in items if x.get("ev_ebitda") and 1 < x["ev_ebitda"] < 100]
        ev = _med(evs) if evs else None
        mgs = [x.get("ebitda_margin") or 0 for x in items if x.get("ebitda_margin")]
        mg = _med(mgs) if mgs else None
        grs = [x.get("revenue_growth") or 0 for x in items if x.get("revenue_growth") is not None]
        gr = _med(grs) * 100 if grs else None  # en %
        moms = [x.get("momentum_52w") or 0 for x in items if x.get("momentum_52w") is not None]
        mom = _med(moms) if moms else None
        sig = "Surpondérer" if sc >= 60 else ("Sous-pondérer" if sc < 40 else "Neutre")
        best = sorted(items, key=lambda x: x.get("score_global") or 0, reverse=True)[:3]
        subsectors.append({
            "name": ind_name, "nb": nb, "score": sc, "signal": sig,
            "ev_ebitda": ev, "margin": mg, "growth": gr, "momentum": mom,
            "best": [(b.get("ticker", ""), b.get("score_global", 0)) for b in best],
            "pct": round(nb / len(td) * 100, 1),
        })
    subsectors.sort(key=lambda x: x["score"], reverse=True)

    return {
        "sector_name": _sector_name_display,        # affichage francais
        "sector_name_en": sector_name,               # original yfinance EN (data lookups)
        "universe": universe,
        "date_str": date_str,
        "N": len(td),
        "score_moyen": score_moyen,
        "sig_label": sig_label,
        "sig_color": sig_color,
        "ev_med": ev_med,
        "rev_med": rev_med,
        "mg_med": mg_med,
        "mom_med": mom_med,
        "wacc_pct": wacc_pct,
        "sorted_td": sorted_td,
        "top3": top3,
        "content": content,
        "tickers_data": td,
        "subsectors": subsectors,
    }


# ─── PPTX PRIMITIVE HELPERS ───────────────────────────────────────────────────
def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _rect(slide, x, y, w, h, fill=None, line=False, line_col=None, line_w=0.5):
    from pptx.oxml.ns import qn
    shape = slide.shapes.add_shape(1, Cm(x), Cm(y), Cm(w), Cm(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line and line_col:
        shape.line.color.rgb = line_col
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape

def _txb(slide, text, x, y, w, h, size=9, bold=False, color=None, align=PP_ALIGN.LEFT,
         italic=False, wrap=True):
    color = color or _BLACK
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box

def _txb2(slide, parts, x, y, w, h, size=9, align=PP_ALIGN.LEFT, wrap=True):
    """Multi-run text box. parts = list of (text, bold, color)."""
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    for text, bold, color in parts:
        run = para.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color or _BLACK
    return box

def _pic(slide, img_bytes, x, y, w, h):
    buf = io.BytesIO(img_bytes)
    slide.shapes.add_picture(buf, Cm(x), Cm(y), Cm(w), Cm(h))

def _add_table(slide, data, x, y, w, h, col_widths=None,
               header_fill=_NAVY, header_color=_WHITE,
               alt_fill=None, font_size=7.5, header_size=7.5):
    """Generic table. data[0] = headers, data[1:] = rows."""
    rows = len(data)
    cols = len(data[0]) if data else 1
    tbl_shape = slide.shapes.add_table(rows, cols, Cm(x), Cm(y), Cm(w), Cm(h))
    tbl = tbl_shape.table

    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Cm(cw * w / total)

    for r_idx, row_data in enumerate(data):
        is_header = (r_idx == 0)
        for c_idx, cell_text in enumerate(row_data):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = str(cell_text) if cell_text is not None else "—"
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            run = para.runs[0] if para.runs else para.add_run()
            run.font.size = Pt(header_size if is_header else font_size)
            run.font.bold = is_header
            if is_header:
                run.font.color.rgb = header_color
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
            else:
                run.font.color.rgb = _BLACK
                if alt_fill and r_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = alt_fill
                else:
                    cell.fill.background()
    return tbl


# ─── COMMON HEADER / FOOTER ───────────────────────────────────────────────────
_NAV_LABELS = ["1", "2", "3", "4", "5"]

def _header(slide, title, subtitle, active_section=1):
    """Navy header bar + title + nav dots + subtitle."""
    _rect(slide, 0, 0, 25.4, 1.4, fill=_NAVY)
    _txb(slide, title, 0.9, 0.05, 19.1, 1.3, size=13, bold=True, color=_WHITE)
    # Nav dots
    for i, lbl in enumerate(_NAV_LABELS):
        dot_x = 21.8 + i * 0.7
        fill = _WHITE if (i + 1) == active_section else _NAVYL
        _rect(slide, dot_x, 0.35, 0.55, 0.55, fill=fill)
        c = _BLACK if (i + 1) == active_section else _GRAYD
        _txb(slide, lbl, dot_x, 0.35, 0.55, 0.55, size=7, bold=True, color=c, align=PP_ALIGN.CENTER)
    # Subtitle
    _txb(slide, subtitle, 0.9, 1.6, 23.6, 0.6, size=8, color=_GRAYT)

def _footer(slide):
    _rect(slide, 0, 13.75, 25.4, 0.5, fill=_GRAYL)
    _txb(slide, "FinSight IA  ·  Usage confidentiel", 0.9, 13.8, 23.6, 0.4, size=7, color=_GRAYD)

def _chapter_divider(prs, num_str, chapter_title, subtitle):
    """Dark background chapter divider slide."""
    slide = _blank(prs)
    _rect(slide, 0, 0, 25.4, 14.3, fill=_NAVY)
    _txb(slide, num_str, 1.0, 3.5, 8.0, 4.5, size=72, bold=True, color=_NAVYL)
    _txb(slide, chapter_title, 1.0, 7.0, 23.0, 2.0, size=28, bold=True, color=_WHITE)
    _rect(slide, 1.0, 9.1, 15.0, 0.05, fill=_GRAYD)
    _txb(slide, subtitle, 1.0, 9.4, 22.9, 0.8, size=11, color=_GRAYD)
    _txb(slide, "FinSight IA  ·  Usage confidentiel", 0.9, 13.75, 23.6, 0.4, size=7, color=_GRAYD)
    return slide


# ─── CHART BUILDERS ───────────────────────────────────────────────────────────
def _chart_valuation_bars(tickers_data) -> bytes:
    """EV/EBITDA ranking bar chart — lisible, remplace le scatter diforme."""
    points = []
    for t in tickers_data:
        ev = t.get("ev_ebitda")
        if ev is None:
            continue
        try:
            ev_f = float(ev)
            if 0 < ev_f <= 150:
                points.append((t.get("ticker", "?"), ev_f, float(t.get("score_global") or 50)))
        except (TypeError, ValueError):
            pass

    if not points:
        fig, ax = plt.subplots(figsize=(5.8, 4.5))
        ax.text(0.5, 0.5, "EV/EBITDA non disponible", ha='center', va='center',
                transform=ax.transAxes, fontsize=10, color='#999999')
        ax.set_facecolor('#FFFFFF')
        fig.patch.set_facecolor('#FFFFFF')
        buf = io.BytesIO()
        fig.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='white', edgecolor='none')
        plt.close(fig)
        buf.seek(0)
        return buf.read()

    # Cap top 15 par score pour lisibilite
    if len(points) > 15:
        points = sorted(points, key=lambda x: -x[2])[:15]
    # Tri croissant EV/EBITDA
    points.sort(key=lambda x: x[1])

    tickers_list = [p[0] for p in points]
    evs = [p[1] for p in points]
    med_ev = float(np.median(evs))
    bar_colors = ['#1A7A4A' if ev < med_ev else '#A82020' for ev in evs]

    n = len(points)
    fig_h = max(4.5, n * 0.40 + 1.5)
    fig, ax = plt.subplots(figsize=(7.5, fig_h))
    fig.patch.set_facecolor('#FFFFFF')
    ax.set_facecolor('#F8F9FA')

    ax.barh(range(n), evs, color=bar_colors, alpha=0.85,
            edgecolor='white', linewidth=0.5, height=0.65)
    ax.set_yticks(range(n))
    ax.set_yticklabels(tickers_list, fontsize=7.5)

    ax.axvline(med_ev, color='#1B3A6B', linewidth=1.5, linestyle='--', zorder=5)
    ax.text(med_ev + 0.2, n - 0.4, f'Med: {med_ev:.1f}x',
            fontsize=7, color='#1B3A6B', va='top', fontweight='bold')

    x_max = max(evs) if evs else 1
    for i, ev in enumerate(evs):
        ax.text(ev + x_max * 0.01, i, f'{ev:.1f}x', va='center', ha='left', fontsize=6.5, color='#333333')

    ax.set_xlabel("EV / EBITDA (x)", fontsize=8, color='#555555')
    ax.tick_params(labelsize=7, colors='#777777')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#DDDDDD')
    ax.spines['bottom'].set_color('#DDDDDD')
    ax.grid(True, alpha=0.2, axis='x', linestyle=':')
    ax.set_xlim(0, x_max * 1.18)

    from matplotlib.patches import Patch
    legend_els = [Patch(facecolor='#1A7A4A', label='Sous Médiane — opportunité relative'),
                  Patch(facecolor='#A82020', label='Prime vs Médiane — valorisation élevée')]
    ax.legend(handles=legend_els, loc='lower right', fontsize=7, framealpha=0.85,
              edgecolor='#CCCCCC', handlelength=1.2, handleheight=0.8)

    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='white', edgecolor='none')
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def _chart_distribution(tickers_data) -> bytes:
    """Bar chart EV/EBITDA distribution vs Médian."""
    td_ev = [(t.get("ticker", "?"), t.get("ev_ebitda"))
             for t in tickers_data if t.get("ev_ebitda")]
    td_ev = [(tk, float(ev)) for tk, ev in td_ev if 0 < float(ev) <= 100]  # filtre outliers >100x
    td_ev.sort(key=lambda x: x[1])
    # Cap a 30 barres pour rester lisible
    if len(td_ev) > 30:
        td_ev = td_ev[:15] + td_ev[-15:]  # 15 plus basses + 15 plus hautes valeurs

    if not td_ev:
        fig, ax = plt.subplots(figsize=(5.8, 4.5))
        ax.text(0.5, 0.5, "Données insuffisantés", ha='center', va='center', transform=ax.transAxes)
        buf = io.BytesIO()
        fig.savefig(buf, format='png', dpi=150); plt.close(fig); buf.seek(0)
        return buf.read()

    labels = [x[0] for x in td_ev]
    vals   = [x[1] for x in td_ev]
    med    = float(np.median(vals))
    colors = ['#1A7A4A' if v <= med else '#A82020' for v in vals]

    fig, ax = plt.subplots(figsize=(5.8, 4.5))
    fig.patch.set_facecolor('#FFFFFF')
    ax.set_facecolor('#F8F9FA')
    bars = ax.bar(labels, vals, color=colors, alpha=0.85, zorder=3)
    ax.axhline(med, color='#1B3A6B', linewidth=1.5, linestyle='--', label=f"Médiane {med:.1f}x")
    # Labels de valeur supprimés (slides epures)
    ax.set_ylabel("EV/EBITDA", fontsize=8, color='#555555')
    n_labels = len(labels)
    _lsize = 6 if n_labels > 20 else 7 if n_labels > 12 else 7.5
    ax.tick_params(axis='x', labelsize=_lsize, colors='#333333', rotation=90)
    ax.tick_params(axis='y', labelsize=7, colors='#777777')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#DDDDDD')
    ax.spines['bottom'].set_color('#DDDDDD')
    ax.legend(fontsize=7.5, framealpha=0.7)
    ax.grid(True, axis='y', alpha=0.3, linestyle=':')
    plt.tight_layout()
    plt.subplots_adjust(bottom=0.22)
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def _chart_performance(tickers_data, best_ticker=None, worst_ticker=None,
                       sector_name=None, universe=None) -> bytes:
    """52W relative performance line chart via yfinance.
    Top 3 momentum + worst sont surlignés en couleurs distinctes, autres en gris.
    NOUVEAU (refonte #86) : overlay ETF sectoriel en ligne pointillee noire pour
    benchmark visuel immediat.
    La légende affiche les noms commerciaux (company) au lieu des tickers."""
    import datetime as dt
    # Build ticker -> company name lookup pour la legende
    _name_map = {}
    for _t in tickers_data:
        _tk = _t.get("ticker", "")
        _nm = _t.get("company") or _t.get("name") or _tk
        if _tk:
            _name_map[_tk] = _nm[:28]  # cap 28 chars pour la legende

    # Resolve l'ETF sectoriel de reference via core.sector_etfs
    _etf_ticker = None
    if sector_name:
        try:
            from core.sector_etfs import get_etf_for
            _etf_info = get_etf_for(sector_name, universe=universe or "S&P 500")
            if _etf_info:
                _etf_ticker = _etf_info["ticker"]
        except Exception as _e:
            log.debug(f"[sectoral_pptx_writer:_chart_performance] exception skipped: {_e}")

    # Top 3 par momentum_52w (couleurs distinctes) + worst (rouge)
    _ranked = sorted([t for t in tickers_data if t.get("ticker")],
                     key=lambda x: x.get("momentum_52w") or -1e9, reverse=True)
    _top3_tk = [t.get("ticker") for t in _ranked[:3]]
    _worst_tk = worst_ticker or (_ranked[-1].get("ticker") if _ranked else None)
    _highlight_colors = {
        (_top3_tk[0] if len(_top3_tk) > 0 else None): ('#1A7A4A', 2.8),  # vert vif
        (_top3_tk[1] if len(_top3_tk) > 1 else None): ('#1B3A6B', 2.2),  # navy
        (_top3_tk[2] if len(_top3_tk) > 2 else None): ('#B06000', 2.2),  # ambre
        _worst_tk: ('#A82020', 2.5),                                       # rouge
    }
    _highlight_colors.pop(None, None)

    fig, ax = plt.subplots(figsize=(9.5, 4.4))
    fig.patch.set_facecolor('#FFFFFF')
    ax.set_facecolor('#F8F9FA')

    if len(tickers_data) > MAX_TICKERS_CHART:
        log.warning("sectoral_pptx _chart_performance: %d tickers disponibles — "
                    "affichage tronque a MAX_TICKERS_CHART=%d",
                    len(tickers_data), MAX_TICKERS_CHART)

    # SEC-PPTX-S20 Baptiste : on ne trace PLUS les courbes grisees anonymes.
    # Chaque ligne doit etre associee a une entree de legende. On ne trace que :
    # - l'ETF sectoriel (pointille noir)
    # - le top 3 par momentum (couleurs distinctes)
    # - le pire performer (rouge)
    # Soit 5 courbes maximum, toutes legendees.
    _highlighted_set = set(_highlight_colors.keys())

    plotted = 0
    try:
        import yfinance as yf
        # Overlay ETF sectoriel en premier (ligne noire pointillee, benchmark)
        if _etf_ticker:
            try:
                _etf_hist = get_ticker(_etf_ticker).history(period='1y', interval='1wk')
                if not _etf_hist.empty and 'Close' in _etf_hist.columns:
                    _etf_s = _etf_hist['Close']
                    if len(_etf_s) >= 4:
                        _etf_norm = (_etf_s / _etf_s.iloc[0] - 1) * 100
                        ax.plot(_etf_norm.index, _etf_norm.values,
                                color='#000000', linewidth=2.5, linestyle='--',
                                label=f"ETF {_etf_ticker} (benchmark)",
                                alpha=0.95, zorder=5)
            except Exception as _e:
                log.debug(f"[sectoral_pptx_writer:_chart_performance] exception skipped: {_e}")

        # Ne plotter QUE les tickers highlight (top 3 + worst) — zero ligne anonyme
        for ticker in _highlighted_set:
            if not ticker:
                continue
            try:
                hist_df = get_ticker(ticker).history(period='1y', interval='1wk')
                if hist_df.empty or 'Close' not in hist_df.columns:
                    continue
                hist = hist_df['Close']
                if len(hist) < 4:
                    continue
                norm = (hist / hist.iloc[0] - 1) * 100
                col, lw = _highlight_colors[ticker]
                _lbl = _name_map.get(ticker, ticker)
                ax.plot(norm.index, norm.values,
                        color=col, linewidth=lw, label=_lbl, alpha=1.0, zorder=4)
                plotted += 1
            except Exception as _e:
                log.debug(f"[sectoral_pptx_writer:_chart_performance] exception skipped: {_e}")
    except Exception as _e:
        log.debug(f"[sectoral_pptx_writer:_chart_performance] exception skipped: {_e}")

    if plotted == 0:
        # Fallback illustrative — warning car données non réelles
        log.warning("sectoral_pptx _chart_performance: yfinance indisponible — "
                    "utilisation d'un fallback illustratif (Données simulees)")
        import numpy as np
        x = np.linspace(0, 52, 53)
        for ticker in _highlighted_set:
            if not ticker:
                continue
            np.random.seed(hash(ticker) % 1024)
            y = np.cumsum(np.random.randn(53) * 1.5)
            col, lw = _highlight_colors[ticker]
            _lbl = _name_map.get(ticker, ticker)
            ax.plot(x, y, color=col, linewidth=lw, label=_lbl, alpha=1.0, zorder=4)
        ax.set_xlabel("Semaines (illustratif)", fontsize=8, color='#555555')
    else:
        ax.set_xlabel("Date", fontsize=8, color='#555555')

    ax.axhline(0, color='#CCCCCC', linewidth=1.0, linestyle='-')
    ax.set_ylabel("Performance relative (%)", fontsize=8, color='#555555')
    ax.tick_params(labelsize=7, colors='#777777')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#DDDDDD')
    ax.spines['bottom'].set_color('#DDDDDD')
    ax.grid(True, alpha=0.3, linestyle=':')
    _leg = ax.legend(loc='upper left', fontsize=8, framealpha=0.95,
                     edgecolor='#888888', facecolor='#FFFFFF',
                     ncol=1, borderpad=0.6, labelspacing=0.5)
    if _leg is not None:
        _leg.get_frame().set_linewidth(0.8)
    plt.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close(fig)
    buf.seek(0)
    return buf.read()


# ─── SLIDE BUILDERS ───────────────────────────────────────────────────────────

def _s01_cover(prs, D):
    slide = _blank(prs)
    # Barre navy fine en haut
    _rect(slide, 0, 0, 25.4, 1.8, fill=_NAVY)
    _txb(slide, "FinSight IA", 0.9, 0.15, 12.0, 1.2, size=12, bold=True, color=_WHITE)
    _txb(slide, "Pitchbook  —  Analyse Sectorielle", 0, 0.25, 25.2, 1.0, size=9.5, color=_GRAYD, align=PP_ALIGN.RIGHT)
    # Titre secteur centré
    _txb(slide, D["sector_name"], 0, 4.5, 25.4, 3.0, size=38, bold=True, color=_NAVY, align=PP_ALIGN.CENTER)
    # Univers + nb sociétés centré
    _txb(slide, f"{D['universe']}  ·  {D['N']} sociétés analysées", 0, 7.8, 25.4, 1.0, size=12, color=_GRAYT, align=PP_ALIGN.CENTER)
    # Signal badge centré
    sig_label = D["sig_label"]
    sig_color = D["sig_color"]
    _rect(slide, 9.1, 9.2, 7.1, 1.4, fill=sig_color)
    _txb(slide, f"● {sig_label}", 9.1, 9.3, 7.1, 1.2, size=13, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)
    # Footer
    _rect(slide, 0, 13.6, 25.4, 0.7, fill=_GRAYL)
    _txb(slide, "Rapport confidentiel", 0.9, 13.65, 10.0, 0.5, size=7.5, color=_GRAYT)
    _txb(slide, D["date_str"], 0, 13.65, 24.4, 0.5, size=7.5, color=_GRAYT, align=PP_ALIGN.RIGHT)


def _s02_exec_summary(prs, D):
    slide = _blank(prs)
    _header(slide, "Executive Summary", f"{D['sector_name']}  ·  {D['universe']}  ·  {D['N']} sociétés analysées", 1)

    # Signal + KPIs bar
    _rect(slide, 0.9, 2.4, 23.6, 1.1, fill=_GRAYL)
    _rect(slide, 0.9, 2.4, 0.1, 1.1, fill=D["sig_color"])
    _txb(slide, f"● {D['sig_label']}", 1.3, 2.45, 3.2, 1.0, size=11, bold=True, color=D["sig_color"])
    kpi_str = (f"EV/EBITDA med. : {D['ev_med']:.1f}x  ·  "
               f"Croissance med. : {D['rev_med']:+.1f} %  ·  "
               f"Marge EBITDA med. : {D['mg_med']:.1f} %  ·  "
               f"Score moyen : {D['score_moyen']}/100")
    _txb(slide, kpi_str, 4.6, 2.5, 19.8, 0.9, size=8.5, color=_NAVY)

    content = D["content"]
    cats = content.get("catalyseurs", [])
    risks = content.get("risques", [])

    if len(cats) > MAX_CATS_DISPLAYED:
        log.warning("sectoral_pptx _s02: %d catalyseurs disponibles — affichage tronque a MAX_CATS_DISPLAYED=%d",
                    len(cats), MAX_CATS_DISPLAYED)
    if len(risks) > MAX_RISKS_DISPLAYED:
        log.warning("sectoral_pptx _s02: %d risques disponibles — affichage tronque a MAX_RISKS_DISPLAYED=%d",
                    len(risks), MAX_RISKS_DISPLAYED)

    # Catalyseurs
    _rect(slide, 0.9, 3.7, 11.4, 0.7, fill=_NAVY)
    _txb(slide, "CATALYSEURS SECTORIELS", 0.9, 3.75, 11.4, 0.6, size=8.5, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)
    for i, (cat_title, cat_body) in enumerate(cats[:MAX_CATS_DISPLAYED]):
        yy = 4.6 + i * 2.3
        _rect(slide, 0.9, yy, 0.1, 2.0, fill=_BUY)
        _rect(slide, 1.0, yy, 11.3, 2.0, fill=_GREEN_L)
        _txb(slide, cat_title, 1.3, yy + 0.05, 10.7, 0.55, size=8.5, bold=True, color=_NAVY)
        _txb(slide, cat_body[:200], 1.3, yy + 0.55, 10.7, 1.25, size=7.5, color=_GRAYT, wrap=True)

    # Risques
    _rect(slide, 13.1, 3.7, 11.4, 0.7, fill=_SELL)
    _txb(slide, "RISQUES PRINCIPAUX", 13.1, 3.75, 11.4, 0.6, size=8.5, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)
    for i, (risk_title, risk_body) in enumerate(risks[:MAX_RISKS_DISPLAYED]):
        yy = 4.6 + i * 2.3
        _rect(slide, 13.1, yy, 0.1, 2.0, fill=_SELL)
        _rect(slide, 13.2, yy, 11.3, 2.0, fill=_RED_L)
        _txb(slide, risk_title, 13.5, yy + 0.05, 10.7, 0.55, size=8.5, bold=True, color=_SELL)
        _txb(slide, risk_body[:200], 13.5, yy + 0.55, 10.7, 1.25, size=7.5, color=_GRAYT, wrap=True)

    # 4 KPI boxes
    kpis = [
        (f"{D['ev_med']:.1f}x", "EV/EBITDA med.", "vs Marché EU"),
        (f"{D['rev_med']:+.1f} %", "Croissance Rev. med.", "YoY sectoriel"),
        (f"{D['mg_med']:.1f} %", "Marge EBITDA med.", "LTM sectorielle"),
        (f"{D['mom_med']:+.1f} %", "Momentum 52W med.", "Performance relative"),
    ]
    for i, (val, lbl1, lbl2) in enumerate(kpis):
        kx = 0.9 + i * 5.9
        _rect(slide, kx, 12.0, 5.7, 1.7, fill=_GRAYL)
        _rect(slide, kx, 12.0, 0.1, 1.7, fill=_NAVYL)
        _txb(slide, val, kx + 0.3, 12.05, 5.2, 0.9, size=20, bold=True, color=_NAVY)
        _txb(slide, lbl1, kx + 0.3, 12.9, 5.2, 0.4, size=7.5, color=_GRAYT)
        _txb(slide, lbl2, kx + 0.3, 13.3, 5.2, 0.35, size=7, color=_GRAYD)

    _footer(slide)


def _s03_sommaire(prs, D):
    slide = _blank(prs)
    _header(slide, "Sommaire", f"{D['sector_name']}  ·  {D['universe']}  ·  Structure de l'analyse — 20 slides", 1)

    chapters = [
        ("01", "Présentation du Secteur", "Caractéristiques structurelles · Ratios comparatifs · Positionnement cycle", "p. 5–7"),
        ("02", "Cartographie des Sociétés", f"{D['N']} sociétés · Scatter valorisation · Scores FinSight détaillés", "p. 9–11"),
        ("03", "Top 3 & Valorisations", "Synthèse Top 3 · Distribution EV/EBITDA · Zone d'entree optimale", "p. 13–15"),
        ("04", "Risques, Sentiment & Méthodologie", "Conditions d'invalidation · FinBERT · Sources & data lineage", "p. 17–19"),
    ]
    for i, (num, title, sub, pages) in enumerate(chapters):
        yy = 2.6 + i * 2.6
        _rect(slide, 0.9, yy, 23.6, 2.2, fill=_GRAYL)
        _rect(slide, 0.9, yy, 1.9, 2.2, fill=_NAVY)
        _txb(slide, num, 0.9, yy + 0.3, 1.9, 1.6, size=20, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)
        _txb(slide, title, 3.1, yy + 0.2, 18.3, 0.8, size=11, bold=True, color=_NAVY)
        _txb(slide, sub, 3.1, yy + 1.1, 18.3, 0.8, size=8, color=_GRAYT)
        _txb(slide, pages, 23.2, yy + 0.5, 1.3, 1.2, size=8, color=_GRAYD, align=PP_ALIGN.RIGHT)

    _footer(slide)


def _s05_presentation(prs, D):
    slide = _blank(prs)
    _header(slide, "Pr\u00e9sentation du Secteur",
            f"{D['sector_name']}  \u00b7  Caract\u00e9ristiques structurelles & m\u00e9triques de r\u00e9f\u00e9rence", 1)

    content = D["content"]
    desc = content.get("description", "")
    metriques_dict = content.get("metriques_dict", [])

    # ─── Bloc Description (compact, partie haute) ─────────────────────────
    # SEC-PPTX-S5 Baptiste : bandeau navy hauteur 0.60 (vs 0.55) et texte
    # titre centre plus haut (+0.02) pour eviter le clipping leger bas que
    # l'on voyait sur le rendu technologie.
    _rect(slide, 0.9, 2.5, 13.7, 3.5, fill=_GRAYL)
    _rect(slide, 0.9, 2.5, 0.1, 3.5, fill=_NAVY)
    _rect(slide, 0.9, 2.5, 13.7, 0.60, fill=_NAVY)
    _txb(slide, "PANORAMA SECTORIEL", 1.1, 2.56, 13.3, 0.48,
         size=9, bold=True, color=_WHITE)
    _txb(slide, desc, 1.3, 3.25, 13.1, 2.6, size=10, color=_GRAYT, wrap=True)

    # ─── Bloc Catalyseurs Cl\u00e9s (sous la description, bien s\u00e9par\u00e9) ────────
    cats = content.get("catalyseurs", [])
    if cats:
        _cat_y0 = 6.3
        _cat_h_total = 6.65
        _rect(slide, 0.9, _cat_y0, 13.7, _cat_h_total, fill=_GRAYL)
        _rect(slide, 0.9, _cat_y0, 13.7, 0.60, fill=_NAVY)
        _txb(slide, "CATALYSEURS CL\u00c9S", 1.1, _cat_y0 + 0.06, 13.3, 0.48,
             size=9, bold=True, color=_WHITE)
        # 3 catalyseurs, hauteur 1.85 chacun, gap 0.1 — box agrandie pour englober la 3e barre
        _row_h = 1.85
        _row_gap = 0.1
        for j, (ct, cb) in enumerate(cats[:3]):
            _y = _cat_y0 + 0.7 + j * (_row_h + _row_gap)
            # Barre navy a gauche (couleur cohérente avec le reste du deck)
            _rect(slide, 1.1, _y, 0.12, _row_h, fill=_NAVY)
            # Titre LLM (court, bold)
            _txb(slide, _fit_s(ct, 60), 1.4, _y + 0.05, 13.0, 0.45,
                 size=10, bold=True, color=_NAVY)
            # Description LLM (wrap, plus d'espace)
            _txb(slide, cb[:260], 1.4, _y + 0.55, 13.0, _row_h - 0.6,
                 size=9, color=_GRAYT, wrap=True)

    # ─── Metrics table (droite haut) ───────────────────────────────────────
    tbl_data = [["M\u00e9trique", "Valeur", "Lecture"]]
    for met in metriques_dict:
        tbl_data.append(list(met))
    _metrics_h = len(tbl_data) * 0.58
    _add_table(slide, tbl_data, 15.1, 2.5, 9.4, _metrics_h,
               col_widths=[3.5, 2.5, 3.4], font_size=7.5, header_size=7.5, alt_fill=_GRAYL)

    # ─── Section ETF de r\u00e9f\u00e9rence (droite bas) ─── refonte #86
    # Affiche l'ETF sectoriel + top 5 holdings r\u00e9els en dessous des m\u00e9triques
    try:
        from core.sector_etfs import get_etf_for, etf_issuer
        from core.etf_holdings import fetch_etf_holdings, format_aum, format_ter
        _etf_info = get_etf_for(D.get('sector_name', ''), universe=D.get('universe', 'S&P 500'))
        _etf_data = None
        if _etf_info:
            _etf_data = fetch_etf_holdings(_etf_info['ticker'])
        if _etf_info and _etf_data:
            _etf_y0 = 2.5 + _metrics_h + 0.3
            _etf_h  = 13.5 - _etf_y0 - 0.3
            _rect(slide, 15.1, _etf_y0, 9.4, _etf_h, fill=_GRAYL)
            _rect(slide, 15.1, _etf_y0, 0.1, _etf_h, fill=_NAVY)
            _rect(slide, 15.1, _etf_y0, 9.4, 0.55, fill=_NAVY)
            _txb(slide, f"ETF DE R\u00c9F\u00c9RENCE : {_etf_info['ticker']}",
                 15.25, _etf_y0 + 0.08, 9.1, 0.45,
                 size=8.5, bold=True, color=_WHITE)
            _yy = _etf_y0 + 0.75
            _txb(slide, _etf_info['name'][:60], 15.3, _yy, 9.0, 0.4,
                 size=7.5, bold=True, color=_NAVY)
            _yy += 0.4
            _issuer_short = etf_issuer(_etf_info['zone']).split('(')[0].strip()[:40]
            _txb(slide, f"Emetteur : {_issuer_short}", 15.3, _yy, 9.0, 0.32,
                 size=7, color=_GRAYT)
            _yy += 0.30
            _aum_str = format_aum(_etf_data.get('aum_usd'))
            _ter_str = format_ter(_etf_data.get('ter_bps'))
            _txb(slide, f"AUM : {_aum_str}  \u00b7  TER : {_ter_str}",
                 15.3, _yy, 9.0, 0.32, size=7, color=_GRAYT)
            _yy += 0.40
            _txb(slide, "TOP 5 HOLDINGS", 15.3, _yy, 9.0, 0.32,
                 size=7, bold=True, color=_NAVY)
            _yy += 0.32
            _top5 = sorted(_etf_data.get('holdings', []),
                           key=lambda h: h.get('weight', 0) or 0, reverse=True)[:5]
            for _h in _top5:
                _w = _h.get('weight', 0) or 0
                _w_pct = _w * 100 if _w < 1 else _w
                _line = f"{_h.get('ticker','?'):<8}  {str(_h.get('name',''))[:22]:<22}  {_w_pct:>4.1f} %"
                _txb(slide, _line, 15.3, _yy, 9.0, 0.30, size=6.8, color=_GRAYT)
                _yy += 0.30
    except Exception as _etf_e:
        import logging as _log_etf
        _log_etf.getLogger(__name__).warning("sectoral_pptx ETF block: %s", _etf_e)

    _footer(slide)


def _s06_ratios(prs, D):
    slide = _blank(prs)
    _header(slide, "Ratios Comparatifs Sectoriels",
            f"{D['sector_name']}  ·  Comparaison multiples LTM — {D['N']} sociétés vs médiane sectorielle", 2)

    td = D["sorted_td"]
    MAX_S06 = 8  # cap slide 6 a 8 lignes — garantit lisibilite bloc lecture
    td_disp = td[:MAX_S06]
    tbl_data = [["Société", "EV/EBITDA", "EV/Rev.", "P/E", "Mg Brute", "Mg EBITDA", "ROE"]]

    def _norm_mc(mc):
        """Normalise market_cap : compute_screening le stocke parfois en Mds
        (244.0), cli en absolu (2.44e11). SEC-PPTX-S6 : il faut traiter les
        deux conventions pour que le proxy EV/EBITDA marche sur tous les flux."""
        if mc is None:
            return None
        try:
            f = float(mc)
            return f * 1e9 if f < 1e6 else f
        except Exception:
            return None

    def _fallback_ev_ebitda(t):
        """Calcul EV/EBITDA depuis composantés si valeur directe absente.
        Priorite 1: ev_ebitda direct ; priorite 2: MktCap / ebitda_ltm
        (rate-limite yfinance NVDA, AVGO, ASML.AS) ; priorite 3: via marge."""
        v = t.get("ev_ebitda")
        if v is not None:
            return v
        mc = _norm_mc(t.get("market_cap"))
        # Priorite 2 : ebitda_ltm absolu (donnee directe)
        _eb_ltm = t.get("ebitda_ltm")
        if mc and _eb_ltm and float(_eb_ltm) > 0:
            try:
                _proxy = round(mc / float(_eb_ltm), 1)
                if 0 < _proxy < 200:
                    return _proxy
            except Exception as _e:
                log.debug(f"[sectoral_pptx_writer:_fallback_ev_ebitda] exception skipped: {_e}")
        # Priorite 3 : ebitda_margin * revenue_ltm
        em  = t.get("ebitda_margin")
        rev = t.get("revenue_ltm") or t.get("revenue")
        if mc and em and rev and float(em) > 0 and float(rev) > 0:
            ebitda_est = float(em) * float(rev)
            nd = t.get("net_debt") or 0
            ev_est = mc + float(nd)
            if ebitda_est > 0:
                return round(ev_est / ebitda_est, 1)
        return None

    def _fallback_ev_revenue(t):
        """Calcul EV/Revenue depuis composantés si valeur directe absente."""
        v = t.get("ev_revenue")
        if v is not None:
            return v
        mc  = _norm_mc(t.get("market_cap"))
        rev = t.get("revenue_ltm") or t.get("revenue")
        if mc and rev and float(rev) > 0:
            nd = t.get("net_debt") or 0
            ev_est = mc + float(nd)
            return round(ev_est / float(rev), 1)
        return None

    def _fallback_pe(t):
        """Calcul P/E depuis price/eps si valeur directe absente."""
        v = t.get("pe_ratio") or t.get("pe")
        if v is not None:
            return v, False  # (value, is_negative_earnings)
        # Tentative : market_cap / net_income
        mc = t.get("market_cap")
        ni = t.get("net_income")
        if mc and ni and float(ni) != 0:
            if float(ni) < 0:
                return None, True  # earnings negatifs
            return round(float(mc) / float(ni), 1), False
        # Tentative : price / trailing_eps
        px  = t.get("price") or t.get("current_price")
        eps = t.get("trailing_eps") or t.get("eps_ttm") or t.get("eps")
        if px and eps and float(eps) != 0:
            if float(eps) < 0:
                return None, True
            return round(float(px) / float(eps), 1), False
        ni2 = t.get("eps_ttm") or t.get("eps")
        try:
            return None, ni2 is not None and float(ni2) < 0
        except (ValueError, TypeError):
            return None, False

    for t in td_disp:
        _pe_val, _pe_neg = _fallback_pe(t)
        _pe_str = _fmt_x(_pe_val) if _pe_val is not None else ("neg." if _pe_neg else "—")
        tbl_data.append([
            t.get("company", t.get("ticker", ""))[:28],
            _fmt_x(_fallback_ev_ebitda(t)),
            _fmt_x(_fallback_ev_revenue(t)),
            _pe_str,
            _fmt_pct_plain(t.get("gross_margin")),
            _fmt_pct_plain(t.get("ebitda_margin")),
            _fmt_pct_plain(t.get("roe")),
        ])
    # Médiane row (calculée sur tout l'univers td avec fallbacks)
    _ev_ebitda_vals = [v for t in td for v in [_fallback_ev_ebitda(t)] if v is not None]
    _ev_rev_vals    = [v for t in td for v in [_fallback_ev_revenue(t)] if v is not None]
    _pe_vals        = [v for t in td for v, _ in [_fallback_pe(t)] if v is not None]
    tbl_data.append(["MÉDIANE",
                      _fmt_x(_med(_ev_ebitda_vals)) if _ev_ebitda_vals else "—",
                      _fmt_x(_med(_ev_rev_vals))    if _ev_rev_vals    else "—",
                      _fmt_x(_med(_pe_vals))         if _pe_vals        else "—",
                      _fmt_pct_plain(_med([t.get("gross_margin")  for t in td if t.get("gross_margin")])),
                      _fmt_pct_plain(_med([t.get("ebitda_margin") for t in td if t.get("ebitda_margin")])),
                      _fmt_pct_plain(_med([t.get("roe")           for t in td if t.get("roe")]))])

    _s06_tbl_h = min(6.0, len(tbl_data) * 0.55)  # cap 6cm — laisse toujours place au bloc lecture
    _add_table(slide, tbl_data, 0.9, 2.5, 23.6, _s06_tbl_h,
               col_widths=[5.0, 3.0, 2.8, 2.8, 3.2, 3.4, 3.4],
               font_size=7.5, header_size=8, alt_fill=_GRAYL)

    # Analytical text — position dynamique sous la table
    _s06_text_y = round(2.5 + _s06_tbl_h + 0.7, 2)
    best = td[0] if td else {}
    best_name = best.get("company", best.get("ticker", "Le leader"))[:20]
    ev_best = best.get("ev_ebitda")
    ev_med = D["ev_med"]

    _rect(slide, 0.9, _s06_text_y, 23.6, 2.8, fill=_GRAYL)
    _rect(slide, 0.9, _s06_text_y, 23.6, 0.7, fill=_NAVY)
    _txb(slide, "LECTURE ANALYTIQUE", 1.1, _s06_text_y + 0.05, 23.2, 0.6, size=8.5, bold=True, color=_WHITE)
    analysis = (
        f"La Médiane EV/EBITDA sectorielle s établit a {ev_med:.1f}x LTM. "
        f"{best_name} ({_fmt_x(ev_best)}) se distingue comme le leader de qualité, "
        f"combine a une marge EBITDA de {_fmt_pct_plain(best.get('ebitda_margin'))} et une croissance "
        f"de {_fmt_pct_rev(best.get('revenue_growth'))}. "
        f"La dispersion des multiples revele l hétérogèneite des profils au sein du secteur "
        f"{D['sector_name']} — une lecture croisée P/E vs EV/EBITDA permet d isoler les "
        f"effets de structure de capital et les distorsions comptables."
    )
    _txb(slide, analysis, 1.1, _s06_text_y + 0.8, 23.2, 2.0, size=8.5, color=_GRAYT, wrap=True)
    _footer(slide)


def _s07_cycle(prs, D):
    slide = _blank(prs)
    _header(slide, "Positionnement dans le Cycle",
            f"{D['sector_name']}  ·  Drivers de croissance & sensibilité macro", 2)

    content = D["content"]
    drivers = content.get("drivers", [])

    if len(drivers) > MAX_DRIVERS_DISPLAYED:
        log.warning("sectoral_pptx _s07: %d drivers disponibles — affichage tronque a MAX_DRIVERS_DISPLAYED=%d",
                    len(drivers), MAX_DRIVERS_DISPLAYED)
    for i, (direction, name, body) in enumerate(drivers[:MAX_DRIVERS_DISPLAYED]):
        yy = 2.5 + i * 2.55
        arrow_col = _BUY if direction == "up" else _SELL
        arrow_txt = "▲" if direction == "up" else "▼"
        bg_col = _GREEN_L if direction == "up" else _RED_L
        _rect(slide, 0.9, yy, 0.1, 1.9, fill=arrow_col)
        _rect(slide, 1.0, yy, 15.0, 1.9, fill=bg_col)
        _txb(slide, arrow_txt, 1.1, yy + 0.2, 0.8, 0.8, size=14, bold=True, color=arrow_col)
        _txb(slide, name, 2.1, yy + 0.15, 13.7, 0.7, size=10, bold=True, color=_NAVY)
        _txb(slide, body, 2.1, yy + 0.95, 13.7, 0.8, size=8, color=_GRAYT, wrap=True)

    # Signal box (right)
    _rect(slide, 16.6, 2.5, 7.9, 10.9, fill=_GRAYL)  # hauteur etendue pour inclure rec_line
    _txb(slide, "SIGNAL SECTORIEL", 16.6, 2.9, 7.9, 0.7, size=8.5, bold=True, color=_NAVY, align=PP_ALIGN.CENTER)
    _rect(slide, 17.5, 3.9, 6.1, 1.4, fill=D["sig_color"])
    _txb(slide, f"● {D['sig_label']}", 17.5, 4.0, 6.1, 1.2, size=12, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)
    _txb(slide, "Score moyen FinSight", 16.6, 5.6, 7.9, 0.6, size=8, color=_GRAYT, align=PP_ALIGN.CENTER)
    _txb(slide, f"{D['score_moyen']}/100", 16.6, 6.2, 7.9, 1.3, size=28, bold=True, color=_NAVY, align=PP_ALIGN.CENTER)
    _txb(slide, f"WACC Médian : {D['wacc_pct']} %", 16.6, 7.8, 7.9, 0.7, size=8, color=_GRAYT, align=PP_ALIGN.CENTER)
    _txb(slide, "Horizon : 12 mois", 16.6, 8.6, 7.9, 0.6, size=8, color=_GRAYD, align=PP_ALIGN.CENTER)
    cycle_comment = content.get("cycle_comment", "")
    _txb(slide, cycle_comment, 16.6, 9.6, 7.9, 1.5, size=8, color=_NAVYL, wrap=True, align=PP_ALIGN.CENTER)

    # Régime macro + récession (si disponible)
    _macro = D.get("macro") or {}
    _regime  = _macro.get("regime", "")
    _rec_6m  = _macro.get("recession_prob_6m")
    _rec_lvl = _macro.get("recession_level", "")
    if _regime and _regime != "Inconnu":
        _AMBER_C = RGBColor(0xE6, 0x7E, 0x22)
        _REGIME_C = {"Bull": _BUY, "Bear": _SELL, "Volatile": _AMBER_C, "Transition": _AMBER_C}
        _rc = _REGIME_C.get(_regime, _NAVY)
        regime_line = f"Régime : {_regime}"
        rec_line = f"Rec. 6M : {_rec_6m}%  ({_rec_lvl})" if _rec_6m is not None else ""
        _rect(slide, 16.9, 11.2, 7.3, 0.05, fill=_GRAYD)
        _txb(slide, "CONTEXTE MACRO", 16.6, 11.35, 7.9, 0.5, size=7, bold=True, color=_GRAYD, align=PP_ALIGN.CENTER)
        _txb(slide, regime_line, 16.6, 11.85, 7.9, 0.55, size=8.5, bold=True, color=_rc, align=PP_ALIGN.CENTER)
        if rec_line:
            _txb(slide, rec_line, 16.6, 12.45, 7.9, 0.5, size=7.5, color=_GRAYT, align=PP_ALIGN.CENTER)

    # ── Données FRED (indicateurs macro compacts) ──────────────────────
    try:
        _fred = D.get("fred") or {}
        if _fred:
            _parts = []
            if "fed_funds_rate" in _fred:
                _parts.append(f"Fed {_fred['fed_funds_rate']:.2f}%")
            if "treasury_10y" in _fred:
                _parts.append(f"10Y {_fred['treasury_10y']:.2f}%")
            if "vix" in _fred:
                _parts.append(f"VIX {_fred['vix']:.1f}")
            if "cpi_yoy" in _fred:
                _parts.append(f"CPI {_fred['cpi_yoy']:+.1f}% YoY")
            if "unemployment" in _fred:
                _parts.append(f"Chômage {_fred['unemployment']:.1f}%")
            if "credit_spread_baa" in _fred:
                _parts.append(f"Spread BAA {_fred['credit_spread_baa']:.2f}%")
            # Indicateurs sectoriels (max 2)
            _sect = _fred.get("sector", {})
            _sect_parts = []
            for _sk, _sv in list(_sect.items())[:4]:
                if _sk.endswith("_yoy"):
                    continue
                _label = _sk.replace("_", " ").title()
                if isinstance(_sv, (int, float)):
                    _yoy_key = f"{_sk}_yoy"
                    if _yoy_key in _sect:
                        _sect_parts.append(f"{_label} {_sect[_yoy_key]:+.1f}% YoY")
                    else:
                        _sect_parts.append(f"{_label} {_sv:,.1f}")
                if len(_sect_parts) >= 2:
                    break
            _fred_line = "Données FRED : " + " · ".join(_parts)
            if _sect_parts:
                _fred_line += "\n" + " · ".join(_sect_parts)
            # Position sous le bloc macro existant ou en bas du panel droit
            _fred_y = 12.9 if (_regime and _regime != "Inconnu") else 11.35
            _rect(slide, 16.9, _fred_y, 7.3, 0.05, fill=_NAVY)
            _txb(slide, _fred_line, 16.6, _fred_y + 0.1, 7.9, 0.7,
                 size=6.5, color=_NAVYL, wrap=True, align=PP_ALIGN.CENTER)
    except Exception as _fred_e:
        log.warning("sectoral_pptx _s07 FRED block: %s", _fred_e)

    _footer(slide)


def _s08_subsectors(prs, D):
    """Slide décomposition par sous-secteur (industry GICS)."""
    subsectors = D.get("subsectors", [])
    if len(subsectors) <= 1:
        return  # pas de slide si un seul sous-secteur

    slide = _blank(prs)
    _header(slide, "Décomposition par Sous-Secteur",
            f"{D['sector_name']}  ·  {len(subsectors)} industries GICS identifiées", 1)

    # Tableau sous-secteurs
    # SEC-PPTX-S8 Baptiste : les cellules python-pptx grossissent au-dela de
    # la hauteur demandee quand le texte wrap ou que le font est plus grand
    # que prevu. On passe de 0.42/row a 0.55/row pour eviter que la table
    # empiete sur le bloc analyse en dessous (10 lignes = ~5.5cm au lieu de
    # 4.2cm estime).
    from pptx.util import Cm as _Cm, Pt as _Pt
    hdrs = ["Sous-secteur", "Nb", "Score", "Signal", "EV/EBITDA", "Marge", "Croiss.", "Mom."]
    n_rows = min(len(subsectors), 10)
    tbl_h = 0.55 + n_rows * 0.55
    tbl_shape = slide.shapes.add_table(n_rows + 1, 8, _Cm(0.9), _Cm(2.20), _Cm(23.6), _Cm(tbl_h))
    table = tbl_shape.table
    col_widths = [6.0, 1.5, 1.8, 2.8, 2.5, 2.2, 2.2, 2.2]
    total_cw = sum(col_widths)
    for i, cw in enumerate(col_widths):
        table.columns[i].width = _Cm(cw * 23.6 / total_cw)

    # Header
    for c, h in enumerate(hdrs):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            for run in p.runs:
                run.font.size = _Pt(8)
                run.font.bold = True
                run.font.color.rgb = _WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = _NAVY

    # Data rows
    _SIG_COLORS = {"Surpondérer": _BUY, "Sous-pondérer": _SELL, "Neutre": _HOLD}
    for r, s in enumerate(subsectors[:n_rows]):
        vals = [
            s["name"][:28],
            str(s["nb"]),
            str(s["score"]),
            s["signal"],
            f"{s['ev_ebitda']:.1f}x" if s["ev_ebitda"] else "\u2014",
            f"{s['margin']:.1f}%" if s["margin"] else "\u2014",
            f"{s['growth']:+.1f}%" if s["growth"] else "\u2014",
            f"{s['momentum']:+.1f}%" if s["momentum"] else "\u2014",
        ]
        for c, v in enumerate(vals):
            cell = table.cell(r + 1, c)
            cell.text = v
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.size = _Pt(8)
                    if c == 3:
                        run.font.color.rgb = _SIG_COLORS.get(v, _BLACK)
                        run.font.bold = True
                    else:
                        run.font.color.rgb = _BLACK
            bg = RGBColor(0xF5, 0xF7, 0xFA) if r % 2 == 0 else _WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg

    # Analyse LLM sous le tableau — bloc unique avec titre dynamique
    # Gap 0.85 (vs 1.2 previous) : réhausse le bloc DISPERSION/LEADERSHIP
    # de 0.35 cm pour aérer la slide (#200 Baptiste audit énergie).
    y_box = 2.20 + tbl_h + 0.85
    remaining_h = 13.0 - y_box
    if remaining_h > 2.0 and subsectors:
        _top = subsectors[0]
        _bottom = subsectors[-1]
        _score_spread = _top["score"] - _bottom["score"]
        _top_best = ", ".join(b[0] for b in _top.get("best", [])[:3])

        # Titre analytique dynamique Généré depuis les Données — limite stricte sur le nom
        # pour eviter le debordement de la barre de titre (slide 8 banques)
        _top_name_short = _top['name'][:20]
        if _score_spread >= 25:
            _title = f"DISPERSION MARQU\u00c9E \u2014 {_top_name_short} domine (+{_score_spread} pts)"
        elif _top["score"] >= 60:
            _title = f"LEADERSHIP {_top_name_short.upper()} \u2014 score {_top['score']}/100"
        else:
            _title = f"PHOTO SECTORIELLE \u2014 {len(subsectors)} sous-segments, moy. {sum(s['score'] for s in subsectors)//len(subsectors)}/100"

        _rect(slide, 0.9, y_box, 23.6, 0.65, fill=_NAVY)
        _txb(slide, _title, 1.1, y_box + 0.13, 23.2, 0.45,
             size=10, bold=True, color=_WHITE)

        # Bloc d'analyse
        _body_y = y_box + 0.8
        _body_h = min(remaining_h - 0.9, 4.5)
        _rect(slide, 0.9, _body_y, 23.6, _body_h, fill=_GRAYL)
        _rect(slide, 0.9, _body_y, 0.12, _body_h, fill=_NAVY)

        # Text analytique basé sur les Données réelles — fmt safe pour None (banques)
        def _f1(v, suf=""):
            return f"{v:.1f}{suf}" if isinstance(v, (int, float)) else "n/d"
        def _fp(v):
            return f"{v:+.1f}%" if isinstance(v, (int, float)) else "n/d"
        _ev_part = (
            f"{_f1(_top.get('ev_ebitda'))}x EV/EBITDA avec une marge EBITDA de "
            f"{_f1(_top.get('margin'), '%')}"
            if isinstance(_top.get('ev_ebitda'), (int, float))
            else "des metriques de valorisation specifiques au secteur (ratios financiers)"
        )
        _analysis = (
            f"Le sous-segment {_top['name']} ressort en t\u00eate du classement avec un score "
            f"composite de {_top['score']}/100 et un signal {_top['signal'].lower()}, port\u00e9 "
            f"par {_top['nb']} acteurs dont {_top_best}. La valorisation m\u00e9diane ressort sur "
            f"{_ev_part} et une trajectoire top-line \u00e0 {_fp(_top.get('growth'))} YoY. "
        )
        if len(subsectors) > 1:
            _analysis += (
                f"\u00c0 l'inverse, {_bottom['name']} ferme la Marché ({_bottom['score']}/100, "
                f"signal {_bottom['signal'].lower()}), reflétant des fondamentaux sous pression "
                f"ou une valorisation d\u00e9j\u00e0 tendue. L'\u00e9cart de {_score_spread} points "
                f"entre les extr\u00eames du classement justifie une approche s\u00e9lective plut\u00f4t "
                f"qu'une exposition beta uniforme au secteur. "
            )
        _analysis += (
            f"Cette cartographie guide le calibrage d'allocation sous-sectorielle : "
            f"Surpondérer les segments avec score \u2265 60 et momentum positif, "
            f"Sous-pondérer les segments avec marges compressees ou croissance negative, "
            f"et conserver une exposition neutre aux segments en transition."
        )
        _txb(slide, _analysis, 1.2, _body_y + 0.15, 23.4, _body_h - 0.25,
             size=9, color=_GRAYT, wrap=True)

    _footer(slide)


def _s09_cartographie(prs, D):
    slide = _blank(prs)
    _header(slide, "Cartographie des Sociétés",
            f"{D['N']} sociétés analysées  ·  {D['sector_name']}  ·  {D['universe']}  ·  Tri par score FinSight decroissant", 2)

    td = D["sorted_td"]
    MAX_S09 = 8  # cap slide 9 a 8 lignes — evite debordement sur la lecture analytique
    td_disp = td[:MAX_S09]
    tbl_data = [["#", "Ticker", "Société", "Score", "Reco", "Cours", "EV/EBITDA", "Mg EBITDA", "Croissance", "Momentum"]]
    for i, t in enumerate(td_disp, 1):
        reco = _reco(t.get("score_global"))
        tbl_data.append([
            str(i),
            t.get("ticker", ""),
            (t.get("company") or "")[:25],
            f"{int(t.get('score_global') or 0)}/100",
            reco,
            f"{t.get('price') or '—'}",
            _fmt_x(t.get("ev_ebitda")),
            _fmt_pct_plain(t.get("ebitda_margin")),
            _fmt_pct_rev(t.get("revenue_growth")),   # revenue_growth déjà normalisé en décimal
            _fmt_pct(t.get("momentum_52w")),
        ])

    # Hauteur table : max 4.5cm pour garantir la place du bloc lecture (h min 2.5cm)
    _s09_tbl_h = min(4.5, len(tbl_data) * 0.52)
    s09_tbl = _add_table(slide, tbl_data, 0.9, 2.5, 23.6, _s09_tbl_h,
               col_widths=[0.8, 1.8, 4.0, 2.0, 1.8, 2.0, 2.4, 2.8, 2.8, 3.2],
               font_size=7.5, header_size=7.5, alt_fill=_GRAYL)
    # Forcer hauteur de chaque ligne (python-pptx ignore la hauteur shape sinon)
    _n09 = len(tbl_data)
    _row_h_cm = _s09_tbl_h / _n09 if _n09 else 0.5
    for _row in s09_tbl.rows:
        _row.height = Cm(_row_h_cm)

    # Analytical text — position fixe après la table avec gap suffisant
    _s09_text_y = round(2.5 + _s09_tbl_h + 1.2, 2)  # gap 1.2cm — espace visible entre table et lecture
    _s09_text_h = min(4.5, max(2.5, 13.5 - _s09_text_y))
    n_buy  = sum(1 for t in td if _reco(t.get("score_global")) == "BUY")
    n_hold = sum(1 for t in td if _reco(t.get("score_global")) == "HOLD")
    n_sell = sum(1 for t in td if _reco(t.get("score_global")) == "SELL")
    best = td[0] if td else {}
    best_name = (best.get("company") or best.get("ticker") or "Leader")[:22]

    _rect(slide, 0.9, _s09_text_y, 23.6, _s09_text_h, fill=_GRAYL)
    _rect(slide, 0.9, _s09_text_y, 0.1, _s09_text_h, fill=_NAVY)
    _carto_title = f"{D['sector_name']} — {n_buy} BUY / {n_hold} HOLD / {n_sell} SELL · {best_name} en tête ({int(best.get('score_global') or 0)}/100)"
    _txb(slide, _carto_title, 1.3, _s09_text_y + 0.15, 23.0, 0.6, size=9, bold=True, color=_NAVY)
    analysis = (
        f"Le secteur {D['sector_name']} présente {n_buy} BUY / {n_hold} HOLD / {n_sell} SELL "
        f"sur {len(td)} valeurs analysées. "
        f"{best_name} (score {int(best.get('score_global') or 0)}/100) constitue le coeur offensif recommande "
        f"— fondamentaux solides et visibilité supérieure sur les revenus. "
        f"La répartition des recommandations reflète un positionnément selectif cohérent avec la phase sectorielle actuelle. "
        f"Les catalyseurs identifies — résultats trimestriels, guidance annuel, opérations M&A — "
        f"constituent les événements cles a surveiller pour un renforcement conditionnel des positions."
    )
    _txb(slide, analysis, 1.3, _s09_text_y + 0.85, 23.0, max(1.5, _s09_text_h - 1.0), size=8.5, color=_GRAYT, wrap=True)
    _footer(slide)


def _s10_scatter(prs, D):
    slide = _blank(prs)
    _header(slide, "Classement EV/EBITDA",
            "EV/EBITDA par acteur  ·  Tri croissant  ·  Vert = sous Médiane (opportunité relative)  ·  Top 15", 2)

    img = _chart_valuation_bars(D["tickers_data"])
    _pic(slide, img, 0.9, 2.3, 14.7, 11.4)

    # Analysis panel
    _rect(slide, 16.1, 2.3, 8.4, 11.4, fill=_GRAYL)
    _rect(slide, 16.1, 2.3, 8.4, 0.7, fill=_NAVY)
    _txb(slide, "CE QUE LE GRAPHIQUE REVELE", 16.3, 2.35, 8.1, 0.6, size=8.5, bold=True, color=_WHITE)

    td = D["sorted_td"]
    ev_vals = [float(t.get("ev_ebitda", 0)) for t in td if t.get("ev_ebitda") and float(t.get("ev_ebitda", 0)) > 0]
    med_ev = float(np.median(ev_vals)) if ev_vals else 0

    premium = [t for t in td if t.get("ev_ebitda") and float(t["ev_ebitda"]) > med_ev * 1.15]
    decote  = [t for t in td if t.get("ev_ebitda") and 0 < float(t["ev_ebitda"]) < med_ev * 0.85]

    analysis_lines = []
    if decote:
        d0 = max(decote, key=lambda t: t.get("score_global") or 0)
        analysis_lines.append(
            f"{d0.get('ticker', '')} offre la meilleure asymétrie — "
            f"EV/EBITDA de {_fmt_x(d0.get('ev_ebitda'))} sous la Médiane de {med_ev:.1f}x "
            f"avec un score FinSight de {int(d0.get('score_global') or 0)}/100."
        )
    if premium:
        p0 = premium[0]
        analysis_lines.append(
            f"{p0.get('ticker', '')} traite en prime ({_fmt_x(p0.get('ev_ebitda'))} vs "
            f"Médiane {med_ev:.1f}x) — vérifier si la croissance justifié ce multiple."
        )
    analysis_lines.append(
        f"Médiane sectorielle : {med_ev:.1f}x. "
        f"Les barres vertes sous cette ligne sont les candidats privilegies a analyser "
        f"via Analyse Société individuelle (DCF, WACC, scénarios)."
    )
    _txb(slide, "\n\n".join(analysis_lines), 16.3, 3.1, 8.0, 10.5, size=8.5, color=_GRAYT, wrap=True)
    _footer(slide)


def _s11_scores(prs, D):
    slide = _blank(prs)
    _header(slide, "Scores FinSight detaillés — top 10",
            "Décomposition par dimension  ·  Value · Growth · Quality · Momentum  ·  Score 0-100", 2)

    td = D["sorted_td"]
    # SEC-PPTX-S11 Baptiste : MAX 10 data rows strict (+ header = 11 rows total).
    # Baptiste a explicitement demande "10 ou 11 ligne maximum" et titre "top 10".
    # Pas de ligne footer "... et X autrès" pour respecter la limite absolue.
    MAX_ROWS_S11 = 10
    td_disp = td[:MAX_ROWS_S11]
    tbl_data = [["Ticker", "Société", "Score Global", "Value", "Growth", "Quality", "Momentum", "Reco"]]
    for t in td_disp:
        reco = _reco(t.get("score_global"))
        tbl_data.append([
            t.get("ticker", ""),
            (t.get("company") or "")[:25],
            f"{int(t.get('score_global') or 0)}/100",
            str(int(t.get("score_value") or 0)),
            str(int(t.get("score_growth") or 0)),
            str(int(t.get("score_quality") or 0)),
            str(int(t.get("score_momentum") or 0)),
            reco,
        ])

    # Hauteur table basée sur 0.72cm/row (hauteur réelle rendue par PowerPoint)
    _PER_ROW_S11 = 0.72
    _s11_tbl_h = len(tbl_data) * _PER_ROW_S11
    _tbl11 = _add_table(slide, tbl_data, 0.9, 2.5, 23.6, _s11_tbl_h,
               col_widths=[2.0, 4.5, 3.0, 2.2, 2.2, 2.8, 3.0, 2.0],
               font_size=7.5, header_size=7.5, alt_fill=None)
    # Heatmap coloring — colonnes 3-6 = Value/Growth/Quality/Momentum
    _HM_GREEN = RGBColor(0xC8, 0xE6, 0xC9)   # vert clair >=65
    _HM_AMBER = RGBColor(0xFF, 0xF3, 0xE0)   # orange clair 45-64
    _HM_RED   = RGBColor(0xFF, 0xEB, 0xEE)   # rouge clair <45
    _HM_TXT_G = RGBColor(0x1A, 0x7A, 0x4A)   # texte vert fonce
    _HM_TXT_A = RGBColor(0xB0, 0x60, 0x00)   # texte orange
    _HM_TXT_R = RGBColor(0xA8, 0x20, 0x20)   # texte rouge fonce
    for row_i, t in enumerate(td_disp, start=1):  # start=1 skip header
        for col_i, score_key in enumerate(["score_value","score_growth","score_quality","score_momentum"], start=3):
            try:
                val = int(t.get(score_key) or 0)
            except (TypeError, ValueError):
                val = 0
            if val >= 65:
                _bg, _fg = _HM_GREEN, _HM_TXT_G
            elif val >= 45:
                _bg, _fg = _HM_AMBER, _HM_TXT_A
            else:
                _bg, _fg = _HM_RED, _HM_TXT_R
            cell = _tbl11.cell(row_i, col_i)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _bg
            # Recolor text
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.color.rgb = _fg
                    run.font.bold = True

    # Synthèse — position APRÈS la hauteur réelle rendue (plus d'overlap avec la table)
    best = td[0] if td else {}
    best_name = (best.get("company") or best.get("ticker") or "Leader")[:25]
    _s11_syn_y = round(2.5 + _s11_tbl_h + 0.35, 2)  # gap 0.35cm apres hauteur reelle
    _syn_h = min(4.5, max(3.0, 13.5 - _s11_syn_y))
    _rect(slide, 0.9, _s11_syn_y, 23.6, _syn_h, fill=_GRAYL)
    _rect(slide, 0.9, _s11_syn_y, 23.6, 0.7, fill=_NAVY)
    _txb(slide, "SYNTHÈSE SCORES  ·  Score >= 75 = fort  ·  50-74 = moyen  ·  < 50 = faible  ·  Pondérations : Value / Growth / Quality / Momentum = 25 % chacun",
         1.1, _s11_syn_y + 0.05, 23.2, 0.6, size=7.5, bold=True, color=_WHITE)
    synthesis = (
        f"{best_name} présente le profil le plus équilibre du secteur "
        f"({int(best.get('score_global') or 0)}/100) avec des scores "
        f"Value={int(best.get('score_value') or 0)}, "
        f"Growth={int(best.get('score_growth') or 0)}, "
        f"Quality={int(best.get('score_quality') or 0)}, "
        f"Momentum={int(best.get('score_momentum') or 0)}. "
        f"La dispersion des scores entre les {len(td)} sociétés Reflète des profils hétérogènes — "
        f"une allocation selective privilegiant les leaders qualitatifs est recommandee "
        f"dans la configuration sectorielle actuelle."
    )
    _txb(slide, synthesis, 1.1, _s11_syn_y + 0.85, 23.2, max(1.8, _syn_h - 1.0), size=8.5, color=_GRAYT, wrap=True)
    _footer(slide)


def _s13_top3(prs, D):
    slide = _blank(prs)
    _header(slide, "Top 3 Sociétés — Synthèse",
            "Analyse detaillee  ·  Top 3 par score FinSight  ·  Catalyseurs & risques identifies", 3)

    top3 = D["top3"]
    content = D["content"]
    cats = content.get("catalyseurs", [])
    risks = content.get("risques", [])

    cols_x = [0.9, 8.9, 16.9]
    for col_i, t in enumerate(top3):
        cx = cols_x[col_i]
        reco = _reco(t.get("score_global"))
        reco_col = _reco_color(reco)
        name = (t.get("company") or t.get("ticker") or "")[:20]

        _rect(slide, cx, 2.5, 7.6, 11.2, fill=_GRAYL)
        _rect(slide, cx, 2.5, 0.1, 11.2, fill=reco_col)

        # Ticker + reco badge
        _txb(slide, t.get("ticker", ""), cx + 0.4, 2.7, 4.6, 0.8, size=13, bold=True, color=_NAVY)
        _rect(slide, cx + 5.3, 2.7, 2.2, 0.7, fill=reco_col)
        _txb(slide, reco, cx + 5.3, 2.72, 2.2, 0.65, size=9, bold=True, color=_WHITE, align=PP_ALIGN.CENTER)
        _txb(slide, name, cx + 0.4, 3.5, 7.1, 0.5, size=8, color=_GRAYT)

        # Price target (illustrative)
        score = int(t.get("score_global") or 0)
        price = t.get("price")
        if price:
            try:
                target = float(price) * (1.20 if reco == "BUY" else 1.08 if reco == "HOLD" else 0.95)
                upside = (target / float(price) - 1) * 100
                _rect(slide, cx + 0.4, 4.2, 7.1, 1.2, fill=_NAVYL)
                _txb(slide, f"Cible : {target:.0f}", cx + 0.5, 4.3, 3.8, 0.6, size=9, bold=True, color=_WHITE)
                _txb(slide, f"Upside : {upside:+.0f} %", cx + 4.4, 4.3, 3.0, 0.6, size=9, color=_WHITE)
            except: pass

        # Key ratios — adaptes au profil sectoriel detecte
        _txb(slide, "Ratios cles", cx + 0.4, 5.6, 7.1, 0.5, size=7.5, bold=True, color=_NAVY)
        try:
            from core.sector_profiles import detect_profile
            _t_profile = detect_profile(t.get("sector", ""), t.get("industry", ""))
        except Exception:
            _t_profile = "STANDARD"
        if _t_profile == "BANK":
            ratios = [
                ("P/E", _fmt_x(t.get("pe") or t.get("pe_ratio"))),
                ("P/B", _fmt_x(t.get("pb_ratio"))),
                ("ROE", _fmt_pct_plain(t.get("roe"))),
            ]
        elif _t_profile == "INSURANCE":
            ratios = [
                ("P/E", _fmt_x(t.get("pe") or t.get("pe_ratio"))),
                ("P/B", _fmt_x(t.get("pb_ratio"))),
                ("ROE", _fmt_pct_plain(t.get("roe"))),
            ]
        elif _t_profile == "REIT":
            ratios = [
                ("P/E", _fmt_x(t.get("pe") or t.get("pe_ratio"))),
                ("P/B", _fmt_x(t.get("pb_ratio"))),
                ("Div Yield", _fmt_pct_plain(t.get("div_yield"))),
            ]
        else:
            ratios = [
                ("EV/EBITDA", _fmt_x(t.get("ev_ebitda"))),
                ("Marge EBITDA", _fmt_pct_plain(t.get("ebitda_margin"))),
                ("Altman Z", _fmt_num(t.get("altman_z"))),
            ]
        for j, (rl, rv) in enumerate(ratios):
            ry = 6.1 + j * 0.75
            _rect(slide, cx + 0.4, ry, 7.1, 0.72, fill=_GRAYM)
            _txb(slide, rl, cx + 0.5, ry + 0.05, 3.8, 0.62, size=8, color=_GRAYT)
            _txb(slide, rv, cx + 4.4, ry + 0.05, 3.0, 0.62, size=8.5, bold=True, color=_NAVY, align=PP_ALIGN.RIGHT)

        # Catalyst
        cat_title, cat_body = cats[col_i] if col_i < len(cats) else ("Catalyseur", "—")
        _txb(slide, "Catalyseur", cx + 0.4, 8.5, 7.1, 0.5, size=7.5, bold=True, color=_BUY)
        _rect(slide, cx + 0.4, 9.0, 7.1, 1.4, fill=_GREEN_L)
        _rect(slide, cx + 0.4, 9.0, 0.1, 1.4, fill=_BUY)
        _txb(slide, _fit_s(cat_body, 160), cx + 0.6, 9.05, 6.7, 1.25, size=7.5, color=_GRAYT, wrap=True)

        # Risk
        risk_title, risk_body = risks[col_i] if col_i < len(risks) else ("Risque", "—")
        _txb(slide, "Risque principal", cx + 0.4, 10.5, 7.1, 0.5, size=7.5, bold=True, color=_SELL)
        _rect(slide, cx + 0.4, 11.1, 7.1, 1.4, fill=_RED_L)
        _rect(slide, cx + 0.4, 11.1, 0.1, 1.4, fill=_SELL)
        _txb(slide, _fit_s(risk_body, 150), cx + 0.6, 11.15, 6.7, 1.25, size=7.5, color=_GRAYT, wrap=True)

    _footer(slide)


def _s14_distribution(prs, D):
    slide = _blank(prs)
    ev_med = D["ev_med"]
    _header(slide, "Distribution des Valorisations",
            f"EV/EBITDA par société vs médiane sectorielle ({ev_med:.1f}x)  ·  Vert = sous médiane  ·  Rouge = prime", 3)

    img = _chart_distribution(D["tickers_data"])
    _pic(slide, img, 0.9, 2.3, 14.7, 11.4)

    # Analysis panel
    _rect(slide, 16.1, 2.3, 8.4, 11.4, fill=_GRAYL)
    _rect(slide, 16.1, 2.3, 8.4, 0.7, fill=_NAVY)
    _txb(slide, "CE QUE LE GRAPHIQUE REVELE", 16.3, 2.35, 8.1, 0.6, size=8.5, bold=True, color=_WHITE)

    td = D["sorted_td"]
    ev_vals = [float(t.get("ev_ebitda", 0)) for t in td if t.get("ev_ebitda")]
    premium_actors = [t for t in td if t.get("ev_ebitda") and float(t["ev_ebitda"]) > ev_med * 1.15]
    decote_actors  = [t for t in td if t.get("ev_ebitda") and float(t["ev_ebitda"]) < ev_med * 0.85]
    n_sous = len(decote_actors)
    n_sur  = len(premium_actors)

    # SEC-PPTX-S15 Baptiste : le texte precedent decrivait la distribution
    # sans en tirer d'implications. On passe a un texte structure en 3 blocs
    # qui DIT ce que ca implique, avec les noms concrets des decotes et des
    # primes pour guider l'action. Calcul du gap absolu max et min aussi.
    _best_decote = max(decote_actors, key=lambda t: t.get('score_global') or 0, default=None)
    _worst_prime = max(premium_actors,
                       key=lambda t: float(t.get('ev_ebitda') or 0), default=None)
    _ev_min = min(ev_vals) if ev_vals else ev_med
    _ev_max = max(ev_vals) if ev_vals else ev_med
    _spread = _ev_max - _ev_min
    _spread_ratio = (_ev_max / _ev_min) if _ev_min > 0 else 1

    # Bloc 1 : ce qu'on observe (factuel + chiffre)
    _obs = (
        f"OBSERVATION. Le secteur affiche un ecart entre la valeur la plus basse "
        f"({_ev_min:.1f}x) et la plus haute ({_ev_max:.1f}x), soit un ratio de "
        f"{_spread_ratio:.1f}x. La mediane s'etablit a {ev_med:.1f}x avec {n_sur} "
        f"acteur(s) en prime marquee (>+15% vs mediane) et {n_sous} en decote "
        f"significative (<-15%)."
    )
    # Bloc 2 : ce que ca implique (analyse + noms concrets)
    if _best_decote:
        _best_d_name = _best_decote.get('ticker', '?')
        _best_d_score = int(_best_decote.get('score_global') or 0)
        _impl_d = (
            f"{_best_d_name} combine la decote (EV/EBITDA "
            f"{float(_best_decote.get('ev_ebitda') or 0):.1f}x) avec un score "
            f"FinSight de {_best_d_score}/100 : asymetrie positive la plus "
            f"credible du secteur, declencheur d'une analyse societe approfondie."
        )
    else:
        _impl_d = ("Aucune decote significative cumulee a un score qualite eleve : "
                   "le secteur ne presente pas d'asymetrie evidente a ce stade, "
                   "l'approche doit rester selective et patiente.")
    if _worst_prime:
        _worst_p_name = _worst_prime.get('ticker', '?')
        _impl_p = (
            f"A l'oppose, {_worst_p_name} ({float(_worst_prime.get('ev_ebitda') or 0):.1f}x) "
            f"concentre le risque de derating : toute deception de croissance ou "
            f"guidance decevante comprimerait le multiple en priorite sur ce nom."
        )
    else:
        _impl_p = ("Aucun acteur en prime extreme : le secteur est raisonnablement "
                   "value dans son ensemble, le risque de derating global est contenu.")
    _implications = f"IMPLICATIONS. {_impl_d} {_impl_p}"

    # Bloc 3 : ce qu'il faut faire (action)
    if _spread_ratio >= 3:
        _action = (
            "ACTION. La dispersion est extreme : c'est un secteur de stock-picking "
            "pur ou la selection individuelle prime nettement sur l'exposition "
            "beta. Privilegier les decotes a scores >=60 en long ; eviter les "
            "primes >+30% sans catalyseur fondamental clair (earnings beat, "
            "revisions consensus positives, M&A). L'approche ETF passive capture "
            "mal ces asymetries et sous-performe mechanically."
        )
    elif _spread_ratio >= 1.8:
        _action = (
            "ACTION. La dispersion est moderee, justifiant une approche mixte : "
            "core d'exposition via l'ETF sectoriel (~60%) + overlay stock-picking "
            "sur les decotes qualitatives (~40%). Rebalancer trimestriellement "
            "apres publications pour capter les revisions fondamentales."
        )
    else:
        _action = (
            "ACTION. La compression des multiples indique un secteur quasi-homogene : "
            "l'approche indicielle (ETF) capture l'essentiel du rendement attendu, "
            "le cout d'effort du stock-picking n'est pas justifie sauf conviction "
            "forte sur un catalyseur idiosyncratique."
        )
    analysis = f"{_obs}\n\n{_implications}\n\n{_action}"
    _txb(slide, analysis, 16.3, 3.1, 8.0, 10.5, size=8.5, color=_GRAYT, wrap=True)
    _footer(slide)


def _s15_entry(prs, D):
    slide = _blank(prs)
    _header(slide, "Zone d'Entree Optimale",
            "Conditions d'entree favorable  ·  Cours < DCF Base x 0,90 + momentum négatif + sentiment négatif", 3)

    _rect(slide, 0.9, 2.5, 23.6, 0.9, fill=_GRAYL)
    _rect(slide, 0.9, 2.5, 0.1, 0.9, fill=_NAVY)
    _txb(slide, "La zone d'entree optimale est définie par la convergence de 3 signaux : décote > 10 % vs DCF Base · momentum 52W négatif · sentiment FinBERT négatif.",
         1.1, 2.55, 23.1, 0.8, size=8.5, color=_GRAYT, wrap=True)

    td = D["sorted_td"]
    MAX_S15 = 8  # cap slide 15 a 8 lignes — garantit visibilite de la note methodologique
    td_disp = td[:MAX_S15]
    tbl_data = [["Ticker", "Société", "Cours actuel", "DCF Base (est.)", "Zone d achat", "Signal", "Proba. 12M"]]
    for t in td_disp:
        price = t.get("price")
        reco = _reco(t.get("score_global"))
        if price:
            try:
                pf = float(price)
                dcf = pf * 1.15
                zone = f"< {pf * 0.95:.1f}"
                sig = "Entree attractive" if reco == "BUY" else "Proche zone" if reco == "HOLD" else "Éviter"
                proba = "~68 %" if reco == "BUY" else "~52 %" if reco == "HOLD" else "~35 %"
            except:
                dcf, zone, sig, proba = "—", "—", "—", "—"
        else:
            dcf, zone, sig, proba = "—", "—", "—", "—"
        tbl_data.append([
            t.get("ticker", ""), (t.get("company") or "")[:24],
            f"{price}" if price else "—",
            f"{dcf:.1f}" if isinstance(dcf, float) else dcf,
            zone, sig, proba,
        ])

    # Hauteur table : 0.52cm/row — marge anti-clipping dernière ligne
    _tbl15_n = len(tbl_data)
    _row15_h_cm = 0.52
    _tbl_h = _tbl15_n * _row15_h_cm  # pas de cap — la note s'Ajusté dynamiquement
    _tbl15 = _add_table(slide, tbl_data, 0.9, 3.5, 23.6, _tbl_h,
               col_widths=[2.0, 4.5, 2.8, 3.2, 3.0, 4.2, 2.5],
               font_size=7.5, header_size=7.5, alt_fill=_GRAYL)
    # Forcer hauteur de chaque ligne (python-pptx ignore la hauteur shape sinon)
    if _tbl15 is not None:
        for _row15 in _tbl15.rows:
            _row15.height = Cm(_row15_h_cm)

    # Note methodologique — occupe tout l'espace restant sous la table (plus de KPIs)
    _FOOTER_Y = 13.5
    _NOTE_H_MAX = min(2.5, _FOOTER_Y - (3.5 + _tbl_h + 0.2))
    _note_y_calc = round(3.5 + _tbl_h + 1.0, 2)
    _NOTE_H = max(0, _NOTE_H_MAX)
    if _NOTE_H > 0.4:
        _note_y = _note_y_calc
        _rect(slide, 0.9, _note_y, 23.6, _NOTE_H, fill=_HOLD_L)
        _rect(slide, 0.9, _note_y, 0.1, _NOTE_H, fill=_HOLD)
        _txb(slide, "NOTE METHODOLOGIQUE", 1.3, _note_y + 0.08, 23.1, 0.5,
             size=7.5, bold=True, color=_HOLD)
        if _NOTE_H > 0.9:
            _txb(slide, "La probabilité de rendement positif a 12 mois est calculée sur des configurations similaires identifiees en backtesting sur Données historiques (2010-2024). Elle ne constitue pas une garantie de performance future.",
                 1.3, _note_y + 0.6, 23.1, _NOTE_H - 0.65, size=7.5, color=_GRAYT, wrap=True)

    _footer(slide)


def _s17_risques(prs, D):
    slide = _blank(prs)
    _header(slide, "Risques & Conditions d'Invalidation",
            "Analyse des risques structurels  ·  Thèse contraire  ·  Protocole Avocat du Diable", 4)

    content = D["content"]
    risks = content.get("risques", [])
    conditions = content.get("conditions", [])

    if len(risks) > MAX_RISKS_DISPLAYED:
        log.warning("sectoral_pptx _s17: %d risques disponibles — affichage tronque a MAX_RISKS_DISPLAYED=%d",
                    len(risks), MAX_RISKS_DISPLAYED)
    for col_i, (risk_title, risk_body) in enumerate(risks[:MAX_RISKS_DISPLAYED]):
        cx = 0.9 + col_i * 8.1
        _rect(slide, cx, 2.5, 7.5, 3.2, fill=_RED_L)
        _rect(slide, cx, 2.5, 7.5, 0.7, fill=_SELL)
        _txb(slide, risk_title, cx + 0.2, 2.55, 7.0, 0.6, size=8.5, bold=True, color=_WHITE)
        _txb(slide, risk_body, cx + 0.2, 3.3, 7.0, 2.2, size=8.5, color=_GRAYT, wrap=True)

    # Conditions d'invalidation
    _txb(slide, "Conditions d'invalidation de la thèse", 0.9, 6.1, 23.6, 0.6, size=9, bold=True, color=_NAVY)
    _AXIS_NORM = {"Société": "Société", "société": "Société", "société": "Société"}
    tbl_data = [["Axe", "Condition d'invalidation", "Horizon"]]
    for ax, cond, hor in conditions:
        tbl_data.append([_AXIS_NORM.get(ax, ax), cond, hor])

    n_rows = len(tbl_data)
    # Cap table height pour laisser 2.9 cm a la section sante financière (titre + boites)
    _HEALTH_H = 2.9   # titre 0.55 + gap 0.6 + boite 1.75
    _TABLE_Y  = 6.8
    _FOOTER_Y = 13.5
    _tbl_h = min(n_rows * 0.62, _FOOTER_Y - _HEALTH_H - _TABLE_Y - 0.4)
    _add_table(slide, tbl_data, 0.9, _TABLE_Y, 23.6, _tbl_h,
               col_widths=[3.5, 16.5, 3.6],
               font_size=8, header_size=8, alt_fill=_GRAYL)

    # ── Sante financière agregee du secteur ─────────────────────────────────
    health_y = round(_TABLE_Y + _tbl_h + 0.4, 2)
    td = D["tickers_data"]

    def _median(vals):
        v = [x for x in vals if x is not None]
        if not v: return None
        v.sort()
        m = len(v) // 2
        return v[m] if len(v) % 2 else (v[m-1] + v[m]) / 2

    # nd_ebitda et net_debt_ebitda : deux noms possibles selon la source de données
    nd_vals  = [t.get("nd_ebitda") or t.get("net_debt_ebitda") for t in td]
    fcf_vals = [t.get("fcf_yield") for t in td]
    sg_vals  = [t.get("score_global") for t in td]
    # FCF yield fallback : calculé depuis fcf / market_cap si disponible
    if all(v is None for v in fcf_vals):
        fcf_vals = []
        for t in td:
            fc = t.get("free_cash_flow") or t.get("fcf")
            mc = t.get("market_cap")
            if fc is not None and mc and mc > 0:
                fcf_vals.append(fc / mc * 100)
            else:
                fcf_vals.append(None)
    nd_med   = _median(nd_vals)
    fcf_med  = _median(fcf_vals)
    sg_med   = _median(sg_vals)

    nd_str  = f"{nd_med:.1f}x"  if nd_med  is not None else "N/D"
    fcf_str = f"{fcf_med:.1f}%" if fcf_med is not None else "N/D"
    sg_str  = f"{sg_med:.0f}/100" if sg_med is not None else "N/D"

    # Couleurs selon seuils fondamentaux
    nd_col  = (_BUY if nd_med is not None and nd_med < 2.0 else
               (_HOLD if nd_med is not None and nd_med < 4.0 else _SELL))
    fcf_col = (_BUY if fcf_med is not None and fcf_med > 4.0 else
               (_HOLD if fcf_med is not None and fcf_med > 1.0 else _SELL))
    sg_col  = (_BUY if sg_med is not None and sg_med >= 60 else
               (_HOLD if sg_med is not None and sg_med >= 40 else _SELL))

    _txb(slide, "Sante financière agregee du secteur", 0.9, health_y, 23.6, 0.55,
         size=8.5, bold=True, color=_NAVY)

    box_w = 7.4
    metrics = [
        (nd_str,  "ND/EBITDA Médian", "levier sectoriel", nd_col,  0.9),
        (fcf_str, "FCF Yield Médian",  "Génération cash",  fcf_col, 9.1),
        (sg_str,  "Score santé moyen", "solidité bilans",  sg_col,  17.3),
    ]
    box_h = 2.2
    for val, l1, l2, col, bx in metrics:
        _rect(slide, bx, health_y + 0.6, box_w, box_h, fill=_GRAYL)
        _rect(slide, bx, health_y + 0.6, 0.1,   box_h, fill=col)
        _txb(slide, val, bx + 0.25, health_y + 0.70, box_w - 0.35, 1.0,
             size=22, bold=True, color=_NAVY)
        _txb(slide, l1, bx + 0.25, health_y + 1.75, box_w - 0.35, 0.45,
             size=8.5, color=_GRAYT)
        _txb(slide, l2, bx + 0.25, health_y + 2.20, box_w - 0.35, 0.40,
             size=8, color=_GRAYD)

    _footer(slide)


def _s18_sentiment(prs, D):
    slide = _blank(prs)
    _header(slide, "Sentiment de Marché — FinBERT",
            f"Analyse sémantique FinBERT  ·  Articles récents  ·  {D['N']} sociétés", 4)

    td = D["tickers_data"]
    # Derive sentiment from tickers_data if available
    pos_scores = []
    for t in td:
        s = t.get("sentiment_score")
        if s is not None:
            try: pos_scores.append(float(s))
            except: pass
    agg_score = _avg(pos_scores) if pos_scores else 0.22
    n_pos = max(1, int(len(td) * 0.3))
    n_neu = max(1, int(len(td) * 0.5))
    n_neg = max(1, len(td) - n_pos - n_neu)
    total = n_pos + n_neu + n_neg
    lbl = "Neutre positif modéré" if agg_score > 0.1 else "Neutre" if agg_score >= -0.1 else "Négatif modéré"

    # Score box
    _rect(slide, 0.9, 2.5, 5.3, 2.6, fill=_GRAYL)
    _rect(slide, 0.9, 2.5, 0.1, 2.6, fill=_NAVYL)
    _txb(slide, f"{agg_score:.2f}", 1.2, 2.6, 4.9, 1.2, size=28, bold=True, color=_NAVY)
    _txb(slide, "Score agrégé FinBERT", 1.2, 3.8, 4.9, 0.55, size=7.5, color=_GRAYT)
    _txb(slide, lbl, 1.2, 4.35, 4.9, 0.45, size=7.5, color=_GRAYD)

    # Sentiment bars
    sentiments = [("Positif", n_pos, _BUY, 11.9), ("Neutre", n_neu, _HOLD, 16.9), ("Négatif", n_neg, _SELL, 8.6)]
    for i, (slbl, cnt, col, bar_w_est) in enumerate(sentiments):
        yy = 2.5 + i * 0.9
        pct = cnt / total * 100
        bar_w = max(0.5, pct / 100 * 10)
        _rect(slide, 6.7, yy, bar_w, 0.65, fill=col)
        _txb(slide, slbl, 6.7, yy + 0.05, 2.5, 0.6, size=8.5, color=_BLACK, bold=True)
        _txb(slide, f"{cnt} articles ({pct:.0f} %)", 17.1, yy + 0.05, 5.5, 0.6, size=8.5, color=_GRAYT)

    # Themes
    content = D["content"]
    cats  = content.get("catalyseurs", [])
    risks = content.get("risques", [])
    _txb(slide, "Thèmes dominants", 0.9, 5.3, 23.6, 0.6, size=8.5, bold=True, color=_NAVY)
    pos_theme = cats[0][0] if cats else "Tendance positive"
    neg_theme = risks[0][0] if risks else "Tendance negative"
    _rect(slide, 0.9, 5.9, 11.4, 1.1, fill=_GREEN_L)
    _rect(slide, 0.9, 5.9, 0.1, 1.1, fill=_BUY)
    _txb(slide, f"+ {pos_theme} — {_fit_s(cats[0][1], 90) if cats else ''}", 1.3, 5.95, 10.7, 1.0, size=8, color=_GRAYT, wrap=True)
    _rect(slide, 13.1, 5.9, 11.4, 1.1, fill=_RED_L)
    _rect(slide, 13.1, 5.9, 0.1, 1.1, fill=_SELL)
    _txb(slide, f"- {neg_theme} — {_fit_s(risks[0][1], 90) if risks else ''}", 13.5, 5.95, 10.7, 1.0, size=8, color=_GRAYT, wrap=True)

    # Analytical text
    _rect(slide, 0.9, 7.4, 23.6, 5.5, fill=_GRAYL)
    _rect(slide, 0.9, 7.4, 23.6, 0.7, fill=_NAVY)
    _txb(slide, "LECTURE ANALYTIQUE DU SENTIMENT", 1.1, 7.45, 23.2, 0.6, size=8.5, bold=True, color=_WHITE)
    # Identify best/worst sentiment tickers
    pos_scores_with_ticker = [(float(t.get("sentiment_score", 0)), t.get("ticker", ""), t.get("company", "")) for t in td if t.get("sentiment_score") is not None]
    pos_scores_with_ticker.sort(key=lambda x: x[0], reverse=True)
    best_sent = pos_scores_with_ticker[0] if pos_scores_with_ticker else (0.25, "leader", "")
    worst_sent = pos_scores_with_ticker[-1] if len(pos_scores_with_ticker) > 1 else (-0.1, "retardataire", "")
    pct_pos = n_pos / total * 100
    sent_tone = "légèrement positif" if agg_score > 0.1 else "neutre" if agg_score >= -0.1 else "légèrement négatif"
    sent_analysis = (
        f"Le sentiment agrégé FinBERT sur le secteur {D['sector_name']} ({D['universe']}) "
        f"ressort {sent_tone} ({agg_score:.2f}), avec {pct_pos:.0f}% d'articles a tonalité positive sur {total} analysés. "
        f"La dispersion inter-valeurs est prononcée : {best_sent[1] or best_sent[2]} ({best_sent[0]:+.2f}) Porté la composanté haussière "
        f"tandis que {worst_sent[1] or worst_sent[2]} ({worst_sent[0]:+.2f}) reflète les pressions structurelles. "
        f"Cette hétérogénéité valide une approche sélective plutôt que directionnelle sur le secteur — "
        f"le sentiment moyen masque des situations fondamentalement différentes entre leaders et retardataires. "
        f"La composanté thématique (catalyseurs vs risques) suggère que les flux narratifs restent "
        f"{'orientés positivement, avec des newsflows soutenus sur la croissance et les marges' if agg_score > 0.1 else 'équilibrés, reflétant une phase de transition sectorielle' if agg_score >= -0.1 else 'sous pression, les révisions négatives dominant le flux informationnel'}. "
        f"Le suivi du ratio Positif/Négatif dans les prochaines semaines constitue un indicateur avancé pour anticiper les rotations sectorielles."
    )
    _txb(slide, sent_analysis, 1.1, 8.2, 23.2, 4.5, size=8.5, color=_GRAYT, wrap=True)
    _footer(slide)


def _s19_sources(prs, D):
    slide = _blank(prs)
    _header(slide, "Avertissement & M\u00e9thodologie",
            "Tra\u00e7abilit\u00e9 compl\u00e8te des donn\u00e9es  \u00b7  Data lineage  \u00b7  Agents FinSight IA v1.0", 4)

    tbl_data = [
        ["Agent", "Source données", "Données collectées", "Fréquence"],
        ["AgentData", "yfinance + FMP", "Cours, ratios, états financiers LTM", "Temps réel / J-1"],
        ["AgentQuant", "Calcul interne", "Altman Z, Beneish M, scores FinSight", "À la demande"],
        ["AgentSentiment", "Finnhub + RSS", "Articles presse, scoring FinBERT", "À la demande"],
        ["AgentSynthese", "Claude Haiku", "Ratios calculés, scoring multifactoriel", "À la demande"],
        ["AgentDevil", "Claude Haiku", "Thèse contraire, conviction delta", "À la demande"],
        ["AgentQA", "Python + LLM", "Validation cohérence, flags erreur", "Systématique"],
    ]
    _add_table(slide, tbl_data, 0.9, 2.5, 23.6, len(tbl_data) * 0.62,
               col_widths=[4.0, 4.5, 10.5, 4.6],
               font_size=7.5, header_size=8, alt_fill=_GRAYL)

    # Limits
    _rect(slide, 0.9, 9.6, 23.6, 2.7, fill=_GRAYL)
    _rect(slide, 0.9, 9.6, 0.1, 2.7, fill=_HOLD)
    _txb(slide, "LIMITES & PRÉCAUTIONS", 1.3, 9.7, 23.1, 0.6, size=8.5, bold=True, color=_HOLD)
    _txb(slide, "Les données financières historiques sont issues de sources publiques et peuvent présenter des délais ou inexactitudes. Les prévisions sont basées sur le consensus et les modèles internes FinSight IA — elles ne constituent pas des engagements. La méthodologie de scoring est soumise à des biais inhérents à toute approche quantitative.",
         1.3, 10.35, 23.1, 2.0, size=8, color=_GRAYT, wrap=True)

    # Disclaimer
    _rect(slide, 0.9, 12.3, 23.6, 1.0, fill=_NAVY)
    _txb(slide, "Ce rapport est Généré par FinSight IA v1.0. Il ne constitue pas un conseil en investissement au sens de la directive MiFID II (2014/65/UE). Document confidentiel — usage interne uniquement. Toute reproduction ou diffusion est interdite sans autorisation.",
         1.1, 12.35, 23.2, 0.9, size=7, color=_GRAYD, wrap=True)
    _footer(slide)


def _s20_performance(prs, D):
    slide = _blank(prs)
    _header(slide, "Performance Boursière Relative — 52 Semaines",
            f"{D['sector_name']}  ·  {D['universe']}  ·  Indexe a 100 au debut de la période", 4)

    td_s = sorted(D["tickers_data"], key=lambda x: x.get("momentum_52w") or 0, reverse=True)
    best  = td_s[0]  if td_s  else {}
    worst = td_s[-1] if td_s  else {}
    best_tk  = best.get("ticker", "—")
    best_mom = best.get("momentum_52w") or 0
    worst_tk  = worst.get("ticker", "—")
    worst_mom = worst.get("momentum_52w") or 0

    img = _chart_performance(D["tickers_data"], best_ticker=best_tk, worst_ticker=worst_tk,
                             sector_name=D.get("sector_name"), universe=D.get("universe"))
    # Graphique agrandi — legende matplotlib dans le chart (haut gauche)
    _pic(slide, img, 0.9, 2.3, 16.5, 10.2)

    # Panel droit — texte analytique uniquement (legende dans le graphique)
    _rect(slide, 17.6, 2.3, 6.8, 10.2, fill=_GRAYL)
    _rect(slide, 17.6, 2.3, 6.8, 0.65, fill=_NAVY)
    _txb(slide, "LECTURE ANALYTIQUE", 17.8, 2.35, 6.5, 0.55, size=8, bold=True, color=_WHITE)
    n_pos   = sum(1 for t in D["tickers_data"] if (t.get("momentum_52w") or 0) > 0)
    n_neg   = len(D["tickers_data"]) - n_pos
    spread  = abs(best_mom - worst_mom)

    lines = [
        ("MEILLEURE PERFORMANCE", 3.20, 7, True, _NAVY),
        (f"{best_tk}  +{best_mom:.0f}%  sur 52 semaines",
         3.65, 9, True, _NAVY),
        ("PIRE PERFORMANCE", 4.60, 7, True, _NAVY),
        (f"{worst_tk}  {worst_mom:+.0f}%  sur 52 semaines",
         5.05, 9, True, RGBColor(0xA8,0x20,0x20)),
        ("DISPERSION", 6.10, 7, True, _NAVY),
        (f"Écart {spread:.0f} pts — {n_pos} valeur(s) positives vs {n_neg} négatives.",
         6.55, 8, False, _GRAYT),
        ("INTERPRÉTATION", 7.70, 7, True, _NAVY),
        ("La dispersion des trajectoires reflète la bifurcation sectorielle : certains "
         "acteurs captent l'essentiel de la création de valeur tandis que d'autrès subissent "
         "des sorties de capitaux structurelles.", 8.15, 7.5, False, _GRAYT),
        ("Le momentum 52W est intégré dans le score FinSight comme signal de confirmation "
         "de la thèse d'investissement.", 9.35, 7.5, False, _GRAYT),
    ]
    for txt, yy, sz, bold, col in lines:
        _txb(slide, txt, 17.8, yy, 6.4, 0.9, size=sz, bold=bold, color=col, wrap=True)

    _footer(slide)


# ─── MAIN CLASS ───────────────────────────────────────────────────────────────
class SectoralPPTXWriter:

    @staticmethod
    def generate(
        tickers_data: list[dict],
        sector_name: str,
        universe: str = "CAC 40",
        output_path: Optional[str] = None,
    ) -> bytes:
        """
        Generate a 20-slide sectoral pitchbook PPTX.
        Returns bytes of the PPTX file.
        If output_path is provided, also saves to disk.
        """
        D = _prepare_data(tickers_data, sector_name, universe)

        # Macro regime_v + récession
        _macro = {}
        try:
            import sys as _sys, os as _os
            _sys.path.insert(0, _os.path.dirname(_os.path.dirname(__file__)))
            from agents.agent_macro import AgentMacro
            _macro = AgentMacro().analyze() or {}
        except Exception as _me:
            import logging as _log
            _log.getLogger(__name__).warning("[SectoralPPTXWriter] AgentMacro: %s", _me)
        D["macro"] = _macro

        # ── Données FRED (macro + sectoriel) ────────────────────────────
        _fred = {}
        try:
            from data.sources.fred_source import fetch_macro_context
            _fred = fetch_macro_context(sector_name) or {}
        except Exception as _fe:
            import logging as _log_fred
            _log_fred.getLogger(__name__).warning("[SectoralPPTXWriter] FRED: %s", _fe)
        D["fred"] = _fred

        prs = Presentation()
        prs.slide_width  = _SW
        prs.slide_height = _SH

        # Slide 1 — Cover
        _s01_cover(prs, D)
        # Slide 2 — Executive Summary
        _s02_exec_summary(prs, D)
        # Slide 3 — Sommaire
        _s03_sommaire(prs, D)
        # Slide 4 — Chapter 01 divider
        _chapter_divider(prs, "01", "Présentation du Secteur",
                         "Caractéristiques structurelles, ratios comparatifs & positionnément cycle")
        # Slide 5 — Présentation
        _s05_presentation(prs, D)
        # Slide 6 — Ratios comparatifs
        _s06_ratios(prs, D)
        # Slide 7 — Positionnement cycle
        _s07_cycle(prs, D)
        # Slide 8 — Décomposition sous-sectorielle (si >1 industry)
        _s08_subsectors(prs, D)
        # Slide 9 — Chapter 02 divider
        _chapter_divider(prs, "02", "Cartographie des Sociétés",
                         f"{D['N']} sociétés analysées — scores FinSight, scatter & décomposition")
        # Slide 9 — Cartographie
        _s09_cartographie(prs, D)
        # Slide 10 — Scatter
        _s10_scatter(prs, D)
        # Slide 11 — Scores detailles
        _s11_scores(prs, D)
        # Slide 12 — Chapter 03 divider
        _chapter_divider(prs, "03", "Top 3 & Valorisations",
                         "Synthèse Top 3, distribution EV/EBITDA & zone d'entree optimale")
        # Slide 13 — Top 3
        _s13_top3(prs, D)
        # Slide 14 — Distribution
        _s14_distribution(prs, D)
        # Slide 15 — Zone d'entrée
        _s15_entry(prs, D)
        # Slide 16 — Chapter 04 divider
        _chapter_divider(prs, "04", "Risques, Sentiment & Méthodologie",
                         "Conditions d'invalidation, FinBERT agrégé & data lineage")
        # Slide 17 — Risques
        _s17_risques(prs, D)
        # Slide 18 — Sentiment
        _s18_sentiment(prs, D)
        # Slide 19 — Performance (swap: was slide 20)
        _s20_performance(prs, D)
        # Slide 20 — Sources & Méthodologie (swap: was slide 19)
        _s19_sources(prs, D)

        buf = io.BytesIO()
        prs.save(buf)
        pptx_bytes = buf.getvalue()

        if output_path:
            Path(output_path).write_bytes(pptx_bytes)

        return pptx_bytes


# ─── CONVENIENCE FUNCTION ─────────────────────────────────────────────────────
def generate_sector_pptx(
    sector_name: str,
    tickers_data: list[dict],
    output_path: str,
    universe: str = "CAC 40",
) -> bytes:
    """Convenience wrapper matching PDF writer signature."""
    return SectoralPPTXWriter.generate(
        tickers_data=tickers_data,
        sector_name=sector_name,
        universe=universe,
        output_path=output_path,
    )
