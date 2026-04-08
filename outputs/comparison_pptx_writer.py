# =============================================================================
# FinSight IA — Comparison PPTX Writer (Pitchbook comparatif 22 slides)
# outputs/comparison_pptx_writer.py
#
# Genere un pitchbook comparatif IB-quality en 22 slides via python-pptx.
# Accepte deux etats pipeline LangGraph + metriques extraites.
#
# Usage :
#   from outputs.comparison_pptx_writer import ComparisonPPTXWriter
#   path = ComparisonPPTXWriter().generate(state_a, state_b, output_path)
# =============================================================================

from __future__ import annotations

import io
import logging
import math
import sys
from datetime import date, datetime
from pathlib import Path
from typing import Any, Dict, List, Optional

log = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Palette couleurs — identique a pptx_writer.py
# ---------------------------------------------------------------------------
NAVY       = "1B3A6B"
NAVY_MID   = "2E5FA3"
NAVY_PALE  = "EEF3FA"
GREEN      = "1A7A4A"
GREEN_PALE = "EAF4EF"
GREEN_MID  = "2E8B57"
RED        = "A82020"
RED_PALE   = "FAF0EF"
AMBER      = "B06000"
AMBER_PALE = "FDF6E8"
WHITE      = "FFFFFF"
GREY_BG    = "F7F8FA"
GREY_TXT   = "555555"
GREY_LIGHT = "888888"
BLACK      = "0D0D0D"

# Couleurs distinctives par societe
COLOR_A     = NAVY_MID   # Societe A — bleu
COLOR_A_PAL = NAVY_PALE
COLOR_B     = GREEN_MID  # Societe B — vert
COLOR_B_PAL = GREEN_PALE

import re as _re

def _x(text) -> str:
    if text is None:
        return ""
    s = str(text)
    # Supprime caracteres invalides XML 1.0
    s = _re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', s)
    # Supprime markdown bold/italic du LLM (**text** et *text*)
    s = _re.sub(r'\*\*(.+?)\*\*', r'\1', s, flags=_re.DOTALL)
    s = _re.sub(r'\*(.+?)\*', r'\1', s, flags=_re.DOTALL)
    # Supprime les asterisques residuels isoles
    s = _re.sub(r'\*+', '', s)
    return s


def rgb(hex_str: str):
    from pptx.dml.color import RGBColor
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ---------------------------------------------------------------------------
# Helpers numeriques
# ---------------------------------------------------------------------------
def _fr(v, dp: int = 1, suffix: str = "") -> str:
    if v is None:
        return "\u2014"
    try:
        return f"{float(v):,.{dp}f}".replace(",", " ").replace(".", ",") + suffix
    except Exception:
        return "\u2014"


def _frpct(v, signed: bool = False) -> str:
    if v is None:
        return "\u2014"
    try:
        fv = float(v)
        s = f"{fv * 100:+.1f}" if signed else f"{fv * 100:.1f}"
        return s.replace(".", ",") + " %"
    except Exception:
        return "\u2014"


def _frpct_margin(v) -> str:
    if v is None:
        return "\u2014"
    try:
        fv = float(v)
        if abs(fv) > 2.0:
            fv /= 100.0
        return f"{fv * 100:.1f}".replace(".", ",") + " %"
    except Exception:
        return "\u2014"


def _frx(v) -> str:
    if v is None:
        return "\u2014"
    try:
        return f"{float(v):.1f}".replace(".", ",") + "x"
    except Exception:
        return "\u2014"


def _frm(v, cur_sym: str = "$") -> str:
    if v is None:
        return "\u2014"
    try:
        v = float(v)
        if abs(v) > 1_000_000_000:
            v = v / 1_000_000
        if cur_sym == "EUR":
            sym_big, sym_small = "Md\u20ac", "M\u20ac"
        else:
            sym_big, sym_small = "Mds" + cur_sym, "M" + cur_sym
        if abs(v) >= 1000:
            return _fr(v / 1000, 1) + " " + sym_big
        return _fr(v, 0) + " " + sym_small
    except Exception:
        return "\u2014"


def _upside(target, current) -> str:
    if not target or not current:
        return "\u2014"
    try:
        return f"{(float(target) / float(current) - 1) * 100:+.1f}".replace(".", ",") + " %"
    except Exception:
        return "\u2014"


def _g(obj, attr, default=None):
    if obj is None:
        return default
    if isinstance(obj, dict):
        return obj.get(attr, default)
    return getattr(obj, attr, default)


def _safe_str(v, default="\u2014") -> str:
    if v is None:
        return default
    s = str(v).strip()
    return s if s else default


def _truncate(s, n: int) -> str:
    s = _safe_str(s, "")
    if len(s) <= n:
        return s
    cut = s[:n]
    last_space = cut.rfind(" ")
    if last_space > n // 2:
        cut = cut[:last_space]
    return cut + "\u2026"


def _fit(s, n: int) -> str:
    """Coupe a n chars sans ajouter '...' — zones ou le debordement est invisible (PPTX clip)."""
    s = _safe_str(s, "")
    if len(s) <= n:
        return s
    cut = s[:n]
    last_space = cut.rfind(" ")
    return cut[:last_space] if last_space > n // 2 else cut


def _fr_date_long(d=None) -> str:
    from datetime import date as _d
    _MONTHS = ["janvier", "fevrier", "mars", "avril", "mai", "juin",
               "juillet", "aout", "septembre", "octobre", "novembre", "decembre"]
    if d is None:
        d = _d.today()
    elif isinstance(d, str):
        try:
            d = _d.fromisoformat(str(d).split("T")[0][:10])
        except Exception:
            return str(d)
    try:
        return f"{d.day} {_MONTHS[d.month - 1]} {d.year}"
    except Exception:
        return str(d)


# ---------------------------------------------------------------------------
# Helpers formes PPTX
# ---------------------------------------------------------------------------
def _set_slide_bg(slide, hex_color: str):
    from pptx.oxml.ns import qn
    from lxml import etree
    cSld = slide._element.find(qn('p:cSld'))
    if cSld is None:
        return
    existing = cSld.find(qn('p:bg'))
    if existing is not None:
        cSld.remove(existing)
    bg_xml = (
        '<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<p:bgPr>'
        f'<a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>'
        '<a:effectLst/>'
        '</p:bgPr>'
        '</p:bg>'
    )
    bg_elem = etree.fromstring(bg_xml)
    spTree = cSld.find(qn('p:spTree'))
    cSld.insert(list(cSld).index(spTree) if spTree is not None else 0, bg_elem)


def _add_text_alpha(slide, x, y, w, h, text, font_size,
                    color_hex="FFFFFF", alpha_val=15000, bold=False):
    from pptx.util import Cm, Pt
    from pptx.oxml.ns import qn
    from lxml import etree
    txBox = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = _x(text)
    run.font.name = "Calibri"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    r_elem = run._r
    rPr = r_elem.find(qn('a:rPr'))
    if rPr is None:
        rPr = etree.Element(qn('a:rPr'))
        r_elem.insert(0, rPr)
    for tag in (qn('a:solidFill'), qn('a:gradFill'), qn('a:noFill'), qn('a:pattFill')):
        for child in list(rPr.findall(tag)):
            rPr.remove(child)
    solidFill = etree.SubElement(rPr, qn('a:solidFill'))
    srgbClr   = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', color_hex.upper())
    alpha_elem = etree.SubElement(srgbClr, qn('a:alpha'))
    alpha_elem.set('val', str(alpha_val))
    rPr.remove(solidFill)
    rPr.insert(0, solidFill)
    return txBox


def add_rect(slide, x, y, w, h, fill_hex, line_hex=None, line_width_pt=0):
    from pptx.util import Cm, Pt
    shape = slide.shapes.add_shape(1, Cm(x), Cm(y), Cm(w), Cm(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_hex)
    if line_hex:
        shape.line.color.rgb = rgb(line_hex)
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, x, y, w, h, text, font_size=10, color_hex=BLACK,
                 bold=False, italic=False, align=None,
                 font_name="Calibri", wrap=True):
    from pptx.util import Cm, Pt
    from pptx.enum.text import PP_ALIGN
    if align is None:
        align = PP_ALIGN.LEFT
    txBox = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = _x(text) if text is not None else "\u2014"
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color_hex)
    return txBox


def _set_cell(cell, text, font_size=8, bold=False, color_hex=BLACK,
              fill_hex=None, align=None, font_name="Calibri", italic=False):
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    if align is None:
        align = PP_ALIGN.CENTER
    if fill_hex:
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(fill_hex)
    tf = cell.text_frame
    for para in tf.paragraphs:
        for run in list(para.runs):
            para._p.remove(run._r)
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = _x(text) if text is not None else "\u2014"
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color_hex)


def add_table(slide, x, y, w, h, num_rows, num_cols, col_widths_pct=None,
              header_data=None, rows_data=None,
              header_fill=NAVY, row_fills=None, border_hex=None):
    from pptx.util import Cm, Emu
    from pptx.enum.text import PP_ALIGN
    total_rows = num_rows + (1 if header_data else 0)
    tbl = slide.shapes.add_table(
        total_rows, num_cols, Cm(x), Cm(y), Cm(w), Cm(h)
    ).table
    total_emu = int(Cm(w))
    if col_widths_pct:
        for ci, pct in enumerate(col_widths_pct):
            tbl.columns[ci].width = int(total_emu * pct)
    else:
        col_w = total_emu // num_cols
        for ci in range(num_cols):
            tbl.columns[ci].width = col_w
    row_offset = 0
    if header_data:
        for ci, val in enumerate(header_data):
            cell = tbl.cell(0, ci)
            _set_cell(cell, val, font_size=7.5, bold=True,
                      color_hex=WHITE, fill_hex=header_fill,
                      align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)
        row_offset = 1
    if rows_data:
        default_fills = [WHITE, GREY_BG]
        for ri, row in enumerate(rows_data):
            if row_fills and ri < len(row_fills):
                fill = row_fills[ri]
            else:
                fill = default_fills[ri % 2]
            for ci, val in enumerate(row):
                cell = tbl.cell(ri + row_offset, ci)
                align = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
                _set_cell(cell, val, font_size=7.5, bold=False,
                          color_hex=BLACK, fill_hex=fill, align=align)
    return tbl


# ---------------------------------------------------------------------------
# Elements communs
# ---------------------------------------------------------------------------
def navy_bar(slide):
    add_rect(slide, 0, 0, 25.4, 1.65, NAVY)


def footer_bar(slide):
    add_rect(slide, 0, 13.39, 25.4, 0.89, NAVY)
    add_text_box(slide, 1.02, 13.44, 23.37, 0.56,
                 "FinSight IA  \u00b7  Usage confidentiel", 7, WHITE)


def section_dots(slide, active: int):
    from pptx.enum.text import PP_ALIGN
    xs = [17.02, 17.65, 18.29, 18.92, 19.56]
    for i, x in enumerate(xs):
        num = i + 1
        fill = WHITE if num == active else NAVY_MID
        add_rect(slide, x, 0.33, 0.51, 0.51, fill)
        color = NAVY if num == active else WHITE
        add_text_box(slide, x + 0.02, 0.34, 0.47, 0.44,
                     str(num), 7, color, bold=True, align=PP_ALIGN.CENTER)


def slide_title(slide, title: str, subtitle: str = ""):
    from pptx.enum.text import PP_ALIGN
    add_text_box(slide, 1.02, 0.33, 15.49, 0.97,
                 title, 13, WHITE, bold=True)
    if subtitle:
        add_text_box(slide, 1.02, 1.73, 23.37, 0.56,
                     subtitle, 9, NAVY_MID)


def kpi_box(slide, x, y, w, h, value, label, sub="",
            fill=NAVY_PALE, accent=NAVY_MID):
    from pptx.enum.text import PP_ALIGN
    add_rect(slide, x, y, w, h, fill)
    add_rect(slide, x, y, 0.13, h, accent)
    val_color = NAVY if fill != NAVY else WHITE
    lbl_color = GREY_TXT if fill != NAVY else WHITE
    sub_color = GREY_LIGHT if fill != NAVY else NAVY_PALE
    add_text_box(slide, x + 0.25, y + 0.15, w - 0.38, 1.0,
                 value, 17, val_color, bold=True)
    add_text_box(slide, x + 0.25, y + 1.1, w - 0.38, 0.50,
                 label, 7.5, lbl_color)
    if sub:
        add_text_box(slide, x + 0.25, y + 1.38, w - 0.38, 0.42,
                     sub, 7, sub_color)


def divider_slide(prs, number_str: str, title: str, subtitle: str):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_bg(slide, NAVY)
    add_rect(slide, 0, 0, 0.3, 14.29, NAVY_MID)
    _add_text_alpha(slide, 1.27, 4.06, 22.86, 4.57,
                    number_str, 80, "FFFFFF", 15000, bold=True)
    add_text_box(slide, 1.27, 5.21, 21.59, 2.29,
                 title, 28, WHITE, bold=True)
    add_rect(slide, 1.27, 7.75, 7.62, 0.01, GREY_LIGHT)
    add_text_box(slide, 1.27, 8.08, 17.78, 0.89,
                 subtitle, 10, "AABBDD")
    add_text_box(slide, 1.02, 13.34, 23.37, 0.56,
                 "FinSight IA  \u00b7  Usage confidentiel", 7, "6677AA")
    return slide


def _company_header_band(slide, tkr_a, tkr_b, name_a="", name_b=""):
    """Bande bicolore sous le titre identifiant les deux societes."""
    lbl_a = f"{tkr_a}" + (f"  {_truncate(name_a, 22)}" if name_a else "")
    lbl_b = f"{tkr_b}" + (f"  {_truncate(name_b, 22)}" if name_b else "")
    add_rect(slide, 1.02, 1.73, 11.44, 0.56, COLOR_A_PAL)
    add_rect(slide, 1.02, 1.73, 0.18, 0.56, COLOR_A)
    add_text_box(slide, 1.35, 1.78, 11.0, 0.46, lbl_a, 8.5, NAVY, bold=True)
    add_rect(slide, 12.94, 1.73, 11.44, 0.56, COLOR_B_PAL)
    add_rect(slide, 12.94, 1.73, 0.18, 0.56, COLOR_B)
    add_text_box(slide, 13.27, 1.78, 11.0, 0.46, lbl_b, 8.5, GREEN, bold=True)


# ---------------------------------------------------------------------------
# LLM Synthesis — appel unique au debut
# ---------------------------------------------------------------------------
def _call_llm(prompt: str, system: str = "", max_tokens: int = 512) -> str:
    """Appel LLM avec fallback anthropic -> mistral."""
    try:
        sys_path_insert = str(Path(__file__).parent.parent)
        if sys_path_insert not in sys.path:
            sys.path.insert(0, sys_path_insert)
        from core.llm_provider import LLMProvider
        for provider in ("anthropic", "mistral", "gemini"):
            try:
                llm = LLMProvider(provider=provider)
                return llm.generate(prompt, system=system, max_tokens=max_tokens)
            except Exception as _e:
                log.debug(f"[cmp_pptx] LLM {provider} failed: {_e}")
                continue
    except Exception as e:
        log.warning(f"[cmp_pptx] LLM unavailable: {e}")
    return ""


def _generate_synthesis(m_a: dict, m_b: dict) -> dict:
    """
    Genere la synthese comparative via LLM.
    Retourne un dict avec les cles : exec_summary, financial_text,
    valuation_text, quality_text, verdict_text, thesis_a, thesis_b.
    """
    tkr_a = m_a.get("ticker_a") or "A"
    tkr_b = m_b.get("ticker_b") or "B"
    name_a = m_a.get("company_name_a") or tkr_a
    name_b = m_b.get("company_name_b") or tkr_b

    def _pct(v):
        if v is None: return "N/A"
        try: return f"{float(v)*100:.1f}%"
        except: return "N/A"
    def _x2(v):
        if v is None: return "N/A"
        try: return f"{float(v):.1f}x"
        except: return "N/A"
    def _n(v):
        if v is None: return "N/A"
        try: return f"{float(v):.1f}"
        except: return "N/A"

    data_str = f"""
{tkr_a} ({name_a}) vs {tkr_b} ({name_b})

VALORISATION:
{tkr_a}: PE={_n(m_a.get('pe_ratio'))}, EV/EBITDA={_x2(m_a.get('ev_ebitda'))}, PB={_n(m_a.get('price_to_book'))}, FCF yield={_pct(m_a.get('fcf_yield'))}
{tkr_b}: PE={_n(m_b.get('pe_ratio'))}, EV/EBITDA={_x2(m_b.get('ev_ebitda'))}, PB={_n(m_b.get('price_to_book'))}, FCF yield={_pct(m_b.get('fcf_yield'))}

PROFITABILITE:
{tkr_a}: EBITDA margin={_pct(m_a.get('ebitda_margin_ltm'))}, ROIC={_pct(m_a.get('roic'))}, ROE={_pct(m_a.get('roe'))}, Rev CAGR 3y={_pct(m_a.get('revenue_cagr_3y'))}
{tkr_b}: EBITDA margin={_pct(m_b.get('ebitda_margin_ltm'))}, ROIC={_pct(m_b.get('roic'))}, ROE={_pct(m_b.get('roe'))}, Rev CAGR 3y={_pct(m_b.get('revenue_cagr_3y'))}

LEVIER:
{tkr_a}: ND/EBITDA={_n(m_a.get('net_debt_ebitda'))}, Beta={_n(m_a.get('beta'))}
{tkr_b}: ND/EBITDA={_n(m_b.get('net_debt_ebitda'))}, Beta={_n(m_b.get('beta'))}

SCORES:
{tkr_a}: FinSight={m_a.get('finsight_score') or 'N/A'}, Piotroski={m_a.get('piotroski_score') or 'N/A'}, Rec={m_a.get('recommendation') or 'N/A'}
{tkr_b}: FinSight={m_b.get('finsight_score') or 'N/A'}, Piotroski={m_b.get('piotroski_score') or 'N/A'}, Rec={m_b.get('recommendation') or 'N/A'}
"""

    system_msg = (
        "Tu es un analyste sell-side senior (style JPMorgan Research). "
        "Reponds en francais, concis et rigoureux. "
        "Pas d'emojis. Structure quoi -> pourquoi -> implications investisseur."
    )

    results = {}

    # Executive summary (60 mots max)
    r = _call_llm(
        f"Redige un executive summary comparatif de 60 mots MAX pour {tkr_a} vs {tkr_b}. "
        f"Inclure : quel titre offre le meilleur rapport qualite/valorisation et pourquoi.\n{data_str}",
        system=system_msg, max_tokens=200
    )
    results["exec_summary"] = r[:500] if r else ""

    # Financial commentary (70 mots max)
    r = _call_llm(
        f"Redige un commentaire financier comparatif de 70 mots MAX : "
        f"tendances P&L de {tkr_a} vs {tkr_b}, differences de marges et croissance, "
        f"implications bilan.\n{data_str}",
        system=system_msg, max_tokens=250
    )
    results["financial_text"] = (r[:600].rsplit(' ', 1)[0] if r and len(r) > 600 else r) or ""

    # Valuation commentary (70 mots max)
    r = _call_llm(
        f"Redige un commentaire de valorisation de 70 mots MAX : "
        f"prime/decote de {tkr_a} vs {tkr_b} sur les multiples cles, "
        f"justification et upside relatif.\n{data_str}",
        system=system_msg, max_tokens=250
    )
    results["valuation_text"] = (r[:600].rsplit(' ', 1)[0] if r and len(r) > 600 else r) or ""

    # Quality commentary (60 mots max)
    r = _call_llm(
        f"Redige un commentaire de qualite financiere de 60 mots MAX : "
        f"solidite bilancielle de {tkr_a} vs {tkr_b}, "
        f"Piotroski, levier, risques de manipulation comptable.\n{data_str}",
        system=system_msg, max_tokens=200
    )
    results["quality_text"] = (r[:500].rsplit(' ', 1)[0] if r and len(r) > 500 else r) or ""

    # Verdict final (80 mots max)
    r = _call_llm(
        f"Verdict comparatif final de 80 mots MAX : quel titre privilegier entre {tkr_a} et {tkr_b}, "
        f"pourquoi (catalyseur cle, valorisation relative, qualite bilan), "
        f"et principal risque de la these.\n{data_str}",
        system=system_msg, max_tokens=300
    )
    results["verdict_text"] = r[:700] if r else ""

    def _split_bull_bear(r, sep="|||"):
        """Split bull/bear text — essaie plusieurs separateurs."""
        if not r:
            return "", ""
        for token in [sep, "\n|||", "|||\n", "---", "\n\nBear", "\n\nBEAR",
                      "\nThese bear", "\nThese Bull", "\nBear"]:
            if token in r:
                parts = r.split(token, 1)
                return parts[0].strip()[:350], parts[1].strip()[:350]
        # Dernier recours : chercher un pattern "bear" dans le texte
        m = _re.search(r'(?i)\n+(?:these? )?bear', r)
        if m:
            return r[:m.start()].strip()[:350], r[m.start():].strip()[:350]
        return r[:350], ""

    # Theses bull/bear A (40 mots chacune)
    r = _call_llm(
        f"Pour {tkr_a}, donne 1 these bull (40 mots) et 1 these bear (40 mots) "
        f"en les separant par '|||'. Ne pas ajouter de titre.\n{data_str}",
        system=system_msg, max_tokens=200
    )
    results["bull_a"], results["bear_a"] = _split_bull_bear(r)

    # Theses bull/bear B (40 mots chacune)
    r = _call_llm(
        f"Pour {tkr_b}, donne 1 these bull (40 mots) et 1 these bear (40 mots) "
        f"en les separant par '|||'. Ne pas ajouter de titre.\n{data_str}",
        system=system_msg, max_tokens=200
    )
    results["bull_b"], results["bear_b"] = _split_bull_bear(r)

    # Narratif cours 52 semaines avec contexte macro (pour PDF page 11 + PPTX)
    _perf_a_1y = m_a.get('perf_1y')
    _perf_b_1y = m_b.get('perf_1y')
    _perf_str = (
        f"{tkr_a} perf 1Y: {_perf_a_1y*100:.1f}% " if _perf_a_1y else f"{tkr_a} perf 1Y: N/D "
    ) + (
        f"| {tkr_b} perf 1Y: {_perf_b_1y*100:.1f}%" if _perf_b_1y else f"| {tkr_b} perf 1Y: N/D"
    )
    _price_ctx = (
        f"Secteur A: {m_a.get('sector_a','N/D')} | Secteur B: {m_b.get('sector_b','N/D')} | "
        f"Beta A: {m_a.get('beta','N/D')} | Beta B: {m_b.get('beta','N/D')} | "
        f"{_perf_str} | "
        f"52W High A: {m_a.get('week52_high','N/D')} | 52W Low A: {m_a.get('week52_low','N/D')} | "
        f"52W High B: {m_b.get('week52_high','N/D')} | 52W Low B: {m_b.get('week52_low','N/D')} | "
        f"PE A: {m_a.get('pe_ratio','N/D')} | PE B: {m_b.get('pe_ratio','N/D')} | "
        f"Earnings Growth A: {m_a.get('eps_growth','N/D')} | Earnings Growth B: {m_b.get('eps_growth','N/D')}"
    )
    r = _call_llm(
        f"Analyse en 120 mots MAX la trajectoire boursiere de {tkr_a} vs {tkr_b} sur 12 mois. "
        f"Donne le CONTEXTE MACRO (taux directeurs, cycle economique, rotation sectorielle) "
        f"qui EXPLIQUE les mouvements. Cite des evenements concrets (resultats trimestriels, "
        f"guidance, M&A, regulation) qui ont pu impacter les cours. "
        f"Termine par les catalyseurs a surveiller pour les 3-6 prochains mois. "
        f"Donnees : {_price_ctx}",
        system=system_msg, max_tokens=400
    )
    results["price_narrative"] = r[:800] if r else ""

    return results


# ---------------------------------------------------------------------------
# Graphiques matplotlib
# ---------------------------------------------------------------------------
def _make_chart_buf(fig) -> io.BytesIO:
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    import matplotlib.pyplot as plt
    plt.close(fig)
    buf.seek(0)
    return buf


def _chart_revenue_ebitda(m_a: dict, m_b: dict, tkr_a: str, tkr_b: str) -> Optional[io.BytesIO]:
    """Graphique CA et marge EBITDA comparatif."""
    try:
        import matplotlib.pyplot as plt
        import matplotlib
        matplotlib.use('Agg')
        import numpy as np

        fig, axes = plt.subplots(1, 2, figsize=(10.5, 3.2))
        fig.patch.set_facecolor('white')

        # Couleurs
        ca = '#2E5FA3'
        cb = '#2E8B57'

        # --- Panneau gauche : Marges EBITDA LTM / N-1 / N-2 ---
        ax1 = axes[0]
        ax1.set_facecolor('#FAFAFA')
        years_lbl = ['N-2', 'N-1', 'LTM']
        vals_a = [
            _pct_val(m_a.get('ebitda_margin_y2')),
            _pct_val(m_a.get('ebitda_margin_y1')),
            _pct_val(m_a.get('ebitda_margin_ltm')),
        ]
        vals_b = [
            _pct_val(m_b.get('ebitda_margin_y2')),
            _pct_val(m_b.get('ebitda_margin_y1')),
            _pct_val(m_b.get('ebitda_margin_ltm')),
        ]
        x = np.arange(len(years_lbl))
        w = 0.32
        bars_a = ax1.bar(x - w/2, vals_a, w, label=tkr_a, color=ca, alpha=0.85)
        bars_b = ax1.bar(x + w/2, vals_b, w, label=tkr_b, color=cb, alpha=0.85)
        for bar in list(bars_a) + list(bars_b):
            h = bar.get_height()
            if h and h > 0:
                ax1.text(bar.get_x() + bar.get_width()/2., h + 0.5,
                         f'{h:.0f}%', ha='center', va='bottom', fontsize=7, color='#333')
        _em_a = _pct_val(m_a.get('ebitda_margin_ltm'))
        _em_b = _pct_val(m_b.get('ebitda_margin_ltm'))
        _em_leader = tkr_a if _em_a > _em_b else (tkr_b if _em_b > _em_a else None)
        _em_delta = abs(_em_a - _em_b)
        _em_title = (f"{_em_leader} plus profitabl. : {_em_a:.0f}% vs {_em_b:.0f}% EBITDA LTM (+{_em_delta:.0f}pts)"
                     if _em_leader else f"Marge EBITDA LTM {_em_a:.0f}% = {_em_b:.0f}%")
        ax1.set_title(_em_title, fontsize=8.5, fontweight='bold', color='#1B3A6B', pad=4)
        ax1.set_xticks(x); ax1.set_xticklabels(years_lbl, fontsize=8)
        ax1.set_ylabel('%', fontsize=8); ax1.legend(fontsize=7.5)
        ax1.spines['top'].set_visible(False); ax1.spines['right'].set_visible(False)
        ax1.tick_params(labelsize=8)

        # --- Panneau droit : Revenue CAGR + ROE + ROIC ---
        ax2 = axes[1]
        ax2.set_facecolor('#FAFAFA')
        metrics = ['Rev CAGR\n3y', 'ROIC', 'ROE']
        vals_a2 = [
            _pct_val(m_a.get('revenue_cagr_3y')),
            _pct_val(m_a.get('roic')),
            _pct_val(m_a.get('roe')),
        ]
        vals_b2 = [
            _pct_val(m_b.get('revenue_cagr_3y')),
            _pct_val(m_b.get('roic')),
            _pct_val(m_b.get('roe')),
        ]
        x2 = np.arange(len(metrics))
        bars_a2 = ax2.bar(x2 - w/2, vals_a2, w, label=tkr_a, color=ca, alpha=0.85)
        bars_b2 = ax2.bar(x2 + w/2, vals_b2, w, label=tkr_b, color=cb, alpha=0.85)
        for bar in list(bars_a2) + list(bars_b2):
            h = bar.get_height()
            if h and h > 0:
                ax2.text(bar.get_x() + bar.get_width()/2., h + 0.3,
                         f'{h:.0f}%', ha='center', va='bottom', fontsize=7, color='#333')
        _roic_a = _pct_val(m_a.get('roic'))
        _roic_b = _pct_val(m_b.get('roic'))
        _r_leader = tkr_a if _roic_a > _roic_b else (tkr_b if _roic_b > _roic_a else None)
        _r_title = (f"ROIC : {tkr_a} {_roic_a:.0f}% vs {tkr_b} {_roic_b:.0f}% — {_r_leader} plus efficient"
                    if _r_leader else f"ROIC equivalent : {_roic_a:.0f}%")
        ax2.set_title(_r_title, fontsize=8.5, fontweight='bold', color='#1B3A6B', pad=4)
        ax2.set_xticks(x2); ax2.set_xticklabels(metrics, fontsize=8)
        ax2.set_ylabel('%', fontsize=8); ax2.legend(fontsize=7.5)
        ax2.spines['top'].set_visible(False); ax2.spines['right'].set_visible(False)
        ax2.tick_params(labelsize=8)

        fig.tight_layout(pad=1.2)
        return _make_chart_buf(fig)
    except Exception as e:
        log.warning(f"[cmp_pptx] chart_revenue_ebitda error: {e}")
        return None


def _chart_multiples(m_a: dict, m_b: dict, tkr_a: str, tkr_b: str) -> Optional[io.BytesIO]:
    """Graphique comparaison multiples de valorisation."""
    try:
        import matplotlib.pyplot as plt
        import matplotlib
        matplotlib.use('Agg')
        import numpy as np

        fig, ax = plt.subplots(figsize=(10.5, 4.5))
        fig.patch.set_facecolor('white')
        ax.set_facecolor('#FAFAFA')

        ca = '#2E5FA3'
        cb = '#2E8B57'
        c_sector = '#BBBBBB'

        metrics_lbl = ['PE', 'EV/EBITDA', 'P/B', 'PEG']
        sector = m_a.get('sector_a', '') or ''
        sec_pe = m_a.get('sector_median_pe') or 20.0
        sec_eveb = m_a.get('sector_median_ev_ebitda') or 14.0

        vals_a = [
            _safe_float(m_a.get('pe_ratio')) or 0,
            _safe_float(m_a.get('ev_ebitda')) or 0,
            _safe_float(m_a.get('price_to_book')) or 0,
            _safe_float(m_a.get('peg_ratio')) or 0,
        ]
        vals_b = [
            _safe_float(m_b.get('pe_ratio')) or 0,
            _safe_float(m_b.get('ev_ebitda')) or 0,
            _safe_float(m_b.get('price_to_book')) or 0,
            _safe_float(m_b.get('peg_ratio')) or 0,
        ]
        sector_refs = [sec_pe, sec_eveb, None, None]

        x = np.arange(len(metrics_lbl))
        w = 0.28
        ax.bar(x - w,   vals_a, w, label=tkr_a, color=ca, alpha=0.85)
        ax.bar(x,       vals_b, w, label=tkr_b, color=cb, alpha=0.85)
        # Mediane sectorielle
        s_vals = [v if v is not None else 0 for v in sector_refs]
        ax.bar(x + w, s_vals, w, label='Med. secteur', color=c_sector, alpha=0.7)

        # Labels valeurs
        for i, (va, vb) in enumerate(zip(vals_a, vals_b)):
            if va is not None and va > 0:
                ax.text(x[i] - w, va + 0.3, f'{va:.1f}', ha='center', va='bottom', fontsize=7, color='#1B3A6B')
            if vb is not None and vb > 0:
                ax.text(x[i], vb + 0.3, f'{vb:.1f}', ha='center', va='bottom', fontsize=7, color='#2E8B57')

        ax.set_xticks(x); ax.set_xticklabels(metrics_lbl, fontsize=9)
        _pe_a = _safe_float(m_a.get('pe_ratio')) or 0
        _pe_b = _safe_float(m_b.get('pe_ratio')) or 0
        if _pe_a and _pe_b:
            _cheap_tkr = tkr_a if _pe_a < _pe_b else tkr_b
            _ev_a = _safe_float(m_a.get('ev_ebitda')) or 0
            _ev_b = _safe_float(m_b.get('ev_ebitda')) or 0
            _m_title = f"{_cheap_tkr} moins cher : PE {_pe_a:.0f}x vs {_pe_b:.0f}x  |  EV/EBITDA {_ev_a:.0f}x vs {_ev_b:.0f}x"
        else:
            _m_title = 'Multiples de Valorisation'
        ax.set_title(_m_title, fontsize=9, fontweight='bold', color='#1B3A6B', pad=5)
        ax.legend(fontsize=8)
        ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
        ax.tick_params(labelsize=8)
        fig.tight_layout(pad=1.0)
        return _make_chart_buf(fig)
    except Exception as e:
        log.warning(f"[cmp_pptx] chart_multiples error: {e}")
        return None


def _chart_finsight_score(m_a: dict, m_b: dict, tkr_a: str, tkr_b: str) -> Optional[io.BytesIO]:
    """Graphique FinSight Score comparatif — barres horizontales."""
    try:
        import matplotlib.pyplot as plt
        import matplotlib
        matplotlib.use('Agg')
        import numpy as np

        fig, ax = plt.subplots(figsize=(11.0, 2.8))
        fig.patch.set_facecolor('white')
        ax.set_facecolor('#FAFAFA')

        fs_a = _safe_float(m_a.get('finsight_score')) or 0
        fs_b = _safe_float(m_b.get('finsight_score')) or 0

        labels = [f'{tkr_a}', f'{tkr_b}']
        values = [fs_a, fs_b]
        colors_bar = ['#2E5FA3', '#2E8B57']

        bars = ax.barh(labels, values, color=colors_bar, height=0.55, alpha=0.9)
        for bar, val in zip(bars, values):
            ax.text(val + 0.5, bar.get_y() + bar.get_height()/2.,
                    f'{val:.0f}/100', ha='left', va='center', fontsize=11,
                    fontweight='bold', color='#333')
        ax.set_xlim(0, 110)
        ax.set_xlabel('Score /100', fontsize=9)
        _fs_winner = tkr_a if fs_a > fs_b else (tkr_b if fs_b > fs_a else None)
        _fs_delta = abs(int(fs_a) - int(fs_b))
        _fs_title = (f"FinSight Score : {_fs_winner} favori ({_fs_delta} pts d'avance)"
                     if _fs_winner else "FinSight Score : egalite")
        ax.set_title(_fs_title, fontsize=10, fontweight='bold', color='#1B3A6B')
        ax.axvline(50, color='#AAAAAA', linestyle='--', linewidth=0.8)
        ax.text(50.5, -0.55, 'Neutre', fontsize=7, color='#888')
        ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
        ax.tick_params(labelsize=9)
        fig.tight_layout(pad=1.0)
        return _make_chart_buf(fig)
    except Exception as e:
        log.warning(f"[cmp_pptx] chart_finsight_score error: {e}")
        return None


def _chart_risk_profile(m_a: dict, m_b: dict, tkr_a: str, tkr_b: str) -> Optional[io.BytesIO]:
    """Radar chart profil de risque."""
    try:
        import matplotlib.pyplot as plt
        import matplotlib
        matplotlib.use('Agg')
        import numpy as np

        categories = ['Beta\n(vol)', 'Levier\n(ND/EBITDA)', 'Qualite\n(Piotros.)', 'Momentum\n(3m)', 'FCF\nYield']
        N = len(categories)

        def _norm(v, lo, hi):
            if v is None: return 0.5
            try:
                return max(0.1, min(1.0, (float(v) - lo) / (hi - lo)))
            except: return 0.5

        # Scores normalises sur [0, 1] (plus haut = meilleur)
        scores_a = [
            1 - _norm(m_a.get('beta') or 1.0, 0.3, 2.5),            # beta: moins = mieux
            1 - _norm(m_a.get('net_debt_ebitda') or 2.0, -1, 6),     # levier: moins = mieux
            _norm(m_a.get('piotroski_score') or 5, 1, 9),             # piotroski: plus = mieux
            _norm((m_a.get('perf_3m') or 0) * 100, -20, 30),         # momentum: plus = mieux
            _norm((m_a.get('fcf_yield') or 0) * 100, 0, 10),         # FCF yield: plus = mieux
        ]
        scores_b = [
            1 - _norm(m_b.get('beta') or 1.0, 0.3, 2.5),
            1 - _norm(m_b.get('net_debt_ebitda') or 2.0, -1, 6),
            _norm(m_b.get('piotroski_score') or 5, 1, 9),
            _norm((m_b.get('perf_3m') or 0) * 100, -20, 30),
            _norm((m_b.get('fcf_yield') or 0) * 100, 0, 10),
        ]

        angles = [n / float(N) * 2 * np.pi for n in range(N)]
        angles += angles[:1]
        scores_a += scores_a[:1]
        scores_b += scores_b[:1]

        fig, ax = plt.subplots(figsize=(4.8, 4.0), subplot_kw=dict(polar=True))
        fig.patch.set_facecolor('white')
        ax.set_facecolor('#FAFAFA')

        ax.plot(angles, scores_a, 'o-', linewidth=2, color='#2E5FA3', label=tkr_a)
        ax.fill(angles, scores_a, alpha=0.18, color='#2E5FA3')
        ax.plot(angles, scores_b, 'o-', linewidth=2, color='#2E8B57', label=tkr_b)
        ax.fill(angles, scores_b, alpha=0.18, color='#2E8B57')

        ax.set_xticks(angles[:-1])
        ax.set_xticklabels(categories, fontsize=7.5)
        ax.set_ylim(0, 1)
        ax.set_yticks([0.25, 0.5, 0.75])
        ax.set_yticklabels([], fontsize=0)
        ax.legend(loc='upper right', bbox_to_anchor=(1.3, 1.15), fontsize=8)
        ax.set_title('Profil de Risque', fontsize=9, fontweight='bold', color='#1B3A6B', pad=12)
        ax.grid(color='#CCCCCC', linestyle='--', linewidth=0.5)
        fig.tight_layout(pad=0.5)
        return _make_chart_buf(fig)
    except Exception as e:
        log.warning(f"[cmp_pptx] chart_risk_profile error: {e}")
        return None


def _chart_ebitda_margins(m_a: dict, m_b: dict, tkr_a: str, tkr_b: str) -> Optional[io.BytesIO]:
    """Graphique marges EBITDA LTM / N-1 / N-2."""
    try:
        import matplotlib.pyplot as plt
        import matplotlib
        matplotlib.use('Agg')
        import numpy as np

        fig, ax = plt.subplots(figsize=(5.5, 3.2))
        fig.patch.set_facecolor('white')
        ax.set_facecolor('#FAFAFA')
        ca = '#2E5FA3'; cb = '#2E8B57'
        years_lbl = ['N-2', 'N-1', 'LTM']
        vals_a = [_pct_val(m_a.get('ebitda_margin_y2')), _pct_val(m_a.get('ebitda_margin_y1')), _pct_val(m_a.get('ebitda_margin_ltm'))]
        vals_b = [_pct_val(m_b.get('ebitda_margin_y2')), _pct_val(m_b.get('ebitda_margin_y1')), _pct_val(m_b.get('ebitda_margin_ltm'))]
        x = np.arange(len(years_lbl)); w = 0.32
        bars_a = ax.bar(x - w/2, vals_a, w, label=tkr_a, color=ca, alpha=0.85)
        bars_b = ax.bar(x + w/2, vals_b, w, label=tkr_b, color=cb, alpha=0.85)
        for bar in list(bars_a) + list(bars_b):
            h = bar.get_height()
            if h and h > 0:
                ax.text(bar.get_x() + bar.get_width()/2., h + 0.5, f'{h:.0f}%',
                        ha='center', va='bottom', fontsize=7, color='#333')
        ax.set_title('Marge EBITDA (%)', fontsize=9, fontweight='bold', color='#1B3A6B', pad=4)
        ax.set_xticks(x); ax.set_xticklabels(years_lbl, fontsize=8)
        ax.set_ylabel('%', fontsize=8); ax.legend(fontsize=7.5)
        ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
        ax.tick_params(labelsize=8)
        fig.tight_layout(pad=1.0)
        return _make_chart_buf(fig)
    except Exception as e:
        log.warning(f"[cmp_pptx] chart_ebitda_margins error: {e}")
        return None


def _chart_growth_returns(m_a: dict, m_b: dict, tkr_a: str, tkr_b: str) -> Optional[io.BytesIO]:
    """Graphique Rev CAGR + ROIC + ROE."""
    try:
        import matplotlib.pyplot as plt
        import matplotlib
        matplotlib.use('Agg')
        import numpy as np

        fig, ax = plt.subplots(figsize=(5.5, 3.2))
        fig.patch.set_facecolor('white')
        ax.set_facecolor('#FAFAFA')
        ca = '#2E5FA3'; cb = '#2E8B57'
        metrics = ['Rev CAGR\n3y', 'ROIC', 'ROE']
        vals_a = [_pct_val(m_a.get('revenue_cagr_3y')), _pct_val(m_a.get('roic')), _pct_val(m_a.get('roe'))]
        vals_b = [_pct_val(m_b.get('revenue_cagr_3y')), _pct_val(m_b.get('roic')), _pct_val(m_b.get('roe'))]
        x = np.arange(len(metrics)); w = 0.32
        bars_a = ax.bar(x - w/2, vals_a, w, label=tkr_a, color=ca, alpha=0.85)
        bars_b = ax.bar(x + w/2, vals_b, w, label=tkr_b, color=cb, alpha=0.85)
        for bar in list(bars_a) + list(bars_b):
            h = bar.get_height()
            if h and h > 0:
                ax.text(bar.get_x() + bar.get_width()/2., h + 0.3, f'{h:.0f}%',
                        ha='center', va='bottom', fontsize=7, color='#333')
        ax.set_title('Croissance & Rentabilite (%)', fontsize=9, fontweight='bold', color='#1B3A6B', pad=4)
        ax.set_xticks(x); ax.set_xticklabels(metrics, fontsize=8)
        ax.set_ylabel('%', fontsize=8); ax.legend(fontsize=7.5)
        ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
        ax.tick_params(labelsize=8)
        fig.tight_layout(pad=1.0)
        return _make_chart_buf(fig)
    except Exception as e:
        log.warning(f"[cmp_pptx] chart_growth_returns error: {e}")
        return None


def _chart_52w_price(tkr_a: str, tkr_b: str) -> Optional[io.BytesIO]:
    """Graphique cours boursiers 52 semaines normalises base 100."""
    try:
        import matplotlib.pyplot as plt
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.dates as mdates
        import yfinance as yf
        from datetime import datetime, timedelta

        end = datetime.today()
        start = end - timedelta(days=370)

        df_a = yf.download(tkr_a, start=start, end=end, progress=False, auto_adjust=True)
        df_b = yf.download(tkr_b, start=start, end=end, progress=False, auto_adjust=True)

        if df_a.empty and df_b.empty:
            return None

        fig, ax = plt.subplots(figsize=(12, 4.5))
        fig.patch.set_facecolor('white')
        ax.set_facecolor('#FAFAFA')
        ca = '#2E5FA3'; cb = '#2E8B57'

        if not df_a.empty:
            close_a = df_a['Close'].squeeze()
            base_a = close_a.iloc[0]
            if base_a and base_a > 0:
                ax.plot(close_a.index, close_a / base_a * 100,
                        color=ca, linewidth=1.8, label=tkr_a)

        if not df_b.empty:
            close_b = df_b['Close'].squeeze()
            base_b = close_b.iloc[0]
            if base_b and base_b > 0:
                ax.plot(close_b.index, close_b / base_b * 100,
                        color=cb, linewidth=1.8, label=tkr_b)

        ax.axhline(y=100, color='#BBBBBB', linestyle='--', linewidth=0.9, alpha=0.7)
        # Titre analytique : qui sur/sous-performe sur 52 semaines
        try:
            _p52_a = ((close_a.iloc[-1] / close_a.iloc[0]) - 1) * 100 if not df_a.empty else None
            _p52_b = ((close_b.iloc[-1] / close_b.iloc[0]) - 1) * 100 if not df_b.empty else None
            if _p52_a is not None and _p52_b is not None:
                _sur = tkr_a if _p52_a > _p52_b else tkr_b
                _52_title = f"52 semaines : {tkr_a} {_p52_a:+.1f}% vs {tkr_b} {_p52_b:+.1f}% — {_sur} surperforme"
            elif _p52_a is not None:
                _52_title = f"52 semaines : {tkr_a} {_p52_a:+.1f}%"
            elif _p52_b is not None:
                _52_title = f"52 semaines : {tkr_b} {_p52_b:+.1f}%"
            else:
                _52_title = 'Performance Relative 52 Semaines (base 100)'
        except Exception:
            _52_title = 'Performance Relative 52 Semaines (base 100)'
        ax.set_title(_52_title, fontsize=10, fontweight='bold', color='#1B3A6B', pad=5)
        ax.set_ylabel('Performance (base 100)', fontsize=8)
        ax.legend(fontsize=9)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.tick_params(labelsize=8)
        ax.xaxis.set_major_formatter(mdates.DateFormatter('%b %y'))
        ax.xaxis.set_major_locator(mdates.MonthLocator(interval=2))
        plt.xticks(rotation=25, ha='right')
        fig.tight_layout(pad=1.0)
        return _make_chart_buf(fig)
    except Exception as e:
        log.warning(f"[cmp_pptx] chart_52w_price error: {e}")
        return None


def _pct_val(v) -> float:
    """Convertit une marge en pourcentage pour les graphiques."""
    if v is None: return 0.0
    try:
        fv = float(v)
        if abs(fv) > 2.0: fv /= 100.0
        return round(fv * 100, 1)
    except: return 0.0


def _safe_float(v) -> Optional[float]:
    if v is None: return None
    try: return float(v)
    except: return None


def _insert_chart(slide, buf: Optional[io.BytesIO], x: float, y: float, w_cm: float, h_cm: float):
    """Insere un graphique BytesIO dans une slide."""
    if buf is None:
        return
    try:
        from pptx.util import Cm
        slide.shapes.add_picture(buf, Cm(x), Cm(y), Cm(w_cm), Cm(h_cm))
    except Exception as e:
        log.warning(f"[cmp_pptx] insert_chart error: {e}")


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def _slide_cover(prs, m_a: dict, m_b: dict):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    tkr_a  = m_a.get('ticker_a') or 'A'
    tkr_b  = m_b.get('ticker_b') or 'B'
    name_a = m_a.get('company_name_a') or tkr_a
    name_b = m_b.get('company_name_b') or tkr_b
    gen_date = _fr_date_long()

    rec_a = (m_a.get('recommendation') or 'HOLD').upper()
    rec_b = (m_b.get('recommendation') or 'HOLD').upper()

    def _rec_c(rec):
        if rec == 'BUY': return GREEN, GREEN_PALE
        if rec == 'SELL': return RED, RED_PALE
        return AMBER, AMBER_PALE

    # Top band navy
    add_rect(slide, 0, 0, 25.4, 2.3, NAVY)
    add_text_box(slide, 1.27, 0.46, 22.86, 0.71, "FinSight IA",
                 10, "8899BB", align=PP_ALIGN.CENTER)
    add_rect(slide, 11.05, 1.32, 3.3, 0.05, WHITE)
    add_text_box(slide, 1.27, 1.52, 22.86, 0.81,
                 "Pitchbook Comparatif  \u2014  Analyse d'Investissement",
                 11, "CCDDEE", align=PP_ALIGN.CENTER)

    # Titre "A vs B"
    title_str = f"{tkr_a}  vs  {tkr_b}"
    add_text_box(slide, 1.27, 4.06, 22.86, 1.8,
                 title_str, 38, NAVY, bold=True, align=PP_ALIGN.CENTER)

    # Noms complets (Name + Ticker)
    add_text_box(slide, 1.27, 5.9, 11.0, 0.71,
                 _truncate(f"{name_a} ({tkr_a})", 44), 11, "888888", align=PP_ALIGN.CENTER)
    add_text_box(slide, 13.13, 5.9, 11.0, 0.71,
                 _truncate(f"{name_b} ({tkr_b})", 44), 11, "888888", align=PP_ALIGN.CENTER)


    # Bandeau recommendation A
    ra_c, ra_pal = _rec_c(rec_a)
    price_a  = m_a.get('share_price')
    tbase_a  = m_a.get('dcf_base')
    up_a     = _upside(tbase_a, price_a)
    cur_a    = "EUR" if (m_a.get('currency_a') or 'USD') == 'EUR' else '$'
    add_rect(slide, 1.27, 7.0, 11.0, 1.1, ra_pal)
    add_rect(slide, 1.27, 7.0, 0.18, 1.1, ra_c)
    add_text_box(slide, 1.6, 7.05, 10.5, 1.0,
                 f"\u25cf {rec_a}  \u00b7  Cible : {_fr(tbase_a, 0)} {cur_a}  \u00b7  Upside : {up_a}",
                 9, NAVY, bold=True, align=PP_ALIGN.CENTER)

    # Bandeau recommendation B
    rb_c, rb_pal = _rec_c(rec_b)
    price_b  = m_b.get('share_price')
    tbase_b  = m_b.get('dcf_base')
    up_b     = _upside(tbase_b, price_b)
    cur_b    = "EUR" if (m_b.get('currency_b') or 'USD') == 'EUR' else '$'
    add_rect(slide, 13.13, 7.0, 11.0, 1.1, rb_pal)
    add_rect(slide, 13.13, 7.0, 0.18, 1.1, rb_c)
    add_text_box(slide, 13.46, 7.05, 10.5, 1.0,
                 f"\u25cf {rec_b}  \u00b7  Cible : {_fr(tbase_b, 0)} {cur_b}  \u00b7  Upside : {up_b}",
                 9, NAVY, bold=True, align=PP_ALIGN.CENTER)

    # Secteur A / B
    sec_a = m_a.get('sector_a') or ''
    sec_b = m_b.get('sector_b') or ''
    add_text_box(slide, 1.27, 8.3, 11.0, 0.56,
                 sec_a, 9, GREY_TXT, align=PP_ALIGN.CENTER)
    add_text_box(slide, 13.13, 8.3, 11.0, 0.56,
                 sec_b, 9, GREY_TXT, align=PP_ALIGN.CENTER)

    # Bottom rule + date
    add_rect(slide, 1.02, 12.4, 23.37, 0.03, "AAAAAA")
    add_text_box(slide, 1.02, 12.65, 11.43, 0.56, "Rapport confidentiel", 8, GREY_TXT)
    add_text_box(slide, 12.95, 12.65, 11.43, 0.56, gen_date, 8, GREY_TXT,
                 align=PP_ALIGN.RIGHT)
    return slide


def _slide_exec_summary(prs, m_a: dict, m_b: dict, synthesis: dict):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    tkr_a = m_a.get('ticker_a') or 'A'
    tkr_b = m_b.get('ticker_b') or 'B'

    navy_bar(slide)
    footer_bar(slide)
    slide_title(slide, "Executive Summary Comparatif")
    section_dots(slide, 1)
    _company_header_band(slide, tkr_a, tkr_b)

    # Bandeau exec summary LLM
    exec_txt = synthesis.get('exec_summary') or ""
    if exec_txt:
        _exec_clean = " ".join(_fit(exec_txt, 360).split())
        add_rect(slide, 1.02, 2.45, 23.37, 1.8, NAVY_PALE)
        add_rect(slide, 1.02, 2.45, 0.13, 1.8, NAVY_MID)
        add_text_box(slide, 1.4, 2.55, 22.8, 1.6, _exec_clean,
                     8.5, NAVY, wrap=True)

    # 4 KPI boxes A vs B — prix, MC, PE, EV/EBITDA
    price_a = m_a.get('share_price')
    price_b = m_b.get('share_price')
    mc_a    = m_a.get('market_cap')
    mc_b    = m_b.get('market_cap')
    pe_a    = m_a.get('pe_ratio')
    pe_b    = m_b.get('pe_ratio')
    ev_a    = m_a.get('ev_ebitda')
    ev_b    = m_b.get('ev_ebitda')

    y_kpi = 4.7
    h_kpi = 1.9
    w_kpi = 5.6
    xs = [1.02, 6.92, 12.81, 18.7]
    vals  = [
        (f"{_fr(price_a, 1)} / {_fr(price_b, 1)}", "Cours actuel A / B", ""),
        (f"{_fr(mc_a, 1)} / {_fr(mc_b, 1)}", "Market Cap (Mds) A / B", ""),
        (f"{_frx(pe_a)} / {_frx(pe_b)}", "PE A / B", ""),
        (f"{_frx(ev_a)} / {_frx(ev_b)}", "EV/EBITDA A / B", ""),
    ]
    for xi, (val, lbl, sub) in zip(xs, vals):
        kpi_box(slide, xi, y_kpi, w_kpi, h_kpi, val, lbl, sub)

    def _fr_nd(v) -> str:
        if v is None: return "\u2014"
        try:
            fv = round(float(v), 1)
            fv = 0.0 if abs(fv) < 0.05 else fv
            return f"{fv:.1f}".replace(".", ",") + "x"
        except Exception:
            return "\u2014"

    # Tableau recapitulatif 2 colonnes
    y_tbl = 6.85
    rows = [
        ("FinSight Score",  str(m_a.get('finsight_score') or '\u2014'), str(m_b.get('finsight_score') or '\u2014')),
        ("Recommandation",  m_a.get('recommendation') or '\u2014', m_b.get('recommendation') or '\u2014'),
        ("Conviction",      _frpct(m_a.get('conviction')), _frpct(m_b.get('conviction'))),
        ("Marge EBITDA LTM", _frpct_margin(m_a.get('ebitda_margin_ltm')), _frpct_margin(m_b.get('ebitda_margin_ltm'))),
        ("ROIC",            _frpct(m_a.get('roic')), _frpct(m_b.get('roic'))),
        ("Rev CAGR 3y",     _frpct(m_a.get('revenue_cagr_3y')), _frpct(m_b.get('revenue_cagr_3y'))),
        ("ND/EBITDA",       _fr_nd(m_a.get('net_debt_ebitda')), _fr_nd(m_b.get('net_debt_ebitda'))),
    ]
    add_table(
        slide, 1.02, y_tbl, 23.37, 5.0,
        num_rows=len(rows), num_cols=3,
        col_widths_pct=[0.42, 0.29, 0.29],
        header_data=["Metrique", tkr_a, tkr_b],
        rows_data=[(r[0], r[1], r[2]) for r in rows],
        header_fill=NAVY, border_hex="DDDDDD"
    )
    return slide


def _slide_sommaire(prs, tkr_a: str, tkr_b: str):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    navy_bar(slide)
    footer_bar(slide)
    slide_title(slide, "Sommaire")

    sections = [
        ("01", "Profil & Identite", "Presentation comparative des deux societes",          5),
        ("02", "Performance Financiere", "P&L, marges, bilan et liquidite compares",        7),
        ("03", "Valorisation", "Multiples, DCF, GBM Monte Carlo et football field",         11),
        ("04", "Qualite & Risques", "Piotroski, levier, beta, VaR et momentum",            15),
        ("05", "Verdict", "FinSight Score, theses d'investissement et verdict final",      18),
    ]

    fills = [WHITE, GREY_BG, WHITE, GREY_BG, WHITE]
    ys    = [2.08, 4.01, 5.94, 7.87, 9.80]

    for i, (num, title, desc, page_n) in enumerate(sections):
        y    = ys[i]
        fill = fills[i]
        add_rect(slide, 1.02, y, 23.37, 1.73, fill)
        add_rect(slide, 1.02, y, 0.13, 1.73, NAVY_MID)
        add_text_box(slide, 1.40, y + 0.15, 1.32, 1.43,
                     num, 17, GREY_LIGHT, bold=True)
        add_text_box(slide, 3.20, y + 0.20, 17.00, 0.71,
                     title, 13, NAVY, bold=True)
        add_text_box(slide, 3.20, y + 0.85, 17.00, 0.71,
                     desc, 8.5, GREY_TXT)
        add_text_box(slide, 20.57, y + 0.45, 3.82, 0.71,
                     str(page_n), 9, NAVY_MID,
                     align=PP_ALIGN.RIGHT)

    return slide


def _slide_profil(prs, m_a: dict, m_b: dict):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    tkr_a = m_a.get('ticker_a') or 'A'
    tkr_b = m_b.get('ticker_b') or 'B'

    navy_bar(slide)
    footer_bar(slide)
    slide_title(slide, "Profil Comparatif")
    section_dots(slide, 1)
    _company_header_band(slide, tkr_a, tkr_b)

    def _profile_col(slide, m, side, tkr):
        if side == 'left':
            x0 = 1.02; lbl_color = NAVY_MID
        else:
            x0 = 12.76; lbl_color = GREEN_MID
        w = 11.0

        # Identite
        name = m.get('company_name_a' if side == 'left' else 'company_name_b') or tkr
        sector = m.get('sector_a' if side == 'left' else 'sector_b') or '\u2014'
        currency = m.get('currency_a' if side == 'left' else 'currency_b') or 'USD'
        cur_sym  = 'EUR' if currency == 'EUR' else '$'

        # En-tete nom (Name + Ticker) — rect plus haut, texte tronque, pas de wrap
        add_rect(slide, x0, 2.45, w, 0.60, COLOR_A_PAL if side == 'left' else COLOR_B_PAL)
        add_rect(slide, x0, 2.45, 0.13, 0.60, lbl_color)
        add_text_box(slide, x0 + 0.22, 2.50, w - 0.3, 0.50,
                     _truncate(f"{name} ({tkr})", 38), 8, lbl_color, bold=True, wrap=False)

        rows_id = [
            ("Secteur",          _truncate(sector, 22)),
            ("Devise",           currency),
            ("Cours",            _fr(m.get('share_price'), 1) + " " + cur_sym),
            ("52W High",         _fr(m.get('week52_high'), 2) + " " + cur_sym),
            ("52W Low",          _fr(m.get('week52_low'), 2) + " " + cur_sym),
            ("Market Cap",       _fr(m.get('market_cap'), 1) + " Mds " + cur_sym),
            ("Valeur Entrep.",   _fr(m.get('enterprise_value'), 1) + " Mds " + cur_sym),
            ("Div. Yield",       _frpct(m.get('dividend_yield'))),
            ("Perf. 1M",         _frpct(m.get('perf_1m'), signed=True)),
            ("Perf. 3M",         _frpct(m.get('perf_3m'), signed=True)),
            ("Perf. 1Y",         _frpct(m.get('perf_1y'), signed=True)),
            ("Proch. Resultat",  m.get('next_earnings_date') or '\u2014'),
        ]

        y_start = 3.18
        row_h   = 0.68
        for ri, (lbl, val) in enumerate(rows_id):
            y = y_start + ri * row_h
            fill = WHITE if ri % 2 == 0 else GREY_BG
            add_rect(slide, x0, y, w, row_h - 0.04, fill)
            add_text_box(slide, x0 + 0.1, y + 0.12, w * 0.52, row_h - 0.15, lbl,
                         7.5, GREY_TXT, italic=True, wrap=False)
            add_text_box(slide, x0 + w * 0.52, y + 0.12, w * 0.45, row_h - 0.15,
                         val, 7.5, BLACK, bold=False, align=PP_ALIGN.RIGHT, wrap=False)

    _profile_col(slide, m_a, 'left', tkr_a)
    _profile_col(slide, m_b, 'right', tkr_b)
    return slide


def _slide_pl(prs, m_a: dict, m_b: dict, synthesis: dict):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    tkr_a = m_a.get('ticker_a') or 'A'
    tkr_b = m_b.get('ticker_b') or 'B'
    cur_a = 'EUR' if (m_a.get('currency_a') or 'USD') == 'EUR' else '$'
    cur_b = 'EUR' if (m_b.get('currency_b') or 'USD') == 'EUR' else '$'

    navy_bar(slide)
    footer_bar(slide)
    slide_title(slide, "Compte de Resultat Compare")
    section_dots(slide, 2)
    _company_header_band(slide, tkr_a, tkr_b)

    # Tableau comparatif P&L
    rows = [
        ("Marge EBITDA LTM",  _frpct_margin(m_a.get('ebitda_margin_ltm')),  _frpct_margin(m_b.get('ebitda_margin_ltm'))),
        ("Marge EBITDA N-1",  _frpct_margin(m_a.get('ebitda_margin_y1')),   _frpct_margin(m_b.get('ebitda_margin_y1'))),
        ("Marge EBITDA N-2",  _frpct_margin(m_a.get('ebitda_margin_y2')),   _frpct_margin(m_b.get('ebitda_margin_y2'))),
        ("Tendance EBITDA",   m_a.get('ebitda_margin_trend') or '\u2014',    m_b.get('ebitda_margin_trend') or '\u2014'),
        ("ROIC",              _frpct(m_a.get('roic')),                       _frpct(m_b.get('roic'))),
        ("ROE",               _frpct(m_a.get('roe')),                        _frpct(m_b.get('roe'))),
        ("Rev CAGR 3y",       _frpct(m_a.get('revenue_cagr_3y')),            _frpct(m_b.get('revenue_cagr_3y'))),
        ("Croissance EPS",    _frpct(m_a.get('eps_growth'), signed=True),    _frpct(m_b.get('eps_growth'), signed=True)),
    ]
    add_table(
        slide, 1.02, 2.45, 23.37, 5.6,
        num_rows=len(rows), num_cols=3,
        col_widths_pct=[0.44, 0.28, 0.28],
        header_data=["Indicateur", tkr_a, tkr_b],
        rows_data=rows,
        header_fill=NAVY, border_hex="DDDDDD"
    )

    # Commentaire analytique sous le tableau (pas de graphiques — réservés slide "Rentabilité")
    fin_txt = synthesis.get('financial_text') or ""
    if fin_txt:
        _fc = " ".join(_fit(fin_txt, 320).split())
        add_rect(slide, 1.02, 8.35, 23.37, 1.6, NAVY_PALE)
        add_rect(slide, 1.02, 8.35, 0.13, 1.6, NAVY_MID)
        add_text_box(slide, 1.4, 8.45, 22.8, 1.5, _fc, 8.5, NAVY, wrap=True)

    return slide


def _slide_marges(prs, m_a: dict, m_b: dict, synthesis: dict):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    tkr_a = m_a.get('ticker_a') or 'A'
    tkr_b = m_b.get('ticker_b') or 'B'

    navy_bar(slide)
    footer_bar(slide)
    slide_title(slide, "Rentabilite & Croissance Comparees")
    section_dots(slide, 2)
    _company_header_band(slide, tkr_a, tkr_b)

    # Commentaire financier LLM
    txt = synthesis.get('financial_text') or ""
    if txt:
        _txt_clean = " ".join(_fit(txt, 280).split())
        add_rect(slide, 1.02, 2.45, 23.37, 1.8, NAVY_PALE)
        add_rect(slide, 1.02, 2.45, 0.13, 1.8, NAVY_MID)
        add_text_box(slide, 1.4, 2.55, 22.8, 1.7, _txt_clean, 8.5, NAVY, wrap=True)

    # KPIs side-by-side
    y_kpi = 4.2
    kw = 5.5; kh = 1.8

    def _kpi_pair(x, v_a, v_b, lbl, sub=""):
        kpi_box(slide, x,      y_kpi, kw, kh, v_a, f"{lbl} — {tkr_a}", sub, COLOR_A_PAL, COLOR_A)
        kpi_box(slide, x + kw + 0.08, y_kpi, kw, kh, v_b, f"{lbl} — {tkr_b}", sub, COLOR_B_PAL, COLOR_B)

    # Row 1
    kpi_box(slide, 1.02,   y_kpi, 5.5, 1.8, _frpct_margin(m_a.get('ebitda_margin_ltm')), f"Marge EBITDA LTM — {tkr_a}", "", COLOR_A_PAL, COLOR_A)
    kpi_box(slide, 7.0,    y_kpi, 5.5, 1.8, _frpct_margin(m_b.get('ebitda_margin_ltm')), f"Marge EBITDA LTM — {tkr_b}", "", COLOR_B_PAL, COLOR_B)
    kpi_box(slide, 12.98,  y_kpi, 5.5, 1.8, _frpct(m_a.get('roic')), f"ROIC — {tkr_a}", "", COLOR_A_PAL, COLOR_A)
    kpi_box(slide, 18.95,  y_kpi, 5.5, 1.8, _frpct(m_b.get('roic')), f"ROIC — {tkr_b}", "", COLOR_B_PAL, COLOR_B)

    # Row 2
    y2 = y_kpi + 2.05
    kpi_box(slide, 1.02,   y2, 5.5, 1.8, _frpct(m_a.get('revenue_cagr_3y')), f"Rev CAGR 3y — {tkr_a}", "", COLOR_A_PAL, COLOR_A)
    kpi_box(slide, 7.0,    y2, 5.5, 1.8, _frpct(m_b.get('revenue_cagr_3y')), f"Rev CAGR 3y — {tkr_b}", "", COLOR_B_PAL, COLOR_B)
    kpi_box(slide, 12.98,  y2, 5.5, 1.8, _frpct(m_a.get('roe')), f"ROE — {tkr_a}", "", COLOR_A_PAL, COLOR_A)
    kpi_box(slide, 18.95,  y2, 5.5, 1.8, _frpct(m_b.get('roe')), f"ROE — {tkr_b}", "", COLOR_B_PAL, COLOR_B)

    # Graphique marges + croissance
    buf_l = _chart_ebitda_margins(m_a, m_b, tkr_a, tkr_b)
    buf_r = _chart_growth_returns(m_a, m_b, tkr_a, tkr_b)
    y_ch = y2 + 2.1
    if buf_l:
        _insert_chart(slide, buf_l, 1.02, y_ch, 11.0, 5.5)
    if buf_r:
        _insert_chart(slide, buf_r, 13.3, y_ch, 11.07, 5.5)

    return slide


def _slide_bilan(prs, m_a: dict, m_b: dict):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    tkr_a = m_a.get('ticker_a') or 'A'
    tkr_b = m_b.get('ticker_b') or 'B'

    navy_bar(slide)
    footer_bar(slide)
    slide_title(slide, "Bilan & Liquidite Compares")
    section_dots(slide, 2)
    _company_header_band(slide, tkr_a, tkr_b)

    rows = [
        ("Net Debt / EBITDA", _fr(m_a.get('net_debt_ebitda'), 1) + "x", _fr(m_b.get('net_debt_ebitda'), 1) + "x"),
        ("Interest Coverage",  _frx(m_a.get('interest_coverage')),       _frx(m_b.get('interest_coverage'))),
        ("Current Ratio",      _frx(m_a.get('current_ratio')),           _frx(m_b.get('current_ratio'))),
        ("Quick Ratio",        _frx(m_a.get('quick_ratio')),             _frx(m_b.get('quick_ratio'))),
        ("Capex / Revenue",   _frpct(m_a.get('capex_to_revenue')),       _frpct(m_b.get('capex_to_revenue'))),
        ("Cash Conversion",   _frx(m_a.get('cash_conversion')),          _frx(m_b.get('cash_conversion'))),
        ("FCF Yield",         _frpct(m_a.get('fcf_yield')),              _frpct(m_b.get('fcf_yield'))),
        ("Altman Z-Score",    _fr(m_a.get('altman_z'), 2),               _fr(m_b.get('altman_z'), 2)),
    ]
    add_table(
        slide, 1.02, 2.45, 23.37, 5.8,
        num_rows=len(rows), num_cols=3,
        col_widths_pct=[0.44, 0.28, 0.28],
        header_data=["Indicateur Bilan", tkr_a, tkr_b],
        rows_data=rows,
        header_fill=NAVY, border_hex="DDDDDD"
    )

    # KPIs bilan
    y_kpi = 8.55
    kpi_box(slide, 1.02,  y_kpi, 5.5, 1.8, _fr(m_a.get('net_debt_ebitda'), 1) + "x", f"Net Debt/EBITDA — {tkr_a}", "", COLOR_A_PAL, COLOR_A)
    kpi_box(slide, 7.0,   y_kpi, 5.5, 1.8, _fr(m_b.get('net_debt_ebitda'), 1) + "x", f"Net Debt/EBITDA — {tkr_b}", "", COLOR_B_PAL, COLOR_B)
    kpi_box(slide, 12.98, y_kpi, 5.5, 1.8, _frpct(m_a.get('fcf_yield')),             f"FCF Yield — {tkr_a}", "", COLOR_A_PAL, COLOR_A)
    kpi_box(slide, 18.95, y_kpi, 5.5, 1.8, _frpct(m_b.get('fcf_yield')),             f"FCF Yield — {tkr_b}", "", COLOR_B_PAL, COLOR_B)

    return slide


def _slide_multiples(prs, m_a: dict, m_b: dict, synthesis: dict):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    tkr_a = m_a.get('ticker_a') or 'A'
    tkr_b = m_b.get('ticker_b') or 'B'

    navy_bar(slide)
    footer_bar(slide)
    slide_title(slide, "Multiples de Valorisation Compares")
    section_dots(slide, 3)
    _company_header_band(slide, tkr_a, tkr_b)

    # Commentaire valorisation
    txt = synthesis.get('valuation_text') or ""
    if txt:
        add_rect(slide, 1.02, 2.45, 23.37, 1.6, NAVY_PALE)
        add_rect(slide, 1.02, 2.45, 0.13, 1.6, NAVY_MID)
        add_text_box(slide, 1.4, 2.55, 22.8, 1.45, _fit(txt, 260), 8.5, NAVY, wrap=True)

    sec_pe   = m_a.get('sector_median_pe') or '\u2014'
    sec_eveb = m_a.get('sector_median_ev_ebitda') or '\u2014'

    rows = [
        ("PE Ratio",          _frx(m_a.get('pe_ratio')),      _frx(m_b.get('pe_ratio')),      _frx(sec_pe) if sec_pe != '\u2014' else '\u2014'),
        ("EV / EBITDA",       _frx(m_a.get('ev_ebitda')),     _frx(m_b.get('ev_ebitda')),     _frx(sec_eveb) if sec_eveb != '\u2014' else '\u2014'),
        ("Price / Book",      _frx(m_a.get('price_to_book')), _frx(m_b.get('price_to_book')), "\u2014"),
        ("PEG Ratio",         _frx(m_a.get('peg_ratio')),     _frx(m_b.get('peg_ratio')),     "\u2014"),
        ("FCF Yield",         _frpct(m_a.get('fcf_yield')),   _frpct(m_b.get('fcf_yield')),   "\u2014"),
        ("Div. Yield",        _frpct(m_a.get('dividend_yield')), _frpct(m_b.get('dividend_yield')), "\u2014"),
    ]
    add_table(
        slide, 1.02, 4.0, 23.37, 4.5,
        num_rows=len(rows), num_cols=4,
        col_widths_pct=[0.36, 0.21, 0.21, 0.22],
        header_data=["Multiple", tkr_a, tkr_b, "Med. Secteur"],
        rows_data=rows,
        header_fill=NAVY, border_hex="DDDDDD"
    )

    # Graphique multiples
    buf = _chart_multiples(m_a, m_b, tkr_a, tkr_b)
    if buf:
        _insert_chart(slide, buf, 1.02, 8.75, 23.37, 4.0)

    return slide


def _slide_dcf(prs, m_a: dict, m_b: dict):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    tkr_a = m_a.get('ticker_a') or 'A'
    tkr_b = m_b.get('ticker_b') or 'B'

    navy_bar(slide)
    footer_bar(slide)
    slide_title(slide, "DCF & Prix Cibles Compares")
    section_dots(slide, 3)
    _company_header_band(slide, tkr_a, tkr_b)

    # DCF A
    cur_a = 'EUR' if (m_a.get('currency_a') or 'USD') == 'EUR' else '$'
    cur_b = 'EUR' if (m_b.get('currency_b') or 'USD') == 'EUR' else '$'

    y_dcf = 2.45
    # Panel A
    add_rect(slide, 1.02, y_dcf, 11.44, 6.8, COLOR_A_PAL)
    add_rect(slide, 1.02, y_dcf, 0.13, 6.8, COLOR_A)
    add_text_box(slide, 1.3, y_dcf + 0.1, 11.0, 0.56, f"DCF — {tkr_a}", 9.5, NAVY_MID, bold=True)

    rows_dcf_a = [
        ("Cours actuel",   _fr(m_a.get('share_price'), 1) + " " + cur_a),
        ("Cible Bear",     _fr(m_a.get('dcf_bear'), 0) + " " + cur_a),
        ("Cible Base",     _fr(m_a.get('dcf_base'), 0) + " " + cur_a),
        ("Cible Bull",     _fr(m_a.get('dcf_bull'), 0) + " " + cur_a),
        ("Upside Base",    _upside(m_a.get('dcf_base'), m_a.get('share_price'))),
        ("Marge Securite", _frpct(m_a.get('margin_of_safety'))),
        ("WACC",           _frpct(m_a.get('wacc'))),
        ("TGR",            _frpct(m_a.get('terminal_growth'))),
        ("MC P10",         _fr(m_a.get('monte_carlo_p10'), 0) + " " + cur_a),
        ("MC P50",         _fr(m_a.get('monte_carlo_p50'), 0) + " " + cur_a),
        ("MC P90",         _fr(m_a.get('monte_carlo_p90'), 0) + " " + cur_a),
    ]
    for ri, (lbl, val) in enumerate(rows_dcf_a):
        y = y_dcf + 0.75 + ri * 0.53
        fill = WHITE if ri % 2 == 0 else COLOR_A_PAL
        add_rect(slide, 1.22, y, 11.04, 0.5, fill)
        add_text_box(slide, 1.35, y + 0.06, 6.0, 0.4, lbl, 8, GREY_TXT, italic=True)
        add_text_box(slide, 7.5, y + 0.06, 4.6, 0.4, val, 8.5, BLACK, bold=False,
                     align=PP_ALIGN.RIGHT)

    # Panel B
    add_rect(slide, 12.94, y_dcf, 11.44, 6.8, COLOR_B_PAL)
    add_rect(slide, 12.94, y_dcf, 0.13, 6.8, COLOR_B)
    add_text_box(slide, 13.22, y_dcf + 0.1, 11.0, 0.56, f"DCF — {tkr_b}", 9.5, GREEN, bold=True)

    rows_dcf_b = [
        ("Cours actuel",   _fr(m_b.get('share_price'), 1) + " " + cur_b),
        ("Cible Bear",     _fr(m_b.get('dcf_bear'), 0) + " " + cur_b),
        ("Cible Base",     _fr(m_b.get('dcf_base'), 0) + " " + cur_b),
        ("Cible Bull",     _fr(m_b.get('dcf_bull'), 0) + " " + cur_b),
        ("Upside Base",    _upside(m_b.get('dcf_base'), m_b.get('share_price'))),
        ("Marge Securite", _frpct(m_b.get('margin_of_safety'))),
        ("WACC",           _frpct(m_b.get('wacc'))),
        ("TGR",            _frpct(m_b.get('terminal_growth'))),
        ("MC P10",         _fr(m_b.get('monte_carlo_p10'), 0) + " " + cur_b),
        ("MC P50",         _fr(m_b.get('monte_carlo_p50'), 0) + " " + cur_b),
        ("MC P90",         _fr(m_b.get('monte_carlo_p90'), 0) + " " + cur_b),
    ]
    for ri, (lbl, val) in enumerate(rows_dcf_b):
        y = y_dcf + 0.75 + ri * 0.53
        fill = WHITE if ri % 2 == 0 else COLOR_B_PAL
        add_rect(slide, 13.14, y, 11.04, 0.5, fill)
        add_text_box(slide, 13.27, y + 0.06, 6.0, 0.4, lbl, 8, GREY_TXT, italic=True)
        add_text_box(slide, 19.4, y + 0.06, 4.4, 0.4, val, 8.5, BLACK, bold=False,
                     align=PP_ALIGN.RIGHT)

    # Verdict entree
    ez_a = m_a.get('entry_zone_ok') or 0
    ez_b = m_b.get('entry_zone_ok') or 0
    y_ev = y_dcf + 7.1
    add_rect(slide, 1.02, y_ev, 11.44, 0.7, GREEN_PALE if ez_a else RED_PALE)
    add_text_box(slide, 1.15, y_ev + 0.1, 11.2, 0.5,
                 ("Zone d'achat : Cours < DCF base" if ez_a else "Hors zone : Cours > DCF base"),
                 8.5, GREEN if ez_a else RED, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 12.94, y_ev, 11.44, 0.7, GREEN_PALE if ez_b else RED_PALE)
    add_text_box(slide, 13.07, y_ev + 0.1, 11.2, 0.5,
                 ("Zone d'achat : Cours < DCF base" if ez_b else "Hors zone : Cours > DCF base"),
                 8.5, GREEN if ez_b else RED, bold=True, align=PP_ALIGN.CENTER)

    return slide


def _slide_monte_carlo(prs, m_a: dict, m_b: dict):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    tkr_a = m_a.get('ticker_a') or 'A'
    tkr_b = m_b.get('ticker_b') or 'B'

    navy_bar(slide)
    footer_bar(slide)
    slide_title(slide, "GBM Monte Carlo & Upside Relatif")
    section_dots(slide, 3)
    _company_header_band(slide, tkr_a, tkr_b)

    # Tableau Monte Carlo
    cur_a = 'EUR' if (m_a.get('currency_a') or 'USD') == 'EUR' else '$'
    cur_b = 'EUR' if (m_b.get('currency_b') or 'USD') == 'EUR' else '$'
    p_a = m_a.get('share_price')
    p_b = m_b.get('share_price')

    def _sigma_fmt(v):
        if v is None: return "\u2014"
        try: return f"{float(v) * 100:.1f}".replace(".", ",") + " %"
        except: return "\u2014"

    rows_mc = [
        ("P10 (scen. pessimiste)",
         _fr(m_a.get('monte_carlo_p10'), 0) + " " + cur_a + "  (" + _upside(m_a.get('monte_carlo_p10'), p_a) + ")",
         _fr(m_b.get('monte_carlo_p10'), 0) + " " + cur_b + "  (" + _upside(m_b.get('monte_carlo_p10'), p_b) + ")"),
        ("P50 (median GBM)",
         _fr(m_a.get('monte_carlo_p50'), 0) + " " + cur_a + "  (" + _upside(m_a.get('monte_carlo_p50'), p_a) + ")",
         _fr(m_b.get('monte_carlo_p50'), 0) + " " + cur_b + "  (" + _upside(m_b.get('monte_carlo_p50'), p_b) + ")"),
        ("P90 (scen. optimiste)",
         _fr(m_a.get('monte_carlo_p90'), 0) + " " + cur_a + "  (" + _upside(m_a.get('monte_carlo_p90'), p_a) + ")",
         _fr(m_b.get('monte_carlo_p90'), 0) + " " + cur_b + "  (" + _upside(m_b.get('monte_carlo_p90'), p_b) + ")"),
        ("Sigma (vol. ann.)", _sigma_fmt(m_a.get('gbm_sigma_annual')), _sigma_fmt(m_b.get('gbm_sigma_annual'))),
        ("Cours actuel S0",  _fr(p_a, 2) + " " + cur_a if p_a else "\u2014",
                             _fr(p_b, 2) + " " + cur_b if p_b else "\u2014"),
    ]
    add_table(
        slide, 1.02, 2.45, 23.37, 3.8,
        num_rows=len(rows_mc), num_cols=3,
        col_widths_pct=[0.36, 0.32, 0.32],
        header_data=["GBM Monte Carlo", tkr_a, tkr_b],
        rows_data=rows_mc,
        header_fill=NAVY, border_hex="DDDDDD"
    )

    # Football field chart — barres horizontales simples
    try:
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        import numpy as np

        fig, ax = plt.subplots(figsize=(10.0, 5.5))
        fig.patch.set_facecolor('white')
        ax.set_facecolor('#FAFAFA')

        methods_a = ['Bear', 'Base', 'Bull']
        vals_a = [
            _safe_float(m_a.get('dcf_bear')) or 0,
            _safe_float(m_a.get('dcf_base')) or 0,
            _safe_float(m_a.get('dcf_bull')) or 0,
        ]
        vals_b = [
            _safe_float(m_b.get('dcf_bear')) or 0,
            _safe_float(m_b.get('dcf_base')) or 0,
            _safe_float(m_b.get('dcf_bull')) or 0,
        ]

        y_pos = np.array([3, 2, 1, -1, -2, -3])
        labels = [f'{tkr_a} Bear', f'{tkr_a} Base', f'{tkr_a} Bull',
                  f'{tkr_b} Bear', f'{tkr_b} Base', f'{tkr_b} Bull']
        all_vals = vals_a + vals_b
        colors_ff = ['#6688BB', '#2E5FA3', '#102040',
                     '#66BB88', '#2E8B57', '#104030']

        for i, (yp, val, lbl, col) in enumerate(zip(y_pos, all_vals, labels, colors_ff)):
            if val:
                ax.barh(yp, val, 0.6, color=col, alpha=0.85)
                ax.text(val + max(all_vals or [1]) * 0.01, yp,
                        f'{val:.0f}', ha='left', va='center', fontsize=8, color='#333')

        # Cours actuels — labels en haut pour éviter débordement
        if p_a:
            ax.axvline(p_a, color='#2E5FA3', linestyle='--', linewidth=1.2, alpha=0.6)
            ax.text(p_a, 3.55, f' {tkr_a}', fontsize=7, color='#2E5FA3', va='bottom', ha='left', clip_on=True)
        if p_b:
            ax.axvline(p_b, color='#2E8B57', linestyle='--', linewidth=1.2, alpha=0.6)
            ax.text(p_b, 3.15, f' {tkr_b}', fontsize=7, color='#2E8B57', va='bottom', ha='left', clip_on=True)

        ax.set_yticks(y_pos)
        ax.set_yticklabels(labels, fontsize=8)
        ax.set_xlabel('Prix cible', fontsize=8)
        ax.set_title('Football Field — DCF Bear/Base/Bull', fontsize=9,
                     fontweight='bold', color='#1B3A6B')
        ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
        ax.tick_params(labelsize=8)
        fig.tight_layout(pad=1.0)
        buf = _make_chart_buf(fig)
        _insert_chart(slide, buf, 1.02, 6.5, 23.37, 6.2)
    except Exception as e:
        log.warning(f"[cmp_pptx] football field error: {e}")

    return slide


def _slide_piotroski(prs, m_a: dict, m_b: dict, synthesis: dict):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    tkr_a = m_a.get('ticker_a') or 'A'
    tkr_b = m_b.get('ticker_b') or 'B'

    navy_bar(slide)
    footer_bar(slide)
    slide_title(slide, "Qualite Financiere & Solidite Bilancielle")
    section_dots(slide, 4)
    _company_header_band(slide, tkr_a, tkr_b)

    # Commentaire qualite LLM
    txt = synthesis.get('quality_text') or ""
    if txt:
        add_rect(slide, 1.02, 2.45, 23.37, 1.8, NAVY_PALE)
        add_rect(slide, 1.02, 2.45, 0.13, 1.8, NAVY_MID)
        add_text_box(slide, 1.4, 2.55, 22.8, 1.6, _fit(txt, 380), 8.5, NAVY, wrap=True)

    def _pio_val(m, key):
        v = m.get(key)
        if v is None: return "\u2014"
        return "1" if int(v) == 1 else "0"

    pio_keys = [
        ("pio_roa_positive",         "ROA positif"),
        ("pio_cfo_positive",         "Cash-flows operations positifs"),
        ("pio_delta_roa",            "Amelioration ROA"),
        ("pio_accruals",             "Qualite earnings (accruals)"),
        ("pio_delta_leverage",       "Reduction levier"),
        ("pio_delta_liquidity",      "Amelioration liquidite"),
        ("pio_no_dilution",          "Pas de dilution"),
        ("pio_delta_gross_margin",   "Amelioration marge brute"),
        ("pio_delta_asset_turnover", "Amelioration rotation actifs"),
    ]
    rows_pio = [(lbl, _pio_val(m_a, key), _pio_val(m_b, key)) for key, lbl in pio_keys]
    # Score total
    sc_a = m_a.get('piotroski_score')
    sc_b = m_b.get('piotroski_score')
    rows_pio.append(("PIOTROSKI F-SCORE TOTAL",
                     str(sc_a) if sc_a is not None else "\u2014",
                     str(sc_b) if sc_b is not None else "\u2014"))

    add_table(
        slide, 1.02, 4.15, 14.0, 8.0,
        num_rows=len(rows_pio), num_cols=3,
        col_widths_pct=[0.60, 0.20, 0.20],
        header_data=["Critere Piotroski", tkr_a, tkr_b],
        rows_data=rows_pio,
        header_fill=NAVY, border_hex="DDDDDD"
    )

    # Tableaux complementaires
    add_rect(slide, 15.5, 4.15, 9.5, 0.6, NAVY)
    add_text_box(slide, 15.65, 4.20, 9.2, 0.5, "Scores Complementaires", 8, WHITE, bold=True)

    rows_sc = [
        ("Beneish M-Score", _fr(m_a.get('beneish_mscore'), 2), _fr(m_b.get('beneish_mscore'), 2)),
        ("Altman Z-Score",  _fr(m_a.get('altman_z'), 2),       _fr(m_b.get('altman_z'), 2)),
        ("Sloan Accruals",  _fr(m_a.get('sloan_accruals'), 3), _fr(m_b.get('sloan_accruals'), 3)),
        ("Cash Conversion", _frx(m_a.get('cash_conversion')), _frx(m_b.get('cash_conversion'))),
    ]
    add_table(
        slide, 15.5, 4.90, 9.5, 3.0,
        num_rows=len(rows_sc), num_cols=3,
        col_widths_pct=[0.50, 0.25, 0.25],
        header_data=["Indicateur", tkr_a, tkr_b],
        rows_data=rows_sc,
        header_fill=NAVY_MID, border_hex="DDDDDD"
    )

    # Legende Beneish + Altman
    add_rect(slide, 15.5, 8.05, 9.5, 4.15, GREY_BG)
    legend_lines = [
        "Beneish M-Score : < -1,78 = faible risque",
        "    manipulation comptable ; > -1,78 = risque",
        "",
        "Altman Z-Score : > 2,99 = sain ;",
        "    1,81-2,99 = zone grise ;",
        "    < 1,81 = risque de detresse",
        "",
        "Sloan Accruals : proche de 0 = qualite",
        "    earnings elevee (cash-based)",
    ]
    y_leg = 8.20
    for line in legend_lines:
        add_text_box(slide, 15.7, y_leg, 9.1, 0.38, line, 7.5, GREY_TXT)
        y_leg += 0.38

    return slide


def _slide_risque(prs, m_a: dict, m_b: dict):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    tkr_a = m_a.get('ticker_a') or 'A'
    tkr_b = m_b.get('ticker_b') or 'B'

    navy_bar(slide)
    footer_bar(slide)
    slide_title(slide, "Profil de Risque & Momentum Compare")
    section_dots(slide, 4)
    _company_header_band(slide, tkr_a, tkr_b)

    # KPIs risque — AAPL colonne gauche, MSFT colonne droite
    y_kpi = 2.45
    kpi_box(slide, 1.02, y_kpi, 5.5, 1.9, _fr(m_a.get('beta'), 2),                   f"Beta — {tkr_a}", "vs marche", COLOR_A_PAL, COLOR_A)
    kpi_box(slide, 7.0,  y_kpi, 5.1, 1.9, _frpct(m_a.get('var_95_1m'), signed=True), f"VaR 95% 1M — {tkr_a}", "",  COLOR_A_PAL, COLOR_A)
    kpi_box(slide, 13.0, y_kpi, 5.5, 1.9, _fr(m_b.get('beta'), 2),                   f"Beta — {tkr_b}", "vs marche", COLOR_B_PAL, COLOR_B)
    kpi_box(slide, 18.95,y_kpi, 5.4, 1.9, _frpct(m_b.get('var_95_1m'), signed=True), f"VaR 95% 1M — {tkr_b}", "",  COLOR_B_PAL, COLOR_B)

    # Tableau performances
    rows_perf = [
        ("Perf. 1 Mois",  _frpct(m_a.get('perf_1m'), signed=True),  _frpct(m_b.get('perf_1m'), signed=True)),
        ("Perf. 3 Mois",  _frpct(m_a.get('perf_3m'), signed=True),  _frpct(m_b.get('perf_3m'), signed=True)),
        ("Perf. 1 An",    _frpct(m_a.get('perf_1y'), signed=True),  _frpct(m_b.get('perf_1y'), signed=True)),
        ("Score Momentum", _fr(m_a.get('momentum_score'), 1) + " /10", _fr(m_b.get('momentum_score'), 1) + " /10"),
        ("Beta",           _fr(m_a.get('beta'), 2), _fr(m_b.get('beta'), 2)),
        ("VaR 95% 1M",    _frpct(m_a.get('var_95_1m')), _frpct(m_b.get('var_95_1m'))),
    ]
    add_table(
        slide, 1.02, 4.65, 12.5, 4.5,
        num_rows=len(rows_perf), num_cols=3,
        col_widths_pct=[0.44, 0.28, 0.28],
        header_data=["Performance / Risque", tkr_a, tkr_b],
        rows_data=rows_perf,
        header_fill=NAVY, border_hex="DDDDDD"
    )

    # Radar chart — positionne dans la colonne B, apres la ligne KPI
    buf = _chart_risk_profile(m_a, m_b, tkr_a, tkr_b)
    if buf:
        _insert_chart(slide, buf, 13.8, 4.5, 10.5, 8.0)

    # 52W range
    y_rng = 9.35
    add_rect(slide, 1.02, y_rng, 12.5, 0.5, NAVY)
    add_text_box(slide, 1.15, y_rng + 0.08, 12.2, 0.35, "Fourchette 52 Semaines", 8, WHITE, bold=True)

    rows_52 = [
        (f"{tkr_a} — 52W Low", _fr(m_a.get('week52_low'), 2), f"{tkr_a} — 52W High", _fr(m_a.get('week52_high'), 2)),
        (f"{tkr_b} — 52W Low", _fr(m_b.get('week52_low'), 2), f"{tkr_b} — 52W High", _fr(m_b.get('week52_high'), 2)),
    ]
    y_52 = y_rng + 0.6
    for ri, (l1, v1, l2, v2) in enumerate(rows_52):
        fill = WHITE if ri % 2 == 0 else GREY_BG
        add_rect(slide, 1.02, y_52 + ri * 0.58, 12.5, 0.55, fill)
        add_text_box(slide, 1.15, y_52 + ri * 0.58 + 0.06, 3.6, 0.45, l1, 7.5, GREY_TXT, wrap=False)
        add_text_box(slide, 4.8,  y_52 + ri * 0.58 + 0.06, 2.4, 0.45, v1, 8, BLACK, bold=True, wrap=False)
        add_text_box(slide, 7.3,  y_52 + ri * 0.58 + 0.06, 3.6, 0.45, l2, 7.5, GREY_TXT, wrap=False)
        add_text_box(slide, 11.0, y_52 + ri * 0.58 + 0.06, 2.3, 0.45, v2, 8, BLACK, bold=True, wrap=False)

    return slide


def _slide_finsight_score(prs, m_a: dict, m_b: dict):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    tkr_a = m_a.get('ticker_a') or 'A'
    tkr_b = m_b.get('ticker_b') or 'B'

    navy_bar(slide)
    footer_bar(slide)
    slide_title(slide, "FinSight Score Compare")
    section_dots(slide, 5)
    _company_header_band(slide, tkr_a, tkr_b)

    fs_a = m_a.get('finsight_score') or 0
    fs_b = m_b.get('finsight_score') or 0

    # Graphique score — pleine largeur
    buf = _chart_finsight_score(m_a, m_b, tkr_a, tkr_b)
    if buf:
        _insert_chart(slide, buf, 1.02, 2.45, 23.37, 4.5)

    # Decomposition score
    y_dec = 7.15
    add_rect(slide, 1.02, y_dec, 23.37, 0.6, NAVY)
    add_text_box(slide, 1.15, y_dec + 0.1, 23.0, 0.45, "Decomposition du Score (4 axes)", 8.5, WHITE, bold=True)

    # 4 axes de score — Value, Growth, Quality, Momentum
    def _score_axis(m, axis) -> str:
        # Proxies simples pour visualisation
        if axis == 'value':
            pe = _safe_float(m.get('pe_ratio')) or 25
            return str(max(0, min(25, round(25 - pe * 0.4))))
        if axis == 'growth':
            cagr = (_safe_float(m.get('revenue_cagr_3y')) or 0) * 100
            return str(max(0, min(25, round(12 + cagr))))
        if axis == 'quality':
            nm = (_safe_float(m.get('ebitda_margin_ltm')) or 0.1)
            if nm > 1: nm /= 100
            return str(max(0, min(25, round(10 + nm * 50))))
        if axis == 'momentum':
            p3m = (_safe_float(m.get('perf_3m')) or 0) * 100
            return str(max(0, min(25, round(12 + p3m * 0.15))))
        return "\u2014"

    axes = [('value', 'Valeur'), ('growth', 'Croissance'), ('quality', 'Qualite'), ('momentum', 'Momentum')]
    y_ax = y_dec + 0.75
    for xi, (axis, lbl) in enumerate(axes):
        xp = 1.02 + xi * 5.84
        w_ax = 5.6
        # Rect englobe label + valeurs (h=1.35 au lieu de 0.48)
        add_rect(slide, xp, y_ax, w_ax, 1.35, GREY_BG)
        add_rect(slide, xp, y_ax, w_ax, 0.04, NAVY_MID)  # ligne accent top
        add_text_box(slide, xp + 0.1, y_ax + 0.08, w_ax - 0.2, 0.35, lbl, 8, GREY_TXT, wrap=False)
        add_text_box(slide, xp + 0.1, y_ax + 0.50, 2.7, 0.65,
                     _score_axis(m_a, axis) + " / 25", 11, NAVY_MID, bold=True)
        add_text_box(slide, xp + 2.9, y_ax + 0.50, 2.5, 0.65,
                     _score_axis(m_b, axis) + " / 25", 11, GREEN_MID, bold=True)

    # KPIs finaux
    y_f = 9.65
    rec_a = (m_a.get('recommendation') or 'HOLD').upper()
    rec_b = (m_b.get('recommendation') or 'HOLD').upper()

    def _fill_rec(r):
        if r == 'BUY': return GREEN_PALE, GREEN
        if r == 'SELL': return RED_PALE, RED
        return AMBER_PALE, AMBER

    fa, aa = _fill_rec(rec_a)
    fb, ab = _fill_rec(rec_b)

    kpi_box(slide, 1.02,  y_f, 5.5, 2.1, str(fs_a) + " / 100", f"FinSight Score — {tkr_a}", "", COLOR_A_PAL, COLOR_A)
    kpi_box(slide, 7.0,   y_f, 5.5, 2.1, str(fs_b) + " / 100", f"FinSight Score — {tkr_b}", "", COLOR_B_PAL, COLOR_B)
    kpi_box(slide, 12.98, y_f, 5.5, 2.1, rec_a, f"Recommandation — {tkr_a}", _frpct(m_a.get('conviction')) + " conviction", fa, aa)
    kpi_box(slide, 18.95, y_f, 5.5, 2.1, rec_b, f"Recommandation — {tkr_b}", _frpct(m_b.get('conviction')) + " conviction", fb, ab)

    return slide


def _slide_theses(prs, m_a: dict, m_b: dict, synthesis: dict):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    tkr_a = m_a.get('ticker_a') or 'A'
    tkr_b = m_b.get('ticker_b') or 'B'

    navy_bar(slide)
    footer_bar(slide)
    slide_title(slide, "Theses d'Investissement Bull / Bear")
    section_dots(slide, 5)
    _company_header_band(slide, tkr_a, tkr_b)

    bull_a = synthesis.get('bull_a') or "Croissance superieure a la moyenne sectorielle avec fort levier operationnel."
    bear_a = synthesis.get('bear_a') or "Valorisation premium susceptible de re-rating si la croissance ralentit."
    bull_b = synthesis.get('bull_b') or "Generation de cash robuste et bilan solide pour soutenir la valorisation."
    bear_b = synthesis.get('bear_b') or "Exposition a un secteur cyclique amplifiant la volatilite des marges."

    # Bande 52W — cours & fourchette
    cur_a = "EUR" if (m_a.get('currency_a') or 'USD') == 'EUR' else '$'
    cur_b = "EUR" if (m_b.get('currency_b') or 'USD') == 'EUR' else '$'
    y_52 = 2.42
    add_rect(slide, 1.02, y_52, 11.44, 0.65, NAVY_PALE)
    add_rect(slide, 1.02, y_52, 0.13, 0.65, COLOR_A)
    add_text_box(slide, 1.25, y_52 + 0.05, 5.3, 0.3,
                 f"Cours : {_fr(m_a.get('share_price'), 1)} {cur_a}", 7.5, NAVY, bold=True)
    add_text_box(slide, 1.25, y_52 + 0.33, 5.3, 0.3,
                 f"52W : {_fr(m_a.get('week52_low'), 1)} - {_fr(m_a.get('week52_high'), 1)} {cur_a}", 7, GREY_TXT)
    add_text_box(slide, 6.8,  y_52 + 0.05, 5.4, 0.3,
                 f"Perf 1Y : {_frpct(m_a.get('perf_1y'), signed=True)}  \u00b7  Div. : {_frpct(m_a.get('dividend_yield'))}", 7, GREY_TXT)

    add_rect(slide, 12.94, y_52, 11.44, 0.65, GREEN_PALE)
    add_rect(slide, 12.94, y_52, 0.13, 0.65, COLOR_B)
    add_text_box(slide, 13.17, y_52 + 0.05, 5.3, 0.3,
                 f"Cours : {_fr(m_b.get('share_price'), 1)} {cur_b}", 7.5, GREEN, bold=True, wrap=False)
    add_text_box(slide, 13.17, y_52 + 0.33, 5.3, 0.3,
                 f"52W : {_fr(m_b.get('week52_low'), 1)} - {_fr(m_b.get('week52_high'), 1)} {cur_b}", 7, GREY_TXT, wrap=False)
    add_text_box(slide, 18.5, y_52 + 0.05, 5.7, 0.30,
                 f"Perf 1Y : {_frpct(m_b.get('perf_1y'), signed=True)}  \u00b7  Div. : {_frpct(m_b.get('dividend_yield'))}", 7, GREY_TXT, wrap=False)

    # Panel A
    y0 = 3.2
    # Bull A
    add_rect(slide, 1.02, y0, 11.44, 0.5, GREEN_PALE)
    add_rect(slide, 1.02, y0, 0.18, 0.5, GREEN)
    add_text_box(slide, 1.35, y0 + 0.07, 11.0, 0.38, f"BULL — {tkr_a}", 8.5, GREEN, bold=True)
    add_rect(slide, 1.02, y0 + 0.55, 11.44, 3.7, GREY_BG)
    add_text_box(slide, 1.15, y0 + 0.65, 11.15, 3.5, _fit(bull_a, 350), 8.5, BLACK, wrap=True)

    # Bear A
    y1 = y0 + 4.5
    add_rect(slide, 1.02, y1, 11.44, 0.5, RED_PALE)
    add_rect(slide, 1.02, y1, 0.18, 0.5, RED)
    add_text_box(slide, 1.35, y1 + 0.07, 11.0, 0.38, f"BEAR — {tkr_a}", 8.5, RED, bold=True)
    add_rect(slide, 1.02, y1 + 0.55, 11.44, 3.5, GREY_BG)
    add_text_box(slide, 1.15, y1 + 0.65, 11.15, 3.3, _fit(bear_a, 350), 8.5, BLACK, wrap=True)

    # Panel B
    # Bull B
    add_rect(slide, 12.94, y0, 11.44, 0.5, GREEN_PALE)
    add_rect(slide, 12.94, y0, 0.18, 0.5, GREEN)
    add_text_box(slide, 13.27, y0 + 0.07, 11.0, 0.38, f"BULL — {tkr_b}", 8.5, GREEN, bold=True)
    add_rect(slide, 12.94, y0 + 0.55, 11.44, 3.7, GREY_BG)
    add_text_box(slide, 13.07, y0 + 0.65, 11.15, 3.5, _fit(bull_b, 350), 8.5, BLACK, wrap=True)

    # Bear B
    add_rect(slide, 12.94, y1, 11.44, 0.5, RED_PALE)
    add_rect(slide, 12.94, y1, 0.18, 0.5, RED)
    add_text_box(slide, 13.27, y1 + 0.07, 11.0, 0.38, f"BEAR — {tkr_b}", 8.5, RED, bold=True)
    add_rect(slide, 12.94, y1 + 0.55, 11.44, 3.5, GREY_BG)
    add_text_box(slide, 13.07, y1 + 0.65, 11.15, 3.3, _fit(bear_b, 350), 8.5, BLACK, wrap=True)

    return slide


def _slide_price_chart(prs, m_a: dict, m_b: dict):
    """Slide 21 — Cours boursiers 52 semaines : texte perf a gauche + graphique a droite."""
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    tkr_a = m_a.get('ticker_a') or 'A'
    tkr_b = m_b.get('ticker_b') or 'B'
    cur_a = "EUR" if (m_a.get('currency_a') or 'USD') == 'EUR' else '$'
    cur_b = "EUR" if (m_b.get('currency_b') or 'USD') == 'EUR' else '$'

    navy_bar(slide)
    footer_bar(slide)
    slide_title(slide, "Performance Boursiere 52 Semaines")
    section_dots(slide, 5)
    _company_header_band(slide, tkr_a, tkr_b)

    # ---------- Colonne gauche — analyse textuelle ----------
    x_txt = 1.02
    w_txt = 11.5
    y0 = 2.45

    # Bloc A
    add_rect(slide, x_txt, y0, w_txt, 0.52, COLOR_A_PAL)
    add_rect(slide, x_txt, y0, 0.13, 0.52, COLOR_A)
    add_text_box(slide, x_txt + 0.22, y0 + 0.10, w_txt - 0.3, 0.35,
                 f"{tkr_a} — Performance & Prix", 8.5, NAVY_MID, bold=True, wrap=False)

    rows_a = [
        ("Cours actuel",   _fr(m_a.get('share_price'), 2) + " " + cur_a),
        ("52W High",       _fr(m_a.get('week52_high'), 2) + " " + cur_a),
        ("52W Low",        _fr(m_a.get('week52_low'),  2) + " " + cur_a),
        ("Perf. 1 Mois",   _frpct(m_a.get('perf_1m'), signed=True)),
        ("Perf. 3 Mois",   _frpct(m_a.get('perf_3m'), signed=True)),
        ("Perf. 1 An",     _frpct(m_a.get('perf_1y'), signed=True)),
        ("Div. Yield",     _frpct(m_a.get('dividend_yield'))),
        ("Volatilite 52S", _frpct(m_a.get('volatility_52w'))),
    ]
    rh = 0.55
    for ri, (lbl, val) in enumerate(rows_a):
        ry = y0 + 0.55 + ri * rh
        fill = WHITE if ri % 2 == 0 else GREY_BG
        add_rect(slide, x_txt, ry, w_txt, rh - 0.04, fill)
        add_text_box(slide, x_txt + 0.1, ry + 0.08, w_txt * 0.52, rh - 0.10, lbl, 7, GREY_TXT, wrap=False)
        add_text_box(slide, x_txt + w_txt * 0.52, ry + 0.08, w_txt * 0.45, rh - 0.10, val,
                     7.5, BLACK, bold=False, align=PP_ALIGN.RIGHT, wrap=False)

    # Bloc B
    y_b = y0 + 0.55 + len(rows_a) * rh + 0.25
    add_rect(slide, x_txt, y_b, w_txt, 0.52, COLOR_B_PAL)
    add_rect(slide, x_txt, y_b, 0.13, 0.52, COLOR_B)
    add_text_box(slide, x_txt + 0.22, y_b + 0.10, w_txt - 0.3, 0.35,
                 f"{tkr_b} — Performance & Prix", 8.5, GREEN_MID, bold=True, wrap=False)

    rows_b = [
        ("Cours actuel",   _fr(m_b.get('share_price'), 2) + " " + cur_b),
        ("52W High",       _fr(m_b.get('week52_high'), 2) + " " + cur_b),
        ("52W Low",        _fr(m_b.get('week52_low'),  2) + " " + cur_b),
        ("Perf. 1 Mois",   _frpct(m_b.get('perf_1m'), signed=True)),
        ("Perf. 3 Mois",   _frpct(m_b.get('perf_3m'), signed=True)),
        ("Perf. 1 An",     _frpct(m_b.get('perf_1y'), signed=True)),
        ("Div. Yield",     _frpct(m_b.get('dividend_yield'))),
        ("Volatilite 52S", _frpct(m_b.get('volatility_52w'))),
    ]
    for ri, (lbl, val) in enumerate(rows_b):
        ry = y_b + 0.55 + ri * rh
        fill = WHITE if ri % 2 == 0 else GREY_BG
        add_rect(slide, x_txt, ry, w_txt, rh - 0.04, fill)
        add_text_box(slide, x_txt + 0.1, ry + 0.08, w_txt * 0.52, rh - 0.10, lbl, 7, GREY_TXT, wrap=False)
        add_text_box(slide, x_txt + w_txt * 0.52, ry + 0.08, w_txt * 0.45, rh - 0.10, val,
                     7.5, BLACK, bold=False, align=PP_ALIGN.RIGHT, wrap=False)

    # ---------- Colonne droite — graphique cours ----------
    x_chart = x_txt + w_txt + 0.5
    w_chart = 25.4 - x_chart - 0.5
    buf = _chart_52w_price(tkr_a, tkr_b)
    if buf:
        _insert_chart(slide, buf, x_chart, 2.45, w_chart, 10.5)
    else:
        add_text_box(slide, x_chart, 6.5, w_chart, 1.0,
                     "Donnees de cours indisponibles", 9, GREY_TXT, wrap=True)

    return slide


def _slide_verdict(prs, m_a: dict, m_b: dict, synthesis: dict):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    tkr_a  = m_a.get('ticker_a') or 'A'
    tkr_b  = m_b.get('ticker_b') or 'B'
    winner = m_a.get('winner') or (tkr_a if (m_a.get('finsight_score') or 0) >= (m_b.get('finsight_score') or 0) else tkr_b)

    navy_bar(slide)
    footer_bar(slide)
    slide_title(slide, "Verdict Final Comparatif")
    section_dots(slide, 5)

    # Bandeau gagnant
    is_a = winner == tkr_a
    w_fill  = COLOR_A if is_a else COLOR_B
    w_label = f"CHOIX PREFERE : {winner}"
    add_rect(slide, 1.02, 1.73, 23.37, 1.0, w_fill)
    add_text_box(slide, 1.27, 1.85, 22.86, 0.76, w_label, 16, WHITE,
                 bold=True, align=PP_ALIGN.CENTER)

    # Texte verdict LLM — compact
    verdict_txt = synthesis.get('verdict_text') or ""
    y_v = 2.95
    add_rect(slide, 1.02, y_v, 23.37, 2.0, NAVY_PALE)
    add_rect(slide, 1.02, y_v, 0.18, 2.0, NAVY_MID)
    _v_clean = " ".join(_fit(verdict_txt, 340).split()) if verdict_txt else \
               f"{winner} est privilegiee sur la base des criteres de valorisation, de qualite bilancielle et de momentum."
    add_text_box(slide, 1.4, y_v + 0.12, 22.7, 1.85, _v_clean, 8.5, NAVY, wrap=True)

    # Scorecard comparative — compact
    y_sc = 5.35
    add_rect(slide, 1.02, y_sc, 23.37, 0.55, NAVY)
    add_text_box(slide, 1.15, y_sc + 0.08, 23.0, 0.42, "Scorecard Comparative", 8.5, WHITE, bold=True)

    def _score_label(a, b):
        if a is None and b is None: return "\u2014", "\u2014"
        fa = _safe_float(a) or 0
        fb = _safe_float(b) or 0
        if fa > fb: return f"\u2714 {a}", str(b) if b is not None else "\u2014"
        if fb > fa: return str(a) if a is not None else "\u2014", f"\u2714 {b}"
        return str(a) if a is not None else "\u2014", str(b) if b is not None else "\u2014"

    sc_rows_raw = [
        ("FinSight Score /100",       m_a.get('finsight_score'), m_b.get('finsight_score')),
        ("Piotroski F-Score /9",      m_a.get('piotroski_score'), m_b.get('piotroski_score')),
        ("Marge EBITDA LTM",          (m_a.get('ebitda_margin_ltm') * 100 if m_a.get('ebitda_margin_ltm') is not None else None),
                                      (m_b.get('ebitda_margin_ltm') * 100 if m_b.get('ebitda_margin_ltm') is not None else None)),
        ("ROIC",                      (m_a.get('roic') * 100 if m_a.get('roic') is not None else None),
                                      (m_b.get('roic') * 100 if m_b.get('roic') is not None else None)),
        ("Rev CAGR 3y",               (m_a.get('revenue_cagr_3y') * 100 if m_a.get('revenue_cagr_3y') is not None else None),
                                      (m_b.get('revenue_cagr_3y') * 100 if m_b.get('revenue_cagr_3y') is not None else None)),
    ]

    sc_rows = []
    for lbl, va, vb in sc_rows_raw:
        sa, sb = _score_label(
            round(float(va), 1) if va is not None else None,
            round(float(vb), 1) if vb is not None else None,
        )
        sc_rows.append((lbl, sa, sb))

    add_table(
        slide, 1.02, y_sc + 0.60, 23.37, 3.8,
        num_rows=len(sc_rows), num_cols=3,
        col_widths_pct=[0.44, 0.28, 0.28],
        header_data=["Critere", tkr_a, tkr_b],
        rows_data=sc_rows,
        header_fill=NAVY_MID, border_hex="DDDDDD"
    )

    # Trajectoire boursiere — narratif avec contexte
    y_traj = 10.0
    perf_a = _frpct(m_a.get('perf_1y'), signed=True)
    perf_b = _frpct(m_b.get('perf_1y'), signed=True)
    cur_a = "EUR" if (m_a.get('currency_a') or 'USD') == 'EUR' else '$'
    cur_b = "EUR" if (m_b.get('currency_b') or 'USD') == 'EUR' else '$'
    price_a = _fr(m_a.get('share_price'), 2)
    price_b = _fr(m_b.get('share_price'), 2)
    # Determiner le leader + contexte
    _p1y_a = m_a.get('perf_1y') or 0
    _p1y_b = m_b.get('perf_1y') or 0
    try: _fa, _fb = float(_p1y_a), float(_p1y_b)
    except: _fa, _fb = 0, 0
    _leader = tkr_a if _fa > _fb else tkr_b
    _beta_a = m_a.get('beta') or 0
    _sec = m_a.get('sector_a') or m_a.get('sector') or ''
    _ctx = ""
    if _sec:
        _ctx = f"Contexte : environnement de taux normalises, discrimination accrue entre profils de qualite dans le secteur {_sec}."
    traj_line1 = (f"{tkr_a} : {price_a} {cur_a} (1Y : {perf_a})  |  "
                  f"{tkr_b} : {price_b} {cur_b} (1Y : {perf_b})  |  "
                  f"Leader : {_leader}")
    add_rect(slide, 1.02, y_traj, 23.37, 1.5, GREY_BG)
    add_rect(slide, 1.02, y_traj, 0.13, 1.5, NAVY_MID)
    add_text_box(slide, 1.35, y_traj + 0.08, 22.8, 0.45, "TRAJECTOIRE BOURSIERE & CONTEXTE",
                 7, NAVY, bold=True, wrap=False)
    add_text_box(slide, 1.35, y_traj + 0.50, 22.8, 0.40, traj_line1, 7.5, GREY_TXT, wrap=False)
    if _ctx:
        add_text_box(slide, 1.35, y_traj + 0.90, 22.8, 0.50,
                     _truncate(_ctx, 180), 7, GREY_TXT, wrap=True)

    return slide


# ---------------------------------------------------------------------------
# Point d'entree principal
# ---------------------------------------------------------------------------

class ComparisonPPTXWriter:
    """
    Genere un pitchbook comparatif IB-quality en 22 slides.

    Usage :
        writer = ComparisonPPTXWriter()
        path   = writer.generate(state_a, state_b, output_path)
    """

    def generate(self, state_a: dict, state_b: dict,
                 output_path: Optional[str] = None) -> str:
        """
        Genere le fichier PPTX comparatif.

        Args:
            state_a, state_b : etats pipeline LangGraph
            output_path      : chemin de sortie (auto-genere si None)

        Returns:
            Chemin du fichier genere.
        """
        from pptx import Presentation
        from pptx.util import Cm

        # 1. Extraire les metriques
        from outputs.comparison_writer import extract_metrics, _fetch_supplements

        def _get_tkr(state, default="A"):
            snap = state.get("raw_data") or state.get("snapshot")
            if snap is not None and not isinstance(snap, (str, dict)):
                try:
                    return snap.ticker or default
                except Exception:
                    pass
            if isinstance(snap, dict):
                t = snap.get("ticker") or snap.get("company_info", {}).get("ticker")
                if t:
                    return t
            return state.get("ticker", default)

        tkr_a = _get_tkr(state_a, "A")
        tkr_b = _get_tkr(state_b, "B")

        log.info(f"[cmp_pptx] generation {tkr_a} vs {tkr_b}")
        supp_a = _fetch_supplements(tkr_a)
        supp_b = _fetch_supplements(tkr_b)
        m_a = extract_metrics(state_a, supp_a)
        m_b = extract_metrics(state_b, supp_b)

        # Garantir que les tickers corrects sont dans m_a/m_b pour la synthese LLM
        m_a["ticker_a"] = tkr_a
        m_b["ticker_b"] = tkr_b
        m_a["company_name_a"] = m_a.get("company_name_a") or tkr_a
        m_b["company_name_b"] = m_b.get("company_name_b") or tkr_b

        # Winner
        fs_a = m_a.get("finsight_score") or 0
        fs_b = m_b.get("finsight_score") or 0
        if fs_a >= fs_b:
            winner = tkr_a; verdict_str = f"{tkr_a} privilege"
        else:
            winner = tkr_b; verdict_str = f"{tkr_b} privilege"
        m_a["winner"] = m_b["winner"] = winner
        m_a["verdict_relative"] = m_b["verdict_relative"] = verdict_str

        # 2. Generer la synthese LLM
        log.info("[cmp_pptx] generation synthese LLM...")
        synthesis = _generate_synthesis(m_a, m_b)

        # 3. Creer la presentation
        prs = Presentation()
        prs.slide_width  = Cm(25.4)
        prs.slide_height = Cm(14.29)

        # Slide 1 — Cover
        _slide_cover(prs, m_a, m_b)
        # Slide 2 — Executive Summary
        _slide_exec_summary(prs, m_a, m_b, synthesis)
        # Slide 3 — Sommaire
        _slide_sommaire(prs, tkr_a, tkr_b)
        # Slide 4 — Divider 01
        divider_slide(prs, "01", "Profil & Identite",
                      f"Presentation comparative de {tkr_a} et {tkr_b}")
        # Slide 5 — Profil comparatif
        _slide_profil(prs, m_a, m_b)
        # Slide 6 — Divider 02
        divider_slide(prs, "02", "Performance Financiere",
                      "Compte de Resultat, Marges, Bilan & Liquidite")
        # Slide 7 — P&L compare
        _slide_pl(prs, m_a, m_b, synthesis)
        # Slide 8 — Marges & Rentabilite
        _slide_marges(prs, m_a, m_b, synthesis)
        # Slide 9 — Bilan & Liquidite
        _slide_bilan(prs, m_a, m_b)
        # Slide 10 — Divider 03
        divider_slide(prs, "03", "Valorisation",
                      "Multiples, DCF, Monte Carlo et Football Field")
        # Slide 11 — Multiples
        _slide_multiples(prs, m_a, m_b, synthesis)
        # Slide 12 — DCF
        _slide_dcf(prs, m_a, m_b)
        # Slide 13 — Monte Carlo + Football Field
        _slide_monte_carlo(prs, m_a, m_b)
        # Slide 14 — Divider 04
        divider_slide(prs, "04", "Qualite & Risques",
                      "Piotroski, Beneish, Altman, Beta, VaR et Momentum")
        # Slide 15 — Qualite financiere
        _slide_piotroski(prs, m_a, m_b, synthesis)
        # Slide 16 — Risque & Momentum
        _slide_risque(prs, m_a, m_b)
        # Slide 17 — Divider 05
        divider_slide(prs, "05", "Verdict",
                      "FinSight Score, Theses et Recommandation Finale")
        # Slide 18 — FinSight Score
        _slide_finsight_score(prs, m_a, m_b)
        # Slide 19 — Theses Bull/Bear
        _slide_theses(prs, m_a, m_b, synthesis)
        # Slide 20 — Cours boursiers 52 semaines
        _slide_price_chart(prs, m_a, m_b)
        # Slide 21 — Verdict final
        _slide_verdict(prs, m_a, m_b, synthesis)

        # 4. Sauvegarder
        if output_path is None:
            out_dir = Path(__file__).parent / "generated" / "cli_tests"
            out_dir.mkdir(parents=True, exist_ok=True)
            output_path = str(out_dir / f"{tkr_a}_vs_{tkr_b}_comparison.pptx")

        prs.save(output_path)
        log.info(f"[cmp_pptx] fichier sauvegarde : {output_path}")
        return output_path

    def generate_bytes(self, state_a: dict, state_b: dict) -> bytes:
        """Genere le fichier en memoire et retourne les bytes."""
        from pptx import Presentation
        from pptx.util import Cm
        import tempfile, os

        tmp = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
        tmp.close()
        try:
            self.generate(state_a, state_b, tmp.name)
            with open(tmp.name, 'rb') as f:
                return f.read()
        finally:
            try:
                os.unlink(tmp.name)
            except Exception:
                pass
