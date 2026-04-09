# =============================================================================
# FinSight IA — Indice Comparison PPTX Writer (15 slides)
# outputs/indice_comparison_pptx_writer.py
#
# Pitchbook comparatif IB-quality pour deux indices boursiers.
# Usage :
#   from outputs.indice_comparison_pptx_writer import IndiceComparisonPPTXWriter
#   buf = IndiceComparisonPPTXWriter.generate(data)   # -> bytes
# =============================================================================
from __future__ import annotations

import io
import logging
import re as _re
from typing import Optional

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import numpy as np

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

log = logging.getLogger(__name__)

# ── Dimensions ────────────────────────────────────────────────────────────────
_SW = Cm(25.4)
_SH = Cm(14.3)

# ── Palette ───────────────────────────────────────────────────────────────────
_NAVY   = RGBColor(0x1B, 0x3A, 0x6B)
_NAVYL  = RGBColor(0x2A, 0x52, 0x98)
_NAVYP  = RGBColor(0xEE, 0xF3, 0xFA)
_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
_GRAYL  = RGBColor(0xF5, 0xF7, 0xFA)
_GRAYM  = RGBColor(0xE8, 0xEC, 0xF0)
_GRAYT  = RGBColor(0x55, 0x55, 0x55)
_GRAYD  = RGBColor(0xAA, 0xAA, 0xAA)
_BLACK  = RGBColor(0x1A, 0x1A, 0x1A)
_BUY    = RGBColor(0x1A, 0x7A, 0x4A)
_SELL   = RGBColor(0xA8, 0x20, 0x20)
_HOLD   = RGBColor(0xB0, 0x60, 0x00)
_BUY_L  = RGBColor(0xE8, 0xF5, 0xEE)
_SELL_L = RGBColor(0xFB, 0xEB, 0xEB)
_HOLD_L = RGBColor(0xFD, 0xF3, 0xE5)
_COL_A  = RGBColor(0x2E, 0x5F, 0xA3)   # bleu fonce
_COL_B  = RGBColor(0x1A, 0x7A, 0x4A)   # vert
_COL_AL = RGBColor(0xEE, 0xF3, 0xFA)
_COL_BL = RGBColor(0xEA, 0xF4, 0xEF)

_MONTHS_FR = {1:"janvier",2:"fevrier",3:"mars",4:"avril",5:"mai",6:"juin",
              7:"juillet",8:"aout",9:"septembre",10:"octobre",11:"novembre",12:"decembre"}


# ── Helpers texte ─────────────────────────────────────────────────────────────

def _x(text) -> str:
    if text is None:
        return ""
    s = _re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', str(text))
    s = _re.sub(r'\*\*(.+?)\*\*', r'\1', s, flags=_re.DOTALL)
    s = _re.sub(r'\*(.+?)\*', r'\1', s, flags=_re.DOTALL)
    s = _re.sub(r'\*+', '', s)
    return s


def _fr_pct(v, signed=False) -> str:
    if v is None:
        return "\u2014"
    try:
        fv = float(v)
        # Valeurs deja en % (>2) ou en decimal (<2)
        if abs(fv) > 2:
            pass  # deja en %
        else:
            fv = fv * 100
        s = f"{fv:+.1f}" if signed else f"{fv:.1f}"
        return s.replace(".", ",") + "\u00a0%"
    except Exception:
        return "\u2014"


def _fr_pct_signed(v) -> str:
    return _fr_pct(v, signed=True)


def _fr_num(v, dp=1) -> str:
    if v is None:
        return "\u2014"
    try:
        return f"{float(v):.{dp}f}".replace(".", ",")
    except Exception:
        return "\u2014"


def _fr_date() -> str:
    import datetime
    d = datetime.date.today()
    return f"{d.day} {_MONTHS_FR[d.month]} {d.year}"


def _sig_c(signal: str) -> RGBColor:
    s = str(signal)
    if "Surp" in s or "Positif" in s:
        return _BUY
    if "Sous" in s or "Negatif" in s:
        return _SELL
    return _HOLD


def _sig_l(signal: str) -> RGBColor:
    s = str(signal)
    if "Surp" in s or "Positif" in s:
        return _BUY_L
    if "Sous" in s or "Negatif" in s:
        return _SELL_L
    return _HOLD_L


# ── Helpers PPTX ──────────────────────────────────────────────────────────────

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, x, y, w, h, fill=None, line_col=None, line_w=0.5):
    shape = slide.shapes.add_shape(1, Cm(x), Cm(y), Cm(w), Cm(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_col:
        shape.line.color.rgb = line_col
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape


def _txb(slide, text, x, y, w, h, size=9, bold=False, color=None,
         align=PP_ALIGN.LEFT, italic=False, wrap=True):
    color = color or _BLACK
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = _x(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def _pic(slide, img_bytes, x, y, w, h):
    buf = io.BytesIO(img_bytes)
    slide.shapes.add_picture(buf, Cm(x), Cm(y), Cm(w), Cm(h))


def _table(slide, data, x, y, w, h, col_widths=None,
           hdr_fill=_NAVY, font_size=7.5):
    rows, cols = len(data), len(data[0]) if data else 1
    tbl_shape = slide.shapes.add_table(rows, cols, Cm(x), Cm(y), Cm(w), Cm(h))
    tbl = tbl_shape.table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Cm(cw * w / total)
    for r, row_data in enumerate(data):
        is_hdr = (r == 0)
        bg = hdr_fill if is_hdr else (_GRAYL if r % 2 == 0 else _WHITE)
        for c, val in enumerate(row_data):
            cell = tbl.cell(r, c)
            cell.text = _x(val) if val is not None else "\u2014"
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            run = para.runs[0] if para.runs else para.add_run()
            run.font.size = Pt(font_size)
            run.font.bold = is_hdr
            run.font.color.rgb = _WHITE if is_hdr else _BLACK
            if is_hdr:
                cell.fill.solid()
                cell.fill.fore_color.rgb = hdr_fill
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg
    return tbl


def _header(slide, title, subtitle=""):
    _rect(slide, 0, 0, 25.4, 1.65, _NAVY)
    _txb(slide, title, 1.02, 0.28, 15.6, 0.97, 13, True, _WHITE)
    if subtitle:
        _txb(slide, subtitle, 1.02, 1.73, 23.37, 0.5, 8.5, False, _NAVYL)


def _footer(slide):
    _rect(slide, 0, 13.39, 25.4, 0.91, _NAVY)
    _txb(slide, "FinSight IA  \u00b7  Usage confidentiel",
         1.02, 13.44, 23.37, 0.56, 7, False, _WHITE)


def _index_band(slide, name_a, name_b):
    """Bande bicolore sous le titre identifiant les deux indices."""
    _rect(slide, 1.02, 1.73, 11.44, 0.52, _COL_AL)
    _rect(slide, 1.02, 1.73, 0.15, 0.52, _COL_A)
    _txb(slide, name_a, 1.3, 1.77, 11.0, 0.44, 8.5, True, _NAVY)
    _rect(slide, 12.94, 1.73, 11.44, 0.52, _COL_BL)
    _rect(slide, 12.94, 1.73, 0.15, 0.52, _COL_B)
    _txb(slide, name_b, 13.22, 1.77, 11.0, 0.44, 8.5, True, _NAVY)


def _divider(prs, num_str, title, subtitle):
    slide = _blank(prs)
    from pptx.oxml.ns import qn
    from lxml import etree
    # Fond navy
    cSld = slide._element.find(qn('p:cSld'))
    existing = cSld.find(qn('p:bg'))
    if existing is not None:
        cSld.remove(existing)
    bg_xml = (
        '<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<p:bgPr><a:solidFill><a:srgbClr val="1B3A6B"/></a:solidFill>'
        '<a:effectLst/></p:bgPr></p:bg>'
    )
    bg_elem = etree.fromstring(bg_xml)
    spTree = cSld.find(qn('p:spTree'))
    cSld.insert(list(cSld).index(spTree) if spTree is not None else 0, bg_elem)

    _rect(slide, 0, 0, 0.3, 14.3, _NAVYL)
    # Numero en filigrane
    txb = slide.shapes.add_textbox(Cm(1.27), Cm(3.8), Cm(22.86), Cm(4.57))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = _x(num_str)
    run.font.size = Pt(80)
    run.font.bold = True
    from pptx.oxml.ns import qn as _qn
    from lxml import etree as _et
    r_elem = run._r
    rPr = r_elem.find(_qn('a:rPr'))
    if rPr is None:
        rPr = _et.Element(_qn('a:rPr'))
        r_elem.insert(0, rPr)
    solidFill = _et.SubElement(rPr, _qn('a:solidFill'))
    srgbClr   = _et.SubElement(solidFill, _qn('a:srgbClr'))
    srgbClr.set('val', 'FFFFFF')
    alpha_elem = _et.SubElement(srgbClr, _qn('a:alpha'))
    alpha_elem.set('val', '15000')
    rPr.remove(solidFill)
    rPr.insert(0, solidFill)

    _txb(slide, title, 1.27, 5.0, 21.59, 2.0, 28, True, _WHITE)
    _rect(slide, 1.27, 7.5, 7.62, 0.01, _GRAYD)
    _txb(slide, subtitle, 1.27, 7.8, 17.78, 0.89, 10, False,
         RGBColor(0xAA, 0xBB, 0xDD))
    _txb(slide, "FinSight IA  \u00b7  Usage confidentiel",
         1.02, 13.34, 23.37, 0.56, 7, False,
         RGBColor(0x66, 0x77, 0xAA))
    return slide


def _kpi_box(slide, x, y, w, h, value, label, color_accent=_NAVYL, bg=None):
    bg = bg or _NAVYP
    _rect(slide, x, y, w, h, RGBColor(bg[0], bg[1], bg[2]) if isinstance(bg, tuple) else bg)
    _rect(slide, x, y, 0.13, h, color_accent)
    _txb(slide, value, x+0.25, y+0.12, w-0.38, h*0.55, 17, True, _NAVY)
    _txb(slide, label, x+0.25, y+h*0.55+0.05, w-0.38, h*0.4, 7, False, _GRAYT)


def _chart_buf(fig) -> bytes:
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close(fig)
    buf.seek(0)
    return buf.read()


# ── Slides ────────────────────────────────────────────────────────────────────

def _slide_cover(prs, d: dict):
    slide = _blank(prs)
    name_a = d.get("name_a", "Indice A")
    name_b = d.get("name_b", "Indice B")
    date_str = d.get("date", _fr_date())

    # Fond navy haut
    _rect(slide, 0, 0, 25.4, 2.5, _NAVY)
    _txb(slide, "FinSight IA", 1.27, 0.46, 22.86, 0.71, 10, False,
         RGBColor(0x88, 0x99, 0xBB), align=PP_ALIGN.CENTER)
    _rect(slide, 10.7, 1.32, 4.0, 0.05, _WHITE)
    _txb(slide, "Pitchbook Comparatif  \u2014  Analyse d'Indice",
         1.27, 1.52, 22.86, 0.81, 11, False,
         RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER)

    # Titre "A vs B"
    _txb(slide, f"{name_a}  vs  {name_b}",
         1.27, 3.8, 22.86, 1.8, 36, True, _NAVY, align=PP_ALIGN.CENTER)

    # Signal A
    sig_a = d.get("signal_a", "Neutre")
    sig_b = d.get("signal_b", "Neutre")
    sc_a  = d.get("score_a", 50)
    sc_b  = d.get("score_b", 50)

    _rect(slide, 1.27, 6.2, 11.0, 1.1, _sig_l(sig_a))
    _rect(slide, 1.27, 6.2, 0.18, 1.1, _sig_c(sig_a))
    _txb(slide, f"\u25cf {sig_a}  \u00b7  Score {sc_a}/100",
         1.6, 6.28, 10.5, 0.9, 9, True, _NAVY, align=PP_ALIGN.CENTER)

    # Signal B
    _rect(slide, 13.13, 6.2, 11.0, 1.1, _sig_l(sig_b))
    _rect(slide, 13.13, 6.2, 0.18, 1.1, _sig_c(sig_b))
    _txb(slide, f"\u25cf {sig_b}  \u00b7  Score {sc_b}/100",
         13.46, 6.28, 10.5, 0.9, 9, True, _NAVY, align=PP_ALIGN.CENTER)

    # Currency / zone
    cur_a = d.get("currency_a", "USD")
    cur_b = d.get("currency_b", "EUR")
    _txb(slide, f"Zone : {cur_a}", 1.27, 7.5, 11.0, 0.56, 8.5, False, _GRAYT,
         align=PP_ALIGN.CENTER)
    _txb(slide, f"Zone : {cur_b}", 13.13, 7.5, 11.0, 0.56, 8.5, False, _GRAYT,
         align=PP_ALIGN.CENTER)

    # Perf 1Y
    p1a = _fr_pct_signed(d.get("perf_1y_a"))
    p1b = _fr_pct_signed(d.get("perf_1y_b"))
    _txb(slide, f"Perf. 1 an : {p1a}", 1.27, 8.3, 11.0, 0.56, 8.5, False, _NAVY,
         align=PP_ALIGN.CENTER)
    _txb(slide, f"Perf. 1 an : {p1b}", 13.13, 8.3, 11.0, 0.56, 8.5, False, _NAVY,
         align=PP_ALIGN.CENTER)

    # Ligne separatrice verticale
    _rect(slide, 12.45, 3.5, 0.02, 7.0, _GRAYM)

    # Date + disclaimer
    _rect(slide, 1.02, 12.4, 23.37, 0.03, _GRAYD)
    _txb(slide, "Rapport confidentiel", 1.02, 12.65, 11.43, 0.56, 8, False, _GRAYT)
    _txb(slide, date_str, 12.95, 12.65, 11.43, 0.56, 8, False, _GRAYT,
         align=PP_ALIGN.RIGHT)


def _slide_overview(prs, d: dict):
    """Vue d'ensemble — 4 KPIs comparatifs + mini-tableau performances."""
    slide = _blank(prs)
    name_a = d.get("name_a", "Indice A")
    name_b = d.get("name_b", "Indice B")
    _header(slide, "Vue d'Ensemble Comparative",
            f"{name_a} vs {name_b}  \u2014  indicateurs cles")
    _footer(slide)
    _index_band(slide, name_a, name_b)

    # 4 KPI boxes
    kpis = [
        (_fr_pct_signed(d.get("perf_1y_a")) + " / " + _fr_pct_signed(d.get("perf_1y_b")),
         "Perf. 1 an (A / B)"),
        (_fr_num(d.get("vol_1y_a")) + " / " + _fr_num(d.get("vol_1y_b")) + "\u00a0%",
         "Volatilite 1 an (A / B)"),
        (_fr_num(d.get("pe_fwd_a"), 1) + "x / " + _fr_num(d.get("pe_fwd_b"), 1) + "x",
         "P/E Forward (A / B)"),
        (str(d.get("score_a", "\u2014")) + " / " + str(d.get("score_b", "\u2014")),
         "Score FinSight (A / B)"),
    ]
    xs = [1.02, 6.92, 12.81, 18.7]
    for i, (val, lbl) in enumerate(kpis):
        _kpi_box(slide, xs[i], 2.55, 5.6, 1.9, val, lbl)

    # Tableau performances
    ytd_a  = _fr_pct_signed(d.get("perf_ytd_a"))
    ytd_b  = _fr_pct_signed(d.get("perf_ytd_b"))
    p1a    = _fr_pct_signed(d.get("perf_1y_a"))
    p1b    = _fr_pct_signed(d.get("perf_1y_b"))
    p3a    = _fr_pct_signed(d.get("perf_3y_a"))
    p3b    = _fr_pct_signed(d.get("perf_3y_b"))
    p5a    = _fr_pct_signed(d.get("perf_5y_a"))
    p5b    = _fr_pct_signed(d.get("perf_5y_b"))

    rows = [
        ["Horizon", name_a[:20], name_b[:20], "Ecart"],
        ["YTD",  ytd_a, ytd_b, _ecart_pct(d.get("perf_ytd_a"), d.get("perf_ytd_b"))],
        ["1 an",  p1a,  p1b,  _ecart_pct(d.get("perf_1y_a"),  d.get("perf_1y_b"))],
        ["3 ans", p3a,  p3b,  _ecart_pct(d.get("perf_3y_a"),  d.get("perf_3y_b"))],
        ["5 ans", p5a,  p5b,  _ecart_pct(d.get("perf_5y_a"),  d.get("perf_5y_b"))],
    ]
    _table(slide, rows, 1.02, 5.1, 23.37, 4.8,
           col_widths=[3, 5, 5, 5])

    # Signal
    sig_a = d.get("signal_a", "Neutre")
    sig_b = d.get("signal_b", "Neutre")
    _rect(slide, 1.02, 10.5, 11.0, 0.85, _sig_l(sig_a))
    _rect(slide, 1.02, 10.5, 0.13, 0.85, _sig_c(sig_a))
    _txb(slide, f"Signal {name_a} : {sig_a} ({d.get('score_a',0)}/100)",
         1.3, 10.62, 10.5, 0.62, 8.5, True, _NAVY)
    _rect(slide, 13.37, 10.5, 11.0, 0.85, _sig_l(sig_b))
    _rect(slide, 13.37, 10.5, 0.13, 0.85, _sig_c(sig_b))
    _txb(slide, f"Signal {name_b} : {sig_b} ({d.get('score_b',0)}/100)",
         13.65, 10.62, 10.5, 0.62, 8.5, True, _NAVY)

    # LLM synthese
    llm_text = d.get("llm", {}).get("overview_read", "")
    if llm_text:
        _rect(slide, 1.02, 11.6, 23.37, 0.38, _NAVY)
        _txb(slide, "Synthese FinSight IA", 1.22, 11.63, 15.0, 0.32, 7.5, True, _WHITE)
        _txb(slide, llm_text, 1.02, 12.08, 23.37, 1.35, 8.5, False, _BLACK, wrap=True)


def _ecart_pct(va, vb) -> str:
    if va is None or vb is None:
        return "\u2014"
    try:
        fa, fb = float(va), float(vb)
        # convert if decimal
        if abs(fa) <= 2: fa *= 100
        if abs(fb) <= 2: fb *= 100
        diff = fa - fb
        s = f"{diff:+.1f}".replace(".", ",") + "\u00a0pp"
        return s
    except Exception:
        return "\u2014"


def _slide_sommaire(prs, d: dict):
    slide = _blank(prs)
    name_a = d.get("name_a", "Indice A")
    name_b = d.get("name_b", "Indice B")
    _header(slide, "Sommaire")
    _footer(slide)

    sections = [
        ("01", "Performance Historique",   "Evolution comparee, YTD / 1Y / 3Y / 5Y"),
        ("02", "Risque Comparatif",         "Volatilite, Sharpe, Max Drawdown"),
        ("03", "Valorisation",              "P/E Forward, P/B, Rendement du dividende, ERP"),
        ("04", "Composition Sectorielle",  "Poids GICS par secteur — ecarts entre indices"),
        ("05", "Constituants",              f"Top 5 valeurs {name_a} et {name_b}"),
        ("06", "Verdict",                   "Signal FinSight, score compare, recommandation"),
    ]

    y0 = 2.4
    dy = 1.6
    for i, (num, title, sub) in enumerate(sections):
        y = y0 + i * dy
        _rect(slide, 1.02, y, 0.7, 0.7, _NAVYL)
        _txb(slide, num, 1.02, y+0.08, 0.7, 0.55, 8.5, True, _WHITE,
             align=PP_ALIGN.CENTER, wrap=False)
        _txb(slide, title, 1.9, y+0.05, 16.0, 0.48, 9.5, True, _NAVY)
        _txb(slide, sub,   1.9, y+0.53, 16.0, 0.38, 8, False, _GRAYT)
        _rect(slide, 1.9, y+0.95, 22.5, 0.01, _GRAYM)


def _slide_perf_chart(prs, d: dict):
    """Performance historique — graphique normalisé + tableau."""
    slide = _blank(prs)
    name_a = d.get("name_a", "Indice A")
    name_b = d.get("name_b", "Indice B")
    _header(slide, "Performance Historique Comparee",
            "Evolution normalisee base 100 sur 1 an")
    _footer(slide)
    _index_band(slide, name_a, name_b)

    ph = d.get("perf_history")
    if not (ph and ph.get("dates") and ph.get("indice_a") and ph.get("indice_b")):
        _txb(slide, "Performance historique non disponible — donnees de cours requises.",
             1.5, 5.5, 14.0, 1.0, size=9, color=_GRAYT, italic=True)
    if ph and ph.get("dates") and ph.get("indice_a") and ph.get("indice_b"):
        try:
            fig, ax = plt.subplots(figsize=(11.5, 4.5))
            fig.patch.set_facecolor('white')
            ax.set_facecolor('#FAFBFD')

            dates_raw = ph["dates"]
            # Decimate for display
            n = len(dates_raw)
            step = max(1, n // 80)
            idx = list(range(0, n, step))
            if idx[-1] != n-1:
                idx.append(n-1)

            ys_a = ph["indice_a"]
            ys_b = ph["indice_b"]
            xs   = list(range(len(idx)))
            lbls = [dates_raw[i][:7] for i in idx]

            ax.plot(xs, [ys_a[i] for i in idx], color='#2E5FA3', linewidth=2,
                    label=name_a[:20], zorder=3)
            ax.plot(xs, [ys_b[i] for i in idx], color='#1A7A4A', linewidth=2,
                    label=name_b[:20], zorder=3)
            ax.axhline(100, color='#CCCCCC', linewidth=0.8, linestyle='--', zorder=1)

            ax.fill_between(xs, [ys_a[i] for i in idx], 100,
                            alpha=0.08, color='#2E5FA3')
            ax.fill_between(xs, [ys_b[i] for i in idx], 100,
                            alpha=0.08, color='#1A7A4A')

            # Ticks
            tick_step = max(1, len(xs)//8)
            ax.set_xticks(xs[::tick_step])
            ax.set_xticklabels(lbls[::tick_step], rotation=30, ha='right', fontsize=7)
            ax.set_ylabel("Base 100", fontsize=8)
            ax.yaxis.set_tick_params(labelsize=7)
            ax.legend(loc='upper left', fontsize=8, framealpha=0.9)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.grid(axis='y', alpha=0.3)
            fig.tight_layout(pad=0.5)

            img = _chart_buf(fig)
            _pic(slide, img, 1.02, 2.4, 15.5, 9.5)
        except Exception as e:
            log.warning(f"[indice_cmp_pptx] chart perf error: {e}")

    # Mini tableau perf
    rows = [
        ["Horizon", name_a[:18], name_b[:18], "Ecart"],
        ["YTD",     _fr_pct_signed(d.get("perf_ytd_a")),
                    _fr_pct_signed(d.get("perf_ytd_b")),
                    _ecart_pct(d.get("perf_ytd_a"), d.get("perf_ytd_b"))],
        ["1 an",    _fr_pct_signed(d.get("perf_1y_a")),
                    _fr_pct_signed(d.get("perf_1y_b")),
                    _ecart_pct(d.get("perf_1y_a"),  d.get("perf_1y_b"))],
        ["3 ans",   _fr_pct_signed(d.get("perf_3y_a")),
                    _fr_pct_signed(d.get("perf_3y_b")),
                    _ecart_pct(d.get("perf_3y_a"),  d.get("perf_3y_b"))],
        ["5 ans",   _fr_pct_signed(d.get("perf_5y_a")),
                    _fr_pct_signed(d.get("perf_5y_b")),
                    _ecart_pct(d.get("perf_5y_a"),  d.get("perf_5y_b"))],
    ]
    _table(slide, rows, 17.0, 2.4, 7.4, 5.5, col_widths=[2, 2, 2, 2])


def _slide_risque(prs, d: dict):
    slide = _blank(prs)
    name_a = d.get("name_a", "Indice A")
    name_b = d.get("name_b", "Indice B")
    _header(slide, "Risque Comparatif",
            "Volatilite, Sharpe ratio et Max Drawdown sur 12 mois")
    _footer(slide)
    _index_band(slide, name_a, name_b)

    vol_a   = d.get("vol_1y_a")
    vol_b   = d.get("vol_1y_b")
    sha_a   = d.get("sharpe_1y_a")
    sha_b   = d.get("sharpe_1y_b")
    dd_a    = d.get("max_dd_a")
    dd_b    = d.get("max_dd_b")

    # 3 KPIs comparatifs
    kpis = [
        (_fr_num(vol_a) + " / " + _fr_num(vol_b) + "\u00a0%",
         "Volatilite annualisee (A / B)",
         "Risque total exprime en ecart-type des rendements quotidiens"),
        (_fr_num(sha_a, 2) + " / " + _fr_num(sha_b, 2),
         "Sharpe ratio 1 an (A / B)",
         "Rendement ajuste du risque — rf : 10Y local"),
        (_fr_pct_signed(dd_a) + " / " + _fr_pct_signed(dd_b),
         "Max Drawdown (A / B)",
         "Perte maximale du plus haut au plus bas sur la periode"),
    ]
    xs = [1.02, 9.22, 17.42]
    for i, (val, lbl, sub) in enumerate(kpis):
        _rect(slide, xs[i], 2.55, 7.8, 2.4, _NAVYP)
        _rect(slide, xs[i], 2.55, 0.15, 2.4, _NAVYL)
        _txb(slide, val, xs[i]+0.3, 2.65, 7.2, 1.1, 15, True, _NAVY)
        _txb(slide, lbl, xs[i]+0.3, 3.65, 7.2, 0.5, 7.5, True, _GRAYT)
        _txb(slide, sub, xs[i]+0.3, 4.1, 7.2, 0.75, 7, False, _GRAYT)

    # Graphique barres radar (vol + sharpe + max_dd normalises)
    try:
        cats = ["Volatilite", "Sharpe", "Max Drawdown"]
        vals_a = [
            float(vol_a or 0),
            float(sha_a or 0),
            abs(float(dd_a or 0)) * 100 if dd_a and abs(float(dd_a)) < 1 else abs(float(dd_a or 0))
        ]
        vals_b = [
            float(vol_b or 0),
            float(sha_b or 0),
            abs(float(dd_b or 0)) * 100 if dd_b and abs(float(dd_b)) < 1 else abs(float(dd_b or 0))
        ]
        x = np.arange(len(cats))
        bw = 0.3
        fig, ax = plt.subplots(figsize=(10, 3.5))
        fig.patch.set_facecolor('white')
        ax.set_facecolor('#FAFBFD')
        b1 = ax.bar(x - bw/2, vals_a, bw, label=name_a[:20], color='#2E5FA3', alpha=0.85)
        b2 = ax.bar(x + bw/2, vals_b, bw, label=name_b[:20], color='#1A7A4A', alpha=0.85)
        ax.set_xticks(x)
        ax.set_xticklabels(cats, fontsize=9)
        ax.legend(fontsize=8, framealpha=0.9)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.grid(axis='y', alpha=0.3)
        # Annotations
        for bar in b1:
            h = bar.get_height()
            ax.annotate(f'{h:.1f}', xy=(bar.get_x() + bar.get_width()/2, h),
                       xytext=(0, 2), textcoords="offset points",
                       ha='center', va='bottom', fontsize=7.5, color='#2E5FA3')
        for bar in b2:
            h = bar.get_height()
            ax.annotate(f'{h:.1f}', xy=(bar.get_x() + bar.get_width()/2, h),
                       xytext=(0, 2), textcoords="offset points",
                       ha='center', va='bottom', fontsize=7.5, color='#1A7A4A')
        fig.tight_layout(pad=0.5)
        img = _chart_buf(fig)
        _pic(slide, img, 1.02, 5.2, 23.37, 6.7)
    except Exception as e:
        log.warning(f"[indice_cmp_pptx] risque chart error: {e}")

    llm_text = d.get("llm", {}).get("risque_read", "")
    if llm_text:
        _rect(slide, 1.02, 12.1, 23.37, 0.38, _NAVY)
        _txb(slide, "Lecture analytique FinSight IA", 1.22, 12.13, 15.0, 0.32, 7.5, True, _WHITE)
        _txb(slide, llm_text, 1.02, 12.58, 23.37, 1.0, 8.5, False, _BLACK, wrap=True)


def _slide_valorisation(prs, d: dict):
    slide = _blank(prs)
    name_a = d.get("name_a", "Indice A")
    name_b = d.get("name_b", "Indice B")
    _header(slide, "Valorisation Comparative",
            "P/E Forward, P/B, Rendement du dividende, ERP")
    _footer(slide)
    _index_band(slide, name_a, name_b)

    pe_a  = d.get("pe_fwd_a")
    pe_b  = d.get("pe_fwd_b")
    pb_a  = d.get("pb_a")
    pb_b  = d.get("pb_b")
    dy_a  = d.get("div_yield_a")
    dy_b  = d.get("div_yield_b")
    erp_a = d.get("erp_a", "\u2014")
    erp_b = d.get("erp_b", "\u2014")

    rows = [
        ["Indicateur", name_a[:22], name_b[:22], "Ecart", "Favorise"],
        ["P/E Forward",
         _fr_num(pe_a, 1) + "x",
         _fr_num(pe_b, 1) + "x",
         _ecart_mult(pe_a, pe_b),
         _favorise_lower(pe_a, pe_b, name_a, name_b)],
        ["P/B (Book Value)",
         _fr_num(pb_a, 1) + "x",
         _fr_num(pb_b, 1) + "x",
         _ecart_mult(pb_a, pb_b),
         _favorise_lower(pb_a, pb_b, name_a, name_b)],
        ["Rendement dividende",
         _fr_pct(dy_a) if dy_a else "\u2014",
         _fr_pct(dy_b) if dy_b else "\u2014",
         _ecart_pct(dy_a, dy_b),
         _favorise_higher(dy_a, dy_b, name_a, name_b)],
        ["ERP (prime de risque)", str(erp_a), str(erp_b), "\u2014", "\u2014"],
    ]
    _table(slide, rows, 1.02, 2.55, 23.37, 6.0,
           col_widths=[5, 4, 4, 4, 5])

    # Graphique barres PE vs PB — uniquement les categories avec au moins une valeur
    try:
        _cats_all = [("P/E Forward", pe_a, pe_b), ("P/B", pb_a, pb_b)]
        _cats_data = [(c, va, vb) for c, va, vb in _cats_all if va is not None or vb is not None]
        if not _cats_data:
            raise ValueError("no valuation data")
        cats = [c[0] for c in _cats_data]
        vals_a = [float(c[1] or 0) for c in _cats_data]
        vals_b = [float(c[2] or 0) for c in _cats_data]
        x = np.arange(len(cats))
        bw = 0.3
        fig, ax = plt.subplots(figsize=(7, 3.0))
        fig.patch.set_facecolor('white')
        ax.set_facecolor('#FAFBFD')
        b1 = ax.bar(x - bw/2, vals_a, bw, label=name_a[:20], color='#2E5FA3', alpha=0.85)
        b2 = ax.bar(x + bw/2, vals_b, bw, label=name_b[:20], color='#1A7A4A', alpha=0.85)
        ax.set_xticks(x)
        ax.set_xticklabels(cats, fontsize=10)
        ax.legend(fontsize=9, framealpha=0.9)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.grid(axis='y', alpha=0.3)
        for bar in b1:
            h = bar.get_height()
            if h > 0:
                ax.annotate(f'{h:.1f}x',
                           xy=(bar.get_x()+bar.get_width()/2, h),
                           xytext=(0, 2), textcoords="offset points",
                           ha='center', va='bottom', fontsize=8.5, color='#2E5FA3')
        for bar in b2:
            h = bar.get_height()
            if h > 0:
                ax.annotate(f'{h:.1f}x',
                           xy=(bar.get_x()+bar.get_width()/2, h),
                           xytext=(0, 2), textcoords="offset points",
                           ha='center', va='bottom', fontsize=8.5, color='#1A7A4A')
        fig.tight_layout(pad=0.5)
        img = _chart_buf(fig)
        _pic(slide, img, 9.0, 9.0, 14.0, 4.0)
    except Exception as e:
        log.warning(f"[indice_cmp_pptx] val chart error: {e}")

    # ERP commentary
    _rect(slide, 1.02, 9.1, 7.5, 3.8, _NAVYP)
    _rect(slide, 1.02, 9.1, 0.13, 3.8, _NAVYL)
    _txb(slide, "Equity Risk Premium",
         1.3, 9.2, 7.0, 0.5, 8.5, True, _NAVY)
    erp_txt = (
        f"{name_a} : ERP = {erp_a}  |  {name_b} : ERP = {erp_b}\n\n"
        f"L'ERP mesure la prime exigee par les investisseurs pour "
        f"detenir des actions plutot que des obligations sans risque. "
        f"Un ERP eleve signale une valorisation plus attractive ou un "
        f"risque percu plus important."
    )
    _txb(slide, erp_txt, 1.3, 9.75, 7.0, 3.0, 7.5, False, _GRAYT, wrap=True)


def _ecart_mult(va, vb) -> str:
    """Ecart en x entre deux multiples."""
    if va is None or vb is None:
        return "\u2014"
    try:
        fa, fb = float(va), float(vb)
        if fb == 0:
            return "\u2014"
        diff = fa - fb
        s = f"{diff:+.1f}".replace(".", ",") + "x"
        return s
    except Exception:
        return "\u2014"


def _favorise_lower(va, vb, na, nb) -> str:
    """Pour les multiples, plus bas = moins cher = favorable."""
    if va is None or vb is None:
        return "\u2014"
    try:
        return na[:14] if float(va) < float(vb) else (nb[:14] if float(vb) < float(va) else "Egal")
    except Exception:
        return "\u2014"


def _favorise_higher(va, vb, na, nb) -> str:
    """Pour rendement dividende, plus haut = plus genereux = favorable."""
    if va is None or vb is None:
        return "\u2014"
    try:
        fa = float(va) * 100 if abs(float(va)) < 1 else float(va)
        fb = float(vb) * 100 if abs(float(vb)) < 1 else float(vb)
        return na[:14] if fa > fb else (nb[:14] if fb > fa else "Egal")
    except Exception:
        return "\u2014"


def _slide_secteurs(prs, d: dict):
    """Composition sectorielle — barres horizontales cote a cote."""
    slide = _blank(prs)
    name_a = d.get("name_a", "Indice A")
    name_b = d.get("name_b", "Indice B")
    _header(slide, "Composition Sectorielle Comparative",
            "Poids par secteur GICS (%) — ecarts entre indices")
    _footer(slide)
    _index_band(slide, name_a, name_b)

    sector_cmp = d.get("sector_comparison", [])

    if sector_cmp:
        try:
            sects = [s[0][:28] for s in sector_cmp[:12]]
            wa    = [float(s[1] or 0) for s in sector_cmp[:12]]
            wb    = [float(s[2] or 0) for s in sector_cmp[:12]]

            fig, ax = plt.subplots(figsize=(14, max(4.0, len(sects) * 0.9)))
            fig.patch.set_facecolor('white')
            ax.set_facecolor('#FAFBFD')

            y   = np.arange(len(sects))
            bh  = 0.35

            ax.barh(y + bh/2, wa, bh, label=name_a[:20], color='#2E5FA3', alpha=0.85)
            ax.barh(y - bh/2, wb, bh, label=name_b[:20], color='#1A7A4A', alpha=0.85)

            ax.set_yticks(y)
            ax.set_yticklabels(sects, fontsize=8)
            ax.invert_yaxis()
            ax.set_xlabel("Poids (%)", fontsize=8)
            ax.legend(loc='lower right', fontsize=8, framealpha=0.9)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.grid(axis='x', alpha=0.3)
            fig.tight_layout(pad=0.5)

            img = _chart_buf(fig)
            _pic(slide, img, 1.02, 2.55, 23.37, 10.5)
        except Exception as e:
            log.warning(f"[indice_cmp_pptx] secteurs chart error: {e}")
    else:
        _txb(slide, "Donnees de composition sectorielle non disponibles.",
             1.02, 4.0, 23.37, 1.0, 9, False, _GRAYT)


def _slide_top5(prs, d: dict, which: str):
    """Top 5 constituants d'un indice."""
    is_a   = (which == "a")
    name   = d.get("name_a" if is_a else "name_b", "Indice")
    top5   = d.get("top5_a"  if is_a else "top5_b", [])
    color  = _COL_A if is_a else _COL_B
    color_l= _COL_AL if is_a else _COL_BL

    slide = _blank(prs)
    _header(slide, f"Top 5 Constituants \u2014 {name}",
            "Societes les plus representees dans l'indice")
    _footer(slide)

    # Bandeau couleur
    _rect(slide, 1.02, 1.73, 23.37, 0.52, color_l)
    _rect(slide, 1.02, 1.73, 0.15, 0.52, color)
    _txb(slide, name, 1.3, 1.77, 23.0, 0.44, 8.5, True, _NAVY)

    rows = [["Societe / Ticker", "Poids (%)", "Secteur"]]
    for item in top5[:5]:
        company  = str(item[0] if len(item) > 0 else "\u2014")[:40]
        ticker   = str(item[1] if len(item) > 1 else "")
        weight   = item[2] if len(item) > 2 else None
        sector   = str(item[3] if len(item) > 3 else "")[:30]
        wt_str   = _fr_num(weight, 1) + "\u00a0%" if weight is not None else "\u2014"
        label    = f"{company} ({ticker})" if ticker else company
        rows.append([label[:48], wt_str, sector])

    while len(rows) < 6:
        rows.append(["\u2014", "\u2014", "\u2014"])

    _table(slide, rows, 1.02, 2.55, 23.37, 7.0,
           col_widths=[12, 3, 8], hdr_fill=color)

    # Note
    _txb(slide, "Source : yfinance  \u00b7  Poids par score composite ou market cap.",
         1.02, 9.7, 23.37, 0.5, 7, False, _GRAYT)


def _slide_score(prs, d: dict):
    """Score FinSight comparatif."""
    slide = _blank(prs)
    name_a = d.get("name_a", "Indice A")
    name_b = d.get("name_b", "Indice B")
    _header(slide, "Signal FinSight Comparatif",
            "Score 0-100 et recommandation")
    _footer(slide)
    _index_band(slide, name_a, name_b)

    sc_a   = int(d.get("score_a", 50))
    sc_b   = int(d.get("score_b", 50))
    sig_a  = d.get("signal_a", "Neutre")
    sig_b  = d.get("signal_b", "Neutre")

    # Gauges textuelles A et B
    for is_a in (True, False):
        x0 = 1.02 if is_a else 13.37
        sc = sc_a if is_a else sc_b
        sig = sig_a if is_a else sig_b
        name = name_a if is_a else name_b
        col  = _sig_c(sig)
        col_l= _sig_l(sig)

        _rect(slide, x0, 2.55, 11.0, 5.5, col_l)
        _rect(slide, x0, 2.55, 0.18, 5.5, col)
        _txb(slide, str(sc), x0+0.4, 3.0, 10.2, 2.5, 60, True, col,
             align=PP_ALIGN.CENTER)
        _txb(slide, "/ 100", x0+0.4, 5.1, 10.2, 0.7, 18, False, _GRAYT,
             align=PP_ALIGN.CENTER)
        _txb(slide, sig, x0+0.4, 5.85, 10.2, 0.8, 14, True, col,
             align=PP_ALIGN.CENTER)
        _txb(slide, f"Signal FinSight  \u2014  {name}",
             x0+0.4, 6.75, 10.2, 0.56, 8, False, _GRAYT,
             align=PP_ALIGN.CENTER)

    # Gagnant
    winner = name_a if sc_a > sc_b else (name_b if sc_b > sc_a else "Egalite")
    _rect(slide, 1.02, 8.5, 23.37, 1.0, _NAVYP)
    _rect(slide, 1.02, 8.5, 0.15, 1.0, _NAVYL)
    _txb(slide, f"Meilleur signal : {winner}  ({sc_a} vs {sc_b})",
         1.3, 8.7, 22.0, 0.7, 11, True, _NAVY)

    # Tableau de bord
    rows = [
        ["Critere", name_a[:20], name_b[:20]],
        ["Score global",  f"{sc_a}/100", f"{sc_b}/100"],
        ["Signal",        sig_a, sig_b],
        ["P/E Forward",   _fr_num(d.get("pe_fwd_a"),1)+"x", _fr_num(d.get("pe_fwd_b"),1)+"x"],
        ["Volatilite",    _fr_num(d.get("vol_1y_a"))+" %", _fr_num(d.get("vol_1y_b"))+" %"],
        ["Perf. 1 an",    _fr_pct_signed(d.get("perf_1y_a")), _fr_pct_signed(d.get("perf_1y_b"))],
    ]
    _table(slide, rows, 1.02, 9.8, 23.37, 2.8, col_widths=[5, 4, 4])


def _slide_verdict(prs, d: dict):
    slide = _blank(prs)
    name_a = d.get("name_a", "Indice A")
    name_b = d.get("name_b", "Indice B")
    _header(slide, "Verdict Final",
            f"Comparaison {name_a} vs {name_b}  \u2014  recommandation")
    _footer(slide)
    _index_band(slide, name_a, name_b)

    sc_a  = int(d.get("score_a", 50))
    sc_b  = int(d.get("score_b", 50))
    sig_a = d.get("signal_a", "Neutre")
    sig_b = d.get("signal_b", "Neutre")
    winner = name_a if sc_a > sc_b else (name_b if sc_b > sc_a else None)
    loser  = name_b if winner == name_a else (name_a if winner == name_b else None)

    # Bandeau verdict
    if winner:
        _rect(slide, 1.02, 2.55, 23.37, 1.4, _BUY_L)
        _rect(slide, 1.02, 2.55, 0.2, 1.4, _BUY)
        _txb(slide, f"\u2713  A priviliegier : {winner}",
             1.45, 2.75, 22.0, 0.8, 13, True, _BUY)
    else:
        _rect(slide, 1.02, 2.55, 23.37, 1.4, _NAVYP)
        _rect(slide, 1.02, 2.55, 0.2, 1.4, _NAVYL)
        _txb(slide, "Profils similaires  \u2014  arbitrage selon objectif investor",
             1.45, 2.9, 22.0, 0.7, 11, True, _NAVY)

    # 4 arguments comparatifs
    args = _build_verdict_args(d, name_a, name_b)
    y0 = 4.4
    dy = 1.6
    n_args = len(args[:4])
    for i, (title, body) in enumerate(args[:4]):
        y = y0 + i * dy
        _rect(slide, 1.02, y, 0.15, 1.3, _NAVYL)
        _txb(slide, title, 1.3, y+0.05, 22.5, 0.42, 8.5, True, _NAVY)
        _txb(slide, body,  1.3, y+0.48, 22.5, 0.75, 8, False, _GRAYT, wrap=True)

    # LLM Recommandation — position dynamique apres les args
    llm_text = d.get("llm", {}).get("verdict_read", "")
    y_llm = y0 + n_args * dy + 0.3  # suit le dernier arg, quel que soit leur nombre
    if llm_text:
        _rect(slide, 1.02, y_llm, 23.37, 0.38, _NAVY)
        _txb(slide, "Recommandation FinSight IA", 1.22, y_llm + 0.03, 15.0, 0.32, 7.5, True, _WHITE)
        _txb(slide, llm_text, 1.02, y_llm + 0.48, 23.37, 1.25, 8.5, False, _BLACK, wrap=True)

    # Disclaimer
    y_disc = min(12.9, y_llm + 1.85)
    _txb(slide, "Analyse FinSight IA  \u00b7  Source : yfinance  "
                "\u00b7  Donnees a titre informatif uniquement  "
                "\u00b7  Ne constitue pas un conseil en investissement.",
         1.02, y_disc, 23.37, 0.4, 6.5, False, _GRAYT)


def _build_verdict_args(d, name_a, name_b):
    args = []
    sc_a = d.get("score_a", 50)
    sc_b = d.get("score_b", 50)
    pe_a = d.get("pe_fwd_a")
    pe_b = d.get("pe_fwd_b")
    p1a  = d.get("perf_1y_a")
    p1b  = d.get("perf_1y_b")
    sha_a= d.get("sharpe_1y_a")
    sha_b= d.get("sharpe_1y_b")
    dy_a = d.get("div_yield_a")
    dy_b = d.get("div_yield_b")

    # Perf
    if p1a is not None and p1b is not None:
        fa = float(p1a) * 100 if abs(float(p1a)) < 2 else float(p1a)
        fb = float(p1b) * 100 if abs(float(p1b)) < 2 else float(p1b)
        winner_p = name_a if fa > fb else name_b
        args.append(("Performance 1 an",
                      f"{name_a} : {fa:+.1f}%  vs  {name_b} : {fb:+.1f}%  "
                      f"\u2014  {winner_p} surperforme de {abs(fa-fb):.1f} pp."))

    # Valorisation
    if pe_a and pe_b:
        pea, peb = float(pe_a), float(pe_b)
        cheaper = name_a if pea < peb else name_b
        args.append(("Valorisation (P/E Forward)",
                      f"{name_a} : {pea:.1f}x  vs  {name_b} : {peb:.1f}x  "
                      f"\u2014  {cheaper} se traite a une decote relative."))

    # Risque / Sharpe
    if sha_a is not None and sha_b is not None:
        sa, sb = float(sha_a), float(sha_b)
        better = name_a if sa > sb else name_b
        args.append(("Rendement ajuste du risque (Sharpe)",
                      f"{name_a} : {sa:.2f}  vs  {name_b} : {sb:.2f}  "
                      f"\u2014  {better} offre un meilleur ratio rendement/risque."))

    # Dividende
    if dy_a and dy_b:
        fa = float(dy_a) * 100 if abs(float(dy_a)) < 1 else float(dy_a)
        fb = float(dy_b) * 100 if abs(float(dy_b)) < 1 else float(dy_b)
        higher = name_a if fa > fb else name_b
        args.append(("Rendement du dividende",
                      f"{name_a} : {fa:.1f}%  vs  {name_b} : {fb:.1f}%  "
                      f"\u2014  {higher} est plus genereux en dividendes."))

    # Score
    if len(args) < 4:
        winner_s = name_a if sc_a > sc_b else (name_b if sc_b > sc_a else "Egalite")
        args.append(("Score FinSight Composite",
                      f"{name_a} : {sc_a}/100  vs  {name_b} : {sc_b}/100  "
                      f"\u2014  signal global : {winner_s}."))

    return args[:4]


# ── Generateur de textes analytiques LLM ─────────────────────────────────────

def _generate_indice_llm(d: dict) -> dict:
    """Un seul appel Mistral pour les textes analytiques du pitchbook indice vs indice."""
    try:
        from core.llm_provider import LLMProvider
        llm = LLMProvider(provider="mistral", model="mistral-small-latest")
        name_a = d.get("name_a", "Indice A")
        name_b = d.get("name_b", "Indice B")
        sig_a  = d.get("signal_a", "Neutre")
        sig_b  = d.get("signal_b", "Neutre")
        sc_a   = d.get("score_a", 50)
        sc_b   = d.get("score_b", 50)
        p1a    = d.get("perf_1y_a", 0)
        p1b    = d.get("perf_1y_b", 0)
        vol_a  = d.get("vol_1y_a", 0)
        vol_b  = d.get("vol_1y_b", 0)
        sha_a  = d.get("sharpe_1y_a", 0)
        sha_b  = d.get("sharpe_1y_b", 0)
        pe_a   = d.get("pe_fwd_a", 0)
        pe_b   = d.get("pe_fwd_b", 0)
        dy_a   = d.get("div_yield_a", 0)
        dy_b   = d.get("div_yield_b", 0)
        prompt = (
            f"Tu es analyste financier senior. Redige des textes courts pour un pitchbook comparatif "
            f"entre deux indices : {name_a} vs {name_b}.\n\n"
            f"Donnees cles :\n"
            f"- {name_a} : Score {sc_a}/100, Signal {sig_a}, Perf.1Y {float(p1a or 0)*100:.1f}%, "
            f"Vol {float(vol_a or 0):.1f}%, Sharpe {float(sha_a or 0):.2f}, "
            f"P/E Fwd {float(pe_a or 0):.1f}x, DivYield {float(dy_a or 0)*100:.1f}%\n"
            f"- {name_b} : Score {sc_b}/100, Signal {sig_b}, Perf.1Y {float(p1b or 0)*100:.1f}%, "
            f"Vol {float(vol_b or 0):.1f}%, Sharpe {float(sha_b or 0):.2f}, "
            f"P/E Fwd {float(pe_b or 0):.1f}x, DivYield {float(dy_b or 0)*100:.1f}%\n\n"
            f"Reponds UNIQUEMENT en JSON valide. Textes en francais sans accents. Max 250 caracteres par champ.\n"
            f'{{\n'
            f'  "overview_read": "Synthese 2 phrases : performance, valorisation et signal global",\n'
            f'  "risque_read": "Lecture risque 2 phrases : volatilite, Sharpe et profil de risque",\n'
            f'  "verdict_read": "Recommandation finale 2-3 phrases : que privilegier et conditions"\n'
            f'}}'
        )
        import json, re
        resp = llm.generate(prompt, max_tokens=400)
        m = re.search(r'\{.*\}', resp, re.DOTALL)
        if m:
            data_out = json.loads(m.group(0))
            log.info("[indice_cmp_pptx] LLM texts OK (%d champs)", len(data_out))
            return data_out
    except Exception as e:
        log.warning("[indice_cmp_pptx] LLM generation failed: %s", e)
    return {}


# ── Classe principale ─────────────────────────────────────────────────────────

class IndiceComparisonPPTXWriter:

    @staticmethod
    def generate(data: dict, output_path: str = None) -> bytes:
        """Genere le PPTX comparatif et retourne les bytes."""
        # Generer les textes LLM une seule fois
        data = dict(data)   # copie pour ne pas muter l'original
        data["llm"] = _generate_indice_llm(data)

        prs = Presentation()
        prs.slide_width  = _SW
        prs.slide_height = _SH

        # Slide 1 — Cover
        _slide_cover(prs, data)
        # Slide 2 — Vue d'ensemble
        _slide_overview(prs, data)
        # Slide 3 — Sommaire
        _slide_sommaire(prs, data)
        # Slide 4 — Divider 01 Performance
        _divider(prs, "01", "Performance Historique",
                 "Evolution comparee sur YTD, 1 an, 3 ans et 5 ans")
        # Slide 5 — Perf chart
        _slide_perf_chart(prs, data)
        # Slide 6 — Divider 02 Risque
        _divider(prs, "02", "Risque Comparatif",
                 "Volatilite, Sharpe ratio et Max Drawdown")
        # Slide 7 — Risque
        _slide_risque(prs, data)
        # Slide 8 — Divider 03 Valorisation
        _divider(prs, "03", "Valorisation",
                 "P/E Forward, P/B, Rendement du dividende et ERP")
        # Slide 9 — Valorisation
        _slide_valorisation(prs, data)
        # Slide 10 — Divider 04 Secteurs
        _divider(prs, "04", "Composition Sectorielle",
                 "Poids GICS par secteur — ecarts entre les deux indices")
        # Slide 11 — Secteurs
        _slide_secteurs(prs, data)
        # Slide 12 — Divider 05 Constituants
        _divider(prs, "05", "Principaux Constituants",
                 f"Top 5 valeurs de chaque indice")
        # Slide 13 — Top 5 Indice A
        _slide_top5(prs, data, "a")
        # Slide 14 — Top 5 Indice B
        _slide_top5(prs, data, "b")
        # Slide 15 — Divider 06 Verdict
        _divider(prs, "06", "Verdict Final",
                 "Signal FinSight, score compare et recommandation")
        # Slide 16 — Score
        _slide_score(prs, data)
        # Slide 17 — Verdict
        _slide_verdict(prs, data)

        buf = io.BytesIO()
        prs.save(buf)
        pptx_bytes = buf.getvalue()

        if output_path:
            from pathlib import Path
            Path(output_path).write_bytes(pptx_bytes)

        return pptx_bytes
