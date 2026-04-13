# =============================================================================
# FinSight IA — PPTX Writer (Pitchbook 20 slides)
# outputs/pptx_writer.py
#
# Généré un pitchbook IB professionnel en 20 slides via python-pptx.
# Utilise le FinSightState (dict) produit par le pipeline LangGraph.
#
# Usage :
#   writer = PPTXWriter()
#   path   = writer.generate(state, output_path)
# =============================================================================

from __future__ import annotations

import io
import logging
import math
import statistics
from datetime import date, datetime
from pathlib import Path
from typing import Any, Dict, List, Optional

log = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Palette couleurs
# ---------------------------------------------------------------------------

NAVY       = "1B3A6B"
NAVY_MID   = "2E5FA3"
NAVY_PALE  = "EEF3FA"
GREEN      = "1A7A4A"
GREEN_PALE = "EAF4EF"
RED        = "A82020"
RED_PALE   = "FAF0EF"
AMBER      = "B06000"
AMBER_PALE = "FDF6E8"
WHITE      = "FFFFFF"
GREY_BG    = "F7F8FA"
GREY_TXT   = "555555"
GREY_LIGHT = "888888"
BLACK      = "0D0D0D"


# ---------------------------------------------------------------------------
# Imports pptx (lazy-guardés dans PPTXWriter.generate)
# ---------------------------------------------------------------------------

def _lazy_imports():
    from pptx import Presentation
    from pptx.util import Cm, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    return Presentation, Cm, Pt, Emu, RGBColor, PP_ALIGN


# ---------------------------------------------------------------------------
# Sanitisation XML — supprimé les caractères interdits en XML 1.0
# Source courante : texte LLM ou API contenant des chars de contrôle
# ---------------------------------------------------------------------------

import re as _re

def _x(text) -> str:
    """Supprimé les caractères invalides XML 1.0 (0x00-0x08, 0x0B-0x0C, 0x0E-0x1F, 0x7F)."""
    if text is None:
        return ""
    s = str(text)
    return _re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', s)


# ---------------------------------------------------------------------------
# Helpers couleur
# ---------------------------------------------------------------------------

def rgb(hex_str: str):
    from pptx.dml.color import RGBColor
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ---------------------------------------------------------------------------
# Helpers mise en forme des nombres
# ---------------------------------------------------------------------------

def _fr(v, dp: int = 1, suffix: str = "") -> str:
    if v is None:
        return "—"
    try:
        return f"{float(v):,.{dp}f}".replace(",", " ").replace(".", ",") + suffix
    except Exception:
        return "—"


def _frm(v, cur_sym: str = "$") -> str:
    if v is None:
        return "—"
    try:
        v = float(v)
        # Guard : LLM peut envoyer des valeurs en EUR/USD absolus (ex: 195_000_000_000)
        # au lieu de millions (195_000). Si > 1e9, on Normalisé en millions.
        if abs(v) > 1_000_000_000:
            v = v / 1_000_000
        if cur_sym == "EUR":
            sym_big, sym_small = "Md\u20ac", "M\u20ac"
        else:
            sym_big, sym_small = "Mds" + cur_sym, "M" + cur_sym
        if abs(v) >= 1000:
            return _fr(v / 1000, 1) + " " + sym_big
        return _fr(v, 0) + " " + sym_small
    except Exception:
        return "—"


def _frn(v) -> str:
    """Formate un montant en millions -> toujours converti en Mds avec 1 décimale."""
    if v is None:
        return "—"
    try:
        v = float(v)
        if abs(v) > 1_000_000_000:  # guard LLM valeurs absolues
            v = v / 1_000_000
        return _fr(v / 1000, 1)
    except Exception:
        return "—"


def _frpct(v, signed: bool = False) -> str:
    if v is None:
        return "—"
    try:
        s = f"{float(v) * 100:+.1f}" if signed else f"{float(v) * 100:.1f}"
        return s.replace(".", ",") + " %"
    except Exception:
        return "—"


def _frpct_margin(v) -> str:
    """Like _frpct but guards against values already in percentage form (e.g. 65.0 vs 0.65).
    Used for gross_margin / ebitda_margin from inconsistent data sources."""
    if v is None:
        return "—"
    try:
        fv = float(v)
        # If |value| > 2.0 (i.e. > 200%), assume already-percentage form — divide by 100.
        # No real gross/EBITDA margin can exceed 200%, so threshold is safe.
        if abs(fv) > 2.0:
            fv = fv / 100.0
        s = f"{fv * 100:.1f}"
        return s.replace(".", ",") + " %"
    except Exception:
        return "—"


def _frx(v) -> str:
    if v is None:
        return "—"
    try:
        return f"{float(v):.1f}".replace(".", ",") + "x"
    except Exception:
        return "—"


def _upside(target, current) -> str:
    if not target or not current:
        return "—"
    try:
        return f"{(float(target) / float(current) - 1) * 100:+.1f}".replace(".", ",") + " %"
    except Exception:
        return "—"


def _g(obj, attr, default=None):
    if obj is None:
        return default
    if isinstance(obj, dict):
        return obj.get(attr, default)
    return getattr(obj, attr, default)


def _valid_years(snap) -> list:
    """Retourne les années triees qui ont au moins une Donnée réelle (exclut les années all-None)."""
    if not (snap and snap.years):
        return []
    result = []
    for y, fy in snap.years.items():
        if fy is None:
            continue
        if any(getattr(fy, attr, None) is not None
               for attr in ("revenue", "cash", "da", "interest_expense")):
            result.append(y)
    return sorted(result, key=lambda y: str(y).replace("_LTM", ""))


# ---------------------------------------------------------------------------
# Helpers XML (transparence, fond slide)
# ---------------------------------------------------------------------------

def _set_slide_bg(slide, hex_color: str):
    """Definit le fond du slide via <p:bgPr> (approche native, comme la référence)."""
    from pptx.oxml.ns import qn
    from lxml import etree
    cSld = slide._element.find(qn('p:cSld'))
    if cSld is None:
        return
    existing = cSld.find(qn('p:bg'))
    if existing is not None:
        cSld.remove(existing)
    bg_xml = (
        '<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<p:bgPr>'
        f'<a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>'
        '<a:effectLst/>'
        '</p:bgPr>'
        '</p:bg>'
    )
    bg_elem = etree.fromstring(bg_xml)
    spTree = cSld.find(qn('p:spTree'))
    cSld.insert(list(cSld).index(spTree) if spTree is not None else 0, bg_elem)


def _add_text_alpha(slide, x, y, w, h, text, font_size,
                    color_hex="FFFFFF", alpha_val=15000, bold=False):
    """Ajouté un textbox avec couleur semi-transparente (alpha en millièmes, 15000=15%)."""
    from pptx.util import Cm, Pt
    from pptx.oxml.ns import qn
    from lxml import etree

    txBox = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = _x(text)
    run.font.name  = "Calibri"
    run.font.size  = Pt(font_size)
    run.font.bold  = bold

    # Injection XML directe (plus fiable que get_or_add_rPr + string parsing)
    r_elem = run._r                          # <a:r>
    rPr = r_elem.find(qn('a:rPr'))
    if rPr is None:
        rPr = etree.Element(qn('a:rPr'))
        r_elem.insert(0, rPr)

    # Supprimer toute couleur existante
    for tag in (qn('a:solidFill'), qn('a:gradFill'), qn('a:noFill'), qn('a:pattFill')):
        for child in list(rPr.findall(tag)):
            rPr.remove(child)

    # Construire solidFill > srgbClr > alpha via SubElement (évite les pb de namespace)
    solidFill  = etree.SubElement(rPr, qn('a:solidFill'))
    srgbClr    = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgbClr.set('val', color_hex.upper())
    alpha_elem = etree.SubElement(srgbClr, qn('a:alpha'))
    alpha_elem.set('val', str(alpha_val))
    # Mettre solidFill en premier dans rPr
    rPr.remove(solidFill)
    rPr.insert(0, solidFill)

    return txBox


# ---------------------------------------------------------------------------
# Helpers formes pptx
# ---------------------------------------------------------------------------

def add_rect(slide, x, y, w, h, fill_hex, line_hex=None, line_width_pt=0):
    from pptx.util import Cm, Pt
    shape = slide.shapes.add_shape(1, Cm(x), Cm(y), Cm(w), Cm(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_hex)
    if line_hex:
        shape.line.color.rgb = rgb(line_hex)
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, x, y, w, h, text, font_size=10, color_hex=BLACK,
                 bold=False, italic=False, align=None,
                 font_name="Calibri", wrap=True):
    from pptx.util import Cm, Pt
    from pptx.enum.text import PP_ALIGN
    if align is None:
        align = PP_ALIGN.LEFT
    txBox = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = _x(text) if text is not None else "—"
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color_hex)
    return txBox


def add_multiline_text_box(slide, x, y, w, h, lines, font_size=8.5,
                           color_hex=BLACK, font_name="Calibri", wrap=True):
    """Text box avec plusieurs paragraphes (list of str)."""
    from pptx.util import Cm, Pt
    from pptx.enum.text import PP_ALIGN
    txBox = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = _x(line) if line is not None else ""
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = rgb(color_hex)
    return txBox


def _set_cell(cell, text, font_size=8, bold=False, color_hex=BLACK,
              fill_hex=None, align=None, font_name="Calibri", italic=False):
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    if align is None:
        align = PP_ALIGN.CENTER
    if fill_hex:
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(fill_hex)
    # Clear existing content then add styled run
    tf = cell.text_frame
    for para in tf.paragraphs:
        for run in list(para.runs):
            para._p.remove(run._r)
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = _x(text) if text is not None else "\u2014"
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color_hex)


def _apply_cell_borders(tbl, total_rows, num_cols, hex_color="DDDDDD", width_pt=0.5):
    from pptx.oxml.ns import qn
    from lxml import etree
    w_emu = str(int(width_pt * 12700))
    ns = 'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
    ln_xml = (
        f'<a:ln {ns} w="{w_emu}">'
        f'<a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>'
        f'</a:ln>'
    )
    for ri in range(total_rows):
        for ci in range(num_cols):
            tc = tbl.cell(ri, ci)._tc
            tcPr = tc.find(qn('a:tcPr'))
            if tcPr is None:
                tcPr = etree.SubElement(tc, qn('a:tcPr'))
            for tag in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
                old = tcPr.find(qn(tag))
                if old is not None:
                    tcPr.remove(old)
                ln = etree.fromstring(ln_xml)
                sub = etree.SubElement(tcPr, qn(tag))
                sub.append(ln[0])          # append <a:solidFill>
                sub.set('w', w_emu)


def add_table(slide, x, y, w, h, num_rows, num_cols, col_widths_pct=None,
              header_data=None, rows_data=None,
              header_fill=NAVY, row_fills=None, border_hex=None):
    """
    Ajoute un tableau python-pptx.
    header_data : list de strings (optionnel)
    rows_data   : list of list of strings
    col_widths_pct : list of floats (somme=1.0) — proportions colonnes
    row_fills   : list de fill_hex pour chaque ligne data (cycling)
    """
    from pptx.util import Cm, Pt, Emu
    from pptx.enum.text import PP_ALIGN

    total_rows = num_rows + (1 if header_data else 0)
    tbl = slide.shapes.add_table(
        total_rows, num_cols, Cm(x), Cm(y), Cm(w), Cm(h)
    ).table

    # Largeurs colonnes
    total_emu = int(Cm(w))
    if col_widths_pct:
        for ci, pct in enumerate(col_widths_pct):
            tbl.columns[ci].width = int(total_emu * pct)
    else:
        col_w = total_emu // num_cols
        for ci in range(num_cols):
            tbl.columns[ci].width = col_w

    row_offset = 0
    # Header
    if header_data:
        for ci, val in enumerate(header_data):
            cell = tbl.cell(0, ci)
            _set_cell(cell, val, font_size=8, bold=True,
                      color_hex=WHITE, fill_hex=header_fill,
                      align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)
        row_offset = 1

    # Data rows
    if rows_data:
        default_fills = [WHITE, GREY_BG]
        for ri, row in enumerate(rows_data):
            if row_fills and ri < len(row_fills):
                fill = row_fills[ri]
            else:
                fill = default_fills[ri % 2]
            for ci, val in enumerate(row):
                cell = tbl.cell(ri + row_offset, ci)
                align = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
                _set_cell(cell, val, font_size=8, bold=False,
                          color_hex=BLACK, fill_hex=fill,
                          align=align)

    if border_hex:
        _apply_cell_borders(tbl, num_rows + (1 if header_data else 0), num_cols, border_hex)

    return tbl


# ---------------------------------------------------------------------------
# Elements communs des slides
# ---------------------------------------------------------------------------

def navy_bar(slide):
    add_rect(slide, 0, 0, 25.4, 1.65, NAVY)


def footer_bar(slide, ticker="", co_name=""):
    add_rect(slide, 0, 13.39, 25.4, 0.89, NAVY)
    add_text_box(slide, 1.02, 13.44, 23.37, 0.56,
                 "FinSight IA  \u00b7  Usage confidentiel", 7, WHITE)


def section_dots(slide, active: int):
    from pptx.enum.text import PP_ALIGN
    xs = [17.02, 17.65, 18.29, 18.92, 19.56]
    for i, x in enumerate(xs):
        num = i + 1
        fill = WHITE if num == active else NAVY_MID
        add_rect(slide, x, 0.33, 0.51, 0.51, fill)
        color = NAVY if num == active else WHITE
        add_text_box(slide, x + 0.02, 0.34, 0.47, 0.44,
                     str(num), 7, color, bold=True,
                     align=PP_ALIGN.CENTER)


def slide_title(slide, title: str, subtitle: str = ""):
    from pptx.enum.text import PP_ALIGN
    add_text_box(slide, 1.02, 0.33, 15.49, 0.97,
                 title, 13, WHITE, bold=True)
    if subtitle:
        add_text_box(slide, 1.02, 1.73, 23.37, 0.56,
                     subtitle, 9, NAVY_MID)


def kpi_box(slide, x, y, w, h, value, label, sub="",
            fill=NAVY_PALE, accent=NAVY_MID):
    from pptx.enum.text import PP_ALIGN
    add_rect(slide, x, y, w, h, fill)
    add_rect(slide, x, y, 0.13, h, accent)
    val_color = NAVY if fill != NAVY else WHITE
    lbl_color = GREY_TXT if fill != NAVY else WHITE
    sub_color = GREY_LIGHT if fill != NAVY else NAVY_PALE
    add_text_box(slide, x + 0.25, y + 0.15, w - 0.38, 1.1,
                 value, 19, val_color, bold=True)
    add_text_box(slide, x + 0.25, y + 1.25, w - 0.38, 0.57,
                 label, 8, lbl_color)
    if sub:
        add_text_box(slide, x + 0.25, y + 1.8, w - 0.38, 0.46,
                     sub, 7, sub_color)


def commentary_box(slide, x, y, w, h, text, accent=NAVY_MID, title=None):
    """Box commentaire avec titre analytique optionnel style JPM/GS.

    Si `title` fourni : barre titre navy en haut + corps texte en dessous.
    Sinon : juste le corps de texte (mode legacy).
    """
    add_rect(slide, x, y, w, h, "F5F7FA")           # fond gris clair
    add_rect(slide, x, y, 0.22, h, accent)           # barre accent gauche
    if title:
        # Titre analytique JPM/GS en haut
        add_text_box(slide, x + 0.30, y + 0.08, w - 0.32, 0.42,
                     title, 9.5, NAVY, bold=True, wrap=False)
        add_text_box(slide, x + 0.30, y + 0.55, w - 0.32, max(h - 0.62, 0.40),
                     text or "\u2014", 8.5, GREY_TXT, wrap=True)
    else:
        add_text_box(slide, x + 0.30, y + 0.08, w - 0.32, max(h - 0.10, 0.40),
                     text or "\u2014", 8.5, GREY_TXT, wrap=True)


def _jpm_title(slide_type: str, ratios=None, snap=None, synthesis=None,
               extra: dict = None) -> str:
    """Génère un titre analytique style JPM/GS dynamique basé sur les données.

    slide_type : "is" (income statement), "bilan", "ratios", "dcf", "peers",
                 "football", "lbo", "risques", "score", "verdict", etc.
    Retourne une string courte (< 100 chars) qui résume l'insight clé de la slide.
    """
    extra = extra or {}
    try:
        # Helpers pour extraction
        ci = snap.company_info if snap else None
        ticker = _g(ci, "ticker", "Société") or "Société"
        latest_yr = None
        if ratios and getattr(ratios, "years", None):
            _lbl = sorted(ratios.years.keys(),
                          key=lambda k: str(k).replace("_LTM", "Z"))[-1]
            latest_yr = ratios.years[_lbl]

        def _gv(field):
            return getattr(latest_yr, field, None) if latest_yr else None

        rev_g = _gv("revenue_growth")
        eb_marg = _gv("ebitda_margin")
        nm = _gv("net_margin")
        roe = _gv("roe")
        roic = _gv("roic")
        nd_eb = _gv("net_debt_ebitda")

        # Conversion en %
        def _pct(v):
            if v is None:
                return None
            try:
                fv = float(v)
                return fv * 100 if abs(fv) <= 2 else fv
            except Exception:
                return None

        rev_g_pct = _pct(rev_g)
        eb_marg_pct = _pct(eb_marg)
        nm_pct = _pct(nm)
        roe_pct = _pct(roe)
        roic_pct = _pct(roic)

        if slide_type == "is":
            # Compte de résultat : focus revenue + marges
            if rev_g_pct is not None and eb_marg_pct is not None:
                if rev_g_pct >= 10 and eb_marg_pct >= 20:
                    return f"{ticker} combine forte croissance ({rev_g_pct:+.1f}%) et marges premium ({eb_marg_pct:.1f}%)"
                elif rev_g_pct >= 5:
                    return f"{ticker} en croissance soutenue ({rev_g_pct:+.1f}% revenue), marge EBITDA {eb_marg_pct:.1f}%"
                elif rev_g_pct < 0:
                    return f"{ticker} en contraction ({rev_g_pct:+.1f}% revenue) — pression sur les marges ({eb_marg_pct:.1f}%)"
                else:
                    return f"{ticker} en plateau ({rev_g_pct:+.1f}% rev), marges stables {eb_marg_pct:.1f}%"
            return f"{ticker} — analyse compte de résultat et tendances de marges"

        if slide_type == "bilan":
            # Bilan : focus levier + liquidité
            if nd_eb is not None:
                try:
                    nd_v = float(nd_eb)
                    if nd_v < 0:
                        return f"{ticker} en cash net positif — bilan robuste, capacité d'investissement intacte"
                    if nd_v < 1:
                        return f"{ticker} faiblement endetté (ND/EBITDA {nd_v:.1f}x) — solidité bilancielle"
                    if nd_v < 3:
                        return f"{ticker} levier modéré (ND/EBITDA {nd_v:.1f}x), structure équilibrée"
                    return f"{ticker} levier élevé (ND/EBITDA {nd_v:.1f}x) — surveiller la capacité de service de la dette"
                except Exception:
                    pass
            return f"{ticker} — analyse de la structure bilancielle"

        if slide_type == "ratios":
            if roe_pct is not None and roic_pct is not None:
                if roe_pct >= 20 and roic_pct >= 15:
                    return f"{ticker} génère ROE {roe_pct:.0f}% / ROIC {roic_pct:.0f}% — création de valeur premium"
                elif roe_pct >= 15:
                    return f"{ticker} ROE {roe_pct:.0f}% — rentabilité solide vs cost of equity"
                else:
                    return f"{ticker} ROE {roe_pct:.0f}% / ROIC {roic_pct:.0f}% — rentabilité à surveiller"
            return f"{ticker} — analyse rentabilité et création de valeur"

        if slide_type == "dcf":
            tbase = _g(synthesis, "target_base") if synthesis else None
            mkt = snap.market if snap else None
            price = _g(mkt, "share_price") if mkt else None
            if tbase and price:
                try:
                    upside = (float(tbase) / float(price) - 1) * 100
                    if upside >= 20:
                        return f"DCF base implique upside de {upside:+.0f}% — valorisation attractive"
                    elif upside >= 0:
                        return f"DCF base proche du cours ({upside:+.0f}%) — valorisation alignée"
                    else:
                        return f"DCF base en dessous du cours ({upside:+.0f}%) — surévaluation potentielle"
                except Exception:
                    pass
            return f"{ticker} — valorisation DCF et sensibilités"

        if slide_type == "peers":
            return f"{ticker} — positionnément vs comparables sectoriels"

        if slide_type == "football":
            return f"{ticker} — football field et fourchette de valorisation"

        if slide_type == "lbo":
            irr = extra.get("irr_base")
            if irr is not None:
                try:
                    irr_v = float(irr) * 100
                    if irr_v >= 20:
                        return f"LBO {ticker} : IRR {irr_v:.1f}% — deal attractif pour sponsor PE top-tier"
                    elif irr_v >= 15:
                        return f"LBO {ticker} : IRR {irr_v:.1f}% — viable pour fonds mid-market"
                    elif irr_v >= 0:
                        return f"LBO {ticker} : IRR {irr_v:.1f}% — sous le seuil PE typique"
                    else:
                        return f"LBO {ticker} : IRR négatif — deal non viable dans les conditions actuelles"
                except Exception:
                    pass
            return f"{ticker} — analyse LBO et returns sponsor"

        if slide_type == "risques":
            return f"{ticker} — risques structurels et conditions d'invalidation"

        if slide_type == "score":
            score = _g(synthesis, "finsight_score") if synthesis else None
            if score:
                try:
                    s = int(score)
                    if s >= 75:
                        return f"{ticker} — Score FinSight {s}/100 — signal d'achat fort"
                    elif s >= 55:
                        return f"{ticker} — Score FinSight {s}/100 — signal modéré favorable"
                    else:
                        return f"{ticker} — Score FinSight {s}/100 — signal de prudence"
                except Exception:
                    pass
            return f"{ticker} — Score FinSight composite"

        if slide_type == "verdict":
            rec = _g(synthesis, "recommendation") if synthesis else None
            if rec:
                return f"{ticker} — verdict final : {str(rec).upper()}"
            return f"{ticker} — synthèse et recommandation finale"

        if slide_type == "capital":
            return f"{ticker} — politique de retour au capital actionnaire"

    except Exception:
        pass

    return ""


def divider_slide(prs, number_str: str, title: str, subtitle: str):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Fond navy via <p:bgPr> (natif, comme la référence — pas de rectangle additionnel)
    _set_slide_bg(slide, NAVY)

    # Left accent bar
    add_rect(slide, 0, 0, 0.3, 14.29, NAVY_MID)

    # Numéro filigrane — blanc à 15 % d'opacité (comme la référence)
    _add_text_alpha(slide, 1.27, 4.06, 22.86, 4.57,
                    number_str, 80, "FFFFFF", 15000, bold=True)

    # Title
    add_text_box(slide, 1.27, 5.21, 21.59, 2.29,
                 title, 28, WHITE, bold=True)

    # Rule
    add_rect(slide, 1.27, 7.75, 7.62, 0.01, GREY_LIGHT)

    # Subtitle
    add_text_box(slide, 1.27, 8.08, 17.78, 0.89,
                 subtitle, 10, "AABBDD")

    # Footer (espaces autour du bullet comme la référence)
    add_text_box(slide, 1.02, 13.34, 23.37, 0.56,
                 "FinSight IA  \u00b7  Usage confidentiel", 7, "6677AA")

    return slide


# ---------------------------------------------------------------------------
# State data extraction helpers
# ---------------------------------------------------------------------------

def _extract_state(state: dict):
    snap      = state.get("raw_data")
    synthesis = state.get("synthesis")
    ratios    = state.get("ratios")
    devil     = state.get("devil")
    sentiment = state.get("sentiment")
    return snap, synthesis, ratios, devil, sentiment


def _ratio(ratios, attr):
    if not ratios:
        return None
    latest = getattr(ratios, "latest_year", None)
    if not latest:
        return None
    yrs = getattr(ratios, "years", {}) or {}
    yr = yrs.get(latest)
    val = getattr(yr, attr, None) if yr else None
    if val is not None:
        return val
    # Fallback: walk years in reverse order to find latest with data
    for y in sorted(yrs.keys(), key=lambda k: str(k).replace("_LTM", ""), reverse=True):
        yr2 = yrs.get(y)
        if yr2 is None:
            continue
        v2 = getattr(yr2, attr, None)
        if v2 is not None:
            return v2
    return None


def _fy(snap, year_key, attr):
    if not snap or not year_key:
        return None
    fy = snap.years.get(year_key) if snap.years else None
    return getattr(fy, attr, None) if fy else None


def _rec_colors(rec: str):
    rec = (rec or "HOLD").upper()
    if rec == "BUY":
        return GREEN_PALE, GREEN
    elif rec == "SELL":
        return RED_PALE, RED
    else:
        return AMBER_PALE, AMBER


def _sent_label_fr(label: str, score: float = 0.0) -> str:
    label = (label or "NEUTRAL").upper()
    if label in ("POSITIVE", "POSITIF"):
        return "Positif"
    elif label in ("NEGATIVE", "NEGATIF"):
        return "Négatif"
    else:
        return "Neutre +" if score > 0 else "Neutre"


def _safe_str(v, default="—") -> str:
    if v is None:
        return default
    s = str(v).strip()
    return s if s else default


def _fr_date_long(d=None) -> str:
    """Converts a date to French long form: '14 mars 2026'."""
    from datetime import date as _date
    _MONTHS = ["janvier", "f\u00e9vrier", "mars", "avril", "mai", "juin",
               "juillet", "ao\u00fbt", "septembre", "octobre", "novembre", "d\u00e9cembre"]
    if d is None:
        d = _date.today()
    elif isinstance(d, str):
        try:
            d = _date.fromisoformat(str(d).split("T")[0][:10])
        except Exception:
            return str(d)
    try:
        return f"{d.day} {_MONTHS[d.month - 1]} {d.year}"
    except Exception:
        return str(d)


_EXCHANGE_MAP = {
    "NMS": "NASDAQ", "NGM": "NASDAQ", "NCM": "NASDAQ", "NAS": "NASDAQ",
    "NYQ": "NYSE",   "NYA": "NYSE",
    "PCX": "NYSE Arca",
    "XPAR": "Euronext Paris", "PAR": "Euronext Paris", "ENX": "Euronext Paris",
    "XLON": "London SE",      "LSE": "London SE",       "IOB": "London SE",
    "XFRA": "Frankfurt",      "FRA": "Frankfurt",        "GER": "Frankfurt",
    "XAMS": "Euronext Amsterdam", "AMS": "Euronext Amsterdam",
    "XBRU": "Euronext Brussels",  "BRU": "Euronext Brussels",
    "XMIL": "Borsa Italiana",     "MIL": "Borsa Italiana",
    "XMAD": "BME",                "MCE": "BME",
    "TDM": "TSX",                 "TSX": "TSX",
}

def _normalize_exchange(ex: str) -> str:
    if not ex:
        return ""
    return _EXCHANGE_MAP.get(ex.upper().strip(), ex)


_TICKER_SUFFIX_MAP = {
    ".PA": "Euronext Paris",  ".L": "London SE",   ".DE": "Frankfurt",
    ".AS": "Euronext Amsterdam", ".BR": "Euronext Brussels",
    ".MI": "Borsa Italiana",  ".MC": "BME",         ".SW": "SIX Swiss",
    ".TO": "TSX",             ".HK": "HKEX",        ".T":  "Tokyo SE",
    ".AX": "ASX",             ".CO": "Nasdaq Copenhagen",
}

def _infer_exchange(ticker: str, exchange: str) -> str:
    """Retourne l'exchange normalisé, avec fallback sur le suffixe du ticker."""
    ex = _normalize_exchange(exchange)
    if ex:
        return ex
    for sfx, exch in _TICKER_SUFFIX_MAP.items():
        if ticker.upper().endswith(sfx.upper()):
            return exch
    return ""


import math as _math

def _cover_layout(co_name: str):
    """
    Calcule dynamiquement la taille de police et les positions y
    pour que le nom, la tagline et la boîte de recommandation
    ne se chevauchent jamais, quelle que soit la longueur du nom.
    """
    n = len(co_name)
    if n <= 15:   fs, cpl = 40, 22
    elif n <= 25: fs, cpl = 34, 28
    elif n <= 40: fs, cpl = 26, 35
    elif n <= 60: fs, cpl = 20, 44
    else:          fs, cpl = 16, 55
    n_lines   = max(1, min(_math.ceil(n / cpl), 5))
    line_h_cm = fs * 0.0353 * 1.35          # pt → cm, + interligne
    name_h    = max(n_lines * line_h_cm + 0.3, 1.4)
    name_y    = 4.06
    tagline_y = name_y + name_h + 0.25
    rec_y     = tagline_y + 0.76 + 0.28
    return fs, name_y, name_h, tagline_y, rec_y


def _truncate(s, n: int) -> str:
    s = _safe_str(s)
    if len(s) <= n:
        return s
    # Cut at last word boundary before n to avoid mid-word truncation
    cut = s[:n]
    last_space = cut.rfind(" ")
    if last_space > n // 2:
        cut = cut[:last_space]
    return cut + "…"


def _fit(s, n: int) -> str:
    """Coupe a n chars sans ajouter '...' — pour les zones ou le débordement est invisible."""
    s = _safe_str(s)
    if len(s) <= n:
        return s
    cut = s[:n]
    last_space = cut.rfind(" ")
    return cut[:last_space] if last_space > n // 2 else cut


def _peer_median(peers: list, attr: str):
    vals = []
    for p in (peers or []):
        v = _g(p, attr)
        if v is not None:
            try:
                vals.append(float(v))
            except Exception:
                pass
    if not vals:
        return None
    return statistics.median(vals)


# ---------------------------------------------------------------------------
# DCF sensitivity helper
# ---------------------------------------------------------------------------

def _dcf_value(tbase, wacc_base, tgr_base, w_delta, t_delta):
    """
    Approximation lineaire de la valeur DCF autour du point de base.
    En l'absence de FCF explicites, on scale le prix cible base
    par le ratio (wacc_ref / wacc_new) * (tgr_new / tgr_ref) simplifie.
    """
    if tbase is None:
        return "—"
    try:
        w = wacc_base + w_delta
        t = tgr_base + t_delta
        if w <= t or w <= 0:
            return "—"
        # Gordon Growth proxy: V ∝ 1 / (w - t)
        base_spread = wacc_base - tgr_base
        new_spread  = w - t
        if base_spread <= 0 or new_spread <= 0:
            return "—"
        val = float(tbase) * base_spread / new_spread
        return _fr(val, 0)
    except Exception:
        return "—"


# ---------------------------------------------------------------------------
# Slide 1 — Cover
# ---------------------------------------------------------------------------

def _slide_cover(prs, snap, synthesis, ratios, devil, sentiment):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    ci       = snap.company_info if snap else None
    mkt      = snap.market if snap else None
    ticker   = _g(ci, "ticker", "—")
    co_name  = _g(ci, "company_name", "—")
    sector   = _g(ci, "sector", "") or ""
    exchange = _infer_exchange(ticker, getattr(ci, "exchange", "") or "" if ci else "")
    currency = _g(ci, "currency", "USD") or "USD"
    cur_sym  = "EUR" if currency == "EUR" else "$"
    gen_date = _fr_date_long(_g(ci, "analysis_date", None) or date.today())
    price    = _g(mkt, "share_price")

    rec      = _g(synthesis, "recommendation", "HOLD") or "HOLD"
    tbase    = _g(synthesis, "target_base")
    rec_fill, rec_accent = _rec_colors(rec)

    # Layout dynamique selon longueur du nom
    fs_name, name_y, name_h, tagline_y, rec_y = _cover_layout(co_name)

    # Top navy band
    add_rect(slide, 0, 0, 25.4, 3.81, NAVY)
    add_text_box(slide, 1.27, 0.46, 22.86, 0.71, "FinSight IA",
                 10, "8899BB", bold=False, align=PP_ALIGN.CENTER)
    add_rect(slide, 11.05, 1.32, 3.3, 0.05, WHITE)
    add_text_box(slide, 1.27, 1.52, 22.86, 0.81,
                 "Pitchbook  \u2014  Analyse d'investissement",
                 11, "CCDDEE", align=PP_ALIGN.CENTER)

    # Company name (taille dynamique, centré)
    add_text_box(slide, 1.27, name_y, 22.86, name_h,
                 co_name, fs_name, NAVY, bold=True, align=PP_ALIGN.CENTER)

    # Tagline ticker · exchange · secteur
    _tagline_parts = [p for p in [ticker, exchange, sector] if p and str(p).strip()]
    add_text_box(slide, 1.27, tagline_y, 22.86, 0.71,
                 "  \u00b7  ".join(_tagline_parts), 11, "888888",
                 align=PP_ALIGN.CENTER)

    # Recommendation box — rectangle solide centré ~40% largeur (gold standard)
    upside_str = _upside(tbase, price)
    rec_w = 10.0   # ~40% de 25.4cm
    rec_h = 1.20
    rec_x = (25.4 - rec_w) / 2
    add_rect(slide, rec_x, rec_y, rec_w, rec_h, rec_accent)
    add_text_box(
        slide, rec_x + 0.15, rec_y + 0.05, rec_w - 0.30, rec_h - 0.10,
        f"\u25cf {rec.upper()}  \u00b7  Prix cible base\u00a0: {_fr(tbase, 0)} {cur_sym}"
        f"  \u00b7  Upside\u00a0: {upside_str}",
        9, WHITE, bold=True, align=PP_ALIGN.CENTER
    )

    # Bottom rule + crédits (8pt)
    add_rect(slide, 1.02, 12.4, 23.37, 0.03, "AAAAAA")
    add_text_box(slide, 1.02, 12.65, 11.43, 0.56,
                 "Rapport confidentiel", 8, GREY_TXT)
    add_text_box(slide, 12.95, 12.65, 11.43, 0.56,
                 gen_date, 8, GREY_TXT, align=PP_ALIGN.RIGHT)

    return slide


# ---------------------------------------------------------------------------
# Slide 2 — Executive Summary
# ---------------------------------------------------------------------------

def _slide_exec_summary(prs, snap, synthesis, ratios, devil, sentiment):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    navy_bar(slide)
    footer_bar(slide)
    slide_title(slide, "Executive Summary")

    ci      = snap.company_info if snap else None
    mkt     = snap.market if snap else None
    currency = _g(ci, "currency", "USD") or "USD"
    cur_sym  = "EUR" if currency == "EUR" else "$"
    price    = _g(mkt, "share_price")

    rec      = _g(synthesis, "recommendation", "HOLD") or "HOLD"
    conv     = _g(synthesis, "conviction", 0.5) or 0.5
    conf     = _g(synthesis, "confidence_score", 0.5) or 0.5
    tbase    = _g(synthesis, "target_base")
    tbear    = _g(synthesis, "target_bear")
    tbull    = _g(synthesis, "target_bull")
    thesis   = _g(synthesis, "thesis", "") or ""
    strengths= _g(synthesis, "strengths", []) or []
    risks_s  = _g(synthesis, "risks", []) or []

    rec_fill, rec_accent = _rec_colors(rec)

    beta     = _g(mkt, "beta_levered")
    rfr      = _g(mkt, "risk_free_rate") or 0.041
    wacc_val = _g(mkt, "wacc") or 0.10

    sent_score  = _g(sentiment, "score", 0.0) or 0.0
    sent_label  = _g(sentiment, "label", "NEUTRAL") or "NEUTRAL"
    sent_articles = _g(sentiment, "articles_analyzed", 0) or 0
    sent_label_display = _sent_label_fr(sent_label, sent_score)
    _s2_is_llm  = _g(sentiment, "engine", "finbert") not in ("finbert", "finbert_fallback", None, "")

    # Rec band
    add_rect(slide, 1.02, 2.08, 23.37, 1.22, rec_fill)
    add_rect(slide, 1.02, 2.08, 0.13, 1.22, rec_accent)
    add_text_box(slide, 1.4, 2.13, 3.05, 1.12,
                 f"● {rec.upper()}", 11, rec_accent, bold=True)
    bull_str = (
        f"Bear : {_fr(tbear, 0)} {cur_sym} ({_upside(tbear, price)})  ·  "
        f"Bull : {_fr(tbull, 0)} {cur_sym} ({_upside(tbull, price)})"
    )
    add_text_box(
        slide, 4.7, 2.13, 19.3, 1.12,
        f"Prix cible base : {_fr(tbase, 0)} {cur_sym}  ·  Upside : {_upside(tbase, price)}  ·  {bull_str}",
        9, NAVY
    )

    # Thesis section header
    add_rect(slide, 1.02, 3.76, 11.56, 0.71, NAVY)
    add_text_box(slide, 1.02, 3.76, 11.56, 0.71,
                 "TH\u00c8SE D'INVESTISSEMENT", 9, WHITE, bold=True)

    # Thesis bullets — espacement dynamique avec corps de texte
    thesis_parts = _split_text(thesis, 3)
    pos_themes   = _g(synthesis, "positive_themes", []) or []
    strength_ys = [4.57, 6.10, 7.63]
    for i, sy in enumerate(strength_ys):
        label = strengths[i] if i < len(strengths) else (thesis_parts[i] if i < len(thesis_parts) else "")
        body  = thesis_parts[i] if i < len(thesis_parts) else ""
        # Fallback body : positive_themes si thesis_parts trop court
        if not body or body == label:
            body = pos_themes[i] if i < len(pos_themes) else ""
        add_rect(slide, 1.02, sy + 0.05, 0.15, 0.36, NAVY_MID)
        add_text_box(slide, 1.4, sy, 10.92, 0.51,
                     _truncate(label, 80), 8.5, NAVY, bold=True)
        add_text_box(slide, 1.4, sy + 0.47, 10.92, 1.05,
                     _fit(body, 190), 7.5, "333333", wrap=True)

    # Risks section header
    add_rect(slide, 13.08, 3.76, 11.3, 0.71, RED)
    add_text_box(slide, 13.08, 3.76, 11.3, 0.71,
                 "RISQUES PRINCIPAUX", 9, WHITE, bold=True)

    # Risques avec corps de texte
    counter_risks = _g(devil, "counter_risks", []) or []
    counter_thesis_txt = _g(devil, "counter_thesis", "") or ""
    risk_bodies  = _split_text(counter_thesis_txt, 3) if counter_thesis_txt else [""] * 3
    neg_themes   = _g(synthesis, "negative_themes", []) or []

    risk_ys = [4.57, 6.10, 7.63]
    for i, ry in enumerate(risk_ys):
        risk_text = (risks_s[i] if i < len(risks_s) else
                     (counter_risks[i] if i < len(counter_risks) else ""))
        # Cascade body : neg_themes → risk_bodies (counter_thesis) → counter_risks → fallback unique
        # Seuil abaissé à 15 chars : neg_themes peut être un court bullet point LLM
        body_r = (neg_themes[i] if i < len(neg_themes) and neg_themes[i] else "")
        if not body_r or len(body_r) < 15:
            rb = risk_bodies[i] if i < len(risk_bodies) else ""
            if rb and len(rb) >= 15:
                body_r = rb
        if not body_r or len(body_r) < 15:
            cr = counter_risks[i] if i < len(counter_risks) else ""
            if cr and len(cr) >= 15:
                body_r = cr
        if not body_r or len(body_r) < 15:
            # Fallback unique : intègre le label du risque pour différencier les 3 blocs
            _lbl = (risk_text.split(":", 1)[-1].strip() if ":" in risk_text
                    else risk_text)
            if _lbl and len(_lbl) >= 10:
                body_r = (f"{_lbl} — impact potentiel sur la valorisation et les marges, "
                          f"surveillance recommandée sur les 6-12 prochains mois.")
            else:
                body_r = ("Impact potentiel sur la thèse d'investissement — "
                          "surveillance recommandée sur les 6-12 prochains mois.")
        add_rect(slide, 13.08, ry + 0.05, 0.15, 0.36, RED)
        add_text_box(slide, 13.46, ry, 10.54, 0.51,
                     _truncate(risk_text, 80), 8.5, NAVY, bold=True)
        add_text_box(slide, 13.46, ry + 0.47, 10.54, 1.05,
                     _fit(body_r, 190), 7.5, "333333", wrap=True)

    # Vertical divider
    add_rect(slide, 12.57, 3.76, 0.03, 4.84, GREY_LIGHT)

    # Horizontal rule
    add_rect(slide, 1.02, 9.17, 23.37, 0.03, "AAAAAA")

    # -------------------------------------------------------------------------
    # Investment Case — Catalyseurs sur toute la largeur (Valorisation Synthétique
    # supprimée — redondante avec slides DCF/Multiples qui suivent)
    # -------------------------------------------------------------------------
    catalysts   = _g(synthesis, "catalysts", []) or []
    peers       = _g(synthesis, "comparable_peers", []) or []
    ev_e        = _ratio(ratios, "ev_ebitda")
    pe_e        = _ratio(ratios, "pe_ratio")
    peer_median_ev_e = _peer_median(peers, "ev_ebitda")
    peer_median_pe   = _peer_median(peers, "pe")

    # --- Catalyseurs section (pleine largeur) — header à y=8.80 (avant 9.26) ---
    add_rect(slide, 1.02, 8.80, 23.37, 0.45, "1A7A4A")
    add_text_box(slide, 1.15, 8.82, 23.20, 0.42,
                 "CATALYSEURS \u00c0 SURVEILLER (12 MOIS)", 7.5, WHITE, bold=True)

    # Affichage 4 catalyseurs en 2 colonnes — boxes agrandies (h=1.10), police 6.5pt
    _cat_layout = [
        (1.02, 9.35),  (12.70, 9.35),   # row 1 : 2 catalyseurs côte à côte
        (1.02, 10.00), (12.70, 10.00),  # row 2
    ]
    for i, (cx, cy) in enumerate(_cat_layout):
        if i < len(catalysts):
            _cat_name = _g(catalysts[i], "title") or _g(catalysts[i], "name") or ""
            _cat_body = _g(catalysts[i], "description") or _g(catalysts[i], "text") or ""
            _cat_txt  = _truncate(f"{_cat_name[:30]} : {_cat_body}", 160) if _cat_body else _truncate(_cat_name, 160)
        else:
            _cat_txt = "\u2014"
        add_rect(slide, cx, cy + 0.05, 0.10, 0.22, "1A7A4A")
        add_text_box(slide, cx + 0.18, cy, 11.42, 0.62,
                     _cat_txt, 6.5, "333333", wrap=True)

    # Horizontal rule (before KPI)
    add_rect(slide, 1.02, 10.75, 23.37, 0.03, "AAAAAA")

    # 4 KPI boxes — y=10.90 pour suivre HR (10.75), h=2.40 pour remplir jusqu'à footer
    kpi_box(slide, 1.02, 10.90, 5.64, 2.40,
            _frx(ev_e), "EV/EBITDA",
            f"vs {_frx(peer_median_ev_e)} med. pairs")
    kpi_box(slide, 6.91, 10.90, 5.64, 2.40,
            _frpct(wacc_val), "WACC",
            f"Beta {_fr(beta, 2) if beta else '—'}  \u00b7  RFR {_frpct(rfr)}")
    kpi_box(slide, 12.80, 10.90, 5.64, 2.40,
            f"{_fr(tbase, 0)} {cur_sym}", "Valeur DCF base",
            f"Upside {_upside(tbase, price)} vs cours")
    kpi_box(slide, 18.69, 10.90, 5.64, 2.40,
            f"{sent_score:+.3f}".replace(".", ","),
            "Sentiment LLM" if _s2_is_llm else "Sentiment FinBERT",
            f"{sent_label_display}  \u00b7  {sent_articles} art.")

    return slide


# ---------------------------------------------------------------------------
# Slide 3 — Sommaire
# ---------------------------------------------------------------------------

def _slide_sommaire(prs, snap=None, synthesis=None):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    navy_bar(slide)
    footer_bar(slide)
    slide_title(slide, "Sommaire")

    # Descriptions dynamiques basées sur les données de la société
    ci       = snap.company_info if snap else None
    co_name  = _g(ci, "company_name", "") or ""
    sector   = _g(ci, "sector", "") or ""
    strengths= _g(synthesis, "strengths", []) or []
    risks    = _g(synthesis, "risks", []) or []
    rec      = _g(synthesis, "recommendation", "") or ""

    def _s1(items, n=2):
        """2 premiers éléments d'une liste, jointure ' · '"""
        parts = [str(x)[:40] for x in items[:n] if x]
        return "  \u00b7  ".join(parts) if parts else ""

    co_desc  = f"{co_name}  \u00b7  {sector}" if co_name and sector else (co_name or "Présentation, mod\u00e8le \u00e9conomique")
    str_desc = _s1(strengths) or "Analyse strat\u00e9gique & positionnément"
    ris_desc = _s1(risks) or "Risques structurels & th\u00e8se contraire"

    sections = [
        ("01", "Company Overview",    _truncate(co_desc, 80),  "4\u20136"),
        ("02", "Analyse Financi\u00e8re",  "Compte de r\u00e9sultat, bilan & liquidit\u00e9, ratios",         "7\u201310"),
        ("03", "Valorisation",        "DCF, comparable peers, Football Field",                 "11\u201314"),
        ("04", "Risques & Strat\u00e9gie", _truncate(ris_desc, 80),  "15\u201316"),
        ("05", "Sentiment & Annexes", "FinBERT, actionnariat, historique & conviction",        "17\u201321"),
    ]
    fills = [WHITE, GREY_BG, WHITE, GREY_BG, WHITE]
    ys    = [2.49, 4.42, 6.35, 8.28, 10.21]

    for i, (num, name, desc, pages) in enumerate(sections):
        y    = ys[i]
        fill = fills[i]
        add_rect(slide, 1.02, y, 23.37, 1.73, fill)
        add_rect(slide, 1.02, y, 0.13, 1.73, NAVY_MID)
        add_text_box(slide, 1.40, y + 0.15, 1.32, 1.43,
                     num, 17, GREY_LIGHT, bold=True)
        add_text_box(slide, 3.20, y + 0.20, 17.00, 0.71,
                     name, 13, NAVY, bold=True)
        add_text_box(slide, 3.20, y + 0.85, 17.00, 0.71,
                     desc, 8.5, GREY_TXT)
        add_text_box(slide, 20.57, y + 0.45, 3.82, 0.71,
                     pages, 9, NAVY_MID,
                     align=PP_ALIGN.RIGHT)

    return slide


# ---------------------------------------------------------------------------
# Slide 5 — Company Overview
# ---------------------------------------------------------------------------

def _slide_company_overview(prs, snap, synthesis, ratios):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    navy_bar(slide)
    footer_bar(slide)
    section_dots(slide, 1)

    ci      = snap.company_info if snap else None
    mkt     = snap.market if snap else None
    ticker  = _g(ci, "ticker", "—")
    co_name = _g(ci, "company_name", "—")
    currency = _g(ci, "currency", "USD") or "USD"
    cur_sym  = "EUR" if currency == "EUR" else "$"
    price    = _g(mkt, "share_price")

    slide_title(slide, "Company Overview",
                f"{co_name}  ·  Présentation & Positionnement stratégique")

    desc     = _g(synthesis, "company_description", "") or ""
    segments = _g(synthesis, "segments", []) or []
    strengths= _g(synthesis, "strengths", []) or []

    # Fallback si synthesis vide : description generique depuis raw data
    _sector_fb = _g(ci, "sector", "N/D") or "N/D"
    if not desc and co_name and co_name != "\u2014":
        desc = (f"{co_name} ({ticker}) opère dans le secteur {_sector_fb}. "
                f"Analyse FinSight IA en cours -- Données de Synthèse non disponibles "
                f"pour cette session. Lancer une nouvelle analyse pour obtenir la "
                f"description complète, les segments et le positionnément stratégique.")

    # Nombre de segments pour la formulation dynamique
    n_seg = len(segments)
    _num_fr = {1:"un",2:"deux",3:"trois",4:"quatre",5:"cinq",6:"six"}
    seg_count_str = _num_fr.get(n_seg, str(n_seg)) if n_seg else ""

    # Description + phrase d'intro dans la même boîte
    if n_seg:
        intro_txt = (
            f" L'entreprise articule son activit\u00e9 autour de "
            f"{seg_count_str} segment{'s' if n_seg > 1 else ''} "
            f"strat\u00e9gique{'s' if n_seg > 1 else ''} :"
        )
        full_desc = _truncate(desc, 700) + intro_txt
    else:
        full_desc = _truncate(desc, 700)

    add_rect(slide, 1.02, 2.67, 13.72, 9.78, GREY_BG)
    add_rect(slide, 1.02, 2.67, 0.13, 9.78, NAVY_MID)
    add_text_box(slide, 1.40, 2.84, 12.95, 4.0,
                 full_desc, 8.5, BLACK, wrap=True)

    if n_seg:
        bullet_y = 7.10
        row_h = 1.95  # hauteur par segment (augmente pour eviter truncation "...")
        for i, seg in enumerate(segments[:4]):
            seg_name = _g(seg, "name", "") or str(seg)
            seg_desc = _g(seg, "description", "") or ""
            by = bullet_y + i * row_h
            add_rect(slide, 1.40, by + 0.12, 0.18, 0.18, NAVY_MID)
            add_text_box(slide, 1.72, by, 12.05, 0.46,
                         _truncate(seg_name, 80), 8.5, BLACK, bold=True)
            if seg_desc:
                add_text_box(slide, 1.72, by + 0.44, 12.05, 1.40,
                             _truncate(seg_desc, 520), 7.5, GREY_TXT, italic=True, wrap=True)
    else:
        # Fallback : liste des strengths si pas de segments
        bullet_y = 7.25
        add_text_box(slide, 1.40, 6.55, 12.95, 0.51,
                     "Positionnement strat\u00e9gique", 9, NAVY, bold=True)
        for i, strength in enumerate(strengths[:4]):
            add_rect(slide, 1.40, bullet_y + i * 0.72 + 0.18, 0.18, 0.18, NAVY_MID)
            add_text_box(slide, 1.72, bullet_y + i * 0.72, 12.05, 0.56,
                         _truncate(str(strength), 80), 8.5, BLACK)

    # 4 KPI boxes on right
    shares    = _g(mkt, "shares_diluted")
    mktcap    = (shares * price / 1000) if (shares and price) else None
    years_sorted = _valid_years(snap)
    latest_yr_key = years_sorted[-1] if years_sorted else None
    rev = _fy(snap, latest_yr_key, "revenue")

    kpi_ys = [2.67, 5.13, 7.59, 10.06]
    kpi_box(slide, 15.44, kpi_ys[0], 8.89, 2.29,
            _frm(mktcap * 1000 if mktcap else None, cur_sym) if mktcap else "—",
            "Capitalisation boursière",
            fill=NAVY, accent=WHITE)
    kpi_box(slide, 15.44, kpi_ys[1], 8.89, 2.29,
            _frm(rev, cur_sym),
            f"Chiffre d'affaires ({latest_yr_key or 'LTM'})")
    kpi_box(slide, 15.44, kpi_ys[2], 8.89, 2.29,
            f"{_fr(price, 2)} {cur_sym}" if price else "—",
            "Cours actuel")
    ev_e = _ratio(ratios, "ev_ebitda")
    pe   = _ratio(ratios, "pe_ratio")
    kpi_box(slide, 15.44, kpi_ys[3], 8.89, 2.29,
            _frx(ev_e),
            "EV/EBITDA (LTM)",
            f"P/E : {_frx(pe)}")

    return slide


# ---------------------------------------------------------------------------
# Slide 6 — Business Model (Segments)
# ---------------------------------------------------------------------------

def _slide_business_model(prs, snap, synthesis):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    navy_bar(slide)
    footer_bar(slide)
    section_dots(slide, 1)
    slide_title(slide, "Mod\u00e8le \u00c9conomique",
                "Segments op\u00e9rationnels et moteurs de croissance")

    segments  = _g(synthesis, "segments", []) or []
    strengths = _g(synthesis, "strengths", []) or []

    # Fallback si pas de segments LLM : construire depuis strengths
    if not segments:
        _fallback_labels = ["Core Business", "Growth Driver", "International", "Innovation"]
        segments = [
            {"name": _fallback_labels[i], "description": strengths[i] if i < len(strengths) else "", "revenue_pct": None}
            for i in range(min(len(strengths), 4) or 3)
        ]

    # Limiter à 4 colonnes max pour la mise en page
    segs_to_show = segments[:4]
    n_cols = len(segs_to_show) or 1

    _palettes = [
        (NAVY_PALE,  NAVY_MID),
        (GREEN_PALE, GREEN),
        (GREY_BG,    GREY_LIGHT),
        (AMBER_PALE, AMBER),
    ]

    total_w = 23.36  # 25.4 - 1.02*2
    col_gap  = 0.30
    col_w    = (total_w - (n_cols - 1) * col_gap) / n_cols
    start_x  = 1.02
    col_h    = 9.78
    col_y    = 2.67

    for i, seg in enumerate(segs_to_show):
        fill, accent = _palettes[i % len(_palettes)]
        seg_name = _g(seg, "name", f"Segment {i+1}") or f"Segment {i+1}"
        seg_desc = _g(seg, "description", "") or ""
        rev_pct  = _g(seg, "revenue_pct")
        pct_str  = f"{float(rev_pct):.0f}%" if rev_pct is not None else "—"

        cx = start_x + i * (col_w + col_gap)
        add_rect(slide, cx, col_y, col_w, col_h, fill)
        add_rect(slide, cx, col_y, col_w, 0.15, accent)
        add_text_box(slide, cx + 0.25, col_y + 0.50, col_w - 0.50, 1.78,
                     pct_str, 22, accent, bold=True)
        add_text_box(slide, cx + 0.25, col_y + 2.29, col_w - 0.50, 0.30,
                     "du CA", 9, GREY_TXT)
        add_rect(slide, cx + 0.25, col_y + 2.79, col_w - 0.50, 0.05, GREY_LIGHT)
        add_text_box(slide, cx + 0.25, col_y + 3.05, col_w - 0.50, 0.71,
                     _truncate(seg_name, 60), 10, NAVY, bold=True, wrap=True)
        # Hauteur description etendue pour utiliser l'espace jusqu'au footer (13.39 - col_y - 3.81)
        desc_h = col_h - 3.81 - 0.30  # marge 0.30 cm avant le bas de la carte
        add_text_box(slide, cx + 0.25, col_y + 3.81, col_w - 0.50, max(desc_h, 5.0),
                     _truncate(seg_desc, 520), 8, GREY_TXT, wrap=True)

    return slide


# ---------------------------------------------------------------------------
# Slide 8 — Compte de Résultat
# ---------------------------------------------------------------------------

def _slide_is(prs, snap, synthesis, ratios):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    navy_bar(slide)
    footer_bar(slide)
    section_dots(slide, 2)

    ci      = snap.company_info if snap else None
    currency= _g(ci, "currency", "USD") or "USD"
    cur_sym = "EUR" if currency == "EUR" else "$"

    years_sorted = _valid_years(snap)
    is_proj  = _g(synthesis, "is_projections", {}) or {}
    proj_keys = list(is_proj.keys())

    yr_range_start = years_sorted[0] if years_sorted else ""
    yr_range_end   = proj_keys[-1] if proj_keys else (years_sorted[-1] if years_sorted else "")

    slide_title(slide, "Analyse Financi\u00e8re \u2014 Compte de R\u00e9sultat",
                f"{currency} millions  \u00b7  {yr_range_start}\u2013{yr_range_end}")

    # Build display columns (up to 3 hist + 2 proj)
    hist_cols = years_sorted[-3:] if len(years_sorted) >= 3 else years_sorted
    proj_cols = proj_keys[:2]
    all_cols  = hist_cols + proj_cols

    header = ["Indicateur"] + [str(c) for c in all_cols]
    n_cols  = len(header)

    # IS rows
    rows_data = []
    for col_k in all_cols:
        pass  # computed below

    def _get_col(col_k):
        if col_k in proj_keys:
            pd = is_proj.get(col_k, {})
            return pd
        else:
            fy = snap.years.get(col_k) if (snap and snap.years) else None
            yr = ratios.years.get(col_k) if (ratios and ratios.years) else None
            return {"fy": fy, "yr": yr}

    row_labels = [
        "Chiffre d'affaires",
        "Croissance YoY",
        "Résultat brut",
        "Marge brute",
        "EBITDA",
        "Marge EBITDA",
        "Résultat net",
        "Marge nette",
    ]

    rows_data = []
    for rl in row_labels:
        row = [rl]
        for col_k in all_cols:
            d = _get_col(col_k)
            if col_k in proj_keys:
                rev   = _g(d, "revenue")
                grow  = _g(d, "revenue_growth")
                gm    = _g(d, "gross_margin")
                ebitda= _g(d, "ebitda")
                em    = _g(d, "ebitda_margin")
                ni    = _g(d, "net_income")
                nm    = _g(d, "net_margin")
            else:
                fy = d.get("fy")
                yr = d.get("yr")
                rev   = getattr(fy, "revenue", None) if fy else None
                grow  = getattr(yr, "revenue_growth", None) if yr else None
                gm    = getattr(yr, "gross_margin", None) if yr else None
                ebitda= getattr(yr, "ebitda", None) if yr else None
                em    = getattr(yr, "ebitda_margin", None) if yr else None
                ni    = getattr(yr, "net_income", None) if yr else None
                nm    = getattr(yr, "net_margin", None) if yr else None

            gp = (float(rev) * float(gm)) if (rev is not None and gm is not None) else None
            if rl == "Chiffre d'affaires":
                row.append(_frn(rev))
            elif rl == "Croissance YoY":
                row.append(_frpct(grow, signed=True))
            elif rl == "Résultat brut":
                row.append(_frn(gp))
            elif rl == "Marge brute":
                row.append(_frpct_margin(gm))
            elif rl == "EBITDA":
                row.append(_frn(ebitda))
            elif rl == "Marge EBITDA":
                row.append(_frpct_margin(em))
            elif rl == "Résultat net":
                row.append(_frn(ni))
            elif rl == "Marge nette":
                row.append(_frpct_margin(nm))
            else:
                row.append("—")
        rows_data.append(row)

    col_w_total = 23.37
    first_col_pct = 0.28
    other_pct     = (1.0 - first_col_pct) / max(len(all_cols), 1)
    col_widths_pct = [first_col_pct] + [other_pct] * len(all_cols)

    tbl_h = min(1.0 + len(row_labels) * 0.71, 5.5)
    is_tbl = add_table(slide, 1.02, 2.54, col_w_total, tbl_h,
              len(row_labels), n_cols,
              col_widths_pct=col_widths_pct,
              header_data=header,
              rows_data=rows_data)

    # Post-process IS table:
    # 1. LTM column (last hist year): DDE8F5 on DATA rows only (header stays NAVY)
    # 2. Sub-metric rows: 7.5pt italic gray (Croissance, Marges)
    # 3. Valeur column: no special treatment needed
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN

    _SUBMAIN = {
        "Chiffre d'affaires", "Résultat brut", "EBITDA", "Résultat net"
    }
    _SUBSUB = {
        "Croissance YoY", "Marge brute", "Marge EBITDA", "Marge nette"
    }

    ltm_col_idx = len(hist_cols)  # table col index of last historical year

    for ri in range(1, len(row_labels) + 1):  # skip header row 0
        row_label = row_labels[ri - 1]
        is_sub = row_label in _SUBSUB

        for ci in range(n_cols):
            cell = is_tbl.cell(ri, ci)
            # LTM column highlight on data rows
            if ci == ltm_col_idx and ltm_col_idx >= 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb("DDE8F5")
            # Sub-metric row styling
            try:
                p = cell.text_frame.paragraphs[0]
                for run in p.runs:
                    if is_sub:
                        run.font.size = Pt(7.5)
                        run.font.italic = True
                        run.font.color.rgb = rgb("333333")
                    else:
                        run.font.size = Pt(8)
                        run.font.bold = True
                        if ci == ltm_col_idx:
                            run.font.color.rgb = rgb(NAVY)
            except Exception:
                pass

    # Commentary — IS-spécifique en priorité (pas valuation_comment qui est hors-sujet)
    fin_comment = _g(synthesis, "financial_commentary", "") or ""
    if not fin_comment.strip():
        # Fallback : générer un commentaire IS déterministe basé sur les chiffres
        latest_yr_c = None
        if ratios and getattr(ratios, "years", None):
            _lbl_c = sorted(ratios.years.keys(),
                            key=lambda k: str(k).replace("_LTM", "Z"))[-1]
            latest_yr_c = ratios.years[_lbl_c]
        if latest_yr_c:
            _rev_g = getattr(latest_yr_c, "revenue_growth", None)
            _eb_m = getattr(latest_yr_c, "ebitda_margin", None)
            _nm = getattr(latest_yr_c, "net_margin", None)
            _gm = getattr(latest_yr_c, "gross_margin", None)
            def _pp(v):
                if v is None: return None
                try:
                    fv = float(v)
                    return fv * 100 if abs(fv) <= 2 else fv
                except: return None
            _rg = _pp(_rev_g)
            _ebmp = _pp(_eb_m)
            _nmp = _pp(_nm)
            _gmp = _pp(_gm)
            parts = []
            if _rg is not None:
                if _rg >= 5:
                    parts.append(f"Croissance revenue de {_rg:+.1f}% sur la dernière période, "
                                 f"reflétant une dynamique commerciale soutenue.")
                elif _rg < 0:
                    parts.append(f"Contraction du chiffre d'affaires de {_rg:+.1f}%, "
                                 f"révélant une pression sur la demande ou un effet base défavorable.")
                else:
                    parts.append(f"Plateau du chiffre d'affaires ({_rg:+.1f}%), "
                                 f"suggérant une phase de maturité ou de transition stratégique.")
            if _gmp is not None:
                parts.append(f"Marge brute de {_gmp:.1f}%, indicateur du pricing power et "
                             f"de la structure de coûts variables.")
            if _ebmp is not None:
                parts.append(f"Marge EBITDA à {_ebmp:.1f}%, mesure de la profitabilité "
                             f"opérationnelle avant amortissements.")
            if _nmp is not None:
                parts.append(f"Marge nette de {_nmp:.1f}%, après prise en compte de la "
                             f"structure financière et de la fiscalité.")
            fin_comment = " ".join(parts) if parts else ""
    if fin_comment.strip():
        _is_title = _jpm_title("is", ratios=ratios, snap=snap, synthesis=synthesis)
        commentary_box(slide, 1.02, 8.94, 23.37, 3.40,
                       _truncate(fin_comment, 700), title=_is_title)

    return slide


# ---------------------------------------------------------------------------
# Slide 9 — Bilan & Liquidité
# ---------------------------------------------------------------------------

def _slide_bilan(prs, snap, synthesis, ratios):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    navy_bar(slide)
    footer_bar(slide)
    section_dots(slide, 2)

    ci      = snap.company_info if snap else None
    currency= _g(ci, "currency", "USD") or "USD"
    cur_sym = "EUR" if currency == "EUR" else "$"

    years_sorted  = _valid_years(snap)
    latest_yr_key = years_sorted[-1] if years_sorted else None
    latest_fy     = snap.years.get(latest_yr_key) if (snap and latest_yr_key) else None

    slide_title(slide, "Bilan & Liquidit\u00e9",
                f"Structure financi\u00e8re  \u00b7  {currency} millions  \u00b7  LTM {latest_yr_key or ''}")

    cash     = getattr(latest_fy, "cash", None) if latest_fy else None
    ltd      = getattr(latest_fy, "long_term_debt", None) if latest_fy else None
    std      = getattr(latest_fy, "short_term_debt", None) if latest_fy else None
    total_debt = ((_safe_float(ltd) or 0) + (_safe_float(std) or 0)) or None
    net_debt_val = ((_safe_float(total_debt) or 0) - (_safe_float(cash) or 0)) if (total_debt is not None or cash is not None) else None

    kpi_box(slide, 1.02, 2.67, 7.37, 2.29, _frm(cash, cur_sym), "Cash & équivalents")
    kpi_box(slide, 1.02, 5.28, 7.37, 2.29, _frm(ltd, cur_sym), "Dette long terme")
    kpi_box(slide, 1.02, 7.90, 7.37, 2.29, _frm(net_debt_val, cur_sym), "Dette nette",
            fill=NAVY, accent=WHITE)

    # Ratio table (7 rows × 4 cols with Signal column)
    cr  = _ratio(ratios, "current_ratio")
    qr  = _ratio(ratios, "quick_ratio")
    de  = _ratio(ratios, "debt_equity")
    ic  = _ratio(ratios, "interest_coverage")
    nde = _ratio(ratios, "net_debt_ebitda")
    az  = _ratio(ratios, "altman_z")
    bm  = _ratio(ratios, "beneish_m")

    def _sig(v, lo_good, hi_warn, reverse=False):
        if v is None: return "\u2014"
        try:
            vf = float(v)
            if not reverse:
                if vf >= lo_good: return "\u2705"
                if vf >= hi_warn: return "\u26a0\ufe0f"
                return "\u274c"
            else:
                if vf <= lo_good: return "\u2705"
                if vf <= hi_warn: return "\u26a0\ufe0f"
                return "\u274c"
        except: return "\u2014"

    az_sig = "\u2705" if (az and float(az) >= 2.99) else ("\u26a0\ufe0f" if (az and float(az) >= 1.81) else ("\u274c" if az else "\u2014"))
    bm_sig = "\u2705" if (bm and float(bm) <= -2.22) else ("\u274c" if bm else "\u2014")

    add_text_box(slide, 9.4, 2.34, 14.99, 0.46,
                 "Ratios de solvabilit\u00e9 & liquidit\u00e9", 9, NAVY, bold=True)

    ratio_rows = [
        ["Ratio liquidité courante", _fr(cr, 2),  "> 1,5",   _sig(cr, 1.5, 1.0)],
        ["Ratio rapide (quick)",     _fr(qr, 2),  "> 1,0",   _sig(qr, 1.0, 0.7)],
        ["D/E",                      _frx(de),    "< 2,0x",  _sig(de, 1.5, 3.0, reverse=True)],
        ["Couverture d'intérêts",    _frx(ic),    "> 3,0x",  _sig(ic, 3.0, 1.5)],
        ["Dette nette / EBITDA",     _frx(nde),   "< 2,0x",  _sig(nde, 2.0, 4.0, reverse=True)],
        ["Altman Z",                 _fr(az, 2),  "> 2,99",  az_sig],
        ["Beneish M",                _fr(bm, 2),  "< -2,22", bm_sig],
    ]
    ratio_tbl = add_table(slide, 9.4, 2.80, 14.99, 5.0,
              len(ratio_rows), 4,
              col_widths_pct=[0.44, 0.21, 0.21, 0.14],
              header_data=["Ratio", "Valeur", "R\u00e9f\u00e9rence", "Signal"],
              rows_data=ratio_rows)

    # Bold les noms d'indicateurs (colonne 0)
    from pptx.util import Pt
    for ri in range(1, len(ratio_rows) + 1):
        try:
            cell = ratio_tbl.cell(ri, 0)
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.bold = True
                run.font.size = Pt(8)
        except Exception:
            pass

    # Commentary bilan : prioriser un texte spécifique au bilan, sinon générer
    fin_comment = _g(synthesis, "balance_sheet_commentary", "") or ""
    if not fin_comment.strip():
        # Fallback : générer un commentaire bilan déterministe
        latest_yr_b = None
        if ratios and getattr(ratios, "years", None):
            _lbl_b = sorted(ratios.years.keys(),
                            key=lambda k: str(k).replace("_LTM", "Z"))[-1]
            latest_yr_b = ratios.years[_lbl_b]
        if latest_yr_b:
            _nd_eb = getattr(latest_yr_b, "net_debt_ebitda", None)
            _cur_r = getattr(latest_yr_b, "current_ratio", None)
            _qck_r = getattr(latest_yr_b, "quick_ratio", None)
            _icr = getattr(latest_yr_b, "interest_coverage", None)
            parts_b = []
            if _nd_eb is not None:
                try:
                    nd_v = float(_nd_eb)
                    if nd_v < 0:
                        parts_b.append(f"Position de cash net positif ({nd_v:.2f}x EBITDA), "
                                       f"conférant une flexibilité financière maximale.")
                    elif nd_v < 1:
                        parts_b.append(f"Levier très faible (ND/EBITDA {nd_v:.2f}x), "
                                       f"laissant une large capacité d'investissement ou de retour au capital.")
                    elif nd_v < 3:
                        parts_b.append(f"Levier modéré (ND/EBITDA {nd_v:.2f}x), "
                                       f"compatible avec un profil investment grade.")
                    else:
                        parts_b.append(f"Levier élevé (ND/EBITDA {nd_v:.2f}x), "
                                       f"contraignant la marge de manœuvre stratégique.")
                except Exception:
                    pass
            if _cur_r is not None:
                try:
                    cr_v = float(_cur_r)
                    parts_b.append(f"Current ratio à {cr_v:.2f}x, mesure de la couverture "
                                   f"des passifs courants par les actifs courants.")
                except Exception:
                    pass
            if _icr is not None:
                try:
                    icr_v = float(_icr)
                    if icr_v >= 5:
                        parts_b.append(f"Couverture des intérêts solide ({icr_v:.1f}x EBIT), "
                                       f"capacité de service de la dette robuste.")
                    elif icr_v >= 2:
                        parts_b.append(f"Couverture des intérêts adéquate ({icr_v:.1f}x EBIT).")
                    else:
                        parts_b.append(f"Couverture des intérêts tendue ({icr_v:.1f}x EBIT) — "
                                       f"surveillance recommandée.")
                except Exception:
                    pass
            fin_comment = " ".join(parts_b) if parts_b else (
                _g(synthesis, "financial_commentary", "") or ""
            )
    if fin_comment.strip():
        _bs_title = _jpm_title("bilan", ratios=ratios, snap=snap, synthesis=synthesis)
        commentary_box(slide, 1.02, 10.75, 23.37, 2.60,
                       fin_comment[:600], title=_bs_title)

    return slide


# ---------------------------------------------------------------------------
# Slide 10 — Ratios Cles
# ---------------------------------------------------------------------------

def _slide_ratios(prs, snap, synthesis, ratios):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    navy_bar(slide)
    footer_bar(slide)
    section_dots(slide, 2)

    years_sorted  = _valid_years(snap)
    latest_yr_key = years_sorted[-1] if years_sorted else None

    slide_title(slide, "Ratios Cl\u00e9s vs. Benchmark Sectoriel",
                f"Positionnement relatif  \u00b7  LTM {latest_yr_key or ''}")

    pe   = _ratio(ratios, "pe_ratio")
    ev_e = _ratio(ratios, "ev_ebitda")
    ev_r = _ratio(ratios, "ev_revenue")
    gm   = _ratio(ratios, "gross_margin")
    em   = _ratio(ratios, "ebitda_margin")
    roe  = _ratio(ratios, "roe")
    az   = _ratio(ratios, "altman_z")
    bm   = _ratio(ratios, "beneish_m")
    roic = _ratio(ratios, "roic")
    capex_rev = _ratio(ratios, "capex_revenue")
    # Si capex_revenue absent des ratios, le calculer depuis snap
    if capex_rev is None and snap and latest_yr_key:
        capex_raw = _fy(snap, latest_yr_key, "capex")
        rev_raw   = _fy(snap, latest_yr_key, "revenue")
        if capex_raw is not None and rev_raw and float(rev_raw) > 0:
            capex_rev = abs(float(capex_raw)) / float(rev_raw)

    def _lecture_pe(v):
        if v is None: return "—"
        try:
            v = float(v)
            if v < 0:    return "Déficit"
            if v < 15:   return "Sous-value"
            if v < 25:   return "Correct"
            if v < 40:   return "Premium"
            return "Élevé"
        except: return "—"

    def _lecture_eve(v):
        if v is None: return "—"
        try:
            v = float(v)
            if v < 8:   return "Bas"
            if v < 14:  return "Correct"
            if v < 20:  return "Premium"
            return "Tr\u00e8s \u00e9lev\u00e9"
        except: return "—"

    def _lecture_z(v):
        if v is None: return "—"
        try:
            v = float(v)
            if v >= 2.99: return "Solide"
            if v >= 1.81: return "Zone grise"
            return "D\u00e9tresse"
        except: return "—"

    def _lecture_bm(v):
        if v is None: return "—"
        try:
            return "Risque manip." if float(v) > -2.22 else "Aucun signal"
        except: return "—"

    def _lecture_evr(v):
        if v is None: return "—"
        try:
            v = float(v)
            if v < 1:   return "Sous-value"
            if v < 5:   return "Correct"
            if v < 10:  return "Premium"
            return "Tr\u00e8s \u00e9lev\u00e9"
        except: return "—"

    def _lecture_gm(v):
        if v is None: return "—"
        try:
            v = float(v) * 100
            if v >= 60: return "Excellent"
            if v >= 40: return "Bon"
            if v >= 20: return "Correct"
            return "Bas"
        except: return "—"

    def _lecture_em(v):
        if v is None: return "—"
        try:
            v = float(v) * 100
            if v >= 30: return "Excellent"
            if v >= 20: return "Bon"
            if v >= 10: return "Correct"
            return "Bas"
        except: return "—"

    def _lecture_roe(v):
        if v is None: return "—"
        try:
            v = float(v) * 100
            if v >= 20: return "Excellent"
            if v >= 15: return "Bon"
            if v >= 5:  return "Correct"
            return "Bas"
        except: return "—"

    def _lecture_roic(v):
        if v is None: return "—"
        try:
            v = float(v) * 100
            if v >= 15: return "Excellent"
            if v >= 10: return "Bon"
            if v >= 5:  return "Correct"
            return "Bas"
        except: return "—"

    def _lecture_capex(v):
        if v is None: return "—"
        try:
            v = abs(float(v)) * 100
            if v < 5:   return "L\u00e9ger"
            if v < 15:  return "Mod\u00e9r\u00e9"
            return "Intensif"
        except: return "—"

    rows = [
        ["P/E",            _frx(pe),       "15\u201325x",  _lecture_pe(pe)],
        ["EV/EBITDA",      _frx(ev_e),     "8\u201314x",   _lecture_eve(ev_e)],
        ["EV/Revenue",     _frx(ev_r),     "1\u20135x",    _lecture_evr(ev_r)],
        ["Marge brute",    _frpct(gm),     "> 40 %",       _lecture_gm(gm)],
        ["Marge EBITDA",   _frpct(em),     "> 20 %",       _lecture_em(em)],
        ["ROE",            _frpct(roe),    "> 15 %",       _lecture_roe(roe)],
        ["ROIC",           _frpct(roic),   "> 10 %",       _lecture_roic(roic)],
        ["CapEx / Revenue",_frpct(capex_rev), "< 15 %",   _lecture_capex(capex_rev)],
        ["Altman Z",       _fr(az, 2),     "> 2,99",       _lecture_z(az)],
        ["Beneish M",      _fr(bm, 2),     "< -2,22",      _lecture_bm(bm)],
    ]

    ratio_tbl = add_table(slide, 1.02, 2.54, 23.37, 7.11,
              len(rows), 4,
              col_widths_pct=[0.25, 0.25, 0.25, 0.25],
              header_data=["Indicateur", "Valeur", "Benchmark", "Lecture"],
              rows_data=rows,
              border_hex="DDDDDD")

    # Per-cell coloring of Lecture column (col 3) + Valeur column navy bold (col 1)
    _GOOD_READS  = {"solide", "correct", "en ligne", "dans la norme", "sous-value", "aucun signal",
                    "bon", "excellent", "leger"}
    _WARN_READS  = {"élevé", "très élevé", "premium", "risque manip.", "détresse", "déficit",
                    "supérieur", "inférieure", "bas", "zone grise", "décote", "intensif"}
    _NEUT_READS  = {"en ligne", "prime technologique", "modéré"}  # amber category

    from pptx.util import Pt
    for ri in range(1, len(rows) + 1):
        # Valeur column (col 1): navy bold
        cell_v = ratio_tbl.cell(ri, 1)
        try:
            p = cell_v.text_frame.paragraphs[0]
            for run in p.runs:
                run.font.color.rgb = rgb(NAVY)
                run.font.bold = True
        except Exception: pass

        # Lecture column (col 3): colored fill + bold
        cell_l = ratio_tbl.cell(ri, 3)
        val_str = (rows[ri - 1][3] or "").strip().lower()
        if val_str in _GOOD_READS:
            cell_l.fill.solid()
            cell_l.fill.fore_color.rgb = rgb(GREEN_PALE)
            lec_tc = GREEN
        elif val_str in _NEUT_READS or val_str in {"supérieur", "supérieure"}:
            cell_l.fill.solid()
            cell_l.fill.fore_color.rgb = rgb(AMBER_PALE)
            lec_tc = AMBER
        elif val_str in _WARN_READS:
            cell_l.fill.solid()
            cell_l.fill.fore_color.rgb = rgb(RED_PALE)
            lec_tc = RED
        else:
            lec_tc = GREY_TXT
        try:
            p = cell_l.text_frame.paragraphs[0]
            for run in p.runs:
                run.font.color.rgb = rgb(lec_tc)
                run.font.bold = True
        except Exception: pass

    ratio_comment = _g(synthesis, "ratio_commentary", "") or ""
    if ratio_comment.strip():
        _ratios_title = _jpm_title("ratios", ratios=ratios, snap=snap, synthesis=synthesis)
        commentary_box(slide, 1.02, 10.08, 23.37, 2.29, ratio_comment, title=_ratios_title)

    return slide


# ---------------------------------------------------------------------------
# Slide 12 — DCF & Scénarios
# ---------------------------------------------------------------------------

def _slide_dcf(prs, snap, synthesis, ratios):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    navy_bar(slide)
    footer_bar(slide)
    section_dots(slide, 3)
    slide_title(slide, "Valorisation \u2014 DCF & Sc\u00e9narios")

    ci      = snap.company_info if snap else None
    mkt     = snap.market if snap else None
    currency= _g(ci, "currency", "USD") or "USD"
    cur_sym = "EUR" if currency == "EUR" else "$"
    price   = _g(mkt, "share_price")
    wacc    = _g(mkt, "wacc") or 0.10
    beta    = _g(mkt, "beta_levered")
    rfr     = _g(mkt, "risk_free_rate") or 0.041
    erp     = _g(mkt, "erp") or 0.055
    tgr     = _g(mkt, "terminal_growth") or 0.03

    tbase   = _g(synthesis, "target_base")
    tbear   = _g(synthesis, "target_bear")
    tbull   = _g(synthesis, "target_bull")

    # LEFT: DCF params table (simplified) + DDE8F5 highlight on valeur intrinseque
    upside_impl = _upside(tbase, price)
    param_rows = [
        ["WACC",                   _frpct(wacc)],
        ["TGR",                    _frpct(tgr)],
        ["Valeur intrins\u00e8que", f"{_fr(tbase, 0)} {cur_sym}" if tbase else "\u2014"],
        ["Cours actuel",           f"{_fr(price, 2)} {cur_sym}" if price else "\u2014"],
        ["Upside implicite",       upside_impl],
    ]
    param_tbl = add_table(slide, 1.02, 2.29, 11.18, 3.81,
              len(param_rows), 2,
              col_widths_pct=[0.55, 0.45],
              header_data=["Param\u00e8tre", "Valeur"],
              rows_data=param_rows,
              border_hex="DDDDDD")

    # Highlight "Valeur intrinsèque" row (row index 3 = data row 2 → tbl row 3)
    from pptx.util import Pt as _Pt_dcf
    try:
        for ci_p in range(2):
            c = param_tbl.cell(3, ci_p)
            c.fill.solid()
            c.fill.fore_color.rgb = rgb("DDE8F5")
            for run in c.text_frame.paragraphs[0].runs:
                run.font.color.rgb = rgb(NAVY)
                run.font.bold = True
    except Exception:
        pass

    # RIGHT: 3-column scénario table (matches reference exactly)
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN as _PA
    sc_tbl_obj = slide.shapes.add_table(4, 3, __import__('pptx.util', fromlist=['Cm']).Cm(12.45), __import__('pptx.util', fromlist=['Cm']).Cm(2.29), __import__('pptx.util', fromlist=['Cm']).Cm(11.94), __import__('pptx.util', fromlist=['Cm']).Cm(3.81)).table
    # Column widths equal
    from pptx.util import Cm as _Cm, Pt as _Pt
    col_w_emu = int(_Cm(11.94)) // 3
    for ci in range(3): sc_tbl_obj.columns[ci].width = col_w_emu

    sc_cfgs = [
        ("Bear",  tbear,  RED,   RED_PALE),
        ("Base",  tbase,  NAVY_MID, NAVY_PALE),
        ("Bull",  tbull,  GREEN, GREEN_PALE),
    ]
    for ci, (lbl, val, accent, pale) in enumerate(sc_cfgs):
        # Row 0: header label
        c = sc_tbl_obj.cell(0, ci)
        c.fill.solid(); c.fill.fore_color.rgb = rgb(accent)
        _set_cell(c, lbl, font_size=10, bold=True, color_hex=WHITE, align=_PA.CENTER)
        # Row 1: price value (large)
        c = sc_tbl_obj.cell(1, ci)
        c.fill.solid(); c.fill.fore_color.rgb = rgb(pale)
        price_str = f"{_fr(val, 0)} {cur_sym}" if val else "\u2014"
        _set_cell(c, price_str, font_size=20, bold=True, color_hex=accent, align=_PA.CENTER)
        # Row 2: upside
        c = sc_tbl_obj.cell(2, ci)
        c.fill.solid(); c.fill.fore_color.rgb = rgb(pale)
        _set_cell(c, _upside(val, price), font_size=12, bold=True, color_hex=accent, align=_PA.CENTER)
        # Row 3: footer label
        c = sc_tbl_obj.cell(3, ci)
        c.fill.solid(); c.fill.fore_color.rgb = rgb(pale)
        _set_cell(c, "Upside / Downside", font_size=7, bold=False, color_hex=GREY_LIGHT, align=_PA.CENTER)

    # ── Monte Carlo DCF — bloc gauche ───────────────────────────────────────
    mc_meta  = getattr(ratios, "meta", {}) or {}
    mc_p10   = mc_meta.get("dcf_mc_p10")
    mc_p50   = mc_meta.get("dcf_mc_p50")
    mc_p90   = mc_meta.get("dcf_mc_p90")
    mc_p2    = mc_meta.get("dcf_mc_p2")
    mc_p98   = mc_meta.get("dcf_mc_p98")
    mc_n     = mc_meta.get("dcf_mc_n_sim")
    mc_label = f"Monte Carlo GBM \u2014 {int(mc_n):,} simulations".replace(",", "\u202f") if mc_n else "Monte Carlo GBM"

    add_text_box(slide, 1.02, 6.65, 11.18, 0.56, mc_label, 9, NAVY, bold=True)

    if mc_p50 is not None:
        _mc_rows = [
            ["Stress Down (P2)",  f"{_fr(mc_p2, 0)} {cur_sym}"  if mc_p2  else "\u2014", _upside(mc_p2,  price)],
            ["Pessimiste (P10)",  f"{_fr(mc_p10, 0)} {cur_sym}" if mc_p10 else "\u2014", _upside(mc_p10, price)],
            ["\u25cf M\u00e9dian (P50)", f"{_fr(mc_p50, 0)} {cur_sym}",                   _upside(mc_p50, price)],
            ["Optimiste (P90)",   f"{_fr(mc_p90, 0)} {cur_sym}" if mc_p90 else "\u2014", _upside(mc_p90, price)],
            ["Stress Up (P98)",   f"{_fr(mc_p98, 0)} {cur_sym}" if mc_p98 else "\u2014", _upside(mc_p98, price)],
        ]
        mc_tbl = add_table(slide, 1.02, 7.37, 11.18, 3.10,
                  len(_mc_rows), 3,
                  col_widths_pct=[0.42, 0.32, 0.26],
                  header_data=["Sc\u00e9nario", "Cours pr\u00e9dit (12M)", "vs Cours"],
                  rows_data=_mc_rows,
                  border_hex="DDDDDD")
        # Highlight médian row
        try:
            for _ci in range(3):
                _c = mc_tbl.cell(3, _ci)
                _c.fill.solid(); _c.fill.fore_color.rgb = rgb("DDE8F5")
                for _run in _c.text_frame.paragraphs[0].runs:
                    _run.font.bold = True; _run.font.color.rgb = rgb(NAVY)
        except Exception:
            pass
    else:
        add_text_box(slide, 1.02, 7.50, 11.18, 2.80,
                     "Simulations GBM non disponibles\npour ce ticker.",
                     9, GREY_LIGHT)

    # ── Sensitivity table — droite ───────────────────────────────────────────
    add_text_box(slide, 12.45, 6.65, 11.94, 0.56,
                 f"Sensibilit\u00e9 WACC \u00d7 TGR \u2014 Valeur ({cur_sym})",
                 9, NAVY, bold=True)

    w_deltas = [-0.02, -0.01, 0.00, 0.01, 0.02]
    t_deltas = [-0.01, -0.005, 0.00, 0.005, 0.01]

    header_s = ["WACC \\ TGR"] + [f"{(tgr + t) * 100:.1f}%".replace(".", ",") for t in t_deltas]
    sens_rows = []
    for w_d in w_deltas:
        w_val = wacc + w_d
        row   = [f"{w_val * 100:.1f}%".replace(".", ",")]
        for t_d in t_deltas:
            row.append(_dcf_value(tbase, wacc, tgr, w_d, t_d))
        sens_rows.append(row)

    sens_tbl = add_table(slide, 12.45, 7.37, 11.94, 3.10,
              len(sens_rows), len(header_s),
              header_data=header_s,
              rows_data=sens_rows,
              border_hex="DDDDDD")

    # Cross-highlight: WACC base row + TGR base column = DDE8F5; intersection = NAVY_MID/blanc
    try:
        base_ri = w_deltas.index(0.00) + 1   # +1 for header row
        base_ci = t_deltas.index(0.00) + 1   # +1 for label col
        n_rows_s = len(sens_rows) + 1
        n_cols_s = len(header_s)
        for ri in range(1, n_rows_s):
            for ci in range(1, n_cols_s):
                if ri == base_ri or ci == base_ci:
                    cell = sens_tbl.cell(ri, ci)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = rgb("DDE8F5")
        # Highlight WACC label cell (col 0) at base row
        try:
            lbl_cell = sens_tbl.cell(base_ri, 0)
            lbl_cell.fill.solid()
            lbl_cell.fill.fore_color.rgb = rgb("DDE8F5")
            for run in lbl_cell.text_frame.paragraphs[0].runs:
                run.font.bold = True
                run.font.color.rgb = rgb(NAVY)
        except Exception:
            pass
        # Intersection: navy_mid + white bold
        base_cell = sens_tbl.cell(base_ri, base_ci)
        base_cell.fill.solid()
        base_cell.fill.fore_color.rgb = rgb(NAVY_MID)
        try:
            p = base_cell.text_frame.paragraphs[0]
            for run in p.runs:
                run.font.color.rgb = rgb(WHITE)
                run.font.bold = True
        except Exception:
            pass
    except Exception:
        pass

    dcf_comment = _g(synthesis, "dcf_commentary", "") or ""
    if dcf_comment.strip():
        _dcf_title = _jpm_title("dcf", ratios=ratios, snap=snap, synthesis=synthesis)
        commentary_box(slide, 1.02, 11.20, 23.37, 1.52, dcf_comment, title=_dcf_title)

    return slide


# ---------------------------------------------------------------------------
# Slide 13 — Comparable Peers
# ---------------------------------------------------------------------------

def _slide_peers(prs, snap, synthesis, ratios):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    navy_bar(slide)
    footer_bar(slide)
    section_dots(slide, 3)

    ci      = snap.company_info if snap else None
    mkt     = snap.market if snap else None
    currency= _g(ci, "currency", "USD") or "USD"
    cur_sym = "EUR" if currency == "EUR" else "$"
    ticker  = _g(ci, "ticker", "—")
    co_name = _g(ci, "company_name", "—")
    price   = _g(mkt, "share_price")
    shares  = _g(mkt, "shares_diluted")
    mktcap  = (shares * price / 1000) if (shares and price) else None

    slide_title(slide, "Comparable Peers",
                f"Analyse par multiples  \u00b7  LTM  \u00b7  Mkt Cap en Mds {currency}")

    ev_e = _ratio(ratios, "ev_ebitda")
    ev_r = _ratio(ratios, "ev_revenue")
    pe   = _ratio(ratios, "pe_ratio")
    gm   = _ratio(ratios, "gross_margin")
    em   = _ratio(ratios, "ebitda_margin")

    peers = _g(synthesis, "comparable_peers", []) or []

    header = ["Soci\u00e9t\u00e9", "Ticker", "Mkt Cap", "EV/EBITDA", "EV/Rev.", "P/E",
              "Mg Brute", "Mg EBITDA"]

    # Target company row — Mkt Cap en Mds (shares en millions * price / 1000 = Mds)
    mktcap_str = _fr(mktcap, 1) if mktcap else "—"
    target_row = [co_name, ticker, mktcap_str,
                  _frx(ev_e), _frx(ev_r), _frx(pe),
                  _frpct_margin(gm), _frpct_margin(em)]

    rows_data = [target_row]
    rows_fills = ["DDE8F5"]

    for peer in peers[:5]:
        pn   = _g(peer, "name", "—") or "—"
        pt   = _g(peer, "ticker", "—") or "—"
        # Ensure proper title case for company names
        pn = str(pn).strip()
        if pn and pn != "—" and pn == pn.lower():
            pn = pn.title()
        p_mc = _g(peer, "market_cap_mds")
        if p_mc is not None:
            try:
                p_mc_f = float(p_mc)
                # Si le LLM a retourné en millions au lieu de milliards (> 1000 Mds improbable)
                if p_mc_f > 10000:
                    p_mc_f = p_mc_f / 1000
                p_mktcap_str = _fr(p_mc_f, 1)
            except Exception:
                p_mktcap_str = "—"
        else:
            p_mktcap_str = "—"
        prow = [
            pn[:30], str(pt).upper() if pt != "—" else "—",
            p_mktcap_str,
            _frx(_g(peer, "ev_ebitda")),
            _frx(_g(peer, "ev_revenue")),
            _frx(_g(peer, "pe")),
            _frpct_margin(_g(peer, "gross_margin")),
            _frpct_margin(_g(peer, "ebitda_margin")),
        ]
        rows_data.append(prow)
        rows_fills.append(WHITE if len(rows_data) % 2 == 0 else GREY_BG)

    # Médian row
    med_eve = _peer_median(peers, "ev_ebitda")
    med_evr = _peer_median(peers, "ev_revenue")
    med_pe  = _peer_median(peers, "pe")
    med_gm  = _peer_median(peers, "gross_margin")
    med_em  = _peer_median(peers, "ebitda_margin")
    median_row = ["Médiane peers", "—", "—",
                  _frx(med_eve), _frx(med_evr), _frx(med_pe),
                  _frpct_margin(med_gm), _frpct_margin(med_em)]
    rows_data.append(median_row)
    rows_fills.append(GREY_BG)

    tbl_h = min(1.0 + len(rows_data) * 0.71, 5.6)
    peers_tbl = add_table(slide, 1.02, 2.54, 23.37, tbl_h,
              len(rows_data), len(header),
              col_widths_pct=[0.20, 0.09, 0.13, 0.12, 0.12, 0.09, 0.13, 0.12],
              header_data=header,
              rows_data=rows_data,
              row_fills=rows_fills)

    # Post-process subject row (row 1): navy bold text
    for ci in range(len(header)):
        cell = peers_tbl.cell(1, ci)
        try:
            p = cell.text_frame.paragraphs[0]
            for run in p.runs:
                run.font.bold = True
                run.font.color.rgb = rgb(NAVY if ci != 1 else NAVY_MID)
                if ci == 1:  # Ticker col: small gray
                    run.font.size = __import__('pptx.util', fromlist=['Pt']).Pt(7.5)
        except Exception:
            pass

    # Post-process peer rows: ticker col (ci=1) small gray
    for ri in range(2, len(rows_data)):  # peer rows (not subject, not median)
        cell = peers_tbl.cell(ri, 1)
        try:
            p = cell.text_frame.paragraphs[0]
            for run in p.runs:
                from pptx.util import Pt as _Pt2
                run.font.size = _Pt2(7.5)
                run.font.color.rgb = rgb(GREY_TXT)
        except Exception:
            pass

    # Médian row (last row): bold=True, NOT italic
    median_ri = len(rows_data)
    for ci in range(len(header)):
        cell = peers_tbl.cell(median_ri, ci)
        try:
            p = cell.text_frame.paragraphs[0]
            for run in p.runs:
                run.font.bold = True
                run.font.italic = False
        except Exception:
            pass

    _p1 = _g(synthesis, "peers_commentary", "") or ""
    _p2 = _g(synthesis, "ratio_commentary", "") or ""
    _p3 = _g(synthesis, "valuation_comment", "") or ""
    peers_comment = " ".join(p for p in [_p1, _p2, _p3] if p.strip())
    if not peers_comment.strip():
        peers_comment = _g(synthesis, "summary", "") or ""
    if peers_comment.strip():
        _peers_title = _jpm_title("peers", ratios=ratios, snap=snap, synthesis=synthesis)
        commentary_box(slide, 1.02, 8.69, 23.37, 4.30, peers_comment, title=_peers_title)

    return slide


# ---------------------------------------------------------------------------
# Slide 14 — Football Field
# ---------------------------------------------------------------------------

def _slide_football_field(prs, snap, synthesis, ratios):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    navy_bar(slide)
    footer_bar(slide)
    section_dots(slide, 3)

    ci       = snap.company_info if snap else None
    mkt      = snap.market if snap else None
    currency = _g(ci, "currency", "USD") or "USD"
    cur_sym  = "EUR" if currency == "EUR" else "$"
    price    = _g(mkt, "share_price")

    slide_title(slide, "Football Field Chart",
                f"Synth\u00e8se des m\u00e9thodes de valorisation  \u00b7  {currency} par action")

    ff     = _g(synthesis, "football_field", []) or []

    # ── Matplotlib Football Field chart ──────────────────────────────────────
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import numpy as np
        _mpl_ok = True
    except Exception:
        _mpl_ok = False

    COLOR_MAP = {
        "DCF - Bear":   "#A82020",
        "DCF - Base":   "#1B3A6B",
        "DCF - Bull":   "#1A7A4A",
        "Comparables":  "#2E5FA3",
        "DDM":          "#B06000",
    }
    DEFAULT_COLORS = ["#1B3A6B", "#2E5FA3", "#1A7A4A", "#A82020", "#B06000",
                      "#6E4FC2", "#888888"]

    ff_methods = []
    ff_lows    = []
    ff_highs   = []
    ff_colors  = []
    for i, item in enumerate(ff):
        label = str(_g(item, "label", "—"))
        low   = _g(item, "range_low")
        high  = _g(item, "range_high")
        if low is None or high is None:
            continue
        lo_f, hi_f = float(low), float(high)
        lo_f, hi_f = min(lo_f, hi_f), max(lo_f, hi_f)   # swap si LLM inverse lo/hi
        if hi_f <= lo_f:
            continue
        ff_methods.append(label)
        ff_lows.append(lo_f)
        ff_highs.append(hi_f)
        ff_colors.append(COLOR_MAP.get(label, DEFAULT_COLORS[i % len(DEFAULT_COLORS)]))

    def _make_ff_buf():
        n = len(ff_methods)
        # Hauteur cible en cm (zone disponible : y=2.54 à y=9.40 avant commentary)
        img_h_cm = min(max(3.8, n * 1.15), 6.8)
        # Dimensions figure proportionnelles à la zone pptx (23.37 x img_h_cm)
        fig_w_in = 12.0
        fig_h_in = img_h_cm * (fig_w_in / 23.37)
        fig, ax = plt.subplots(figsize=(fig_w_in, fig_h_in))
        y = np.arange(n)
        all_v = ff_lows + ff_highs
        plage = (max(all_v) - min(all_v)) if len(all_v) > 1 else max(all_v)
        offset = max(plage * 0.04, 3.0)
        for i, (lo, hi, col) in enumerate(zip(ff_lows, ff_highs, ff_colors)):
            ax.barh(y[i], hi - lo, left=lo, height=0.52, color=col, alpha=0.87, zorder=3)
            ax.text((lo + hi) / 2, y[i], f"{int((lo + hi) / 2)}",
                    va='center', ha='center', fontsize=9, color='white',
                    fontweight='bold', zorder=4)
            ax.text(lo - offset, y[i], f"{int(lo)}", va='center', ha='right',
                    fontsize=8, color='#444', clip_on=False)
            ax.text(hi + offset, y[i], f"{int(hi)}", va='center', ha='left',
                    fontsize=8, color='#444', clip_on=False)
        if price and ff_lows:
            ax.axvline(x=price, color='#B06000', linewidth=2.0,
                       linestyle='--', zorder=5)
            price_lbl = f"{cur_sym}{price:.0f}"
            ax.text(price, n - 0.05, f"Cours : {price_lbl}",
                    fontsize=8, color='#B06000', fontweight='bold',
                    ha='center', va='top', clip_on=False,
                    bbox=dict(boxstyle='round,pad=0.3', facecolor='white',
                              alpha=0.95, edgecolor='#B06000', linewidth=0.8))
        ax.set_yticks(y)
        ax.set_yticklabels(ff_methods, fontsize=9, color='#222')
        margin = max(25, plage * 0.20)
        ax.set_xlim(min(all_v) - margin, max(all_v) + margin)
        ax.set_xlabel(f"Valeur par action ({currency})", fontsize=9, color='#555')
        for sp in ['top', 'right']:
            ax.spines[sp].set_visible(False)
        ax.spines['left'].set_color('#D0D5DD')
        ax.spines['bottom'].set_color('#D0D5DD')
        ax.tick_params(axis='x', labelsize=8.5)
        ax.tick_params(axis='y', length=0)
        ax.set_facecolor('white')
        fig.patch.set_facecolor('white')
        ax.grid(axis='x', alpha=0.3, color='#D0D5DD', zorder=0)
        plt.tight_layout(pad=1.4)
        buf = io.BytesIO()
        fig.savefig(buf, format='png', dpi=180, bbox_inches='tight')
        plt.close(fig)
        buf.seek(0)
        return buf

    if _mpl_ok and ff_methods:
        try:
            chart_buf = _make_ff_buf()
            n_ff = len(ff_methods)
            img_h_cm = min(max(3.8, n_ff * 1.15), 6.8)
            from pptx.util import Cm
            img_top  = Cm(2.54)
            img_left = Cm(1.02)
            img_w    = Cm(23.37)
            img_h    = Cm(img_h_cm)
            slide.shapes.add_picture(chart_buf, img_left, img_top, img_w, img_h)
            commentary_y = 2.54 + img_h_cm + 0.35
        except Exception as e:
            log.warning("FF chart embed failed: %s", e)
            _ff_fallback(slide, ff, price, cur_sym, currency)
            commentary_y = 9.80
    else:
        _ff_fallback(slide, ff, price, cur_sym, currency)
        commentary_y = 9.80

    _ff_c1 = _g(synthesis, "dcf_commentary", "") or ""
    _ff_c2 = _g(synthesis, "valuation_comment", "") or ""
    dcf_comment = " ".join(c for c in [_ff_c1, _ff_c2] if c.strip())
    if not dcf_comment.strip():
        dcf_comment = _g(synthesis, "peers_commentary", "") or ""
    if dcf_comment.strip():
        _ff_h = min(4.0, 13.25 - commentary_y)
        _ff_title = _jpm_title("football", ratios=ratios, snap=snap, synthesis=synthesis)
        commentary_box(slide, 1.02, commentary_y, 23.37, _ff_h, dcf_comment, title=_ff_title)

    return slide


def _ff_fallback(slide, ff, price, cur_sym, currency):
    """Fallback table when matplotlib unavailable."""
    header = ["Méthode", "Fourchette basse", "Fourchette haute", "Point central"]
    row_fills_map = {
        "DCF - Bear": RED_PALE,
        "DCF - Base": NAVY_PALE,
        "DCF - Bull": GREEN_PALE,
    }
    rows_data  = []
    rows_fills = []
    for item in ff:
        label = str(_g(item, "label", "—"))
        low   = _g(item, "range_low")
        high  = _g(item, "range_high")
        mid   = _g(item, "midpoint")
        rows_data.append([label, _fr(low, 0), _fr(high, 0), _fr(mid, 0)])
        rows_fills.append(row_fills_map.get(label, GREY_BG))
    rows_data.append([f"Cours actuel ({_fr(price, 2)})", "—", "—", _fr(price, 2)])
    rows_fills.append(NAVY_PALE)
    if not rows_data:
        rows_data  = [["Aucune donnée disponible", "—", "—", "—"]]
        rows_fills = [GREY_BG]
    tbl_h = min(1.0 + len(rows_data) * 0.71, 6.0)
    add_table(slide, 1.02, 2.54, 23.37, tbl_h,
              len(rows_data), 4,
              col_widths_pct=[0.35, 0.20, 0.20, 0.25],
              header_data=header,
              rows_data=rows_data,
              row_fills=rows_fills)


# ---------------------------------------------------------------------------
# Slide IB/PE — Multiples Historiques 5 ans
# ---------------------------------------------------------------------------

def _slide_multiples_historiques(prs, snap, synthesis, ratios):
    """P/E et EV/EBITDA sur 5 ans — line chart matplotlib + commentary."""
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    navy_bar(slide)
    footer_bar(slide)
    section_dots(slide, 3)

    ci       = snap.company_info if snap else None
    ticker   = _g(ci, "ticker", "—")
    currency = _g(ci, "currency", "USD") or "USD"

    slide_title(slide, "Multiples de Valorisation Historiques",
                f"{ticker} \u00b7 P/E \u0026 EV/EBITDA sur 5 ans \u00b7 Standards march\u00e9")

    # ── Collecte des données historiques ─────────────────────────────────────
    years_data = []
    if ratios and getattr(ratios, "years", None):
        _yrs = ratios.years
        for lbl in sorted(_yrs.keys(), key=lambda k: str(k).replace("_LTM","Z")):
            yr = _yrs[lbl]
            pe   = getattr(yr, "pe_ratio",  None)
            ev_eb = getattr(yr, "ev_ebitda", None)
            pb   = getattr(yr, "pb_ratio",  None)
            # Exclure années sans Données utiles (yfinance manque de Données historiques lointaines)
            if pe is None and ev_eb is None and pb is None:
                continue
            years_data.append({
                "label": lbl.replace("_LTM", " LTM"),
                "pe":    pe,
                "ev_eb": ev_eb,
                "pb":    pb,
            })
    years_data = years_data[-5:]  # 5 dernières années max

    labels  = [d["label"]  for d in years_data]
    pe_vals = [d["pe"]     for d in years_data]
    ev_vals = [d["ev_eb"]  for d in years_data]
    pb_vals = [d["pb"]     for d in years_data]

    pe_clean  = [v for v in pe_vals  if v is not None]
    ev_clean  = [v for v in ev_vals  if v is not None]
    has_data  = len(pe_clean) > 1 or len(ev_clean) > 1

    # ── Matplotlib single-axis chart (P/E + EV/EBITDA sur même échelle) ──────
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import numpy as np

        fig, ax1 = plt.subplots(figsize=(11.5, 3.8))

        x = np.arange(len(labels))
        pe_plot = [v if v is not None else np.nan for v in pe_vals]
        ev_plot = [v if v is not None else np.nan for v in ev_vals]

        if any(v == v for v in pe_plot):  # has non-nan
            ax1.plot(x, pe_plot, color="#1B3A6B", linewidth=2.2, marker='o',
                     markersize=6, label="P/E", zorder=4)
            ax1.fill_between(x, pe_plot, alpha=0.08, color="#1B3A6B")

        if any(v == v for v in ev_plot):
            ax1.plot(x, ev_plot, color="#1A7A4A", linewidth=2.2, marker='s',
                     markersize=6, linestyle='--', label="EV/EBITDA", zorder=4)

        # Auto-scale y pour inclure toutes les séries avec marge
        _all_vals = [v for v in pe_plot + ev_plot if v == v]
        if _all_vals:
            _margin = (max(_all_vals) - min(_all_vals)) * 0.12 or 2.0
            ax1.set_ylim(max(0, min(_all_vals) - _margin), max(_all_vals) + _margin)

        ax1.set_xticks(x)
        ax1.set_xticklabels(labels, fontsize=9)
        ax1.set_ylabel("Multiple (x)", fontsize=9, color="#333")
        ax1.tick_params(axis='y', labelsize=8)
        ax1.tick_params(axis='x', labelsize=8.5)
        ax1.legend(fontsize=9, loc='upper center', bbox_to_anchor=(0.5, 1.12),
                   ncol=2, framealpha=0.95, edgecolor='#DDDDDD',
                   handlelength=1.8, handletextpad=0.5)

        ax1.spines['top'].set_visible(False)
        ax1.spines['right'].set_visible(False)
        ax1.spines['left'].set_color('#D0D5DD')
        ax1.spines['bottom'].set_color('#D0D5DD')
        ax1.set_facecolor('white')
        fig.patch.set_facecolor('white')
        ax1.grid(axis='y', alpha=0.25, color='#D0D5DD', zorder=0)

        # Titre analytique : tendance observée
        if pe_clean and len(pe_clean) >= 2:
            delta_pe = pe_clean[-1] - pe_clean[0]
            _dir = "expansion" if delta_pe > 0 else "compression"
            ax1.set_title(
                f"P/E : {_dir} de {abs(delta_pe):.1f}x sur la p\u00e9riode  |  "
                f"EV/EBITDA actuel : {_fr(ev_clean[-1] if ev_clean else None, 1)}x",
                fontsize=9, color='#333', pad=6
            )

        plt.tight_layout(pad=1.2)
        buf_mc = io.BytesIO()
        fig.savefig(buf_mc, format='png', dpi=180, bbox_inches='tight')
        plt.close(fig)
        buf_mc.seek(0)

        from pptx.util import Cm
        img_h_cm = 7.0
        slide.shapes.add_picture(buf_mc, Cm(1.02), Cm(2.40), Cm(15.80), Cm(img_h_cm))
        _chart_ok = True
    except Exception as _e:
        log.warning("multiples_historiques chart: %s", _e)
        _chart_ok = False

    # ── Tableau synthèse valeurs — droite ────────────────────────────────────
    tbl_x, tbl_y, tbl_w = 17.20, 2.40, 7.20
    add_text_box(slide, tbl_x, tbl_y, tbl_w, 0.50,
                 "SYNTHÈSE DES MULTIPLES", 8, NAVY, bold=True)
    if years_data:
        snap_rows = []
        for d in years_data:
            snap_rows.append([
                d["label"],
                _fr(d["pe"],   1) + "x" if d["pe"]   is not None else "—",
                _fr(d["ev_eb"],1) + "x" if d["ev_eb"] is not None else "—",
                _fr(d["pb"],   1) + "x" if d["pb"]    is not None else "—",
            ])
        add_table(slide, tbl_x, tbl_y + 0.65, tbl_w, min(0.55 * len(snap_rows) + 0.60, 6.5),
                  len(snap_rows), 4,
                  col_widths_pct=[0.30, 0.22, 0.26, 0.22],
                  header_data=["Année", "P/E", "EV/EB", "P/B"],
                  rows_data=snap_rows,
                  border_hex="DDDDDD")

    # ── Commentary analytique ─────────────────────────────────────────────────
    _pe_last   = pe_clean[-1]  if pe_clean  else None
    _ev_last   = ev_clean[-1]  if ev_clean  else None
    _pe_first  = pe_clean[0]   if pe_clean  else None
    _ev_first  = ev_clean[0]   if ev_clean  else None
    if _pe_last and _pe_first:
        _delta = _pe_last - _pe_first
        _dir   = "expansion multiple" if _delta > 0 else "compression multiple"
        _comment = (
            f"Le P/E affiche une {_dir} de {abs(_delta):.1f}x sur la p\u00e9riode ({_fr(_pe_first,1)}x \u2192 {_fr(_pe_last,1)}x), "
            f"signal d\u2019une r\u00e9\u00e9valuation {'favorable' if _delta > 0 else 'négative'} du profil de croissance. "
            f"L\u2019EV/EBITDA {'se stabilise' if _ev_last and _ev_first and abs(_ev_last-_ev_first)<2 else 'diverge'} "
            f"\u00e0 {_fr(_ev_last,1)}x vs {_fr(_ev_first,1)}x historique."
        )
    else:
        _comment = _g(synthesis, "ratio_commentary", "") or "Donn\u00e9es historiques insuffisantés pour \u00e9tablir une tendance des multiples."
    _mh_c2 = _g(synthesis, "financial_commentary", "") or ""
    if _mh_c2.strip():
        _comment = _comment + " " + _mh_c2 if _comment.strip() else _mh_c2
    _mh_title = _jpm_title("ratios", ratios=ratios, snap=snap, synthesis=synthesis)
    commentary_box(slide, 1.02, 10.00, 23.37, 3.10, _comment, title=_mh_title)

    return slide


# ---------------------------------------------------------------------------
# Slide IB/PE — Capital Returns & FCF
# ---------------------------------------------------------------------------

def _slide_capital_returns(prs, snap, synthesis, ratios):
    """FCF yield, dividendes et retour total aux actionnaires — niveau IB."""
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    navy_bar(slide)
    footer_bar(slide)
    section_dots(slide, 3)

    ci       = snap.company_info if snap else None
    mkt      = snap.market if snap else None
    ticker   = _g(ci, "ticker", "—")
    currency = _g(ci, "currency", "USD") or "USD"
    cur_sym  = "EUR" if currency == "EUR" else "$"
    price    = _g(mkt, "share_price")

    slide_title(slide, "Capital Returns \u0026 Free Cash Flow",
                f"{ticker} \u00b7 FCF yield, dividendes & allocation du capital")

    # ── Collecte données ─────────────────────────────────────────────────────
    cap_rows = []
    if ratios and getattr(ratios, "years", None):
        for lbl in sorted(ratios.years.keys(), key=lambda k: str(k).replace("_LTM","Z")):
            yr = ratios.years[lbl]
            fcf       = getattr(yr, "fcf", None)
            fcf_yield = getattr(yr, "fcf_yield", None)
            div_pay   = getattr(yr, "dividends_paid_abs", None)
            ebitda    = getattr(yr, "ebitda", None)
            capex_r   = getattr(yr, "capex_ratio", None)
            div_pout  = getattr(yr, "dividend_payout", None)
            # Exclure années sans Données utiles
            if all(v is None for v in (fcf, fcf_yield, ebitda)):
                continue
            cap_rows.append({
                "label":     lbl.replace("_LTM", " LTM"),
                "fcf":       fcf,
                "fcf_yield": fcf_yield,
                "div_pay":   div_pay,
                "ebitda":    ebitda,
                "capex_r":   capex_r,
                "div_pout":  div_pout,
            })
    cap_rows = cap_rows[-5:]

    # ── KPI boxes top ─────────────────────────────────────────────────────────
    latest_yr = cap_rows[-1] if cap_rows else {}
    _kpis = [
        ("FCF Yield",     _frpct(latest_yr.get("fcf_yield"))  if latest_yr.get("fcf_yield") else "—"),
        ("FCF (" + ("LTM" if cap_rows else "N/A") + ")",
                          _frm(latest_yr.get("fcf"), cur_sym) if latest_yr.get("fcf") else "—"),
        ("Capex / Rev",   _frpct(latest_yr.get("capex_r"))    if latest_yr.get("capex_r") else "—"),
        ("Div. Payout",   _frpct(latest_yr.get("div_pout"))   if latest_yr.get("div_pout") is not None else "0%"),
    ]
    kpi_w = 5.60
    for i, (lbl, val) in enumerate(_kpis):
        bx = 1.02 + i * (kpi_w + 0.18)
        add_rect(slide, bx, 2.29, kpi_w, 1.52, NAVY_PALE)
        add_rect(slide, bx, 2.29, 0.18, 1.52, NAVY)
        add_text_box(slide, bx + 0.26, 2.42, kpi_w - 0.32, 0.65,
                     val, 16, NAVY, bold=True)
        add_text_box(slide, bx + 0.26, 3.05, kpi_w - 0.32, 0.65,
                     lbl, 8, GREY_TXT)

    # ── FCF bar chart ──────────────────────────────────────────────────────────
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import numpy as np

        labels_c = [d["label"] for d in cap_rows]
        fcf_c    = [d["fcf"]   if d["fcf"] is not None else 0 for d in cap_rows]
        div_c    = [abs(d["div_pay"]) if d["div_pay"] is not None else 0 for d in cap_rows]

        fig, ax = plt.subplots(figsize=(9.0, 3.2))
        x = np.arange(len(labels_c))
        bw = 0.35
        bars_fcf = ax.bar(x - bw/2, [f/1000 if f else 0 for f in fcf_c], bw,
                          label="FCF (Mds)", color="#1A7A4A", alpha=0.85)
        bars_div = ax.bar(x + bw/2, [d/1000 if d else 0 for d in div_c], bw,
                          label="Dividendes versés (Mds)", color="#2E5FA3", alpha=0.85)

        ax.set_xticks(x)
        ax.set_xticklabels(labels_c, fontsize=8.5)
        ax.set_ylabel(f"Mds {cur_sym}", fontsize=8.5)
        ax.legend(fontsize=8, loc='upper left', framealpha=0.9)
        ax.tick_params(labelsize=8)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.spines['left'].set_color('#D0D5DD')
        ax.spines['bottom'].set_color('#D0D5DD')
        ax.set_facecolor('white')
        fig.patch.set_facecolor('white')
        ax.grid(axis='y', alpha=0.25, color='#D0D5DD')

        # Valeurs au dessus des barres
        for bar in bars_fcf:
            h = bar.get_height()
            if h != 0:
                ax.text(bar.get_x() + bar.get_width()/2, h + abs(h)*0.02,
                        f"{h:.1f}", ha='center', va='bottom', fontsize=7.5, color='#1A7A4A')
        for bar in bars_div:
            h = bar.get_height()
            if h != 0:
                ax.text(bar.get_x() + bar.get_width()/2, h + abs(h)*0.02,
                        f"{h:.1f}", ha='center', va='bottom', fontsize=7.5, color='#2E5FA3')

        plt.tight_layout(pad=1.0)
        buf_cr = io.BytesIO()
        fig.savefig(buf_cr, format='png', dpi=180, bbox_inches='tight')
        plt.close(fig)
        buf_cr.seek(0)

        from pptx.util import Cm
        slide.shapes.add_picture(buf_cr, Cm(1.02), Cm(4.20), Cm(14.50), Cm(5.80))
    except Exception as _e:
        log.warning("capital_returns chart: %s", _e)

    # ── Tableau allocation du capital — droite ────────────────────────────────
    add_text_box(slide, 16.20, 4.20, 8.20, 0.50,
                 "ALLOCATION DU CAPITAL", 8, NAVY, bold=True)
    if cap_rows:
        alloc_rows = []
        for d in cap_rows:
            alloc_rows.append([
                d["label"],
                _frm(d["fcf"],   cur_sym) if d["fcf"]   is not None else "—",
                _frpct(d["fcf_yield"])    if d["fcf_yield"] is not None else "—",
                _frpct(d["div_pout"])     if d["div_pout"]  is not None else "0%",
                _frpct(d["capex_r"])      if d["capex_r"]   is not None else "—",
            ])
        add_table(slide, 16.20, 4.85, 8.20, min(0.55 * len(alloc_rows) + 0.60, 5.0),
                  len(alloc_rows), 5,
                  col_widths_pct=[0.22, 0.20, 0.19, 0.19, 0.20],
                  header_data=["Ann.", "FCF", "Yield", "Div%", "Capex%"],
                  rows_data=alloc_rows,
                  border_hex="DDDDDD")

    # ── Commentary ───────────────────────────────────────────────────────────
    fcf_vals = [d["fcf"] for d in cap_rows if d["fcf"] is not None]
    fy_vals  = [d["fcf_yield"] for d in cap_rows if d["fcf_yield"] is not None]
    if fcf_vals and len(fcf_vals) >= 2:
        _delta_fcf = fcf_vals[-1] - fcf_vals[0]
        _dir       = "croissance" if _delta_fcf > 0 else "contraction"
        _fy_str    = _frpct(fy_vals[-1]) if fy_vals else "—"
        _comment   = (
            f"Le FCF affiche une {_dir} de {_frm(abs(_delta_fcf), cur_sym)} sur la p\u00e9riode, "
            f"avec un FCF yield actuel de {_fy_str} — "
            f"{'attractif pour un acheteur institutionnel' if fy_vals and fy_vals[-1] and float(fy_vals[-1]) > 0.04 else 'à surveiller au regard de la valorisation de marché'}. "
            f"L\u2019allocation vers les dividendes "
            f"{'est maintenue malgré la pression sur les marges' if fcf_vals[-1] and fcf_vals[-1] > 0 else 'reflète un levier financier élevé'}."
        )
    else:
        _comment = "Donn\u00e9es FCF insuffisantés pour \u00e9tablir une tendance de l\u2019allocation du capital sur la p\u00e9riode."
    _cr_c2 = _g(synthesis, "financial_commentary", "") or ""
    if _cr_c2.strip():
        _comment = _comment + " " + _cr_c2 if _comment.strip() else _cr_c2
    _cr_title = _jpm_title("capital", ratios=ratios, snap=snap, synthesis=synthesis)
    commentary_box(slide, 1.02, 10.30, 23.37, 2.80, _comment, title=_cr_title)

    return slide


# ---------------------------------------------------------------------------
# Slide IB/PE — Analyse LBO
# ---------------------------------------------------------------------------

def _slide_lbo(prs, snap, synthesis, ratios):
    """Analyse LBO — entry/exit multiples, dette, IRR/MOIC — niveau PE."""
    import math as _math_lbo
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    navy_bar(slide)
    footer_bar(slide)
    section_dots(slide, 3)

    ci       = snap.company_info if snap else None
    mkt      = snap.market if snap else None
    ticker   = _g(ci, "ticker", "—")
    currency = _g(ci, "currency", "USD") or "USD"
    cur_sym  = "EUR" if currency == "EUR" else "$"

    slide_title(slide, "Analyse LBO \u2014 Leveraged Buyout",
                f"{ticker} \u00b7 Attractivit\u00e9 PE \u00b7 Multiples d\u2019entr\u00e9e vs IRR cible")

    # ── Données LBO ───────────────────────────────────────────────────────────
    latest_yr = None
    if ratios and getattr(ratios, "years", None):
        _lbl = sorted(ratios.years.keys(), key=lambda k: str(k).replace("_LTM","Z"))[-1]
        latest_yr = ratios.years[_lbl]

    ebitda_raw   = getattr(latest_yr, "ebitda",    None) if latest_yr else None
    net_debt_raw = getattr(latest_yr, "net_debt",  None) if latest_yr else None
    fcf_raw      = getattr(latest_yr, "fcf",       None) if latest_yr else None
    rev_growth   = getattr(latest_yr, "revenue_growth", None) if latest_yr else None

    ebitda   = float(ebitda_raw)   if ebitda_raw   is not None else None
    net_debt = float(net_debt_raw) if net_debt_raw is not None else None
    fcf      = float(fcf_raw)      if fcf_raw      is not None else None

    # Croissance EBITDA projetée (5 ans) : revenue_growth ou 5% par défaut
    _g_rate = float(rev_growth) if rev_growth and rev_growth > 0.001 else 0.05
    _g_rate = min(_g_rate, 0.20)  # cap 20%

    # ── Hypothèses LBO ────────────────────────────────────────────────────────
    entry_multiples = [8.0, 10.0, 12.0, 14.0]
    exit_multiples  = [8.0, 10.0, 12.0]
    leverage_ratio  = 4.0  # Dette = 4x EBITDA (hypothèse PE classique)
    hold_years      = 5
    debt_repay_pct  = 0.50  # 50% de la dette remboursée en 5 ans

    # ── Table IRR ─────────────────────────────────────────────────────────────
    def _lbo_irr(entry_mult, exit_mult):
        if ebitda is None or ebitda <= 0:
            return None, None
        entry_ev     = ebitda * entry_mult
        entry_debt   = min(ebitda * leverage_ratio, entry_ev * 0.65)
        entry_equity = max(entry_ev - entry_debt, 1)

        # EBITDA projeté dans 5 ans
        ebitda_exit  = ebitda * ((1 + _g_rate) ** hold_years)
        exit_ev      = ebitda_exit * exit_mult
        exit_debt    = entry_debt * (1 - debt_repay_pct)
        exit_equity  = max(exit_ev - exit_debt, 0.01)

        moic = exit_equity / entry_equity
        irr  = moic ** (1 / hold_years) - 1
        return irr, moic

    # Tableau entry × exit : IRR
    irr_rows = []
    for em in entry_multiples:
        row = [f"{em:.0f}x"]
        for xm in exit_multiples:
            irr, moic = _lbo_irr(em, xm)
            if irr is not None:
                row.append(f"{irr*100:.1f}% / {moic:.1f}x")
            else:
                row.append("—")
        irr_rows.append(row)

    exit_hdrs = [f"Exit {xm:.0f}x" for xm in exit_multiples]
    tbl_irr = add_table(
        slide, 1.02, 2.50, 13.50, min(0.65 * len(irr_rows) + 0.65, 4.50),
        len(irr_rows), len(exit_multiples) + 1,
        col_widths_pct=[0.22] + [0.26] * len(exit_multiples),
        header_data=["Entry mult."] + exit_hdrs,
        rows_data=irr_rows,
        border_hex="DDDDDD"
    )

    # ── Color-code les cellules selon IRR ≥ 20% = vert, 15-20% = ambre, <15% = rouge ──
    try:
        for ri, em in enumerate(entry_multiples):
            for ci_idx, xm in enumerate(exit_multiples):
                irr, _ = _lbo_irr(em, xm)
                cell = tbl_irr.cell(ri + 1, ci_idx + 1)
                cell.fill.solid()
                if irr is not None:
                    if irr >= 0.20:
                        cell.fill.fore_color.rgb = rgb(GREEN_PALE)
                    elif irr >= 0.15:
                        cell.fill.fore_color.rgb = rgb(AMBER_PALE)
                    else:
                        cell.fill.fore_color.rgb = rgb(RED_PALE)
    except Exception:
        pass

    # Note sous le tableau
    add_text_box(slide, 1.02, 7.20, 13.50, 0.50,
                 f"Hypotheses : Levier {leverage_ratio:.0f}x EBITDA \u00b7 {hold_years} ans holding "
                 f"\u00b7 {debt_repay_pct*100:.0f}% remboursement dette \u00b7 EBITDA +{_g_rate*100:.0f}%/an",
                 7.5, GREY_TXT)
    add_text_box(slide, 1.02, 7.65, 13.50, 0.50,
                 "Lecture : IRR / MOIC  \u25cf vert \u2265 20%  \u25cf ambre 15-20%  \u25cf rouge < 15%",
                 7.5, GREY_LIGHT)

    # ── Hypothèses clés — droite ──────────────────────────────────────────────
    add_text_box(slide, 15.20, 2.50, 9.20, 0.50,
                 "HYPOTHÈSES & STRUCTURE LBO", 8, NAVY, bold=True)

    _hypo_rows = [
        ["EBITDA LTM",       _frm(ebitda, cur_sym) if ebitda else "—"],
        ["FCF LTM",          _frm(fcf, cur_sym)    if fcf    else "—"],
        ["Levier PE",        f"{leverage_ratio:.0f}x EBITDA"],
        ["Dette entr\u00e9e",  _frm(ebitda * leverage_ratio if ebitda else None, cur_sym)],
        ["EBITDA exit (est.)",_frm(ebitda * ((1+_g_rate)**hold_years) if ebitda else None, cur_sym)],
        ["Holding",          f"{hold_years} ans"],
        ["IRR cible PE",     "\u2265 20%"],
    ]
    add_table(slide, 15.20, 3.10, 9.20, min(0.55 * len(_hypo_rows) + 0.60, 5.50),
              len(_hypo_rows), 2,
              col_widths_pct=[0.52, 0.48],
              header_data=["Param\u00e8tre", "Valeur"],
              rows_data=_hypo_rows,
              border_hex="DDDDDD")

    # ── Note méthodologique ───────────────────────────────────────────────────
    add_rect(slide, 15.20, 8.60, 9.20, 2.80, AMBER_PALE)
    add_rect(slide, 15.20, 8.60, 0.20, 2.80, AMBER)
    add_text_box(slide, 15.50, 8.70, 8.80, 0.45,
                 "NOTE MÉTHODOLOGIQUE", 7.5, AMBER, bold=True)
    add_text_box(slide, 15.50, 9.15, 8.80, 2.20,
                 "Analyse indicative uniquement. Le LBO assume un rachat "
                 "100% equity par un fonds PE. La dette est remboursée par "
                 "le FCF de la cible. L'IRR ne tient pas compte des frais "
                 "de transaction ni de la structure fiscale optimisée. "
                 "Un IRR \u2265 20% est le seuil standard d'un fonds PE tier-1.",
                 8, GREY_TXT)

    # ── Commentary — agent LLM dedie si disponible ──────────────────────────
    irr_base, moic_base = _lbo_irr(10.0, 10.0)
    irr_bull, _         = _lbo_irr(8.0, 12.0)
    irr_bear, _         = _lbo_irr(14.0, 8.0)

    _comment = ""
    try:
        from agents.agent_lbo import generate_lbo_texts
        # Construire un dict lbo_data minimal a partir du calcul Python
        # (l'agent fera lui-meme ses propres references si besoin)
        _eligible = (
            ebitda is not None and ebitda > 0
            and (rev_growth or 0) >= 0
            and (net_debt is None or (net_debt / max(ebitda, 1)) < 4.0)
        )
        _lbo_data_min = {
            "eligible": _eligible,
            "mega_flag": "standard",
            "irr_base": irr_base,
            "moic_base": moic_base,
            "irr_bull": irr_bull,
            "irr_bear": irr_bear,
            "leverage_exit": leverage_ratio * (1 - debt_repay_pct),
            "equity_entry": ebitda * 10.0 - (ebitda * leverage_ratio) if ebitda else None,
            "equity_exit": (ebitda * ((1+_g_rate)**hold_years) * 10.0 - ebitda * leverage_ratio * (1-debt_repay_pct)) if ebitda else None,
            "company_name": _g(ci, "company_name", ticker),
            "ticker": ticker,
        }
        _metrics_min = {
            "ebitda_margin_ltm": getattr(latest_yr, "ebitda_margin", None) if latest_yr else None,
            "roic": getattr(latest_yr, "roic", None) if latest_yr else None,
            "net_debt_ebitda": (net_debt / max(ebitda, 1)) if (net_debt and ebitda) else None,
            "company_name_a": _g(ci, "company_name", ticker),
            "ticker_a": ticker,
            "sector": _g(ci, "sector", ""),
        }
        _texts = generate_lbo_texts(_lbo_data_min, _metrics_min)
        _comment = _texts.get("returns_text") or _texts.get("eligibility_text") or ""
    except Exception as _e:
        _comment = ""

    if not _comment and irr_base is not None:
        _lbo_signal = "attractive" if irr_base >= 0.20 else ("limite" if irr_base >= 0.15 else "insuffisanté")
        _comment = (
            f"À 10x EBITDA d'entrée / 10x de sortie, le LBO génère un IRR de {irr_base*100:.1f}% "
            f"(MOIC {moic_base:.1f}x en {hold_years} ans) — attractivité PE jugée {_lbo_signal}. "
            f"Un fonds tier-1 ciblant 20%+ d'IRR doit entrer en dessous de {next((em for em in entry_multiples if (_lbo_irr(em,10.0)[0] or 0)>=0.20), entry_multiples[-1]):.0f}x EBITDA. "
            f"La génération de FCF ({_frm(fcf, cur_sym) if fcf else '—'}) est le principal levier de remboursement de la dette."
        )
    if not _comment:
        _comment = (
            "EBITDA non disponible — analyse LBO indicative impossible. "
            "Veuillez relancer l'analyse avec des données financières complètes."
        )
    _lbo_one_title = _jpm_title("lbo", ratios=ratios, snap=snap, synthesis=synthesis,
                                extra={"irr_base": irr_base})
    commentary_box(slide, 1.02, 11.60, 23.37, 1.60, _comment, title=_lbo_one_title)

    return slide


# ═════════════════════════════════════════════════════════════════════════════
# 3 SLIDES LBO INSTITUTIONNELLES (cadre / returns / stress)
# ═════════════════════════════════════════════════════════════════════════════

def _build_lbo_pack(snap, synthesis, ratios) -> dict:
    """Construit le pack de données LBO + textes LLM pour les 3 slides.

    Réplique simplifiée de la mécanique du modèle Excel `outputs/lbo_model.py`,
    mais en pur Python pour ne pas dépendre du recalcul Excel.

    Returns:
        dict avec : eligible, mega_flag, hypotheses, sources_uses, returns,
        sensitivity (5×5), scenarios (bull/base/bear), debt_schedule, llm_texts
    """
    # ── Extraction métriques ──
    ci = snap.company_info if snap else None
    ticker = _g(ci, "ticker", "—")
    company = _g(ci, "company_name", ticker)
    sector = _g(ci, "sector", "")
    currency = _g(ci, "currency", "USD") or "USD"
    cur_sym = "EUR" if currency == "EUR" else "$"

    latest_yr = None
    if ratios and getattr(ratios, "years", None):
        _lbl = sorted(ratios.years.keys(),
                      key=lambda k: str(k).replace("_LTM", "Z"))[-1]
        latest_yr = ratios.years[_lbl]

    ebitda = float(getattr(latest_yr, "ebitda", 0) or 0) if latest_yr else 0.0
    ebit = float(getattr(latest_yr, "ebit", ebitda * 0.85) or 0) if latest_yr else 0.0
    fcf = float(getattr(latest_yr, "fcf", 0) or 0) if latest_yr else 0.0
    net_debt = float(getattr(latest_yr, "net_debt", 0) or 0) if latest_yr else 0.0
    rev = float(getattr(latest_yr, "revenue", 0) or 0) if latest_yr else 0.0
    rev_growth = getattr(latest_yr, "revenue_growth", None) if latest_yr else None
    ebitda_margin = float(getattr(latest_yr, "ebitda_margin", 0) or 0) if latest_yr else 0.0
    roic = getattr(latest_yr, "roic", None) if latest_yr else None

    # Market cap & EV
    mkt = snap.market if snap else None
    market_cap = float(_g(mkt, "market_cap", 0) or 0)
    ev_market = market_cap + net_debt

    # ── Critères d'éligibilité ──
    cash_conv = (fcf / ebitda) if ebitda > 0 else 0
    nd_ebitda = (net_debt / ebitda) if ebitda > 0 else 99.0
    ebitda_marg_pct = ebitda_margin * 100 if abs(ebitda_margin) <= 2 else ebitda_margin

    elig_marge = ebitda_marg_pct > 15
    elig_cash = cash_conv > 0.6
    elig_levier = nd_ebitda < 3.5
    elig_ebitda = ebitda > 0
    eligible = all([elig_marge, elig_cash, elig_levier, elig_ebitda])

    # Mega-deal flag (en millions $)
    if ev_market > 100_000:
        mega_flag = "theorique"
    elif ev_market > 50_000:
        mega_flag = "mega"
    else:
        mega_flag = "standard"

    # ── Hypothèses LBO standards ──
    multiple_marche = ev_market / ebitda if ebitda > 0 else 14.0
    # Multiple d'entrée = max(14, multiple_Marché * 1.1) — réaliste
    entry_mult = max(14.0, round(multiple_marche * 1.1, 1)) if multiple_marche > 0 else 14.0
    leverage = 5.0       # Total leverage
    senior_lvg = 3.5
    mezz_lvg = 1.5
    senior_rate = 0.08
    mezz_rate = 0.10
    exit_mult = entry_mult  # conservatisme
    hold_years = 5
    fees_pct = 0.025
    cash_min_pct = 0.05

    # ── Sources & Uses ──
    ev_deal = entry_mult * ebitda
    senior_debt = senior_lvg * ebitda
    mezz_debt = mezz_lvg * ebitda
    purchase_eq = ev_deal - net_debt
    refi = max(net_debt, 0)
    fees_amount = ev_deal * fees_pct
    min_cash = ev_deal * cash_min_pct
    total_uses = purchase_eq + refi + fees_amount + min_cash
    sponsor_eq = total_uses - senior_debt - mezz_debt
    total_sources = sponsor_eq + senior_debt + mezz_debt

    # ── Projections + Returns base case ──
    g_rate = float(rev_growth) if rev_growth and rev_growth > 0.001 else 0.05
    g_rate = min(max(g_rate, 0.0), 0.20)
    ebitda_y5 = ebitda * ((1 + g_rate) ** hold_years)

    # Debt schedule simplifié : senior remboursée à 50% sur 5 ans, mezz bullet
    senior_y5 = senior_debt * 0.5
    mezz_y5 = mezz_debt * ((1 + 0.05) ** hold_years)  # PIK 5%/an
    cash_y5 = (fcf if fcf > 0 else ebitda * 0.5) * hold_years * 0.3  # 30% du FCF capitalisé
    net_debt_y5 = max(senior_y5 + mezz_y5 - cash_y5, 0)
    leverage_y5 = net_debt_y5 / ebitda_y5 if ebitda_y5 > 0 else 0

    ev_exit = ebitda_y5 * exit_mult
    equity_exit = ev_exit - net_debt_y5
    moic_base = (equity_exit / sponsor_eq) if sponsor_eq > 0 else 0
    irr_base = (moic_base ** (1 / hold_years) - 1) if moic_base > 0 else -0.99

    # ── Sensibilité 5×5 (entry × exit) ──
    sens_grid = []
    entry_offsets = [-2, -1, 0, 1, 2]
    exit_offsets = [-2, -1, 0, 1, 2]
    for off_e in entry_offsets:
        row = []
        em = entry_mult + off_e
        if em <= 0:
            row = [None] * 5
            sens_grid.append(row)
            continue
        ev_d = em * ebitda
        senior_d = senior_lvg * ebitda
        mezz_d = mezz_lvg * ebitda
        purchase_d = ev_d - net_debt
        total_u = purchase_d + max(net_debt, 0) + ev_d * fees_pct + ev_d * cash_min_pct
        sponsor_d = total_u - senior_d - mezz_d
        for off_x in exit_offsets:
            xm = exit_mult + off_x
            if xm <= 0 or sponsor_d <= 0:
                row.append(None)
                continue
            ev_x = ebitda_y5 * xm
            eq_x = ev_x - net_debt_y5
            moic = eq_x / sponsor_d if sponsor_d > 0 else 0
            irr = (moic ** (1 / hold_years) - 1) if moic > 0 else -0.99
            row.append(irr)
        sens_grid.append(row)

    # ── Scénarios bull/base/bear ──
    def _scenario(rev_adj, marg_adj, exit_adj):
        eb = ebitda * (1 + marg_adj)
        eb_y5 = eb * ((1 + g_rate + rev_adj) ** hold_years)
        ev_x = eb_y5 * (exit_mult + exit_adj)
        nd_y5 = max(senior_y5 + mezz_y5 - cash_y5, 0)
        eq_x = ev_x - nd_y5
        moic = eq_x / sponsor_eq if sponsor_eq > 0 else 0
        irr = (moic ** (1 / hold_years) - 1) if moic > 0 else -0.99
        lvg_x = nd_y5 / eb_y5 if eb_y5 > 0 else 0
        icr = (eb_y5 * 0.85) / max(senior_y5 * senior_rate + mezz_y5 * 0.05, 1)
        return {"irr": irr, "moic": moic, "leverage": lvg_x, "icr": icr}

    scenarios = {
        "bull": _scenario(0.03, 0.02, 1),
        "base": {"irr": irr_base, "moic": moic_base,
                 "leverage": leverage_y5,
                 "icr": (ebit * ((1 + g_rate) ** hold_years)) / max(senior_y5 * senior_rate + mezz_y5 * 0.05, 1)},
        "bear": _scenario(-0.03, -0.02, -1),
    }

    # ── Debt schedule simplifié pour affichage ──
    debt_schedule = []
    senior_bal = senior_debt
    mezz_bal = mezz_debt
    cash_bal = 0
    for yr in range(hold_years + 1):
        ebitda_yr = ebitda * ((1 + g_rate) ** yr)
        if yr > 0:
            interest_senior = senior_bal * senior_rate
            interest_mezz = mezz_bal * 0.05
            mezz_bal *= 1.05  # PIK
            mandatory = senior_bal * 0.01
            fcf_yr = ebitda_yr * 0.85 - interest_senior - interest_mezz
            sweep = max(min((fcf_yr - mandatory) * 0.75, senior_bal - mandatory), 0)
            senior_bal = max(senior_bal - mandatory - sweep, 0)
            cash_bal += max(fcf_yr - mandatory - sweep, 0)
        net_d = senior_bal + mezz_bal - cash_bal
        debt_schedule.append({
            "year": yr, "ebitda": ebitda_yr,
            "senior": senior_bal, "mezz": mezz_bal,
            "cash": cash_bal, "net_debt": net_d,
            "leverage": net_d / ebitda_yr if ebitda_yr > 0 else 0
        })

    # ── Génération textes LLM via agent_lbo ──
    lbo_data_for_llm = {
        "eligible": eligible, "mega_flag": mega_flag,
        "irr_base": irr_base, "moic_base": moic_base,
        "irr_bull": scenarios["bull"]["irr"],
        "irr_bear": scenarios["bear"]["irr"],
        "leverage_exit": leverage_y5,
        "equity_entry": sponsor_eq,
        "equity_exit": equity_exit,
        "company_name": company, "ticker": ticker,
    }
    metrics_min = {
        "ebitda_margin_ltm": ebitda_margin,
        "roic": roic,
        "net_debt_ebitda": nd_ebitda,
        "company_name_a": company,
        "ticker_a": ticker,
        "sector": sector,
    }
    try:
        from agents.agent_lbo import generate_lbo_texts
        llm_texts = generate_lbo_texts(lbo_data_for_llm, metrics_min)
    except Exception as _e:
        llm_texts = {
            "eligibility_text": "",
            "hypotheses_text": "",
            "returns_text": "",
            "risks_text": "",
        }

    return {
        "ticker": ticker,
        "company": company,
        "sector": sector,
        "cur_sym": cur_sym,
        "eligible": eligible,
        "mega_flag": mega_flag,
        "ebitda": ebitda,
        "ebitda_y5": ebitda_y5,
        "ev_market": ev_market,
        "ev_deal": ev_deal,
        "multiple_Marché": multiple_marche,
        "entry_mult": entry_mult,
        "exit_mult": exit_mult,
        "leverage": leverage,
        "hold_years": hold_years,
        "sources_uses": {
            "sponsor_equity": sponsor_eq,
            "senior_debt": senior_debt,
            "mezz_debt": mezz_debt,
            "total_sources": total_sources,
            "purchase_equity": purchase_eq,
            "refi": refi,
            "fees": fees_amount,
            "min_cash": min_cash,
            "total_uses": total_uses,
        },
        "returns": {
            "sponsor_equity": sponsor_eq,
            "equity_exit": equity_exit,
            "irr_base": irr_base,
            "moic_base": moic_base,
        },
        "sensitivity": {
            "entry_mults": [entry_mult + o for o in entry_offsets],
            "exit_mults": [exit_mult + o for o in exit_offsets],
            "grid": sens_grid,
        },
        "scénarios": scenarios,
        "debt_schedule": debt_schedule,
        "llm_texts": llm_texts,
    }


def _slide_lbo_cadre(prs, snap, pack: dict):
    """Slide LBO 1/3 — Cadre du deal (Sources & Uses + éligibilité)."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    navy_bar(slide)
    footer_bar(slide)
    section_dots(slide, 3)

    ticker = pack["ticker"]
    company = pack["company"]
    cur = pack["cur_sym"]

    slide_title(slide, "Modèle LBO — Cadre du Deal",
                f"{company} ({ticker})  ·  Sources & Uses, hypothèses et éligibilité")

    # ── Bandeau profil + éligibilité ──
    elig_color = GREEN_PALE if pack["eligible"] else RED_PALE
    elig_accent = GREEN if pack["eligible"] else RED
    elig_label = "✓ PROFIL LBO-ABLE" if pack["eligible"] else "✗ PROFIL NON ÉLIGIBLE"
    elig_sub = ("Société éligible aux critères LBO mid-market standard"
                if pack["eligible"]
                else "Critères non satisfaits — analyse à titre théorique")
    add_rect(slide, 1.02, 2.45, 23.37, 1.20, elig_color)
    add_rect(slide, 1.02, 2.45, 0.20, 1.20, elig_accent)
    add_text_box(slide, 1.40, 2.55, 22.0, 0.55,
                 elig_label, 13, elig_accent, bold=True)
    add_text_box(slide, 1.40, 3.10, 22.0, 0.50,
                 elig_sub, 9, NAVY_MID)

    # Mega-deal flag
    if pack["mega_flag"] != "standard":
        mega_text = ("⚠⚠ Scénario théorique — EV > 100 Md$, hors marché LBO historique"
                     if pack["mega_flag"] == "theorique"
                     else "⚠ Mega-deal — EV 50-100 Md$, syndicate sponsorship requis")
        add_rect(slide, 1.02, 3.80, 23.37, 0.55, AMBER_PALE)
        add_rect(slide, 1.02, 3.80, 0.20, 0.55, AMBER)
        add_text_box(slide, 1.40, 3.88, 22.0, 0.42,
                     mega_text, 9.5, AMBER, bold=True)
        sources_y = 4.55
    else:
        sources_y = 4.00

    # ── Sources & Uses (côte à côte) ──
    su = pack["sources_uses"]
    _PP_CENTER = __import__("pptx").enum.text.PP_ALIGN.CENTER

    # Header SOURCES (gauche) — text_box centré verticalement dans le rect
    add_rect(slide, 1.02, sources_y, 11.40, 0.55, NAVY)
    add_text_box(slide, 1.02, sources_y + 0.08, 11.40, 0.42,
                 "SOURCES (financement)", 10.5, WHITE, bold=True,
                 align=_PP_CENTER, wrap=False)

    sources_rows = [
        ["Sponsor Equity", _frm(su["sponsor_equity"], cur),
         f"{(su['sponsor_equity']/su['total_sources']*100 if su['total_sources']>0 else 0):.1f} %"],
        ["Senior Term Loan B", _frm(su["senior_debt"], cur),
         f"{(su['senior_debt']/su['total_sources']*100 if su['total_sources']>0 else 0):.1f} %"],
        ["Mezzanine", _frm(su["mezz_debt"], cur),
         f"{(su['mezz_debt']/su['total_sources']*100 if su['total_sources']>0 else 0):.1f} %"],
        ["TOTAL SOURCES", _frm(su["total_sources"], cur), "100,0 %"],
    ]
    add_table(slide, 1.02, sources_y + 0.62, 11.40, 3.0,
              4, 3,
              col_widths_pct=[0.50, 0.30, 0.20],
              header_data=None,
              rows_data=sources_rows,
              border_hex="DDDDDD")

    # Header USES (droite) — text_box centré verticalement dans le rect
    add_rect(slide, 12.94, sources_y, 11.46, 0.55, NAVY)
    add_text_box(slide, 12.94, sources_y + 0.08, 11.46, 0.42,
                 "USES (utilisation)", 10.5, WHITE, bold=True,
                 align=_PP_CENTER, wrap=False)

    uses_rows = [
        ["Purchase price equity", _frm(su["purchase_equity"], cur)],
        ["Refinancement dette", _frm(su["refi"], cur)],
        ["Transaction fees", _frm(su["fees"], cur)],
        ["Cash minimum BS", _frm(su["min_cash"], cur)],
        ["TOTAL USES", _frm(su["total_uses"], cur)],
    ]
    add_table(slide, 12.94, sources_y + 0.62, 11.46, 3.5,
              5, 2,
              col_widths_pct=[0.65, 0.35],
              header_data=None,
              rows_data=uses_rows,
              border_hex="DDDDDD")

    # ── Box LLM hypothèses (footer) — texte adapté selon éligibilité ──
    if not pack["eligible"]:
        # Texte explicatif pour profil NON éligible
        reasons = []
        from pptx.util import Cm as _Cm
        latest_yr_l = None
        if hasattr(snap, "ratios"):
            pass  # ratios passés via pack ailleurs
        # Reconstruire les raisons depuis les métriques connues du pack
        ebitda = pack.get("ebitda", 0)
        ev_market = pack.get("ev_market", 0)
        ebit = ebitda * 0.85 if ebitda else 0
        if ebitda <= 0:
            reasons.append("EBITDA négatif ou nul — la société ne génère pas de cash flow opérationnel "
                           "suffisant pour servir une dette LBO")
        # Inférer les autres raisons depuis les sources_uses
        sponsor_eq = su.get("sponsor_equity", 0)
        senior = su.get("senior_debt", 0)
        if ebitda > 0:
            margin_estimate = ebitda / max(ev_market / 14, 1)  # rough
            if (sponsor_eq / max(su.get("total_sources", 1), 1)) > 0.7:
                reasons.append("Sponsor equity représenterait > 70% du financement — leverage insuffisant "
                               "pour atteindre les rendements PE typiques (>20% IRR)")
        if not reasons:
            reasons.append("Combinaison des critères (marge EBITDA, cash conversion, levier actuel) ne permet pas "
                           "un montage LBO standard mid-market")

        hypotheses_text = (
            f"Profil non éligible au LBO mid-market standard. Raisons identifiées : "
            f"{' ; '.join(reasons)}. "
            f"L'analyse ci-dessus est fournie à titre théorique uniquement — un sponsor PE n'engagerait "
            f"pas de capital sur cette cible dans les conditions actuelles. Les chiffres affichés "
            f"(IRR, MOIC, sources/uses) traduisent ce que donnerait un montage forcé avec les hypothèses "
            f"standards, ils ne représentent pas une recommandation."
        )
    else:
        hypotheses_text = pack["llm_texts"].get("hypotheses_text") or (
            f"Multiple d'entrée {pack['entry_mult']:.1f}x EBITDA "
            f"(vs marché {pack['multiple_Marché']:.1f}x). Leverage 5x EBITDA "
            f"(Senior 3,5x à 8% + Mezz 1,5x à 10%). Sortie conservatrice à "
            f"{pack['exit_mult']:.1f}x. Hypothèses opérationnelles importées du DCF."
        )

    _lbo_title = _jpm_title("lbo", snap=snap, extra={"irr_base": pack.get("returns", {}).get("irr_base")})
    commentary_box(slide, 1.02, 11.50, 23.37, 1.80, hypotheses_text, title=_lbo_title)
    return slide


def _slide_lbo_returns(prs, snap, pack: dict):
    """Slide LBO 2/3 — Returns Sponsor (KPIs + heatmap sensibilité)."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    navy_bar(slide)
    footer_bar(slide)
    section_dots(slide, 3)

    ticker = pack["ticker"]
    company = pack["company"]
    cur = pack["cur_sym"]

    slide_title(slide, "Modèle LBO — Returns Sponsor",
                f"{company} ({ticker})  ·  IRR, MOIC et sensibilité multiples")

    # ── 4 KPIs en haut ──
    ret = pack["returns"]
    irr_base = ret["irr_base"]
    moic_base = ret["moic_base"]

    _irr_fill = GREEN_PALE if irr_base >= 0.20 else (AMBER_PALE if irr_base >= 0.15 else RED_PALE)
    _irr_accent = GREEN if irr_base >= 0.20 else (AMBER if irr_base >= 0.15 else RED)
    kpi_box(slide, 1.02, 2.45, 5.60, 2.30,
            f"{irr_base*100:.1f}%", "IRR Sponsor", "Base case 5 ans",
            fill=_irr_fill, accent=_irr_accent)
    kpi_box(slide, 6.92, 2.45, 5.60, 2.30,
            f"{moic_base:.2f}x", "MOIC", "Money-on-invested",
            fill=NAVY_PALE, accent=NAVY_MID)
    kpi_box(slide, 12.81, 2.45, 5.60, 2.30,
            _frm(ret["sponsor_equity"], cur), "Equity entry", "Sponsor au closing",
            fill=NAVY_PALE, accent=NAVY_MID)
    kpi_box(slide, 18.71, 2.45, 5.60, 2.30,
            _frm(ret["equity_exit"], cur), "Equity exit Y5", "Vente sponsor",
            fill=NAVY_PALE, accent=NAVY_MID)

    # ── Heatmap sensibilité IRR 5×5 ──
    add_text_box(slide, 1.02, 5.10, 23.37, 0.45,
                 "SENSIBILITÉ IRR — Multiple d'entrée × Multiple de sortie",
                 9, NAVY, bold=True)

    sens = pack["sensitivity"]
    grid = sens["grid"]
    ent_mults = sens["entry_mults"]
    ex_mults = sens["exit_mults"]

    # Header row : multiples sortie
    head_row = ["Entry ↓ / Exit →"] + [f"{m:.1f}x" for m in ex_mults]
    sens_rows = []
    for i, em in enumerate(ent_mults):
        row = [f"{em:.1f}x"]
        for j in range(5):
            v = grid[i][j] if i < len(grid) and j < len(grid[i]) else None
            row.append(f"{v*100:+.1f}%" if v is not None else "—")
        sens_rows.append(row)

    tbl_sens = add_table(
        slide, 1.02, 5.65, 16.50, 5.20,
        len(sens_rows), 6,
        col_widths_pct=[0.20, 0.16, 0.16, 0.16, 0.16, 0.16],
        header_data=head_row,
        rows_data=sens_rows,
        border_hex="DDDDDD"
    )

    # Heatmap colors
    try:
        for ri in range(len(grid)):
            for ci_idx in range(5):
                v = grid[ri][ci_idx] if ri < len(grid) and ci_idx < len(grid[ri]) else None
                if v is None:
                    continue
                cell = tbl_sens.cell(ri + 1, ci_idx + 1)
                cell.fill.solid()
                if v >= 0.25:
                    cell.fill.fore_color.rgb = rgb(GREEN_PALE)
                elif v >= 0.15:
                    cell.fill.fore_color.rgb = rgb(AMBER_PALE)
                elif v >= 0:
                    cell.fill.fore_color.rgb = rgb("FDF3E5")
                else:
                    cell.fill.fore_color.rgb = rgb(RED_PALE)
    except Exception:
        pass

    # ── Tableau hypothèses retenues (droite) ──
    add_text_box(slide, 17.85, 5.10, 6.55, 0.45,
                 "HYPOTHÈSES BASE", 9, NAVY, bold=True)
    hypo_rows = [
        ["Multiple entrée", f"{pack['entry_mult']:.1f}x"],
        ["Multiple sortie", f"{pack['exit_mult']:.1f}x"],
        ["Leverage total", f"{pack['leverage']:.1f}x"],
        ["Période", f"{pack['hold_years']} ans"],
        ["EBITDA Y5", _frm(pack["ebitda_y5"], cur)],
    ]
    add_table(slide, 17.85, 5.65, 6.55, 4.20,
              5, 2,
              col_widths_pct=[0.55, 0.45],
              header_data=None,
              rows_data=hypo_rows,
              border_hex="DDDDDD")

    # ── Box LLM lecture des returns (footer) ──
    returns_text = pack["llm_texts"].get("returns_text") or (
        f"IRR sponsor base {irr_base*100:.1f}% sur {pack['hold_years']} ans, "
        f"MOIC {moic_base:.2f}x. "
        f"{'Très attractif (>20%)' if irr_base >= 0.20 else ('Acceptable (15-20%)' if irr_base >= 0.15 else 'Sous le seuil PE typique')}. "
        f"La sensibilité aux multiples révèle les zones de robustesse."
    )
    _lbo_ret_title = _jpm_title("lbo", snap=snap, extra={"irr_base": irr_base})
    commentary_box(slide, 1.02, 11.20, 23.37, 2.10, returns_text, title=_lbo_ret_title)
    return slide


def _slide_lbo_stress(prs, snap, pack: dict):
    """Slide LBO 3/3 — Debt Schedule + Scénarios + Risques."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    navy_bar(slide)
    footer_bar(slide)
    section_dots(slide, 3)

    ticker = pack["ticker"]
    company = pack["company"]
    cur = pack["cur_sym"]

    slide_title(slide, "Modèle LBO — Debt Schedule & Stress Test",
                f"{company} ({ticker})  ·  Évolution leverage et scénarios bull/bear")

    # ── Debt Schedule (Y0-Y5) ──
    add_text_box(slide, 1.02, 2.45, 23.37, 0.45,
                 "DEBT SCHEDULE — Évolution du leverage sur 5 ans", 9, NAVY, bold=True)

    sched = pack["debt_schedule"]
    # Unité dans le header, pas dans chaque cellule
    _unit = f"Mds{cur}" if cur != "EUR" else "Md€"
    head = ["Année", f"EBITDA ({_unit})", f"Senior ({_unit})", f"Mezz ({_unit})",
            f"Cash ({_unit})", f"Net Debt ({_unit})", "Lvg (×)"]

    def _val_only(v):
        """Retourne juste le nombre en Mds (float) sans suffixe."""
        if v is None: return "—"
        try:
            fv = float(v)
            if abs(fv) > 1_000_000_000:
                fv = fv / 1_000_000  # normalise valeurs absolues → millions
            # fv est en millions, on convertit en Mds si ≥ 1000 M
            if abs(fv) >= 1000:
                return _fr(fv / 1000, 1)
            return _fr(fv / 1000, 2)  # < 1 Md → 2 décimales
        except:
            return "—"

    sched_rows = []
    for d in sched:
        sched_rows.append([
            f"Y{d['year']}",
            _val_only(d["ebitda"]),
            _val_only(d["senior"]),
            _val_only(d["mezz"]),
            _val_only(d["cash"]),
            _val_only(d["net_debt"]),
            f"{d['leverage']:.2f}x",
        ])

    add_table(slide, 1.02, 3.00, 23.37, 4.30,
              len(sched_rows), 7,
              col_widths_pct=[0.10, 0.15, 0.15, 0.15, 0.15, 0.15, 0.15],
              header_data=head,
              rows_data=sched_rows,
              border_hex="DDDDDD")

    # ── Scénarios Bull / Base / Bear (3 colonnes) ──
    add_text_box(slide, 1.02, 7.55, 23.37, 0.45,
                 "SCÉNARIOS DE STRESS — Bull / Base / Bear", 9, NAVY, bold=True)

    sc = pack["scénarios"]
    col_w = 7.62
    col_xs = [1.02, 8.97, 16.92]
    col_labels = ["BULL", "BASE", "BEAR"]
    col_keys = ["bull", "base", "bear"]
    col_colors = [GREEN, NAVY_MID, RED]
    col_pales = [GREEN_PALE, NAVY_PALE, RED_PALE]

    for i, (lbl, key, col_c, col_p) in enumerate(zip(col_labels, col_keys, col_colors, col_pales)):
        cx = col_xs[i]
        # Header de colonne — text_box aligné EXACTEMENT sur le rect
        add_rect(slide, cx, 8.10, col_w, 0.55, col_c)
        add_text_box(slide, cx, 8.20, col_w, 0.36,
                     lbl, 11, WHITE, bold=True,
                     align=__import__("pptx").enum.text.PP_ALIGN.CENTER, wrap=False)
        # Body
        add_rect(slide, cx, 8.65, col_w, 2.80, col_p)
        s = sc[key]
        irr_v = s["irr"] * 100
        rows_text = [
            ("IRR Sponsor", f"{irr_v:+.1f}%"),
            ("MOIC", f"{s['moic']:.2f}x"),
            ("Leverage exit", f"{s['leverage']:.2f}x"),
            ("ICR exit", f"{s['icr']:.1f}x"),
        ]
        for ri, (lab, val) in enumerate(rows_text):
            ry = 8.75 + ri * 0.65
            add_text_box(slide, cx + 0.30, ry, col_w * 0.55, 0.45,
                         lab, 9, GREY_TXT)
            add_text_box(slide, cx + col_w * 0.55, ry, col_w * 0.40, 0.45,
                         val, 11, col_c, bold=True, align=__import__("pptx").enum.text.PP_ALIGN.RIGHT)

    # ── Box LLM risques (footer) — texte tronqué 380 chars + box élargie ──
    risks_text = pack["llm_texts"].get("risks_text") or (
        f"Le scénario bear traduit la sensibilité du deal à la compression des marges, "
        f"à la hausse des taux et au multiple compression. "
        f"La thèse s'invalide si EBITDA -200 bps ou multiple sortie -2×."
    )
    risks_text = _truncate(risks_text, 380)
    _stress_irr = None
    try:
        _stress_irr = pack["scénarios"]["base"]["irr"]
    except Exception:
        pass
    _stress_title = _jpm_title("lbo", snap=snap, extra={"irr_base": _stress_irr})
    commentary_box(slide, 1.02, 11.45, 23.37, 1.85, risks_text, title=_stress_title)
    return slide


# ---------------------------------------------------------------------------
# Slide 16 — Risques & Conditions d'Invalidation
# ---------------------------------------------------------------------------

def _slide_risques(prs, snap, synthesis, devil, extra_scores: dict = None):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    navy_bar(slide)
    footer_bar(slide)
    section_dots(slide, 4)
    slide_title(slide, "Risques & Conditions d'Invalidation",
                "Analyse des risques structurels  \u00b7  Th\u00e8se contraire")

    counter_thesis = _g(devil, "counter_thesis", "") or ""
    counter_risks  = _g(devil, "counter_risks", []) or []
    risks_s        = _g(synthesis, "risks", []) or []
    invalidation   = _g(synthesis, "invalidation_list", []) or []

    risk_sources  = counter_risks if counter_risks else risks_s
    ct_parts      = _split_text(counter_thesis, 3)
    neg_themes_s  = _g(synthesis, "negative_themes", []) or []

    card_configs = [
        (RED_PALE,   RED,        "Risque 1"),
        (NAVY_PALE,  NAVY_MID,   "Risque 2"),
        (AMBER_PALE, AMBER,      "Risque 3"),
    ]
    card_w = 7.62
    card_h = 5.33
    card_y = 2.67
    gaps   = [1.02, 8.97, 16.92]

    for i, (fill, accent, label) in enumerate(card_configs):
        cx    = gaps[i]
        risk  = str(risk_sources[i]) if i < len(risk_sources) else label
        body  = ct_parts[i] if i < len(ct_parts) else ""
        # Fallback 1 : negative_themes du synthesis
        if not body.strip():
            body = str(neg_themes_s[i]) if i < len(neg_themes_s) else ""
        # Fallback 2 : utiliser le texte du risque lui-meme comme corps de carte
        # (le titre affiche le debut, le corps affiche l'intégralité — plus de carte vide)
        if not body.strip() and risk and risk != label:
            body = risk
        # Fallback 3 : synthesis.risks complet si counter_risks ne fournit que des titres courts
        if not body.strip() and i < len(risks_s):
            body = str(risks_s[i])
        # Fallback ultime : message generique localise
        if not body.strip():
            _generic = [
                "Risque macro-économique : remontée des taux ou ralentissement de la demande mondiale.",
                "Risque sectoriel : intensification de la concurrence ou disruption technologique.",
                "Risque spécifique : détérioration des marges ou exécution stratégique.",
            ]
            body = _generic[i] if i < len(_generic) else "Analyse en cours."
        add_rect(slide, cx, card_y, card_w, card_h, fill)
        add_rect(slide, cx, card_y, card_w, 0.15, accent)
        add_text_box(slide, cx + 0.30, card_y + 0.30, card_w - 0.60, 1.10,
                     _truncate(risk, 220), 9, accent, bold=True, wrap=True)
        add_text_box(slide, cx + 0.30, card_y + 1.60, card_w - 0.60, 3.50,
                     _fit(body, 400), 8, GREY_TXT, wrap=True)

    # Invalidation table
    add_rect(slide, 1.02, 8.33, 23.37, 0.03, "AAAAAA")
    add_text_box(slide, 1.02, 8.53, 23.37, 0.46,
                 "Conditions d'invalidation de la thèse", 9, NAVY, bold=True)

    if not invalidation:
        inv_rows = [
            ["Macro",     "Remontée des taux directeurs > 200bps",              "6-12 mois"],
            ["Sectoriel", "Rupture technologique remettant en cause le modèle", "12-18 mois"],
            ["Société",   "Détérioration matérielle des marges opérationnelles","2-3 trim."],
        ]
    else:
        _AXIS_NORM = {"Société": "Société", "société": "Société", "société": "Société",
                      "Sectoriel": "Sectoriel", "Macro": "Macro"}
        inv_rows = []
        for item in invalidation[:3]:
            axis  = str(_g(item, "axis",      "—"))
            axis  = _AXIS_NORM.get(axis, axis)
            cond  = str(_g(item, "condition", "—"))
            horiz = str(_g(item, "horizon",   "—"))
            inv_rows.append([axis, cond, horiz])

    inv_tbl = add_table(slide, 1.02, 9.19, 23.37, 3.05,
              len(inv_rows), 3,
              col_widths_pct=[0.15, 0.65, 0.20],
              header_data=["Axe", "Condition d'invalidation", "Horizon"],
              rows_data=inv_rows)

    # Bold Axe column (col 0)
    for ri in range(1, len(inv_rows) + 1):
        cell = inv_tbl.cell(ri, 0)
        try:
            p = cell.text_frame.paragraphs[0]
            for run in p.runs:
                run.font.bold = True
                run.font.color.rgb = rgb(NAVY)
        except Exception:
            pass

    # --- Bande scores compact : 2 lignes (12.30 → 13.33) ----------------------
    # Ligne 1 : Detresse · M&A · Régime macroeconomique
    # Ligne 2 : Qualité Earnings · Structure Capital · Soutenabilite Dividende
    # Zone disponible : 12.24 (fin table) → 13.39 (footer) = 1.15cm
    if extra_scores:
        _dist  = extra_scores.get('composite_distress') or {}
        _ma    = extra_scores.get('ma_score')           or {}
        _macro = extra_scores.get('macro')              or {}
        _eq    = extra_scores.get('earnings_quality')   or {}
        _cs    = extra_scores.get('capital_structure')  or {}
        _ds    = extra_scores.get('dividend_sustainability') or {}

        _CMAP = {
            'Sain': "1A7A4A", 'Saine': "1A7A4A",
            'Modéré': AMBER, 'Moderé': AMBER, 'Modérée': AMBER,
            'Vigilance': AMBER, 'Critique': RED,
            'Très attractive': "1A7A4A", 'Attractive': "1A7A4A",
            'Moderate': AMBER, 'Peu attractive': RED,
            'Bull': "1A7A4A", 'Bear': RED,
            'Volatile': AMBER, 'Transition': AMBER,
            'Excellente': "1A7A4A", 'Correcte': "1B3A6B",
            'Faible': RED, 'N/D': "B0B0B0",
            'Très soutenable': "1A7A4A", 'Soutenable': "1A7A4A",
            'Tendu': AMBER, 'Insoutenable': RED,
        }

        strip_h = 0.50
        gap     = 0.03
        bx_w    = 7.79
        bx_x    = [1.02, 9.01, 17.01]

        # --- Ligne 1 : Distress / M&A / Régime ---
        row1_y  = 12.05
        row1_items = []
        if _dist.get('score') is not None:
            d_lbl = _dist.get('label', '—')
            row1_items.append((f"Détresse : {_dist['score']}/100  {d_lbl}",
                               _CMAP.get(d_lbl, NAVY)))
        if _ma.get('score') is not None:
            m_lbl = _ma.get('label', '—')
            row1_items.append((f"M&A : {_ma['score']}/100  {m_lbl}",
                               _CMAP.get(m_lbl, NAVY)))
        regime = _macro.get('regime_v', '')
        rec_6m = _macro.get('récession_prob_6m')
        if regime and regime != 'Inconnu':
            rec_part = f"  Rec.6M:{rec_6m}%" if rec_6m is not None else ''
            row1_items.append((f"Régime : {regime_v}{rec_part}", _CMAP.get(regime, NAVY)))

        for i, (txt, col) in enumerate(row1_items[:3]):
            x = bx_x[i]
            add_rect(slide, x, row1_y, bx_w, strip_h, "F0F4F8")
            add_rect(slide, x, row1_y, 0.12, strip_h, col)
            add_text_box(slide, x + 0.22, row1_y + 0.03, bx_w - 0.30, strip_h - 0.06,
                         txt, 7.0, NAVY, bold=False)

        # --- Ligne 2 : Earnings Quality / Structure Capital / Dividende ---
        row2_y = row1_y + strip_h + gap
        row2_items = []

        eq_lbl = _eq.get('label')
        if eq_lbl and eq_lbl != 'N/D':
            cc = _eq.get('cash_conversion')
            cc_str = f" ({cc:.2f}x)" if cc is not None else ''
            row2_items.append((f"Earnings Qualité : {eq_lbl}{cc_str}",
                               _CMAP.get(eq_lbl, NAVY)))

        cs_lbl = _cs.get('label')
        if cs_lbl and cs_lbl != 'N/D':
            ratio = _cs.get('short_term_ratio')
            r_str = f" ({ratio*100:.0f}% CT)" if ratio is not None else ''
            row2_items.append((f"Struct. Capital : {cs_lbl}{r_str}",
                               _CMAP.get(cs_lbl, NAVY)))

        ds_lbl = _ds.get('label')
        if ds_lbl and ds_lbl not in ('N/D', 'Sans dividende'):
            cov = _ds.get('fcf_coverage')
            c_str = f" ({cov:.1f}x)" if cov is not None else ''
            row2_items.append((f"Dividende : {ds_lbl}{c_str}",
                               _CMAP.get(ds_lbl, NAVY)))
        elif not row2_items:
            row2_items = []  # no items → skip row 2

        for i, (txt, col) in enumerate(row2_items[:3]):
            x = bx_x[i]
            add_rect(slide, x, row2_y, bx_w, strip_h, "EBF0F7")
            add_rect(slide, x, row2_y, 0.12, strip_h, col)
            add_text_box(slide, x + 0.22, row2_y + 0.03, bx_w - 0.30, strip_h - 0.06,
                         txt, 7.0, NAVY, bold=False)

    return slide


# ---------------------------------------------------------------------------
# Slide 18 — Sentiment FinBERT
# ---------------------------------------------------------------------------

def _slide_sentiment(prs, snap, synthesis, sentiment):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    navy_bar(slide)
    footer_bar(slide)
    section_dots(slide, 5)

    sent_score    = _g(sentiment, "score", 0.0) or 0.0
    sent_label    = _g(sentiment, "label", "NEUTRAL") or "NEUTRAL"
    sent_articles = _g(sentiment, "articles_analyzed", 0) or 0
    sent_conf     = _g(sentiment, "confidence", 0.0) or 0.0
    sent_breakdown= _g(sentiment, "breakdown", {}) or {}

    sent_label_display = _sent_label_fr(sent_label, sent_score)
    rec = _g(synthesis, "recommendation", "HOLD") or "HOLD"
    sent_meta   = _g(sentiment, "meta", {}) or {}
    engine_used = sent_meta.get("engine", "finbert")
    _is_llm     = engine_used == "llm_groq"
    engine_lbl  = "Modèle LLM Groq" if _is_llm else "FinBERT (fallback)"

    _all_neutral = (sent_articles > 0 and
                    (_g(sent_breakdown, "neutral_count", 0) or 0) >= sent_articles * 0.9)
    _warn_finbert = _all_neutral and not _is_llm

    slide_title(
        slide,
        "Sentiment de March\u00e9 \u2014 " + ("LLM Groq" if _is_llm else "FinBERT"),
        f"Analyse s\u00e9mantique  \u00b7  {sent_articles} articles  \u00b7  7 derniers jours"
        + ("  \u00b7  Classification LLM" if _is_llm else "  \u00b7  English-only fallback"),
    )

    # 4 KPI cards
    kpi_box(slide, 1.02,  2.67, 7.11, 2.79,
            sent_label_display, "Orientation globale",
            f"Score agr\u00e9g\u00e9 : {sent_score:+.3f}".replace(".", ","))
    kpi_box(slide, 8.64,  2.67, 5.08, 2.79,
            str(sent_articles), "Articles analys\u00e9s", "Fen\u00eatre : 7 jours")
    kpi_box(slide, 14.22, 2.67, 4.57, 2.79,
            _frpct(sent_conf), "Confiance IA", engine_lbl)
    _, rec_accent = _rec_colors(rec)
    kpi_box(slide, 19.30, 2.67, 4.83, 2.79,
            rec.upper(), "Coh\u00e9rence", "Align\u00e9 avec la th\u00e8se")

    # Breakdown table
    add_text_box(slide, 1.02, 5.79, 23.37, 0.46,
                 "D\u00e9tail par orientation", 9, NAVY_MID, bold=True)

    pos_val  = _g(sent_breakdown, "avg_positive",  None)
    neg_val  = _g(sent_breakdown, "avg_negative",  None)
    neut_val = _g(sent_breakdown, "avg_neutral",   None)
    pos_cnt  = _g(sent_breakdown, "positive_count", None)
    neg_cnt  = _g(sent_breakdown, "negative_count", None)
    neu_cnt  = _g(sent_breakdown, "neutral_count",  None)

    def _fmt_score(v):
        if v is None: return "\u2014"
        try:    return f"{float(v):+.3f}".replace(".", ",")
        except: return str(v)

    def _cnt(v):
        if v is None: return "\u2014"
        return str(int(v))

    pos_themes  = _g(synthesis, "positive_themes", []) or []
    neg_themes  = _g(synthesis, "negative_themes", []) or []
    # Limiter à 1 thème par ligne pour éviter l'explosion de hauteur du tableau
    pos_theme_str  = _truncate(str(pos_themes[0]), 240) if pos_themes else "Catalyseurs, croissance, r\u00e9sultats"
    neg_theme_str  = _truncate(str(neg_themes[0]), 240) if neg_themes else "Risques macro, concurrence, dette"

    neut_theme_str = ("Actualit\u00e9 sectorielle g\u00e9n\u00e9rale"
                      if _is_llm else "Articles non-anglais \u2014 LLM indisponible ce run"
                      if _all_neutral else "Actualit\u00e9 sectorielle g\u00e9n\u00e9rale")

    # Si count = 0 → score et thème non significatifs, on les masque
    break_rows = [
        ["Positif",
         _cnt(pos_cnt),
         _fmt_score(pos_val) if (pos_cnt or 0) > 0 else "\u2014",
         pos_theme_str       if (pos_cnt or 0) > 0 else "Aucun article class\u00e9 positif"],
        ["Neutre",
         _cnt(neu_cnt),
         _fmt_score(neut_val) if (neu_cnt or 0) > 0 else "\u2014",
         neut_theme_str],
        ["N\u00e9gatif",
         _cnt(neg_cnt),
         _fmt_score(neg_val) if (neg_cnt or 0) > 0 else "\u2014",
         neg_theme_str       if (neg_cnt or 0) > 0 else "Aucun article class\u00e9 n\u00e9gatif"],
    ]
    sent_row_fills = [GREEN_PALE, WHITE, RED_PALE]
    tbl_y    = 6.48
    tbl_h_s  = 3.20
    sent_tbl = add_table(slide, 1.02, tbl_y, 23.37, tbl_h_s,
              3, 4,
              col_widths_pct=[0.15, 0.15, 0.20, 0.50],
              header_data=["Orientation", "Articles", "Score moyen", "Th\u00e8mes principaux"],
              rows_data=break_rows,
              row_fills=sent_row_fills)

    # Bold + colored text per sentiment row
    _sent_colors = [GREEN, GREY_TXT, RED]
    from pptx.util import Pt as _Pt3
    for ri in range(1, 4):
        tc_color = _sent_colors[ri - 1]
        for ci in range(2):  # bold color only on first 2 cols
            cell = sent_tbl.cell(ri, ci)
            try:
                p = cell.text_frame.paragraphs[0]
                for run in p.runs:
                    run.font.bold = True
                    run.font.color.rgb = rgb(tc_color)
            except Exception:
                pass

    # Bannière avertissement si FinBERT fallback + tous neutres
    if _warn_finbert:
        add_rect(slide, 1.02, tbl_y + tbl_h_s + 0.25, 23.37, 0.72, "FFF3CD")
        add_rect(slide, 1.02, tbl_y + tbl_h_s + 0.25, 0.13,  0.72, "E6A817")
        add_text_box(
            slide, 1.35, tbl_y + tbl_h_s + 0.30, 23.0, 0.62,
            "\u26a0  Classification LLM indisponible ce run \u2014 FinBERT English-only class\u00e9 tous les articles FR en neutre. "
            "Les scores ci-dessus ne sont pas significatifs. "
            "Relancer l\u2019analyse lorsque Groq est disponible.",
            7.5, "7D5A00", wrap=True
        )
        # Titres des articles bruts (jusqu'à 8) pour context
        samples = _g(sentiment, "samples", []) or []
        all_news_titles = [s.get("headline", "") for s in samples if s.get("headline")]
        if not all_news_titles:
            # Fallback: reconstruire depuis breakdown si disponible
            all_news_titles = []
        if all_news_titles:
            add_text_box(slide, 1.02, tbl_y + tbl_h_s + 1.15, 23.37, 0.40,
                         "Titrès des articles analys\u00e9s (bruts, non class\u00e9s) :",
                         8, NAVY_MID, bold=True)
            art_txt = "\n".join(f"\u2022 {h[:110]}" for h in all_news_titles[:6])
            add_text_box(slide, 1.02, tbl_y + tbl_h_s + 1.60, 23.37, 1.80,
                         art_txt, 7.5, "333333", wrap=True)
        comment_y = tbl_y + tbl_h_s + (3.65 if all_news_titles else 1.80)
    else:
        comment_y = max(tbl_y + tbl_h_s + 1.20, 11.5)

    # Commentaire sentiment — LLM Groq en priorité, fallback agrégé
    val_comment = _g(sentiment, "llm_commentary", "") or ""
    if not val_comment.strip():
        lbl_fr = _sent_label_fr(sent_label, sent_score)
        if _warn_finbert:
            val_comment = (
                f"Analyse de sentiment indisponible ce run : le mod\u00e8le LLM Groq n\u2019a pas pu classifier "
                f"les {sent_articles} articles. FinBERT \u00e9tant English-only, tous les titrès "
                f"fran\u00e7ais ont \u00e9t\u00e9 class\u00e9s neutrès par d\u00e9faut. "
                f"Relancer l\u2019analyse avec Groq actif pour obtenir un score exploitable."
            )
        else:
            val_comment = (
                f"Le sentiment {lbl_fr.lower()} (score {sent_score:+.3f}) est bas\u00e9 sur "
                f"{sent_articles} articles analys\u00e9s sur 7 jours. "
                f"Confiance IA\u00a0: {sent_conf:.0%}. "
                f"Coh\u00e9rence avec la recommandation {rec}\u00a0: "
                + ("valid\u00e9e." if (sent_score > 0.05 and rec == "BUY") or
                                      (sent_score < -0.05 and rec == "SELL") or
                                      (abs(sent_score) <= 0.05 and rec == "HOLD")
                   else "surveiller.")
            )
    comment_y = min(comment_y, 11.8)
    _comment_h = min(2.0, 13.25 - comment_y)  # eviter chevauchement footer (y=13.39)
    add_text_box(slide, 1.02, comment_y, 23.37, _comment_h,
                 _truncate(val_comment, 420), 8.5, "333333", wrap=True)

    return slide


# ---------------------------------------------------------------------------
# Slide 19 — Actionnariat
# ---------------------------------------------------------------------------

def _slide_actionnariat(prs, snap, synthesis):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    navy_bar(slide)
    footer_bar(slide)
    section_dots(slide, 5)

    ci       = snap.company_info if snap else None
    gen_date = _fr_date_long(_g(ci, "analysis_date", None) or date.today())

    slide_title(slide, "Actionnariat & Structure du Capital",
                f"R\u00e9partition de l'actionnariat  \u00b7  Au {gen_date}")

    # ----------------------------------------------------------------
    # Données actionnaires depuis snap.institutional_holders (yfinance)
    # ----------------------------------------------------------------
    holders = (getattr(snap, "institutional_holders", None) or []) if snap else []

    # Fallback si aucune donnée réelle
    if not holders:
        holders = [
            {"name": "Institutionnels passifs",  "type": "Institutionnel", "pct": 24.0, "style": "Passif"},
            {"name": "Institutionnels actifs",   "type": "Institutionnel", "pct": 12.0, "style": "Actif"},
            {"name": "Insiders & dirigeants",    "type": "Insiders",       "pct": 8.0,  "style": "Insider"},
            {"name": "Autrès (Retail)",          "type": "Retail & autres","pct": 56.0, "style": "\u2014"},
        ]

    # ----------------------------------------------------------------
    # Tableau gauche : Top actionnaires (Actionnaire | Type | % | Style)
    # ----------------------------------------------------------------
    holder_rows = []
    for h in holders:
        pct_val = h.get("pct")
        pct_str = f"{float(pct_val):,.1f}\u00a0%".replace(",", "\u00a0") if pct_val is not None else "\u2014"
        holder_rows.append([
            str(h.get("name", "\u2014"))[:35],
            str(h.get("type", "\u2014")),
            pct_str,
            str(h.get("style", "\u2014")),
        ])

    n_rows      = len(holder_rows)
    tbl_h       = min(0.71 + n_rows * 0.61, 7.0)

    add_text_box(slide, 1.02, 2.08, 13.97, 0.56,
                 "Top actionnaires", 9, NAVY_MID, bold=True)
    holders_tbl = add_table(slide, 1.02, 2.69, 14.73, tbl_h,
              n_rows, 4,
              col_widths_pct=[0.38, 0.26, 0.20, 0.16],
              header_data=["Actionnaire", "Type", "D\u00e9tention\u00a0%", "Style"],
              rows_data=holder_rows,
              border_hex="DDDDDD")

    # Bold + couleur sur la colonne Style
    from pptx.util import Pt as _Pt_h
    _style_colors = {"Passif": "2E5FA3", "Actif": "1A7A4A", "Insider": "B06000"}
    for ri in range(1, n_rows + 1):
        try:
            cell = holders_tbl.cell(ri, 3)
            for run in cell.text_frame.paragraphs[0].runs:
                style_val = holder_rows[ri - 1][3]
                col = _style_colors.get(style_val, GREY_TXT)
                run.font.color.rgb = rgb(col)
                run.font.bold = True
                run.font.size = _Pt_h(7.5)
        except Exception:
            pass

    # ----------------------------------------------------------------
    # Tableau droit : Répartition par type (dynamique depuis les données)
    # ----------------------------------------------------------------
    pct_passif  = sum(h.get("pct") or 0 for h in holders if h.get("style") == "Passif")
    pct_actif   = sum(h.get("pct") or 0 for h in holders if h.get("style") == "Actif")
    pct_insider = sum(h.get("pct") or 0 for h in holders if h.get("style") == "Insider")
    pct_retail  = sum(h.get("pct") or 0 for h in holders if h.get("style") == "\u2014")

    def _pct_str(v):
        return f"{v:.1f}\u00a0%" if v > 0 else "\u2014"

    type_rows = []
    if pct_passif  > 0: type_rows.append(["Institutionnel passif",  _pct_str(pct_passif)])
    if pct_actif   > 0: type_rows.append(["Institutionnel actif",   _pct_str(pct_actif)])
    if pct_insider > 0: type_rows.append(["Insider",                _pct_str(pct_insider)])
    if pct_retail  > 0: type_rows.append(["Retail & autres",        _pct_str(pct_retail)])
    if not type_rows:
        type_rows = [["Donn\u00e9es indisponibles", "\u2014"]]

    type_tbl_h = min(0.71 + len(type_rows) * 0.61, tbl_h)

    add_text_box(slide, 16.00, 2.08, 8.38, 0.56,
                 "R\u00e9partition par type", 9, NAVY_MID, bold=True)
    add_table(slide, 16.00, 2.69, 8.38, type_tbl_h,
              len(type_rows), 2,
              col_widths_pct=[0.60, 0.40],
              header_data=["Type d'actionnaire", "Part estim\u00e9e"],
              rows_data=type_rows,
              border_hex="DDDDDD")

    # ----------------------------------------------------------------
    # Commentaire
    # ----------------------------------------------------------------
    commentary_y = 2.69 + tbl_h + 0.35
    _act_c1 = _g(synthesis, "valuation_comment", "") or ""
    _act_c2 = _g(synthesis, "peers_commentary", "") or ""
    _act_c3 = (_g(synthesis, "thesis", "") or "").replace(" | ", " ")
    _act_c4 = _g(synthesis, "summary", "") or ""
    comment_txt = " ".join(c for c in [_act_c1, _act_c2, _act_c3, _act_c4] if c.strip())
    if comment_txt.strip():
        _act_h = min(13.25 - commentary_y, 5.50)
        if _act_h > 0.5:
            _act_title = _jpm_title("verdict", snap=snap, synthesis=synthesis)
            commentary_box(slide, 1.02, commentary_y, 23.37, _act_h, comment_txt, title=_act_title)

    return slide


# ---------------------------------------------------------------------------
# Slide 20 — Historique de Cours
# ---------------------------------------------------------------------------

def _slide_historique(prs, snap, synthesis):
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    navy_bar(slide)
    footer_bar(slide)
    section_dots(slide, 5)

    ci       = snap.company_info if snap else None
    mkt      = snap.market if snap else None
    ticker   = _g(ci, "ticker", "—")
    currency = _g(ci, "currency", "USD") or "USD"
    cur_sym  = "EUR" if currency == "EUR" else "$"
    exchange = getattr(ci, "exchange", "") or "" if ci else ""
    gen_date = _fr_date_long(_g(ci, "analysis_date", None) or date.today())
    price    = _g(mkt, "share_price")

    _exch_part = f"  \u00b7  {exchange}" if exchange else ""
    slide_title(slide, f"Historique de Cours \u2014 52 Semaines",
                f"{ticker}{_exch_part}  \u00b7  {currency}  \u00b7  Au {gen_date}")

    history = snap.stock_history if (snap and snap.stock_history) else []
    prices  = []
    for pt in history:
        p = _g(pt, "price")
        if p is not None:
            try:
                prices.append(float(p))
            except Exception:
                pass

    p_high = max(prices) if prices else price
    p_low  = min(prices) if prices else price
    p_first = prices[0] if prices else price
    p_last  = prices[-1] if prices else price
    perf_52w = ((p_last / p_first) - 1) if (p_first and p_last and p_first != 0) else None

    # 4 KPI boxes
    kpi_box(slide, 1.02,  2.67, 5.33, 2.08,
            f"{_fr(price, 2)} {cur_sym}", "Cours actuel")
    kpi_box(slide, 6.73,  2.67, 5.33, 2.08,
            f"{_fr(p_high, 2)} {cur_sym}", "Plus haut 52 sem.")
    kpi_box(slide, 12.45, 2.67, 5.33, 2.08,
            f"{_fr(p_low, 2)} {cur_sym}", "Plus bas 52 sem.")
    kpi_box(slide, 18.16, 2.67, 6.10, 2.08,
            _frpct(perf_52w, signed=True),
            "Performance 52 sem.")

    # Titre analytique basé sur la performance
    _perf_label = _frpct(perf_52w, signed=True) if perf_52w is not None else ""
    _trend_txt = ("en hausse" if perf_52w and perf_52w > 0.05
                  else ("en repli" if perf_52w and perf_52w < -0.05 else "stable"))
    _hist_subtitle = (f"Cours {_trend_txt} de {_perf_label} sur 12 mois"
                      if _perf_label else "Évolution du cours — 12 derniers mois")
    add_text_box(slide, 1.02, 5.13, 23.37, 0.46,
                 _hist_subtitle, 9, NAVY, bold=True)

    # Graphique line chart matplotlib — à gauche, box LLM macro à droite
    chart_x = 1.02
    chart_y = 5.84
    chart_w = 15.50
    chart_h = 5.08

    _chart_inserted = False
    if prices and len(prices) >= 2:
        try:
            import matplotlib
            matplotlib.use('Agg')
            import matplotlib.pyplot as plt
            import io as _io

            months_lbl = [_g(pt, "month", "") or "" for pt in history[-len(prices):]]

            fig_w_in = 12.0
            fig_h_in = chart_h * (fig_w_in / chart_w)
            fig, ax = plt.subplots(figsize=(fig_w_in, fig_h_in))
            fig.patch.set_facecolor('#F5F7FA')
            ax.set_facecolor('#F5F7FA')

            x_idx = list(range(len(prices)))
            ax.plot(x_idx, prices, color='#2E5FA3', linewidth=2.2, zorder=3)
            ax.fill_between(x_idx, prices, min(prices) * 0.98,
                            color='#2E5FA3', alpha=0.08)

            # Marqueur dernier cours
            ax.scatter([x_idx[-1]], [prices[-1]], color='#2E5FA3', s=40, zorder=4)
            ax.annotate(f"{prices[-1]:.0f} {currency}",
                        xy=(x_idx[-1], prices[-1]),
                        xytext=(-10, 6), textcoords='offset points',
                        fontsize=8, color='#1B3A6B', fontweight='bold')

            # Plus haut / plus bas
            idx_hi = prices.index(p_high)
            idx_lo = prices.index(p_low)
            ax.annotate(f"H: {p_high:.0f}", xy=(idx_hi, p_high),
                        xytext=(0, 8), textcoords='offset points',
                        fontsize=7, color='#1B5E20',
                        arrowprops=dict(arrowstyle='-', color='#1B5E20', lw=0.7))
            ax.annotate(f"B: {p_low:.0f}", xy=(idx_lo, p_low),
                        xytext=(0, -14), textcoords='offset points',
                        fontsize=7, color='#B71C1C',
                        arrowprops=dict(arrowstyle='-', color='#B71C1C', lw=0.7))

            ax.set_xticks(x_idx)
            ax.set_xticklabels(months_lbl, fontsize=7.5, color='#555', rotation=20, ha='right')
            ax.tick_params(axis='y', labelsize=8, labelcolor='#555')
            ax.set_ylabel(currency, fontsize=8, color='#555')
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['left'].set_color('#D0D5DD')
            ax.spines['bottom'].set_color('#D0D5DD')
            ax.grid(axis='y', alpha=0.3, color='#D0D5DD', linewidth=0.5)

            plt.tight_layout(pad=0.5)
            _buf = _io.BytesIO()
            fig.savefig(_buf, format='png', dpi=150, bbox_inches='tight')
            plt.close(fig)
            _buf.seek(0)

            from pptx.util import Cm as _Cm
            slide.shapes.add_picture(_buf, _Cm(chart_x), _Cm(chart_y), _Cm(chart_w), _Cm(chart_h))
            _chart_inserted = True
        except Exception as _e:
            log.warning(f"[pptx] historique line chart error: {_e}")

    if not _chart_inserted:
        add_rect(slide, chart_x, chart_y, chart_w, chart_h, GREY_BG,
                 line_hex="D0D0D0", line_width_pt=0.5)
        add_text_box(slide, chart_x + 8.0, chart_y + 2.0, 7.37, 1.0,
                     "Historique de cours non disponible", 10, GREY_TXT)

    # Box LLM macro à droite du chart — contexte expliquant les fluctuations
    _macro_ctx = _g(synthesis, "macro_context", "") or _g(synthesis, "market_context", "") or ""
    if not _macro_ctx.strip():
        _macro_ctx = (
            f"L'évolution du cours sur 52 semaines reflète les dynamiques macroéconomiques "
            f"(politique monétaire, cycle économique) et les catalyseurs spécifiques au titre. "
            f"Les phases de surperformance et de correction identifiées sur le graphique "
            f"doivent être croisées avec le calendrier des publications de résultats "
            f"et les révisions de consensus pour distinguer le signal du bruit."
        )
    _macro_title = f"Contexte macro et catalyseurs — {ticker}"
    commentary_box(slide, 17.00, 5.84, 7.39, 5.08,
                   _truncate(_macro_ctx, 500), title=_macro_title)

    # Commentary en bas — verdict
    thesis_s = _g(synthesis, "summary", "") or _g(synthesis, "thesis", "") or ""
    if thesis_s.strip():
        _hist_title = _jpm_title("verdict", snap=snap, synthesis=synthesis)
        commentary_box(slide, 1.02, 11.20, 23.37, 2.15, _fit(thesis_s, 450), title=_hist_title)

    return slide


# ---------------------------------------------------------------------------
# Slide 21 — Conviction Tracker & Investment Thesis Summary
# ---------------------------------------------------------------------------

def _slide_conviction_tracker(prs, snap, synthesis, ratios, devil, sentiment):
    """Slide 21 : Conviction Tracker JPM-style — visualisation de la Thèse d investissement."""
    from pptx.enum.text import PP_ALIGN
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    navy_bar(slide)
    footer_bar(slide)
    section_dots(slide, 5)

    ci       = snap.company_info if snap else None
    ticker   = _g(ci, "ticker", "—")
    rec      = (_g(synthesis, "recommendation") or "HOLD").upper()
    conv_raw = _g(synthesis, "conviction") or 0.5
    try: conv = float(conv_raw)
    except: conv = 0.5
    conv_pct = int(conv * 100)
    gen_date = _fr_date_long(_g(ci, "analysis_date", None) or date.today())

    slide_title(slide, "Conviction & Th\u00e8se d\u2019Investissement",
                f"{ticker}  \u00b7  FinSight IA  \u00b7  Rapport {gen_date}")

    # ── Bloc recommandation central ─────────────────────────────────────────
    _REC_COL = {"BUY": ("1A7A4A", "E8F5E9"), "SELL": ("A82020", "FFEBEE")}.get(rec, ("B06000", "FFF8E1"))
    add_rect(slide, 0.90, 2.30, 7.50, 3.50, _REC_COL[1], line_hex=_REC_COL[0], line_width_pt=2)
    add_text_box(slide, 0.90, 2.40, 7.50, 1.10, rec, 36, _REC_COL[0], bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, 0.90, 3.55, 7.50, 0.55, f"Conviction : {conv_pct}%", 11, GREY_TXT, align=PP_ALIGN.CENTER)
    add_text_box(slide, 0.90, 4.15, 7.50, 0.55, gen_date, 8.5, GREY_TXT, align=PP_ALIGN.CENTER)

    # ── Conviction meter (barre horizontale) ────────────────────────────────
    bar_x, bar_y, bar_w, bar_h = 0.90, 6.10, 7.50, 0.55
    add_rect(slide, bar_x, bar_y, bar_w, bar_h, "E8ECF0")
    add_rect(slide, bar_x, bar_y, bar_w * conv, bar_h, _REC_COL[0])
    add_text_box(slide, bar_x, bar_y - 0.40, bar_w, 0.38, "NIVEAU DE CONVICTION", 7, NAVY, bold=True)

    # ── Devil's Advocate delta ───────────────────────────────────────────────
    devil_conv = _g(devil, "devil_conviction") or 0.5
    delta = conv - (float(devil_conv) if devil_conv else 0.5)
    delta_lbl = f"+{delta*100:.0f}% vs Avocat du Diable" if delta >= 0 else f"{delta*100:.0f}% vs Avocat du Diable"
    delta_col = "1A7A4A" if delta >= 0 else "A82020"
    add_rect(slide, 0.90, 7.00, 7.50, 0.75, "F5F7FA")
    add_text_box(slide, 1.00, 7.05, 7.30, 0.65, f"Delta conviction : {delta_lbl}", 8.5, delta_col, bold=True)

    # ── Thesis Summary ───────────────────────────────────────────────────────
    thesis_raw = _g(synthesis, "thesis", "") or ""
    thesis_parts = [p.strip() for p in thesis_raw.split(" | ") if p.strip()] if thesis_raw else []
    # Fallback si Thèse vide — construit depuis les ratios disponibles
    if not thesis_parts:
        _co = _g(ci, "company_name", ticker) or ticker
        _sector = _g(ci, "sector", "") or ""
        _p1 = f"Recommandation {rec} pour {_co}" + (f", secteur {_sector}" if _sector else "") + "."
        _fund_parts = []
        if ratios:
            _roic = _g(ratios, "roic")
            _roe  = _g(ratios, "roe")
            _nm   = _g(ratios, "net_margin")
            try:
                if _roic is not None: _fund_parts.append(f"ROIC {float(_roic)*100:.1f}%")
            except Exception: pass
            try:
                if _roe is not None: _fund_parts.append(f"ROE {float(_roe)*100:.1f}%")
            except Exception: pass
            try:
                if _nm is not None: _fund_parts.append(f"marge nette {float(_nm)*100:.1f}%")
            except Exception: pass
        _p2 = ("Fondamentaux : " + ", ".join(_fund_parts) + ".") if _fund_parts else None
        thesis_parts = [_p1]
        if _p2:
            thesis_parts.append(_p2)
        thesis_parts.append("Relancer l'analyse pour obtenir la Thèse complète avec Synthèse IA.")
    add_rect(slide, 9.20, 2.30, 15.20, 0.55, NAVY)
    add_text_box(slide, 9.35, 2.35, 15.0, 0.45, "TH\u00c8SE D\u2019INVESTISSEMENT", 8, WHITE, bold=True)
    y_th = 2.95
    for i, part in enumerate(thesis_parts[:3]):
        add_rect(slide, 9.20, y_th, 0.10, 0.80, _REC_COL[0])
        add_text_box(slide, 9.45, y_th + 0.05, 14.85, 0.75, _fit(part, 150), 8.5, GREY_TXT, wrap=True)
        y_th += 0.95

    # ── Catalyseurs / Risques ────────────────────────────────────────────────
    pos_themes = _g(synthesis, "positive_themes") or []
    neg_themes = _g(synthesis, "negative_themes") or []
    y_mid = max(y_th + 0.30, 5.95)

    add_rect(slide, 9.20, y_mid, 7.40, 0.50, "1A7A4A")
    add_text_box(slide, 9.35, y_mid + 0.05, 7.2, 0.40, "CATALYSEURS BULLS", 7.5, WHITE, bold=True)
    y_c = y_mid + 0.60
    for th in (pos_themes[:2] if pos_themes else ["N/D"]):
        _c_txt = _fit(str(th), 75)
        add_text_box(slide, 9.45, y_c, 7.10, 1.20, f"\u2022 {_c_txt}", 7, GREY_TXT, wrap=True)
        y_c += 1.25

    add_rect(slide, 17.00, y_mid, 6.80, 0.50, "A82020")
    add_text_box(slide, 17.10, y_mid + 0.05, 6.6, 0.40, "RISQUES BEARS", 7.5, WHITE, bold=True)
    y_r = y_mid + 0.60
    for th in (neg_themes[:2] if neg_themes else ["N/D"]):
        _r_txt = _fit(str(th), 75)
        add_text_box(slide, 17.10, y_r, 6.6, 1.20, f"\u2022 {_r_txt}", 7, GREY_TXT, wrap=True)
        y_r += 1.25

    # ── Invalidation conditions ──────────────────────────────────────────────
    inv_list = _g(synthesis, "invalidation_list") or []
    if inv_list:
        y_inv = max(y_c, y_r) + 0.30
        if y_inv < 11.5:
            _inv_h = min(1.55, 12.90 - y_inv)
            add_rect(slide, 9.20, y_inv, 15.20, _inv_h, "FFF3CD")
            add_rect(slide, 9.20, y_inv, 0.10, _inv_h, "B06000")
            inv_str = "  \u00b7  ".join(
                f"{_g(it,'axis','?')}: {str(_g(it,'condition',''))[:90]}" for it in inv_list[:2])
            add_text_box(slide, 9.45, y_inv + 0.05, 14.75, _inv_h - 0.10,
                         f"\u26a0 Conditions d\u2019invalidation : {inv_str}", 7.5, "7A5000", wrap=True)

    # ── Note methodologique ──────────────────────────────────────────────────
    add_text_box(slide, 0.90, 12.95, 23.60, 0.35,
                 "FinSight IA Généré cette analyse a un instant T. Pour le suivi dans le temps, "
                 "comparer les rapports successifs. La conviction Reflète les Données disponibles a la date d'analyse.",
                 7, GREY_TXT, wrap=True)

    return slide


# ---------------------------------------------------------------------------
# Utility
# ---------------------------------------------------------------------------

def _split_text(text: str, n: int) -> list:
    """Divise un texte en n parties. Priorité au séparateur ' | ' (format LLM devil)."""
    if not text:
        return [""] * n
    text = text.strip()
    # Format explicite avec ' | ' (agent_devil)
    if " | " in text:
        parts = [p.strip() for p in text.split(" | ") if p.strip()]
        # Compléter si moins de n parties
        while len(parts) < n:
            parts.append("")
        return parts[:n]
    # Sinon découpage par phrases
    for sep in (". ", "! ", "? "):
        text = text.replace(sep, ".|")
    parts = [p.strip() for p in text.split("|") if p.strip()]
    if not parts:
        return [text[:len(text) // n] for _ in range(n)]
    chunk = max(1, len(parts) // n)
    result = []
    for i in range(n):
        start = i * chunk
        end   = start + chunk if i < n - 1 else len(parts)
        result.append(". ".join(parts[start:end]))
    return result


def _safe_float(v) -> Optional[float]:
    if v is None:
        return None
    try:
        return float(v)
    except Exception:
        return None


# ---------------------------------------------------------------------------
# PPTXWriter — classe principale
# ---------------------------------------------------------------------------

class PPTXWriter:
    """
    Genere un pitchbook IB en 20 slides a partir du FinSightState (dict).
    """

    def generate(self, state: dict, output_path: str) -> str:
        try:
            from pptx import Presentation
            from pptx.util import Cm, Emu
        except ImportError:
            raise RuntimeError("python-pptx requis : pip install python-pptx")

        snap, synthesis, ratios, devil, sentiment = _extract_state(state)

        # Scores additionnels (Composite Distress, M&A, Macro)
        _ticker_for_scores = (snap.company_info.ticker if snap and snap.company_info else None) or state.get('ticker', '')
        _yr_for_scores = None
        if snap and ratios and ratios.years:
            _latest_lbl = sorted(ratios.years.keys())[-1] if ratios.years else None
            _yr_for_scores = ratios.years.get(_latest_lbl)
        _extra_scores = {}
        try:
            from agents.agent_quant import (compute_composite_distress, compute_ma_score,
                                            compute_earnings_quality, compute_capital_structure,
                                            compute_dividend_sustainability)
            if _yr_for_scores is not None:
                _extra_scores['composite_distress']      = compute_composite_distress(_yr_for_scores)
                _extra_scores['ma_score']                = compute_ma_score(_yr_for_scores)
                _extra_scores['earnings_quality']        = compute_earnings_quality(_yr_for_scores)
                _extra_scores['capital_structure']       = compute_capital_structure(_yr_for_scores)
                _extra_scores['dividend_sustainability'] = compute_dividend_sustainability(_yr_for_scores)
        except Exception as _ese:
            log.warning("[PPTXWriter] distress/ma_score: %s", _ese)
        try:
            from agents.agent_macro import AgentMacro
            _extra_scores['macro'] = AgentMacro().analyze()
        except Exception as _ese:
            log.warning("[PPTXWriter] macro: %s", _ese)

        # Fallback minimal si snap absent
        if snap is None:
            log.warning("PPTXWriter: raw_data absent dans le state — fallback vide")
            from data.models import FinancialSnapshot, CompanyInfo, MarketData
            snap = FinancialSnapshot(
                ticker="N/A",
                company_info=CompanyInfo(company_name="N/A", ticker="N/A"),
                years={},
                market=MarketData(),
                stock_history=[],
            )

        prs = Presentation()
        prs.slide_width  = Cm(25.4)
        prs.slide_height = Cm(14.29)

        # --- Slide 1: Cover ---
        _slide_cover(prs, snap, synthesis, ratios, devil, sentiment)

        # --- Slide 2: Executive Summary ---
        _slide_exec_summary(prs, snap, synthesis, ratios, devil, sentiment)

        # --- Slide 3: Sommaire ---
        _slide_sommaire(prs, snap, synthesis)

        # --- Slide 4: Divider Company Overview ---
        divider_slide(prs, "01", "Company Overview",
                      "Pr\u00e9sentation, mod\u00e8le \u00e9conomique & donn\u00e9es de march\u00e9")

        # --- Slide 5: Company Overview ---
        _slide_company_overview(prs, snap, synthesis, ratios)

        # --- Slide 6: Business Model ---
        _slide_business_model(prs, snap, synthesis)

        # --- Slide 7: Divider Analyse Financière ---
        divider_slide(prs, "02", "Analyse Financi\u00e8re",
                      "Compte de r\u00e9sultat, bilan, liquidit\u00e9 & ratios")

        # --- Slide 8: Compte de Résultat ---
        _slide_is(prs, snap, synthesis, ratios)

        # --- Slide 9: Bilan & Liquidité ---
        _slide_bilan(prs, snap, synthesis, ratios)

        # --- Slide 10: Ratios Cles ---
        _slide_ratios(prs, snap, synthesis, ratios)

        # --- Slide 11: Divider Valorisation ---
        divider_slide(prs, "03", "Valorisation",
                      "DCF, comparable peers & Football Field Chart")

        # --- Slide 12: DCF & Scénarios ---
        _slide_dcf(prs, snap, synthesis, ratios)

        # --- Slide 13: Comparable Peers ---
        _slide_peers(prs, snap, synthesis, ratios)

        # --- Slide 14: Football Field ---
        _slide_football_field(prs, snap, synthesis, ratios)

        # --- Slides IB/PE : Multiples historiques + Capital Returns + LBO ---
        try:
            _slide_multiples_historiques(prs, snap, synthesis, ratios)
        except Exception as _e_mh:
            log.error("[PPTXWriter] _slide_multiples_historiques FAILED: %s", _e_mh, exc_info=True)
        try:
            _slide_capital_returns(prs, snap, synthesis, ratios)
        except Exception as _e_cr:
            log.error("[PPTXWriter] _slide_capital_returns FAILED: %s", _e_cr, exc_info=True)
        # ── 3 slides LBO institutionnelles (cadre / returns / stress) ───
        try:
            _lbo_pack = _build_lbo_pack(snap, synthesis, ratios)
            _slide_lbo_cadre(prs, snap, _lbo_pack)
            _slide_lbo_returns(prs, snap, _lbo_pack)
            _slide_lbo_stress(prs, snap, _lbo_pack)
        except Exception as _e_lbo:
            log.error("[PPTXWriter] LBO slides FAILED: %s", _e_lbo, exc_info=True)
            # Fallback : ancienne slide LBO unique
            try:
                _slide_lbo(prs, snap, synthesis, ratios)
            except Exception as _e_lbo2:
                log.error("[PPTXWriter] _slide_lbo fallback FAILED: %s", _e_lbo2, exc_info=True)

        # --- Divider Risques ---
        divider_slide(prs, "04", "Risques & Strat\u00e9gie",
                      "Avocat du diable & conditions d'invalidation")

        # --- Slide 16: Risques ---
        _slide_risques(prs, snap, synthesis, devil, extra_scores=_extra_scores)

        # --- Slide 17: Divider Sentiment ---
        divider_slide(prs, "05", "Sentiment & Annexes",
                      "FinBERT, actionnariat & historique de cours")

        # --- Slide 18: Sentiment ---
        _slide_sentiment(prs, snap, synthesis, sentiment)

        # --- Slide 19: Actionnariat ---
        _slide_actionnariat(prs, snap, synthesis)

        # --- Slide 20: Historique de Cours ---
        _slide_historique(prs, snap, synthesis)

        # --- Slide 21: Conviction Tracker ---
        _slide_conviction_tracker(prs, snap, synthesis, ratios, devil, sentiment)

        # Save
        out_path = Path(output_path)
        out_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out_path))
        log.info("PPTXWriter: pitchbook sauvegarde -> %s", out_path)
        return str(out_path)
