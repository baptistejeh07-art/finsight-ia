"""
PPTX writer PME — 10 slides style restitution client pour CGP/expert-comptable.

Slides :
  S1  : Couverture
  S2  : Identité & dirigeants
  S3  : Chiffres clés (KPI)
  S4  : SIG évolution
  S5  : Rentabilité
  S6  : Solidité
  S7  : Efficacité
  S8  : Benchmark sectoriel
  S9  : Scoring santé + bankabilité
  S10 : Synthèse + disclaimer
"""

from __future__ import annotations

import logging
from pathlib import Path

from core.pappers.analytics import PmeAnalysis, YearAccounts
from core.pappers.benchmark import BenchmarkResult
from core.pappers.bodacc_client import BodaccSummary

log = logging.getLogger(__name__)

NAVY_HEX = 0x1B2A4A
NAVY_LIGHT_HEX = 0x3A5288


def _fmt_eur(v):
    if v is None: return "—"
    abs_v = abs(v)
    if abs_v >= 1e9: return f"{v/1e9:.1f} Md€"
    if abs_v >= 1e6: return f"{v/1e6:.1f} M€"
    if abs_v >= 1e3: return f"{v/1e3:.0f} k€"
    return f"{v:.0f} €"


def _fmt_pct(v, d=1):
    if v is None: return "—"
    return f"{v*100:.{d}f} %"


def _fmt_x(v, d=2):
    if v is None: return "—"
    return f"{v:.{d}f}x"


def write_pme_pptx(
    output_path: str | Path,
    yearly_accounts: list[YearAccounts],
    analysis: PmeAnalysis,
    benchmark: BenchmarkResult,
    bodacc: BodaccSummary | None,
    siren: str,
    denomination: str,
    profile_name: str = "",
    language: str = "fr",
    currency: str = "EUR",
) -> Path:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dgm.color import RGBColor as _
    except ImportError:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
        except ImportError as e:
            raise RuntimeError("python-pptx requis : pip install python-pptx") from e

    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    # Helper i18n
    from core.i18n import t as _i18n_t, ratio_label as _i18n_ratio, sig_label as _i18n_sig, scoring_label as _i18n_score, normalize_language
    _lang = normalize_language(language)
    def _t(key, default=None):
        return _i18n_t(_lang, key, default)
    def _rl(key):
        return _i18n_ratio(key, _lang)
    def _sl(key):
        return _i18n_sig(key, _lang)
    def _sc(key):
        return _i18n_score(key, _lang)

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def _add_title(slide, text, top=Inches(0.3)):
        tb = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.3), Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x1B, 0x2A, 0x4A)

    def _add_kpi_card(slide, left, top, w, h, label, value, color=None):
        # Fond blanc + bordure
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xF5)
        box.line.color.rgb = RGBColor(0xE5, 0xE5, 0xE5)
        box.shadow.inherit = False
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.1)
        # label
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(0x73, 0x73, 0x73)
        # value
        p2 = tf.add_paragraph()
        p2.text = value
        p2.font.size = Pt(18)
        p2.font.bold = True
        if color:
            p2.font.color.rgb = color
        else:
            p2.font.color.rgb = RGBColor(0x17, 0x17, 0x17)

    def _add_bullet(slide, left, top, w, h, items: list[str], size=Pt(14)):
        tb = slide.shapes.add_textbox(left, top, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "• " + item
            p.font.size = size
            p.font.color.rgb = RGBColor(0x40, 0x40, 0x40)

    last_year = max(analysis.ratios_by_year.keys()) if analysis.ratios_by_year else None

    # ─── S1 Couverture ───
    s = prs.slides.add_slide(blank_layout)
    # Background navy
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x1B, 0x2A, 0x4A)
    bg.line.fill.background()
    # Title
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = denomination
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER
    # Subtitle
    tb2 = s.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.7), Inches(0.5))
    p = tb2.text_frame.paragraphs[0]
    p.text = f"Analyse financière — Contrôle de gestion · FinSight IA"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0xC4, 0xD5, 0xE8)
    p.alignment = PP_ALIGN.CENTER
    # Meta
    tb3 = s.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5))
    p = tb3.text_frame.paragraphs[0]
    p.text = f"SIREN {siren} · Profil sectoriel : {profile_name or analysis.profile.name}"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0x8A, 0xAE, 0xDB)
    p.alignment = PP_ALIGN.CENTER

    # ─── S2 Identité ───
    s = prs.slides.add_slide(blank_layout)
    _add_title(s, _t("report.identity_governance"))
    lines = [
        f"Dénomination : {denomination}",
        f"SIREN : {siren}",
        f"Profil sectoriel : {analysis.profile.name}",
    ]
    if bodacc:
        lines.append(f"Annonces BODACC : {bodacc.total_annonces} (procédures collectives : {len(bodacc.procedures_collectives)})")
        lines.append(f"Société radiée : {'Oui' if bodacc.radie else 'Non'}")
    _add_bullet(s, Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.5), lines)

    # ─── S3 KPI ───
    s = prs.slides.add_slide(blank_layout)
    _add_title(s, _t("report.key_figures"))
    if last_year:
        sig = analysis.sig_by_year[last_year]
        ca = next((y.chiffre_affaires for y in yearly_accounts if y.annee == last_year), None)
        kpis = [
            (_sl("chiffre_affaires"), _fmt_eur(ca)),
            (_sl("ebe"), _fmt_eur(sig.ebe)),
            (_sl("resultat_net"), _fmt_eur(sig.resultat_net)),
            (_sl("caf"), _fmt_eur(sig.capacite_autofinancement)),
            (_sc("health_score"), f"{analysis.health_score:.0f}/100" if analysis.health_score else "—"),
            ("Altman Z", f"{analysis.altman_z:.2f}" if analysis.altman_z else "—"),
        ]
        for i, (label, value) in enumerate(kpis):
            col = i % 3
            row = i // 3
            left = Inches(0.8 + col * 4.1)
            top = Inches(1.5 + row * 2.4)
            _add_kpi_card(s, left, top, Inches(3.8), Inches(2.0), label, value)

    # ─── S4 SIG ───
    s = prs.slides.add_slide(blank_layout)
    _add_title(s, _t("report.sig"))
    years = sorted(analysis.sig_by_year.keys())
    if years:
        sig_items = []
        for y in years[-3:]:  # 3 dernières années
            sig = analysis.sig_by_year[y]
            sig_items.append(f"{y} · VA {_fmt_eur(sig.valeur_ajoutee)} · EBE {_fmt_eur(sig.ebe)} · REX {_fmt_eur(sig.resultat_exploitation)} · RN {_fmt_eur(sig.resultat_net)}")
        _add_bullet(s, Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.5), sig_items, Pt(13))

    # ─── S5 Rentabilité ───
    s = prs.slides.add_slide(blank_layout)
    _add_title(s, _t("report.profitability"))
    if last_year:
        r = analysis.ratios_by_year[last_year]
        items = [
            (_rl("marge_ebitda"), _fmt_pct(r.marge_ebitda)),
            (_rl("marge_nette"), _fmt_pct(r.marge_nette)),
            (_rl("roce"), _fmt_pct(r.roce)),
            (_rl("roe"), _fmt_pct(r.roe)),
        ]
        for i, (label, value) in enumerate(items):
            col = i % 2
            row = i // 2
            _add_kpi_card(s, Inches(0.8 + col * 6.2), Inches(1.5 + row * 2.4),
                          Inches(5.9), Inches(2.0), label, value)

    # ─── S6 Solidité ───
    s = prs.slides.add_slide(blank_layout)
    _add_title(s, _t("report.solvency"))
    if last_year:
        r = analysis.ratios_by_year[last_year]
        items = [
            (_rl("dette_nette_ebitda"), _fmt_x(r.dette_nette_ebitda)),
            (_rl("couverture_interets"), _fmt_x(r.couverture_interets)),
            (_rl("autonomie_financiere"), _fmt_pct(r.autonomie_financiere)),
            (_rl("bfr_jours_ca"), f"{r.bfr_jours_ca:.0f} j" if r.bfr_jours_ca is not None else "—"),
        ]
        for i, (label, value) in enumerate(items):
            col = i % 2
            row = i // 2
            _add_kpi_card(s, Inches(0.8 + col * 6.2), Inches(1.5 + row * 2.4),
                          Inches(5.9), Inches(2.0), label, value)

    # ─── S7 Efficacité ───
    s = prs.slides.add_slide(blank_layout)
    _add_title(s, _t("report.efficiency"))
    if last_year:
        r = analysis.ratios_by_year[last_year]
        items = [
            (_rl("dso_jours"), f"{r.dso_jours:.0f} j" if r.dso_jours else "—"),
            (_rl("dpo_jours"), f"{r.dpo_jours:.0f} j" if r.dpo_jours else "—"),
            (_rl("ca_par_employe"), _fmt_eur(r.ca_par_employe)),
            (_rl("charges_perso_ca"), _fmt_pct(r.charges_perso_ca)),
        ]
        for i, (label, value) in enumerate(items):
            col = i % 2
            row = i // 2
            _add_kpi_card(s, Inches(0.8 + col * 6.2), Inches(1.5 + row * 2.4),
                          Inches(5.9), Inches(2.0), label, value)

    # ─── S8 Benchmark ───
    s = prs.slides.add_slide(blank_layout)
    _add_title(s, f"{_t('report.sector_positioning')} ({_t('common.source')} : {benchmark.source})")
    lines = []
    if benchmark.forces:
        lines.append("Forces : " + ", ".join(benchmark.forces))
    if benchmark.faiblesses:
        lines.append("Faiblesses : " + ", ".join(benchmark.faiblesses))
    if not lines:
        lines.append("Positionnement dans la médiane sur la plupart des ratios.")
    _add_bullet(s, Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.5), lines)

    # ─── S9 Scoring ───
    s = prs.slides.add_slide(blank_layout)
    _add_title(s, _t("report.scoring_health"))
    items = [
        ("Altman Z-Score", f"{analysis.altman_z:.2f}" if analysis.altman_z else "—"),
        ("Verdict", analysis.altman_verdict),
        ("Score santé FinSight", f"{analysis.health_score:.0f}/100" if analysis.health_score else "—"),
        ("Score bankabilité", f"{analysis.bankability_score:.0f}/100" if analysis.bankability_score else "—"),
        ("Dette additionnelle accessible (cible 3×EBITDA)", _fmt_eur(analysis.debt_capacity_estimate)),
    ]
    for i, (label, value) in enumerate(items):
        _add_kpi_card(s, Inches(0.8), Inches(1.3 + i * 0.9), Inches(11.7), Inches(0.8),
                      label, value)

    # ─── S10 Synthèse ───
    s = prs.slides.add_slide(blank_layout)
    _add_title(s, _t("report.synthesis_disclaimer"))
    lines = [
        f"Cette analyse porte sur {denomination} (SIREN {siren}).",
        f"Profil sectoriel appliqué : {analysis.profile.name}.",
        f"Source comptable : données Pappers (liasse fiscale) ou saisie manuelle.",
        "Avertissement MiFID II : cette analyse ne constitue pas un conseil en investissement, fiscal ou juridique.",
        "Les ratios, benchmarks et scores sont fournis à titre indicatif et ne se substituent pas à une expertise comptable certifiée.",
    ]
    _add_bullet(s, Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.5), lines, Pt(12))

    prs.save(output_path)
    return output_path
